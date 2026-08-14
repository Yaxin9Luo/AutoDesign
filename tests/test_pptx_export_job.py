from __future__ import annotations

import hashlib
import json
import os
from pathlib import Path
from types import SimpleNamespace
import sys
import tempfile
import threading
import time
import unittest
from unittest import mock

from autodesign.process_supervision import ProcessLedger, process_identity, process_is_alive
from autodesign.run_control import CancellationToken, RunCancelled


def _digest_tree(root: Path) -> str:
    digest = hashlib.sha256()
    for path in sorted(item for item in root.rglob("*") if item.is_file()):
        digest.update(str(path.relative_to(root)).encode())
        digest.update(path.read_bytes())
    return digest.hexdigest()


def _settings(command: str) -> SimpleNamespace:
    return SimpleNamespace(
        code_editor_cmd=command,
        code_editor_harness="custom",
        code_editor_model=None,
        code_editor_timeout_s=10,
        code_editor_max_attempts=1,
        harness_api_key=None,
    )


class PptxExportJobTests(unittest.TestCase):
    def test_happy_path_keeps_attempts_and_output_inside_derived_run(self) -> None:
        from autodesign.pptx_export_job import run_pptx_export_job

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source_dir = root / "runs" / "source" / "final"
            source_dir.mkdir(parents=True)
            source = source_dir / "poster.html"
            source.write_text('<main data-w="1600" data-h="900">Poster</main>', encoding="utf-8")
            (source_dir / "figure.png").write_bytes(b"figure")
            source_digest = _digest_tree(source_dir)
            run_dir = root / "runs" / "pptx-derived"
            agent = root / "pptx_agent.py"
            agent.write_text(
                """
from pathlib import Path
import os
import sys
from pptx import Presentation
from pptx.util import Inches
sys.stdin.read()
Path('env-secret.txt').write_text(os.environ.get('AUTODESIGN_UNRELATED_SECRET', '<absent>'))
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])
for index in range(3):
    slide.shapes.add_textbox(Inches(index), Inches(index), Inches(2), Inches(1))
prs.save('export.pptx')
Path('export_done.json').write_text('{\"ok\": true}')
""".strip(),
                encoding="utf-8",
            )

            with mock.patch.dict(
                "os.environ",
                {"AUTODESIGN_UNRELATED_SECRET": "must-not-leak"},
                clear=False,
            ):
                result = run_pptx_export_job(
                    run_id="pptx-derived",
                    run_dir=run_dir,
                    source_html=source,
                    artifact={"artifact_type": "landing", "name": "My Artifact", "canvas": {}},
                    settings=_settings(f"{sys.executable} {agent}"),
                    cancellation_token=CancellationToken.never("pptx-derived"),
                )

            output = Path(result["pptx_path"])
            self.assertTrue(output.is_file())
            self.assertTrue(output.is_relative_to(run_dir))
            self.assertTrue(Path(result["attempts"][0]["attempt_dir"]).is_relative_to(run_dir))
            self.assertEqual(_digest_tree(source_dir), source_digest)
            self.assertFalse((source_dir / "exports").exists())
            self.assertEqual(result["canvas"], {"w": 1600, "h": 900})
            self.assertEqual(
                set(result),
                {"run_id", "pptx_path", "manifest_path", "attempts", "canvas"},
            )
            self.assertEqual(
                (run_dir / "pptx-export" / "attempt_01" / "env-secret.txt").read_text(),
                "<absent>",
            )

    def test_staging_rejects_file_and_directory_symlinks(self) -> None:
        from autodesign.pptx_export_job import PptxExportJobError, run_pptx_export_job

        for link_kind in ("file", "directory"):
            with self.subTest(link_kind=link_kind), tempfile.TemporaryDirectory() as tmp:
                root = Path(tmp)
                source_dir = root / "runs" / "source" / "final"
                source_dir.mkdir(parents=True)
                source = source_dir / "poster.html"
                source.write_text('<main data-w="1600" data-h="900">Poster</main>', encoding="utf-8")
                external = root / ("external.txt" if link_kind == "file" else "external-dir")
                if link_kind == "file":
                    external.write_text("outside", encoding="utf-8")
                else:
                    external.mkdir()
                    (external / "outside.txt").write_text("outside", encoding="utf-8")
                link = source_dir / ("linked.txt" if link_kind == "file" else "linked-dir")
                try:
                    link.symlink_to(external, target_is_directory=link_kind == "directory")
                except (OSError, NotImplementedError):
                    self.skipTest("symlinks unavailable")

                with self.assertRaises(PptxExportJobError) as raised:
                    run_pptx_export_job(
                        run_id=f"pptx-link-{link_kind}",
                        run_dir=root / "runs" / f"pptx-link-{link_kind}",
                        source_html=source,
                        artifact={"artifact_type": "landing", "name": "Unsafe", "canvas": {}},
                        settings=_settings(f"{sys.executable} -c pass"),
                        cancellation_token=CancellationToken.never(f"pptx-link-{link_kind}"),
                    )

                self.assertEqual(raised.exception.code, "unsafe_source_tree")

    def test_agent_symlink_outputs_are_never_validated_or_promoted(self) -> None:
        from pptx import Presentation
        from pptx.util import Inches

        from autodesign.pptx_export_job import PptxExportJobError, run_pptx_export_job

        for linked_name in ("export.pptx", "export_done.json"):
            with self.subTest(linked_name=linked_name), tempfile.TemporaryDirectory() as tmp:
                root = Path(tmp)
                source_dir = root / "runs" / "source" / "final"
                source_dir.mkdir(parents=True)
                source = source_dir / "poster.html"
                source.write_text('<main data-w="1600" data-h="900">Poster</main>', encoding="utf-8")
                external_pptx = root / "outside.pptx"
                presentation = Presentation()
                slide = presentation.slides.add_slide(presentation.slide_layouts[6])
                for index in range(3):
                    slide.shapes.add_textbox(Inches(index), Inches(index), Inches(2), Inches(1))
                presentation.save(external_pptx)
                external_done = root / "outside-done.json"
                external_done.write_text("{}", encoding="utf-8")
                agent = root / "symlink_agent.py"
                agent.write_text(
                    f"""
from pathlib import Path
import shutil, sys
sys.stdin.read()
if {linked_name!r} == 'export.pptx':
    Path('export.pptx').symlink_to({str(external_pptx)!r})
    Path('export_done.json').write_text('{{}}')
else:
    shutil.copy2({str(external_pptx)!r}, 'export.pptx')
    Path('export_done.json').symlink_to({str(external_done)!r})
""".strip(),
                    encoding="utf-8",
                )
                run_dir = root / "runs" / f"pptx-output-link-{linked_name.replace('.', '-')}"

                with self.assertRaises(PptxExportJobError) as raised:
                    run_pptx_export_job(
                        run_id=run_dir.name,
                        run_dir=run_dir,
                        source_html=source,
                        artifact={"artifact_type": "landing", "name": "Unsafe output", "canvas": {}},
                        settings=_settings(f"{sys.executable} {agent}"),
                        cancellation_token=CancellationToken.never(run_dir.name),
                    )

                self.assertEqual(raised.exception.code, "pptx_agent_export_failed")
                self.assertFalse(any((run_dir / "exports").glob("*.pptx")))

    def test_cancelled_agent_kills_descendants_and_never_publishes_pptx(self) -> None:
        from autodesign.pptx_export_job import run_pptx_export_job

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source_dir = root / "runs" / "source" / "final"
            source_dir.mkdir(parents=True)
            source = source_dir / "poster.html"
            source.write_text('<main data-w="1600" data-h="900">Poster</main>', encoding="utf-8")
            source_digest = _digest_tree(source_dir)
            run_dir = root / "runs" / "pptx-cancel"
            agent = root / "slow_pptx_agent.py"
            agent.write_text(
                """
from pathlib import Path
import subprocess, sys, time
sys.stdin.read()
child = subprocess.Popen([sys.executable, '-c',
    \"import pathlib,time; time.sleep(1.2); pathlib.Path('descendant-finished').write_text('late')\"])
Path('child.pid').write_text(str(child.pid))
time.sleep(1.4)
Path('export.pptx').write_bytes(b'late-pptx')
Path('export_done.json').write_text('{}')
""".strip(),
                encoding="utf-8",
            )
            cancelled = threading.Event()
            token = CancellationToken(store=None, run_id="pptx-cancel", signal_event=cancelled)
            outcome: list[BaseException | dict[str, object]] = []

            def execute() -> None:
                try:
                    outcome.append(
                        run_pptx_export_job(
                            run_id="pptx-cancel",
                            run_dir=run_dir,
                            source_html=source,
                            artifact={"artifact_type": "landing", "name": "Cancelled", "canvas": {}},
                            settings=_settings(f"{sys.executable} {agent}"),
                            cancellation_token=token,
                        )
                    )
                except BaseException as exc:
                    outcome.append(exc)

            worker = threading.Thread(target=execute)
            worker.start()
            child_pid_path = run_dir / "pptx-export" / "attempt_01" / "child.pid"
            deadline = time.monotonic() + 5
            while time.monotonic() < deadline and not child_pid_path.exists():
                time.sleep(0.02)
            self.assertTrue(child_pid_path.exists(), "agent fixture did not spawn its descendant")
            child_pid = int(child_pid_path.read_text())
            cancelled.set()
            worker.join(timeout=8)

            self.assertFalse(worker.is_alive())
            self.assertEqual(len(outcome), 1)
            self.assertIsInstance(outcome[0], RunCancelled)
            for record in ProcessLedger(run_dir).read().processes:
                self.assertFalse(process_is_alive(record.identity), record)
            with self.assertRaises(ProcessLookupError):
                process_identity(child_pid)
            self.assertFalse(any((run_dir / "exports").glob("*.pptx")))
            self.assertEqual(_digest_tree(source_dir), source_digest)
            cancelled_snapshot = _digest_tree(run_dir)
            time.sleep(1.5)
            self.assertFalse(any(run_dir.rglob("descendant-finished")))
            self.assertFalse(any((run_dir / "exports").glob("*.pptx")))
            self.assertEqual(_digest_tree(run_dir), cancelled_snapshot)

    def test_pptx_cancel_during_promotion_leaves_no_final_or_manifest(self) -> None:
        from autodesign.pptx_export_job import run_pptx_export_job

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source_dir = root / "runs" / "source" / "final"
            source_dir.mkdir(parents=True)
            source = source_dir / "poster.html"
            source.write_text('<main data-w="1600" data-h="900">Poster</main>', encoding="utf-8")
            agent = root / "pptx_agent.py"
            agent.write_text(
                """
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
import sys
sys.stdin.read()
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])
for index in range(3):
    slide.shapes.add_textbox(Inches(index), Inches(index), Inches(2), Inches(1))
prs.save('export.pptx')
Path('export_done.json').write_text('{}')
""".strip(),
                encoding="utf-8",
            )
            run_dir = root / "runs" / "pptx-cancel-promotion"
            output = run_dir / "exports" / "My-Artifact.pptx"
            output.parent.mkdir(parents=True)
            output.write_bytes(b"existing-final")
            cancellation = threading.Event()
            token = CancellationToken(store=None, run_id="pptx-cancel-promotion", signal_event=cancellation)
            real_replace = os.replace

            def cancel_after_replace(source: object, destination: object) -> None:
                real_replace(source, destination)
                if Path(destination) == output:
                    cancellation.set()

            with mock.patch("autodesign.pptx_export_job.os.replace", side_effect=cancel_after_replace):
                with self.assertRaises(RunCancelled):
                    run_pptx_export_job(
                        run_id="pptx-cancel-promotion",
                        run_dir=run_dir,
                        source_html=source,
                        artifact={"artifact_type": "landing", "name": "My Artifact", "canvas": {}},
                        settings=_settings(f"{sys.executable} {agent}"),
                        cancellation_token=token,
                    )

            self.assertEqual(output.read_bytes(), b"existing-final")
            self.assertFalse(output.with_suffix(".agent-export.json").exists())

    def test_pptx_cancel_before_manifest_restores_complete_previous_delivery(self) -> None:
        from autodesign.pptx_export_job import run_pptx_export_job

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source_dir = root / "runs" / "source" / "final"
            source_dir.mkdir(parents=True)
            source = source_dir / "poster.html"
            source.write_text('<main data-w="1600" data-h="900">Poster</main>', encoding="utf-8")
            agent = root / "pptx_agent.py"
            agent.write_text(
                """
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
import sys
sys.stdin.read()
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])
for index in range(3):
    slide.shapes.add_textbox(Inches(index), Inches(index), Inches(2), Inches(1))
prs.save('export.pptx')
Path('export_done.json').write_text('{}')
""".strip(),
                encoding="utf-8",
            )
            real_replace = os.replace

            for previous_delivery in (False, True):
                for cancellation_boundary in ("pptx_replace", "manifest_replace"):
                    with self.subTest(
                        previous_delivery=previous_delivery,
                        cancellation_boundary=cancellation_boundary,
                    ):
                        suffix = f"{int(previous_delivery)}-{cancellation_boundary}"
                        run_dir = root / "runs" / f"pptx-delivery-{suffix}"
                        output = run_dir / "exports" / "My-Artifact.pptx"
                        manifest = output.with_suffix(".agent-export.json")
                        output.parent.mkdir(parents=True)
                        if previous_delivery:
                            output.write_bytes(b"previous-pptx")
                            manifest.write_bytes(b'{"generation":"previous"}\n')
                        cancellation = threading.Event()
                        token = CancellationToken(
                            store=None,
                            run_id=run_dir.name,
                            signal_event=cancellation,
                        )
                        target = output if cancellation_boundary == "pptx_replace" else manifest

                        def cancel_after_boundary(source_path: object, destination_path: object) -> None:
                            real_replace(source_path, destination_path)
                            if Path(destination_path) == target:
                                cancellation.set()

                        with mock.patch(
                            "autodesign.pptx_export_job.os.replace",
                            side_effect=cancel_after_boundary,
                        ):
                            with self.assertRaises(RunCancelled):
                                run_pptx_export_job(
                                    run_id=run_dir.name,
                                    run_dir=run_dir,
                                    source_html=source,
                                    artifact={
                                        "artifact_type": "landing",
                                        "name": "My Artifact",
                                        "canvas": {},
                                    },
                                    settings=_settings(f"{sys.executable} {agent}"),
                                    cancellation_token=token,
                                )

                        if previous_delivery:
                            self.assertEqual(output.read_bytes(), b"previous-pptx")
                            self.assertEqual(
                                manifest.read_bytes(),
                                b'{"generation":"previous"}\n',
                            )
                            self.assertEqual(
                                {path.name for path in output.parent.iterdir()},
                                {output.name, manifest.name},
                            )
                        else:
                            self.assertFalse(output.exists())
                            self.assertFalse(manifest.exists())
                            self.assertEqual(list(output.parent.iterdir()), [])

    def test_pptx_cancel_at_complete_rolls_back_delivery(self) -> None:
        from autodesign.pptx_export_job import run_pptx_export_job

        class CancelAtComplete(CancellationToken):
            def raise_if_cancelled(self, phase: str) -> None:
                if phase == "pptx_export.complete":
                    raise RunCancelled(self.run_id, phase)

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            source_dir = root / "runs" / "source" / "final"
            source_dir.mkdir(parents=True)
            source = source_dir / "poster.html"
            source.write_text('<main data-w="1600" data-h="900">Poster</main>', encoding="utf-8")
            agent = root / "pptx_agent.py"
            agent.write_text(
                """
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches
import sys
sys.stdin.read()
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])
for index in range(3):
    slide.shapes.add_textbox(Inches(index), Inches(index), Inches(2), Inches(1))
prs.save('export.pptx')
Path('export_done.json').write_text('{}')
""".strip(),
                encoding="utf-8",
            )

            for previous_delivery in (False, True):
                with self.subTest(previous_delivery=previous_delivery):
                    run_dir = root / "runs" / f"pptx-complete-{int(previous_delivery)}"
                    output = run_dir / "exports" / "My-Artifact.pptx"
                    manifest = output.with_suffix(".agent-export.json")
                    output.parent.mkdir(parents=True)
                    if previous_delivery:
                        output.write_bytes(b"previous-pptx")
                        manifest.write_bytes(b'{"generation":"previous"}\n')

                    with self.assertRaises(RunCancelled):
                        run_pptx_export_job(
                            run_id=run_dir.name,
                            run_dir=run_dir,
                            source_html=source,
                            artifact={
                                "artifact_type": "landing",
                                "name": "My Artifact",
                                "canvas": {},
                            },
                            settings=_settings(f"{sys.executable} {agent}"),
                            cancellation_token=CancelAtComplete(
                                store=None,
                                run_id=run_dir.name,
                            ),
                        )

                    if previous_delivery:
                        self.assertEqual(output.read_bytes(), b"previous-pptx")
                        self.assertEqual(
                            manifest.read_bytes(),
                            b'{"generation":"previous"}\n',
                        )
                        self.assertEqual(
                            {path.name for path in output.parent.iterdir()},
                            {output.name, manifest.name},
                        )
                    else:
                        self.assertEqual(list(output.parent.iterdir()), [])

    @unittest.skipUnless(hasattr(os, "fork"), "hard-crash recovery requires fork")
    def test_pptx_hard_crash_recovery_is_idempotent_across_transaction_phases(self) -> None:
        import autodesign.pptx_export_job as job
        from autodesign.run_control import durable_replace_json as real_durable_replace_json

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            real_replace = os.replace

            for previous_state in ("none", "output_only", "manifest_only", "both"):
                for crash_phase in (
                    "after_pptx_replace",
                    "after_manifest_replace",
                    "after_committed_journal",
                ):
                    with self.subTest(
                        previous_state=previous_state,
                        crash_phase=crash_phase,
                    ):
                        export_dir = root / f"delivery-{previous_state}-{crash_phase}"
                        export_dir.mkdir()
                        output = export_dir / "My-Artifact.pptx"
                        manifest = output.with_suffix(".agent-export.json")
                        staged = export_dir / ".My-Artifact.partial.pptx"
                        staged.write_bytes(b"new-pptx")
                        if previous_state in {"output_only", "both"}:
                            output.write_bytes(b"previous-pptx")
                        if previous_state in {"manifest_only", "both"}:
                            manifest.write_bytes(b'{"generation":"previous"}\n')

                        child_pid = os.fork()
                        if child_pid == 0:
                            def crash_after_replace(source_path: object, destination_path: object) -> None:
                                real_replace(source_path, destination_path)
                                destination = Path(destination_path)
                                if crash_phase == "after_pptx_replace" and destination == output:
                                    os._exit(91)
                                if crash_phase == "after_manifest_replace" and destination == manifest:
                                    os._exit(92)

                            def crash_after_journal(path: Path, payload: object) -> Path:
                                result = real_durable_replace_json(path, payload)
                                if (
                                    crash_phase == "after_committed_journal"
                                    and isinstance(payload, dict)
                                    and payload.get("state") == "committed"
                                ):
                                    os._exit(93)
                                return result

                            try:
                                with mock.patch(
                                    "autodesign.pptx_export_job.os.replace",
                                    side_effect=crash_after_replace,
                                ), mock.patch(
                                    "autodesign.pptx_export_job.durable_replace_json",
                                    side_effect=crash_after_journal,
                                    create=True,
                                ):
                                    job._publish_pptx_delivery(
                                        staged_pptx=staged,
                                        output=output,
                                        manifest_path=manifest,
                                        manifest={
                                            "run_id": "new-run",
                                            "generation": "new",
                                        },
                                        token=CancellationToken.never("new-run"),
                                    )
                            except BaseException:
                                os._exit(94)
                            os._exit(95)

                        _, status = os.waitpid(child_pid, 0)
                        self.assertTrue(os.WIFEXITED(status), status)
                        expected_exit = {
                            "after_pptx_replace": 91,
                            "after_manifest_replace": 92,
                            "after_committed_journal": 93,
                        }[crash_phase]
                        self.assertEqual(os.WEXITSTATUS(status), expected_exit)

                        recovered = job.recover_pptx_delivery_transactions(export_dir)
                        self.assertEqual(len(recovered), 1)
                        self.assertEqual(job.recover_pptx_delivery_transactions(export_dir), ())

                        committed = crash_phase == "after_committed_journal"
                        if committed:
                            self.assertEqual(output.read_bytes(), b"new-pptx")
                            self.assertEqual(
                                json.loads(manifest.read_text(encoding="utf-8"))["generation"],
                                "new",
                            )
                        else:
                            if previous_state in {"output_only", "both"}:
                                self.assertEqual(output.read_bytes(), b"previous-pptx")
                            else:
                                self.assertFalse(output.exists())
                            if previous_state in {"manifest_only", "both"}:
                                self.assertEqual(
                                    manifest.read_bytes(),
                                    b'{"generation":"previous"}\n',
                                )
                            else:
                                self.assertFalse(manifest.exists())

                        expected_names = {output.name, manifest.name} if committed else set()
                        if not committed and previous_state in {"output_only", "both"}:
                            expected_names.add(output.name)
                        if not committed and previous_state in {"manifest_only", "both"}:
                            expected_names.add(manifest.name)
                        self.assertEqual(
                            {path.name for path in export_dir.iterdir()},
                            expected_names,
                        )

    @unittest.skipUnless(hasattr(os, "fork"), "hard-crash recovery requires fork")
    def test_pptx_recovery_cleanup_is_restartable_after_each_unlink(self) -> None:
        import autodesign.pptx_export_job as job
        from autodesign.run_control import durable_replace_json as real_durable_replace_json

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            cleanup_names = (
                "staged_pptx_name",
                "staged_manifest_name",
                "output_backup_name",
                "manifest_backup_name",
                "journal",
            )

            for recovery_state in ("prepared", "committed"):
                crash_points = cleanup_names
                if recovery_state == "prepared":
                    crash_points = ("rolled_back_state", *crash_points)
                for crash_point in crash_points:
                    with self.subTest(
                        recovery_state=recovery_state,
                        crash_point=crash_point,
                    ):
                        export_dir = root / f"{recovery_state}-{crash_point}"
                        export_dir.mkdir()
                        output = export_dir / "My-Artifact.pptx"
                        manifest = output.with_suffix(".agent-export.json")
                        staged_pptx = export_dir / ".My-Artifact.partial.pptx"
                        staged_manifest = export_dir / ".My-Artifact.agent-export.partial"
                        output_backup = export_dir / ".My-Artifact.pptx.tx.rollback"
                        manifest_backup = export_dir / ".My-Artifact.agent-export.json.tx.rollback"
                        journal_path = export_dir / ".My-Artifact.delivery-transaction.json"
                        output.write_bytes(
                            b"new-pptx" if recovery_state == "committed" else b"interrupted-pptx"
                        )
                        manifest.write_bytes(
                            b'{"generation":"new"}\n'
                            if recovery_state == "committed"
                            else b'{"generation":"interrupted"}\n'
                        )
                        staged_pptx.write_bytes(b"staged-pptx")
                        staged_manifest.write_bytes(b'{"generation":"staged"}\n')
                        output_backup.write_bytes(b"previous-pptx")
                        manifest_backup.write_bytes(b'{"generation":"previous"}\n')
                        payload = {
                            "version": 1,
                            "state": recovery_state,
                            "output_name": output.name,
                            "manifest_name": manifest.name,
                            "staged_pptx_name": staged_pptx.name,
                            "staged_manifest_name": staged_manifest.name,
                            "output_backup_name": output_backup.name,
                            "manifest_backup_name": manifest_backup.name,
                            "had_output": True,
                            "had_manifest": True,
                        }
                        real_durable_replace_json(journal_path, payload)
                        crash_target = {
                            "staged_pptx_name": staged_pptx,
                            "staged_manifest_name": staged_manifest,
                            "output_backup_name": output_backup,
                            "manifest_backup_name": manifest_backup,
                            "journal": journal_path,
                        }.get(crash_point)

                        child_pid = os.fork()
                        if child_pid == 0:
                            real_unlink = Path.unlink

                            def crash_after_rollback_state(path: Path, data: object) -> Path:
                                result = real_durable_replace_json(path, data)
                                if (
                                    crash_point == "rolled_back_state"
                                    and isinstance(data, dict)
                                    and data.get("state") == "rolled_back"
                                ):
                                    os._exit(101)
                                return result

                            def crash_after_unlink(
                                path: Path,
                                *args: object,
                                **kwargs: object,
                            ) -> None:
                                real_unlink(path, *args, **kwargs)
                                if crash_target is not None and Path(path) == crash_target:
                                    os._exit(102)

                            try:
                                with mock.patch(
                                    "autodesign.pptx_export_job.durable_replace_json",
                                    side_effect=crash_after_rollback_state,
                                ), mock.patch.object(
                                    Path,
                                    "unlink",
                                    autospec=True,
                                    side_effect=crash_after_unlink,
                                ):
                                    job.recover_pptx_delivery_transactions(export_dir)
                            except BaseException:
                                os._exit(103)
                            os._exit(104)

                        _, status = os.waitpid(child_pid, 0)
                        self.assertTrue(os.WIFEXITED(status), status)
                        self.assertEqual(
                            os.WEXITSTATUS(status),
                            101 if crash_point == "rolled_back_state" else 102,
                        )

                        job.recover_pptx_delivery_transactions(export_dir)
                        self.assertEqual(job.recover_pptx_delivery_transactions(export_dir), ())
                        self.assertEqual(job.recover_pptx_delivery_transactions(export_dir), ())
                        expected_pptx = (
                            b"new-pptx" if recovery_state == "committed" else b"previous-pptx"
                        )
                        expected_manifest = (
                            b'{"generation":"new"}\n'
                            if recovery_state == "committed"
                            else b'{"generation":"previous"}\n'
                        )
                        self.assertEqual(output.read_bytes(), expected_pptx)
                        self.assertEqual(manifest.read_bytes(), expected_manifest)
                        self.assertEqual(
                            {path.name for path in export_dir.iterdir()},
                            {output.name, manifest.name},
                        )


if __name__ == "__main__":
    unittest.main()
