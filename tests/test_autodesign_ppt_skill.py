from __future__ import annotations

import base64
import hashlib
import importlib.util
import json
import os
import re
import shutil
import subprocess
import sys
import tempfile
import threading
import time
import unittest
import zipfile
from pathlib import Path
from types import SimpleNamespace
from unittest import mock


REPO_ROOT = Path(__file__).resolve().parents[1]
SKILL_ROOT = REPO_ROOT / "agent_skills" / "autodesign-ppt"
HARNESS_PATH = SKILL_ROOT / "scripts" / "ppt_harness.py"
EXPORTER_PATH = SKILL_ROOT / "scripts" / "export_pptx.py"
SETUP_PATH = SKILL_ROOT / "scripts" / "setup_ppt.py"
PPT_LOCK_PATH = SKILL_ROOT / "scripts" / "requirements-ppt.lock"
SAFE_NAVIGATION_SCRIPT = "(()=>{const s=[...document.querySelectorAll('.deck-slide')];const i=()=>Math.max(0,s.findIndex(x=>'#'+x.id===location.hash));const g=n=>{const x=s[Math.min(s.length-1,Math.max(0,n))];if(x){location.hash=x.id;x.scrollIntoView({block:'start'})}};addEventListener('keydown',e=>{if(e.key==='ArrowLeft'){e.preventDefault();g(i()-1)}else if(e.key==='ArrowRight'){e.preventDefault();g(i()+1)}});addEventListener('hashchange',()=>{const x=s[i()];if(x)x.scrollIntoView({block:'start'})})})();"


def _load_script(name: str, path: Path):
    if not path.is_file():
        return None
    spec = importlib.util.spec_from_file_location(name, path)
    assert spec is not None and spec.loader is not None
    module = importlib.util.module_from_spec(spec)
    sys.modules[name] = module
    spec.loader.exec_module(module)
    return module


def _png_bytes() -> bytes:
    return base64.b64decode(
        "iVBORw0KGgoAAAANSUhEUgAAAAIAAAACCAIAAAD91JpzAAAAFElEQVR4nGP4z8DA"
        "wMDAxMDAwAIAAwEBAAHiIw0AAAAASUVORK5CYII="
    )


def _deck_html(slide_count: int = 18, *, remote_asset: bool = False) -> str:
    slides: list[str] = []
    for index in range(1, slide_count + 1):
        slide_id = f"slide-{index:02d}"
        image_src = "https://example.com/figure.png" if remote_asset and index == 2 else "assets/pixel.png"
        content = [
            (
                f'<h1 data-pptx-kind="text" data-pptx-x="120" data-pptx-y="90" '
                f'data-pptx-w="1680" data-pptx-h="120" data-font-size="54" '
                f'data-claim-id="claim-{index:02d}" data-source-ids="ev-001">'
                "The paper presents a grounded research finding."
                "</h1>"
            ),
            (
                '<p data-pptx-kind="text" data-pptx-x="120" data-pptx-y="260" '
                'data-pptx-w="760" data-pptx-h="180" data-font-size="28" '
                'data-source-ids="ev-001">A concise source-backed explanation.</p>'
            ),
            (
                '<div data-pptx-kind="shape" data-shape="rect" data-pptx-x="120" '
                'data-pptx-y="500" data-pptx-w="760" data-pptx-h="12" '
                'data-fill="#6B3FA0"></div>'
            ),
        ]
        if index == 2:
            content.append(
                f'<img data-pptx-kind="image" data-pptx-x="1000" data-pptx-y="260" '
                f'data-pptx-w="700" data-pptx-h="500" data-source-ids="visual-explicit-001" '
                f'src="{image_src}" alt="Source figure">'
            )
        if index == 3:
            content.append(
                '<table data-pptx-kind="table" data-pptx-x="980" data-pptx-y="250" '
                'data-pptx-w="740" data-pptx-h="420" data-source-ids="ev-001">'
                '<tr><th>Method</th><th>Result</th></tr>'
                '<tr><td>AutoDesign</td><td>Grounded</td></tr>'
                '</table>'
            )
        slides.append(
            f'<section class="deck-slide" id="{slide_id}" data-slide-id="{slide_id}" '
            f'data-slide-index="{index}" data-slide-role="evidence" data-section="paper-talk" '
            f'data-assertion-title="The paper presents a grounded research finding." '
            f'data-source-ids="ev-001" data-speaker-notes="[Sources] ev-001 [Talk] Explain slide {index}." '
            'data-width="1920" data-height="1080" data-background="#F7F7F3">'
            + "".join(content)
            + "</section>"
        )
    return "".join(
        [
            "<!doctype html><html><head><meta charset=\"utf-8\"><style>",
            "html,body{margin:0;background:#efefec}.deck-slide{position:relative;width:1920px;height:1080px;overflow:hidden;background:#F7F7F3}",
            "@media print{.deck-slide{display:block!important;break-after:page;page-break-after:always}}",
            "</style></head><body>",
            f'<main id="deck" data-autodesign-artifact-root="deck" data-slide-count="{slide_count}" data-width="1920" data-height="1080">',
            *slides,
            "</main>",
            f"<script data-autodesign-navigation>{SAFE_NAVIGATION_SCRIPT}</script>",
            "</body></html>",
        ]
    )


def _deck_html_for_plan(plan: dict[str, object]) -> str:
    html = _deck_html(int(plan["slide_count"]))
    for slide in plan["slides"]:
        index = int(slide["slide_index"])
        html = html.replace(
            f'data-slide-index="{index}" data-slide-role="evidence" '
            'data-section="paper-talk" '
            'data-assertion-title="The paper presents a grounded research finding."',
            f'data-slide-index="{index}" data-slide-role="{slide["role"]}" '
            f'data-section="{slide["chapter"]}" '
            f'data-assertion-title="{slide["assertion_title"]}"',
            1,
        ).replace(
            "The paper presents a grounded research finding.</h1>",
            f'{slide["assertion_title"]}</h1>',
            1,
        ).replace(
            f"[Sources] ev-001 [Talk] Explain slide {index}.",
            str(slide["speaker_note_intent"]),
            1,
        )
    return html


class AutoDesignPptSkillTests(unittest.TestCase):
    def setUp(self) -> None:
        self.temporary = tempfile.TemporaryDirectory()
        self.root = Path(self.temporary.name)
        self.harness = _load_script("portable_ppt_harness", HARNESS_PATH)
        self.exporter = _load_script("portable_ppt_exporter", EXPORTER_PATH)
        self.setup = _load_script("portable_ppt_setup", SETUP_PATH)

    def tearDown(self) -> None:
        self.temporary.cleanup()

    def _require(self, module: object | None, path: Path) -> object:
        self.assertIsNotNone(module, f"missing standalone PPT implementation: {path}")
        return module

    def _write_fixture(self, slide_count: int = 18, *, remote_asset: bool = False) -> Path:
        artifact = self.root / "artifact"
        (artifact / "assets").mkdir(parents=True)
        (artifact / "assets" / "pixel.png").write_bytes(_png_bytes())
        html = artifact / "deck.html"
        html.write_text(_deck_html(slide_count, remote_asset=remote_asset), encoding="utf-8")
        return html

    def _story_plan(
        self,
        harness: object,
        *,
        slide_count: int = 18,
        evidence_ref: str = "ev-001",
    ) -> dict[str, object]:
        return {
            "format_version": 1,
            "slides": [
                {
                    "slide_id": f"slide-{index:02d}",
                    "role": role,
                    "evidence_refs": [evidence_ref],
                }
                for index, (role, _title, _job) in enumerate(
                    harness._arc_for_count(slide_count),
                    start=1,
                )
            ],
        }

    def _role_rich_evidence(self) -> str:
        return (
            "The paper title, authors, affiliation, and thesis identify the work. "
            "A roadmap and outline give an overview. The problem is an unresolved "
            "challenge whose motivation and significance explain its importance. "
            "Related work leaves a prior-work gap. The contributions introduce three "
            "advances. The method overview presents a framework and architecture; its "
            "core mechanism coordinates each module through an algorithmic procedure. "
            "A loss equation defines the objective. Experiments use a dataset, benchmark, "
            "and evaluation metric. Performance, accuracy, improvement, and score describe "
            "the primary results. Robustness and generalization hold across conditions. "
            "An ablation removes each variant. Qualitative case-study examples visualize "
            "representative behavior. Limitations and failure modes define uncertainty. "
            "The implications connect application to practice. A takeaway summary states "
            "the key finding, and the conclusion closes with discussion and future work. "
            "Additional evidence and findings support each figure and table."
        )

    def _conditioned_story_fixture(
        self,
    ) -> tuple[dict[str, str], list[str], dict[str, str]]:
        role_refs = {
            "cover": "ev-001",
            "outline": "ev-002",
            "problem": "ev-003",
            "motivation": "ev-004",
            "prior-gap": "ev-005",
            "contributions": "ev-006",
            "method-overview": "ev-007",
            "mechanism": "ev-008",
            "objective": "ev-009",
            "setup": "ev-010",
            "primary-results": "ev-011",
            "results-deep-dive": "ev-012",
            "evidence-analysis": "ev-013",
            "qualitative": "ev-014",
            "limitations": "ev-015",
            "implications": "ev-016",
            "takeaways": "ev-017",
            "closing": "ev-018",
        }
        evidence_texts = {
            "ev-001": "The paper title, authors, affiliation, and thesis identify the work.",
            "ev-002": "The roadmap gives an overview of the talk.",
            "ev-003": "The research problem is an unresolved design challenge.",
            "ev-004": "The motivation and significance explain why the work is important.",
            "ev-005": "Related work leaves a clear prior-work gap.",
            "ev-006": "The paper contributions introduce three advances.",
            "ev-007": "The method overview presents the framework and architecture.",
            "ev-008": "The core mechanism coordinates modules through an algorithmic procedure.",
            "ev-009": "The objective is defined by a loss equation.",
            "ev-010": "The experiment uses a dataset, benchmark, and evaluation metric.",
            "ev-011": "Primary performance improves accuracy and the reported score.",
            "ev-012": "A secondary result breakdown gives additional per-category findings.",
            "ev-013": "Error analysis interprets the observed trend and supporting evidence.",
            "ev-014": "Qualitative case-study examples visualize representative behavior.",
            "ev-015": "Limitations and failure modes define uncertainty in the current scope.",
            "ev-016": "The implications connect the application to research practice.",
            "ev-017": "The takeaway summary states the key finding.",
            "ev-018": "The conclusion and discussion close with future work.",
        }
        roles = list(role_refs)
        return evidence_texts, roles, role_refs

    def _initialize_run(self, brief: str = "Create a conference deck.") -> Path:
        harness = self._require(self.harness, HARNESS_PATH)
        paper = self.root / "paper.md"
        paper.write_text(
            "# Grounded paper\n\n" + self._role_rich_evidence() + "\n",
            encoding="utf-8",
        )
        run = self.root / "run"
        harness._command_init(
            SimpleNamespace(
                run_dir=run,
                source=paper,
                extra_asset=[],
                reference_image=[],
                archive_sha256=None,
            )
        )
        slide_count = harness._explicit_slide_count(brief) or 18
        story_plan = self.root / "story-plan.json"
        story_plan.write_text(
            json.dumps(
                self._story_plan(harness, slide_count=slide_count),
                ensure_ascii=False,
            ),
            encoding="utf-8",
        )
        harness._command_plan(
            SimpleNamespace(
                run_dir=run,
                brief=brief,
                slide_count=None,
                story_plan=story_plan,
                visual_allocations=None,
            )
        )
        return run

    def test_planner_defaults_paper_decks_to_exactly_18_slides(self) -> None:
        harness = self._require(self.harness, HARNESS_PATH)
        plan = harness.build_deck_plan(
            "Create a conference deck from this paper.",
            ["ev-001"],
            evidence_texts={"ev-001": self._role_rich_evidence()},
            story_plan=self._story_plan(harness),
        )
        self.assertEqual(plan["slide_count"], 18)
        self.assertEqual(plan["count_source"], "academic_default")
        self.assertEqual(len(plan["slides"]), 18)
        self.assertEqual([item["slide_id"] for item in plan["slides"]], [f"slide-{i:02d}" for i in range(1, 19)])
        self.assertEqual(
            [item["role"] for item in plan["slides"]],
            [item[0] for item in harness._ACADEMIC_ARC],
        )

    def test_explicit_user_slide_count_overrides_the_18_slide_default(self) -> None:
        harness = self._require(self.harness, HARNESS_PATH)
        for brief, expected in (
            ("Make 12 slides.", 12),
            ("请做12张幻灯片。", 12),
            ("请生成12张PPT。", 12),
            ("请生成 22 页 PPT。", 22),
            ("我想要 22 页的 PPT。", 22),
            ("Create a 15-page deck", 15),
        ):
            with self.subTest(brief=brief):
                plan = harness.build_deck_plan(
                    brief,
                    ["ev-001"],
                    evidence_texts={"ev-001": self._role_rich_evidence()},
                    story_plan=self._story_plan(harness, slide_count=expected),
                )
                self.assertEqual(plan["slide_count"], expected)
                self.assertEqual(plan["count_source"], "explicit_user")
                self.assertEqual(len(plan["slides"]), expected)

    def test_chinese_numeral_count_requires_a_real_deck_target(self) -> None:
        harness = self._require(self.harness, HARNESS_PATH)
        for brief, expected in (
            ("请做十二页幻灯片。", 12),
            ("请制作二十五张PPT。", 25),
            ("请生成六十页演示文稿。", 60),
            ("请做三十页。", 30),
            ("Make it 14 pages.", 14),
        ):
            with self.subTest(brief=brief):
                chinese = harness.build_deck_plan(
                    brief,
                    ["ev-001"],
                    evidence_texts={"ev-001": self._role_rich_evidence()},
                    story_plan=self._story_plan(harness, slide_count=expected),
                )
                self.assertEqual(chinese["slide_count"], expected)
                self.assertEqual(chinese["count_source"], "explicit_user")

        for brief in (
            "Turn this 12-page paper into a conference deck.",
            "Create a presentation from this 12-page manuscript.",
            "Make slides from the 12-page preprint.",
            "Summarize this 25-page article as slides.",
            "Turn the 30-page PDF into a deck.",
        ):
            with self.subTest(source_metadata=brief):
                source_metadata = harness.build_deck_plan(
                    brief,
                    ["ev-001"],
                    evidence_texts={"ev-001": self._role_rich_evidence()},
                    story_plan=self._story_plan(harness),
                )
                self.assertEqual(source_metadata["slide_count"], 18)
                self.assertEqual(
                    source_metadata["count_source"],
                    "academic_default",
                )

    def test_default_planner_assigns_evidence_by_slide_role_not_extraction_order(self) -> None:
        harness = self._require(self.harness, HARNESS_PATH)
        evidence_texts = {
            "ev-001": "Results show a 9.4 percent improvement over strong baselines.",
            "ev-002": "Limitations include latency, missing modalities, and future work.",
            "ev-003": "Our method uses a planner, refinement loop, and modular architecture.",
            "ev-004": "The problem is that existing systems fail on long-horizon design tasks.",
            "ev-005": "The paper title and authors establish its identity.",
            "ev-006": "The roadmap previews the research argument.",
            "ev-007": "The motivation and significance explain why the work matters.",
            "ev-008": "Related work leaves a clear gap.",
            "ev-009": "The paper's contributions introduce three advances.",
            "ev-010": "The core mechanism coordinates specialized modules.",
            "ev-011": "The loss equation defines the training objective.",
            "ev-012": "Experiments use a dataset, benchmark, and evaluation metric.",
            "ev-013": "Robustness and generalization hold across conditions.",
            "ev-014": "The ablation removes each variant in turn.",
            "ev-015": "Qualitative examples visualize representative cases.",
            "ev-016": "The implications connect the application to practice.",
            "ev-017": "The takeaway summary distills the key finding.",
            "ev-018": "The discussion closes with future work.",
        }
        plan = harness.build_deck_plan(
            "Create a conference deck.",
            list(evidence_texts),
            evidence_texts=evidence_texts,
        )
        refs_by_role = {
            str(slide["role"]): slide["evidence_refs"] for slide in plan["slides"]
        }
        self.assertEqual(refs_by_role["problem"], ["ev-004"])
        self.assertEqual(refs_by_role["method-overview"], ["ev-003"])
        self.assertEqual(refs_by_role["primary-results"], ["ev-001"])
        self.assertEqual(refs_by_role["limitations"], ["ev-002"])
        self.assertEqual(plan["evidence_assignment_source"], "semantic_default")

    def test_default_planner_requests_a_story_plan_instead_of_misassigning_evidence(self) -> None:
        harness = self._require(self.harness, HARNESS_PATH)
        with self.assertRaisesRegex(
            harness.PptHarnessError,
            "could not semantically assign evidence",
        ):
            harness.build_deck_plan(
                "Create a conference deck.",
                ["ev-001", "ev-002"],
                evidence_texts={
                    "ev-001": "Our method uses a modular architecture and planner.",
                    "ev-002": "Results improve performance over the baseline.",
                },
            )

    def test_default_planner_requires_distinctive_ablation_evidence(self) -> None:
        harness = self._require(self.harness, HARNESS_PATH)
        for evidence in (
            [
                {"id": "ev-001", "text": "The method improves the result."},
                {"id": "ev-002", "text": "The result supports the method."},
            ],
            [{"id": "ev-001", "text": ""}],
        ):
            with (
                self.subTest(evidence=evidence),
                self.assertRaisesRegex(
                    harness.PptHarnessError,
                    "role ablation; provide --story-plan",
                ),
            ):
                harness._semantic_evidence_ref(
                    "ablation",
                    "Ablation",
                    "Isolate which components create the gain",
                    evidence,
                )

    def test_conditional_roles_reject_generic_single_token_evidence(self) -> None:
        harness = self._require(self.harness, HARNESS_PATH)
        weak_evidence = {
            "ablation": (
                "The framework operates without external supervision.",
                "The architecture uses a model variant.",
                "We remove a module from the configuration; accuracy is the metric.",
            ),
            "robustness": (
                "The appendix reports the variance definition.",
                "The architecture targets generalization.",
                "Across datasets, accuracy is the evaluation metric.",
            ),
            "qualitative": (
                "For example, the repository includes visualization tooling.",
            ),
        }
        for role, examples in weak_evidence.items():
            for text in examples:
                with (
                    self.subTest(role=role, text=text),
                    self.assertRaisesRegex(
                        harness.PptHarnessError,
                        f"role {role}; provide --story-plan",
                    ),
                ):
                    harness._semantic_evidence_ref(
                        role,
                        role.title(),
                        "Use only distinctive experimental evidence",
                        [{"id": "ev-001", "text": text}],
                    )

    def test_conditional_roles_accept_explicit_or_combined_experimental_evidence(self) -> None:
        harness = self._require(self.harness, HARNESS_PATH)
        supported_evidence = {
            "ablation": (
                "The ablation study removes the router and reports an accuracy drop.",
                "Comparing the model with and without the routing module reduces accuracy by two points.",
            ),
            "robustness": (
                "The robustness evaluation remains stable under distribution shift.",
                "Across datasets, accuracy remains stable under noisy inputs.",
                "我们不仅评估鲁棒性，还报告了扰动下保持稳定的结果。",
                "Robustness results remain stable and motivate future work.",
            ),
            "qualitative": (
                "Qualitative analysis shows representative failure examples.",
                "A representative case study illustrates the model's failure behavior.",
            ),
        }
        for role, examples in supported_evidence.items():
            for text in examples:
                with self.subTest(role=role, text=text):
                    self.assertEqual(
                        harness._semantic_evidence_ref(
                            role,
                            role.title(),
                            "Use only distinctive experimental evidence",
                            [{"id": "ev-001", "text": text}],
                        ),
                        "ev-001",
                    )

    def test_conditional_roles_reject_negated_or_future_work_absence(self) -> None:
        harness = self._require(self.harness, HARNESS_PATH)
        absent_evidence = {
            "ablation": (
                "We do not conduct an ablation study.",
                "Without an ablation study, component effects remain unknown.",
                "Ablation studies are left to future work.",
                "No component removal is performed; accuracy drops from 82% to 79%.",
                "我们没有进行消融实验。",
            ),
            "robustness": (
                "We do not evaluate robustness.",
                "The paper lacks robustness evaluation.",
                "Robustness evaluation is left for future work.",
                "No shifted dataset is evaluated; accuracy is 78.5% versus 79.0%.",
                "我们未评估鲁棒性。",
            ),
            "qualitative": (
                "No qualitative analysis is provided.",
                "Without qualitative analysis, no cases can be interpreted.",
                "Qualitative evidence will be considered in future work.",
                "No representative case is shown.",
                "论文未提供定性分析。",
            ),
        }
        for role, examples in absent_evidence.items():
            for text in examples:
                with (
                    self.subTest(role=role, text=text),
                    self.assertRaisesRegex(
                        harness.PptHarnessError,
                        f"role {role}; provide --story-plan",
                    ),
                ):
                    harness._semantic_evidence_ref(
                        role,
                        role.title(),
                        "Do not treat stated absence as evidence",
                        [{"id": "ev-001", "text": text}],
                    )

    def test_conditional_roles_accept_a_complete_positive_clause_beside_absence(self) -> None:
        harness = self._require(self.harness, HARNESS_PATH)
        mixed_evidence = {
            "ablation": (
                "The ablation study removes the router and accuracy drops by 2 points. "
                "Additional ablations are left to future work."
            ),
            "robustness": (
                "Under distribution shift, accuracy remains stable. "
                "Additional robustness evaluation is left to future work."
            ),
            "qualitative": (
                "A representative case study shows the failure behavior. "
                "Additional qualitative analysis is future work."
            ),
        }
        for role, text in mixed_evidence.items():
            with self.subTest(role=role):
                self.assertEqual(
                    harness._semantic_evidence_ref(
                        role,
                        role.title(),
                        "Use the independent positive clause",
                        [{"id": "ev-001", "text": text}],
                    ),
                    "ev-001",
                )

    def test_conditional_roles_do_not_join_signals_across_clauses(self) -> None:
        harness = self._require(self.harness, HARNESS_PATH)
        split_evidence = {
            "ablation": (
                "We remove a routing module from preprocessing. "
                "Accuracy improves for an unrelated full-model baseline."
            ),
            "robustness": (
                "We evaluate a shifted dataset. "
                "Accuracy remains stable on the original benchmark."
            ),
            "qualitative": (
                "We include a representative case study. "
                "The model shows improved quantitative accuracy."
            ),
        }
        for role, text in split_evidence.items():
            with (
                self.subTest(role=role),
                self.assertRaisesRegex(
                    harness.PptHarnessError,
                    f"role {role}; provide --story-plan",
                ),
            ):
                harness._semantic_evidence_ref(
                    role,
                    role.title(),
                    "Do not join evidence across clauses",
                    [{"id": "ev-001", "text": text}],
                )

    def test_conditional_roles_accept_numeric_observed_comparisons(self) -> None:
        harness = self._require(self.harness, HARNESS_PATH)
        measured_evidence = {
            "ablation": (
                "Removing the routing module changes accuracy from 82.1% to 79.4%.",
                "W/o routing module: 79.4%; full model: 82.1%.",
                "W/o routing module accuracy: 79.4%; full model accuracy: 82.1%.",
                "W/o routing module accuracy: 79.4%; full model: 82.1%.",
                (
                    "W/o routing module uses 2 stages and reports accuracy: 79.4%; "
                    "full model uses 4 stages and reports accuracy: 82.1%."
                ),
                (
                    "W/o routing module accuracy for the 2-stage configuration is 79.4%; "
                    "full model accuracy for the 4-stage configuration is 82.1%."
                ),
                (
                    "W/o routing module accuracy is 79.4% across two stages; "
                    "full model accuracy is 82.1% across four stages."
                ),
                "无监督方法的消融结果显示，准确率从82.1%下降到79.4%。",
                "缺乏路由模块时，准确率从82.1%下降到79.4%。",
            ),
            "robustness": (
                "On the shifted dataset, accuracy is 78.5% versus 79.0% in-domain.",
            ),
        }
        for role, examples in measured_evidence.items():
            for text in examples:
                with self.subTest(role=role, text=text):
                    self.assertEqual(
                        harness._semantic_evidence_ref(
                            role,
                            role.title(),
                            "Accept an observed numeric comparison",
                            [{"id": "ev-001", "text": text}],
                        ),
                        "ev-001",
                    )

    def test_conditional_roles_reject_structural_counts_as_labeled_outcomes(self) -> None:
        harness = self._require(self.harness, HARNESS_PATH)
        false_comparisons = (
            (
                "W/o routing module uses 2 stages and 4 attention heads; "
                "full model is the reference."
            ),
            (
                "W/o routing module uses 2 stages at 50% capacity; "
                "full routing module uses 4 stages at 100% capacity."
            ),
            "W/o routing module accuracy: 79.4%; full model latency: 82.1%.",
            "W/o routing module mapping: 2; full routing module mapping: 4.",
            "W/o routing module uses 2 metrics; full routing module uses 4 metrics.",
            (
                "W/o routing module uses 2 stages and reports accuracy; "
                "full model uses 4 stages and reports accuracy."
            ),
            (
                "W/o routing module uses 2 heads and reports accuracy; "
                "full model uses 4 heads and reports accuracy."
            ),
            (
                "W/o routing module uses 2 blocks and reports accuracy; "
                "full model uses 4 blocks and reports accuracy."
            ),
            (
                "W/o routing module has depth 2 and reports accuracy; "
                "full model has depth 4 and reports accuracy."
            ),
            (
                "W/o routing module depth is 2 and reports accuracy; "
                "full model depth is 4 and reports accuracy."
            ),
            (
                "W/o routing module head count was 2 and reports accuracy; "
                "full model head count was 4 and reports accuracy."
            ),
            (
                "W/o routing module block count equals 2 and reports accuracy; "
                "full model block count equals 4 and reports accuracy."
            ),
            (
                "W/o routing module has 2 heads and reports accuracy; "
                "full model has 4 heads and reports accuracy."
            ),
        )
        for text in false_comparisons:
            with (
                self.subTest(text=text),
                self.assertRaisesRegex(
                    harness.PptHarnessError,
                    "role ablation; provide --story-plan",
                ),
            ):
                harness._semantic_evidence_ref(
                    "ablation",
                    "Ablation",
                    "Do not treat architecture counts as measured outcomes",
                    [{"id": "ev-001", "text": text}],
                )

    def test_conditional_roles_reject_chinese_absence_but_not_method_terms(self) -> None:
        harness = self._require(self.harness, HARNESS_PATH)
        absent_evidence = {
            "ablation": "本文缺乏消融实验。",
            "robustness": "本文缺乏鲁棒性评估。",
            "qualitative": "本文无定性分析。",
        }
        for role, text in absent_evidence.items():
            with (
                self.subTest(role=role),
                self.assertRaisesRegex(
                    harness.PptHarnessError,
                    f"role {role}; provide --story-plan",
                ),
            ):
                harness._semantic_evidence_ref(
                    role,
                    role.title(),
                    "Reject a source statement that declares evidence absent",
                    [{"id": "ev-001", "text": text}],
                )

    def test_default_planner_uses_supported_numeric_experimental_roles(self) -> None:
        harness = self._require(self.harness, HARNESS_PATH)
        evidence_texts, roles, _role_refs = self._conditioned_story_fixture()
        evidence_texts["ev-012"] = (
            "On the shifted dataset, accuracy is 78.5% versus 79.0% in-domain."
        )
        roles[11] = "robustness"
        roles[12] = "ablation"
        comparisons = (
            (
                "W/o routing module accuracy for the 2-stage configuration is 79.4%; "
                "full model accuracy for the 4-stage configuration is 82.1%."
            ),
            (
                "W/o routing module accuracy is 79.4% across two stages; "
                "full model accuracy is 82.1% across four stages."
            ),
        )
        for text in comparisons:
            evidence_texts["ev-013"] = text
            with self.subTest(text=text):
                plan = harness.build_deck_plan(
                    "Create a conference deck.",
                    list(evidence_texts),
                    evidence_texts=evidence_texts,
                )
                self.assertEqual(
                    [str(slide["role"]) for slide in plan["slides"]],
                    roles,
                )

    def test_default_planner_rejects_structural_counts_as_ablation_outcomes(self) -> None:
        harness = self._require(self.harness, HARNESS_PATH)
        false_comparisons = (
            (
                "W/o routing module uses 2 stages and 4 attention heads; "
                "full model is the reference."
            ),
            (
                "W/o routing module uses 2 stages and reports accuracy; "
                "full model uses 4 stages and reports accuracy."
            ),
            (
                "W/o routing module uses 2 heads and reports accuracy; "
                "full model uses 4 heads and reports accuracy."
            ),
            (
                "W/o routing module uses 2 blocks and reports accuracy; "
                "full model uses 4 blocks and reports accuracy."
            ),
            (
                "W/o routing module has depth 2 and reports accuracy; "
                "full model has depth 4 and reports accuracy."
            ),
            (
                "W/o routing module depth is 2 and reports accuracy; "
                "full model depth is 4 and reports accuracy."
            ),
            (
                "W/o routing module head count was 2 and reports accuracy; "
                "full model head count was 4 and reports accuracy."
            ),
            (
                "W/o routing module block count equals 2 and reports accuracy; "
                "full model block count equals 4 and reports accuracy."
            ),
            (
                "W/o routing module has 2 heads and reports accuracy; "
                "full model has 4 heads and reports accuracy."
            ),
        )
        for text in false_comparisons:
            evidence_texts, _roles, _role_refs = self._conditioned_story_fixture()
            evidence_texts["ev-013"] = text
            with (
                self.subTest(text=text),
                self.assertRaisesRegex(
                    harness.PptHarnessError,
                    "role ablation; provide --story-plan",
                ),
            ):
                harness.build_deck_plan(
                    "Create a conference deck.",
                    list(evidence_texts),
                    evidence_texts=evidence_texts,
                )

    def test_default_planner_rejects_chinese_absence_statements(self) -> None:
        harness = self._require(self.harness, HARNESS_PATH)
        evidence_texts, _roles, _role_refs = self._conditioned_story_fixture()
        evidence_texts["ev-012"] = "本文缺乏鲁棒性评估。"
        evidence_texts["ev-013"] = "本文缺乏消融实验。"
        evidence_texts["ev-014"] = "本文无定性分析。"
        with self.assertRaisesRegex(
            harness.PptHarnessError,
            "role robustness; provide --story-plan",
        ):
            harness.build_deck_plan(
                "Create a conference deck.",
                list(evidence_texts),
                evidence_texts=evidence_texts,
            )

    def test_default_planner_uses_independent_positive_clauses(self) -> None:
        harness = self._require(self.harness, HARNESS_PATH)
        evidence_texts, roles, _role_refs = self._conditioned_story_fixture()
        evidence_texts["ev-012"] = (
            "Under distribution shift, accuracy remains stable. "
            "Additional robustness evaluation is left to future work."
        )
        evidence_texts["ev-013"] = (
            "The ablation study removes the router and accuracy drops by 2 points. "
            "Additional ablations are left to future work."
        )
        evidence_texts["ev-014"] = (
            "A representative case study shows the failure behavior. "
            "Additional qualitative analysis is future work."
        )
        roles[11] = "robustness"
        roles[12] = "ablation"
        plan = harness.build_deck_plan(
            "Create a conference deck.",
            list(evidence_texts),
            evidence_texts=evidence_texts,
        )
        self.assertEqual(
            [str(slide["role"]) for slide in plan["slides"]],
            roles,
        )

    def test_default_planner_substitutes_unsupported_experimental_roles(self) -> None:
        harness = self._require(self.harness, HARNESS_PATH)
        evidence_texts, expected_roles, role_refs = self._conditioned_story_fixture()
        plan = harness.build_deck_plan(
            "Create a conference deck.",
            list(evidence_texts),
            evidence_texts=evidence_texts,
        )
        roles = [str(slide["role"]) for slide in plan["slides"]]
        self.assertEqual(len(roles), 18)
        self.assertEqual(roles, expected_roles)
        self.assertNotIn("robustness", roles)
        self.assertNotIn("ablation", roles)
        self.assertEqual(len(roles), len(set(roles)))
        self.assertEqual(
            {
                str(slide["role"]): list(slide["evidence_refs"])
                for slide in plan["slides"]
            },
            {role: [source_id] for role, source_id in role_refs.items()},
        )

    def test_host_story_plan_accepts_supported_conditioned_substitutions(self) -> None:
        harness = self._require(self.harness, HARNESS_PATH)
        evidence_texts, roles, role_refs = self._conditioned_story_fixture()
        story_plan = {
            "format_version": 1,
            "slides": [
                {
                    "slide_id": f"slide-{index:02d}",
                    "role": role,
                    "evidence_refs": [role_refs[role]],
                }
                for index, role in enumerate(roles, start=1)
            ],
        }
        plan = harness.build_deck_plan(
            "Create a conference deck.",
            list(evidence_texts),
            evidence_texts=evidence_texts,
            story_plan=story_plan,
        )
        self.assertEqual(plan["evidence_assignment_source"], "host_story_plan")
        self.assertEqual(
            [str(slide["role"]) for slide in plan["slides"]],
            roles,
        )

    def test_host_story_plan_rejects_an_ungrounded_conditional_role(self) -> None:
        harness = self._require(self.harness, HARNESS_PATH)
        for weak_ablation in (
            "The method improves the result.",
            "The architecture uses a model variant.",
            "The framework operates without external supervision.",
        ):
            evidence_texts, roles, role_refs = self._conditioned_story_fixture()
            evidence_texts["ev-019"] = weak_ablation
            roles[12] = "ablation"
            role_refs["ablation"] = "ev-019"
            story_plan = {
                "format_version": 1,
                "slides": [
                    {
                        "slide_id": f"slide-{index:02d}",
                        "role": role,
                        "evidence_refs": [role_refs[role]],
                    }
                    for index, role in enumerate(roles, start=1)
                ],
            }
            with (
                self.subTest(weak_ablation=weak_ablation),
                self.assertRaisesRegex(
                    harness.PptHarnessError,
                    "role ablation.*evidence|evidence.*role ablation",
                ),
            ):
                harness.build_deck_plan(
                    "Create a conference deck.",
                    list(evidence_texts),
                    evidence_texts=evidence_texts,
                    story_plan=story_plan,
                )

    def test_host_story_plan_rejects_negated_conditional_roles(self) -> None:
        harness = self._require(self.harness, HARNESS_PATH)
        cases = (
            (11, "robustness", "The paper lacks robustness evaluation."),
            (12, "ablation", "Without an ablation study, effects remain unknown."),
            (13, "qualitative", "Without qualitative analysis, cases are unavailable."),
        )
        for slot, role, text in cases:
            evidence_texts, roles, role_refs = self._conditioned_story_fixture()
            evidence_texts["ev-019"] = text
            roles[slot] = role
            role_refs[role] = "ev-019"
            with (
                self.subTest(role=role),
                self.assertRaisesRegex(
                    harness.PptHarnessError,
                    f"role {role}.*evidence|evidence.*role {role}",
                ),
            ):
                harness.build_deck_plan(
                    "Create a conference deck.",
                    list(evidence_texts),
                    evidence_texts=evidence_texts,
                    story_plan={
                        "format_version": 1,
                        "slides": [
                            {
                                "slide_id": f"slide-{index:02d}",
                                "role": planned_role,
                                "evidence_refs": [role_refs[planned_role]],
                            }
                            for index, planned_role in enumerate(roles, start=1)
                        ],
                    },
                )

    def test_host_story_plan_accepts_numeric_experimental_comparisons(self) -> None:
        harness = self._require(self.harness, HARNESS_PATH)
        cases = (
            (
                11,
                "robustness",
                "On the shifted dataset, accuracy is 78.5% versus 79.0% in-domain.",
            ),
            (
                12,
                "ablation",
                "W/o routing module accuracy: 79.4%; full model accuracy: 82.1%.",
            ),
            (
                12,
                "ablation",
                (
                    "W/o routing module accuracy for the 2-stage configuration is 79.4%; "
                    "full model accuracy for the 4-stage configuration is 82.1%."
                ),
            ),
            (
                12,
                "ablation",
                (
                    "W/o routing module accuracy is 79.4% across two stages; "
                    "full model accuracy is 82.1% across four stages."
                ),
            ),
        )
        for slot, role, text in cases:
            evidence_texts, roles, role_refs = self._conditioned_story_fixture()
            evidence_texts["ev-019"] = text
            roles[slot] = role
            role_refs[role] = "ev-019"
            plan = harness.build_deck_plan(
                "Create a conference deck.",
                list(evidence_texts),
                evidence_texts=evidence_texts,
                story_plan={
                    "format_version": 1,
                    "slides": [
                        {
                            "slide_id": f"slide-{index:02d}",
                            "role": planned_role,
                            "evidence_refs": [role_refs[planned_role]],
                        }
                        for index, planned_role in enumerate(roles, start=1)
                    ],
                },
            )
            self.assertEqual(str(plan["slides"][slot]["role"]), role)

    def test_host_story_plan_rejects_structural_counts_as_ablation_outcomes(self) -> None:
        harness = self._require(self.harness, HARNESS_PATH)
        false_comparisons = (
            (
                "W/o routing module uses 2 stages and 4 attention heads; "
                "full model is the reference."
            ),
            (
                "W/o routing module uses 2 stages and reports accuracy; "
                "full model uses 4 stages and reports accuracy."
            ),
            (
                "W/o routing module uses 2 heads and reports accuracy; "
                "full model uses 4 heads and reports accuracy."
            ),
            (
                "W/o routing module uses 2 blocks and reports accuracy; "
                "full model uses 4 blocks and reports accuracy."
            ),
            (
                "W/o routing module has depth 2 and reports accuracy; "
                "full model has depth 4 and reports accuracy."
            ),
            (
                "W/o routing module depth is 2 and reports accuracy; "
                "full model depth is 4 and reports accuracy."
            ),
            (
                "W/o routing module head count was 2 and reports accuracy; "
                "full model head count was 4 and reports accuracy."
            ),
            (
                "W/o routing module block count equals 2 and reports accuracy; "
                "full model block count equals 4 and reports accuracy."
            ),
            (
                "W/o routing module has 2 heads and reports accuracy; "
                "full model has 4 heads and reports accuracy."
            ),
        )
        for text in false_comparisons:
            evidence_texts, roles, role_refs = self._conditioned_story_fixture()
            evidence_texts["ev-019"] = text
            roles[12] = "ablation"
            role_refs["ablation"] = "ev-019"
            with (
                self.subTest(text=text),
                self.assertRaisesRegex(
                    harness.PptHarnessError,
                    "role ablation.*evidence|evidence.*role ablation",
                ),
            ):
                harness.build_deck_plan(
                    "Create a conference deck.",
                    list(evidence_texts),
                    evidence_texts=evidence_texts,
                    story_plan={
                        "format_version": 1,
                        "slides": [
                            {
                                "slide_id": f"slide-{index:02d}",
                                "role": planned_role,
                                "evidence_refs": [role_refs[planned_role]],
                            }
                            for index, planned_role in enumerate(roles, start=1)
                        ],
                    },
                )

    def test_host_story_plan_rejects_chinese_absence_statements(self) -> None:
        harness = self._require(self.harness, HARNESS_PATH)
        cases = (
            (11, "robustness", "本文缺乏鲁棒性评估。"),
            (12, "ablation", "本文缺乏消融实验。"),
            (13, "qualitative", "本文无定性分析。"),
        )
        for slot, role, text in cases:
            evidence_texts, roles, role_refs = self._conditioned_story_fixture()
            evidence_texts["ev-019"] = text
            roles[slot] = role
            role_refs[role] = "ev-019"
            with (
                self.subTest(role=role),
                self.assertRaisesRegex(
                    harness.PptHarnessError,
                    f"role {role}.*evidence|evidence.*role {role}",
                ),
            ):
                harness.build_deck_plan(
                    "Create a conference deck.",
                    list(evidence_texts),
                    evidence_texts=evidence_texts,
                    story_plan={
                        "format_version": 1,
                        "slides": [
                            {
                                "slide_id": f"slide-{index:02d}",
                                "role": planned_role,
                                "evidence_refs": [role_refs[planned_role]],
                            }
                            for index, planned_role in enumerate(roles, start=1)
                        ],
                    },
                )

    def test_host_story_plan_uses_independent_positive_clauses(self) -> None:
        harness = self._require(self.harness, HARNESS_PATH)
        cases = (
            (
                11,
                "robustness",
                "Under distribution shift, accuracy remains stable. "
                "Additional robustness evaluation is left to future work.",
            ),
            (
                12,
                "ablation",
                "The ablation study removes the router and accuracy drops by 2 points. "
                "Additional ablations are left to future work.",
            ),
            (
                13,
                "qualitative",
                "A representative case study shows the failure behavior. "
                "Additional qualitative analysis is future work.",
            ),
        )
        for slot, role, text in cases:
            evidence_texts, roles, role_refs = self._conditioned_story_fixture()
            evidence_texts["ev-019"] = text
            roles[slot] = role
            role_refs[role] = "ev-019"
            plan = harness.build_deck_plan(
                "Create a conference deck.",
                list(evidence_texts),
                evidence_texts=evidence_texts,
                story_plan={
                    "format_version": 1,
                    "slides": [
                        {
                            "slide_id": f"slide-{index:02d}",
                            "role": planned_role,
                            "evidence_refs": [role_refs[planned_role]],
                        }
                        for index, planned_role in enumerate(roles, start=1)
                    ],
                },
            )
            self.assertEqual(str(plan["slides"][slot]["role"]), role)

    def test_host_story_plan_does_not_join_signals_across_clauses(self) -> None:
        harness = self._require(self.harness, HARNESS_PATH)
        cases = (
            (
                11,
                "robustness",
                "We evaluate a shifted dataset. Accuracy remains stable on the original benchmark.",
            ),
            (
                12,
                "ablation",
                "We remove a routing module. Accuracy improves for an unrelated baseline.",
            ),
            (
                13,
                "qualitative",
                "We include a representative case study. Quantitative accuracy improves.",
            ),
        )
        for slot, role, text in cases:
            evidence_texts, roles, role_refs = self._conditioned_story_fixture()
            evidence_texts["ev-019"] = text
            roles[slot] = role
            role_refs[role] = "ev-019"
            with (
                self.subTest(role=role),
                self.assertRaisesRegex(
                    harness.PptHarnessError,
                    f"role {role}.*evidence|evidence.*role {role}",
                ),
            ):
                harness.build_deck_plan(
                    "Create a conference deck.",
                    list(evidence_texts),
                    evidence_texts=evidence_texts,
                    story_plan={
                        "format_version": 1,
                        "slides": [
                            {
                                "slide_id": f"slide-{index:02d}",
                                "role": planned_role,
                                "evidence_refs": [role_refs[planned_role]],
                            }
                            for index, planned_role in enumerate(roles, start=1)
                        ],
                    },
                )

    def test_host_story_plan_rejects_wrong_slot_and_duplicate_substitutions(self) -> None:
        harness = self._require(self.harness, HARNESS_PATH)
        evidence_texts, roles, role_refs = self._conditioned_story_fixture()

        wrong_slot = list(roles)
        wrong_slot[11] = "method-detail"
        evidence_texts["ev-019"] = "A source-backed method detail explains one component."
        role_refs["method-detail"] = "ev-019"

        duplicated = list(roles)
        duplicated[12] = "results-deep-dive"

        for label, candidate_roles, expected_error in (
            ("wrong phase", wrong_slot, "not a supported substitution"),
            ("duplicated role", duplicated, "duplicated"),
        ):
            with (
                self.subTest(label=label),
                self.assertRaisesRegex(harness.PptHarnessError, expected_error),
            ):
                harness.build_deck_plan(
                    "Create a conference deck.",
                    list(evidence_texts),
                    evidence_texts=evidence_texts,
                    story_plan={
                        "format_version": 1,
                        "slides": [
                            {
                                "slide_id": f"slide-{index:02d}",
                                "role": role,
                                "evidence_refs": [role_refs[role]],
                            }
                            for index, role in enumerate(candidate_roles, start=1)
                        ],
                    },
                )

    def test_host_story_plan_is_validated_before_becoming_immutable(self) -> None:
        harness = self._require(self.harness, HARNESS_PATH)
        roles = [
            "cover", "outline", "problem", "motivation", "prior-gap",
            "contributions", "method-overview", "mechanism", "objective", "setup",
            "primary-results", "robustness", "ablation", "qualitative", "limitations",
            "implications", "takeaways", "closing",
        ]
        host_slides = [
            {
                "slide_id": f"slide-{index:02d}",
                "role": role,
                "evidence_refs": ["ev-002" if role == "method-overview" else "ev-001"],
            }
            for index, role in enumerate(roles, start=1)
        ]
        story_plan = {"format_version": 1, "slides": host_slides}
        plan = harness.build_deck_plan(
            "Create a conference deck.",
            ["ev-001", "ev-002"],
            evidence_texts={
                "ev-001": self._role_rich_evidence(),
                "ev-002": "The method uses a planner and iterative architecture.",
            },
            story_plan=story_plan,
        )
        self.assertEqual(plan["evidence_assignment_source"], "host_story_plan")
        self.assertEqual(plan["slides"][6]["evidence_refs"], ["ev-002"])

        invalid = json.loads(json.dumps(story_plan))
        invalid["slides"][6]["evidence_refs"] = ["ev-999"]
        with self.assertRaisesRegex(harness.PptHarnessError, "unknown evidence"):
            harness.build_deck_plan(
                "Create a conference deck.",
                ["ev-001", "ev-002"],
                evidence_texts={
                    "ev-001": self._role_rich_evidence(),
                    "ev-002": "The method uses a planner and iterative architecture.",
                },
                story_plan=invalid,
            )

    def test_cjk_source_anchor_is_sentence_aware_and_display_bounded(self) -> None:
        harness = self._require(self.harness, HARNESS_PATH)
        source = (
            "这是论文的核心方法，它通过规划器和迭代反馈提升设计质量。"
            + "后续细节" * 400
            + " "
            + self._role_rich_evidence()
        )
        plan = harness.build_deck_plan(
            "制作学术演示文稿。",
            ["ev-001"],
            evidence_texts={"ev-001": source},
            story_plan=self._story_plan(harness),
        )
        title = str(plan["slides"][0]["assertion_title"])
        self.assertIn("这是论文的核心方法", title)
        self.assertNotIn("后续细节后续细节后续细节", title)
        self.assertLessEqual(len(title), 100)

    def test_source_anchor_does_not_split_emoji_modifiers_or_flag_pairs(self) -> None:
        harness = self._require(self.harness, HARNESS_PATH)
        modifier_anchor = harness._bounded_source_anchor(
            "甲" * 71 + "👍🏽" + "乙" * 40
        )
        flag_anchor = harness._bounded_source_anchor(
            "甲" * 71 + "🇨🇳" + "乙" * 40
        )
        self.assertIn("👍🏽", modifier_anchor)
        self.assertNotIn("👍…", modifier_anchor)
        self.assertIn("🇨🇳", flag_anchor)
        self.assertNotIn("🇨…", flag_anchor)

    def test_saved_plan_is_hash_bound_and_snapshotted_into_each_attempt(self) -> None:
        harness = self._require(self.harness, HARNESS_PATH)
        run = self._initialize_run()
        begin = harness._command_begin(SimpleNamespace(run_dir=run))
        attempt = begin["attempt_id"]
        snapshot = run / "attempts" / attempt / "artifact" / "provenance" / "plan.json"
        self.assertTrue(snapshot.is_file())
        self.assertEqual(snapshot.read_bytes(), (run / "plan.json").read_bytes())
        binding = json.loads((run / "plan-binding.json").read_text(encoding="utf-8"))
        state = json.loads((run / "run.json").read_text(encoding="utf-8"))
        self.assertEqual(state["ppt_plan_sha256"], binding["plan_sha256"])

        plan = json.loads((run / "plan.json").read_text(encoding="utf-8"))
        plan["slide_count"] = 7
        (run / "plan.json").write_text(json.dumps(plan), encoding="utf-8")
        with self.assertRaisesRegex(Exception, "plan.*(?:hash|binding|changed)"):
            harness._resume(run)

        (run / "plan-binding.json").write_text(
            json.dumps(harness._plan_binding(plan)), encoding="utf-8"
        )
        with self.assertRaisesRegex(Exception, "plan.*(?:hash|binding|changed)"):
            harness._resume(run)

    def test_html_contract_accepts_exact_ids_assertions_sources_notes_canvas_and_navigation(self) -> None:
        exporter = self._require(self.exporter, EXPORTER_PATH)
        html = self._write_fixture()
        report = exporter.validate_deck_html(html, expected_slide_count=18)
        self.assertTrue(report["passed"], report)
        self.assertEqual(report["slide_ids"], [f"slide-{i:02d}" for i in range(1, 19)])
        self.assertEqual(report["assertion_title_count"], 18)
        self.assertEqual(report["speaker_note_count"], 18)
        self.assertEqual(report["canvas"], {"width": 1920, "height": 1080})
        self.assertTrue(report["keyboard_navigation"])

    def test_html_contract_rejects_remote_assets_and_wrong_slide_count(self) -> None:
        exporter = self._require(self.exporter, EXPORTER_PATH)
        html = self._write_fixture(17, remote_asset=True)
        report = exporter.validate_deck_html(html, expected_slide_count=18)
        self.assertFalse(report["passed"])
        codes = {item["code"] for item in report["issues"]}
        self.assertIn("slide_count_mismatch", codes)
        self.assertIn("remote_asset", codes)

    def test_html_contract_rejects_remote_links_and_css_generated_text(self) -> None:
        exporter = self._require(self.exporter, EXPORTER_PATH)
        html = self._write_fixture()
        (html.parent / "theme").write_text(
            '.deck-slide::after{content:"Raster-only conclusion"}'
            '@import url("https://example.com/remote.css");',
            encoding="utf-8",
        )
        text = html.read_text(encoding="utf-8")
        text = text.replace(
            "The paper presents a grounded research finding.</h1>",
            '<a href="https://example.com/paper">The paper presents a grounded research finding.</a></h1>',
            1,
        ).replace(
            "</head>",
            '<link rel="stylesheet" href="theme"></head>',
            1,
        )
        html.write_text(text, encoding="utf-8")
        report = exporter.validate_deck_html(html, expected_slide_count=18)
        codes = {item["code"] for item in report["issues"]}
        self.assertFalse(report["passed"])
        self.assertIn("remote_asset", codes)
        self.assertIn("css_generated_text", codes)

    def test_html_contract_rejects_symlink_and_hardlink_assets(self) -> None:
        exporter = self._require(self.exporter, EXPORTER_PATH)
        outside = self.root / "outside.png"
        outside.write_bytes(_png_bytes())
        for link_kind in ("symlink", "hardlink"):
            with self.subTest(link_kind=link_kind):
                artifact = self.root / link_kind
                (artifact / "assets").mkdir(parents=True)
                link = artifact / "assets" / "pixel.png"
                if link_kind == "symlink":
                    link.symlink_to(outside)
                else:
                    os.link(outside, link)
                html = artifact / "deck.html"
                html.write_text(_deck_html(), encoding="utf-8")
                report = exporter.validate_deck_html(html, expected_slide_count=18)
                self.assertFalse(report["passed"])
                self.assertIn(
                    "unsafe_local_asset",
                    {item["code"] for item in report["issues"]},
                )

    def test_slide_metadata_notes_and_all_native_text_are_source_mapped(self) -> None:
        exporter = self._require(self.exporter, EXPORTER_PATH)
        harness = self._require(self.harness, HARNESS_PATH)
        html = self._write_fixture()
        text = html.read_text(encoding="utf-8")
        text = text.replace(' data-slide-id="slide-01"', "", 1)
        text = text.replace('data-slide-index="1"', 'data-slide-index="999"', 1)
        text = text.replace(
            "[Sources] ev-001 [Talk] Explain slide 1.",
            "[Sources] ev-999 [Talk] Accuracy is 99.9%.",
            1,
        )
        text = text.replace("<td>Grounded</td>", "<td>Invented 88.8%</td>", 1)
        html.write_text(text, encoding="utf-8")
        deck = exporter.parse_deck_html(html)
        report = exporter.validate_deck_html(html, expected_slide_count=18)
        codes = {item["code"] for item in report["issues"]}
        self.assertFalse(report["passed"])
        self.assertIn("slide_id_contract", codes)
        self.assertIn("slide_index_contract", codes)
        self.assertIn("speaker_note_sources", codes)

        claims = exporter.claims_from_deck(deck)
        claim_text = "\n".join(str(item.get("text", "")) for item in claims)
        self.assertIn("A concise source-backed explanation.", claim_text)
        self.assertIn("Invented 88.8%", claim_text)
        self.assertIn("Accuracy is 99.9%", claim_text)

        plan = harness.build_deck_plan(
            "Create a conference deck.",
            ["ev-001"],
            evidence_texts={"ev-001": self._role_rich_evidence()},
            story_plan=self._story_plan(harness),
        )
        plan_gate = harness.validate_deck_against_plan(deck, plan)
        self.assertFalse(plan_gate["passed"])
        self.assertTrue(plan_gate["issues"])

        conforming = _deck_html_for_plan(plan)
        html.write_text(conforming, encoding="utf-8")
        accepted_gate = harness.validate_deck_against_plan(
            exporter.parse_deck_html(html), plan
        )
        self.assertTrue(accepted_gate["passed"], accepted_gate)

    def test_plan_gate_exactly_binds_assertion_native_claim_sources_and_notes(self) -> None:
        exporter = self._require(self.exporter, EXPORTER_PATH)
        harness = self._require(self.harness, HARNESS_PATH)
        plan = harness.build_deck_plan(
            "Create a conference deck.",
            ["ev-001"],
            evidence_texts={"ev-001": self._role_rich_evidence()},
            story_plan=self._story_plan(harness),
        )
        html = self._write_fixture()
        conforming = _deck_html_for_plan(plan)
        planned_assertion = str(plan["slides"][0]["assertion_title"])

        mutations = {
            "assertion text": conforming.replace(
                f">{planned_assertion}</h1>",
                ">Unplanned opening</h1>",
                1,
            ),
            "visible text sources": conforming.replace(
                'data-source-ids="ev-001">A concise source-backed explanation.',
                'data-source-ids="ev-999">A concise source-backed explanation.',
                1,
            ),
            "table sources": conforming.replace(
                'data-pptx-h="420" data-source-ids="ev-001">',
                'data-pptx-h="420" data-source-ids="ev-999">',
                1,
            ),
            "speaker note intent": conforming.replace(
                str(plan["slides"][0]["speaker_note_intent"]),
                "[Sources] ev-001 [Talk] A different unplanned talk track.",
                1,
            ),
        }
        for label, authored in mutations.items():
            with self.subTest(label=label):
                html.write_text(authored, encoding="utf-8")
                report = harness.validate_deck_against_plan(
                    exporter.parse_deck_html(html), plan
                )
                self.assertFalse(report["passed"], report)

        html.write_text(conforming, encoding="utf-8")
        deck = exporter.parse_deck_html(html)
        self.assertTrue(harness.validate_deck_against_plan(deck, plan)["passed"])
        claims = exporter.claims_from_deck(deck)
        self.assertTrue(all(claim["source_ids"] == ["ev-001"] for claim in claims))

    def test_html_contract_rejects_visible_text_that_would_be_rasterized_in_pptx(self) -> None:
        exporter = self._require(self.exporter, EXPORTER_PATH)
        html = self._write_fixture()
        text = html.read_text(encoding="utf-8")
        text = text.replace(
            "</section>",
            "<aside>This important result was never tagged as editable.</aside></section>",
            1,
        )
        html.write_text(text, encoding="utf-8")
        report = exporter.validate_deck_html(html, expected_slide_count=18)
        self.assertFalse(report["passed"])
        self.assertIn("untagged_visible_text", {item["code"] for item in report["issues"]})

    def test_html_contract_rejects_arbitrary_executable_script(self) -> None:
        exporter = self._require(self.exporter, EXPORTER_PATH)
        html = self._write_fixture()
        html.write_text(
            html.read_text(encoding="utf-8").replace(
                "</body>", "<script>fetch('https://example.com')</script></body>"
            ),
            encoding="utf-8",
        )
        report = exporter.validate_deck_html(html, expected_slide_count=18)
        self.assertFalse(report["passed"])
        self.assertIn("unsafe_script", {item["code"] for item in report["issues"]})

    def test_slide_audit_variants_isolate_every_slide_and_preserve_local_assets(self) -> None:
        harness = self._require(self.harness, HARNESS_PATH)
        html = self._write_fixture()
        (html.parent / "theme.css").write_text(".deck-slide{color:#171717}\n", encoding="utf-8")
        html.write_text(
            html.read_text(encoding="utf-8").replace(
                "</head>", '<link rel="stylesheet" href="theme.css"></head>'
            ),
            encoding="utf-8",
        )
        variant_dir = self.root / "variants"
        variants = harness.create_slide_audit_variants(html, variant_dir, 18)
        self.assertEqual(len(variants), 18)
        for index, path in enumerate(variants, start=1):
            text = path.read_text(encoding="utf-8")
            self.assertIn(f"#slide-{index:02d}", text)
            self.assertTrue((path.parent / "assets" / "pixel.png").is_file())
            self.assertTrue((path.parent / "theme.css").is_file())
            self.assertNotIn("http://", text)
            self.assertNotIn("https://", text)
            audit_css = re.search(
                r'<style data-autodesign-slide-audit>(.*?)</style>',
                text,
                flags=re.DOTALL,
            )
            self.assertIsNotNone(audit_css)
            self.assertNotRegex(
                audit_css.group(1),
                rf"#{re.escape(f'slide-{index:02d}')}[^}}]*(?:width|height):",
            )

    def test_computed_canvas_gate_rejects_wrong_authored_slide_root_size(self) -> None:
        harness = self._require(self.harness, HARNESS_PATH)
        measurements = [
            {
                "slide_id": f"slide-{index:02d}",
                "computed_width": 100 if index == 4 else 1920,
                "computed_height": 100 if index == 4 else 1080,
                "offset_width": 100 if index == 4 else 1920,
                "offset_height": 100 if index == 4 else 1080,
                "rect_width": 100 if index == 4 else 1920,
                "rect_height": 100 if index == 4 else 1080,
            }
            for index in range(1, 19)
        ]
        report = harness.validate_computed_slide_canvases(measurements, 18)
        self.assertFalse(report["passed"], report)
        self.assertIn("slide-04", "\n".join(report["issues"]))

    def test_pptx_export_reopens_with_editable_text_table_image_shape_and_notes(self) -> None:
        exporter = self._require(self.exporter, EXPORTER_PATH)
        html = self._write_fixture()
        output = self.root / "deck.pptx"
        report = exporter.export_deck_to_pptx(html, output)
        self.assertTrue(output.is_file())
        self.assertTrue(report["passed"], report)
        self.assertEqual(report["slide_count"], 18)
        self.assertEqual(report["notes_count"], 18)
        self.assertGreaterEqual(report["editable_text_shapes"], 36)
        self.assertGreaterEqual(report["table_shapes"], 1)
        self.assertGreaterEqual(report["picture_shapes"], 1)
        self.assertGreaterEqual(report["editable_shape_count"], 18)
        self.assertEqual(report["slide_size_inches"], [13.333333, 7.5])
        with zipfile.ZipFile(output) as archive:
            self.assertIn("ppt/slides/slide18.xml", archive.namelist())
            self.assertIn("ppt/notesSlides/notesSlide18.xml", archive.namelist())
            slide_xml = archive.read("ppt/slides/slide3.xml")
            self.assertIn(b"<a:tbl>", slide_xml)
            self.assertIn(b"The paper presents a grounded research finding.", slide_xml)

    def test_pptx_reopen_enforces_native_table_image_text_shape_and_notes_contract(self) -> None:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        exporter = self._require(self.exporter, EXPORTER_PATH)
        html = self._write_fixture()
        deck = exporter.parse_deck_html(html)
        native_contract = exporter.native_object_contract(deck)
        original = self.root / "native-contract.pptx"
        exporter.export_deck_to_pptx(html, original)

        mutations = ("table", "image", "shape", "notes")
        for mutation in mutations:
            with self.subTest(mutation=mutation):
                presentation = Presentation(str(original))
                removed = False
                for slide in presentation.slides:
                    if mutation == "notes" and not removed:
                        slide.notes_slide.notes_text_frame.text = ""
                        removed = True
                        break
                    for shape in list(slide.shapes):
                        name = str(getattr(shape, "name", ""))
                        selected = (
                            mutation == "table" and getattr(shape, "has_table", False)
                        ) or (
                            mutation == "image"
                            and shape.shape_type == MSO_SHAPE_TYPE.PICTURE
                            and not name.startswith("background:")
                        ) or (
                            mutation == "shape" and ":shape:" in name
                        )
                        if selected:
                            shape._element.getparent().remove(shape._element)
                            removed = True
                            break
                    if removed:
                        break
                self.assertTrue(removed)
                mutated = self.root / f"missing-{mutation}.pptx"
                presentation.save(str(mutated))
                report = exporter.inspect_pptx(
                    mutated,
                    expected_slide_count=18,
                    native_contract=native_contract,
                )
                self.assertFalse(report["passed"], report)
                self.assertIn(
                    f"pptx_native_{mutation}",
                    {issue["code"] for issue in report["issues"]},
                )

    def test_pptx_font_sizes_use_the_same_144_pixel_canvas_scale_as_positions(self) -> None:
        from pptx import Presentation

        exporter = self._require(self.exporter, EXPORTER_PATH)
        html = self._write_fixture()
        output = self.root / "scaled-fonts.pptx"
        exporter.export_deck_to_pptx(html, output)
        presentation = Presentation(str(output))
        title = next(
            shape
            for shape in presentation.slides[0].shapes
            if getattr(shape, "has_text_frame", False)
            and "The paper presents" in shape.text_frame.text
        )
        points = title.text_frame.paragraphs[0].runs[0].font.size.pt
        self.assertAlmostEqual(points, 27.0, places=2)

    def test_pptx_validation_rejects_whole_slide_rasterization_without_editable_overlays(self) -> None:
        exporter = self._require(self.exporter, EXPORTER_PATH)
        html = self._write_fixture()
        text = html.read_text(encoding="utf-8")
        text = text.replace('data-pptx-kind="text"', 'data-ignored-kind="text"')
        text = text.replace(
            '<img data-pptx-kind="image" data-pptx-x="1000" data-pptx-y="260" data-pptx-w="700" data-pptx-h="500"',
            '<img data-pptx-kind="image" data-pptx-x="0" data-pptx-y="0" data-pptx-w="1920" data-pptx-h="1080"',
        )
        html.write_text(text, encoding="utf-8")
        report = exporter.validate_deck_html(html, expected_slide_count=18)
        self.assertFalse(report["passed"])
        self.assertIn("missing_editable_text", {item["code"] for item in report["issues"]})

    def test_contact_sheet_html_references_every_slide_preview(self) -> None:
        harness = self._require(self.harness, HARNESS_PATH)
        previews = []
        for index in range(1, 19):
            path = self.root / "previews" / f"slide-{index:02d}.png"
            path.parent.mkdir(exist_ok=True)
            path.write_bytes(_png_bytes())
            previews.append(path)
        sheet = harness.write_contact_sheet_html(previews, self.root / "contact" / "index.html")
        text = sheet.read_text(encoding="utf-8")
        self.assertEqual(text.count("<img "), 18)
        self.assertIn("Slide 18", text)

    def test_interrupted_validation_discards_only_unverified_qa_scratch(self) -> None:
        harness = self._require(self.harness, HARNESS_PATH)
        attempt = self.root / "attempt"
        scratch = attempt / "qa" / "deck"
        scratch.mkdir(parents=True)
        (scratch / "partial.json").write_text("{}\n", encoding="utf-8")
        cleaned = harness.prepare_qa_directory(attempt)
        self.assertEqual(cleaned, scratch.resolve())
        self.assertTrue(cleaned.is_dir())
        self.assertEqual(list(cleaned.iterdir()), [])

        file_attempt = self.root / "file-attempt"
        file_scratch = file_attempt / "qa" / "deck"
        file_scratch.parent.mkdir(parents=True)
        file_scratch.write_text("interrupted\n", encoding="utf-8")
        cleaned_file = harness.prepare_qa_directory(file_attempt)
        self.assertTrue(cleaned_file.is_dir())
        self.assertEqual(list(cleaned_file.iterdir()), [])

    def test_validation_rejects_unlisted_or_hardlinked_delivery_files(self) -> None:
        harness = self._require(self.harness, HARNESS_PATH)
        run = self._initialize_run()
        attempt = harness._command_begin(SimpleNamespace(run_dir=run))["attempt_id"]
        artifact = run / "attempts" / attempt / "artifact"
        (artifact / "assets").mkdir(exist_ok=True)
        (artifact / "assets" / "pixel.png").write_bytes(_png_bytes())
        (artifact / "deck.html").write_text(_deck_html(), encoding="utf-8")
        (artifact / "deck.pdf").write_bytes(b"pdf")
        (artifact / "deck.pptx").write_bytes(b"pptx")
        outside = self.root / "private.txt"
        outside.write_text("must not be delivered\n", encoding="utf-8")
        os.link(outside, artifact / "private.txt")
        args = SimpleNamespace(
            run_dir=run,
            attempt=attempt,
            browser_cache=None,
            ppt_cache=None,
            offline_browser=True,
            offline_ppt=True,
        )
        with (
            mock.patch.object(harness.portable, "write_source_map"),
            mock.patch.object(
                harness,
                "validate_deck_against_plan",
                return_value={"name": "deck_plan", "passed": True, "issues": []},
                create=True,
            ),
            mock.patch.object(
                harness,
                "_run_visual_gate",
                return_value={"name": "visual_provenance", "passed": True, "issues": []},
            ),
            mock.patch.object(
                harness,
                "render_and_validate_deck",
                return_value={"passed": True, "checks": [], "preview_paths": []},
            ),
            mock.patch.object(harness.portable, "record_deterministic_result") as record,
            self.assertRaisesRegex(harness.PptHarnessError, "unexpected artifact|hardlink"),
        ):
            harness._command_validate(args)
        record.assert_not_called()

        (artifact / "private.txt").unlink()
        (artifact / "notes.json").write_text(
            '{"format_version":1,"slides":[]}\n', encoding="utf-8"
        )
        allowed = harness._artifact_delivery_paths(
            artifact, harness.exporter.parse_deck_html(artifact / "deck.html"), require_outputs=True
        )
        self.assertIn("artifact/provenance/plan.json", allowed)
        self.assertEqual(
            set(allowed),
            {
                "artifact/assets/pixel.png",
                "artifact/deck.html",
                "artifact/deck.pdf",
                "artifact/deck.pptx",
                "artifact/notes.json",
                "artifact/provenance/plan.json",
            },
        )

    def test_validate_hash_binds_plan_all_native_claims_and_delivery_into_review(self) -> None:
        harness = self._require(self.harness, HARNESS_PATH)
        run = self._initialize_run()
        attempt = harness._command_begin(SimpleNamespace(run_dir=run))["attempt_id"]
        artifact = run / "attempts" / attempt / "artifact"
        (artifact / "assets").mkdir(exist_ok=True)
        (artifact / "assets" / "pixel.png").write_bytes(_png_bytes())
        plan = json.loads((run / "plan.json").read_text(encoding="utf-8"))
        authored = _deck_html_for_plan(plan).replace(
            "A concise source-backed explanation.",
            "AutoDesign reports a grounded result.",
        )
        (artifact / "deck.html").write_text(authored, encoding="utf-8")

        def fake_render(_html, *, qa_dir, **_kwargs):
            (artifact / "deck.pdf").write_bytes(b"pdf")
            (artifact / "deck.pptx").write_bytes(b"pptx")
            preview = qa_dir / "slide-01.png"
            preview.parent.mkdir(parents=True, exist_ok=True)
            preview.write_bytes(_png_bytes())
            return {
                "passed": True,
                "checks": [],
                "preview_paths": [str(preview)],
            }

        args = SimpleNamespace(
            run_dir=run,
            attempt=attempt,
            browser_cache=None,
            ppt_cache=None,
            offline_browser=True,
            offline_ppt=True,
        )
        with (
            mock.patch.object(
                harness,
                "_run_visual_gate",
                return_value={"name": "visual_provenance", "passed": True, "issues": []},
            ),
            mock.patch.object(harness, "render_and_validate_deck", side_effect=fake_render),
        ):
            result = harness._command_validate(args)
        self.assertTrue(result["passed"], result)

        deterministic = json.loads(
            (run / "attempts" / attempt / "qa" / "deterministic.json").read_text(
                encoding="utf-8"
            )
        )
        self.assertIn("artifact/provenance/plan.json", deterministic["artifact_hashes"])
        source_map = json.loads(
            (run / "attempts" / attempt / "provenance" / "source-map.json").read_text(
                encoding="utf-8"
            )
        )
        self.assertEqual(len(source_map["claims"]), 18 * 3 + 1)
        self.assertIn("slide_table", {claim["claim_type"] for claim in source_map["claims"]})
        self.assertEqual(
            sum(claim["claim_type"] == "speaker_notes" for claim in source_map["claims"]),
            18,
        )
        context = harness._command_review_context(
            SimpleNamespace(run_dir=run, attempt=attempt)
        )
        self.assertEqual(context["artifact_hashes"], deterministic["artifact_hashes"])
        self.assertIn("artifact/provenance/plan.json", context["artifact_hashes"])

    def test_rendered_comparison_rejects_a_blank_powerpoint_render(self) -> None:
        from PIL import Image, ImageDraw

        exporter = self._require(self.exporter, EXPORTER_PATH)
        canonical = self.root / "canonical.png"
        rendered = self.root / "rendered.png"
        image = Image.new("RGB", (1920, 1080), "white")
        draw = ImageDraw.Draw(image)
        draw.rectangle((100, 100, 850, 300), fill="black")
        draw.rectangle((100, 420, 1700, 900), fill="#6B3FA0")
        image.save(canonical)
        Image.new("RGB", (1920, 1080), "white").save(rendered)
        report = exporter.compare_rendered_slides([canonical], [rendered])
        self.assertFalse(report["passed"], report)
        self.assertLess(report["slide_metrics"][0]["edge_recall"], 0.1)

    def test_libreoffice_render_requires_pixel_comparison_when_office_is_available(self) -> None:
        harness = self._require(self.harness, HARNESS_PATH)
        with (
            mock.patch.object(
                harness.exporter,
                "render_pptx_with_libreoffice",
                return_value={"performed": True, "page_count": 18, "pdf": "rendered.pdf"},
            ),
            mock.patch.object(harness.shutil, "which", return_value=None),
        ):
            report = harness._optional_office_comparison(
                self.root / "deck.pptx",
                self.root / "previews",
                self.root / "qa",
                18,
                object(),
            )
        self.assertTrue(report["performed"])
        self.assertFalse(report["comparison_performed"])
        self.assertFalse(report["passed"])

    def test_setup_runtime_is_versioned_outside_the_installed_skill(self) -> None:
        setup = self._require(self.setup, SETUP_PATH)
        spec = setup.runtime_spec(cache_root=self.root / "cache")
        self.assertIn("python-pptx-1.0.2", spec.cache_key)
        self.assertTrue(str(spec.cache_dir).startswith(str(self.root / "cache")))
        self.assertFalse(str(spec.cache_dir).startswith(str(SKILL_ROOT)))
        self.assertEqual(setup.PINNED_PACKAGES["python-pptx"], "1.0.2")
        with self.assertRaisesRegex(setup.PptRuntimeError, "outside the installed Skill"):
            setup.runtime_spec(cache_root=SKILL_ROOT / "generated-cache")

    def test_ppt_runtime_uses_an_artifact_hash_lock_and_require_hashes(self) -> None:
        setup = self._require(self.setup, SETUP_PATH)
        self.assertTrue(PPT_LOCK_PATH.is_file())
        lock_text = PPT_LOCK_PATH.read_text(encoding="utf-8")
        for name, version in setup.PINNED_PACKAGES.items():
            self.assertIn(f"{name}=={version}".lower(), lock_text.lower())
        self.assertGreaterEqual(lock_text.count("--hash=sha256:"), len(setup.PINNED_PACKAGES))
        spec = setup.runtime_spec(cache_root=self.root / "cache")
        self.assertEqual(
            spec.package_lock_sha256,
            hashlib.sha256(PPT_LOCK_PATH.read_bytes()).hexdigest(),
        )

        commands: list[list[str]] = []

        def completed(command, **_kwargs):
            commands.append([str(item) for item in command])
            return subprocess.CompletedProcess(command, 0, "", "")

        with (
            mock.patch.object(setup.subprocess, "run", side_effect=completed),
            mock.patch.object(setup, "_runtime_tree_sha256", return_value="a" * 64, create=True),
        ):
            setup._install(self.root / "staging", spec)
        pip_command = commands[1]
        self.assertIn("--require-hashes", pip_command)
        self.assertIn("--no-deps", pip_command)
        self.assertIn("--no-compile", pip_command)
        self.assertIn("--requirement", pip_command)
        self.assertNotIn("python-pptx==1.0.2", pip_command)

    def test_ppt_runtime_rejects_same_version_cache_content_tampering(self) -> None:
        setup = self._require(self.setup, SETUP_PATH)
        spec = setup.runtime_spec(cache_root=self.root / "cache")
        cache = spec.cache_dir
        python = cache / ("venv/Scripts/python.exe" if os.name == "nt" else "venv/bin/python")
        python.parent.mkdir(parents=True)
        python.write_text("#!/bin/sh\nexit 0\n", encoding="utf-8")
        python.chmod(0o755)
        site = (
            cache / "venv" / "Lib" / "site-packages"
            if os.name == "nt"
            else cache
            / "venv"
            / "lib"
            / f"python{sys.version_info.major}.{sys.version_info.minor}"
            / "site-packages"
        )
        site.mkdir(parents=True)
        module = site / "pptx.py"
        module.write_text("VERSION = 1\n", encoding="utf-8")
        digest = setup._runtime_tree_sha256(cache)
        setup._atomic_json(cache / setup._STATE_FILE, setup._state_payload(spec, digest))
        module.write_text("VERSION = 2\n", encoding="utf-8")
        with (
            mock.patch.object(setup, "_probe"),
            self.assertRaisesRegex(setup.PptRuntimeError, "content hash|tree hash|tampered"),
        ):
            setup.inspect_runtime(cache, spec)

    def test_ppt_runtime_fails_closed_when_bytecode_cache_appears(self) -> None:
        setup = self._require(self.setup, SETUP_PATH)
        spec = setup.runtime_spec(cache_root=self.root / "cache")
        cache = spec.cache_dir
        site = (
            cache / "venv" / "Lib" / "site-packages"
            if os.name == "nt"
            else cache
            / "venv"
            / "lib"
            / f"python{sys.version_info.major}.{sys.version_info.minor}"
            / "site-packages"
        )
        site.mkdir(parents=True)
        (site / "pptx.py").write_text("VERSION = 1\n", encoding="utf-8")
        cache_dir = site / "__pycache__"
        cache_dir.mkdir()
        (cache_dir / "pptx.cpython-test.pyc").write_bytes(b"executable bytecode")
        with self.assertRaisesRegex(setup.PptRuntimeError, "bytecode|__pycache__"):
            setup._runtime_tree_sha256(cache)

    def test_ppt_runtime_lock_prevents_dual_holders_after_liveness_false_negative(self) -> None:
        setup = self._require(self.setup, SETUP_PATH)
        lock = self.root / "runtime.lock"
        first_entered = threading.Event()
        release_first = threading.Event()
        second_entered = threading.Event()
        errors: list[BaseException] = []

        def first() -> None:
            try:
                with setup._runtime_lock(lock, 2.0):
                    first_entered.set()
                    release_first.wait(2.0)
            except BaseException as error:
                errors.append(error)

        def second() -> None:
            try:
                first_entered.wait(2.0)
                with setup._runtime_lock(lock, 2.0):
                    second_entered.set()
            except BaseException as error:
                errors.append(error)

        with mock.patch.object(setup, "_process_alive", return_value=False):
            left = threading.Thread(target=first)
            right = threading.Thread(target=second)
            left.start()
            right.start()
            self.assertTrue(first_entered.wait(1.0))
            time.sleep(0.15)
            self.assertFalse(second_entered.is_set())
            release_first.set()
            left.join(2.0)
            right.join(2.0)
        self.assertEqual(errors, [])
        self.assertTrue(second_entered.is_set())

    def test_ppt_runtime_advisory_lock_serializes_independent_processes(self) -> None:
        lock = self.root / "runtime.lock"
        first_entered = self.root / "first-entered"
        second_entered = self.root / "second-entered"
        release_first = self.root / "release-first"
        release_second = self.root / "release-second"
        release_second.write_text("ready\n", encoding="utf-8")
        helper = self.root / "lock-holder.py"
        helper.write_text(
            """
import importlib.util
import sys
import time
from pathlib import Path

module_path, lock_path, entered_path, release_path = map(Path, sys.argv[1:])
spec = importlib.util.spec_from_file_location("ppt_lock_holder", module_path)
module = importlib.util.module_from_spec(spec)
sys.modules[spec.name] = module
spec.loader.exec_module(module)
module._process_alive = lambda _pid: False
with module._runtime_lock(lock_path, 5.0):
    entered_path.write_text("entered\\n", encoding="utf-8")
    deadline = time.monotonic() + 4.0
    while not release_path.exists():
        if time.monotonic() >= deadline:
            raise SystemExit("timed out waiting for release")
        time.sleep(0.02)
""".strip()
            + "\n",
            encoding="utf-8",
        )
        environment = {**os.environ, "PYTHONDONTWRITEBYTECODE": "1"}
        first = subprocess.Popen(
            [
                sys.executable,
                "-B",
                str(helper),
                str(SETUP_PATH),
                str(lock),
                str(first_entered),
                str(release_first),
            ],
            env=environment,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
        )
        second = None
        try:
            deadline = time.monotonic() + 3.0
            while not first_entered.exists() and time.monotonic() < deadline:
                time.sleep(0.02)
            self.assertTrue(first_entered.exists())
            second = subprocess.Popen(
                [
                    sys.executable,
                    "-B",
                    str(helper),
                    str(SETUP_PATH),
                    str(lock),
                    str(second_entered),
                    str(release_second),
                ],
                env=environment,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True,
            )
            time.sleep(0.25)
            self.assertFalse(second_entered.exists())
            release_first.write_text("release\n", encoding="utf-8")
            first_output = first.communicate(timeout=5)
            second_output = second.communicate(timeout=5)
            self.assertEqual(first.returncode, 0, "".join(first_output))
            self.assertEqual(second.returncode, 0, "".join(second_output))
            self.assertTrue(second_entered.exists())
        finally:
            for process in (first, second):
                if process is not None and process.poll() is None:
                    process.kill()
                    process.communicate()

    def test_cli_help_does_not_write_into_the_installed_skill(self) -> None:
        installed = self.root / "installed" / "autodesign-ppt"
        shutil.copytree(SKILL_ROOT, installed)
        for generated in installed.rglob("__pycache__"):
            shutil.rmtree(generated)
        before = {
            path.relative_to(installed).as_posix(): path.read_bytes()
            for path in installed.rglob("*")
            if path.is_file()
        }
        environment = dict(os.environ)
        environment.pop("PYTHONDONTWRITEBYTECODE", None)
        completed = subprocess.run(
            [sys.executable, str(installed / "scripts" / "ppt_harness.py"), "--help"],
            cwd=self.root,
            env=environment,
            text=True,
            capture_output=True,
            check=False,
        )
        self.assertEqual(completed.returncode, 0, completed.stdout + completed.stderr)
        after = {
            path.relative_to(installed).as_posix(): path.read_bytes()
            for path in installed.rglob("*")
            if path.is_file()
        }
        self.assertEqual(after, before)
        self.assertEqual(list(installed.rglob("__pycache__")), [])

    def test_passing_review_cannot_bypass_the_bound_minimum_scores(self) -> None:
        harness = self._require(self.harness, HARNESS_PATH)
        for invalid_score in (3, float("nan")):
            with self.subTest(invalid_score=invalid_score):
                scores = {
                    dimension: 4 for dimension in harness.REVIEW_RUBRIC["dimensions"]
                }
                scores["source_fidelity"] = invalid_score
                review_path = self.root / "review.json"
                review_path.write_text(
                    json.dumps({"verdict": "pass", "dimension_scores": scores}),
                    encoding="utf-8",
                )
                args = SimpleNamespace(run_dir=self.root, attempt="01", review=review_path)
                with (
                    mock.patch.object(harness, "_resume", return_value={}),
                    mock.patch.object(harness.portable, "record_semantic_review") as record,
                    self.assertRaisesRegex(harness.PptHarnessError, "at least 4"),
                ):
                    harness._command_record_review(args)
                record.assert_not_called()

    def test_resume_rejects_a_persisted_passing_review_below_the_bound_minimum(self) -> None:
        harness = self._require(self.harness, HARNESS_PATH)
        attempt = "01"
        review_path = self.root / "attempts" / attempt / "qa" / "semantic-review.json"
        review_path.parent.mkdir(parents=True)
        review_path.write_text(
            json.dumps(
                {
                    "verdict": "pass",
                    "dimension_scores": {
                        dimension: 1 for dimension in harness.REVIEW_RUBRIC["dimensions"]
                    },
                }
            ),
            encoding="utf-8",
        )
        with (
            mock.patch.object(
                harness.portable,
                "resume_run",
                return_value={"state": "semantic_passed", "active_attempt": attempt},
            ),
            mock.patch.object(harness, "_verify_plan_binding"),
            self.assertRaisesRegex(harness.PptHarnessError, "at least 4|minimum"),
        ):
            harness._resume(self.root)

    def test_resume_rejects_a_final_delivery_file_with_an_external_hardlink(self) -> None:
        harness = self._require(self.harness, HARNESS_PATH)
        delivered = self.root / "final" / "deck.pptx"
        delivered.parent.mkdir(parents=True)
        delivered.write_bytes(b"pptx")
        os.link(delivered, self.root / "outside-copy.pptx")
        with (
            mock.patch.object(
                harness.portable,
                "resume_run",
                return_value={"state": "finalized", "active_attempt": "01"},
            ),
            mock.patch.object(harness, "_verify_plan_binding"),
            self.assertRaisesRegex(harness.PptHarnessError, "hardlink|link count"),
        ):
            harness._resume(self.root)

    def test_final_delivery_exact_set_symlink_and_hardlink_are_rejected(self) -> None:
        harness = self._require(self.harness, HARNESS_PATH)
        final = self.root / "final"
        final.mkdir()
        deck = final / "deck.pptx"
        deck.write_bytes(b"pptx")
        harness.portable.atomic_write_json(
            final / "delivery-manifest.json",
            {
                "format_version": 1,
                "attempt_id": "01",
                "verification_status": "verified",
                "files": {"deck.pptx": harness.portable.sha256_file(deck)},
            },
        )
        self.assertEqual(
            harness.portable._verify_delivery(final)["files"],
            {"deck.pptx": harness.portable.sha256_file(deck)},
        )

        extra = final / "unlisted.txt"
        extra.write_text("extra\n", encoding="utf-8")
        with self.assertRaisesRegex(Exception, "exactly match"):
            harness.portable._verify_delivery(final)
        extra.unlink()

        outside = self.root / "outside.txt"
        outside.write_text("outside\n", encoding="utf-8")
        symlink = final / "linked.txt"
        symlink.symlink_to(outside)
        with self.assertRaisesRegex(Exception, "symlink|non-regular"):
            harness.portable._verify_delivery(final)
        symlink.unlink()

        os.link(deck, self.root / "external-hardlink.pptx")
        with self.assertRaisesRegex(harness.PptHarnessError, "hardlink|link count"):
            harness._verify_final_delivery_links(self.root)

    @unittest.skipUnless(
        os.environ.get("AUTODESIGN_SKILL_REAL_PPT") == "1",
        "set AUTODESIGN_SKILL_REAL_PPT=1 for the full browser/PDF/PPTX smoke",
    )
    def test_real_validate_pipeline_produces_18_previews_pdf_contact_sheet_and_reopen_report(self) -> None:
        harness = self._require(self.harness, HARNESS_PATH)
        html = self._write_fixture()
        result = harness.render_and_validate_deck(
            html,
            expected_slide_count=18,
            qa_dir=self.root / "qa",
            browser_cache=Path(os.environ["AUTODESIGN_SKILL_BROWSER_CACHE"]),
            ppt_cache=self.root / "ppt-cache",
            offline_browser=True,
        )
        self.assertTrue(result["passed"], result)
        self.assertEqual(len(result["preview_paths"]), 19)
        self.assertEqual(result["pdf_page_count"], 18)
        self.assertTrue(Path(result["contact_sheet"]).is_file())
        self.assertTrue(Path(result["pptx_path"]).is_file())
        self.assertEqual(result["pptx_validation"]["slide_count"], 18)
        if shutil.which("soffice"):
            self.assertTrue(result["rendered_pptx_comparison"]["performed"])
            self.assertTrue(result["rendered_pptx_comparison"]["comparison_performed"])
            self.assertTrue(result["rendered_pptx_comparison"]["passed"])
            self.assertEqual(result["rendered_pptx_comparison"]["page_count"], 18)


if __name__ == "__main__":
    unittest.main()
