#!/usr/bin/env python3
"""Validate tagged HTML slides and export native editable PowerPoint files."""

from __future__ import annotations

import argparse
import json
import math
import os
import re
import shutil
import subprocess
import sys
import tempfile
import uuid
import zipfile
from dataclasses import dataclass, field
from html.parser import HTMLParser
from pathlib import Path
from stat import S_ISREG
from typing import Any, Iterable, Mapping, Sequence
from urllib.parse import unquote, urlsplit


SLIDE_WIDTH_PX = 1920
SLIDE_HEIGHT_PX = 1080
SLIDE_WIDTH_IN = 13.333333333333334
SLIDE_HEIGHT_IN = 7.5
PX_PER_INCH = 144.0
_SLIDE_ID = re.compile(r"^slide-(\d{2})$")
_HEX_COLOR = re.compile(r"^#?([0-9a-fA-F]{6})$")
_REMOTE_SCHEMES = {"http", "https", "ws", "wss", "ftp", "file", "javascript"}
_URL_ATTRIBUTES = {
    "action",
    "archive",
    "background",
    "cite",
    "codebase",
    "data",
    "formaction",
    "href",
    "longdesc",
    "manifest",
    "ping",
    "poster",
    "profile",
    "src",
    "usemap",
    "xlink:href",
}
SAFE_NAVIGATION_SCRIPT = "(()=>{const s=[...document.querySelectorAll('.deck-slide')];const i=()=>Math.max(0,s.findIndex(x=>'#'+x.id===location.hash));const g=n=>{const x=s[Math.min(s.length-1,Math.max(0,n))];if(x){location.hash=x.id;x.scrollIntoView({block:'start'})}};addEventListener('keydown',e=>{if(e.key==='ArrowLeft'){e.preventDefault();g(i()-1)}else if(e.key==='ArrowRight'){e.preventDefault();g(i()+1)}});addEventListener('hashchange',()=>{const x=s[i()];if(x)x.scrollIntoView({block:'start'})})})();"


class PptContractError(RuntimeError):
    """The deck cannot be delivered as a truthful editable PowerPoint."""


def _external_output(path: Path | str) -> Path:
    target = Path(path).expanduser().resolve(strict=False)
    package = Path(__file__).resolve().parent.parent
    try:
        target.relative_to(package)
    except ValueError:
        return target
    raise PptContractError("generated output must stay outside the installed Skill")


@dataclass
class ExportElement:
    tag: str
    kind: str
    attrs: dict[str, str]
    text_parts: list[str] = field(default_factory=list)
    rows: list[list[str]] = field(default_factory=list)
    current_row: list[str] | None = None
    current_cell: list[str] | None = None

    @property
    def text(self) -> str:
        return " ".join(" ".join(self.text_parts).split())


@dataclass
class SlideModel:
    slide_id: str
    attrs: dict[str, str]
    elements: list[ExportElement] = field(default_factory=list)
    text_parts: list[str] = field(default_factory=list)


@dataclass
class DeckModel:
    html_path: Path
    root_attrs: dict[str, str]
    slides: list[SlideModel]
    resources: list[tuple[str, str, str]]
    stylesheets: list[str]
    raw_html: str
    root_count: int = 0


class _DeckParser(HTMLParser):
    def __init__(self, html_path: Path, raw_html: str) -> None:
        super().__init__(convert_charrefs=True)
        self.model = DeckModel(html_path, {}, [], [], [], raw_html)
        self._slide: SlideModel | None = None
        self._element: ExportElement | None = None
        self._element_depth = 0

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        values = {name.lower(): value or "" for name, value in attrs}
        lowered = tag.lower()
        if values.get("data-autodesign-artifact-root") == "deck":
            self.model.root_attrs = values
            self.model.root_count += 1
        classes = set(values.get("class", "").split())
        if lowered == "section" and "deck-slide" in classes:
            slide_id = values.get("data-slide-id") or values.get("id") or ""
            self._slide = SlideModel(slide_id, values)
            self.model.slides.append(self._slide)
        kind = values.get("data-pptx-kind", "").strip().lower()
        if self._slide is not None and kind:
            self._element = ExportElement(lowered, kind, values)
            self._slide.elements.append(self._element)
            self._element_depth = 1
        elif self._element is not None:
            self._element_depth += 1
        if self._element is not None and lowered == "tr":
            self._element.current_row = []
        if self._element is not None and lowered in {"td", "th"}:
            self._element.current_cell = []
        if (
            lowered == "link"
            and "stylesheet" in values.get("rel", "").lower().split()
            and values.get("href", "").strip()
        ):
            self.model.stylesheets.append(values["href"].strip())
        for attribute, value in values.items():
            if attribute in _URL_ATTRIBUTES and value.strip():
                self.model.resources.append((lowered, attribute, value.strip()))
            elif attribute == "srcset":
                for candidate in value.split(","):
                    reference = candidate.strip().split(maxsplit=1)[0]
                    if reference:
                        self.model.resources.append((lowered, attribute, reference))
        if lowered in {"area", "base", "br", "col", "embed", "hr", "img", "input", "link", "meta", "source", "track", "wbr"}:
            self._element = None
            self._element_depth = 0

    def handle_startendtag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        self.handle_starttag(tag, attrs)
        self.handle_endtag(tag)

    def handle_data(self, data: str) -> None:
        if self._slide is not None:
            self._slide.text_parts.append(data)
        if self._element is None:
            return
        if self._element.current_cell is not None:
            self._element.current_cell.append(data)
        else:
            self._element.text_parts.append(data)

    def handle_endtag(self, tag: str) -> None:
        lowered = tag.lower()
        if self._element is not None and lowered in {"td", "th"}:
            cell = " ".join(" ".join(self._element.current_cell or []).split())
            if self._element.current_row is not None:
                self._element.current_row.append(cell)
            self._element.current_cell = None
        if self._element is not None and lowered == "tr":
            if self._element.current_row is not None:
                self._element.rows.append(self._element.current_row)
            self._element.current_row = None
        if self._element is not None:
            self._element_depth -= 1
            if self._element_depth <= 0:
                self._element = None
                self._element_depth = 0
        if lowered == "section" and self._slide is not None:
            self._slide = None


def _atomic_json(path: Path, value: dict[str, Any]) -> None:
    path = _external_output(path)
    path.parent.mkdir(parents=True, exist_ok=True)
    temporary = path.with_name(f".{path.name}.tmp-{uuid.uuid4().hex}")
    try:
        with temporary.open("w", encoding="utf-8", newline="\n") as handle:
            json.dump(value, handle, ensure_ascii=True, indent=2, sort_keys=True)
            handle.write("\n")
            handle.flush()
            os.fsync(handle.fileno())
        os.replace(temporary, path)
    finally:
        temporary.unlink(missing_ok=True)


def parse_deck_html(html_path: Path | str) -> DeckModel:
    path = Path(html_path).expanduser().resolve(strict=True)
    raw = path.read_text(encoding="utf-8")
    parser = _DeckParser(path, raw)
    parser.feed(raw)
    parser.close()
    return parser.model


def _issue(code: str, message: str, *, slide_id: str = "") -> dict[str, str]:
    value = {"code": code, "message": message}
    if slide_id:
        value["slide_id"] = slide_id
    return value


def _positive_number(value: str, field_name: str) -> float:
    try:
        number = float(value)
    except (TypeError, ValueError) as error:
        raise PptContractError(f"{field_name} must be numeric") from error
    if not math.isfinite(number) or number < 0:
        raise PptContractError(f"{field_name} must be finite and non-negative")
    return number


def _bbox(element: ExportElement) -> tuple[float, float, float, float]:
    values = tuple(
        _positive_number(element.attrs.get(name, ""), name)
        for name in ("data-pptx-x", "data-pptx-y", "data-pptx-w", "data-pptx-h")
    )
    x, y, width, height = values
    if width <= 0 or height <= 0:
        raise PptContractError("editable element width and height must be positive")
    if x + width > SLIDE_WIDTH_PX + 1 or y + height > SLIDE_HEIGHT_PX + 1:
        raise PptContractError("editable element is outside the 1920x1080 canvas")
    return values


def _source_ids(value: str) -> list[str]:
    return [item for item in re.split(r"[\s,;]+", value.strip()) if item]


def _normalized_text(parts: Iterable[str]) -> str:
    return " ".join(" ".join(parts).split())


def resolve_local_dependency(root: Path, base: Path, reference: str) -> Path | None:
    parsed = urlsplit(reference.strip().strip("'\""))
    if parsed.scheme or parsed.netloc or reference.startswith("//"):
        raise PptContractError(f"asset is not local: {reference}")
    if not parsed.path and parsed.fragment:
        return None
    relative = Path(unquote(parsed.path))
    if not parsed.path or relative.is_absolute() or ".." in relative.parts:
        raise PptContractError(f"asset escapes the deck workspace: {reference}")
    workspace = root.resolve(strict=True)
    parent = base.resolve(strict=True)
    try:
        parent.relative_to(workspace)
    except ValueError as error:
        raise PptContractError(f"asset base escapes the deck workspace: {reference}") from error
    unresolved = parent / relative
    cursor = workspace
    try:
        parts = unresolved.relative_to(workspace).parts
    except ValueError as error:
        raise PptContractError(f"asset escapes the deck workspace: {reference}") from error
    for part in parts:
        cursor = cursor / part
        if cursor.is_symlink():
            raise PptContractError(f"local asset is a symlink: {reference}")
    candidate = unresolved.resolve(strict=True)
    try:
        candidate.relative_to(workspace)
    except ValueError as error:
        raise PptContractError(f"asset escapes the deck workspace: {reference}") from error
    status = candidate.lstat()
    if not S_ISREG(status.st_mode) or status.st_nlink > 1:
        raise PptContractError(f"local asset is missing or unsafe: {reference}")
    return candidate


def _local_asset(deck: DeckModel, reference: str) -> Path:
    candidate = resolve_local_dependency(
        deck.html_path.parent, deck.html_path.parent, reference
    )
    if candidate is None:
        raise PptContractError(f"local asset is missing or unsafe: {reference}")
    return candidate


def parse_speaker_notes(value: str) -> tuple[list[str], str] | None:
    match = re.fullmatch(
        r"\s*\[Sources\]\s*(?P<sources>.*?)\s+\[Talk\]\s*(?P<talk>\S.*?)\s*",
        value,
        flags=re.DOTALL,
    )
    if not match:
        return None
    sources = _source_ids(match.group("sources"))
    talk = " ".join(match.group("talk").split())
    if not sources or not talk:
        return None
    return sources, talk


def validate_deck_html(
    html_path: Path | str, *, expected_slide_count: int
) -> dict[str, Any]:
    deck = parse_deck_html(html_path)
    issues: list[dict[str, str]] = []
    expected_ids = [f"slide-{index:02d}" for index in range(1, expected_slide_count + 1)]
    actual_ids = [slide.slide_id for slide in deck.slides]
    if len(deck.slides) != expected_slide_count:
        issues.append(_issue("slide_count_mismatch", f"expected {expected_slide_count}, found {len(deck.slides)}"))
    exact_id_attributes = all(
        slide.attrs.get("id") == expected
        and slide.attrs.get("data-slide-id") == expected
        for slide, expected in zip(deck.slides, expected_ids)
    )
    if (
        actual_ids != expected_ids
        or len(set(actual_ids)) != len(actual_ids)
        or not exact_id_attributes
    ):
        issues.append(_issue("slide_id_contract", "slide ids must be unique and contiguous slide-01..slide-N"))
    root = deck.root_attrs
    if deck.root_count != 1:
        issues.append(_issue("deck_root_count", "deck must contain exactly one artifact root"))
    if root.get("data-slide-count") != str(expected_slide_count):
        issues.append(_issue("root_slide_count", "deck root data-slide-count does not match the plan"))
    if root.get("data-width") != "1920" or root.get("data-height") != "1080":
        issues.append(_issue("root_canvas", "deck root must declare a 1920x1080 canvas"))
    compact_css = re.sub(r"\s+", "", deck.raw_html.lower())
    if "width:1920px" not in compact_css or "height:1080px" not in compact_css:
        issues.append(_issue("css_canvas", "deck CSS must define 1920x1080 slide frames"))
    scripts = re.findall(
        r"<script\b(?P<attrs>[^>]*)>(?P<body>.*?)</script\s*>",
        deck.raw_html,
        flags=re.IGNORECASE | re.DOTALL,
    )
    navigation = (
        len(scripts) == 1
        and scripts[0][0].strip() == "data-autodesign-navigation"
        and scripts[0][1].strip() == SAFE_NAVIGATION_SCRIPT
    )
    if not navigation:
        issues.append(
            _issue(
                "unsafe_script",
                "deck permits only the exact bundled ArrowLeft/ArrowRight hash-navigation script",
            )
        )
    if re.search(r"\son[a-z]+\s*=", deck.raw_html, flags=re.IGNORECASE):
        issues.append(_issue("event_handler_attribute", "inline event-handler attributes are forbidden"))
    if re.search(r"<(?:audio|embed|iframe|object|video)\b", deck.raw_html, flags=re.IGNORECASE):
        issues.append(_issue("unsafe_embedded_content", "embedded executable or media content is forbidden"))
    css_queue: list[tuple[Path, str]] = [
        (deck.html_path.parent, css)
        for css in re.findall(
            r"<style\b[^>]*>(.*?)</style\s*>",
            deck.raw_html,
            flags=re.IGNORECASE | re.DOTALL,
        )
    ]
    css_queue.extend(
        (deck.html_path.parent, match.group(1))
        for match in re.finditer(
            r"\sstyle\s*=\s*['\"](.*?)['\"]",
            deck.raw_html,
            flags=re.IGNORECASE | re.DOTALL,
        )
    )
    seen_css_files: set[Path] = set()
    for reference in deck.stylesheets:
        try:
            stylesheet = resolve_local_dependency(
                deck.html_path.parent, deck.html_path.parent, reference
            )
        except (OSError, PptContractError):
            continue
        if stylesheet is not None and stylesheet not in seen_css_files:
            seen_css_files.add(stylesheet)
            try:
                css_queue.append(
                    (stylesheet.parent, stylesheet.read_text(encoding="utf-8"))
                )
            except (OSError, UnicodeError) as error:
                issues.append(_issue("unsafe_local_asset", str(error)))
    css_url = re.compile(
        r"url\(\s*(?:(['\"])(.*?)\1|([^)'\"\s]+))\s*\)",
        flags=re.IGNORECASE,
    )
    css_import = re.compile(
        r"@import\s+(?:url\(\s*)?(?:(['\"])(.*?)\1|([^)'\";\s]+))\s*\)?",
        flags=re.IGNORECASE,
    )
    generated_content_reported = False
    position = 0
    while position < len(css_queue):
        base, css = css_queue[position]
        position += 1
        if not generated_content_reported and re.search(
            r"(?:^|[;{])\s*content\s*:\s*(?!none\b|normal\b|['\"]{2})[^;}]+",
            css,
            flags=re.IGNORECASE,
        ):
            issues.append(
                _issue(
                    "css_generated_text",
                    "CSS generated content would become raster-only text in PowerPoint",
                )
            )
            generated_content_reported = True
        import_references = {
            match.group(2) or match.group(3) for match in css_import.finditer(css)
        }
        references = {
            *(match.group(2) or match.group(3) for match in css_url.finditer(css)),
            *import_references,
        }
        for reference in references:
            try:
                dependency = resolve_local_dependency(
                    deck.html_path.parent, base, reference
                )
            except (OSError, PptContractError) as error:
                parsed = urlsplit(reference)
                code = (
                    "remote_asset"
                    if parsed.scheme.lower() in _REMOTE_SCHEMES
                    or reference.startswith("//")
                    else "unsafe_local_asset"
                )
                issues.append(_issue(code, str(error)))
                continue
            if (
                reference in import_references
                and dependency is not None
                and dependency not in seen_css_files
            ):
                seen_css_files.add(dependency)
                try:
                    css_queue.append(
                        (dependency.parent, dependency.read_text(encoding="utf-8"))
                    )
                except (OSError, UnicodeError) as error:
                    issues.append(_issue("unsafe_local_asset", str(error)))

    assertion_count = 0
    note_count = 0
    claim_ids: set[str] = set()
    for index, slide in enumerate(deck.slides, start=1):
        slide_id = slide.slide_id or f"slide-at-{index}"
        attrs = slide.attrs
        if attrs.get("data-width") != "1920" or attrs.get("data-height") != "1080":
            issues.append(_issue("slide_canvas", "slide must declare 1920x1080", slide_id=slide_id))
        if attrs.get("data-slide-index") != str(index):
            issues.append(
                _issue(
                    "slide_index_contract",
                    f"slide index must be {index}",
                    slide_id=slide_id,
                )
            )
        for required in ("data-slide-role", "data-section", "data-assertion-title"):
            if not attrs.get(required, "").strip():
                issues.append(_issue("missing_slide_metadata", f"missing {required}", slide_id=slide_id))
        if attrs.get("data-assertion-title", "").strip():
            assertion_count += 1
        if not _source_ids(attrs.get("data-source-ids", "")):
            issues.append(_issue("missing_slide_sources", "slide must cite source ids", slide_id=slide_id))
        slide_sources = _source_ids(attrs.get("data-source-ids", ""))
        notes = attrs.get("data-speaker-notes", "").strip()
        if notes:
            note_count += 1
        parsed_notes = parse_speaker_notes(notes)
        if parsed_notes is None:
            issues.append(_issue("speaker_note_contract", "notes require [Sources] and [Talk]", slide_id=slide_id))
        elif parsed_notes[0] != slide_sources:
            issues.append(
                _issue(
                    "speaker_note_sources",
                    "speaker-note source IDs must exactly match slide source IDs",
                    slide_id=slide_id,
                )
            )
        editable_text = sum(
            1 for element in slide.elements if element.kind == "text" and element.text
        )
        editable_visible_parts: list[str] = []
        slide_claim_count = 0
        for element in slide.elements:
            if element.kind not in {"text", "image", "shape", "table"}:
                issues.append(_issue("unknown_editable_kind", f"unknown data-pptx-kind {element.kind}", slide_id=slide_id))
                continue
            try:
                x, y, width, height = _bbox(element)
            except PptContractError as error:
                issues.append(_issue("invalid_editable_bbox", str(error), slide_id=slide_id))
                continue
            if element.kind == "text":
                editable_visible_parts.append(element.text)
                if not element.text:
                    issues.append(_issue("empty_editable_text", "editable text must not be empty", slide_id=slide_id))
                size = element.attrs.get("data-font-size", "")
                if size:
                    try:
                        if _positive_number(size, "data-font-size") < 14:
                            issues.append(_issue("small_typography", "editable text is below 14 px", slide_id=slide_id))
                    except PptContractError as error:
                        issues.append(_issue("invalid_typography", str(error), slide_id=slide_id))
                claim_id = element.attrs.get("data-claim-id", "").strip()
                if not _source_ids(element.attrs.get("data-source-ids", "")):
                    issues.append(
                        _issue(
                            "missing_text_sources",
                            "every visible text element must cite source IDs",
                            slide_id=slide_id,
                        )
                    )
                if claim_id:
                    slide_claim_count += 1
                    if claim_id in claim_ids:
                        issues.append(_issue("duplicate_claim_id", f"duplicate claim id {claim_id}", slide_id=slide_id))
                    claim_ids.add(claim_id)
                    if not _source_ids(element.attrs.get("data-source-ids", "")):
                        issues.append(_issue("ungrounded_claim", f"claim {claim_id} has no source ids", slide_id=slide_id))
            elif element.kind == "image":
                reference = element.attrs.get("src", "").strip()
                try:
                    _local_asset(deck, reference)
                except (OSError, PptContractError) as error:
                    code = (
                        "remote_asset"
                        if urlsplit(reference).scheme.lower() in _REMOTE_SCHEMES
                        or reference.startswith("//")
                        else "unsafe_local_asset"
                    )
                    issues.append(_issue(code, str(error), slide_id=slide_id))
                if not _source_ids(element.attrs.get("data-source-ids", "")):
                    issues.append(_issue("missing_image_sources", "source image must cite source ids", slide_id=slide_id))
                if width * height >= SLIDE_WIDTH_PX * SLIDE_HEIGHT_PX * 0.9 and editable_text == 0:
                    issues.append(_issue("whole_slide_rasterization", "a full-slide image cannot replace editable content", slide_id=slide_id))
            elif element.kind == "table":
                editable_visible_parts.extend(
                    cell for row in element.rows for cell in row
                )
                if not element.rows or not any(any(cell.strip() for cell in row) for row in element.rows):
                    issues.append(_issue("empty_editable_table", "editable table must contain cells", slide_id=slide_id))
                if not _source_ids(element.attrs.get("data-source-ids", "")):
                    issues.append(
                        _issue(
                            "missing_table_sources",
                            "every native table must cite source IDs",
                            slide_id=slide_id,
                        )
                    )
        if editable_text == 0:
            issues.append(_issue("missing_editable_text", "every slide requires native editable text", slide_id=slide_id))
        if slide_claim_count == 0:
            issues.append(_issue("missing_slide_claim", "every slide requires a source-bound claim element", slide_id=slide_id))
        if _normalized_text(slide.text_parts) != _normalized_text(editable_visible_parts):
            issues.append(
                _issue(
                    "untagged_visible_text",
                    "all visible slide text must be explicitly tagged for native PPTX export",
                    slide_id=slide_id,
                )
            )

    for tag, attribute, reference in deck.resources:
        if tag == "img" and any(
            element.attrs.get("src") == reference
            for slide in deck.slides
            for element in slide.elements
            if element.kind == "image"
        ):
            continue
        parsed = urlsplit(reference)
        if parsed.scheme.lower() in _REMOTE_SCHEMES or reference.startswith("//"):
            issues.append(_issue("remote_asset", f"remote {tag}[{attribute}] is forbidden"))
        elif parsed.scheme:
            issues.append(_issue("unsafe_asset_scheme", f"unsupported asset scheme: {parsed.scheme}"))
        elif not reference.startswith("#"):
            try:
                resolve_local_dependency(
                    deck.html_path.parent, deck.html_path.parent, reference
                )
            except (OSError, PptContractError) as error:
                issues.append(_issue("unsafe_local_asset", str(error)))

    return {
        "format_version": 1,
        "passed": not issues,
        "expected_slide_count": expected_slide_count,
        "actual_slide_count": len(deck.slides),
        "slide_ids": actual_ids,
        "assertion_title_count": assertion_count,
        "speaker_note_count": note_count,
        "canvas": {"width": SLIDE_WIDTH_PX, "height": SLIDE_HEIGHT_PX},
        "keyboard_navigation": navigation,
        "issues": issues,
    }


def claims_from_deck(deck: DeckModel) -> list[dict[str, Any]]:
    claims: list[dict[str, Any]] = []
    for slide in deck.slides:
        for element_index, element in enumerate(slide.elements, start=1):
            if element.kind == "text":
                text = element.text
            elif element.kind == "table":
                text = " | ".join(
                    cell for row in element.rows for cell in row if cell.strip()
                )
            else:
                continue
            claim_id = element.attrs.get("data-claim-id", "").strip() or (
                f"{slide.slide_id}-{element.kind}-{element_index:02d}"
            )
            claims.append(
                {
                    "id": claim_id,
                    "text": text,
                    "source_ids": _source_ids(element.attrs.get("data-source-ids", "")),
                    "claim_type": f"slide_{element.kind}",
                }
            )
        parsed_notes = parse_speaker_notes(
            slide.attrs.get("data-speaker-notes", "")
        )
        if parsed_notes is not None:
            source_ids, talk = parsed_notes
            claims.append(
                {
                    "id": f"{slide.slide_id}-speaker-notes",
                    "text": talk,
                    "source_ids": source_ids,
                    "claim_type": "speaker_notes",
                }
            )
    return claims


def notes_from_deck(deck: DeckModel) -> list[dict[str, Any]]:
    return [
        {
            "slide_id": slide.slide_id,
            "source_ids": _source_ids(slide.attrs.get("data-source-ids", "")),
            "speaker_notes": slide.attrs.get("data-speaker-notes", ""),
        }
        for slide in deck.slides
    ]


def _rgb(value: str, default: str = "000000"):
    from pptx.dml.color import RGBColor

    match = _HEX_COLOR.fullmatch(value.strip()) if value else None
    return RGBColor.from_string(match.group(1).upper() if match else default)


def _inches(value: float):
    from pptx.util import Inches

    return Inches(value / PX_PER_INCH)


def _set_shape_metadata(shape: Any, *, name: str, source_ids: Iterable[str]) -> None:
    shape.name = name[:255]
    description = "source_ids=" + ",".join(source_ids)
    nodes = shape._element.xpath(".//p:cNvPr")
    if nodes:
        nodes[0].set("descr", description[:1024])


def _add_text(slide: Any, element: ExportElement, slide_id: str) -> None:
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Pt

    x, y, width, height = _bbox(element)
    shape = slide.shapes.add_textbox(_inches(x), _inches(y), _inches(width), _inches(height))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = _inches(6)
    frame.margin_right = _inches(6)
    frame.margin_top = _inches(4)
    frame.margin_bottom = _inches(4)
    anchors = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}
    frame.vertical_anchor = anchors.get(element.attrs.get("data-valign", "top"), MSO_ANCHOR.TOP)
    paragraph = frame.paragraphs[0]
    paragraph.text = element.text
    alignments = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}
    paragraph.alignment = alignments.get(element.attrs.get("data-align", "left"), PP_ALIGN.LEFT)
    run = paragraph.runs[0]
    font = run.font
    font.name = element.attrs.get("data-font-family", "Arial").split(",")[0].strip(" '\"") or "Arial"
    font.size = Pt(
        _positive_number(element.attrs.get("data-font-size", "28"), "data-font-size")
        * 72.0
        / PX_PER_INCH
    )
    font.bold = element.attrs.get("data-bold", "").lower() in {"1", "true", "yes", "bold"}
    font.italic = element.attrs.get("data-italic", "").lower() in {"1", "true", "yes", "italic"}
    font.color.rgb = _rgb(element.attrs.get("data-color", "#171717"), "171717")
    fill = element.attrs.get("data-fill", "").strip()
    if fill and fill.lower() != "none":
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(fill, "FFFFFF")
    else:
        shape.fill.background()
    shape.line.fill.background()
    _set_shape_metadata(
        shape,
        name=f"{slide_id}:text:{element.attrs.get('data-claim-id', 'editable')}",
        source_ids=_source_ids(element.attrs.get("data-source-ids", "")),
    )


def _add_shape(slide: Any, element: ExportElement, slide_id: str) -> None:
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR

    x, y, width, height = _bbox(element)
    shape_kind = element.attrs.get("data-shape", "rect").lower()
    if shape_kind == "line":
        shape = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            _inches(x),
            _inches(y),
            _inches(x + width),
            _inches(y + height),
        )
        shape.line.color.rgb = _rgb(element.attrs.get("data-stroke", "#171717"), "171717")
        shape.line.width = _inches(max(1.0, _positive_number(element.attrs.get("data-stroke-width", "2"), "data-stroke-width")))
    else:
        shape_type = MSO_AUTO_SHAPE_TYPE.OVAL if shape_kind == "ellipse" else MSO_AUTO_SHAPE_TYPE.RECTANGLE
        shape = slide.shapes.add_shape(shape_type, _inches(x), _inches(y), _inches(width), _inches(height))
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(element.attrs.get("data-fill", "#D7D9D2"), "D7D9D2")
        stroke = element.attrs.get("data-stroke", "").strip()
        if stroke and stroke.lower() != "none":
            shape.line.color.rgb = _rgb(stroke)
        else:
            shape.line.fill.background()
    _set_shape_metadata(shape, name=f"{slide_id}:shape:{shape_kind}", source_ids=_source_ids(element.attrs.get("data-source-ids", "")))


def _add_image(slide: Any, deck: DeckModel, element: ExportElement, slide_id: str) -> None:
    x, y, width, height = _bbox(element)
    source = _local_asset(deck, element.attrs.get("src", ""))
    shape = slide.shapes.add_picture(str(source), _inches(x), _inches(y), _inches(width), _inches(height))
    _set_shape_metadata(shape, name=f"{slide_id}:image:{source.name}", source_ids=_source_ids(element.attrs.get("data-source-ids", "")))


def _add_table(slide: Any, element: ExportElement, slide_id: str) -> None:
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Pt

    x, y, width, height = _bbox(element)
    rows = element.rows
    columns = max((len(row) for row in rows), default=0)
    if not rows or columns <= 0:
        raise PptContractError("editable table has no cells")
    shape = slide.shapes.add_table(len(rows), columns, _inches(x), _inches(y), _inches(width), _inches(height))
    table = shape.table
    for row_index, row in enumerate(rows):
        for column_index in range(columns):
            cell = table.cell(row_index, column_index)
            cell.text = row[column_index] if column_index < len(row) else ""
            cell.margin_left = _inches(5)
            cell.margin_right = _inches(5)
            cell.margin_top = _inches(3)
            cell.margin_bottom = _inches(3)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _rgb("#EEECE7" if row_index == 0 else "#FFFFFF")
            paragraph = cell.text_frame.paragraphs[0]
            paragraph.alignment = PP_ALIGN.LEFT
            for run in paragraph.runs:
                run.font.name = element.attrs.get("data-font-family", "Arial")
                run.font.size = Pt(
                    _positive_number(element.attrs.get("data-font-size", "20"), "data-font-size")
                    * 72.0
                    / PX_PER_INCH
                )
                run.font.bold = row_index == 0
                run.font.color.rgb = _rgb(element.attrs.get("data-color", "#171717"), "171717")
    _set_shape_metadata(shape, name=f"{slide_id}:table:editable", source_ids=_source_ids(element.attrs.get("data-source-ids", "")))


def _slide_background(slide: SlideModel) -> str:
    value = slide.attrs.get("data-background", "#FFFFFF")
    return value if _HEX_COLOR.fullmatch(value) else "#FFFFFF"


def make_text_free_backgrounds(
    deck: DeckModel,
    preview_dir: Path | str,
    output_dir: Path | str,
) -> dict[str, Path]:
    from PIL import Image, ImageDraw

    previews = Path(preview_dir).resolve(strict=True)
    output = _external_output(output_dir)
    output.mkdir(parents=True, exist_ok=True)
    backgrounds: dict[str, Path] = {}
    for slide in deck.slides:
        preview = previews / f"{slide.slide_id}.png"
        if not preview.is_file():
            raise PptContractError(f"missing HTML preview for {slide.slide_id}")
        with Image.open(preview) as source:
            image = source.convert("RGB").resize((SLIDE_WIDTH_PX, SLIDE_HEIGHT_PX))
        draw = ImageDraw.Draw(image)
        fill = _slide_background(slide)
        for element in slide.elements:
            x, y, width, height = _bbox(element)
            padding = 3
            draw.rectangle(
                (
                    max(0, round(x) - padding),
                    max(0, round(y) - padding),
                    min(SLIDE_WIDTH_PX, round(x + width) + padding),
                    min(SLIDE_HEIGHT_PX, round(y + height) + padding),
                ),
                fill=fill,
            )
        target = output / f"{slide.slide_id}.png"
        image.save(target, format="PNG", optimize=True)
        backgrounds[slide.slide_id] = target
    return backgrounds


def native_object_contract(deck: DeckModel) -> dict[str, Any]:
    """Describe the native objects that must survive PPTX export and reopen."""

    slides: list[dict[str, Any]] = []
    for slide in deck.slides:
        shape_types = [
            element.attrs.get("data-shape", "rect").strip().lower()
            for element in slide.elements
            if element.kind == "shape"
        ]
        if any(kind not in {"rect", "ellipse", "line"} for kind in shape_types):
            raise PptContractError("native PPTX object contract has an unknown shape type")
        slides.append(
            {
                "slide_id": slide.slide_id,
                "text_shapes": sum(
                    element.kind == "text" and bool(element.text)
                    for element in slide.elements
                ),
                "table_shapes": sum(
                    element.kind == "table" for element in slide.elements
                ),
                "image_shapes": sum(
                    element.kind == "image" for element in slide.elements
                ),
                "shape_count": len(shape_types),
                "shape_types": shape_types,
                "speaker_notes": slide.attrs.get("data-speaker-notes", "").strip(),
            }
        )
    return {"format_version": 1, "slides": slides}


def _validate_native_contract(
    native_contract: Mapping[str, Any], expected_slide_count: int
) -> list[Mapping[str, Any]]:
    slides = native_contract.get("slides")
    if native_contract.get("format_version") != 1 or not isinstance(slides, list):
        raise PptContractError("native PPTX object contract is invalid")
    if len(slides) != expected_slide_count or any(
        not isinstance(slide, Mapping) for slide in slides
    ):
        raise PptContractError("native PPTX object contract slide count is invalid")
    expected_ids = [f"slide-{index:02d}" for index in range(1, expected_slide_count + 1)]
    if [slide.get("slide_id") for slide in slides] != expected_ids:
        raise PptContractError("native PPTX object contract slide order is invalid")
    for slide in slides:
        for field_name in (
            "text_shapes",
            "table_shapes",
            "image_shapes",
            "shape_count",
        ):
            value = slide.get(field_name)
            if isinstance(value, bool) or not isinstance(value, int) or value < 0:
                raise PptContractError(
                    f"native PPTX object contract has invalid {field_name}"
                )
        shape_types = slide.get("shape_types")
        if (
            not isinstance(shape_types, list)
            or len(shape_types) != slide["shape_count"]
            or any(kind not in {"rect", "ellipse", "line"} for kind in shape_types)
        ):
            raise PptContractError("native PPTX object contract has invalid shape types")
        if not isinstance(slide.get("speaker_notes"), str):
            raise PptContractError("native PPTX object contract has invalid speaker notes")
    return slides


def inspect_pptx(
    path: Path | str,
    *,
    expected_slide_count: int,
    native_contract: Mapping[str, Any] | None = None,
) -> dict[str, Any]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE

    target = Path(path).resolve(strict=True)
    presentation = Presentation(str(target))
    issues: list[dict[str, str]] = []
    if len(presentation.slides) != expected_slide_count:
        issues.append(_issue("pptx_slide_count", f"expected {expected_slide_count}, found {len(presentation.slides)}"))
    width = round(presentation.slide_width / 914400, 6)
    height = round(presentation.slide_height / 914400, 6)
    if abs(width - SLIDE_WIDTH_IN) > 0.001 or abs(height - SLIDE_HEIGHT_IN) > 0.001:
        issues.append(_issue("pptx_canvas", f"unexpected slide size {width}x{height} inches"))
    notes_count = 0
    editable_text = 0
    tables = 0
    pictures = 0
    editable_shapes = 0
    expected_native = (
        _validate_native_contract(native_contract, expected_slide_count)
        if native_contract is not None
        else None
    )
    observed_native: list[dict[str, Any]] = []
    for index, slide in enumerate(presentation.slides, start=1):
        slide_id = f"slide-{index:02d}"
        slide_text = 0
        slide_tables = 0
        slide_images = 0
        slide_shape_types: list[str] = []
        notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            notes_count += 1
        if "[Sources]" not in notes or "[Talk]" not in notes:
            issues.append(_issue("pptx_notes", "speaker notes are missing [Sources]/[Talk]", slide_id=slide_id))
        for shape in slide.shapes:
            name = str(getattr(shape, "name", ""))
            if getattr(shape, "has_table", False):
                tables += 1
                slide_tables += 1
                editable_shapes += 1
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                pictures += 1
                if not name.startswith("background:"):
                    slide_images += 1
                    editable_shapes += 1
            elif getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
                editable_text += 1
                editable_shapes += 1
                slide_text += 1
            else:
                editable_shapes += 1
                if shape.shape_type == MSO_SHAPE_TYPE.LINE:
                    slide_shape_types.append("line")
                elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                    if shape.auto_shape_type == MSO_AUTO_SHAPE_TYPE.RECTANGLE:
                        slide_shape_types.append("rect")
                    elif shape.auto_shape_type == MSO_AUTO_SHAPE_TYPE.OVAL:
                        slide_shape_types.append("ellipse")
                    else:
                        slide_shape_types.append("unknown-auto-shape")
                else:
                    slide_shape_types.append(f"unknown-{shape.shape_type}")
        if slide_text == 0:
            issues.append(_issue("pptx_missing_editable_text", "slide has no editable text overlay", slide_id=slide_id))
        observed = {
            "slide_id": slide_id,
            "text_shapes": slide_text,
            "table_shapes": slide_tables,
            "image_shapes": slide_images,
            "shape_count": len(slide_shape_types),
            "shape_types": slide_shape_types,
            "speaker_notes": notes,
        }
        observed_native.append(observed)
        if expected_native is not None and index <= len(expected_native):
            expected = expected_native[index - 1]
            for field_name, issue_code in (
                ("text_shapes", "pptx_native_text"),
                ("table_shapes", "pptx_native_table"),
                ("image_shapes", "pptx_native_image"),
            ):
                if observed[field_name] != expected[field_name]:
                    issues.append(
                        _issue(
                            issue_code,
                            f"expected {expected[field_name]} {field_name}, found {observed[field_name]}",
                            slide_id=slide_id,
                        )
                    )
            if (
                observed["shape_count"] != expected["shape_count"]
                or observed["shape_types"] != expected["shape_types"]
            ):
                issues.append(
                    _issue(
                        "pptx_native_shape",
                        "native shape count or types differ from the canonical deck contract",
                        slide_id=slide_id,
                    )
                )
            if observed["speaker_notes"] != expected["speaker_notes"].strip():
                issues.append(
                    _issue(
                        "pptx_native_notes",
                        "speaker notes differ from the canonical deck contract",
                        slide_id=slide_id,
                    )
                )
    if notes_count != expected_slide_count:
        issues.append(_issue("pptx_notes_count", f"expected {expected_slide_count} note pages, found {notes_count}"))
    return {
        "format_version": 1,
        "passed": not issues,
        "slide_count": len(presentation.slides),
        "slide_size_inches": [width, height],
        "notes_count": notes_count,
        "editable_text_shapes": editable_text,
        "table_shapes": tables,
        "picture_shapes": pictures,
        "editable_shape_count": editable_shapes,
        "native_object_contract": observed_native,
        "issues": issues,
    }


def export_deck_to_pptx(
    html_path: Path | str,
    output_path: Path | str,
    *,
    preview_dir: Path | str | None = None,
    background_dir: Path | str | None = None,
) -> dict[str, Any]:
    from pptx import Presentation
    from pptx.util import Inches

    deck = parse_deck_html(html_path)
    object_contract = native_object_contract(deck)
    validation = validate_deck_html(deck.html_path, expected_slide_count=len(deck.slides))
    if not validation["passed"]:
        raise PptContractError("HTML deck contract failed before PPTX export")
    backgrounds: dict[str, Path] = {}
    if preview_dir is not None:
        if background_dir is None:
            raise PptContractError("background_dir is required with preview_dir")
        backgrounds = make_text_free_backgrounds(deck, preview_dir, background_dir)

    presentation = Presentation()
    presentation.slide_width = Inches(SLIDE_WIDTH_IN)
    presentation.slide_height = Inches(SLIDE_HEIGHT_IN)
    while presentation.slides:
        relationship_id = presentation.slides._sldIdLst[0].rId
        presentation.part.drop_rel(relationship_id)
        del presentation.slides._sldIdLst[0]
    blank = presentation.slide_layouts[6]
    for slide_model in deck.slides:
        slide = presentation.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = _rgb(_slide_background(slide_model), "FFFFFF")
        background = backgrounds.get(slide_model.slide_id)
        if background is not None:
            picture = slide.shapes.add_picture(str(background), 0, 0, presentation.slide_width, presentation.slide_height)
            _set_shape_metadata(picture, name=f"background:{slide_model.slide_id}", source_ids=[])
        for element in slide_model.elements:
            if element.kind == "text":
                _add_text(slide, element, slide_model.slide_id)
            elif element.kind == "shape":
                _add_shape(slide, element, slide_model.slide_id)
            elif element.kind == "image":
                _add_image(slide, deck, element, slide_model.slide_id)
            elif element.kind == "table":
                _add_table(slide, element, slide_model.slide_id)
        slide.notes_slide.notes_text_frame.text = slide_model.attrs.get("data-speaker-notes", "")
    presentation.core_properties.title = "Source-grounded conference deck"
    presentation.core_properties.subject = "Editable HTML-derived presentation"
    output = _external_output(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    temporary = output.with_name(f".{output.name}.tmp-{uuid.uuid4().hex}")
    try:
        presentation.save(str(temporary))
        with zipfile.ZipFile(temporary) as archive:
            bad = archive.testzip()
            if bad:
                raise PptContractError(f"PPTX archive member failed CRC: {bad}")
        os.replace(temporary, output)
    finally:
        temporary.unlink(missing_ok=True)
    return inspect_pptx(
        output,
        expected_slide_count=len(deck.slides),
        native_contract=object_contract,
    )


def render_pptx_with_libreoffice(
    pptx_path: Path | str,
    output_dir: Path | str,
    *,
    timeout_seconds: int = 180,
) -> dict[str, Any]:
    executable = shutil.which("soffice") or shutil.which("libreoffice")
    if not executable:
        return {"performed": False, "reason": "LibreOffice renderer unavailable"}
    pptx = Path(pptx_path).resolve(strict=True)
    output = _external_output(output_dir)
    output.mkdir(parents=True, exist_ok=True)
    rendered = output / f"{pptx.stem}.pdf"
    if rendered.is_symlink():
        raise PptContractError("LibreOffice output PDF must not be a symlink")
    rendered.unlink(missing_ok=True)
    with tempfile.TemporaryDirectory(prefix="autodesign-ppt-office-") as profile:
        command = [
            executable,
            f"-env:UserInstallation={Path(profile).as_uri()}",
            "--headless",
            "--convert-to",
            "pdf",
            "--outdir",
            str(output),
            str(pptx),
        ]
        result = subprocess.run(command, capture_output=True, text=True, timeout=timeout_seconds, check=False)
    if result.returncode != 0 or not rendered.is_file():
        detail = (result.stderr or result.stdout or "no renderer output").strip()[-1000:]
        raise PptContractError(f"LibreOffice could not reopen/render deck.pptx: {detail}")
    page_count = pdf_page_count(rendered)
    return {"performed": True, "pdf": str(rendered), "page_count": page_count}


def pdf_page_count(path: Path | str) -> int:
    executable = shutil.which("pdfinfo")
    if not executable:
        raise PptContractError("pdfinfo is required to verify PDF page count")
    result = subprocess.run([executable, str(Path(path).resolve(strict=True))], capture_output=True, text=True, timeout=30, check=False)
    if result.returncode != 0:
        raise PptContractError("pdfinfo could not inspect the generated PDF")
    match = re.search(r"(?m)^Pages:\s*(\d+)\s*$", result.stdout)
    if not match:
        raise PptContractError("pdfinfo did not report a page count")
    return int(match.group(1))


def compare_rendered_slides(
    canonical_paths: Sequence[Path | str],
    rendered_paths: Sequence[Path | str],
) -> dict[str, Any]:
    from PIL import Image, ImageChops, ImageFilter, ImageStat

    if len(canonical_paths) != len(rendered_paths) or not canonical_paths:
        raise PptContractError("rendered comparison requires equal non-empty slide sets")
    scores: list[float] = []
    metrics: list[dict[str, float]] = []
    for canonical_path, rendered_path in zip(canonical_paths, rendered_paths):
        with Image.open(canonical_path) as canonical_image:
            canonical = canonical_image.convert("RGB").resize((640, 360))
        with Image.open(rendered_path) as rendered_image:
            rendered = rendered_image.convert("RGB").resize((640, 360))
        difference = ImageChops.difference(canonical, rendered)
        mean = sum(ImageStat.Stat(difference).mean) / 3.0
        similarity = round(max(0.0, 1.0 - mean / 255.0), 6)
        scores.append(similarity)
        canonical_edges = canonical.convert("L").filter(ImageFilter.FIND_EDGES)
        rendered_edges = rendered.convert("L").filter(ImageFilter.FIND_EDGES)
        # FIND_EDGES marks the outer frame even for a blank image. Removing it
        # keeps the metric focused on actual slide content.
        for edge_image in (canonical_edges, rendered_edges):
            edge_image.paste(0, (0, 0, edge_image.width, 2))
            edge_image.paste(0, (0, edge_image.height - 2, edge_image.width, edge_image.height))
            edge_image.paste(0, (0, 0, 2, edge_image.height))
            edge_image.paste(0, (edge_image.width - 2, 0, edge_image.width, edge_image.height))
        canonical_mask = canonical_edges.point(lambda value: 255 if value >= 24 else 0)
        rendered_mask = rendered_edges.point(lambda value: 255 if value >= 24 else 0)
        rendered_nearby = rendered_mask.filter(ImageFilter.MaxFilter(9))
        canonical_data = canonical_mask.get_flattened_data()
        rendered_data = rendered_nearby.get_flattened_data()
        canonical_pixels = sum(1 for value in canonical_data if value)
        overlap_pixels = sum(
            1
            for canonical_value, rendered_value in zip(
                canonical_data, rendered_data
            )
            if canonical_value and rendered_value
        )
        edge_recall = (
            overlap_pixels / canonical_pixels if canonical_pixels else 1.0
        )
        metrics.append(
            {
                "similarity": similarity,
                "edge_recall": round(edge_recall, 6),
            }
        )
    minimum = min(scores)
    minimum_edge_recall = min(metric["edge_recall"] for metric in metrics)
    return {
        "performed": True,
        "slide_scores": scores,
        "slide_metrics": metrics,
        "minimum_similarity": minimum,
        "minimum_edge_recall": minimum_edge_recall,
        "mean_similarity": round(sum(scores) / len(scores), 6),
        "passed": minimum >= 0.65 and minimum_edge_recall >= 0.2,
    }


def _command_export(args: argparse.Namespace) -> int:
    report = export_deck_to_pptx(
        args.html,
        args.output,
        preview_dir=args.preview_dir,
        background_dir=args.background_dir,
    )
    _atomic_json(args.report, report)
    return 0 if report["passed"] else 2


def _command_inspect(args: argparse.Namespace) -> int:
    report = inspect_pptx(args.pptx, expected_slide_count=args.expected_slide_count)
    _atomic_json(args.report, report)
    return 0 if report["passed"] else 2


def _command_compare(args: argparse.Namespace) -> int:
    canonical = [
        args.canonical_dir / f"slide-{index:02d}.png"
        for index in range(1, args.expected_slide_count + 1)
    ]
    rendered = sorted(args.rendered_dir.glob("slide-*.png"))
    if any(not path.is_file() for path in canonical):
        raise PptContractError("canonical slide previews are incomplete")
    if len(rendered) != args.expected_slide_count:
        raise PptContractError(
            f"expected {args.expected_slide_count} rendered PPTX pages, found {len(rendered)}"
        )
    report = compare_rendered_slides(canonical, rendered)
    _atomic_json(args.report, report)
    return 0 if report["passed"] else 2


def main(argv: Sequence[str] | None = None) -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    commands = parser.add_subparsers(dest="command", required=True)
    export = commands.add_parser("export")
    export.add_argument("--html", type=Path, required=True)
    export.add_argument("--output", type=Path, required=True)
    export.add_argument("--report", type=Path, required=True)
    export.add_argument("--preview-dir", type=Path)
    export.add_argument("--background-dir", type=Path)
    export.set_defaults(handler=_command_export)
    inspect = commands.add_parser("inspect")
    inspect.add_argument("--pptx", type=Path, required=True)
    inspect.add_argument("--expected-slide-count", type=int, required=True)
    inspect.add_argument("--report", type=Path, required=True)
    inspect.set_defaults(handler=_command_inspect)
    compare = commands.add_parser("compare")
    compare.add_argument("--canonical-dir", type=Path, required=True)
    compare.add_argument("--rendered-dir", type=Path, required=True)
    compare.add_argument("--expected-slide-count", type=int, required=True)
    compare.add_argument("--report", type=Path, required=True)
    compare.set_defaults(handler=_command_compare)
    args = parser.parse_args(argv)
    try:
        return int(args.handler(args))
    except (OSError, UnicodeError, ValueError, PptContractError, subprocess.SubprocessError) as error:
        print(f"ERROR: {error}", file=sys.stderr)
        return 1


if __name__ == "__main__":
    raise SystemExit(main())
