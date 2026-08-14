"""Core cancellable HTML-to-PPTX derived job."""

from __future__ import annotations

import json
import os
from pathlib import Path
import re
import secrets
import shlex
import shutil
import stat
import sys
from typing import Any
from urllib.parse import urlsplit

from bs4 import BeautifulSoup

from .agents.external_author_process import (
    ExternalAuthorProcessRequest,
    run_external_author_process,
)
from .config import Settings, code_editor_command_for_harness, harness_subprocess_env
from .run_control import CancellationToken, durable_replace_json
from .util.io import atomic_write_json


class PptxExportJobError(RuntimeError):
    """A derived PPTX export failed without depending on the Web transport."""

    def __init__(self, code: str, message: str, *, diagnostics: dict[str, Any] | None = None) -> None:
        self.code = code
        self.diagnostics = diagnostics or {}
        super().__init__(message)


def safe_export_stem(raw: str) -> str:
    clean = re.sub(r"[^a-zA-Z0-9._-]+", "-", raw.strip()).strip("-._")
    return (clean or "artifact")[:80]


def html_export_canvas_size(
    source: Path,
    artifact: dict[str, Any],
    canvas: dict[str, Any],
) -> tuple[int, int]:
    if artifact.get("artifact_type") == "poster":
        authored = _authored_paper_poster_size(source)
        if authored is not None:
            return authored
    width = _positive_int(canvas.get("w")) or _num_from_html_attr(source, "data-w") or 1440
    height = _positive_int(canvas.get("h")) or _num_from_html_attr(source, "data-h") or 900
    return max(1, width), max(1, height)


def run_pptx_export_job(
    *,
    run_id: str,
    run_dir: Path,
    source_html: Path,
    artifact: dict[str, Any],
    settings: Settings,
    cancellation_token: CancellationToken,
    artifact_name: str | None = None,
) -> dict[str, Any]:
    """Export one source HTML snapshot while keeping all writes in the derived run.

    The exact result keys are ``run_id``, ``pptx_path``, ``manifest_path``,
    ``attempts``, and ``canvas``. Transport adapters must normalize them to
    their private/public protocol explicitly.
    """

    token = cancellation_token
    run_dir = Path(run_dir)
    requested_source = Path(source_html)
    token.raise_if_cancelled("pptx_export.start")
    recover_pptx_delivery_transactions(run_dir / "exports")
    if requested_source.is_symlink():
        raise PptxExportJobError(
            "unsafe_source_tree",
            "source HTML must not be a symlink",
        )
    source = requested_source.resolve()
    if not source.is_file() or source.is_symlink():
        raise PptxExportJobError("missing_source_html", f"source HTML does not exist: {source}")
    source_root = source.parent.resolve()
    canvas = artifact.get("canvas") if isinstance(artifact.get("canvas"), dict) else {}
    width, height = html_export_canvas_size(source, artifact, canvas)
    command = _resolved_command(settings)
    timeout_s = max(1, int(getattr(settings, "code_editor_timeout_s", 600) or 600))
    max_attempts = max(1, int(getattr(settings, "code_editor_max_attempts", 2) or 2))
    attempts_root = run_dir / "pptx-export"
    _mkdir(attempts_root, token, "pptx_export.mkdir_attempts")

    attempts: list[dict[str, Any]] = []
    repair_feedback: dict[str, Any] | None = None
    for attempt_index in range(1, max_attempts + 1):
        token.raise_if_cancelled("pptx_export.attempt_start")
        attempt_dir = attempts_root / f"attempt_{attempt_index:02d}"
        _mkdir(attempt_dir, token, "pptx_export.mkdir_attempt")
        _stage_inputs(source, source_root, attempt_dir, token)
        if repair_feedback is not None:
            _atomic_json(
                attempt_dir / "validation_feedback.json",
                repair_feedback,
                token,
                "pptx_export.write_feedback",
            )
        prompt = _build_prompt(
            source_name=source.name,
            artifact=artifact,
            width=width,
            height=height,
            repair_feedback=repair_feedback,
        )
        _write_text(
            attempt_dir / "export_prompt.md",
            prompt,
            token,
            "pptx_export.write_prompt",
        )
        invocation = _invoke_agent(
            run_id=run_id,
            run_dir=run_dir,
            attempt_index=attempt_index,
            attempt_dir=attempt_dir,
            command=command,
            prompt=prompt,
            timeout_s=timeout_s,
            settings=settings,
            token=token,
        )
        token.raise_if_cancelled("pptx_export.after_agent")
        validation = _validate_pptx(attempt_dir / "export.pptx", attempt_dir)
        attempt_record = {
            "attempt": attempt_index,
            "attempt_dir": str(attempt_dir),
            "invocation": invocation,
            "validation": validation,
        }
        attempts.append(attempt_record)
        _atomic_json(
            attempt_dir / "export_attempt_result.json",
            attempt_record,
            token,
            "pptx_export.write_attempt_result",
        )
        if invocation.get("status") == "ok" and validation.get("ok") is True:
            stem = safe_export_stem(artifact_name or str(artifact.get("name") or source.stem or "artifact"))
            export_dir = run_dir / "exports"
            _mkdir(export_dir, token, "pptx_export.mkdir_output")
            recover_pptx_delivery_transactions(export_dir)
            output = export_dir / f"{stem}.pptx"
            partial = export_dir / f".{stem}.partial.pptx"
            _copy_file(
                attempt_dir / "export.pptx",
                partial,
                token,
                "pptx_export.copy_validated_output",
            )
            manifest = {
                "render_mode": "external_code_editor_pptx_export",
                "source_html": str(source),
                "run_id": run_id,
                "code_editor_harness": getattr(settings, "code_editor_harness", "codex"),
                "attempts": attempts,
                "canvas": {"w": width, "h": height},
            }
            manifest_path = output.with_suffix(".agent-export.json")
            _publish_pptx_delivery(
                staged_pptx=partial,
                output=output,
                manifest_path=manifest_path,
                manifest=manifest,
                token=token,
            )
            return {
                "run_id": run_id,
                "pptx_path": str(output),
                "manifest_path": str(manifest_path),
                "attempts": attempts,
                "canvas": {"w": width, "h": height},
            }
        repair_feedback = {
            "reason": invocation.get("reason") if invocation.get("status") != "ok" else "validation_failed",
            "invocation": invocation,
            "validation": validation,
        }

    raise PptxExportJobError(
        "pptx_agent_export_failed",
        "PowerPoint export agent did not produce a valid .pptx",
        diagnostics={"attempts": attempts[-2:]},
    )


def _resolved_command(settings: Settings) -> tuple[str, ...]:
    configured = str(getattr(settings, "code_editor_cmd", "") or "").strip()
    harness = str(getattr(settings, "code_editor_harness", "codex") or "codex").strip()
    command = configured or code_editor_command_for_harness(
        harness,
        getattr(settings, "code_editor_model", None),
    )
    try:
        parts = tuple(shlex.split(command))
    except ValueError as exc:
        raise PptxExportJobError("invalid_agent_command", f"could not parse code-editor command: {exc}") from exc
    if not parts:
        raise PptxExportJobError("missing_agent_command", "code-editor command is not configured")
    return parts


def _invoke_agent(
    *,
    run_id: str,
    run_dir: Path,
    attempt_index: int,
    attempt_dir: Path,
    command: tuple[str, ...],
    prompt: str,
    timeout_s: int,
    settings: Settings,
    token: CancellationToken,
) -> dict[str, Any]:
    output = attempt_dir / "export.pptx"
    done_marker = attempt_dir / "export_done.json"
    for target in (output, done_marker):
        token.raise_if_cancelled("pptx_export.reset_output.before")
        try:
            target.unlink()
        except FileNotFoundError:
            pass
        token.raise_if_cancelled("pptx_export.reset_output.after")
    env = harness_subprocess_env(
        _minimal_process_environment(),
        harness=str(getattr(settings, "code_editor_harness", "") or ""),
        api_key=getattr(settings, "harness_api_key", None),
    )
    author_python = (
        env.get("AUTODESIGN_AUTHOR_PYTHON", "").strip()
        or env.get("DESIGN_ANYTHING_AUTHOR_PYTHON", "").strip()
        or sys.executable
    )
    env["AUTODESIGN_AUTHOR_PYTHON"] = author_python
    env.setdefault("DESIGN_ANYTHING_AUTHOR_PYTHON", author_python)
    token.raise_if_cancelled("pptx_export.before_agent_spawn")
    result = run_external_author_process(
        ExternalAuthorProcessRequest(
            run_id=run_id,
            attempt=attempt_index,
            command=command,
            cwd=attempt_dir,
            prompt=prompt,
            timeout_s=timeout_s,
            stdout_path=attempt_dir / ".export_agent.stdout.log",
            stderr_path=attempt_dir / ".export_agent.stderr.log",
            env=env,
            completion_requested=lambda: _agent_output_completion_reason(
                output,
                done_marker,
                attempt_dir,
            ),
            interruption_requested=token.is_cancelled,
            poll_interval_s=0.025,
            run_dir=run_dir,
            cancellation_token=token,
            sensitive_values=(str(getattr(settings, "harness_api_key", "") or ""),),
        )
    )
    token.raise_if_cancelled("pptx_export.after_agent_process")
    output_safe = _is_regular_file_within(output, attempt_dir)
    marker_safe = _is_regular_file_within(done_marker, attempt_dir)
    status = "ok" if result.status == "ok" and output_safe and marker_safe else "error"
    reason = result.reason
    if status != "ok" and (output.exists() or done_marker.exists()) and not (output_safe and marker_safe):
        reason = "unsafe_agent_output"
    elif status != "ok" and output.exists() and not done_marker.exists():
        reason = "missing_done_marker"
    elif status != "ok" and not output.exists():
        reason = "missing_export_pptx"
    return {
        "status": status,
        "reason": reason,
        "returncode": result.returncode,
        "timed_out": result.timed_out,
        "elapsed_s": round(result.elapsed_s, 3),
        "stdout_excerpt": result.stdout[-1400:],
        "stderr_excerpt": result.stderr[-900:],
        "output_safety": {
            "pptx_regular_within_attempt": output_safe,
            "done_marker_regular_within_attempt": marker_safe,
        },
    }


def _stage_inputs(
    source: Path,
    source_root: Path,
    attempt_dir: Path,
    token: CancellationToken,
) -> None:
    _assert_source_entry(source, source_root)
    _copy_file(source, attempt_dir / "current.html", token, "pptx_export.stage_current")
    if source.name != "current.html":
        _copy_file(source, attempt_dir / source.name, token, "pptx_export.stage_named_source")
    for child in source.parent.iterdir():
        token.raise_if_cancelled("pptx_export.stage_sibling")
        if child.name == "exports":
            continue
        _assert_source_entry(child, source_root)
        try:
            if child.resolve() == source:
                continue
        except OSError:
            pass
        target = attempt_dir / child.name
        if child.is_dir():
            _copy_tree(child, target, source_root, token)
        elif child.is_file():
            try:
                if child.stat().st_size <= 50 * 1024 * 1024:
                    _copy_file(child, target, token, "pptx_export.stage_file")
            except OSError:
                continue


def _minimal_process_environment() -> dict[str, str]:
    allowed = {
        "PATH",
        "HOME",
        "TMPDIR",
        "TMP",
        "TEMP",
        "LANG",
        "LC_ALL",
        "LC_CTYPE",
        "SYSTEMROOT",
        "COMSPEC",
        "PATHEXT",
        "WINDIR",
        "USERPROFILE",
        "APPDATA",
        "LOCALAPPDATA",
        "XDG_CACHE_HOME",
        "XDG_CONFIG_HOME",
        "SSL_CERT_FILE",
        "SSL_CERT_DIR",
        "REQUESTS_CA_BUNDLE",
        "CURL_CA_BUNDLE",
        "HTTP_PROXY",
        "HTTPS_PROXY",
        "ALL_PROXY",
        "NO_PROXY",
        "http_proxy",
        "https_proxy",
        "all_proxy",
        "no_proxy",
        "PLAYWRIGHT_BROWSERS_PATH",
        "PLAYWRIGHT_NODEJS_PATH",
        "CHROME_PATH",
    }
    return {
        name: value
        for name, value in os.environ.items()
        if name in allowed and not (_is_proxy_name(name) and _proxy_has_credentials(value))
    }


def _is_proxy_name(name: str) -> bool:
    return name.upper() in {"HTTP_PROXY", "HTTPS_PROXY", "ALL_PROXY"}


def _proxy_has_credentials(value: str) -> bool:
    parsed = urlsplit(value if "://" in value else f"//{value}")
    return parsed.username is not None or parsed.password is not None


def _copy_tree(
    source: Path,
    destination: Path,
    source_root: Path,
    token: CancellationToken,
) -> None:
    _assert_source_entry(source, source_root)
    _mkdir(destination, token, "pptx_export.stage_directory")
    for child in source.iterdir():
        token.raise_if_cancelled("pptx_export.stage_tree")
        if child.name in {"__pycache__", ".DS_Store"}:
            continue
        _assert_source_entry(child, source_root)
        target = destination / child.name
        if child.is_dir():
            _copy_tree(child, target, source_root, token)
        elif child.is_file():
            try:
                if child.stat().st_size <= 50 * 1024 * 1024:
                    _copy_file(child, target, token, "pptx_export.stage_tree_file")
            except OSError:
                continue


def _assert_source_entry(path: Path, source_root: Path) -> None:
    if path.is_symlink():
        raise PptxExportJobError(
            "unsafe_source_tree",
            f"source tree contains a symlink: {path.name}",
        )
    try:
        resolved = path.resolve(strict=True)
        mode = path.lstat().st_mode
    except OSError as exc:
        raise PptxExportJobError(
            "unsafe_source_tree",
            f"source tree entry is not stable: {path.name}",
        ) from exc
    if not _path_inside(resolved, source_root.resolve()):
        raise PptxExportJobError(
            "unsafe_source_tree",
            f"source tree entry escapes its canonical root: {path.name}",
        )
    if not (stat.S_ISREG(mode) or stat.S_ISDIR(mode)):
        raise PptxExportJobError(
            "unsafe_source_tree",
            f"source tree contains a non-regular entry: {path.name}",
        )


def _is_regular_file_within(path: Path, root: Path) -> bool:
    try:
        if path.is_symlink() or not stat.S_ISREG(path.lstat().st_mode):
            return False
        resolved = path.resolve(strict=True)
    except OSError:
        return False
    return _path_inside(resolved, root.resolve())


def _agent_output_completion_reason(
    output: Path,
    done_marker: Path,
    attempt_dir: Path,
) -> str | None:
    output_safe = _is_regular_file_within(output, attempt_dir)
    marker_safe = _is_regular_file_within(done_marker, attempt_dir)
    if output_safe and marker_safe:
        return "done_marker"
    if _path_entry_exists(output) and not output_safe:
        return "unsafe_agent_output"
    if _path_entry_exists(done_marker) and not marker_safe:
        return "unsafe_agent_output"
    return None


def _path_entry_exists(path: Path) -> bool:
    try:
        path.lstat()
        return True
    except OSError:
        return False


def _path_inside(path: Path, root: Path) -> bool:
    try:
        path.relative_to(root)
        return True
    except ValueError:
        return False


def _build_prompt(
    *,
    source_name: str,
    artifact: dict[str, Any],
    width: int,
    height: int,
    repair_feedback: dict[str, Any] | None,
) -> str:
    name = str(artifact.get("name") or source_name or "AutoDesign artifact").strip()
    repair_block = ""
    if repair_feedback is not None:
        repair_block = (
            "\nThis is a repair attempt. Read `validation_feedback.json` first, then fix the export. "
            "Do not repeat a conversion strategy that created invalid or unreadable PPTX output.\n"
        )
    return f"""Export the current AutoDesign HTML artifact to an editable PowerPoint file.

Artifact name: {name}
Source HTML files:
- current.html
- {source_name}

Canvas size: {width} x {height} CSS pixels.

Output contract:
- Write `export.pptx`.
- Write `export_done.json` with a concise JSON summary.
- Keep all work inside this directory.
- Use this Python interpreter exactly: {sys.executable}
- Do not install packages, create virtual environments, or run setup commands.

Quality contract:
- Match the browser-rendered poster/design as closely as practical.
- Set the PowerPoint slide size to the same aspect ratio and logical size as the HTML canvas.
- Avoid naive text-node extraction that creates overlapping text boxes.
- Prefer section-level and table-cell-level extraction.
- Keep ordinary titles, body text, labels, section bars, metric chips, and tables editable.
- Use images only for real source figures, charts, logos, or diagrams already present in the HTML.
- Preserve visual hierarchy, colors, panel boundaries, spacing, and a useful one-slide composition.
- Do not use remote, file:, script, iframe, event-handler, or unsafe URLs.

Suggested implementation:
1. Use Playwright to render `current.html` locally and inspect computed boxes at {width}x{height}.
2. Use python-pptx to create a single blank slide at the matching dimensions.
3. Rebuild semantic blocks, tables, images, and editable text at their rendered positions.
4. Save `export.pptx`, reopen it, and verify at least one slide and useful shapes.
{repair_block}
"""


def _validate_pptx(path: Path, attempt_dir: Path) -> dict[str, Any]:
    errors: list[str] = []
    warnings: list[str] = []
    if not _is_regular_file_within(path, attempt_dir):
        if path.exists() or path.is_symlink():
            return {
                "ok": False,
                "errors": ["export.pptx is not a regular file owned by the attempt"],
                "warnings": warnings,
            }
        return {"ok": False, "errors": ["export.pptx missing"], "warnings": warnings}
    try:
        size = path.stat().st_size
    except OSError as exc:
        return {"ok": False, "errors": [f"export.pptx unreadable: {exc}"], "warnings": warnings}
    if size < 2048:
        errors.append("export.pptx is too small")
    try:
        from pptx import Presentation

        presentation = Presentation(str(path))
        slide_count = len(presentation.slides)
        shape_count = sum(len(slide.shapes) for slide in presentation.slides)
        if slide_count < 1:
            errors.append("export.pptx has no slides")
        if shape_count < 3:
            warnings.append("export.pptx has very few shapes")
    except Exception as exc:  # noqa: BLE001
        errors.append(f"export.pptx could not be opened: {exc}")
        slide_count = 0
        shape_count = 0
    return {
        "ok": not errors,
        "errors": errors,
        "warnings": warnings,
        "bytes": size,
        "slide_count": slide_count,
        "shape_count": shape_count,
    }


def _authored_paper_poster_size(source: Path) -> tuple[int, int] | None:
    try:
        text = source.read_text(encoding="utf-8", errors="ignore")
        document = BeautifulSoup(text, "html.parser")
    except OSError:
        return None
    root = document.select_one(".paper-poster")
    if root is None or document.select_one(".canvas") is not None:
        return None
    width = _positive_int(root.get("data-w"))
    height = _positive_int(root.get("data-h"))
    if width and height:
        return width, height
    style = str(root.get("style") or "")
    width = _style_px(style, "width") or _css_px(text, ".paper-poster", "width")
    height = _style_px(style, "height") or _css_px(text, ".paper-poster", "height")
    return (width or 3072, height or 1536)


def _positive_int(raw: Any) -> int | None:
    try:
        value = int(float(str(raw)))
    except (TypeError, ValueError):
        return None
    return value if value > 0 else None


def _num_from_html_attr(path: Path, attribute: str) -> int | None:
    try:
        text = path.read_text(encoding="utf-8", errors="ignore")
    except OSError:
        return None
    match = re.search(rf"\b{re.escape(attribute)}=[\"']([0-9.]+)[\"']", text)
    return _positive_int(match.group(1)) if match else None


def _style_px(style: str, property_name: str) -> int | None:
    match = re.search(rf"(?:^|;)\s*{re.escape(property_name)}\s*:\s*([0-9.]+)px", style, re.IGNORECASE)
    return _positive_int(match.group(1)) if match else None


def _css_px(text: str, selector: str, property_name: str) -> int | None:
    match = re.search(
        rf"{re.escape(selector)}\s*\{{[^}}]*\b{re.escape(property_name)}\s*:\s*([0-9.]+)px",
        text,
        re.IGNORECASE | re.DOTALL,
    )
    return _positive_int(match.group(1)) if match else None


def _mkdir(path: Path, token: CancellationToken, phase: str) -> None:
    token.raise_if_cancelled(f"{phase}.before")
    path.mkdir(parents=True, exist_ok=True)
    token.raise_if_cancelled(f"{phase}.after")


def _write_text(path: Path, content: str, token: CancellationToken, phase: str) -> None:
    token.raise_if_cancelled(f"{phase}.before")
    path.write_text(content, encoding="utf-8")
    token.raise_if_cancelled(f"{phase}.after")


def _copy_file(source: Path, destination: Path, token: CancellationToken, phase: str) -> None:
    token.raise_if_cancelled(f"{phase}.before")
    shutil.copy2(source, destination)
    token.raise_if_cancelled(f"{phase}.after")


def _publish_pptx_delivery(
    *,
    staged_pptx: Path,
    output: Path,
    manifest_path: Path,
    manifest: dict[str, Any],
    token: CancellationToken,
) -> None:
    """Publish one PPTX/manifest delivery and roll back the pair together."""

    transaction = f"{os.getpid()}-{secrets.token_hex(10)}"
    staged_manifest = manifest_path.with_name(
        f".{manifest_path.name}.{transaction}.partial"
    )
    output_backup = output.with_name(f".{output.name}.{transaction}.rollback")
    manifest_backup = manifest_path.with_name(
        f".{manifest_path.name}.{transaction}.rollback"
    )
    journal_path = output.with_name(
        f".{output.stem}.delivery-transaction.json"
    )
    had_output = output.is_file()
    had_manifest = manifest_path.is_file()
    journal_payload = {
        "version": 1,
        "state": "prepared",
        "output_name": output.name,
        "manifest_name": manifest_path.name,
        "staged_pptx_name": staged_pptx.name,
        "staged_manifest_name": staged_manifest.name,
        "output_backup_name": output_backup.name,
        "manifest_backup_name": manifest_backup.name,
        "had_output": had_output,
        "had_manifest": had_manifest,
    }
    try:
        _atomic_json(
            staged_manifest,
            manifest,
            token,
            "pptx_export.write_manifest",
        )
        if had_output:
            shutil.copy2(output, output_backup)
        if had_manifest:
            shutil.copy2(manifest_path, manifest_backup)
        for durable_file in (
            staged_pptx,
            staged_manifest,
            *(path for path in (output_backup, manifest_backup) if path.is_file()),
        ):
            _fsync_file(durable_file)
        _fsync_directory(output.parent)
        durable_replace_json(journal_path, journal_payload)
        token.raise_if_cancelled("pptx_export.delivery.before_pptx_replace")
        os.replace(staged_pptx, output)
        _fsync_directory(output.parent)
        token.raise_if_cancelled("pptx_export.delivery.after_pptx_replace")
        os.replace(staged_manifest, manifest_path)
        _fsync_directory(output.parent)
        token.raise_if_cancelled("pptx_export.delivery.after_manifest_replace")
        token.raise_if_cancelled("pptx_export.complete")
        durable_replace_json(
            journal_path,
            {**journal_payload, "state": "committed"},
        )
    except BaseException:
        if journal_path.is_file():
            _recover_pptx_delivery_journal(journal_path)
        else:
            _unlink_delivery_paths(
                staged_pptx,
                staged_manifest,
                output_backup,
                manifest_backup,
            )
        raise
    else:
        _recover_pptx_delivery_journal(journal_path)


def recover_pptx_delivery_transactions(export_dir: Path) -> tuple[Path, ...]:
    """Recover durable PPTX delivery journals before export or Web startup."""

    root = Path(export_dir)
    if not root.is_dir():
        return ()
    recovered: list[Path] = []
    for journal_path in sorted(root.glob(".*.delivery-transaction.json")):
        _recover_pptx_delivery_journal(journal_path)
        recovered.append(journal_path)
    return tuple(recovered)


def _recover_pptx_delivery_journal(journal_path: Path) -> None:
    root = journal_path.parent
    try:
        payload = json.loads(journal_path.read_text(encoding="utf-8"))
    except (OSError, ValueError) as exc:
        raise PptxExportJobError(
            "pptx_delivery_recovery_failed",
            f"could not read PPTX delivery journal: {journal_path.name}",
        ) from exc
    if not isinstance(payload, dict) or payload.get("version") != 1:
        raise PptxExportJobError(
            "pptx_delivery_recovery_failed",
            f"unsupported PPTX delivery journal: {journal_path.name}",
        )
    state = payload.get("state")
    if state not in {"prepared", "rolled_back", "committed"}:
        raise PptxExportJobError(
            "pptx_delivery_recovery_failed",
            f"invalid PPTX delivery journal state: {state!r}",
        )
    paths = {
        key: _journal_local_path(root, payload, key)
        for key in (
            "output_name",
            "manifest_name",
            "staged_pptx_name",
            "staged_manifest_name",
            "output_backup_name",
            "manifest_backup_name",
        )
    }
    had_output = payload.get("had_output")
    had_manifest = payload.get("had_manifest")
    if type(had_output) is not bool or type(had_manifest) is not bool:
        raise PptxExportJobError(
            "pptx_delivery_recovery_failed",
            "PPTX delivery journal has invalid prior-file flags",
        )
    output = paths["output_name"]
    manifest = paths["manifest_name"]
    output_backup = paths["output_backup_name"]
    manifest_backup = paths["manifest_backup_name"]
    if state == "prepared":
        for existed, backup in (
            (had_output, output_backup),
            (had_manifest, manifest_backup),
        ):
            if existed and not _is_regular_file_within(backup, root):
                raise PptxExportJobError(
                    "pptx_delivery_recovery_failed",
                    f"PPTX delivery rollback backup is missing: {backup.name}",
                )
        _restore_delivery_file(output, output_backup, existed=had_output)
        _restore_delivery_file(manifest, manifest_backup, existed=had_manifest)
        _fsync_directory(root)
        payload = {**payload, "state": "rolled_back"}
        durable_replace_json(journal_path, payload)
        state = "rolled_back"
    if state == "rolled_back":
        for existed, target in (
            (had_output, output),
            (had_manifest, manifest),
        ):
            if (existed and not _is_regular_file_within(target, root)) or (
                not existed and _path_entry_exists(target)
            ):
                raise PptxExportJobError(
                    "pptx_delivery_recovery_failed",
                    "rolled-back PPTX delivery is incomplete",
                )
    elif not (
        _is_regular_file_within(output, root)
        and _is_regular_file_within(manifest, root)
    ):
        raise PptxExportJobError(
            "pptx_delivery_recovery_failed",
            "committed PPTX delivery is incomplete",
        )
    _unlink_delivery_paths(
        paths["staged_pptx_name"],
        paths["staged_manifest_name"],
        output_backup,
        manifest_backup,
    )
    _fsync_directory(root)
    journal_path.unlink(missing_ok=True)
    _fsync_directory(root)


def _journal_local_path(root: Path, payload: dict[str, Any], key: str) -> Path:
    raw = payload.get(key)
    if not isinstance(raw, str) or not raw or Path(raw).name != raw:
        raise PptxExportJobError(
            "pptx_delivery_recovery_failed",
            f"PPTX delivery journal has invalid {key}",
        )
    return root / raw


def _restore_delivery_file(target: Path, backup: Path, *, existed: bool) -> None:
    if not existed:
        target.unlink(missing_ok=True)
        return
    recovery = target.with_name(f".{target.name}.recovering")
    shutil.copy2(backup, recovery)
    _fsync_file(recovery)
    os.replace(recovery, target)


def _unlink_delivery_paths(*paths: Path) -> None:
    for path in paths:
        path.unlink(missing_ok=True)


def _fsync_file(path: Path) -> None:
    with path.open("rb") as handle:
        os.fsync(handle.fileno())


def _fsync_directory(directory: Path) -> None:
    if os.name == "nt":
        return
    descriptor = os.open(directory, os.O_RDONLY)
    try:
        os.fsync(descriptor)
    finally:
        os.close(descriptor)


def _atomic_json(path: Path, data: Any, token: CancellationToken, phase: str) -> None:
    token.raise_if_cancelled(f"{phase}.before")
    atomic_write_json(path, data)
    token.raise_if_cancelled(f"{phase}.after")
