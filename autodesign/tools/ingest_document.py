"""ingest_document — v1.2 paper2any entry point.

v1.1 asked a VLM to locate figure bboxes on rasterized PDF pages. That
was unreliable: the model returned half-page screenshots, clipped
diagrams, and hallucinated "figures" on text-only pages.

v1.2 separates two concerns:

1. **Figure localization** is now done by **pymupdf directly** —
   `doc.extract_image(xref)` pulls embedded raster images at native
   resolution (e.g. 1890×1211 PNGs the paper author uploaded), and
   `page.get_drawings()` + proximity clustering + 300 dpi rendering
   catches vector-drawn architecture diagrams. No VLM guessing.

2. **Reading / matching** is still VLM work — but the default model
   is now Qwen-VL-Max via OpenRouter (~5× cheaper and faster than
   Claude Sonnet for this non-reasoning workload). Structure extraction
   is text-only by default and receives title / authors / abstract /
   sections / figures / tables JSON. `figures` now carries only
   `{caption, page, description}` — no bbox or idx (we already have
   bboxes from pymupdf).

   Candidate figure crops are sent only to caption matching, where the
   VLM decides which caption matches and whether the crop is a real
   figure. Table crops are sent only to table parsing.

The `rendered_layers` record shape is unchanged (see the downstream
hydration helpers in `autodesign/tools/composite.py`): callers reference
ingested figures by the stable
`layer_id` we register (e.g. `ingest_fig_01`), and the renderer
hydrates `src_path` / `bbox` / `caption` from the layer registry.

Markdown and image branches are untouched.
"""

from __future__ import annotations

import json
import os
import re
import shutil
from concurrent.futures import FIRST_COMPLETED, ThreadPoolExecutor, wait
from pathlib import Path
from typing import Any, Callable, Mapping

import fitz  # pymupdf
from PIL import Image, ImageDraw, ImageFont

from ..agents.deck_outline_agent import DeckOutlineAgent
from ..agents.paper_memory_agent import PaperMemoryAgent
from ..config import effective_poster_harness_mode
from ._contract import ToolContext, obs_error, obs_ok
from ..schema import ToolResultRecord
from ..util.canvas_planner import refine_canvas_plan_from_ingest
from ..util.deck_planner import (
    fallback_deck_plan,
    should_refine_deck_plan,
)
from ..util.io import atomic_write_json, sha256_file
from ..util.logging import log
from ..util.run_paths import resolve_run_dir
from ..util.academic_palette import (
    active_academic_color_system,
    academic_color_system_options,
    explicit_academic_color_system,
    rank_academic_color_system_options,
    select_academic_color_system,
)
from ..util.pipeline_cache import (
    cache_entry_dir,
    pipeline_cache_enabled,
    stable_cache_key,
)
from ..util.poster_plan_contract import (
    build_poster_plan_contract,
    preflight_poster_plan_contract,
)
from ..util.pdf import (
    PdfFigureCandidate,
    PdfTableCandidate,
    ScannedPdfError,
    dedup_raster_vector,
    dedup_tables_against_figures,
    detect_scanned_pdf,
    discover_captioned_visual_groups,
    extract_embedded_rasters,
    extract_page_text,
    extract_table_candidates,
    extract_vector_clusters,
    page_count,
    pdf_figure_crop_quality_flags,
    pdf_table_crop_quality_flags,
    recover_caption_anchored_visuals,
    render_page_png,
)
from ..util.paper_visual_storyboard import build_paper_visual_storyboard
from ..util.source_visual_eligibility import (
    classify_source_visual,
    constrain_optional_source_visual_ids,
)
from ..util.paper_source_sanity import (
    PaperSourceInputError,
    PaperSourceSanityError,
    assert_valid_paper_source_pdf,
)
from ..util.paper_memory import (
    build_paper_memory,
    compact_paper_memory_for_planner,
    merge_paper_memories,
    read_paper_memory_cache,
    write_paper_memory_cache,
    paper_memory_cache_key,
    write_paper_memory_run_artifacts,
)
from ..util.paper_memory_dossier import (
    read_paper_memory_dossier_cache,
    write_paper_memory_dossier_cache,
    write_paper_memory_dossier_run_artifacts,
)
from ..util.vlm import VlmImage, vlm_call_json
from .discover_paper_resources import (
    discover_paper_resources_for_context,
    should_auto_discover_paper_resources,
)


# Max PDF bytes we accept in one call (belt-and-suspenders — pymupdf
# itself can open almost anything, but ingest touches every page and
# we want to fail fast on pathological inputs rather than spin).
_MAX_PDF_BYTES = 80 * 1024 * 1024  # 80 MB
_MAX_PDF_PAGES = 100
_INGEST_EXECUTOR_CANCEL_POLL_S = 0.02

_MD_IMG_RE = re.compile(r"!\[([^\]]*)\]\(([^)\s]+)(?:\s+\"[^\"]*\")?\)")
_TEXT_UNIT_BUCKETS = ("problem", "method", "evidence", "limitations", "takeaways")

_VISUAL_ROLE_PRIORITY = {
    "method": 5,
    "evidence": 4,
    "table": 4,
    "qualitative": 3,
    "fallback": 1,
}
_LOW_INFORMATION_VISUAL_TERMS = (
    "logo",
    "publisher",
    "watermark",
    "header",
    "footer",
    "copyright",
    "decorative",
    "border",
    "equation",
)
_SEVERE_CROP_CURATION_FLAGS = frozenset({
    "algorithm_caption_leak",
    "body_text_leak",
    "caption_in_crop",
    "caption_strip_leak",
    "edge_visual_remnant",
    "figure_caption_leak",
    "header_band_leak",
    "multi_caption_leak",
    "other_caption_in_crop",
    "page_like_table_crop",
    "table_body_text_leak",
    "section_heading_leak",
    "running_header_leak",
    "partial_visual_crop",
    "neighbor_asset_leak",
    "page_like_figure_crop",
    "table_fragment_crop",
    "table_without_structure",
})
_SELECTED_VISUAL_BLOCKING_FLAGS = _SEVERE_CROP_CURATION_FLAGS | frozenset({
    "image_payload_unavailable",
    "low_information_visual",
    "low_value_example_crop",
    "unlocated_raster_component",
})
_SELECTED_VISUAL_CONDITIONAL_FLAGS = frozenset({
    "high_edge_whitespace",
    "low_caption_confidence",
    "low_detail_visual_content",
    "mostly_white_visual",
    "no_caption",
    "source_page_unknown",
})
_DENSE_SYNTHESIS_TARGETS_PATH = (
    Path(__file__).resolve().parents[2]
    / "skills"
    / "poster"
    / "visual_recipe"
    / "assets"
    / "dense_synthesis_targets.json"
)
_RESEARCH_SYNTHESIS_PROFILE = "research_synthesis_dense"
_EDITORIAL_FLOW_PROFILE = "conference_editorial_flow"
_DENSE_PAPER_POSTER_PROFILES = {
    _RESEARCH_SYNTHESIS_PROFILE,
    _EDITORIAL_FLOW_PROFILE,
}

# Max parallel caption-matching calls. VLM calls are HTTP-bound; 4–6 in
# flight keeps wall time tight without tripping OpenRouter rate limits.
_CAPTION_MATCH_PARALLELISM = 6
_CHEAP_CAPTION_MATCH_MAX_CANDIDATES = 48
_CHEAP_TABLE_PARSE_MAX_CANDIDATES = 10
_DOGFOOD_CAPTION_MATCH_MAX_CANDIDATES = 64
_DOGFOOD_TABLE_PARSE_MAX_CANDIDATES = 12
_HARNESS_TABLE_PARSE_TIMEOUT_S = 180.0
_HARNESS_TABLE_PARSE_MAX_RETRIES = 1
# Confidence floor for accepting a VLM caption match — below this we
# still keep the figure but flag it; `is_real_figure=false` drops it
# regardless of confidence.
_CAPTION_MATCH_MIN_CONFIDENCE = 0.35
_CAPTION_GEOMETRY_HIGH_CONFIDENCE_MIN_SCORE = 0.65
# DPI for the optional cover-page image handed to the VLM during
# structure extraction. This path is disabled by default; set
# INGEST_STRUCTURE_COVER_IMAGE=1 only when deliberately debugging visual
# title/logo grounding. Normal structure extraction is text-only.
_STRUCTURE_PAGE_DPI = 144
# Total text budget across all pages sent to the structure extractor.
# ~60k chars fits comfortably in the ~16k-token context window Qwen
# uses for this call while leaving headroom for the JSON response.
_STRUCTURE_TOTAL_TEXT_CAP = 60_000
# Scanned-PDF OCR fallback (v1.2.5). When `detect_scanned_pdf` returns
# True we render each page at this DPI and hand the PNG to the VLM for
# text extraction — 200 dpi is the sweet spot: dense enough for small
# body text, not so heavy that a 40-page doc burns minutes. We still
# cap at `_MAX_PDF_PAGES` so runaway OCR cost is impossible.
_OCR_PAGE_DPI = 200
_OCR_PAGE_PARALLELISM = 6
_OCR_PER_PAGE_TIMEOUT_S = 120.0

_REFERENCE_SECTION_RE = re.compile(
    r"^\s*(?:\d+(?:\.\d+)*\.?\s+)?(?:references|bibliography|literature\s+cited|works\s+cited|参考文献)\s*$",
    re.IGNORECASE,
)
_APPENDIX_SECTION_RE = re.compile(
    r"^\s*(?:\d+(?:\.\d+)*\.?\s+)?(?:appendix|appendices|supplementary\s+(?:material|materials|information|appendix)|supplement|supplements|附录)\b(?:\s*[:.\-]?\s*.*)?$",
    re.IGNORECASE,
)
_REFERENCE_ENTRY_RE = re.compile(
    r"(\b(?:19|20)\d{2}[a-z]?\b|\bdoi\b|\barxiv\b|https?://|^\s*\[\d+\]|^\s*\d+\.\s+[A-Z])",
    re.IGNORECASE,
)
_IDENTITY_HINT_RE = re.compile(
    r"\b(author|authors|team|affiliation|affiliations|institution|institutions|university|lab|labs|research|inc\.?|corp\.?|llc|team|团队|作者|单位|机构)\b",
    re.IGNORECASE,
)


def _int_env(name: str, default: int) -> int:
    raw = os.getenv(name, "").strip()
    if not raw:
        return default
    try:
        return int(raw)
    except ValueError:
        return default


def _float_env(name: str, default: float) -> float:
    raw = os.getenv(name, "").strip()
    if not raw:
        return default
    try:
        return float(raw)
    except ValueError:
        return default


def _bool_env(name: str, default: bool = False) -> bool:
    raw = os.getenv(name, "").strip().lower()
    if not raw:
        return default
    if raw in {"1", "true", "yes", "on", "y"}:
        return True
    if raw in {"0", "false", "no", "off", "n"}:
        return False
    return default


def _structure_cover_image_enabled() -> bool:
    return _bool_env("INGEST_STRUCTURE_COVER_IMAGE", False)


def _ingest_vlm_parallelism(ctx: ToolContext, default: int) -> int:
    explicit = _int_env("INGEST_VLM_PARALLELISM", 0)
    if explicit <= 0:
        explicit = _int_env("INGEST_CAPTION_PARALLELISM", 0)
    if explicit > 0:
        return max(1, explicit)
    if _poster_harness_mode(ctx) == "cheap":
        return 2
    model = str(getattr(ctx.settings, "ingest_model", "") or "").lower()
    if model.startswith("longcat-"):
        return 1
    return max(1, default)


def _raise_if_ingest_cancelled(ctx: ToolContext | None, phase: str) -> None:
    checker = getattr(ctx, "raise_if_cancelled", None)
    if callable(checker):
        checker(phase)


def _ingest_cancellation_token(ctx: ToolContext) -> Any:
    return getattr(ctx, "cancellation_token", None)


def _shutdown_ingest_executor(
    executor: ThreadPoolExecutor,
    *,
    cancelled: bool,
) -> None:
    if not cancelled:
        executor.shutdown(wait=True)
        return
    try:
        executor.shutdown(wait=False, cancel_futures=True)
    except TypeError:  # pragma: no cover - Python < 3.9 compatibility
        executor.shutdown(wait=False)


def _run_bounded_ingest_pool(
    *,
    items: list[tuple[int, Any]],
    worker: Callable[[int, Any], Any],
    parallelism: int,
    ctx: ToolContext,
    phase: str,
) -> list[tuple[int, Any, Exception | None]]:
    """Run at most ``parallelism`` ingest tasks with cancellation polling."""
    _raise_if_ingest_cancelled(ctx, f"{phase}.before_executor")
    if not items:
        return []

    executor = ThreadPoolExecutor(max_workers=parallelism)
    in_flight: dict[Any, int] = {}
    outcomes: list[tuple[int, Any, Exception | None]] = []
    next_item = 0

    def submit_until_full() -> None:
        nonlocal next_item
        while next_item < len(items) and len(in_flight) < parallelism:
            _raise_if_ingest_cancelled(ctx, f"{phase}.before_submit")
            item_index, payload = items[next_item]
            future = executor.submit(worker, item_index, payload)
            in_flight[future] = item_index
            next_item += 1
            _raise_if_ingest_cancelled(ctx, f"{phase}.after_submit")

    try:
        submit_until_full()
        while in_flight:
            _raise_if_ingest_cancelled(ctx, f"{phase}.before_wait")
            done, _pending = wait(
                tuple(in_flight),
                timeout=_INGEST_EXECUTOR_CANCEL_POLL_S,
                return_when=FIRST_COMPLETED,
            )
            _raise_if_ingest_cancelled(ctx, f"{phase}.after_wait")
            for future in done:
                item_index = in_flight.pop(future)
                _raise_if_ingest_cancelled(ctx, f"{phase}.before_result")
                try:
                    result = future.result()
                    error = None
                except Exception as exc:  # individual ingest failures degrade
                    result = None
                    error = exc
                _raise_if_ingest_cancelled(ctx, f"{phase}.after_result")
                outcomes.append((item_index, result, error))
                _raise_if_ingest_cancelled(ctx, f"{phase}.after_result_buffer")
            submit_until_full()
    except BaseException:
        for future in in_flight:
            future.cancel()
        _shutdown_ingest_executor(executor, cancelled=True)
        raise
    else:
        _shutdown_ingest_executor(executor, cancelled=False)
        _raise_if_ingest_cancelled(ctx, f"{phase}.after_executor")
        return outcomes


def _poster_harness_mode(ctx: ToolContext | None = None) -> str:
    settings = getattr(ctx, "settings", None) if ctx is not None else None
    return effective_poster_harness_mode(settings)


def _is_harness_mode(ctx: ToolContext | None = None) -> bool:
    return _poster_harness_mode(ctx) in {"cheap", "quality", "dogfood"}


def _paper_source_scope() -> str:
    raw = os.getenv("PAPER_SOURCE_SCOPE", "main_body").strip().lower()
    aliases = {
        "body": "main_body",
        "main": "main_body",
        "main-only": "main_body",
        "body_plus_supplement": "body_plus_appendix",
        "body+appendix": "body_plus_appendix",
        "appendix": "body_plus_appendix",
        "full": "all",
    }
    value = aliases.get(raw, raw)
    if value not in {"main_body", "body_plus_appendix", "all"}:
        return "main_body"
    return value


def _paper_body_page_window(page_texts: list[str]) -> dict[str, int | str | None]:
    """Return the PDF page window worth ingesting for poster synthesis.

    References, appendices, and supplementary pages are usually not poster
    source material. Author/team identity is handled separately, so this window
    intentionally describes content pages only.
    """
    total = len(page_texts)
    source_scope = _paper_source_scope()
    references_start: int | None = None
    appendix_start: int | None = None
    for idx, text in enumerate(page_texts, start=1):
        lines = [re.sub(r"\s+", " ", line).strip() for line in text.splitlines()]
        nonempty = [line for line in lines if line]
        if not nonempty:
            continue
        for pos, line in enumerate(nonempty):
            if (
                appendix_start is None
                and _APPENDIX_SECTION_RE.match(line)
                and _appendix_heading_starts_supplement(
                    nonempty,
                    pos,
                    page_num=idx,
                    total_pages=total,
                )
            ):
                appendix_start = idx
            if (
                references_start is None
                and _REFERENCE_SECTION_RE.match(line)
                and _reference_heading_starts_bibliography(
                    nonempty,
                    pos,
                    page_num=idx,
                    total_pages=total,
                )
            ):
                references_start = idx
        if references_start is not None and appendix_start is not None:
            break
    cutoff_candidates: list[tuple[int, str]] = []
    if source_scope != "all" and references_start is not None:
        cutoff_candidates.append((references_start, "references"))
    if source_scope == "main_body" and appendix_start is not None:
        cutoff_candidates.append((appendix_start, "appendix_or_supplement"))
    cutoff_start, cutoff_reason = min(cutoff_candidates) if cutoff_candidates else (None, None)
    body_page_count = max(1, int(cutoff_start) - 1) if cutoff_start else total
    return {
        "total_page_count": total,
        "body_page_count": body_page_count,
        "references_start_page": references_start,
        "appendix_start_page": appendix_start,
        "cutoff_start_page": cutoff_start,
        "cutoff_reason": cutoff_reason,
        "source_scope": source_scope,
        "ignored_reference_page_count": (
            max(0, total - references_start + 1)
            if source_scope != "all" and references_start is not None
            else 0
        ),
        "ignored_non_body_page_count": max(0, total - body_page_count),
    }


def _enforce_pdf_page_cap(page_window: Mapping[str, Any]) -> None:
    source_scope = str(page_window.get("source_scope") or _paper_source_scope())
    total_pages = int(page_window.get("total_page_count") or 0)
    body_pages = int(page_window.get("body_page_count") or total_pages)
    measured_pages = total_pages if source_scope == "all" else body_pages
    label = "pages" if source_scope == "all" else "body pages"
    if measured_pages > _MAX_PDF_PAGES:
        raise RuntimeError(
            f"PDF has {measured_pages} {label} (cap {_MAX_PDF_PAGES}). Trim the document."
        )


def _appendix_heading_starts_supplement(
    lines: list[str],
    pos: int,
    *,
    page_num: int,
    total_pages: int,
) -> bool:
    if page_num < 2:
        return False
    min_body_pages = 2 if total_pages <= 6 else 3
    if page_num <= min_body_pages:
        return False
    heading = lines[pos].lower()
    near_page_top = pos <= min(12, max(4, int(len(lines) * 0.15)))
    standalone_heading = heading.strip(" .:-") in {
        "appendix",
        "appendices",
        "supplement",
        "supplements",
        "supplementary material",
        "supplementary materials",
        "supplementary information",
        "supplementary appendix",
        "附录",
    }
    if "appendix" in heading or "appendices" in heading or "supplement" in heading or "附录" in heading:
        return standalone_heading or near_page_top
    late_page = page_num >= max(4, int(total_pages * 0.35))
    return late_page and near_page_top


def _reference_heading_starts_bibliography(
    lines: list[str],
    pos: int,
    *,
    page_num: int,
    total_pages: int,
) -> bool:
    """Return True when a reference heading is the bibliography boundary.

    Real papers often put `References` after a conclusion/acknowledgement block
    halfway down the page, so scanning only the first page lines misses the
    boundary. At the same time, tables of contents and prose can contain a
    standalone `References` line early in a document. Require a plausible body
    page plus citation-like entries after the heading before truncating.
    """
    if page_num < 2:
        return False
    min_body_pages = 2 if total_pages <= 6 else 3
    if page_num <= min_body_pages:
        following_score = _reference_entry_score(lines[pos + 1: pos + 36])
        return following_score >= 4

    following = lines[pos + 1: pos + 46]
    following_score = _reference_entry_score(following)
    if following_score >= 2:
        return True

    # Late-document headings are usually enough when there is at least one
    # citation-like line after them; this catches pages where PyMuPDF breaks
    # reference entries across many short lines.
    late_page = page_num >= max(4, int(total_pages * 0.35))
    if late_page and following_score >= 1:
        return True

    return False


def _reference_entry_score(lines: list[str]) -> int:
    score = 0
    for line in lines:
        if not line:
            continue
        if _REFERENCE_SECTION_RE.match(line):
            continue
        if _REFERENCE_ENTRY_RE.search(line):
            score += 1
            continue
        # Bibliographies often split author lists and titles over lines; commas
        # plus initials are a useful weak signal after a reference heading.
        if len(line) >= 24 and line.count(",") >= 2 and re.search(r"\b[A-Z]\.", line):
            score += 1
    return score


def _identity_metadata_snippets(
    page_texts: list[str],
    page_window: dict[str, Any],
    *,
    max_chars: int = 6000,
) -> str:
    """Return author/team/affiliation text outside the poster content window.

    The structure extractor gets body text for scientific content. Identity is
    supplied separately so author/team lines survive main-body truncation but
    references/appendix prose does not become poster claim material.
    """
    if not page_texts or max_chars <= 0:
        return ""
    body_pages = int(page_window.get("body_page_count") or len(page_texts))
    blocks: list[str] = []
    used = 0
    for page_num, text in enumerate(page_texts, start=1):
        lines = [re.sub(r"\s+", " ", line).strip() for line in text.splitlines()]
        nonempty = [line for line in lines if line]
        if not nonempty:
            continue
        if page_num <= 2:
            excerpt = _front_matter_identity_excerpt(nonempty)
        else:
            if page_num <= body_pages and not _page_has_identity_hint(nonempty):
                continue
            if _page_reference_dominated(nonempty) and not _page_has_strong_identity_header(nonempty):
                continue
            excerpt = _identity_hint_excerpt(nonempty)
        if not excerpt:
            continue
        block = f"[IDENTITY PAGE {page_num}]\n{excerpt.strip()}\n"
        remaining = max_chars - used
        if remaining <= 0:
            break
        if len(block) > remaining:
            block = block[:remaining].rstrip() + "\n[identity snippets truncated]\n"
        blocks.append(block)
        used += len(block)
        if used >= max_chars:
            break
    return "\n".join(blocks).strip()


def _front_matter_identity_excerpt(lines: list[str], *, max_chars: int = 2600) -> str:
    selected: list[str] = []
    used = 0
    for line in lines[:100]:
        low = line.lower().strip()
        if selected and re.match(r"^(?:abstract|keywords?|1\.?\s+introduction|introduction)\b", low):
            break
        if _REFERENCE_SECTION_RE.match(line):
            break
        if used + len(line) + 1 > max_chars:
            break
        selected.append(line)
        used += len(line) + 1
    return "\n".join(selected)


def _identity_hint_excerpt(lines: list[str], *, context_lines: int = 4) -> str:
    selected: list[str] = []
    seen: set[str] = set()
    for pos, line in enumerate(lines):
        if not (_IDENTITY_HINT_RE.search(line) or "@" in line):
            continue
        start = max(0, pos - context_lines)
        end = min(len(lines), pos + context_lines + 1)
        for candidate in lines[start:end]:
            if _REFERENCE_SECTION_RE.match(candidate):
                continue
            key = candidate.lower()
            if key in seen:
                continue
            seen.add(key)
            selected.append(candidate)
    return "\n".join(selected[:80])


def _page_has_identity_hint(lines: list[str]) -> bool:
    return any(_IDENTITY_HINT_RE.search(line) or "@" in line for line in lines[:120])


def _page_has_strong_identity_header(lines: list[str]) -> bool:
    return any(
        re.search(r"\b(authors?|affiliations?|team)\b|团队|作者|单位|机构", line, re.IGNORECASE)
        for line in lines[:120]
    )


def _page_reference_dominated(lines: list[str]) -> bool:
    sample = lines[:140]
    score = _reference_entry_score(sample)
    return score >= 8 and score >= max(4, len(sample) // 8)


def _ingest_candidate_cap(
    env_name: str,
    cheap_default: int,
    *,
    ctx: ToolContext | None = None,
    dogfood_default: int | None = None,
) -> int:
    explicit = _int_env(env_name, 0)
    if explicit > 0:
        return explicit
    mode = _poster_harness_mode(ctx)
    if mode == "cheap":
        return cheap_default
    if mode == "dogfood" and dogfood_default is not None:
        return dogfood_default
    return 0


def _table_parse_timeout_s(ctx: ToolContext) -> float:
    default = float(getattr(ctx.settings, "ingest_http_timeout", 600.0) or 600.0)
    explicit = _float_env("INGEST_TABLE_PARSE_TIMEOUT_S", 0.0)
    if explicit > 0:
        return explicit
    if _is_harness_mode(ctx):
        return min(default, _HARNESS_TABLE_PARSE_TIMEOUT_S)
    return default


def _table_parse_max_retries(ctx: ToolContext | None = None) -> int | None:
    explicit = _int_env("INGEST_TABLE_PARSE_MAX_RETRIES", -1)
    if explicit >= 0:
        return explicit
    if _is_harness_mode(ctx):
        return _HARNESS_TABLE_PARSE_MAX_RETRIES
    return None


_OCR_PROMPT = """\
You are an OCR engine. The image is ONE page from a PDF. Extract every
readable word exactly as it appears, preserving reading order, newlines
between paragraphs, and bullets / numbered lists. Ignore page numbers,
running headers, and watermarks unless they carry content.

Output **a single fenced JSON code block, nothing else**:

```json
{"text": "<extracted text with \\n newlines>"}
```

Rules:
- Use `\\n\\n` between paragraphs, `\\n` between lines inside a paragraph.
- Do NOT translate. Preserve the original language.
- Do NOT invent words the image doesn't show.
- If the page is blank or unreadable, return `{"text": ""}`.
"""


_INGEST_STRUCTURE_PROMPT = """\
You are a document-structure extractor for AutoDesign. You will be
given extracted text of the paper body before Appendix/Supplement/
References/Bibliography (pages are marked with [PAGE N] headers).
An optional cover-page image may be included only for title/logo
grounding; do not rely on it for figure/table evidence. Optional
identity/metadata snippets may also be included for title, authors,
team, venue, and affiliation recovery only. Return a STRICT JSON
manifest that downstream tools consume verbatim.

Output **a single fenced JSON code block, nothing else**:

```json
{
  "title": "<paper/doc title>",
  "authors": ["<Author Name>", ...],
  "affiliations": ["<institution / lab / company affiliation>", ...],
  "venue": "<conference / publication / null>",
  "abstract": "<2-4 sentence abstract>",
  "sections": [
    {"idx": 1, "heading": "Introduction",
     "summary": "<2-3 sentences of what this section argues>",
     "key_points": ["<one-liner>", "<one-liner>", ...]},
    ...
  ],
  "figures": [
    {"caption": "<figure's full caption text, as it appears in the doc>",
     "page": <1-indexed page where the figure is anchored>,
     "description": "<1 sentence describing what's in the figure; used later to match it to a crop>"},
    ...
  ],
  "tables": [
    {"caption": "<table caption>", "page": <int>},
    ...
  ],
  "key_quotes": ["<memorable line from the doc>", ...]
}
```

Rules:
- Titles: use the human-facing title, not the first line.
- Sections: include each top-level heading (or the logical equivalent
  if the doc doesn't have explicit headings). Max ~10 sections; if the
  doc has more, collapse aggressively (group related subsections).
- Figures: only include REAL visual figures (diagrams, charts,
  screenshots, photos). Ignore logos, page headers, decorative borders,
  page numbers, watermarks, and inline equation renders. Captions as
  they literally appear in the doc. Include sub-panels only if they
  have their own caption line.
- Ignore references/bibliography content. It is citation metadata, not
  poster content.
- Use identity/metadata snippets only for title, authors, venue, and
  affiliations. Do not turn identity snippets into sections, figures,
  tables, key_quotes, or poster claims.
- Pages: 1-indexed.
- Empty lists are fine. Don't guess.
- No extra prose outside the fenced JSON block.
"""


_TABLE_PARSE_PROMPT = """\
You are a table parser. You will see ONE image cropped from a PDF,
plus the raw cell text that pymupdf's table finder guessed for the
same region (the cell splits may be WRONG — trust the image). You
will also see a short list of table-caption candidates pulled from
the same paper so you can match the table to its caption.

Your job:

1. Decide whether the image is actually a data table — a grid of
   rows/columns with comparable values. Diagrams, figure panels,
   math-equation arrays, OCR example screenshots, text paragraphs,
   and decorative layout artifacts are NOT data tables.

2. If it is a table, output clean structured data: pick the header
   row, expand merged cells so every (row, col) has a value, and
   preserve numeric values as-is (no rounding, no re-formatting).
   Short dash / em-dash entries (—, -) stay as the literal string "—".

3. Match it to one of the caption candidates if any fits; otherwise
   return `matched_idx=null` and set a short `title` you extracted
   from the image.

Output **a single fenced JSON code block, nothing else**:

```json
{
  "is_table": <true | false>,
  "matched_idx": <int index into the caption candidate list, or null>,
  "title": "<short title or caption; empty string when unknown>",
  "short_title": "<≤15 chars, ≤3 words; empty if not a table>",
  "headers": ["<col1>", "<col2>", ...],
  "rows": [
    ["<r1c1>", "<r1c2>", ...],
    ["<r2c1>", "<r2c2>", ...],
    ...
  ],
  "col_highlight_rule": ["", "max", "max", "min", ...],
  "reason": "<short explanation>"
}
```

Rules:
- `short_title` is a terse label (≤ 15 chars, ≤ 3 words) for tight
  bboxes where the full `title` / caption won't fit. Summarize the
  table's SUBJECT, not its number: e.g. "Benchmarks", "Ablations",
  "消融实验". Match paper language. Empty string when `is_table=false`.
- Every row in `rows` must have the same length as `headers` (pad
  with "—" if necessary). If you are not confident about the header
  row, leave `headers: []` and put everything in `rows` (first row
  will be treated as header downstream).
- **Two-row headers**: if the table has a parent header spanning
  sub-columns (e.g. "Understanding" over MMMU / MathVista / etc.),
  flatten into a single row using "Parent / Child" format (e.g.
  "Understanding / MMMU"). Do NOT emit a two-row header.
- Do NOT invent rows/columns not visible in the image.
- If `is_table=false`, set `headers` and `rows` to empty lists.
- `col_highlight_rule`: same length as `headers`. For each column,
  emit `"max"` when higher values are better (accuracy, F1, win
  rate), `"min"` when lower is better (loss, error rate, latency),
  or `""` for label / non-numeric / ambiguous columns. The
  downstream renderer bolds the winning row per column — so emitting
  this honestly for benchmark tables is high leverage.
"""


_CAPTION_MATCH_PROMPT = """\
You are a figure↔caption matcher. You will see ONE image cropped from
a PDF plus a short list of figure-caption candidates pulled from the
same paper. Pick the caption that belongs to the image, or report that
the image is not a real figure.

Output **a single fenced JSON code block, nothing else**:

```json
{
  "matched_idx": <int index into the candidate list, or null>,
  "confidence": <float 0.0–1.0>,
  "is_real_figure": <true | false>,
  "short_caption": "<≤15 chars, ≤3 words; empty if not a real figure>",
  "sub_panels": [],
  "reason": "<short explanation>"
}
```

Rules:
- `matched_idx` indexes the candidate list the user gives you (0-based).
- If none of the captions match OR the image is a logo, page header,
  publisher mark, decorative border, watermark, or equation render
  (not a real figure), set `matched_idx=null` and
  `is_real_figure=false`.
- If the image is a real figure but its caption isn't in the
  candidate list (the list was truncated), set `matched_idx=null` and
  `is_real_figure=true`. Downstream will keep it with an empty caption.
- Prefer confidence ≥ 0.7 when you're sure; otherwise be honest and
  use a lower value.
- `short_caption` is a terse label (≤ 15 chars, ≤ 3 words — think
  poster footer where space is tight: "Architecture", "Scaling
  curves", "Ablation"). It summarizes what the image SHOWS, NOT the
  figure number. Match the paper's language (EN/中文). Empty string
  when `is_real_figure=false`.
- `sub_panels` is deprecated for AutoDesign ingest. Always return `[]`.
  Downstream posters use complete source figures and complete tables,
  not partial crops of a chart, label, waveform, example, or sub-panel.
"""


def ingest_document(args: dict[str, Any], *, ctx: ToolContext) -> ToolResultRecord:
    if ctx.state.get("paper_source_sanity_required") is True:
        sanity_paths = _paper_source_paths_for_preflight(args, ctx)
        if not sanity_paths:
            return obs_error(
                "paper source sanity could not verify the original PDF for this reused ingest run",
                category="validation",
                payload={
                    "issue_id": "paper_source_sanity_unverifiable",
                    "repair_route": "replace_paper_source",
                },
            )
        for paper_path in sanity_paths:
            try:
                assert_valid_paper_source_pdf(paper_path)
            except (PaperSourceSanityError, PaperSourceInputError) as e:
                return obs_error(
                    str(e),
                    category="validation",
                    payload={
                        "issue_id": e.issue_id,
                        "repair_route": e.repair_route,
                        "source_file": str(paper_path),
                        "paper_source_sanity": e.report,
                    },
                )
    reused = _reuse_ingest_run_if_requested(args, ctx)
    if reused is not None:
        return reused

    raw = args.get("file_paths")
    if not raw or not isinstance(raw, list):
        return obs_error(
            "ingest_document needs 'file_paths': list[str]",
            category="validation",
        )

    summaries: list[dict[str, Any]] = []
    new_summaries: list[dict[str, Any]] = []

    for fp_str in raw:
        fp = Path(str(fp_str)).expanduser()
        if not fp.is_absolute():
            fp = fp.resolve()
        if not fp.exists():
            return obs_error(f"file not found: {fp}", category="not_found")
        if not fp.is_file():
            return obs_error(f"not a regular file: {fp}", category="validation")

        ext = fp.suffix.lower()
        log("ingest.start", file=str(fp), ext=ext, bytes=fp.stat().st_size)

        if ext == ".pdf" and ctx.state.get("paper_source_sanity_required") is True:
            try:
                assert_valid_paper_source_pdf(fp)
            except (PaperSourceSanityError, PaperSourceInputError) as e:
                log(
                    "ingest.paper_source_sanity.rejected",
                    file=str(fp),
                    issue_id=e.issue_id,
                    repair_route=e.repair_route,
                )
                return obs_error(
                    str(e),
                    category="validation",
                    payload={
                        "issue_id": e.issue_id,
                        "repair_route": e.repair_route,
                        "source_file": str(fp),
                        "paper_source_sanity": e.report,
                    },
                )

        cached = _cached_ingest_summary(fp, ctx.state.get("ingested") or [])
        cached_from_state = cached is not None
        if cached is None and ext == ".pdf":
            cached = _cached_pdf_ingest_summary(fp, ctx)
        if cached is not None:
            if ext == ".pdf":
                cached = _ensure_pdf_summary_paper_memory(fp, cached, ctx)
            summaries.append(cached)
            if not cached_from_state:
                new_summaries.append(cached)
            log(
                "ingest.cache_hit",
                file=str(fp),
                type=cached.get("type"),
                registered=len(cached.get("registered_layer_ids") or []),
            )
            continue

        try:
            if ext == ".pdf":
                s = _ingest_pdf(fp, ctx)
            elif ext == ".docx":
                s = _ingest_docx(fp, ctx)
            elif ext == ".pptx":
                s = _ingest_pptx(fp, ctx)
            elif ext in (".md", ".markdown", ".txt"):
                s = _ingest_markdown(fp, ctx)
            elif ext in (".png", ".jpg", ".jpeg", ".webp"):
                s = _ingest_image(fp, ctx)
            else:
                return obs_error(
                    f"unsupported file type {ext!r}; supported: "
                    ".pdf, .docx, .pptx, .md/.markdown/.txt, "
                    ".png/.jpg/.jpeg/.webp",
                    category="unsupported_format",
                )
        except ScannedPdfError as e:
            return obs_error(f"ingest failed on {fp.name}: {e}", category="parse_error")
        except RuntimeError as e:
            return obs_error(f"ingest failed on {fp.name}: {e}", category="parse_error")

        summaries.append(s)
        new_summaries.append(s)
        if ext == ".pdf":
            _write_pdf_ingest_cache(fp, ctx, s)

    ctx.state.setdefault("ingested", []).extend(new_summaries)
    log("ingest.done", files=len(summaries), new_files=len(new_summaries),
        total_figures=sum(len(s.get("registered_layer_ids", [])) for s in summaries))

    # Build a structured payload with per-figure metadata. The policy needs
    # this to pick figures meaningfully (caption + dims + source page = the
    # actual environment state). NOT prose; the policy can iterate it.
    rendered_state = ctx.state.get("rendered_layers")
    if not isinstance(rendered_state, dict):
        rendered_state = {}
        ctx.state["rendered_layers"] = rendered_state
    rendered = rendered_state
    prior_canvas_plan = ctx.state.get("canvas_plan")
    refined_canvas_plan = refine_canvas_plan_from_ingest(
        prior_canvas_plan if isinstance(prior_canvas_plan, dict) else None,
        summaries,
        rendered,
        brief=str(ctx.state.get("run_brief") or ""),
    )
    if refined_canvas_plan is not None and refined_canvas_plan != prior_canvas_plan:
        ctx.state["canvas_plan"] = refined_canvas_plan
        atomic_write_json(ctx.run_dir / "canvas_plan.json", refined_canvas_plan)
        ingest_shape = refined_canvas_plan.get("ingest_shape") or {}
        log(
            "canvas_plan.refined",
            preset_id=refined_canvas_plan.get("preset_id"),
            lock_level=refined_canvas_plan.get("lock_level"),
            source=refined_canvas_plan.get("source"),
            n_visuals=ingest_shape.get("n_visuals"),
            wide_visuals=ingest_shape.get("wide_visuals"),
            motion_or_sequence_visuals=ingest_shape.get("motion_or_sequence_visuals"),
        )
    files_payload: list[dict[str, Any]] = []
    figures_payload: list[dict[str, Any]] = []
    tables_payload: list[dict[str, Any]] = []
    contact_sheets: list[str] = []
    recommended_figures_payload: dict[str, list[str]] = {}
    visual_candidate_scores_payload: list[dict[str, Any]] = []
    recommended_text_units_payload: dict[str, list[dict[str, Any]]] = (
        _empty_text_unit_buckets()
    )
    figure_catalog_summaries: list[dict[str, Any]] = []
    paper_visual_layer_ids: list[str] = []
    paper_memories: list[dict[str, Any]] = []

    for s in summaries:
        f = Path(s["file"]).name
        t = s["type"]
        provenance_figure_ids = [
            str(v) for v in list(s.get("registered_figure_ids") or s.get("registered_layer_ids") or [])
            if str(v or "").strip()
        ]
        figure_ids = [
            str(v) for v in provenance_figure_ids
            if str(v or "").strip()
            and not _hide_from_planner_visual_catalog(str(v), rendered.get(str(v)) or {})
        ]
        provenance_table_ids = [
            str(v) for v in list(s.get("registered_table_ids") or [])
            if str(v or "").strip()
        ]
        table_ids = [
            str(v) for v in provenance_table_ids
            if str(v or "").strip()
            and not _hide_from_planner_visual_catalog(str(v), rendered.get(str(v)) or {})
        ]
        paper_visual_layer_ids.extend(str(v) for v in list(provenance_figure_ids) + list(provenance_table_ids))

        file_entry: dict[str, Any] = {
            "name": f,
            "type": t,
            "n_figures": len(figure_ids),
            "n_tables": len(table_ids),
        }
        if t in ("pdf", "docx", "pptx"):
            m = s.get("manifest") or {}
            file_entry["title"] = m.get("title")
            file_entry["n_sections"] = len(m.get("sections") or [])
            if t == "pdf":
                file_entry["authors"] = list(m.get("authors") or [])
                if isinstance(m.get("_page_window"), dict):
                    file_entry["body_window"] = m.get("_page_window")
                if s.get("contact_sheet_path"):
                    file_entry["contact_sheet_path"] = s.get("contact_sheet_path")
                    contact_sheets.append(str(s.get("contact_sheet_path")))
                if isinstance(s.get("recommended_figures"), dict):
                    _merge_recommended_visual_buckets(
                        recommended_figures_payload,
                        s.get("recommended_figures") or {},
                    )
                if isinstance(s.get("figure_catalog_summary"), dict):
                    file_entry["figure_catalog_summary"] = s.get("figure_catalog_summary")
                    figure_catalog_summaries.append(s.get("figure_catalog_summary"))
                if isinstance(s.get("visual_candidate_scores"), list):
                    visual_candidate_scores_payload.extend(
                        list(s.get("visual_candidate_scores") or [])
                    )
                if isinstance(s.get("recommended_text_units"), dict):
                    file_entry["recommended_text_units"] = s.get("recommended_text_units")
                    _merge_text_unit_buckets(
                        recommended_text_units_payload,
                        s.get("recommended_text_units") or {},
                    )
                if isinstance(s.get("paper_memory"), dict):
                    file_entry["paper_memory"] = {
                        "cache_key": s["paper_memory"].get("cache_key"),
                        "chunk_count": s["paper_memory"].get("chunk_count"),
                        "categories": s["paper_memory"].get("categories") or {},
                    }
                    paper_memories.append(s["paper_memory"])
            elif t in ("docx", "pptx"):
                visual_ids = list(figure_ids) + list(table_ids)
                ranked_visual_ids: list[str] = []
                if visual_ids:
                    visual_scores = _annotate_visual_curation(visual_ids, rendered)
                    visual_candidate_scores_payload.extend(visual_scores)
                    ranked_visual_ids = [
                        str(item.get("layer_id") or "")
                        for item in visual_scores
                        if str(item.get("layer_id") or "")
                    ]
                    non_pdf_recommendations = _recommend_paper_visuals(
                        ranked_visual_ids,
                        rendered,
                    )
                    _merge_recommended_visual_buckets(
                        recommended_figures_payload,
                        non_pdf_recommendations,
                    )
                    figure_catalog_summary = _figure_catalog_summary(
                        ranked_visual_ids,
                        rendered,
                    )
                    file_entry["figure_catalog_summary"] = figure_catalog_summary
                    figure_catalog_summaries.append(figure_catalog_summary)
                text_units = _recommend_paper_text_units(
                    m,
                    rendered,
                    ranked_visual_ids,
                    ctx.state.get("claim_graph"),
                )
                file_entry["recommended_text_units"] = text_units
                _merge_text_unit_buckets(recommended_text_units_payload, text_units)
        elif t == "markdown":
            file_entry["n_chars"] = s.get("n_chars")
        elif t == "image":
            file_entry["width"] = s.get("width")
            file_entry["height"] = s.get("height")
        files_payload.append(file_entry)

        # Show every ranked source visual to the designer. References are still
        # filtered at the PDF page-window level, but the main-body catalog
        # should not be trimmed to save prompt tokens.
        ranked = _rank_figure_ids_for_planner(figure_ids, rendered)
        file_entry["n_figures_shown"] = len(ranked)

        for fid in ranked:
            rec = rendered.get(fid) or {}
            figures_payload.append({
                "layer_id": fid,
                "source_file": f,
                "source_page": rec.get("source_page"),
                "source_ref": rec.get("source_ref"),
                "image_size": rec.get("image_size"),
                "caption": rec.get("caption"),
                "caption_short": rec.get("caption_short"),
                "sha256": rec.get("sha256"),
                "source_bbox_pdf_points": rec.get("source_bbox_pdf_points"),
                "source_pdf_sha256": rec.get("source_pdf_sha256"),
            } | _curation_payload(rec))
        for tid in table_ids:
            rec = rendered.get(tid) or {}
            tables_payload.append({
                "layer_id": tid,
                "source_file": f,
                "source_page": rec.get("source_page"),
                "n_rows": len(rec.get("rows") or []),
                "n_cols": len(rec.get("headers") or []) or (
                    len((rec.get("rows") or [[]])[0])
                ),
                "caption": rec.get("caption") or rec.get("title"),
                "caption_short": rec.get("caption_short"),
                "source_bbox_pdf_points": rec.get("source_bbox_pdf_points"),
                "source_pdf_sha256": rec.get("source_pdf_sha256"),
            } | _curation_payload(rec))

    visual_candidate_scores_payload = _filter_visual_candidate_scores_for_planner(
        visual_candidate_scores_payload,
        rendered,
    )
    recommended_figures_payload = _filter_recommended_visual_buckets_for_planner(
        recommended_figures_payload,
        rendered,
        visual_candidate_scores_payload,
    )

    paper_visual_provenance = _build_paper_visual_provenance(
        layer_ids=paper_visual_layer_ids,
        rendered=rendered,
        run_dir=ctx.run_dir,
        summaries=summaries,
    )
    if paper_visual_provenance.get("assets"):
        ctx.state["paper_visual_provenance"] = paper_visual_provenance
        atomic_write_json(ctx.run_dir / "paper_visual_provenance.json", paper_visual_provenance)
        log(
            "paper_visual_provenance.created",
            assets=len(paper_visual_provenance.get("assets") or []),
            source_docs=len(paper_visual_provenance.get("source_documents") or []),
        )
    paper_memory = merge_paper_memories(paper_memories)
    if paper_memory:
        ctx.state["paper_memory"] = paper_memory
        write_paper_memory_run_artifacts(ctx.run_dir, paper_memory)
        log(
            "paper_memory.run_payload",
            chunks=paper_memory.get("chunk_count"),
            documents=len(paper_memory.get("documents") or []) or 1,
        )
    paper_manifest = {}
    for summary in summaries:
        if isinstance(summary, dict) and summary.get("type") == "pdf":
            paper_manifest = summary.get("manifest") or {}
            break
    paper_memory_dossier = _ensure_paper_memory_dossier(
        ctx=ctx,
        paper_memory=paper_memory,
        paper_manifest=paper_manifest,
        paper_visual_provenance=paper_visual_provenance,
        recommended_text_units=recommended_text_units_payload,
        recommended_figures=recommended_figures_payload,
    )
    if paper_memory_dossier:
        _raise_if_ingest_cancelled(ctx, "ingest.paper_memory_dossier.before_state_merge")
        ctx.state["paper_memory_dossier"] = paper_memory_dossier
        _raise_if_ingest_cancelled(ctx, "ingest.paper_memory_dossier.before_run_write")
        write_paper_memory_dossier_run_artifacts(ctx.run_dir, paper_memory_dossier)
        _raise_if_ingest_cancelled(ctx, "ingest.paper_memory_dossier.after_run_write")
    paper_visual_storyboard = build_paper_visual_storyboard(
        manifest=paper_manifest,
        recommended_text_units=recommended_text_units_payload,
        recommended_figures=recommended_figures_payload,
        visual_candidate_scores=visual_candidate_scores_payload,
        paper_visual_provenance=paper_visual_provenance,
        canvas_plan=ctx.state.get("canvas_plan") if isinstance(ctx.state.get("canvas_plan"), dict) else None,
    )
    if paper_visual_storyboard:
        ctx.state["paper_visual_storyboard"] = paper_visual_storyboard
        atomic_write_json(ctx.run_dir / "paper_visual_storyboard.json", paper_visual_storyboard)
        log(
            "paper_visual_storyboard.created",
            selected=len(paper_visual_storyboard.get("selected_assets") or []),
            target=paper_visual_storyboard.get("target_visual_count"),
        )
    paper_resources: dict[str, Any] = {}
    if should_auto_discover_paper_resources(ctx) and paper_manifest:
        try:
            paper_resources = discover_paper_resources_for_context(
                ctx,
                title=str(paper_manifest.get("title") or "") or None,
                authors=list(paper_manifest.get("authors") or []),
                max_results=12,
                search_web=True,
            )
        except Exception as e:  # noqa: BLE001 - resource search must not break ingest
            paper_resources = {
                "kind": "paper_resource_manifest",
                "version": 1,
                "resources": [],
                "resource_chips": [],
                "warnings": [f"{type(e).__name__}: {e}"],
            }
            ctx.state["paper_resources"] = paper_resources
            atomic_write_json(ctx.run_dir / "paper_resource_manifest.json", paper_resources)
            log("paper_resources.discover_failed", error=f"{type(e).__name__}: {e}")

    prior_deck_plan = ctx.state.get("deck_plan")
    raw_brief = str(ctx.state.get("raw_user_brief") or ctx.state.get("run_brief") or "")
    if should_refine_deck_plan(
        prior_deck_plan if isinstance(prior_deck_plan, dict) else None,
        raw_brief=raw_brief,
    ):
        if getattr(ctx.settings, "enable_deck_outline", True):
            try:
                _raise_if_ingest_cancelled(ctx, "ingest.deck_outline.before_plan")
                refined_deck_plan = DeckOutlineAgent(ctx.settings).plan(
                    raw_brief=raw_brief,
                    enhanced_brief=str(ctx.state.get("run_brief") or ""),
                    base_plan=prior_deck_plan if isinstance(prior_deck_plan, dict) else None,
                    summaries=summaries,
                    rendered_layers=rendered,
                    figures_payload=figures_payload,
                    tables_payload=tables_payload,
                    claim_graph=ctx.state.get("claim_graph"),
                    cancellation_token=_ingest_cancellation_token(ctx),
                )
                _raise_if_ingest_cancelled(ctx, "ingest.deck_outline.after_plan")
            except Exception as e:  # noqa: BLE001
                log(
                    "deck_outline.failed",
                    error=f"{type(e).__name__}: {e}",
                    fallback="deterministic",
                )
                refined_deck_plan = fallback_deck_plan(
                    prior_deck_plan if isinstance(prior_deck_plan, dict) else None,
                    raw_brief=raw_brief,
                    summaries=summaries,
                    rendered_layers=rendered,
                    claim_graph=ctx.state.get("claim_graph"),
                    reason=f"deck outline agent failed: {type(e).__name__}",
                )
        else:
            refined_deck_plan = fallback_deck_plan(
                prior_deck_plan if isinstance(prior_deck_plan, dict) else None,
                raw_brief=raw_brief,
                summaries=summaries,
                rendered_layers=rendered,
                claim_graph=ctx.state.get("claim_graph"),
                reason="deck outline agent disabled",
            )
        if refined_deck_plan and refined_deck_plan != prior_deck_plan:
            _raise_if_ingest_cancelled(ctx, "ingest.deck_outline.before_state_merge")
            ctx.state["deck_plan"] = refined_deck_plan
            _raise_if_ingest_cancelled(ctx, "ingest.deck_outline.before_plan_write")
            atomic_write_json(ctx.run_dir / "deck_plan.json", refined_deck_plan)
            _raise_if_ingest_cancelled(ctx, "ingest.deck_outline.after_plan_write")
            _raise_if_ingest_cancelled(ctx, "ingest.deck_outline.before_refined_log")
            log(
                "deck_plan.refined",
                source=refined_deck_plan.get("source"),
                status=refined_deck_plan.get("status"),
                slide_count=refined_deck_plan.get("slide_count"),
                lock_level=refined_deck_plan.get("lock_level"),
            )
            _raise_if_ingest_cancelled(ctx, "ingest.deck_outline.after_refined_log")

    poster_content_brief = _build_poster_content_brief(
        summaries=summaries,
        rendered=rendered,
        recommended_figures=recommended_figures_payload,
        recommended_text_units=recommended_text_units_payload,
        visual_candidate_scores=visual_candidate_scores_payload,
        paper_visual_provenance=paper_visual_provenance,
        paper_visual_storyboard=paper_visual_storyboard,
        canvas_plan=ctx.state.get("canvas_plan"),
        raw_brief=str(ctx.state.get("raw_user_brief") or ctx.state.get("run_brief") or ""),
        required_color_system=ctx.state.get("required_color_system"),
    )
    if poster_content_brief:
        poster_content_brief = _sanitize_poster_content_brief_visual_eligibility(
            poster_content_brief,
            rendered,
        )
        ctx.state["poster_content_brief"] = poster_content_brief
        atomic_write_json(ctx.run_dir / "poster_content_brief.json", poster_content_brief)
        log(
            "poster_content_brief.created",
            sections=len(poster_content_brief.get("sections") or []),
            target_visuals=(
                poster_content_brief.get("visual_selection") or {}
            ).get("target_visual_count"),
            selected_visuals=len(
                (poster_content_brief.get("visual_selection") or {}).get("primary_visual_ids") or []
            ),
        )
    poster_plan_contract = build_poster_plan_contract(
        poster_content_brief,
        canvas_plan=ctx.state.get("canvas_plan"),
        rendered_layers=rendered,
    )
    if poster_plan_contract:
        ctx.state["poster_plan_contract"] = poster_plan_contract
        atomic_write_json(ctx.run_dir / "poster_plan_contract.json", poster_plan_contract)
        log(
            "poster_plan_contract.created",
            selected_visuals=len(poster_plan_contract.get("selected_visuals") or []),
            required_sections=len(poster_plan_contract.get("required_sections") or []),
            layout_archetype=poster_plan_contract.get("layout_archetype"),
        )
    poster_contract_preflight: dict[str, Any] = {}
    if poster_plan_contract:
        poster_contract_preflight = preflight_poster_plan_contract(
            poster_plan_contract,
            poster_content_brief,
            rendered_layers=rendered,
            canvas_plan=ctx.state.get("canvas_plan"),
        )
        ctx.state["poster_contract_preflight"] = poster_contract_preflight
        atomic_write_json(ctx.run_dir / "poster_contract_preflight.json", poster_contract_preflight)
        log(
            "poster_contract_preflight.done",
            status=poster_contract_preflight.get("status"),
            findings=len(poster_contract_preflight.get("findings") or []),
            p0=sum(
                1 for finding in (poster_contract_preflight.get("findings") or [])
                if str(finding.get("severity")).upper() == "P0"
            ),
        )

    return obs_ok({
        "n_files": len(summaries),
        "n_figures": len(figures_payload),
        "n_tables": len(tables_payload),
        "files": files_payload,
        "figures": figures_payload,
        "tables": tables_payload,
        "contact_sheet_path": contact_sheets[0] if contact_sheets else None,
        "contact_sheets": contact_sheets,
        "recommended_figures": recommended_figures_payload,
        "visual_candidate_scores": visual_candidate_scores_payload,
        "paper_visual_provenance": _compact_paper_visual_provenance_for_planner(
            paper_visual_provenance,
            paper_visual_storyboard=paper_visual_storyboard,
        ),
        "paper_visual_storyboard": _compact_paper_visual_storyboard_for_planner(
            paper_visual_storyboard
        ),
        "recommended_text_units": recommended_text_units_payload,
        "paper_memory": compact_paper_memory_for_planner(paper_memory),
        "paper_memory_dossier": _compact_paper_memory_dossier_for_planner(
            paper_memory_dossier
        ),
        "paper_resources": _compact_paper_resources_for_planner(paper_resources),
        "figure_catalog_summary": (
            figure_catalog_summaries[0]
            if len(figure_catalog_summaries) == 1
            else {"files": figure_catalog_summaries}
        ),
        "poster_content_brief": poster_content_brief,
        "poster_plan_contract": poster_plan_contract,
        "poster_contract_preflight": poster_contract_preflight,
        "canvas_plan": ctx.state.get("canvas_plan"),
        "deck_plan": ctx.state.get("deck_plan"),
    })


def _load_ingest_state_from_dir(
    ctx: ToolContext,
    source_dir: Path,
    *,
    reused_canvas_plan_override: dict[str, Any] | None = None,
) -> dict[str, Any] | ToolResultRecord:
    """Deserialize an existing run's ingest artifacts into ctx.state.

    Used by both the copy-path (``_reuse_ingest_run_if_requested`` with a
    foreign ``source_dir``) and the in-place-path (``--resume-run`` where
    ``source_dir == ctx.run_dir`` and materials are already on disk).

    Returns a dict of loaded payloads on success, or a ``ToolResultRecord``
    error on validation failure (missing required file, no visual assets).
    """
    required = [
        "paper_visual_provenance.json",
        "paper_memory.json",
        "paper_visual_storyboard.json",
        "poster_content_brief.json",
        "poster_plan_contract.json",
    ]
    missing = [name for name in required if not (source_dir / name).exists()]
    if missing:
        return obs_error(
            f"reuse ingest run is missing required artifacts: {', '.join(missing)}",
            category="validation",
            payload={
                "issue_id": "reuse_ingest_artifacts_missing",
                "source_run_dir": str(source_dir),
                "missing": missing,
            },
        )

    layers_src = source_dir / "layers"
    if not layers_src.exists() or not layers_src.is_dir():
        return obs_error(
            f"reuse ingest run is missing layers directory: {layers_src}",
            category="validation",
            payload={"issue_id": "reuse_ingest_layers_missing"},
        )

    def read_json(name: str) -> dict[str, Any]:
        data = json.loads((source_dir / name).read_text(encoding="utf-8"))
        return data if isinstance(data, dict) else {}

    paper_visual_provenance = read_json("paper_visual_provenance.json")
    paper_memory = read_json("paper_memory.json")
    paper_memory_dossier = (
        read_json("paper_memory_dossier.json")
        if (source_dir / "paper_memory_dossier.json").exists()
        else {}
    )
    paper_visual_storyboard = read_json("paper_visual_storyboard.json")
    poster_content_brief = read_json("poster_content_brief.json")
    poster_plan_contract = read_json("poster_plan_contract.json")
    poster_content_brief = _normalize_reused_poster_content_brief(
        poster_content_brief,
        ctx,
        paper_memory=paper_memory,
        poster_plan_contract=poster_plan_contract,
    )
    poster_content_brief = _apply_storyboard_tiers_to_poster_content_brief(
        poster_content_brief,
        paper_visual_storyboard,
    )
    poster_contract_preflight = (
        read_json("poster_contract_preflight.json")
        if (source_dir / "poster_contract_preflight.json").exists() else {}
    )
    reused_canvas_plan = (
        reused_canvas_plan_override
        if isinstance(reused_canvas_plan_override, dict) and reused_canvas_plan_override
        else ctx.state.get("canvas_plan")
        if isinstance(ctx.state.get("canvas_plan"), dict)
        else read_json("canvas_plan.json") if (source_dir / "canvas_plan.json").exists() else {}
    )

    rendered_layers = _rendered_layers_from_reused_provenance(
        paper_visual_provenance,
        run_dir=ctx.run_dir,
        layers_dir=ctx.layers_dir,
    )
    if not rendered_layers:
        return obs_error(
            "reuse ingest run has no reusable paper visual assets.",
            category="validation",
            payload={
                "issue_id": "reuse_ingest_no_visual_assets",
                "source_run_dir": str(source_dir),
            },
        )

    paper_visual_storyboard = _sanitize_paper_visual_storyboard_for_rendered(
        paper_visual_storyboard,
        rendered_layers,
    )
    poster_content_brief = _apply_storyboard_tiers_to_poster_content_brief(
        poster_content_brief,
        paper_visual_storyboard,
    )
    poster_content_brief = _sanitize_poster_content_brief_visual_eligibility(
        poster_content_brief,
        rendered_layers,
    )

    rebuilt_contract = build_poster_plan_contract(
        poster_content_brief,
        canvas_plan=reused_canvas_plan,
        rendered_layers=rendered_layers,
    )
    if rebuilt_contract:
        poster_plan_contract = rebuilt_contract
        poster_contract_preflight = preflight_poster_plan_contract(
            poster_plan_contract,
            poster_content_brief,
            rendered_layers=rendered_layers,
        )
        log(
            "ingest.reuse_contract_rebuilt",
            selected_visuals=len(poster_plan_contract.get("selected_visuals") or []),
            min_visual_count=(
                (poster_plan_contract.get("density_targets") or {}).get("min_visual_count")
                if isinstance(poster_plan_contract.get("density_targets"), dict) else None
            ),
            target_visual_count=(
                (poster_plan_contract.get("density_targets") or {}).get("target_visual_count")
                if isinstance(poster_plan_contract.get("density_targets"), dict) else None
            ),
        )

    ctx.state.setdefault("ingested", [])
    if reused_canvas_plan:
        ctx.state["canvas_plan"] = reused_canvas_plan
    ctx.state["rendered_layers"] = rendered_layers
    ctx.state["paper_visual_provenance"] = paper_visual_provenance
    ctx.state["paper_memory"] = paper_memory
    ctx.state["paper_memory_dossier"] = paper_memory_dossier
    ctx.state["paper_visual_storyboard"] = paper_visual_storyboard
    ctx.state["poster_content_brief"] = poster_content_brief
    ctx.state["poster_plan_contract"] = poster_plan_contract
    if poster_contract_preflight:
        ctx.state["poster_contract_preflight"] = poster_contract_preflight
    summary = _reused_ingest_summary(
        paper_memory=paper_memory,
        paper_visual_provenance=paper_visual_provenance,
        rendered_layers=rendered_layers,
    )
    if summary:
        ctx.state["ingested"] = [summary]

    return {
        "paper_visual_provenance": paper_visual_provenance,
        "paper_memory": paper_memory,
        "paper_memory_dossier": paper_memory_dossier,
        "paper_visual_storyboard": paper_visual_storyboard,
        "poster_content_brief": poster_content_brief,
        "poster_plan_contract": poster_plan_contract,
        "poster_contract_preflight": poster_contract_preflight,
        "reused_canvas_plan": reused_canvas_plan,
        "rendered_layers": rendered_layers,
        "summary": summary,
    }


def _paper_source_paths_for_preflight(args: dict[str, Any], ctx: ToolContext) -> list[Path]:
    raw_paths = args.get("file_paths")
    if isinstance(raw_paths, list):
        paths = [Path(str(value)).expanduser().resolve() for value in raw_paths]
        current_pdfs = [path for path in paths if path.suffix.lower() == ".pdf" and path.is_file()]
        if current_pdfs:
            return current_pdfs

    reuse_value = str(
        ctx.state.get("reuse_ingest_run") or args.get("reuse_ingest_run") or ""
    ).strip()
    if not reuse_value:
        return []
    source_dir = resolve_run_dir(ctx.settings.out_dir, reuse_value)
    resume_state = source_dir / "resume_state.json"
    if not resume_state.is_file():
        return []
    try:
        payload = json.loads(resume_state.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError):
        return []
    paths = [Path(str(value)).expanduser() for value in (payload.get("attachments") or [])]
    return [path for path in paths if path.suffix.lower() == ".pdf" and path.is_file()]


def _reuse_ingest_run_if_requested(args: dict[str, Any], ctx: ToolContext) -> ToolResultRecord | None:
    raw = (
        ctx.state.get("reuse_ingest_run")
        if isinstance(ctx.state, dict) else None
    ) or args.get("reuse_ingest_run")
    reuse_value = str(raw or "").strip()
    if not reuse_value:
        return None

    source_dir = resolve_run_dir(ctx.settings.out_dir, reuse_value)
    if not source_dir.exists() or not source_dir.is_dir():
        return obs_error(
            f"reuse ingest run not found: {source_dir}",
            category="not_found",
            payload={"issue_id": "reuse_ingest_run_not_found"},
        )

    in_place = source_dir == ctx.run_dir.resolve()
    if not in_place:
        # Copy pass — pull the source run's artifacts into this new run_dir
        # so downstream code (and _stage_inputs per attempt) can rely on them.
        source_layers = source_dir / "layers"
        if not source_layers.exists() or not source_layers.is_dir():
            return obs_error(
                f"reuse ingest run is missing layers directory: {source_layers}",
                category="validation",
                payload={"issue_id": "reuse_ingest_layers_missing"},
            )
        shutil.copytree(source_layers, ctx.layers_dir, dirs_exist_ok=True)
        for name in (
            "paper_visual_provenance.json",
            "paper_memory.json",
            "paper_memory_dossier.json",
            "paper_visual_storyboard.json",
            "poster_content_brief.json",
            "poster_plan_contract.json",
            "poster_contract_preflight.json",
            "canvas_plan.json",
        ):
            src = source_dir / name
            if src.exists():
                shutil.copy2(src, ctx.run_dir / name)
        for name in ("paper_memory.md", "paper_memory_dossier.md"):
            src = source_dir / name
            if src.exists():
                shutil.copy2(src, ctx.run_dir / name)
        evidence_src = source_dir / "paper_evidence_packs"
        if evidence_src.exists() and evidence_src.is_dir():
            shutil.copytree(evidence_src, ctx.run_dir / "paper_evidence_packs", dirs_exist_ok=True)

    loaded = _load_ingest_state_from_dir(ctx, ctx.run_dir if in_place else source_dir)
    if isinstance(loaded, ToolResultRecord):
        return loaded

    # Persist normalized artifacts under the new run root (in-place callers
    # already have them on disk but rewriting keeps normalization + rebuilt
    # contract consistent with what ctx.state now advertises).
    for name, data in (
        ("paper_visual_provenance.json", loaded["paper_visual_provenance"]),
        ("paper_memory.json", loaded["paper_memory"]),
        ("paper_memory_dossier.json", loaded["paper_memory_dossier"]),
        ("paper_visual_storyboard.json", loaded["paper_visual_storyboard"]),
        ("poster_content_brief.json", loaded["poster_content_brief"]),
        ("canvas_plan.json", loaded["reused_canvas_plan"]),
    ):
        if data:
            atomic_write_json(ctx.run_dir / name, data)
    if loaded["poster_plan_contract"]:
        atomic_write_json(ctx.run_dir / "poster_plan_contract.json", loaded["poster_plan_contract"])
    if loaded["poster_contract_preflight"]:
        atomic_write_json(ctx.run_dir / "poster_contract_preflight.json", loaded["poster_contract_preflight"])

    paper_visual_provenance = loaded["paper_visual_provenance"]
    paper_visual_storyboard = loaded["paper_visual_storyboard"]
    paper_memory = loaded["paper_memory"]
    paper_memory_dossier = loaded["paper_memory_dossier"]
    poster_content_brief = loaded["poster_content_brief"]
    poster_plan_contract = loaded["poster_plan_contract"]
    poster_contract_preflight = loaded["poster_contract_preflight"]
    rendered_layers = loaded["rendered_layers"]
    summary = loaded["summary"]

    figures_payload, tables_payload = _reused_figures_tables_payload(rendered_layers)
    contact_sheets = [
        str(path)
        for path in sorted(ctx.layers_dir.glob("ingest_contact_sheet*"))
        if path.is_file()
    ]

    log(
        "ingest.reuse",
        source_run_dir=str(source_dir),
        in_place=in_place,
        assets=len(rendered_layers),
        figures=len(figures_payload),
        tables=len(tables_payload),
    )
    return obs_ok({
        "reused": True,
        "source_run_dir": str(source_dir),
        "n_files": 1 if summary else 0,
        "n_figures": len(figures_payload),
        "n_tables": len(tables_payload),
        "files": [{
            "type": "pdf",
            "path": summary.get("path"),
            "title": ((summary.get("manifest") or {}).get("title") if summary else None),
            "n_figures_shown": len(figures_payload),
        }] if summary else [],
        "figures": figures_payload,
        "tables": tables_payload,
        "contact_sheet_path": contact_sheets[0] if contact_sheets else None,
        "contact_sheets": contact_sheets,
        "recommended_figures": [],
        "visual_candidate_scores": [],
        "paper_visual_provenance": _compact_paper_visual_provenance_for_planner(
            paper_visual_provenance,
            paper_visual_storyboard=paper_visual_storyboard,
        ),
        "paper_visual_storyboard": _compact_paper_visual_storyboard_for_planner(
            paper_visual_storyboard
        ),
        "recommended_text_units": {},
        "paper_memory": compact_paper_memory_for_planner(paper_memory),
        "paper_memory_dossier": _compact_paper_memory_dossier_for_planner(
            paper_memory_dossier
        ),
        "paper_resources": {},
        "figure_catalog_summary": {},
        "poster_content_brief": poster_content_brief,
        "poster_plan_contract": poster_plan_contract,
        "poster_contract_preflight": poster_contract_preflight,
        "canvas_plan": ctx.state.get("canvas_plan"),
        "deck_plan": ctx.state.get("deck_plan"),
    })


def _paper_poster_aesthetic_contract() -> dict[str, str]:
    return {
        "canvas_policy": "white or near-white academic canvas; no tinted full-board backgrounds",
        "palette_usage_policy": (
            "Use the selected palette sparingly. The identity header uses the "
            "fixed white/near-white treatment with a single top accent rule "
            "only."
        ),
        "header_surface_policy": (
            "Use the fixed identity-header style: white/near-white header with "
            "a single top accent rule only. Do not use bottom header rules, "
            "filled title bands, four-sided outlines, or mixed header styles "
            "for new paper posters."
        ),
        "section_surface_policy": (
            "Use compact filled primary section heading bands with white text, "
            "while keeping section bodies and panel interiors white/neutral. Do "
            "not fill panels with secondary tints or wrap each panel in a colored box."
        ),
        "table_surface_policy": (
            "Native tables use white cells and booktabs-like horizontal rules; no "
            "decorative zebra striping or saturated headers."
        ),
        "source_wrapper_policy": "Source figure/table wrappers stay white with transparent borders only; keep the DOM box for measurement but do not show visible wrapper outlines or shadows.",
        "color_dominance_policy": "Most of the poster should read as paper content and whitespace, not color blocks.",
    }


def _complete_academic_color_system_options(
    *,
    raw_brief: str,
    manifest: dict[str, Any],
    recommended_text_units: dict[str, Any] | None = None,
) -> list[dict[str, Any]]:
    ranked = rank_academic_color_system_options(
        raw_brief=raw_brief,
        manifest=manifest,
        recommended_text_units=recommended_text_units,
    )
    out: list[dict[str, Any]] = []
    seen: set[str] = set()
    for option in [*ranked, *academic_color_system_options()]:
        palette_id = str(option.get("palette_id") or "").strip()
        if palette_id and palette_id not in seen:
            out.append(option)
            seen.add(palette_id)
    return out


def _normalize_reused_poster_content_brief(
    poster_content_brief: dict[str, Any],
    ctx: ToolContext,
    *,
    paper_memory: dict[str, Any] | None = None,
    poster_plan_contract: dict[str, Any] | None = None,
) -> dict[str, Any]:
    if not isinstance(poster_content_brief, dict) or not poster_content_brief:
        return {}
    out = dict(poster_content_brief)
    background_contract = out.get("background_contract")
    if isinstance(background_contract, dict) and "logo_policy" in background_contract:
        out["background_contract"] = {
            key: value
            for key, value in background_contract.items()
            if key != "logo_policy"
        }
    raw_brief = str(
        (ctx.state or {}).get("raw_user_brief")
        or (ctx.state or {}).get("run_brief")
        or ""
    )
    metadata = paper_memory.get("metadata") if isinstance(paper_memory, dict) and isinstance(paper_memory.get("metadata"), dict) else {}
    affiliations = [
        str(item).strip()
        for item in (
            out.get("affiliations")
            or out.get("institutions")
            or metadata.get("affiliations")
            or metadata.get("institutions")
            or []
        )
        if str(item or "").strip()
    ][:16]
    if affiliations:
        out["affiliations"] = affiliations
        out["institutions"] = affiliations
    required_color_system = dict((ctx.state or {}).get("required_color_system") or {})
    explicit_color_system = explicit_academic_color_system(raw_brief=raw_brief)
    recommended_color_system = required_color_system or explicit_color_system or active_academic_color_system(
        out,
        poster_plan_contract,
        raw_brief=raw_brief,
        manifest=_reused_color_system_manifest(out, paper_memory or {}),
    )
    out["recommended_color_system"] = recommended_color_system
    out["color_system"] = recommended_color_system
    if required_color_system:
        out["required_color_system"] = required_color_system
    else:
        out.pop("required_color_system", None)
    color_roles = (
        recommended_color_system.get("roles")
        if isinstance(recommended_color_system.get("roles"), dict)
        else {}
    )
    out["background_contract"] = {
        **(
            out.get("background_contract")
            if isinstance(out.get("background_contract"), dict)
            else {}
        ),
        "default": "native white/cream academic canvas",
        "use_generated_background": False,
        "allowed_fill": _take_unique_values([
            "#FFFFFF",
            "#FAFDFB",
            "#FBFBF7",
            color_roles.get("background"),
        ]),
        "structure": "mostly white surfaces, fixed white identity header with a single top accent rule, filled primary section heading bands, thin neutral separators",
    }
    out["background_contract"].pop("logo_policy", None)
    out["aesthetic_contract"] = _paper_poster_aesthetic_contract()
    out["color_system_options"] = _complete_academic_color_system_options(
        raw_brief=raw_brief,
        manifest=_reused_color_system_manifest(out, paper_memory or {}),
    )
    out["institution_color_signals"] = _institution_color_signals(
        affiliations,
        out.get("authors") if isinstance(out.get("authors"), list) else [],
    )
    explicit_profile = _reference_profile_env() or _reference_profile_from_brief(raw_brief)
    if explicit_profile == _RESEARCH_SYNTHESIS_PROFILE:
        return out
    if str(out.get("reference_profile") or "") != _RESEARCH_SYNTHESIS_PROFILE:
        return out

    visual_selection = out.get("visual_selection") if isinstance(out.get("visual_selection"), dict) else {}
    role_buckets = visual_selection.get("role_buckets") if isinstance(visual_selection.get("role_buckets"), dict) else {}
    selected_visuals = [
        str(item)
        for item in (visual_selection.get("primary_visual_ids") or [])
        if str(item or "").strip()
    ]
    normalized_buckets = {
        str(key): [str(item) for item in (value or []) if str(item or "").strip()]
        for key, value in role_buckets.items()
    }
    out["reference_profile"] = _EDITORIAL_FLOW_PROFILE
    out["native_reference_targets"] = _compact_editorial_flow_targets({})
    out["reference_archetype_skeleton"] = _editorial_flow_reference_archetype_skeleton()
    out["panel_plan"] = _editorial_flow_panel_plan(
        recommended_figures=normalized_buckets,
        selected_visuals=selected_visuals,
    )
    out["editorial_column_plan"] = out["panel_plan"]
    return out


def _reused_color_system_manifest(
    poster_content_brief: dict[str, Any],
    paper_memory: dict[str, Any],
) -> dict[str, Any]:
    metadata = paper_memory.get("metadata") if isinstance(paper_memory.get("metadata"), dict) else {}
    manifest = {
        "title": metadata.get("title") or poster_content_brief.get("title") or "",
        "authors": metadata.get("authors") or poster_content_brief.get("authors") or [],
        "affiliations": (
            metadata.get("affiliations")
            or metadata.get("institutions")
            or poster_content_brief.get("affiliations")
            or poster_content_brief.get("institutions")
            or []
        ),
        "venue": metadata.get("venue") or poster_content_brief.get("venue") or "",
        "abstract": metadata.get("abstract") or poster_content_brief.get("abstract") or "",
    }
    keywords = []
    for value in (metadata.get("keywords"), poster_content_brief.get("keywords")):
        if isinstance(value, list):
            keywords.extend(str(item) for item in value if str(item or "").strip())
    if keywords:
        manifest["keywords"] = keywords

    pieces: list[str] = []
    for key in ("title", "venue", "abstract", "reference_profile", "layout_archetype"):
        _append_palette_context(pieces, poster_content_brief.get(key))
    for section in poster_content_brief.get("sections") or []:
        if isinstance(section, dict):
            for key in ("title", "heading", "label", "purpose", "summary", "text"):
                _append_palette_context(pieces, section.get(key))
    for panel in poster_content_brief.get("panel_plan") or []:
        if isinstance(panel, dict):
            for key in ("slot_id", "purpose", "text_budget"):
                _append_palette_context(pieces, panel.get(key))
    if pieces:
        manifest["summary"] = ". ".join(pieces[:80])
    return manifest


def _append_palette_context(pieces: list[str], value: Any) -> None:
    text = str(value or "").strip()
    if text:
        pieces.append(text[:500])


def _brief_selected_source_visual_ids(visual_selection: dict[str, Any]) -> list[str]:
    ids: list[str] = []
    for key in (
        "high_priority_visual_ids",
        "primary_visual_ids",
        "storyboard_primary_asset_ids",
        "storyboard_selected_asset_ids",
    ):
        ids.extend(str(value or "").strip() for value in (visual_selection.get(key) or []))
    if not ids and isinstance(visual_selection.get("visual_records"), list):
        ids.extend(
            str(item.get("layer_id") or item.get("asset_id") or "").strip()
            for item in visual_selection.get("visual_records") or []
            if isinstance(item, dict)
        )
    return _take_unique_values([
        item for item in ids
        if item.startswith(("ingest_fig_", "ingest_table_"))
    ])


def _refresh_brief_visual_unit_metrics(poster_content_brief: dict[str, Any]) -> dict[str, Any]:
    if not isinstance(poster_content_brief, dict) or not poster_content_brief:
        return poster_content_brief
    out = dict(poster_content_brief)
    visual_selection = dict(out.get("visual_selection") or {})
    policy = dict(out.get("source_asset_policy") or {})
    selected_count = len(_brief_selected_source_visual_ids(visual_selection))
    unit_target = _safe_int(
        visual_selection.get("poster_visual_unit_target")
        or policy.get("poster_visual_unit_target"),
        default=0,
    )
    if unit_target <= 0:
        return out
    shortfall = max(0, unit_target - selected_count)
    tasks = (
        visual_selection.get("supplemental_native_visual_tasks")
        or out.get("supplemental_native_visual_tasks")
        or policy.get("supplemental_native_visual_tasks")
        or []
    )
    task_count = len(tasks) if isinstance(tasks, list) else _safe_int(
        visual_selection.get("supplemental_native_visual_task_count")
        or policy.get("supplemental_native_visual_task_count"),
        default=0,
    )
    for target in (visual_selection, policy):
        target["selected_source_visual_count"] = selected_count
        target["poster_visual_unit_target"] = unit_target
        target["source_visual_shortfall"] = shortfall
        target["supplemental_native_visual_task_count"] = task_count
    out["visual_selection"] = visual_selection
    out["source_asset_policy"] = policy
    return out


def _apply_storyboard_tiers_to_poster_content_brief(
    poster_content_brief: dict[str, Any],
    paper_visual_storyboard: dict[str, Any],
) -> dict[str, Any]:
    if not isinstance(poster_content_brief, dict) or not poster_content_brief:
        return {}
    storyboard = _compact_paper_visual_storyboard_for_planner(paper_visual_storyboard)
    if not storyboard:
        return poster_content_brief
    out = dict(poster_content_brief)
    out["visual_storyboard"] = storyboard
    visual_selection = dict(out.get("visual_selection") or {})
    selected = [
        str(item.get("asset_id") or "").strip()
        for item in list(storyboard.get("selected_assets") or [])
        if isinstance(item, dict) and str(item.get("asset_id") or "").strip()
    ]
    primary = [
        str(item.get("asset_id") or "").strip()
        for item in list(storyboard.get("primary_assets") or [])
        if isinstance(item, dict) and str(item.get("asset_id") or "").strip()
    ]
    secondary = [
        str(item.get("asset_id") or "").strip()
        for item in list(storyboard.get("secondary_assets") or [])
        if isinstance(item, dict) and str(item.get("asset_id") or "").strip()
    ]
    reserve = [
        str(item.get("asset_id") or "").strip()
        for item in list(storyboard.get("reserve_assets") or [])
        if isinstance(item, dict) and str(item.get("asset_id") or "").strip()
    ]
    rejected = [
        str(item.get("asset_id") or "").strip()
        for item in list(storyboard.get("rejected_assets") or [])
        if isinstance(item, dict) and str(item.get("asset_id") or "").strip()
    ]
    if selected:
        visual_selection["storyboard_selected_asset_ids"] = selected
    if primary:
        visual_selection["primary_visual_ids"] = primary
        visual_selection["storyboard_primary_asset_ids"] = primary
    if secondary:
        visual_selection["secondary_visual_ids"] = secondary
        visual_selection["storyboard_secondary_asset_ids"] = secondary
    if reserve:
        visual_selection["reserve_visual_ids"] = reserve
        visual_selection["storyboard_reserve_asset_ids"] = reserve
    if rejected:
        visual_selection["forbidden_visual_ids"] = rejected
        visual_selection["storyboard_rejected_asset_ids"] = rejected
    if primary:
        current_target = _safe_int(visual_selection.get("target_visual_count"), default=0)
        visual_selection["target_visual_count"] = max(current_target, len(primary))
        visual_selection["max_visual_count"] = max(
            _safe_int(visual_selection.get("max_visual_count"), default=0),
            len(primary) + len(secondary),
            visual_selection["target_visual_count"],
        )
    elif storyboard.get("target_visual_count"):
        current_target = _safe_int(visual_selection.get("target_visual_count"), default=0)
        visual_selection["target_visual_count"] = max(
            current_target,
            _safe_int(storyboard.get("target_visual_count"), default=0),
        )
    if "selected_source_visual_count" in visual_selection:
        visual_selection["selected_source_visual_count"] = len(primary or selected)
    if "source_visual_shortfall" in visual_selection:
        unit_target = _safe_int(visual_selection.get("poster_visual_unit_target"), default=0)
        visual_selection["source_visual_shortfall"] = max(0, unit_target - len(primary or selected))
    out["visual_selection"] = visual_selection
    policy = dict(out.get("source_asset_policy") or {})
    if primary:
        policy["primary_assets_mandatory"] = primary
    if secondary:
        policy["secondary_assets_optional"] = secondary
    if reserve:
        policy["reserve_assets_replacements_only"] = reserve
    if rejected:
        policy["forbidden_source_ids"] = rejected
    out["source_asset_policy"] = policy
    return _refresh_brief_visual_unit_metrics(out)


def _sanitize_paper_visual_storyboard_for_rendered(
    paper_visual_storyboard: dict[str, Any],
    rendered_layers: dict[str, dict[str, Any]],
) -> dict[str, Any]:
    if not isinstance(paper_visual_storyboard, dict) or not paper_visual_storyboard:
        return {}
    rendered_layers = rendered_layers if isinstance(rendered_layers, dict) else {}
    tiers = {
        str(layer_id): str(classify_source_visual(str(layer_id), rec).get("visual_selection_tier") or "rejected")
        for layer_id, rec in rendered_layers.items()
        if str(layer_id).startswith(("ingest_fig_", "ingest_table_")) and isinstance(rec, dict)
    }
    blocked = {
        str(layer_id)
        for layer_id, tier in tiers.items()
        if tier == "rejected"
    }
    reserve_only = {layer_id for layer_id, tier in tiers.items() if tier == "reserve_unmatched"}
    out = dict(paper_visual_storyboard)
    minimum_count = _storyboard_minimum_source_visual_count(out)
    optional_ids = constrain_optional_source_visual_ids(
        [
            str(item.get("asset_id") or item.get("layer_id") or "")
            for key in ("secondary_assets", "selected_assets", "reserve_assets")
            for item in list(out.get(key) or [])
            if isinstance(item, dict)
        ],
        rendered_layers,
        minimum_count=minimum_count,
    )
    allowed_reserve_ids = {
        layer_id for layer_id in optional_ids
        if tiers.get(layer_id) == "reserve_unmatched"
    }
    blocked_rejections = [
        {
            "asset_id": layer_id,
            "reason": "designer-selected-ineligible source asset: "
            + ", ".join(_selected_visual_reject_reasons(layer_id, rendered_layers.get(layer_id) or {})),
        }
        for layer_id in sorted(blocked)
    ]

    def keep_items(key: str) -> list[dict[str, Any]]:
        return [
            item for item in list(out.get(key) or [])
            if isinstance(item, dict)
            and str(item.get("asset_id") or item.get("layer_id") or "").strip()
            and str(item.get("asset_id") or item.get("layer_id") or "").strip() not in blocked
            and (
                str(item.get("asset_id") or item.get("layer_id") or "").strip() not in reserve_only
                or (
                    key in {"selected_assets", "secondary_assets"}
                    and str(item.get("asset_id") or item.get("layer_id") or "").strip() in allowed_reserve_ids
                )
            )
        ]

    for key in ("selected_assets", "primary_assets", "secondary_assets", "reserve_assets"):
        out[key] = keep_items(key)
    rejected = [
        item for item in list(out.get("rejected_assets") or [])
        if isinstance(item, dict)
        and tiers.get(str(item.get("asset_id") or item.get("layer_id") or "").strip(), "rejected") == "rejected"
    ]
    existing = {
        str(item.get("asset_id") or item.get("layer_id") or "")
        for item in rejected
        if isinstance(item, dict)
    }
    rejected.extend(item for item in blocked_rejections if str(item.get("asset_id") or "") not in existing)
    out["rejected_assets"] = rejected
    out["primary_asset_count"] = len(out.get("primary_assets") or [])
    out["secondary_asset_count"] = len(out.get("secondary_assets") or [])
    metrics = dict(out.get("metrics") or {})
    metrics["selected_asset_count"] = len(out.get("selected_assets") or [])
    metrics["primary_asset_count"] = len(out.get("primary_assets") or [])
    metrics["secondary_asset_count"] = len(out.get("secondary_assets") or [])
    metrics["planner_ineligible_asset_count"] = len(blocked)
    out["metrics"] = metrics
    return out


def _storyboard_minimum_source_visual_count(storyboard: dict[str, Any] | None) -> int:
    if not isinstance(storyboard, dict):
        return 0
    metrics = storyboard.get("metrics") if isinstance(storyboard.get("metrics"), dict) else {}
    capacity = metrics.get("capacity") if isinstance(metrics.get("capacity"), dict) else {}
    return max(0, _safe_int(capacity.get("minimum_count"), default=0))


def _sanitize_poster_content_brief_visual_eligibility(
    poster_content_brief: dict[str, Any],
    rendered_layers: dict[str, dict[str, Any]],
) -> dict[str, Any]:
    if not isinstance(poster_content_brief, dict) or not poster_content_brief:
        return {}
    rendered_layers = rendered_layers if isinstance(rendered_layers, dict) else {}
    existing_visual_selection = (
        poster_content_brief.get("visual_selection")
        if isinstance(poster_content_brief.get("visual_selection"), dict)
        else {}
    )
    existing_policy = (
        poster_content_brief.get("source_asset_policy")
        if isinstance(poster_content_brief.get("source_asset_policy"), dict)
        else {}
    )
    existing_forbidden = set(_take_unique_values([
        *[str(v) for v in (existing_visual_selection.get("forbidden_visual_ids") or [])],
        *[str(v) for v in (existing_visual_selection.get("storyboard_rejected_asset_ids") or [])],
        *[str(v) for v in (existing_policy.get("forbidden_source_ids") or [])],
    ]))
    tiers = {
        str(layer_id): str(classify_source_visual(str(layer_id), rec).get("visual_selection_tier") or "rejected")
        for layer_id, rec in rendered_layers.items()
        if str(layer_id).startswith(("ingest_fig_", "ingest_table_")) and isinstance(rec, dict)
    }
    visual_storyboard = (
        poster_content_brief.get("visual_storyboard")
        if isinstance(poster_content_brief.get("visual_storyboard"), dict)
        else {}
    )
    storyboard_rejected = {
        str(item.get("asset_id") or item.get("layer_id") or "").strip()
        for item in visual_storyboard.get("rejected_assets") or []
        if isinstance(item, dict)
        and str(item.get("asset_id") or item.get("layer_id") or "").strip()
    }
    storyboard_rejected.update(
        str(asset_id or "").strip()
        for asset_id in visual_storyboard.get("rejected_asset_ids") or []
        if str(asset_id or "").strip()
    )
    allowed_primary = {
        layer_id for layer_id, tier in tiers.items()
        if tier == "eligible" and layer_id not in storyboard_rejected
    }
    blocked = {layer_id for layer_id, tier in tiers.items() if tier == "rejected"}
    minimum_count = _storyboard_minimum_source_visual_count(
        poster_content_brief.get("visual_storyboard")
        if isinstance(poster_content_brief.get("visual_storyboard"), dict)
        else {}
    )
    optional_candidates = _take_unique_values([
        *[str(v) for v in (existing_visual_selection.get("secondary_visual_ids") or [])],
        *[str(v) for v in (existing_visual_selection.get("reserve_visual_ids") or [])],
        *[str(v) for v in (existing_visual_selection.get("storyboard_selected_asset_ids") or [])],
        *[str(v) for v in (existing_visual_selection.get("storyboard_secondary_asset_ids") or [])],
        *[str(v) for v in (existing_visual_selection.get("storyboard_reserve_asset_ids") or [])],
        *[str(v) for v in (existing_policy.get("secondary_assets_optional") or [])],
        *[str(v) for v in (existing_policy.get("reserve_assets_replacements_only") or [])],
    ])
    constrained_optional = set(constrain_optional_source_visual_ids(
        optional_candidates,
        rendered_layers,
        minimum_count=minimum_count,
    ))
    allowed_optional = {
        layer_id for layer_id, tier in tiers.items()
        if layer_id not in storyboard_rejected
        and (tier == "eligible" or (tier == "reserve_unmatched" and layer_id in constrained_optional))
    }

    def clean_ids(values: Any, *, allow_reserve: bool = False) -> list[str]:
        allowed = allowed_optional if allow_reserve else allowed_primary
        out: list[str] = []
        for value in values or []:
            item = str(value or "").strip()
            if item and item in allowed and item not in out:
                out.append(item)
        return out

    out = dict(poster_content_brief)
    visual_selection = dict(out.get("visual_selection") or {})
    for key in (
        "primary_visual_ids",
        "high_priority_visual_ids",
        "storyboard_primary_asset_ids",
    ):
        if key in visual_selection:
            visual_selection[key] = clean_ids(visual_selection.get(key))
    for key in (
        "secondary_visual_ids",
        "reserve_visual_ids",
        "storyboard_selected_asset_ids",
        "storyboard_secondary_asset_ids",
        "storyboard_reserve_asset_ids",
    ):
        if key in visual_selection:
            visual_selection[key] = clean_ids(visual_selection.get(key), allow_reserve=True)
    forbidden = _take_unique_values([
        *[
            value for value in existing_forbidden
            if not value.startswith(("ingest_fig_", "ingest_table_"))
        ],
        *sorted(storyboard_rejected),
        *sorted(blocked),
    ])
    visual_selection["forbidden_visual_ids"] = forbidden
    visual_selection["storyboard_rejected_asset_ids"] = forbidden
    if isinstance(visual_selection.get("role_buckets"), dict):
        visual_selection["role_buckets"] = {
            str(bucket): clean_ids(values)
            for bucket, values in visual_selection.get("role_buckets", {}).items()
        }
    if isinstance(visual_selection.get("visual_records"), list):
        visual_selection["visual_records"] = [
            item for item in visual_selection.get("visual_records") or []
            if isinstance(item, dict)
            and str(item.get("layer_id") or "").strip() in allowed_optional
        ]
    if isinstance(visual_selection.get("source_asset_records"), list):
        visual_selection["source_asset_records"] = [
            item for item in visual_selection.get("source_asset_records") or []
            if isinstance(item, dict)
            and str(item.get("asset_id") or "").strip() in allowed_optional
        ]
    out["visual_selection"] = visual_selection

    policy = dict(out.get("source_asset_policy") or {})
    if "primary_assets_mandatory" in policy:
        policy["primary_assets_mandatory"] = clean_ids(policy.get("primary_assets_mandatory"))
    for key in ("secondary_assets_optional", "reserve_assets_replacements_only"):
        if key in policy:
            policy[key] = clean_ids(policy.get(key), allow_reserve=True)
    policy["forbidden_source_ids"] = _take_unique_values([
        *[
            value for value in existing_forbidden
            if not value.startswith(("ingest_fig_", "ingest_table_"))
        ],
        *sorted(blocked),
    ])
    out["source_asset_policy"] = policy

    if isinstance(out.get("visual_storyboard"), dict):
        out["visual_storyboard"] = _sanitize_paper_visual_storyboard_for_rendered(
            out.get("visual_storyboard") or {},
            rendered_layers,
        )
    for list_key in ("sections", "panel_plan", "editorial_column_plan"):
        if isinstance(out.get(list_key), list):
            cleaned_items = []
            for item in out.get(list_key) or []:
                if not isinstance(item, dict):
                    cleaned_items.append(item)
                    continue
                item_out = dict(item)
                if "visual_ids" in item_out:
                    item_out["visual_ids"] = clean_ids(item_out.get("visual_ids"))
                cleaned_items.append(item_out)
            out[list_key] = cleaned_items
    return _refresh_brief_visual_unit_metrics(out)


def _rendered_layers_from_reused_provenance(
    provenance: dict[str, Any],
    *,
    run_dir: Path,
    layers_dir: Path,
) -> dict[str, dict[str, Any]]:
    rendered: dict[str, dict[str, Any]] = {}
    for asset in provenance.get("assets") or []:
        if not isinstance(asset, dict):
            continue
        asset_id = str(asset.get("asset_id") or "").strip()
        output_file = str(asset.get("output_file") or "").strip()
        if not asset_id or not output_file:
            continue
        rel_path = Path(output_file)
        src_path = (run_dir / rel_path).resolve() if not rel_path.is_absolute() else rel_path
        if not src_path.exists():
            fallback = layers_dir / f"img_{asset_id}.png"
            if fallback.exists():
                src_path = fallback.resolve()
            else:
                continue
        width = _safe_int(asset.get("output_width_px"), default=0)
        height = _safe_int(asset.get("output_height_px"), default=0)
        rec = {
            "layer_id": asset_id,
            "name": str(asset.get("caption_short") or asset_id),
            "kind": "table" if str(asset.get("kind") or "").lower() == "table" else "image",
            "source": "ingested_pdf",
            "source_id": asset_id,
            "src_path": str(src_path),
            "png_path": str(src_path),
            "caption": asset.get("caption_full") or asset.get("caption_short") or "",
            "caption_short": asset.get("caption_short") or "",
            "caption_full": asset.get("caption_full") or "",
            "source_page": asset.get("source_page"),
            "source_pdf": asset.get("source_pdf"),
            "source_pdf_sha256": asset.get("source_pdf_sha256"),
            "source_bbox_pdf_points": asset.get("source_bbox_pdf_points"),
            "sha256": asset.get("output_sha256"),
            "visual_role": asset.get("visual_role"),
            "visual_score": asset.get("visual_score"),
            "curation_reason": asset.get("curation_reason"),
            "curation_flags": asset.get("curation_flags") or [],
            "crop_quality_flags": asset.get("crop_quality_flags") or [],
            "placement_quality_flags": asset.get("placement_quality_flags") or [],
            "caption_association_method": asset.get("caption_association_method") or "",
            "caption_confidence": asset.get("caption_confidence"),
            "extract_strategy": asset.get("extract_strategy"),
            "material_quality": _json_clone(asset.get("material_quality"))
            if isinstance(asset.get("material_quality"), dict)
            else {},
            "protected_anchor": bool(asset.get("protected_anchor")),
            "anchor_kind": asset.get("anchor_kind"),
            "anchor_label": asset.get("anchor_label"),
            "anchor_reason": asset.get("anchor_reason"),
            "captioned_source_group": bool(asset.get("captioned_source_group")),
            "source_group_id": asset.get("source_group_id"),
            "source_group_kind": asset.get("source_group_kind"),
            "source_group_label": asset.get("source_group_label"),
            "source_group_caption": asset.get("source_group_caption"),
            "source_group_source": asset.get("source_group_source"),
            "table_parse_status": asset.get("table_parse_status"),
            "output_file": f"layers/{src_path.name}",
        }
        if width and height:
            rec["image_size"] = f"{width}x{height}"
            rec["width"] = width
            rec["height"] = height
        if str(rec.get("extract_strategy") or "").lower() in {"raster", "embedded"}:
            rec["placement_quality_flags"] = _take_unique_values([
                *[str(v) for v in (rec.get("placement_quality_flags") or [])],
                *[str(v) for v in (rec.get("crop_quality_flags") or [])],
            ])
            rec["crop_quality_flags"] = []
        if rec["kind"] == "table":
            table_meta = asset.get("table_metadata") if isinstance(asset.get("table_metadata"), dict) else {}
            if table_meta:
                for key in ("headers", "rows", "col_highlight_rule", "title"):
                    if key in table_meta:
                        rec[key] = _json_clone(table_meta.get(key))
                if table_meta.get("table_visual_source"):
                    rec["table_visual_source"] = table_meta.get("table_visual_source")
            rec.setdefault("headers", [])
            rec.setdefault("rows", [])
            rec.setdefault("col_highlight_rule", [])
        rec.update(_visual_eligibility_payload(asset_id, rec, asset))
        rendered[asset_id] = rec
    return rendered


def _reused_ingest_summary(
    *,
    paper_memory: dict[str, Any],
    paper_visual_provenance: dict[str, Any],
    rendered_layers: dict[str, dict[str, Any]],
) -> dict[str, Any]:
    metadata = paper_memory.get("metadata") if isinstance(paper_memory.get("metadata"), dict) else {}
    source_file = str(paper_memory.get("source_file") or "").strip()
    if not source_file:
        docs = paper_visual_provenance.get("source_documents") or []
        if docs and isinstance(docs[0], dict):
            source_file = str(docs[0].get("source_file") or docs[0].get("name") or "").strip()
    chunks = paper_memory.get("chunks") if isinstance(paper_memory.get("chunks"), list) else []
    raw_text = "\n\n".join(
        str(chunk.get("text") or "")
        for chunk in chunks
        if isinstance(chunk, dict) and str(chunk.get("text") or "").strip()
    )
    manifest = {
        "title": metadata.get("title") or "",
        "authors": metadata.get("authors") or [],
        "venue": metadata.get("venue") or "",
        "abstract": metadata.get("abstract") or "",
        "raw_text": raw_text,
    }
    table_ids = [
        layer_id for layer_id, rec in sorted(rendered_layers.items())
        if isinstance(rec, dict) and rec.get("kind") == "table"
    ]
    figure_ids = [
        layer_id for layer_id, rec in sorted(rendered_layers.items())
        if isinstance(rec, dict) and rec.get("kind") != "table"
    ]
    return {
        "type": "pdf",
        "path": source_file,
        "manifest": manifest,
        "raw_text": raw_text,
        "registered_layer_ids": sorted(rendered_layers),
        "registered_figure_ids": figure_ids,
        "registered_table_ids": table_ids,
        "paper_memory_cache_key": paper_memory.get("cache_key"),
        "body_window": paper_memory.get("body_window"),
    }


def _reused_figures_tables_payload(
    rendered_layers: dict[str, dict[str, Any]],
) -> tuple[list[dict[str, Any]], list[dict[str, Any]]]:
    figures: list[dict[str, Any]] = []
    tables: list[dict[str, Any]] = []
    for layer_id, rec in sorted(rendered_layers.items()):
        payload = {
            "layer_id": layer_id,
            "source_file": rec.get("source_pdf"),
            "source_page": rec.get("source_page"),
            "source_ref": rec.get("source_ref"),
            "image_size": rec.get("image_size"),
            "caption": rec.get("caption"),
            "caption_short": rec.get("caption_short"),
            "sha256": rec.get("sha256"),
            "source_bbox_pdf_points": rec.get("source_bbox_pdf_points"),
            "source_pdf_sha256": rec.get("source_pdf_sha256"),
            "visual_role": rec.get("visual_role"),
            "visual_score": rec.get("visual_score"),
            "curation_reason": rec.get("curation_reason"),
            "curation_flags": rec.get("curation_flags") or [],
        }
        if rec.get("kind") == "table":
            rows = rec.get("rows") if isinstance(rec.get("rows"), list) else []
            headers = rec.get("headers") if isinstance(rec.get("headers"), list) else []
            payload["n_rows"] = len(rows)
            payload["n_cols"] = len(headers) or (len(rows[0]) if rows and isinstance(rows[0], list) else 0)
            tables.append(payload)
        else:
            figures.append(payload)
    return figures, tables


# ───────────────────────────── PDF branch ──────────────────────────────

def _ingest_pdf(fp: Path, ctx: ToolContext) -> dict[str, Any]:
    _raise_if_ingest_cancelled(ctx, "ingest.pdf.start")
    size = fp.stat().st_size
    if size > _MAX_PDF_BYTES:
        raise RuntimeError(
            f"PDF too large ({size / 1_048_576:.1f} MB > "
            f"{_MAX_PDF_BYTES / 1_048_576:.0f} MB). Split the document."
        )
    _raise_if_ingest_cancelled(ctx, "ingest.pdf.before_page_count")
    pages = page_count(fp)
    _raise_if_ingest_cancelled(ctx, "ingest.pdf.after_page_count")
    if _paper_source_scope() == "all" and pages > _MAX_PDF_PAGES:
        raise RuntimeError(
            f"PDF has {pages} pages (cap {_MAX_PDF_PAGES}). Trim the document."
        )

    _raise_if_ingest_cancelled(ctx, "ingest.pdf.before_open")
    doc = fitz.open(fp)
    is_scanned = False
    identity_text = ""
    page_window: dict[str, Any] = {
        "total_page_count": pages,
        "body_page_count": pages,
        "references_start_page": None,
        "appendix_start_page": None,
        "cutoff_start_page": None,
        "cutoff_reason": None,
        "source_scope": _paper_source_scope(),
        "ignored_reference_page_count": 0,
        "ignored_non_body_page_count": 0,
    }
    try:
        _raise_if_ingest_cancelled(ctx, "ingest.pdf.after_open")
        _raise_if_ingest_cancelled(ctx, "ingest.pdf.before_scanned_detection")
        is_scanned = detect_scanned_pdf(doc)
        _raise_if_ingest_cancelled(ctx, "ingest.pdf.after_scanned_detection")
        if is_scanned and pages > _MAX_PDF_PAGES:
            raise RuntimeError(
                f"Scanned PDF has {pages} pages (cap {_MAX_PDF_PAGES}). Trim the document."
            )
        if is_scanned:
            # v1.2.5 — OCR fallback via Qwen-VL-Max. Scanned PDFs have
            # no embedded rasters (the pages ARE the images) so figure
            # extraction is skipped; we render each page and OCR it so
            # structure extraction still has something to chew on.
            all_page_texts = _ocr_scanned_pdf(fp, doc, ctx)
            _raise_if_ingest_cancelled(ctx, "ingest.pdf.after_ocr")
            page_window = _paper_body_page_window(all_page_texts)
            _enforce_pdf_page_cap(page_window)
            identity_text = _identity_metadata_snippets(all_page_texts, page_window)
            log(
                "ingest.pdf.body_window",
                file=fp.name,
                total_pages=page_window["total_page_count"],
                body_pages=page_window["body_page_count"],
                references_start_page=page_window["references_start_page"],
                appendix_start_page=page_window["appendix_start_page"],
                cutoff_reason=page_window["cutoff_reason"],
                source_scope=page_window["source_scope"],
                ignored_reference_pages=page_window["ignored_reference_page_count"],
                ignored_non_body_pages=page_window["ignored_non_body_page_count"],
            )
            page_texts = all_page_texts[: int(page_window["body_page_count"] or len(all_page_texts))]
            candidates: list[PdfFigureCandidate] = []
            table_candidates: list[PdfTableCandidate] = []
        else:
            # 0. Page text is cheap and lets us stop before appendix,
            # supplement, and references before rendering/caption matching.
            _raise_if_ingest_cancelled(ctx, "ingest.pdf.before_page_text")
            all_page_texts = extract_page_text(doc)
            _raise_if_ingest_cancelled(ctx, "ingest.pdf.after_page_text")
            page_window = _paper_body_page_window(all_page_texts)
            _enforce_pdf_page_cap(page_window)
            identity_text = _identity_metadata_snippets(all_page_texts, page_window)
            body_page_count = int(page_window["body_page_count"] or len(all_page_texts))
            page_texts = all_page_texts[:body_page_count]
            log(
                "ingest.pdf.body_window",
                file=fp.name,
                total_pages=page_window["total_page_count"],
                body_pages=page_window["body_page_count"],
                references_start_page=page_window["references_start_page"],
                appendix_start_page=page_window["appendix_start_page"],
                cutoff_reason=page_window["cutoff_reason"],
                source_scope=page_window["source_scope"],
                ignored_reference_pages=page_window["ignored_reference_page_count"],
                ignored_non_body_pages=page_window["ignored_non_body_page_count"],
            )

            # 1. pymupdf figure candidates (no LLM).
            candidates = []
            _raise_if_ingest_cancelled(ctx, "ingest.pdf.before_raster_extract")
            candidates.extend(extract_embedded_rasters(doc, ctx.layers_dir, max_page=body_page_count))
            _raise_if_ingest_cancelled(ctx, "ingest.pdf.after_raster_extract")
            _raise_if_ingest_cancelled(ctx, "ingest.pdf.before_vector_extract")
            candidates.extend(extract_vector_clusters(doc, ctx.layers_dir, max_page=body_page_count))
            _raise_if_ingest_cancelled(ctx, "ingest.pdf.after_vector_extract")
            candidates = dedup_raster_vector(candidates)
            log("ingest.pdf.candidates", file=fp.name,
                raster=sum(1 for c in candidates if c.strategy == "raster"),
                vector=sum(1 for c in candidates if c.strategy == "vector"))

            # 1b. PyMuPDF table candidates (localization only; VLM parses
            # cells). Figure-overlap dedup runs after caption-anchor recovery
            # so early paper tables can be protected before local heuristics
            # drop them.
            _raise_if_ingest_cancelled(ctx, "ingest.pdf.before_table_extract")
            table_candidates = extract_table_candidates(
                doc, ctx.layers_dir, max_page=body_page_count,
            )
            _raise_if_ingest_cancelled(ctx, "ingest.pdf.after_table_extract")
            log("ingest.pdf.table_candidates.raw",
                file=fp.name, n=len(table_candidates))

    finally:
        doc.close()

    cover_png: Path | None = None
    if _structure_cover_image_enabled():
        # Opt-in compatibility/debug path for visual title/logo grounding.
        # Figure/table evidence still flows through their own crop-specific
        # VLM stages below.
        cover_png = ctx.layers_dir / "ingest_page_001.png"
        try:
            _raise_if_ingest_cancelled(ctx, "ingest.pdf.before_cover_render")
            render_page_png(fp, 1, cover_png, dpi=_STRUCTURE_PAGE_DPI)
            _raise_if_ingest_cancelled(ctx, "ingest.pdf.after_cover_render")
        except Exception as e:
            _raise_if_ingest_cancelled(ctx, "ingest.pdf.before_cover_render_failure_log")
            log("ingest.pdf.render_fail", page=1, error=str(e))
            _raise_if_ingest_cancelled(ctx, "ingest.pdf.after_cover_render_failure_log")
            cover_png = None
    else:
        log(
            "ingest.pdf.structure.cover_image_skipped",
            file=fp.name,
            reason="text_only_default",
        )

    structure_degraded = False
    try:
        manifest = _extract_structure(page_texts, cover_png, ctx, fp, identity_text=identity_text)
        if not isinstance(manifest, dict):
            raise RuntimeError(
                f"structure extractor returned {type(manifest).__name__}"
            )
    except Exception as e:  # noqa: BLE001 - keep ingest alive on VLM/API faults
        structure_degraded = True
        log(
            "ingest.pdf.structure.degraded",
            file=fp.name,
            error=f"{type(e).__name__}: {e}"[:500],
            fallback="local_text_and_candidate_figures",
        )
        manifest = _fallback_pdf_manifest(page_texts, fp, identity_text=identity_text)
    _normalize_manifest_lists(manifest)
    if identity_text and not manifest.get("authors"):
        fallback_authors = _fallback_authors_from_identity(identity_text)
        if fallback_authors:
            manifest["authors"] = fallback_authors
    if identity_text and not manifest.get("affiliations"):
        fallback_affiliations = _fallback_affiliations_from_identity(identity_text)
        if fallback_affiliations:
            manifest["affiliations"] = fallback_affiliations
    manifest["_structure_degraded"] = bool(structure_degraded)
    manifest["_page_window"] = page_window
    manifest["page_count"] = page_window.get("total_page_count")
    manifest["body_page_count"] = page_window.get("body_page_count")
    manifest["references_start_page"] = page_window.get("references_start_page")
    manifest["appendix_start_page"] = page_window.get("appendix_start_page")
    manifest["cutoff_start_page"] = page_window.get("cutoff_start_page")
    manifest["cutoff_reason"] = page_window.get("cutoff_reason")
    manifest["source_scope"] = page_window.get("source_scope")
    manifest["ignored_reference_page_count"] = page_window.get("ignored_reference_page_count")
    manifest["ignored_non_body_page_count"] = page_window.get("ignored_non_body_page_count")

    if not is_scanned:
        try:
            _raise_if_ingest_cancelled(ctx, "ingest.pdf.before_caption_recovery_open")
            with fitz.open(fp) as recovery_doc:
                _raise_if_ingest_cancelled(ctx, "ingest.pdf.after_caption_recovery_open")
                _raise_if_ingest_cancelled(ctx, "ingest.pdf.before_caption_discovery")
                discovered_caption_groups = discover_captioned_visual_groups(
                    recovery_doc,
                    manifest=manifest,
                    max_page=int(page_window.get("body_page_count") or len(page_texts) or len(recovery_doc)),
                )
                _raise_if_ingest_cancelled(ctx, "ingest.pdf.after_caption_discovery")
                _enrich_manifest_with_pdf_caption_groups(manifest, discovered_caption_groups)
                _raise_if_ingest_cancelled(ctx, "ingest.pdf.before_caption_crop_recovery")
                candidates, table_candidates = recover_caption_anchored_visuals(
                    recovery_doc,
                    ctx.layers_dir,
                    manifest=manifest,
                    page_texts=page_texts,
                    figures=candidates,
                    tables=table_candidates,
                    max_page=int(page_window.get("body_page_count") or len(page_texts) or len(recovery_doc)),
                )
                _raise_if_ingest_cancelled(ctx, "ingest.pdf.after_caption_crop_recovery")
        except Exception as e:  # noqa: BLE001 - anchor recovery is best-effort
            _raise_if_ingest_cancelled(ctx, "ingest.pdf.before_caption_recovery_failure_log")
            log(
                "ingest.pdf.caption_anchor.recovery_failed",
                file=fp.name,
                error=f"{type(e).__name__}: {e}"[:500],
            )
            _raise_if_ingest_cancelled(ctx, "ingest.pdf.after_caption_recovery_failure_log")
    pre_dedup = len(table_candidates)
    table_candidates = dedup_tables_against_figures(
        table_candidates, candidates,
    )
    log("ingest.pdf.table_candidates",
        file=fp.name, n=len(table_candidates),
        dropped_by_figure_overlap=pre_dedup - len(table_candidates),
        protected=sum(1 for t in table_candidates if getattr(t, "protected_anchor", False)))

    # 3. Caption matching — per figure candidate, parallelized.
    registered_layer_ids: list[str] = []
    if candidates:
        candidates = _budget_figure_candidates(candidates, ctx)
        _raise_if_ingest_cancelled(ctx, "ingest.pdf.before_caption_match")
        matches = (
            _fallback_caption_matches(candidates, manifest)
            if structure_degraded
            else _match_captions_parallel(candidates, manifest, ctx)
        )
        _raise_if_ingest_cancelled(ctx, "ingest.pdf.after_caption_match")
        _raise_if_ingest_cancelled(ctx, "ingest.pdf.before_figure_registration")
        registered_layer_ids = _register_candidates(
            candidates=candidates, matches=matches, ctx=ctx, pdf_path=fp,
        )
        _raise_if_ingest_cancelled(ctx, "ingest.pdf.after_figure_registration")
    else:
        log("ingest.pdf.no_candidates", file=fp.name,
            note="pymupdf returned 0 figures; VLM manifest may still be useful")

    # 4. Table parsing — per table candidate, parallelized. VLM reads
    # the bbox image + pymupdf's raw cell guess and returns clean
    # structured rows/headers. Rejects diagrams/equations/etc. that
    # find_tables() misclassified.
    registered_table_ids: list[str] = []
    if table_candidates:
        table_candidates = _budget_table_candidates(table_candidates, ctx)
        if structure_degraded:
            log(
                "ingest.pdf.table_parse.with_degraded_structure",
                file=fp.name,
                n_candidates=len(table_candidates),
                reason="using fallback manifest",
            )
        _raise_if_ingest_cancelled(ctx, "ingest.pdf.before_table_parse")
        parsed = _parse_tables_parallel(table_candidates, manifest, ctx)
        _raise_if_ingest_cancelled(ctx, "ingest.pdf.after_table_parse")
        _raise_if_ingest_cancelled(ctx, "ingest.pdf.before_table_registration")
        registered_table_ids = _register_tables(
            candidates=table_candidates,
            parsed=parsed,
            ctx=ctx,
            pdf_path=fp,
            manifest=manifest,
        )
        _raise_if_ingest_cancelled(ctx, "ingest.pdf.after_table_registration")

    _raise_if_ingest_cancelled(ctx, "ingest.pdf.before_post_registration_merge")
    all_registered = registered_layer_ids + registered_table_ids
    rendered = ctx.state.get("rendered_layers") or {}
    visual_candidate_scores = _annotate_visual_curation(all_registered, rendered)
    ranked_visual_ids = _rank_visual_ids_for_planner(all_registered, rendered)
    contact_sheet_path = _build_ingest_contact_sheet(
        ranked_visual_ids,
        rendered,
        ctx,
        fp,
    )
    recommended_figures = _recommend_paper_visuals(ranked_visual_ids, rendered)
    figure_catalog_summary = _figure_catalog_summary(ranked_visual_ids, rendered)
    recommended_text_units = _recommend_paper_text_units(
        manifest,
        rendered,
        ranked_visual_ids,
        ctx.state.get("claim_graph"),
    )
    memory_key = paper_memory_cache_key(
        pdf_path=fp,
        body_window=page_window,
        manifest=manifest,
        rendered_layers=rendered if isinstance(rendered, dict) else {},
        registered_layer_ids=all_registered,
        recommended_text_units=recommended_text_units,
    )
    paper_memory = read_paper_memory_cache(ctx.settings, memory_key)
    memory_cache_hit = paper_memory is not None
    if paper_memory is None:
        paper_memory = build_paper_memory(
            pdf_path=fp,
            page_texts=page_texts,
            manifest=manifest,
            body_window=page_window,
            rendered_layers=rendered if isinstance(rendered, dict) else {},
            registered_layer_ids=all_registered,
            recommended_text_units=recommended_text_units,
        )
        write_paper_memory_cache(ctx.settings, paper_memory)
    else:
        paper_memory = dict(paper_memory)
        paper_memory["source_file"] = str(fp)
    log(
        "paper_memory.cache_hit" if memory_cache_hit else "paper_memory.created",
        file=fp.name,
        cache_key=str(paper_memory.get("cache_key") or "")[:12],
        chunks=paper_memory.get("chunk_count"),
    )
    # v2.7 — keep the verbatim page text (already computed for manifest
    # extraction) so the composite-stage provenance validator can
    # substring-match LayerNode.evidence_quote against it. Cost: ~200 KB
    # on a 40-page paper, negligible vs the 5+ MB run-dir.
    return {
        "file": str(fp), "type": "pdf", "manifest": manifest,
        "page_window": page_window,
        "registered_layer_ids": all_registered,
        "registered_figure_ids": registered_layer_ids,
        "registered_table_ids": registered_table_ids,
        "contact_sheet_path": str(contact_sheet_path) if contact_sheet_path else None,
        "recommended_figures": recommended_figures,
        "visual_candidate_scores": visual_candidate_scores,
        "recommended_text_units": recommended_text_units,
        "paper_memory": paper_memory,
        "figure_catalog_summary": figure_catalog_summary,
        "raw_text": "\n\n".join(page_texts),
        "summary": f"{manifest.get('title', '?')} — "
                   f"{len(registered_layer_ids)} figure(s), "
                   f"{len(registered_table_ids)} table(s), "
                   f"{len(manifest.get('sections', []))} section(s), "
                   f"body pages {page_window.get('body_page_count')}/"
                   f"{page_window.get('total_page_count')}",
    }


def _ocr_scanned_pdf(
    fp: Path, doc: "fitz.Document", ctx: ToolContext,
) -> list[str]:
    """Render each page of a scanned PDF at `_OCR_PAGE_DPI` and ask
    Qwen-VL-Max (or whichever VLM `settings.ingest_model` points at) to
    OCR it. Returns one string per page (1-indexed: index 0 = page 1).

    Runs pages in parallel via ThreadPoolExecutor so a 40-page doc takes
    ~8 s wall time at 6 workers instead of 40 × 1 s serial. OCR failures
    on individual pages degrade to empty strings — we don't block the
    whole doc on one bad page, and the structure extractor handles
    partial text fine.
    """
    import time as _time

    _raise_if_ingest_cancelled(ctx, "ingest.ocr.start")
    n_pages = len(doc)
    parallelism = _ingest_vlm_parallelism(ctx, _OCR_PAGE_PARALLELISM)
    log("ingest.pdf.ocr.start", file=fp.name, pages=n_pages,
        dpi=_OCR_PAGE_DPI, parallelism=parallelism,
        model=ctx.settings.ingest_model)
    _raise_if_ingest_cancelled(ctx, "ingest.ocr.after_start_log")
    t0 = _time.monotonic()

    page_pngs: list[Path | None] = []
    for i in range(n_pages):
        _raise_if_ingest_cancelled(ctx, "ingest.ocr.before_page_render")
        out_path = ctx.layers_dir / f"ingest_ocr_page_{i + 1:03d}.png"
        try:
            pix = doc[i].get_pixmap(dpi=_OCR_PAGE_DPI)
            _raise_if_ingest_cancelled(ctx, "ingest.ocr.after_page_render")
            _raise_if_ingest_cancelled(ctx, "ingest.ocr.before_page_save")
            pix.save(str(out_path))
            _raise_if_ingest_cancelled(ctx, "ingest.ocr.after_page_save")
            page_pngs.append(out_path)
        except Exception as e:
            _raise_if_ingest_cancelled(ctx, "ingest.ocr.before_render_failure_log")
            log("ingest.pdf.ocr.render_fail", page=i + 1, error=str(e))
            _raise_if_ingest_cancelled(ctx, "ingest.ocr.after_render_failure_log")
            page_pngs.append(None)

    def ocr_one(idx: int, _payload: Any) -> str:
        _raise_if_ingest_cancelled(ctx, "ingest.ocr.worker.start")
        png = page_pngs[idx]
        if png is None:
            return ""
        try:
            _raise_if_ingest_cancelled(ctx, "ingest.ocr.worker.before_image_read")
            image = VlmImage.from_path(png)
            _raise_if_ingest_cancelled(ctx, "ingest.ocr.worker.after_image_read")
            _raise_if_ingest_cancelled(ctx, "ingest.ocr.worker.before_vlm")
            result = vlm_call_json(
                settings=ctx.settings,
                model=ctx.settings.ingest_model,
                system=_OCR_PROMPT,
                user_text=f"Page {idx + 1} of {n_pages}. OCR it.",
                images=[image],
                max_tokens=4096,
                timeout_s=_OCR_PER_PAGE_TIMEOUT_S,
                cancellation_token=_ingest_cancellation_token(ctx),
            )
            _raise_if_ingest_cancelled(ctx, "ingest.ocr.worker.after_vlm")
            return str(result.get("text") or "")
        except Exception as e:
            _raise_if_ingest_cancelled(ctx, "ingest.ocr.worker.before_failure_log")
            log("ingest.pdf.ocr.page_fail", page=idx + 1, error=str(e))
            _raise_if_ingest_cancelled(ctx, "ingest.ocr.worker.after_failure_log")
            return ""

    results: dict[int, str] = {}
    outcomes = _run_bounded_ingest_pool(
        items=[(i, None) for i in range(n_pages)],
        worker=ocr_one,
        parallelism=parallelism,
        ctx=ctx,
        phase="ingest.ocr.pool",
    )
    for idx, text, error in outcomes:
        _raise_if_ingest_cancelled(ctx, "ingest.ocr.before_result_merge")
        results[idx] = "" if error is not None else str(text or "")
        _raise_if_ingest_cancelled(ctx, "ingest.ocr.after_result_merge")

    _raise_if_ingest_cancelled(ctx, "ingest.ocr.before_ordered_result")
    page_texts = [results.get(i, "") for i in range(n_pages)]
    total_chars = sum(len(t) for t in page_texts)
    _raise_if_ingest_cancelled(ctx, "ingest.ocr.before_done_log")
    log("ingest.pdf.ocr.done", file=fp.name,
        wall_s=round(_time.monotonic() - t0, 1),
        total_chars=total_chars,
        pages_with_text=sum(1 for t in page_texts if t.strip()))
    _raise_if_ingest_cancelled(ctx, "ingest.ocr.after_done_log")
    return page_texts


def _extract_structure(
    page_texts: list[str],
    cover_png: Path | None,
    ctx: ToolContext,
    fp: Path,
    *,
    identity_text: str = "",
) -> dict[str, Any]:
    import time as _time

    _raise_if_ingest_cancelled(ctx, "ingest.structure.start")
    # Concatenate body text with [PAGE N] headers; appendix/supplement and
    # references/bibliography pages are removed before this call.
    body_lines: list[str] = []
    used = 0
    for page_num, text in enumerate(page_texts, start=1):
        chunk = f"\n[PAGE {page_num}]\n{text.strip()}"
        if used + len(chunk) > _STRUCTURE_TOTAL_TEXT_CAP:
            body_lines.append(
                f"\n[remaining {len(page_texts) - page_num + 1} "
                f"pages omitted — cap {_STRUCTURE_TOTAL_TEXT_CAP} chars]"
            )
            break
        body_lines.append(chunk)
        used += len(chunk)
    full_text = "".join(body_lines)

    images = (
        [VlmImage.from_path(cover_png)]
        if cover_png is not None and _structure_cover_image_enabled()
        else []
    )
    visual_context_note = (
        " An optional cover-page image is attached for title/logo grounding only."
        if images
        else ""
    )
    user_text = (
        "Below is the extracted body text of the paper before Appendix/"
        "Supplement/References/Bibliography."
        f"{visual_context_note} "
        "Extract the structured manifest as JSON per the system prompt. "
        "Return STRICT JSON only.\n\n"
        f"{full_text}"
    )
    if identity_text.strip():
        user_text += (
            "\n\nIdentity snippets for title/authors/team/school-institution-company "
            "names only. Do not use this section for sections, figures, "
            "tables, key_quotes, or poster claims.\n\n"
            f"{identity_text.strip()}"
        )

    _raise_if_ingest_cancelled(ctx, "ingest.structure.before_request_log")
    log("ingest.pdf.structure.request",
        file=fp.name, text_chars=used, n_pages=len(page_texts),
        cover_image=bool(images),
        model=ctx.settings.ingest_model,
        timeout_s=ctx.settings.ingest_http_timeout)
    _raise_if_ingest_cancelled(ctx, "ingest.structure.after_request_log")
    t0 = _time.monotonic()
    _raise_if_ingest_cancelled(ctx, "ingest.structure.before_vlm")
    manifest = vlm_call_json(
        settings=ctx.settings,
        model=ctx.settings.ingest_model,
        system=_INGEST_STRUCTURE_PROMPT,
        user_text=user_text,
        images=images,
        max_tokens=8192,
        cancellation_token=_ingest_cancellation_token(ctx),
    )
    _raise_if_ingest_cancelled(ctx, "ingest.structure.after_vlm")
    _raise_if_ingest_cancelled(ctx, "ingest.structure.before_response_log")
    log("ingest.pdf.structure.response",
        file=fp.name, wall_s=round(_time.monotonic() - t0, 1))
    _raise_if_ingest_cancelled(ctx, "ingest.structure.after_response_log")
    return manifest


def _normalize_manifest_lists(manifest: dict[str, Any]) -> None:
    """Coerce None → [] on list-shaped keys. Claude/Qwen occasionally
    emit `"figures": null` for papers they find no figures in, which
    breaks downstream len() calls."""
    for key in ("sections", "figures", "tables", "authors", "affiliations", "key_quotes"):
        if not isinstance(manifest.get(key), list):
            manifest[key] = []


def _fallback_pdf_manifest(
    page_texts: list[str],
    fp: Path,
    *,
    identity_text: str = "",
) -> dict[str, Any]:
    """Local-only PDF manifest used when the VLM structure pass fails.

    Figure localization already happened before the structure request, so
    aborting the whole run here wastes useful candidate images. This fallback
    gives the planner enough title/section/caption scaffolding to continue.
    """
    title = _fallback_title(page_texts, fp)
    abstract = _fallback_abstract(page_texts)
    sections = _fallback_sections(page_texts, abstract)
    figures = _fallback_captions(page_texts, kind="figure")
    tables = _fallback_captions(page_texts, kind="table")
    key_quotes = _fallback_key_quotes(page_texts)
    authors = _fallback_authors_from_identity(identity_text)
    affiliations = _fallback_affiliations_from_identity(identity_text)
    log(
        "ingest.pdf.fallback_manifest",
        file=fp.name,
        title_chars=len(title),
        authors=len(authors),
        affiliations=len(affiliations),
        sections=len(sections),
        figures=len(figures),
        tables=len(tables),
    )
    return {
        "title": title,
        "authors": authors,
        "affiliations": affiliations,
        "venue": None,
        "abstract": abstract,
        "sections": sections,
        "figures": figures,
        "tables": tables,
        "key_quotes": key_quotes,
    }


def _fallback_authors_from_identity(identity_text: str) -> list[str]:
    if not identity_text.strip():
        return []
    lines = [_clean_space(line) for line in identity_text.splitlines()]
    lines = [line for line in lines if line and not line.startswith("[IDENTITY PAGE")]
    authors: list[str] = []
    collecting = False
    for line in lines:
        low = line.lower().strip(":")
        if re.fullmatch(r"(?:authors?|author list|team|contributors?)", low):
            collecting = True
            continue
        if collecting:
            if re.match(r"^(?:abstract|keywords?|introduction|references|appendix)\b", low):
                collecting = False
                continue
            if _looks_affiliation_or_contact_line(line):
                continue
            authors.extend(_split_fallback_author_line(line))
            if len(authors) >= 20:
                return _dedupe_preserve_order(authors)[:20]
    if not authors:
        for line in lines:
            if re.search(r"\bteam\b|团队", line, re.IGNORECASE) and 3 <= len(line) <= 100:
                return [line]
    return _dedupe_preserve_order(authors)[:20]


def _split_fallback_author_line(line: str) -> list[str]:
    if not line or _looks_affiliation_or_contact_line(line):
        return []
    cleaned = re.sub(r"\s+(?:and|&)\s+", ", ", line)
    if not re.search(r"[,;·|]", cleaned):
        return [cleaned] if _looks_authorish_line(cleaned) else []
    parts = [
        _clean_space(part)
        for part in re.split(r"\s*(?:,|;|·|\|)\s*", cleaned)
        if _clean_space(part)
    ]
    return [part for part in parts if _looks_authorish_line(part)]


_AFFILIATION_KEYWORD_RE = re.compile(
    r"\b(?:university|college|institute|institution|laborator(?:y|ies)|lab|labs|"
    r"department|school of|research|academy|center|centre|corporation|company|"
    r"inc\.?|corp\.?|llc|ltd|team)\b|大学|学院|研究院|研究所|实验室|公司|机构|团队",
    re.IGNORECASE,
)
_KNOWN_ORGANIZATION_RE = re.compile(
    r"\b(?:Arizona State University|Georgia State University|Stanford|MIT|"
    r"Carnegie Mellon|CMU|Berkeley|Harvard|Caltech|Princeton|Tsinghua|"
    r"Peking University|ETH|EPFL|Meta|Google|Microsoft|OpenAI|"
    r"DeepMind|NVIDIA|Apple|Amazon|Adobe|ByteDance|Tencent|Alibaba)\b",
    re.IGNORECASE,
)


def _fallback_affiliations_from_identity(identity_text: str) -> list[str]:
    if not identity_text.strip():
        return []
    lines = [_clean_space(line) for line in identity_text.splitlines()]
    lines = [line for line in lines if line and not line.startswith("[IDENTITY PAGE")]
    affiliations: list[str] = []
    collecting = False
    for line in lines:
        low = line.lower().strip(":")
        if re.fullmatch(r"(?:affiliations?|institutions?|schools?|organizations?|单位|机构)", low):
            collecting = True
            continue
        if collecting:
            if re.fullmatch(r"(?:authors?|author list|abstract|keywords?|introduction|references|appendix)", low):
                collecting = False
                continue
            affiliations.extend(_split_fallback_affiliation_line(line))
            if len(affiliations) >= 16:
                return _dedupe_preserve_order(affiliations)[:16]
    if not affiliations:
        for line in lines[:80]:
            affiliations.extend(_split_fallback_affiliation_line(line))
            if len(affiliations) >= 16:
                break
    return _dedupe_preserve_order(affiliations)[:16]


def _split_fallback_affiliation_line(line: str) -> list[str]:
    if not _looks_affiliationish_line(line):
        return []
    if _looks_contact_only_line(line):
        return []
    parts = [
        _clean_space(part)
        for part in re.split(r"\s*(?:;|\||·|•|/)\s*", line)
        if _clean_space(part)
    ]
    out: list[str] = []
    for part in parts or [line]:
        cleaned = re.sub(r"^\d+\s*", "", part).strip(" ,")
        cleaned = re.sub(r"\s*\{[^{}]*\}\s*", " ", cleaned)
        cleaned = _clean_space(cleaned)
        if not cleaned or len(cleaned) > 120 or _looks_contact_only_line(cleaned):
            continue
        known = _KNOWN_ORGANIZATION_RE.search(cleaned)
        if known and re.search(r"\bteam\b|团队", cleaned, re.IGNORECASE):
            out.append(known.group(0))
        elif _looks_affiliationish_line(cleaned):
            out.append(cleaned)
    return out


def _looks_affiliationish_line(line: str) -> bool:
    text = _clean_space(line)
    if not (2 <= len(text) <= 180):
        return False
    low = text.lower()
    if any(token in low for token in ("abstract", "keywords", "figure", "table", "references")):
        return False
    if _looks_contact_only_line(text):
        return False
    return bool(_AFFILIATION_KEYWORD_RE.search(text) or _KNOWN_ORGANIZATION_RE.search(text))


def _looks_contact_only_line(line: str) -> bool:
    low = line.lower()
    if "http://" in low or "https://" in low or "www." in low:
        return True
    email_count = len(re.findall(r"[\w.+-]+@[\w.-]+\.\w+", line))
    return email_count > 0 and not (_AFFILIATION_KEYWORD_RE.search(line) or _KNOWN_ORGANIZATION_RE.search(line))


def _institution_color_signals(affiliations: list[str], authors: list[str] | None = None) -> dict[str, Any]:
    organizations = _dedupe_preserve_order([
        org
        for org in (_organization_signal_name(item) for item in affiliations)
        if org
    ])
    if not organizations:
        organizations = _dedupe_preserve_order([
            org
            for org in (_organization_signal_name(item) for item in (authors or []))
            if org
        ])
    strength = "none"
    if len(organizations) == 1:
        strength = "strong"
    elif len(organizations) > 1:
        strength = "mixed"
    return {
        "organizations": organizations[:12],
        "signal_strength": strength,
        "source": "paper_identity",
        "selection_guidance": (
            "Use these organization names only as soft color-association signals. "
            "Do not fetch logos, official brand colors, or brand assets; when signals "
            "are mixed or unclear, prioritize paper domain, source visual harmony, "
            "and academic readability."
        ),
    }


def _organization_signal_name(value: Any) -> str:
    text = _clean_space(str(value or ""))
    if not text:
        return ""
    known = _KNOWN_ORGANIZATION_RE.search(text)
    if known:
        return known.group(0)
    if not _looks_affiliationish_line(text):
        return ""
    text = re.sub(r"^\d+\s*", "", text)
    return text[:80].strip(" ,")


def _looks_authorish_line(line: str) -> bool:
    if not (2 <= len(line) <= 100):
        return False
    low = line.lower()
    if any(token in low for token in ("http", "www.", "abstract", "keyword", "figure", "table")):
        return False
    if re.search(r"\b(?:university|institute|laboratory|department|school of|college|inc\.?|corp\.?|llc)\b", low):
        return False
    if "@" in line:
        return False
    return bool(re.search(r"[A-Za-z\u4e00-\u9fff]", line))


def _looks_affiliation_or_contact_line(line: str) -> bool:
    low = line.lower()
    return (
        "@" in line
        or bool(_AFFILIATION_KEYWORD_RE.search(low))
        or bool(re.search(r"\b(?:email|correspondence)\b", low))
    )


def _dedupe_preserve_order(items: list[str]) -> list[str]:
    out: list[str] = []
    seen: set[str] = set()
    for item in items:
        key = item.lower()
        if key in seen:
            continue
        seen.add(key)
        out.append(item)
    return out


def _fallback_title(page_texts: list[str], fp: Path) -> str:
    first_page = page_texts[0] if page_texts else ""
    skip_prefixes = (
        "arxiv", "preprint", "doi", "http", "www.", "published",
        "accepted", "proceedings",
    )
    lines = [_clean_space(raw) for raw in first_page.splitlines()]
    for idx, line in enumerate(lines):
        if not (8 <= len(line) <= 180):
            continue
        low = line.lower()
        if low.startswith(skip_prefixes):
            continue
        if re.fullmatch(r"[\d\s.\-]+", line):
            continue
        parts = [line]
        for nxt in lines[idx + 1: idx + 4]:
            nxt_low = nxt.lower()
            if not nxt or nxt_low == "abstract" or "@" in nxt:
                break
            if nxt_low.startswith(skip_prefixes):
                break
            if len(nxt) < 18 and not nxt.endswith(("-", ":")):
                break
            parts.append(nxt)
            joined = _clean_space(" ".join(parts))
            if len(joined) >= 60 and not nxt.endswith(("-", ":")):
                break
        return _clean_space(" ".join(parts))[:180]
    return fp.stem


def _fallback_abstract(page_texts: list[str]) -> str:
    text = "\n".join(page_texts)
    match = re.search(
        r"(?is)\babstract\b\s*[:.\-]?\s*(.*?)(?=\n\s*(?:1\s+|1\.\s*)?"
        r"introduction\b|\n\s*2\s+|\Z)",
        text,
    )
    if not match:
        return ""
    return _clean_space(match.group(1))[:1400]


def _fallback_sections(
    page_texts: list[str],
    abstract: str,
) -> list[dict[str, Any]]:
    sections: list[dict[str, Any]] = []
    seen: set[str] = set()
    heading_re = re.compile(
        r"^\s*((?:\d+|[IVX]+)(?:\.\d+)*)[\s.]+"
        r"([A-Z][A-Za-z0-9 ,:/()&+\-]{2,90})\s*$"
    )
    for page_num, text in enumerate(page_texts, start=1):
        for raw in text.splitlines():
            line = _clean_space(raw)
            match = heading_re.match(line)
            if not match:
                continue
            title = _clean_space(match.group(2))
            low = title.lower()
            if low.startswith(("figure", "table", "appendix")):
                continue
            key = low
            if key in seen:
                continue
            seen.add(key)
            sections.append({"title": title, "page": page_num, "summary": ""})
            if len(sections) >= 16:
                return sections
    if not sections and abstract:
        sections.append({"title": "Abstract", "page": 1, "summary": abstract[:500]})
    return sections


def _fallback_captions(page_texts: list[str], *, kind: str) -> list[dict[str, Any]]:
    if kind == "figure":
        label_re = r"(?:Figure|Fig\.?)"
        out_key = "fig"
    else:
        label_re = r"Table"
        out_key = "table"
    cap_re = re.compile(
        rf"\b({label_re})\s+([A-Za-z]?\d+(?:\.\d+)?)\s*[:.]?\s+"
        rf"(.{{12,700}}?)(?=\b(?:Figure|Fig\.?|Table)\s+"
        rf"[A-Za-z]?\d+(?:\.\d+)?\s*[:.]?\s+|\Z)",
        re.IGNORECASE | re.DOTALL,
    )
    results: list[dict[str, Any]] = []
    seen: set[str] = set()
    for page_num, text in enumerate(page_texts, start=1):
        page_flat = _clean_space(text)
        for match in cap_re.finditer(page_flat):
            number = match.group(2)
            body = _clean_space(match.group(3))[:500]
            if not body:
                continue
            caption = f"{match.group(1)} {number}: {body}"
            key = caption.lower()[:180]
            if key in seen:
                continue
            seen.add(key)
            results.append({
                "id": f"{out_key}_{len(results) + 1}",
                "page": page_num,
                "caption": caption,
                "description": body[:240],
            })
            if len(results) >= 40:
                return results
    return results


def _fallback_key_quotes(page_texts: list[str]) -> list[str]:
    quotes: list[str] = []
    for text in page_texts[:8]:
        for raw in text.splitlines():
            line = _clean_space(raw)
            if 90 <= len(line) <= 260:
                quotes.append(line)
                if len(quotes) >= 6:
                    return quotes
    return quotes


def _fallback_caption_matches(
    candidates: list[PdfFigureCandidate],
    manifest: dict[str, Any],
) -> dict[int, dict[str, Any]]:
    figures = list(manifest.get("figures") or [])
    used: set[int] = set()
    matches: dict[int, dict[str, Any]] = {}
    for cand_idx, cand in enumerate(candidates):
        if _candidate_has_source_group(cand):
            best_idx = _source_group_manifest_index(cand, figures, used)
            if best_idx is not None:
                used.add(best_idx)
            caption = str(getattr(cand, "source_group_caption", "") or "")
            if best_idx is not None:
                caption = str(figures[best_idx].get("caption") or caption)
            matches[cand_idx] = {
                "matched_idx": best_idx,
                "confidence": 0.95,
                "is_real_figure": True,
                "reason": "bound from PDF captioned source group",
                "caption_text": caption,
                "short_caption": _fallback_short_caption(caption),
                "sub_panels": [],
                "caption_association_method": "captioned_group",
            }
            continue
        best_idx: int | None = None
        best_distance = 999
        for fig_idx, fig in enumerate(figures):
            if fig_idx in used:
                continue
            try:
                page = int(fig.get("page") or 0)
            except (TypeError, ValueError):
                page = 0
            distance = abs(page - cand.page) if page > 0 else 99
            if distance < best_distance:
                best_distance = distance
                best_idx = fig_idx
        caption = ""
        if best_idx is not None and best_distance <= 1:
            used.add(best_idx)
            caption = str(figures[best_idx].get("caption") or "")
        matches[cand_idx] = {
            "matched_idx": best_idx,
            "confidence": 0.25 if caption else 0.0,
            "is_real_figure": True,
            "reason": "local fallback after structure extraction failed",
            "caption_text": caption,
            "short_caption": _fallback_short_caption(caption),
            "sub_panels": [],
            "caption_association_method": "geometry_fallback" if caption else "unmatched",
        }
    return matches


def _fallback_short_caption(caption: str) -> str:
    text = re.sub(
        r"(?i)\b(?:figure|fig\.?|table)\s+[A-Za-z]?\d+(?:\.\d+)?\s*[:.]?",
        "",
        caption,
    )
    words = re.findall(r"[A-Za-z0-9][A-Za-z0-9+/_-]*", text)
    if not words:
        return ""
    return " ".join(words[:3])[:28]


def _budget_figure_candidates(
    candidates: list[PdfFigureCandidate],
    ctx: ToolContext,
) -> list[PdfFigureCandidate]:
    cap = _ingest_candidate_cap(
        "INGEST_CAPTION_MAX_CANDIDATES",
        _CHEAP_CAPTION_MATCH_MAX_CANDIDATES,
        ctx=ctx,
        dogfood_default=_DOGFOOD_CAPTION_MATCH_MAX_CANDIDATES,
    )
    if cap <= 0 or len(candidates) <= cap:
        return candidates
    keep_indices = _budget_source_candidate_indices(
        candidates,
        cap,
        page_of=lambda cand: cand.page,
        area_of=lambda cand: cand.width_px * cand.height_px,
    )
    kept = [cand for i, cand in enumerate(candidates) if i in keep_indices]
    dropped = [cand for i, cand in enumerate(candidates) if i not in keep_indices]
    preferred_count = sum(1 for cand in candidates if _candidate_has_source_group(cand))
    _unlink_unbudgeted_candidate_images(cand.path for cand in dropped)
    log(
        "ingest.pdf.caption_match.budget",
        mode=_poster_harness_mode(ctx) or "default",
        original=len(candidates),
        kept=len(kept),
        dropped=len(dropped),
        max_candidates=cap,
        preferred_source_groups=preferred_count,
        preferred_source_groups_kept=sum(1 for cand in kept if _candidate_has_source_group(cand)),
        parallelism=_ingest_vlm_parallelism(ctx, _CAPTION_MATCH_PARALLELISM),
    )
    return kept


def _budget_table_candidates(
    candidates: list[PdfTableCandidate],
    ctx: ToolContext,
) -> list[PdfTableCandidate]:
    cap = _ingest_candidate_cap(
        "INGEST_TABLE_PARSE_MAX_CANDIDATES",
        _CHEAP_TABLE_PARSE_MAX_CANDIDATES,
        ctx=ctx,
        dogfood_default=_DOGFOOD_TABLE_PARSE_MAX_CANDIDATES,
    )
    if cap <= 0 or len(candidates) <= cap:
        return candidates
    keep_indices = _budget_source_candidate_indices(
        candidates,
        cap,
        page_of=lambda cand: cand.page,
        area_of=lambda cand: cand.width_px * cand.height_px,
    )
    kept = [cand for i, cand in enumerate(candidates) if i in keep_indices]
    dropped = [cand for i, cand in enumerate(candidates) if i not in keep_indices]
    preferred_count = sum(1 for cand in candidates if _candidate_has_source_group(cand))
    _unlink_unbudgeted_candidate_images(cand.image_path for cand in dropped)
    log(
        "ingest.pdf.table_parse.budget",
        mode=_poster_harness_mode(ctx) or "default",
        original=len(candidates),
        kept=len(kept),
        dropped=len(dropped),
        max_candidates=cap,
        preferred_source_groups=preferred_count,
        preferred_source_groups_kept=sum(1 for cand in kept if _candidate_has_source_group(cand)),
        parallelism=_ingest_vlm_parallelism(ctx, _CAPTION_MATCH_PARALLELISM),
    )
    return kept


def _budget_candidate_indices(
    n_candidates: int,
    cap: int,
    *,
    page_of: Any,
    area_of: Any,
) -> set[int]:
    if cap <= 0 or n_candidates <= cap:
        return set(range(n_candidates))
    early_quota = max(1, int(round(cap * 0.55)))
    early = sorted(
        range(n_candidates),
        key=lambda i: (int(page_of(i) or 0), -int(area_of(i) or 0), i),
    )[:early_quota]
    keep = set(early)
    for idx in sorted(range(n_candidates), key=lambda i: (int(area_of(i) or 0), -i), reverse=True):
        keep.add(idx)
        if len(keep) >= cap:
            break
    return keep


def _budget_source_candidate_indices(
    candidates: list[Any],
    cap: int,
    *,
    page_of: Any,
    area_of: Any,
) -> set[int]:
    if cap <= 0 or len(candidates) <= cap:
        return set(range(len(candidates)))
    keep: set[int] = set()
    preferred = [
        i for i, cand in enumerate(candidates)
        if _candidate_has_source_group(cand)
    ]
    for idx in sorted(
        preferred,
        key=lambda i: _source_candidate_budget_key(
            candidates[i],
            page=_safe_int(page_of(candidates[i]), default=999),
            area=_safe_int(area_of(candidates[i]), default=0),
        ),
        reverse=True,
    ):
        keep.add(idx)
        if len(keep) >= cap:
            return keep
    baseline = _budget_candidate_indices(
        len(candidates),
        cap,
        page_of=lambda i: page_of(candidates[i]),
        area_of=lambda i: area_of(candidates[i]),
    )
    for idx in sorted(
        baseline,
        key=lambda i: _source_candidate_budget_key(
            candidates[i],
            page=_safe_int(page_of(candidates[i]), default=999),
            area=_safe_int(area_of(candidates[i]), default=0),
        ),
        reverse=True,
    ):
        keep.add(idx)
        if len(keep) >= cap:
            return keep
    for idx in sorted(
        range(len(candidates)),
        key=lambda i: _source_candidate_budget_key(
            candidates[i],
            page=_safe_int(page_of(candidates[i]), default=999),
            area=_safe_int(area_of(candidates[i]), default=0),
        ),
        reverse=True,
    ):
        keep.add(idx)
        if len(keep) >= cap:
            break
    return keep


def _source_candidate_budget_key(candidate: Any, *, page: int, area: int) -> tuple[int, int, int, int, int]:
    label = str(
        getattr(candidate, "source_group_label", "")
        or getattr(candidate, "anchor_label", "")
        or ""
    ).strip()
    caption = str(getattr(candidate, "source_group_caption", "") or "").strip()
    return (
        1 if _candidate_has_source_group(candidate) else 0,
        1 if caption else 0,
        1 if label else 0,
        min(area, 12_000_000),
        -page,
    )


def _candidate_has_source_group(candidate: Any) -> bool:
    if bool(getattr(candidate, "captioned_source_group", False)):
        return True
    if str(getattr(candidate, "source_group_id", "") or "").strip():
        return True
    if str(getattr(candidate, "source_group_label", "") or "").strip():
        return True
    if str(getattr(candidate, "source_group_caption", "") or "").strip():
        return True
    if not bool(getattr(candidate, "protected_anchor", False)):
        return False
    return bool(
        str(getattr(candidate, "anchor_label", "") or "").strip()
        or str(getattr(candidate, "anchor_reason", "") or "").strip()
    )


def _unlink_unbudgeted_candidate_images(paths: Any) -> None:
    for path in paths:
        try:
            Path(path).unlink()
        except OSError:
            pass


def _clean_space(s: str) -> str:
    return re.sub(r"\s+", " ", str(s or "")).strip()


def _match_captions_parallel(
    candidates: list[PdfFigureCandidate],
    manifest: dict[str, Any],
    ctx: ToolContext,
) -> dict[int, dict[str, Any]]:
    """Run caption matching for each candidate in a thread pool.

    Returns a dict keyed by candidate index → match result dict
    (`matched_idx`, `confidence`, `is_real_figure`, `reason`, and
    `caption_text` filled in from the manifest for convenience).
    """
    _raise_if_ingest_cancelled(ctx, "ingest.caption_match.start")
    all_figs = list(manifest.get("figures", []))
    results: dict[int, dict[str, Any]] = {}
    used_figures: set[int] = set()

    for cand_idx, candidate in enumerate(candidates):
        _raise_if_ingest_cancelled(ctx, "ingest.caption_match.before_local_merge")
        if not _candidate_has_source_group(candidate):
            continue
        matched_idx = _source_group_manifest_index(candidate, all_figs, used_figures)
        if matched_idx is not None:
            used_figures.add(matched_idx)
        caption = str(getattr(candidate, "source_group_caption", "") or "").strip()
        if matched_idx is not None:
            caption = str(all_figs[matched_idx].get("caption") or caption)
        results[cand_idx] = {
            "matched_idx": matched_idx,
            "confidence": 0.95,
            "is_real_figure": True,
            "reason": "bound from PDF captioned source group",
            "caption_text": caption,
            "short_caption": _fallback_short_caption(caption),
            "sub_panels": [],
            "caption_association_method": "captioned_group",
        }
        _raise_if_ingest_cancelled(ctx, "ingest.caption_match.after_local_merge")

    geometry_options = _caption_geometry_options(candidates, all_figs, used_figures)
    for cand_idx, fig_idx, score in _high_confidence_geometry_assignments(
        geometry_options,
        used_figures,
    ):
        _raise_if_ingest_cancelled(ctx, "ingest.caption_match.before_geometry_merge")
        if cand_idx in results or fig_idx in used_figures:
            continue
        used_figures.add(fig_idx)
        caption = str(all_figs[fig_idx].get("caption") or "")
        results[cand_idx] = {
            "matched_idx": fig_idx,
            "confidence": round(max(0.70, min(0.94, score)), 3),
            "is_real_figure": True,
            "reason": "high-confidence same-page caption geometry",
            "caption_text": caption,
            "short_caption": _fallback_short_caption(caption),
            "sub_panels": [],
            "caption_association_method": "geometry",
        }
        _raise_if_ingest_cancelled(ctx, "ingest.caption_match.after_geometry_merge")

    pending = [idx for idx in range(len(candidates)) if idx not in results]
    if not all_figs:
        for idx in pending:
            _raise_if_ingest_cancelled(ctx, "ingest.caption_match.before_empty_merge")
            results[idx] = _unmatched_caption_result("no captions in manifest")
            _raise_if_ingest_cancelled(ctx, "ingest.caption_match.after_empty_merge")
        _raise_if_ingest_cancelled(ctx, "ingest.caption_match.before_summary")
        _log_caption_association_summary(results)
        _raise_if_ingest_cancelled(ctx, "ingest.caption_match.after_summary")
        return results

    parallelism = _ingest_vlm_parallelism(ctx, _CAPTION_MATCH_PARALLELISM)
    _raise_if_ingest_cancelled(ctx, "ingest.caption_match.before_start_log")
    log("ingest.pdf.caption_match.start",
        n_candidates=len(candidates), n_vlm_candidates=len(pending), parallelism=parallelism,
        model=ctx.settings.ingest_model)
    _raise_if_ingest_cancelled(ctx, "ingest.caption_match.after_start_log")
    raw_matches: dict[int, dict[str, Any]] = {}

    def match_one(i: int, candidate: PdfFigureCandidate) -> dict[str, Any]:
        _raise_if_ingest_cancelled(ctx, "ingest.caption_match.worker.start")
        return _match_one_caption(i, candidate, all_figs, ctx)

    outcomes = _run_bounded_ingest_pool(
        items=[(i, candidates[i]) for i in pending],
        worker=match_one,
        parallelism=parallelism,
        ctx=ctx,
        phase="ingest.caption_match.pool",
    )
    for i, match, error in outcomes:
        _raise_if_ingest_cancelled(ctx, "ingest.caption_match.before_raw_merge")
        if error is not None:
            _raise_if_ingest_cancelled(ctx, "ingest.caption_match.before_failure_log")
            log("ingest.pdf.caption_match_fail", cand_idx=i, error=str(error))
            _raise_if_ingest_cancelled(ctx, "ingest.caption_match.after_failure_log")
            match = _unmatched_caption_result(f"match call failed: {error}")
        raw_matches[i] = match
        _raise_if_ingest_cancelled(ctx, "ingest.caption_match.after_raw_merge")

    vlm_proposals = sorted(
        (
            (_safe_float(match.get("confidence"), default=0.0), i, int(match["matched_idx"]))
            for i, match in raw_matches.items()
            if isinstance(match.get("matched_idx"), int)
            and _safe_float(match.get("confidence"), default=0.0) >= _CAPTION_MATCH_MIN_CONFIDENCE
            and bool(match.get("is_real_figure", True))
        ),
        key=lambda item: (-item[0], item[1], item[2]),
    )
    for _confidence, cand_idx, fig_idx in vlm_proposals:
        _raise_if_ingest_cancelled(ctx, "ingest.caption_match.before_vlm_merge")
        if cand_idx in results or fig_idx in used_figures:
            continue
        used_figures.add(fig_idx)
        match = raw_matches[cand_idx]
        match["caption_association_method"] = "vlm"
        results[cand_idx] = match
        _raise_if_ingest_cancelled(ctx, "ingest.caption_match.after_vlm_merge")

    fallback_order = sorted(
        (i for i in pending if i not in results),
        key=lambda i: (
            -max((score for _fig_idx, score in geometry_options.get(i, [])), default=0.0),
            i,
        ),
    )
    for i in fallback_order:
        _raise_if_ingest_cancelled(ctx, "ingest.caption_match.before_fallback_merge")
        match = raw_matches.get(i) or _unmatched_caption_result("VLM returned no result")
        fallback_idx = _unique_geometry_fallback_index(
            i,
            geometry_options,
            used_figures,
        )
        if fallback_idx is not None and bool(match.get("is_real_figure", True)):
            used_figures.add(fallback_idx)
            caption = str(all_figs[fallback_idx].get("caption") or "")
            reason = str(match.get("reason") or "VLM did not return a usable match")
            results[i] = {
                "matched_idx": fallback_idx,
                "confidence": 0.55,
                "is_real_figure": True,
                "reason": f"{reason}; unique same-page geometry fallback"[:300],
                "caption_text": caption,
                "short_caption": _fallback_short_caption(caption),
                "sub_panels": [],
                "caption_association_method": "geometry_fallback",
            }
        else:
            match["matched_idx"] = None
            match["confidence"] = 0.0
            match["caption_text"] = ""
            match["short_caption"] = ""
            match["caption_association_method"] = "unmatched"
            results[i] = match
        _raise_if_ingest_cancelled(ctx, "ingest.caption_match.after_fallback_merge")
    _raise_if_ingest_cancelled(ctx, "ingest.caption_match.before_summary")
    _log_caption_association_summary(results)
    _raise_if_ingest_cancelled(ctx, "ingest.caption_match.after_summary")
    return results


def _unmatched_caption_result(reason: str) -> dict[str, Any]:
    return {
        "matched_idx": None,
        "confidence": 0.0,
        "is_real_figure": True,
        "reason": reason,
        "caption_text": "",
        "short_caption": "",
        "sub_panels": [],
        "caption_association_method": "unmatched",
    }


def _source_group_manifest_index(
    candidate: PdfFigureCandidate,
    figures: list[dict[str, Any]],
    used: set[int],
) -> int | None:
    label = str(getattr(candidate, "source_group_label", "") or getattr(candidate, "anchor_label", ""))
    label_match = re.search(r"(?i)(?:figure|fig\.?)?\s*([A-Za-z]?\d+(?:\.\d+)?)", label)
    wanted_label = label_match.group(1).lower() if label_match else ""
    wanted_caption = _clean_space(str(getattr(candidate, "source_group_caption", "") or "")).lower()
    for idx, figure in enumerate(figures):
        if idx in used:
            continue
        caption = _clean_space(str(figure.get("caption") or ""))
        if wanted_caption and caption.lower() == wanted_caption:
            return idx
        caption_match = re.search(r"(?i)\b(?:figure|fig\.?)\s*([A-Za-z]?\d+(?:\.\d+)?)", caption)
        if wanted_label and caption_match and caption_match.group(1).lower() == wanted_label:
            return idx
    return None


def _caption_geometry_options(
    candidates: list[PdfFigureCandidate],
    figures: list[dict[str, Any]],
    used_figures: set[int],
) -> dict[int, list[tuple[int, float]]]:
    options: dict[int, list[tuple[int, float]]] = {}
    for cand_idx, candidate in enumerate(candidates):
        if candidate.bbox_pt is None:
            continue
        candidate_rect = _bbox4(candidate.bbox_pt)
        if candidate_rect is None:
            continue
        for fig_idx, figure in enumerate(figures):
            if fig_idx in used_figures or _safe_int(figure.get("page"), default=0) != candidate.page:
                continue
            caption_rect = _manifest_caption_bbox(figure)
            if caption_rect is None:
                continue
            overlap = _bbox_horizontal_overlap(candidate_rect, caption_rect)
            distance = _bbox_vertical_distance(candidate_rect, caption_rect)
            if overlap < 0.20 or distance > 180.0:
                continue
            score = 0.65 * overlap + 0.35 * max(0.0, 1.0 - distance / 180.0)
            options.setdefault(cand_idx, []).append((fig_idx, score))
    for values in options.values():
        values.sort(key=lambda item: item[1], reverse=True)
    return options


def _high_confidence_geometry_assignments(
    options: dict[int, list[tuple[int, float]]],
    used_figures: set[int],
) -> list[tuple[int, int, float]]:
    confident: list[tuple[int, int, float, float]] = []
    for cand_idx, values in options.items():
        available = [item for item in values if item[0] not in used_figures]
        if not available:
            continue
        best_idx, best_score = available[0]
        second_score = available[1][1] if len(available) > 1 else 0.0
        margin = best_score - second_score
        if best_score >= _CAPTION_GEOMETRY_HIGH_CONFIDENCE_MIN_SCORE and margin >= 0.15:
            confident.append((cand_idx, best_idx, best_score, margin))
    confident.sort(key=lambda item: (-item[3], -item[2], item[0], item[1]))
    assigned_candidates: set[int] = set()
    assigned_figures = set(used_figures)
    out: list[tuple[int, int, float]] = []
    for cand_idx, fig_idx, score, _margin in confident:
        if cand_idx in assigned_candidates or fig_idx in assigned_figures:
            continue
        assigned_candidates.add(cand_idx)
        assigned_figures.add(fig_idx)
        out.append((cand_idx, fig_idx, score))
    return out


def _unique_geometry_fallback_index(
    cand_idx: int,
    geometry_options: dict[int, list[tuple[int, float]]],
    used_figures: set[int],
) -> int | None:
    geometric = [fig_idx for fig_idx, _ in geometry_options.get(cand_idx, []) if fig_idx not in used_figures]
    if len(geometric) == 1:
        return geometric[0]
    return None


def _manifest_caption_bbox(figure: dict[str, Any]) -> tuple[float, float, float, float] | None:
    for key in ("caption_bbox_pdf_points", "caption_rect", "caption_bbox"):
        bbox = _bbox4(figure.get(key))
        if bbox is not None:
            return bbox
    return None


def _bbox4(value: Any) -> tuple[float, float, float, float] | None:
    if not isinstance(value, (list, tuple)) or len(value) < 4:
        return None
    try:
        x0, y0, x1, y1 = [float(item) for item in value[:4]]
    except (TypeError, ValueError):
        return None
    if x1 <= x0 or y1 <= y0:
        return None
    return x0, y0, x1, y1


def _bbox_horizontal_overlap(
    left: tuple[float, float, float, float],
    right: tuple[float, float, float, float],
) -> float:
    overlap = max(0.0, min(left[2], right[2]) - max(left[0], right[0]))
    return overlap / max(1.0, min(left[2] - left[0], right[2] - right[0]))


def _bbox_vertical_distance(
    left: tuple[float, float, float, float],
    right: tuple[float, float, float, float],
) -> float:
    if right[1] >= left[3]:
        return right[1] - left[3]
    if left[1] >= right[3]:
        return left[1] - right[3]
    return 0.0


def _log_caption_association_summary(results: dict[int, dict[str, Any]]) -> None:
    methods = {method: 0 for method in ("captioned_group", "geometry", "vlm", "geometry_fallback", "unmatched")}
    failures: dict[str, int] = {}
    for result in results.values():
        method = str(result.get("caption_association_method") or "unmatched")
        methods[method] = methods.get(method, 0) + 1
        if method in {"geometry_fallback", "unmatched"}:
            reason = str(result.get("reason") or "unknown")[:160]
            failures[reason] = failures.get(reason, 0) + 1
    log("ingest.pdf.caption_match.summary", methods=methods, failure_reasons=failures)


def _match_one_caption(
    cand_idx: int,
    candidate: PdfFigureCandidate,
    all_figs: list[dict[str, Any]],
    ctx: ToolContext,
) -> dict[str, Any]:
    _raise_if_ingest_cancelled(ctx, "ingest.caption_match.worker.before_prepare")
    # Filter captions to the candidate's page ± 1, fall back to whole
    # manifest if the page window is empty (handles figures whose caption
    # ends up on the next page due to column overflow).
    near: list[tuple[int, dict[str, Any]]] = [
        (i, f) for i, f in enumerate(all_figs)
        if abs(int(f.get("page", 0)) - candidate.page) <= 1
    ]
    pool = near if near else list(enumerate(all_figs))
    # Indices in `pool` are *relative* to the pool list we show the VLM;
    # we need to remap back to the full `all_figs` index.
    local_to_global = {local_i: global_i for local_i, (global_i, _) in enumerate(pool)}

    lines = []
    for local_i, (_, fig) in enumerate(pool):
        cap = (fig.get("caption") or "").replace("\n", " ")[:240]
        lines.append(f"  [{local_i}] (p.{fig.get('page', '?')}) {cap}")
    user_text = (
        f"Candidate captions near page {candidate.page} "
        f"(strategy: {candidate.strategy}):\n"
        + "\n".join(lines)
        + "\n\nReturn the JSON described in the system prompt."
    )

    _raise_if_ingest_cancelled(ctx, "ingest.caption_match.worker.before_image_read")
    image = VlmImage.from_path(candidate.path)
    _raise_if_ingest_cancelled(ctx, "ingest.caption_match.worker.after_image_read")
    _raise_if_ingest_cancelled(ctx, "ingest.caption_match.worker.before_vlm")
    result = vlm_call_json(
        settings=ctx.settings,
        model=ctx.settings.ingest_model,
        system=_CAPTION_MATCH_PROMPT,
        user_text=user_text,
        images=[image],
        max_tokens=512,
        cancellation_token=_ingest_cancellation_token(ctx),
    )
    _raise_if_ingest_cancelled(ctx, "ingest.caption_match.worker.after_vlm")

    # Remap local index → global manifest index, attach caption text.
    local_idx = result.get("matched_idx")
    global_idx = None
    caption_text = ""
    if isinstance(local_idx, int) and local_idx in local_to_global:
        global_idx = local_to_global[local_idx]
        caption_text = str(all_figs[global_idx].get("caption", ""))

    sub_panels_raw = result.get("sub_panels") or []
    if isinstance(sub_panels_raw, list) and sub_panels_raw:
        _raise_if_ingest_cancelled(ctx, "ingest.caption_match.worker.before_subpanel_log")
        log(
            "ingest.pdf.sub_panel.vlm_ignored",
            candidate=candidate.path.name,
            requested=len(sub_panels_raw),
            reason="partial_crops_disabled",
        )
        _raise_if_ingest_cancelled(ctx, "ingest.caption_match.worker.after_subpanel_log")

    _raise_if_ingest_cancelled(ctx, "ingest.caption_match.worker.before_return")
    return {
        "matched_idx": global_idx,
        "confidence": float(result.get("confidence", 0.0) or 0.0),
        "is_real_figure": bool(result.get("is_real_figure", True)),
        "reason": str(result.get("reason", ""))[:200],
        "caption_text": caption_text,
        # v2.3 — VLM-generated ≤15-char label for tight poster/deck slots.
        # Clipped defensively; empty string on fake figures or missing field.
        "short_caption": str(result.get("short_caption", "") or "")[:40],
        "sub_panels": [],
    }


def _register_sub_panels(
    *,
    parent_layer_id: str,
    parent_path: Path,
    parent_caption: str,
    panels: list[dict[str, Any]],
    ctx: ToolContext,
    pdf_path: Path,
    source_page: int | None,
    source_pdf_sha256: str | None = None,
    parent_bbox_pdf_points: tuple[float, float, float, float] | None = None,
) -> list[str]:
    """Crop each VLM-detected sub-panel out of the parent figure PNG and
    register as its own `rendered_layers` entry. Returns the list of newly
    created layer_ids (empty list when `panels` is empty or all panels
    fail validation).

    The sub-panel layer_id is `{parent}_{label}` (e.g. `ingest_fig_02_a`)
    so the naming convention carries the parent relationship without
    needing a new schema field. Failed crops (bbox outside source, Pillow
    error) are skipped silently with a log event.
    """
    if not panels:
        return []

    from PIL import Image as _Image

    try:
        parent_img = _Image.open(parent_path)
        parent_img.load()  # force decode so we can close the file handle
    except Exception as e:
        log("ingest.pdf.sub_panel.open_fail",
            parent=parent_layer_id, error=str(e))
        return []

    pw, ph = parent_img.size
    registered: list[str] = []
    for panel in panels:
        label = panel.get("label", "")
        x0, y0, x1, y1 = panel["bbox"]
        # Clamp bbox to image bounds defensively.
        x0 = max(0, min(x0, pw - 1))
        y0 = max(0, min(y0, ph - 1))
        x1 = max(x0 + 1, min(x1, pw))
        y1 = max(y0 + 1, min(y1, ph))
        if (x1 - x0) < 20 or (y1 - y0) < 20:
            log("ingest.pdf.sub_panel.too_small",
                parent=parent_layer_id, label=label,
                bbox=[x0, y0, x1, y1])
            continue

        sub_layer_id = f"{parent_layer_id}_{label}"
        # Avoid collisions if the VLM returned dup labels.
        if sub_layer_id in ctx.state["rendered_layers"]:
            continue

        sub_path = ctx.layers_dir / f"img_{sub_layer_id}.png"
        try:
            crop = parent_img.crop((x0, y0, x1, y1))
            crop.save(sub_path, "PNG", optimize=True)
        except Exception as e:
            log("ingest.pdf.sub_panel.crop_fail",
                parent=parent_layer_id, label=label, error=str(e))
            continue

        cw = x1 - x0
        ch = y1 - y0
        ctx.state["rendered_layers"][sub_layer_id] = {
            "layer_id": sub_layer_id,
            "name": f"{parent_layer_id}_{label}",
            "kind": "image",
            "z_index": 5,
            "bbox": None,
            "src_path": str(sub_path),
            "aspect_ratio": _aspect_from_dims(cw, ch),
            "image_size": f"{cw}x{ch}",
            "sha256": sha256_file(sub_path),
            "source": "ingested_pdf",
            "source_file": str(pdf_path),
            "source_pdf": pdf_path.name,
            "source_pdf_sha256": source_pdf_sha256,
            "source_page": source_page,
            "source_bbox_pdf_points": _sub_panel_pdf_bbox(
                parent_bbox_pdf_points,
                (pw, ph),
                (x0, y0, x1, y1),
            ),
            "caption": panel.get("caption", "") or parent_caption,
            "caption_short": panel.get("short_caption", ""),
            "extract_strategy": "sub_panel",
            "caption_confidence": 0.8,   # VLM-provided bboxes we trust more than auto-cluster
            # naming convention — no schema field; planner reads the `_a`/`_b`
            # suffix + `extract_strategy="sub_panel"` marker.
            "parent_layer_id": parent_layer_id,
        }
        registered.append(sub_layer_id)

    try:
        parent_img.close()
    except Exception:
        pass

    log("ingest.pdf.sub_panel.register",
        parent=parent_layer_id,
        requested=len(panels), kept=len(registered))
    return registered


def _sub_panel_pdf_bbox(
    parent_bbox_pdf_points: tuple[float, float, float, float] | None,
    parent_size_px: tuple[int, int],
    panel_bbox_px: tuple[int, int, int, int],
) -> list[float] | None:
    """Map a sub-panel crop bbox from parent PNG pixels back to PDF points."""
    if not parent_bbox_pdf_points:
        return None
    pw, ph = parent_size_px
    if pw <= 0 or ph <= 0:
        return None
    px0, py0, px1, py1 = panel_bbox_px
    bx0, by0, bx1, by1 = parent_bbox_pdf_points
    bw = bx1 - bx0
    bh = by1 - by0
    if bw <= 0 or bh <= 0:
        return None
    return [
        round(bx0 + (px0 / pw) * bw, 2),
        round(by0 + (py0 / ph) * bh, 2),
        round(bx0 + (px1 / pw) * bw, 2),
        round(by0 + (py1 / ph) * bh, 2),
    ]


def _refine_vector_candidate_bbox(
    candidate: PdfFigureCandidate,
    match: dict[str, Any],
    doc: fitz.Document,
    ctx: ToolContext | None = None,
) -> PdfFigureCandidate | None:
    _raise_if_ingest_cancelled(ctx, "ingest.vector_refine.start")
    if candidate.strategy != "vector" or candidate.bbox_pt is None:
        return candidate
    if candidate.page < 1 or candidate.page > len(doc):
        return candidate

    page = doc[candidate.page - 1]
    original = fitz.Rect(candidate.bbox_pt) & page.rect
    if original.is_empty:
        return candidate

    caption = str(match.get("caption_text") or "").strip()
    caption_block = _matched_caption_block(page, original, caption)
    upper_noise_trim_y = _vector_upper_page_noise_trim_y(page, original)
    if caption_block is None and _count_figure_caption_blocks(page, original) >= 2:
        log(
            "ingest.pdf.vector_crop.reject_multi_caption",
            page=candidate.page,
            path=candidate.path.name,
            bbox=[round(v, 2) for v in original],
            reason="unmatched vector cluster spans multiple figure captions",
        )
        return None
    text_boundary = _vector_lower_text_boundary(page, original, caption, caption_block)
    if caption_block is None and text_boundary is None and upper_noise_trim_y is None:
        return candidate

    refined = fitz.Rect(original)
    if upper_noise_trim_y is not None and original.y0 + 1 < upper_noise_trim_y < original.y1 - 40:
        refined.y0 = max(refined.y0, upper_noise_trim_y)

    # When a caption/prose block sits below the drawing, use it as the lower
    # boundary. This catches vector clusters whose bbox merged the plot with
    # the figure caption and following body paragraph.
    lower_boundary = text_boundary or caption_block
    if lower_boundary is not None and original.y0 + 20 < lower_boundary.y0 <= original.y1 + 120:
        refined.y1 = min(refined.y1, lower_boundary.y0 - 4)

    # Single-column figures can be merged with the other text column, but a
    # short centered caption is common under a complete wide figure.
    if caption_block is not None and _caption_supports_horizontal_crop(original, caption_block, page.rect):
        refined.x0 = max(refined.x0, caption_block.x0 - 4)
        refined.x1 = min(refined.x1, caption_block.x1 + 4)

    refined &= page.rect
    if not _refined_bbox_is_useful(original, refined):
        return candidate

    try:
        _raise_if_ingest_cancelled(ctx, "ingest.vector_refine.before_render")
        pix = page.get_pixmap(clip=refined, dpi=300)
        _raise_if_ingest_cancelled(ctx, "ingest.vector_refine.after_render")
        _raise_if_ingest_cancelled(ctx, "ingest.vector_refine.before_save")
        pix.save(str(candidate.path))
        _raise_if_ingest_cancelled(ctx, "ingest.vector_refine.after_save")
    except Exception as e:  # noqa: BLE001 - keep original crop on render faults
        _raise_if_ingest_cancelled(ctx, "ingest.vector_refine.before_failure_log")
        log(
            "ingest.pdf.vector_crop.refine_fail",
            page=candidate.page,
            path=candidate.path.name,
            error=f"{type(e).__name__}: {e}"[:200],
        )
        _raise_if_ingest_cancelled(ctx, "ingest.vector_refine.after_failure_log")
        return candidate

    _raise_if_ingest_cancelled(ctx, "ingest.vector_refine.before_refined_log")
    log(
        "ingest.pdf.vector_crop.refined",
        page=candidate.page,
        path=candidate.path.name,
        old_bbox=[round(v, 2) for v in original],
        new_bbox=[round(v, 2) for v in refined],
        old_size=f"{candidate.width_px}x{candidate.height_px}",
        new_size=f"{pix.width}x{pix.height}",
        upper_noise_trim_y=round(upper_noise_trim_y, 2) if upper_noise_trim_y is not None else None,
        text_boundary=[round(v, 2) for v in lower_boundary] if lower_boundary is not None else None,
    )
    _raise_if_ingest_cancelled(ctx, "ingest.vector_refine.after_refined_log")
    return PdfFigureCandidate(
        page=candidate.page,
        bbox_pt=(round(refined.x0, 2), round(refined.y0, 2),
                 round(refined.x1, 2), round(refined.y1, 2)),
        path=candidate.path,
        width_px=pix.width,
        height_px=pix.height,
        strategy=candidate.strategy,
        xref=candidate.xref,
    )


def _vector_upper_page_noise_trim_y(
    page: fitz.Page,
    candidate_bbox: fitz.Rect,
) -> float | None:
    """Return a top crop boundary below running page headers/rules.

    Vector clustering can attach a paper's running title separator to a
    nearby chart. Keep this conservative: only trim page-top, page-wide text
    or rules that overlap the candidate's top edge.
    """
    page_rect = page.rect
    max_top_y = min(82.0, page_rect.height * 0.13)
    trim_candidates: list[float] = []

    for block in page.get_text("blocks"):
        rect = fitz.Rect(block[:4])
        if rect.y0 > max_top_y or rect.y1 < candidate_bbox.y0 - 12:
            continue
        if rect.y0 > candidate_bbox.y0 + 42:
            continue
        if _rect_horizontal_overlap(candidate_bbox, rect) < 0.20:
            continue
        text = _clean_space(str(block[4] or ""))
        if not _looks_like_running_page_header_text(text, rect, page_rect):
            continue
        trim_candidates.append(rect.y1 + 3.0)

    for drawing in page.get_drawings():
        rect_obj = drawing.get("rect")
        if rect_obj is None:
            continue
        rect = fitz.Rect(rect_obj)
        if not _looks_like_running_page_header_rule(rect, page_rect):
            continue
        if rect.y1 < candidate_bbox.y0 - 12 or rect.y0 > candidate_bbox.y0 + 42:
            continue
        if _rect_horizontal_overlap(candidate_bbox, rect) < 0.20:
            continue
        trim_candidates.append(rect.y1 + 3.0)

    if not trim_candidates:
        return None
    trim_y = max(trim_candidates)
    if trim_y <= candidate_bbox.y0 + 1:
        return None
    if trim_y - candidate_bbox.y0 > max(48.0, candidate_bbox.height * 0.28):
        return None
    return trim_y


def _looks_like_running_page_header_text(
    text: str,
    rect: fitz.Rect,
    page_rect: fitz.Rect,
) -> bool:
    cleaned = _clean_space(text)
    if not cleaned:
        return False
    if rect.width < page_rect.width * 0.35:
        return False
    if rect.height > 32:
        return False
    if re.match(r"(?i)^\s*(?:figure|fig\.?|table)\s+[A-Za-z]?\d+", cleaned):
        return False
    if _looks_like_code_or_formula_text(cleaned):
        return False
    words = re.findall(r"[A-Za-z][A-Za-z-]{2,}", cleaned)
    if len(words) < 3:
        return False
    alpha = sum(1 for ch in cleaned if ch.isalpha())
    if alpha < 12:
        return False
    return True


def _looks_like_running_page_header_rule(
    rect: fitz.Rect,
    page_rect: fitz.Rect,
) -> bool:
    if rect.width <= 0:
        return False
    if rect.y0 > min(84.0, page_rect.height * 0.13):
        return False
    if rect.width < page_rect.width * 0.45:
        return False
    if rect.height > 3.0:
        return False
    return True


def _matched_caption_block(
    page: fitz.Page,
    candidate_bbox: fitz.Rect,
    caption: str,
) -> fitz.Rect | None:
    label = _figure_caption_label(caption)
    if not label:
        return None
    label_re = re.compile(rf"(?i)^\s*(?:figure|fig\.?)\s*{re.escape(label)}\b")
    best: tuple[float, fitz.Rect] | None = None
    for block in page.get_text("blocks"):
        rect = fitz.Rect(block[:4])
        text = _clean_space(str(block[4] or ""))
        if not label_re.search(text):
            continue
        caption_evidence = _caption_header_bonus(text, caption, label)
        overlap_evidence = _caption_text_overlap_score(text, caption)
        if caption_evidence <= 0 and overlap_evidence < 0.24:
            continue
        if rect.x1 < candidate_bbox.x0 - 80 or rect.x0 > candidate_bbox.x1 + 80:
            continue
        if rect.y0 < candidate_bbox.y0 - 20 or rect.y0 > candidate_bbox.y1 + 120:
            continue
        overlap = _rect_horizontal_overlap(candidate_bbox, rect)
        distance = abs(rect.y0 - candidate_bbox.y1)
        score = overlap * 1000.0 - distance
        score += caption_evidence + overlap_evidence * 2000.0
        if best is None or score > best[0]:
            best = (score, rect)
    return best[1] if best is not None else None


def _vector_lower_text_boundary(
    page: fitz.Page,
    candidate_bbox: fitz.Rect,
    caption: str,
    caption_block: fitz.Rect | None,
) -> fitz.Rect | None:
    label = _figure_caption_label(caption)
    min_y = candidate_bbox.y0 + min(72.0, max(36.0, candidate_bbox.height * 0.10))
    best: tuple[float, fitz.Rect] | None = None
    for block in page.get_text("blocks"):
        rect = fitz.Rect(block[:4])
        if rect.x1 < candidate_bbox.x0 - 12 or rect.x0 > candidate_bbox.x1 + 12:
            continue
        if rect.y0 < min_y or rect.y0 > candidate_bbox.y1 + 120:
            continue
        if _rect_horizontal_overlap(candidate_bbox, rect) < 0.35:
            continue
        text = _clean_space(str(block[4] or ""))
        if not _looks_like_external_figure_text(text, rect, candidate_bbox, label):
            continue
        score = -rect.y0
        header_bonus = _caption_header_bonus(text, caption, label)
        if caption_block is not None and abs(rect.y0 - caption_block.y0) < 2:
            header_bonus += 5000.0
        score += header_bonus
        if best is None or score > best[0]:
            best = (score, rect)
    return best[1] if best is not None else None


def _caption_header_bonus(text: str, caption: str, label: str) -> float:
    if not label:
        return 0.0
    cleaned = _clean_space(text)
    bonus = 0.0
    if re.match(rf"(?i)^\s*(?:figure|fig\.?)\s*{re.escape(label)}\s*[:.)-]", cleaned):
        bonus += 8000.0
    caption_prefix = _caption_prefix(caption)
    if caption_prefix and _clean_space(cleaned).lower().startswith(caption_prefix):
        bonus += 12000.0
    return bonus


def _caption_prefix(caption: str) -> str:
    cleaned = _clean_space(caption).lower()
    if not cleaned:
        return ""
    return cleaned[: min(72, len(cleaned))]


def _caption_text_overlap_score(text: str, caption: str) -> float:
    text_tokens = _caption_match_tokens(text)
    caption_tokens = _caption_match_tokens(caption)
    if not text_tokens or not caption_tokens:
        return 0.0
    return len(text_tokens & caption_tokens) / max(1, min(len(text_tokens), len(caption_tokens)))


def _caption_match_tokens(text: str) -> set[str]:
    stop = {
        "figure", "fig", "the", "and", "for", "with", "into", "from", "that",
        "this", "their", "are", "was", "were", "been", "have", "has", "our",
    }
    return {
        token.lower()
        for token in re.findall(r"[A-Za-z][A-Za-z0-9-]{2,}", text or "")
        if token.lower() not in stop
    }


def _looks_like_external_figure_text(
    text: str,
    rect: fitz.Rect,
    candidate_bbox: fitz.Rect,
    label: str,
) -> bool:
    cleaned = _clean_space(text)
    if not cleaned:
        return False
    if label and re.match(rf"(?i)^\s*(?:figure|fig\.?)\s*{re.escape(label)}\s*[:.)-]", cleaned):
        return True
    if len(cleaned) < 56:
        return False
    if _looks_like_code_or_formula_text(cleaned):
        return False
    if rect.width < min(150.0, candidate_bbox.width * 0.38):
        return False
    words = re.findall(r"[A-Za-z][A-Za-z-]+", cleaned)
    if len(words) < 14:
        return False
    alpha = sum(1 for ch in cleaned if ch.isalpha())
    digit = sum(1 for ch in cleaned if ch.isdigit())
    if alpha < max(24, digit * 2):
        return False
    if not re.search(r"[a-z][.!?](?:\s|$)", cleaned):
        return False
    return any(ch in cleaned for ch in ".:,;()")


def _looks_like_code_or_formula_text(text: str) -> bool:
    cleaned = _clean_space(text)
    if not cleaned:
        return False
    code_markers = (
        "image_encoder", "text_encoder", "np.", "logits", "softmax",
        "cross_entropy", "normalize", "return ", "def ", "lambda",
    )
    if any(marker in cleaned for marker in code_markers):
        return True
    symbol_count = sum(1 for ch in cleaned if ch in "_=#[]{}<>\\/")
    alpha = sum(1 for ch in cleaned if ch.isalpha())
    return symbol_count >= 4 and symbol_count / max(1, alpha) > 0.06


def _figure_caption_label(caption: str) -> str:
    match = re.search(r"(?i)\b(?:figure|fig\.?)\s*([0-9]+[A-Za-z]?)\b", caption or "")
    return match.group(1) if match else ""


def _caption_label_for_kind(caption: str, kind: str) -> str:
    prefix = r"table" if kind == "table" else r"fig(?:ure)?\.?"
    match = re.search(rf"(?i)\b{prefix}\s*([0-9]+[A-Za-z]?)\b", caption or "")
    return match.group(1) if match else ""


def _enrich_manifest_with_pdf_caption_groups(
    manifest: dict[str, Any],
    groups: list[Any],
) -> None:
    """Persist PDF-discovered caption geometry for deterministic matching."""
    for group in groups:
        kind = "table" if str(getattr(group, "kind", "")).lower() == "table" else "figure"
        key = "tables" if kind == "table" else "figures"
        label = str(getattr(group, "label", "") or "").lower()
        page = _safe_int(getattr(group, "page", 0), default=0)
        caption = str(getattr(group, "caption_text", "") or "").strip()
        caption_bbox = list(getattr(group, "caption_rect", ()) or ())
        items = manifest.setdefault(key, [])
        if not isinstance(items, list):
            items = []
            manifest[key] = items
        matched: dict[str, Any] | None = None
        for item in items:
            if not isinstance(item, dict):
                continue
            item_label = _caption_label_for_kind(str(item.get("caption") or ""), kind).lower()
            item_page = _safe_int(item.get("page"), default=0)
            if label and item_label == label and (not page or not item_page or item_page == page):
                matched = item
                break
        if matched is None:
            matched = {
                "caption": caption,
                "page": page,
                "description": "",
            }
            items.append(matched)
        elif caption and not str(matched.get("caption") or "").strip():
            matched["caption"] = caption
        if len(caption_bbox) == 4:
            matched["caption_bbox_pdf_points"] = caption_bbox
        matched["caption_source"] = "pdf_caption_block"


def _source_group_metadata_for_candidate(
    *,
    kind: str,
    candidate: Any,
    manifest: dict[str, Any] | None = None,
    caption: str = "",
    source: str = "",
) -> dict[str, Any]:
    kind = "table" if str(kind or "").lower() == "table" else "figure"
    manifest_match = _manifest_caption_for_source_candidate(kind, candidate, manifest or {})
    caption_text = (
        str(getattr(candidate, "source_group_caption", "") or "").strip()
        or str(caption or "").strip()
        or str(manifest_match.get("caption") or "").strip()
    )
    label = (
        str(getattr(candidate, "source_group_label", "") or "").strip()
        or str(getattr(candidate, "anchor_label", "") or "").strip()
        or _caption_label_for_kind(caption_text, kind)
    )
    page = _safe_int(getattr(candidate, "page", None), default=0)
    group_id = str(getattr(candidate, "source_group_id", "") or "").strip()
    if not group_id:
        if label:
            group_id = f"{kind}:p{page or 0}:{label.lower()}"
        else:
            bbox = getattr(candidate, "bbox_pt", None)
            if isinstance(bbox, tuple) and len(bbox) >= 4:
                group_id = (
                    f"{kind}:p{page or 0}:bbox:"
                    f"{round(float(bbox[0]) / 40)}:"
                    f"{round(float(bbox[1]) / 40)}:"
                    f"{round(float(bbox[2]) / 40)}:"
                    f"{round(float(bbox[3]) / 40)}"
                )
            elif page:
                group_id = f"{kind}:p{page}:candidate"
    source_text = (
        str(getattr(candidate, "source_group_source", "") or "").strip()
        or str(manifest_match.get("source") or "").strip()
        or str(source or "").strip()
        or str(getattr(candidate, "anchor_reason", "") or "").strip()
    )
    return {
        "source_group_id": group_id,
        "source_group_kind": kind,
        "source_group_label": label,
        "source_group_caption": caption_text,
        "source_group_source": source_text,
    }


def _manifest_caption_for_source_candidate(
    kind: str,
    candidate: Any,
    manifest: dict[str, Any],
) -> dict[str, str]:
    key = "tables" if kind == "table" else "figures"
    items = list((manifest or {}).get(key) or [])
    if not items:
        return {}
    label = str(getattr(candidate, "anchor_label", "") or "").strip().lower()
    page = _safe_int(getattr(candidate, "page", None), default=0)
    best: tuple[int, dict[str, Any]] | None = None
    for item in items:
        if not isinstance(item, dict):
            continue
        caption = str(item.get("caption") or item.get("title") or "")
        item_label = _caption_label_for_kind(caption, kind).lower()
        try:
            item_page = int(item.get("page") or 0)
        except (TypeError, ValueError):
            item_page = 0
        score = 0
        if label and item_label == label:
            score += 6
        if page and item_page:
            distance = abs(page - item_page)
            if distance == 0:
                score += 4
            elif distance == 1:
                score += 2
        if score <= 0:
            continue
        if best is None or score > best[0]:
            best = (score, item)
    if best is None:
        return {}
    matched = best[1]
    return {
        "caption": str(matched.get("caption") or matched.get("title") or ""),
        "source": f"manifest.{key}",
    }


def _caption_header_bonus_for_kind(text: str, caption: str, label: str, kind: str) -> float:
    kind = "table" if str(kind or "").lower() == "table" else "figure"
    if kind != "table":
        return _caption_header_bonus(text, caption, label)
    if not label:
        return 0.0
    cleaned = _clean_space(text)
    bonus = 0.0
    if re.match(rf"(?i)^\s*table\s*{re.escape(label)}\s*[:.)-]", cleaned):
        bonus += 8000.0
    caption_prefix = _caption_prefix(caption)
    if caption_prefix and cleaned.lower().startswith(caption_prefix):
        bonus += 12000.0
    return bonus


def _foreign_caption_label_in_text(text: str, *, label: str, kind: str) -> bool:
    current_kind = "table" if str(kind or "").lower() == "table" else "figure"
    current_label = str(label or "").strip().lower()
    for match in re.finditer(r"(?i)\b(figure|fig\.?|table)\s*([0-9]+[A-Za-z]?)\b", text or ""):
        found_kind = "table" if match.group(1).lower().startswith("table") else "figure"
        found_label = match.group(2).lower()
        if found_kind != current_kind or (current_label and found_label != current_label):
            return True
    return False


def _looks_like_running_page_header_text(text: str, rect: fitz.Rect, page_rect: fitz.Rect) -> bool:
    cleaned = _clean_space(text)
    if not cleaned:
        return False
    edge_band = max(48.0, page_rect.height * 0.075)
    if rect.y1 > page_rect.y0 + edge_band and rect.y0 < page_rect.y1 - edge_band:
        return False
    if rect.width < min(page_rect.width * 0.20, 150.0):
        return False
    words = re.findall(r"[A-Za-z][A-Za-z-]+", cleaned)
    if not 3 <= len(words) <= 22:
        return False
    if re.search(r"(?i)\b(?:figure|fig\.?|table)\s+[0-9]+", cleaned):
        return False
    return len(cleaned) <= 180


def _crop_has_page_furniture_geometry(page: fitz.Page, bbox: fitz.Rect) -> bool:
    page_rect = page.rect
    edge_band = max(48.0, page_rect.height * 0.075)
    crop_edge_zone = max(24.0, bbox.height * 0.16)

    def near_page_edge(rect: fitz.Rect) -> bool:
        return rect.y1 <= page_rect.y0 + edge_band or rect.y0 >= page_rect.y1 - edge_band

    def near_crop_edge(rect: fitz.Rect) -> bool:
        return rect.y1 <= bbox.y0 + crop_edge_zone or rect.y0 >= bbox.y1 - crop_edge_zone

    for block in page.get_text("dict").get("blocks", []):
        if block.get("type") != 1:
            continue
        try:
            rect = fitz.Rect(block.get("bbox"))
        except Exception:
            continue
        if not rect.intersects(bbox) or _rect_horizontal_overlap(bbox, rect) < 0.06:
            continue
        if not near_page_edge(rect) or not near_crop_edge(rect):
            continue
        if rect.height <= 42.0 or rect.width <= page_rect.width * 0.24:
            return True

    for drawing in page.get_drawings():
        rect_obj = drawing.get("rect")
        if rect_obj is None:
            continue
        try:
            rect = fitz.Rect(rect_obj)
        except Exception:
            continue
        if rect.is_empty or not rect.intersects(bbox):
            continue
        if _rect_horizontal_overlap(bbox, rect) < 0.08:
            continue
        if near_page_edge(rect) and near_crop_edge(rect) and _looks_like_running_page_header_rule(rect, page_rect):
            return True
    return False


def _looks_like_section_heading_text(text: str, rect: fitz.Rect, candidate_bbox: fitz.Rect) -> bool:
    cleaned = _clean_space(text)
    if not cleaned:
        return False
    if re.search(r"(?i)\b(?:figure|fig\.?|table)\s+[0-9]+", cleaned):
        return False
    words = re.findall(r"[A-Za-z][A-Za-z-]+", cleaned)
    if not 1 <= len(words) <= 12:
        return False
    if rect.width < min(120.0, candidate_bbox.width * 0.24):
        return False
    if cleaned.endswith((".", ",", ";", ":")) and not re.match(r"^\s*\d+(?:\.\d+)*\s+\S+", cleaned):
        return False
    if re.match(r"^\s*(?:[0-9]+(?:\.[0-9]+)*|[IVX]+)\s+[A-Z][A-Za-z-]+", cleaned):
        return True
    title_words = [word for word in words if word[:1].isupper()]
    return len(words) >= 2 and len(title_words) >= max(2, len(words) - 1)


def _looks_like_document_section_heading_text(text: str) -> bool:
    cleaned = _clean_space(text).strip(" .:-")
    if not cleaned:
        return False
    lowered = cleaned.lower()
    if lowered in {
        "abstract",
        "introduction",
        "related work",
        "background",
        "method",
        "methods",
        "experiments",
        "results",
        "discussion",
        "limitations",
        "conclusion",
        "conclusions",
        "references",
        "appendix",
    }:
        return True
    return re.match(r"^\s*(?:[0-9]+(?:\.[0-9]+)*|[IVX]+)\s+[A-Z][A-Za-z-]+", cleaned) is not None


def _pdf_crop_quality_flags(
    page: fitz.Page,
    bbox: fitz.Rect,
    caption: str,
    *,
    kind: str = "figure",
) -> list[str]:
    if bbox.is_empty:
        return []
    kind = "table" if str(kind or "").lower() == "table" else "figure"
    label = _caption_label_for_kind(caption, kind)
    caption_blocks = 0
    paragraph_blocks = 0
    foreign_caption_blocks = 0
    running_header_blocks = 0
    section_heading_blocks = 0
    for block in page.get_text("blocks"):
        rect = fitz.Rect(block[:4])
        if not rect.intersects(bbox):
            continue
        if _rect_horizontal_overlap(bbox, rect) < 0.35:
            continue
        text = _clean_space(str(block[4] or ""))
        if not text:
            continue
        if label and _caption_header_bonus_for_kind(text, caption, label, kind) > 0:
            caption_blocks += 1
            continue
        if _foreign_caption_label_in_text(text, label=label, kind=kind):
            foreign_caption_blocks += 1
            continue
        if _looks_like_running_page_header_text(text, rect, page.rect):
            running_header_blocks += 1
            continue
        if _looks_like_section_heading_text(text, rect, bbox):
            if kind == "table" or _looks_like_document_section_heading_text(text):
                section_heading_blocks += 1
            continue
        if _looks_like_external_figure_text(text, rect, bbox, label):
            paragraph_blocks += 1

    flags: list[str] = []
    if kind == "figure":
        for raw_flag in pdf_figure_crop_quality_flags(page, bbox, caption_rect=None):
            flag = str(raw_flag or "").strip()
            if flag and flag not in flags:
                flags.append(flag)
    if _crop_has_page_furniture_geometry(page, bbox):
        flags.append("running_header_leak")
    if caption_blocks:
        flags.append("caption_in_crop")
    if paragraph_blocks:
        flags.append("table_body_text_leak" if kind == "table" else "body_text_leak")
    if section_heading_blocks:
        flags.append("section_heading_leak")
    if running_header_blocks and "running_header_leak" not in flags:
        flags.append("running_header_leak")
    if foreign_caption_blocks:
        flags.append("neighbor_asset_leak")
    crop_area_frac = (bbox.width * bbox.height) / max(1.0, page.rect.width * page.rect.height)
    if paragraph_blocks >= 2 or (paragraph_blocks >= 1 and crop_area_frac >= 0.10):
        flags.append("page_like_table_crop" if kind == "table" else "page_like_figure_crop")
    return flags


def _table_crop_quality_flags(page: fitz.Page, bbox: fitz.Rect, caption: str) -> list[str]:
    flags: list[str] = []
    for source_flags in (
        _pdf_crop_quality_flags(page, bbox, caption, kind="table"),
        pdf_table_crop_quality_flags(page, bbox, caption_rect=None),
    ):
        for raw_flag in source_flags:
            flag = str(raw_flag or "").strip()
            if not flag:
                continue
            if flag not in flags:
                flags.append(flag)
            if flag == "body_text_leak" and "table_body_text_leak" not in flags:
                flags.append("table_body_text_leak")
    return flags


def _recompute_pdf_record_crop_quality_flags(
    pdf_doc: fitz.Document,
    layer_id: str,
    rec: dict[str, Any],
) -> list[str]:
    raw_bbox = rec.get("source_bbox_pdf_points")
    if not isinstance(raw_bbox, (list, tuple)) or len(raw_bbox) < 4:
        return list(rec.get("crop_quality_flags") or [])
    page_num = _safe_int(rec.get("source_page"), default=0)
    if not 1 <= page_num <= len(pdf_doc):
        return list(rec.get("crop_quality_flags") or [])
    try:
        bbox = fitz.Rect(raw_bbox[:4])
    except Exception:
        return list(rec.get("crop_quality_flags") or [])
    if bbox.is_empty:
        return list(rec.get("crop_quality_flags") or [])
    page = pdf_doc[page_num - 1]
    caption = str(rec.get("caption") or rec.get("title") or "")
    if str(rec.get("kind") or "").lower() == "table" or str(layer_id).startswith("ingest_table_"):
        return _table_crop_quality_flags(page, bbox, caption)
    return _pdf_crop_quality_flags(page, bbox, caption, kind="figure")


def _refresh_pdf_record_curation_from_current_rules(
    *,
    layer_id: str,
    rec: dict[str, Any],
    pdf_doc: fitz.Document | None = None,
) -> None:
    if pdf_doc is not None:
        rec["crop_quality_flags"] = _recompute_pdf_record_crop_quality_flags(pdf_doc, layer_id, rec)
    rec.update(_score_visual_candidate(layer_id, rec))


def _count_figure_caption_blocks(page: fitz.Page, candidate_bbox: fitz.Rect) -> int:
    count = 0
    for block in page.get_text("blocks"):
        rect = fitz.Rect(block[:4])
        if not rect.intersects(candidate_bbox):
            continue
        text = _clean_space(str(block[4] or ""))
        matches = re.findall(r"(?i)\b(?:figure|fig\.?)\s+[0-9]+\b", text)
        count += len(matches)
    return count


def _rect_horizontal_overlap(a: fitz.Rect, b: fitz.Rect) -> float:
    x0 = max(a.x0, b.x0)
    x1 = min(a.x1, b.x1)
    if x1 <= x0:
        return 0.0
    return (x1 - x0) / max(1.0, min(a.width, b.width))


def _caption_supports_horizontal_crop(
    original: fitz.Rect,
    caption_block: fitz.Rect,
    page_rect: fitz.Rect,
) -> bool:
    if caption_block.width >= page_rect.width * 0.55:
        return False
    if original.width <= caption_block.width * 1.25:
        return False
    if _rect_horizontal_overlap(original, caption_block) < 0.45:
        return False
    original_center = (original.x0 + original.x1) / 2.0
    caption_center = (caption_block.x0 + caption_block.x1) / 2.0
    if abs(original_center - caption_center) <= max(36.0, original.width * 0.12):
        return False
    return True


def _refined_bbox_is_useful(original: fitz.Rect, refined: fitz.Rect) -> bool:
    if refined.is_empty:
        return False
    if refined.width < 80 or refined.height < 60:
        return False
    original_area = max(1.0, original.width * original.height)
    refined_area = refined.width * refined.height
    if refined_area / original_area < 0.15:
        return False
    return (
        abs(refined.x0 - original.x0) > 1
        or abs(refined.y0 - original.y0) > 1
        or abs(refined.x1 - original.x1) > 1
        or abs(refined.y1 - original.y1) > 1
    )


def _register_candidates(
    *,
    candidates: list[PdfFigureCandidate],
    matches: dict[int, dict[str, Any]],
    ctx: ToolContext,
    pdf_path: Path,
) -> list[str]:
    """Apply caption matches + fake-figure filter. Register survivors
    in `ctx.state["rendered_layers"]` with the downstream contract
    schema, then rename their PNGs to `img_{layer_id}.png` so the
    layers directory stays tidy."""
    _raise_if_ingest_cancelled(ctx, "ingest.figure_register.start")
    registered: list[str] = []
    accepted_candidates = 0
    ignored_sub_panels = 0
    next_idx = 1
    source_pdf_sha256 = sha256_file(pdf_path)
    _raise_if_ingest_cancelled(ctx, "ingest.figure_register.after_source_hash")
    vector_doc: fitz.Document | None = None
    if any(c.bbox_pt is not None for c in candidates):
        try:
            vector_doc = fitz.open(pdf_path)
        except Exception as e:  # noqa: BLE001
            log("ingest.pdf.vector_crop.open_fail", file=pdf_path.name, error=str(e)[:200])

    try:
        for cand_idx, cand in enumerate(candidates):
            _raise_if_ingest_cancelled(ctx, "ingest.figure_register.before_candidate")
            match = matches.get(cand_idx, {})
            is_real = bool(match.get("is_real_figure", True))
            if not is_real:
                log("ingest.pdf.reject_fake",
                    page=cand.page, strategy=cand.strategy,
                    reason=match.get("reason", ""), path=cand.path.name)
                # Clean up the rejected PNG so we don't bloat the run dir.
                try:
                    _raise_if_ingest_cancelled(ctx, "ingest.figure_register.before_reject_unlink")
                    cand.path.unlink()
                    _raise_if_ingest_cancelled(ctx, "ingest.figure_register.after_reject_unlink")
                except OSError:
                    pass
                continue

            if vector_doc is not None and cand.strategy == "vector":
                refined = _refine_vector_candidate_bbox(cand, match, vector_doc, ctx)
                if refined is None:
                    try:
                        _raise_if_ingest_cancelled(ctx, "ingest.figure_register.before_refine_unlink")
                        cand.path.unlink()
                        _raise_if_ingest_cancelled(ctx, "ingest.figure_register.after_refine_unlink")
                    except OSError:
                        pass
                    continue
                cand = refined
            crop_quality_flags: list[str] = []
            if vector_doc is not None and cand.bbox_pt is not None and 1 <= cand.page <= len(vector_doc):
                crop_quality_flags = _pdf_crop_quality_flags(
                    vector_doc[cand.page - 1],
                    fitz.Rect(cand.bbox_pt),
                    str(match.get("caption_text") or ""),
                )
            placement_quality_flags = crop_quality_flags if cand.strategy in {"raster", "embedded"} else []
            if placement_quality_flags:
                crop_quality_flags = []

            layer_id = f"ingest_fig_{next_idx:02d}"
            next_idx += 1
            final_path = ctx.layers_dir / f"img_{layer_id}.png"
            try:
                _raise_if_ingest_cancelled(ctx, "ingest.figure_register.before_rename")
                if cand.path.resolve() != final_path.resolve():
                    shutil.move(str(cand.path), str(final_path))
                _raise_if_ingest_cancelled(ctx, "ingest.figure_register.after_rename")
            except OSError as e:
                _raise_if_ingest_cancelled(ctx, "ingest.figure_register.before_rename_failure_log")
                log("ingest.pdf.rename_fail",
                    layer_id=layer_id, error=str(e))
                _raise_if_ingest_cancelled(ctx, "ingest.figure_register.after_rename_failure_log")
                final_path = cand.path

            _raise_if_ingest_cancelled(ctx, "ingest.figure_register.before_state_merge")
            ctx.state["rendered_layers"][layer_id] = {
                "layer_id": layer_id,
                "name": f"figure_{next_idx - 1}",
                "kind": "image",
                "z_index": 5,
                "bbox": None,
                "src_path": str(final_path),
                "aspect_ratio": _aspect_from_dims(cand.width_px, cand.height_px),
                "image_size": f"{cand.width_px}x{cand.height_px}",
                "sha256": sha256_file(final_path),
                "source": "ingested_pdf",
                "source_file": str(pdf_path),
                "source_pdf": pdf_path.name,
                "source_pdf_sha256": source_pdf_sha256,
                "source_page": cand.page,
                "source_bbox_pdf_points": list(cand.bbox_pt) if cand.bbox_pt else None,
                "source_image_xref": cand.xref,
                "caption": match.get("caption_text", "") or getattr(cand, "source_group_caption", ""),
                # v2.3 — VLM's ≤15-char label for tight bboxes; planner picks
                # between `caption` (full) and `caption_short` based on slot size.
                "caption_short": match.get("short_caption", ""),
                "extract_strategy": cand.strategy,   # for debugging
                "protected_anchor": bool(getattr(cand, "protected_anchor", False)),
                "anchor_kind": getattr(cand, "anchor_kind", "") or "",
                "anchor_label": getattr(cand, "anchor_label", "") or "",
                "anchor_reason": getattr(cand, "anchor_reason", "") or "",
                **_candidate_source_group_payload(
                    cand,
                    kind="figure",
                    caption=str(match.get("caption_text") or ""),
                    source=str(match.get("caption_association_method") or ""),
                ),
                "caption_confidence": match.get("confidence", 0.0),
                "caption_association_method": match.get("caption_association_method", "unmatched"),
                "crop_quality_flags": crop_quality_flags,
                "placement_quality_flags": placement_quality_flags,
            }
            _raise_if_ingest_cancelled(ctx, "ingest.figure_register.after_state_merge")
            registered.append(layer_id)
            accepted_candidates += 1

            sub_panels = match.get("sub_panels") or []
            if isinstance(sub_panels, list) and sub_panels:
                ignored_sub_panels += len(sub_panels)
                _raise_if_ingest_cancelled(ctx, "ingest.figure_register.before_subpanel_log")
                log(
                    "ingest.pdf.sub_panel.ignored",
                    parent=layer_id,
                    requested=len(sub_panels),
                    reason="partial_crops_disabled",
                )
                _raise_if_ingest_cancelled(ctx, "ingest.figure_register.after_subpanel_log")
    finally:
        if vector_doc is not None:
            vector_doc.close()

    _raise_if_ingest_cancelled(ctx, "ingest.figure_register.before_done_log")
    log("ingest.pdf.register",
        input_candidates=len(candidates),
        kept=len(registered),
        kept_originals=accepted_candidates,
        sub_panels_added=0,
        ignored_sub_panels=ignored_sub_panels,
        dropped=max(0, len(candidates) - accepted_candidates))
    _raise_if_ingest_cancelled(ctx, "ingest.figure_register.after_done_log")
    return registered


# ───────────────────────── Table parsing ───────────────────────────────

def _parse_tables_parallel(
    candidates: list[PdfTableCandidate],
    manifest: dict[str, Any],
    ctx: ToolContext,
) -> dict[int, dict[str, Any]]:
    """Run VLM parse per table candidate in a thread pool.

    Returns a dict: candidate idx → {is_table, headers, rows, title,
    matched_idx, caption_text, reason}. Callers use it to decide
    whether to register the candidate as a `kind="table"` layer.
    """
    _raise_if_ingest_cancelled(ctx, "ingest.table_parse.start")
    results: dict[int, dict[str, Any]] = {}
    parallelism = _ingest_vlm_parallelism(ctx, _CAPTION_MATCH_PARALLELISM)
    timeout_s = _table_parse_timeout_s(ctx)
    max_retries = _table_parse_max_retries(ctx)
    _raise_if_ingest_cancelled(ctx, "ingest.table_parse.before_start_log")
    log("ingest.pdf.table_parse.start",
        n_candidates=len(candidates), parallelism=parallelism,
        model=ctx.settings.ingest_model,
        timeout_s=timeout_s,
        max_retries=max_retries)
    _raise_if_ingest_cancelled(ctx, "ingest.table_parse.after_start_log")

    def parse_one(i: int, candidate: PdfTableCandidate) -> dict[str, Any]:
        _raise_if_ingest_cancelled(ctx, "ingest.table_parse.worker.start")
        return _parse_one_table(i, candidate, manifest, ctx)

    outcomes = _run_bounded_ingest_pool(
        items=list(enumerate(candidates)),
        worker=parse_one,
        parallelism=parallelism,
        ctx=ctx,
        phase="ingest.table_parse.pool",
    )
    for i, result, error in outcomes:
        _raise_if_ingest_cancelled(ctx, "ingest.table_parse.before_result_merge")
        if error is not None:
            _raise_if_ingest_cancelled(ctx, "ingest.table_parse.before_failure_log")
            log("ingest.pdf.table_parse_fail", cand_idx=i, error=str(error))
            _raise_if_ingest_cancelled(ctx, "ingest.table_parse.after_failure_log")
            result = {
                "is_table": False, "headers": [], "rows": [],
                "title": "", "matched_idx": None,
                "caption_text": "", "reason": f"parse failed: {error}",
            }
        results[i] = result
        _raise_if_ingest_cancelled(ctx, "ingest.table_parse.after_result_merge")
    return results


def _parse_one_table(
    cand_idx: int,
    candidate: PdfTableCandidate,
    manifest: dict[str, Any],
    ctx: ToolContext,
) -> dict[str, Any]:
    _raise_if_ingest_cancelled(ctx, "ingest.table_parse.worker.before_prepare")
    # Filter caption candidates to same page ± 1 (tables sometimes
    # caption on the following page for bottom-of-page tables). Include
    # entries with missing/unknown `page` so appendix tables the VLM
    # failed to page-number aren't silently excluded from the pool.
    all_tables = list(manifest.get("tables", []))
    near: list[tuple[int, dict[str, Any]]] = []
    for i, t in enumerate(all_tables):
        raw_page = t.get("page")
        if raw_page is None:
            near.append((i, t))
            continue
        try:
            pnum = int(raw_page)
        except (TypeError, ValueError):
            near.append((i, t))
            continue
        if abs(pnum - candidate.page) <= 1:
            near.append((i, t))
    pool = near if near else list(enumerate(all_tables))
    local_to_global = {local_i: global_i for local_i, (global_i, _) in enumerate(pool)}

    cap_lines = []
    for local_i, (_, tbl) in enumerate(pool):
        cap = (tbl.get("caption") or "").replace("\n", " ")[:240]
        cap_lines.append(f"  [{local_i}] (p.{tbl.get('page', '?')}) {cap}")
    cap_block = "\n".join(cap_lines) or "  (no caption candidates on this page)"

    # Show the VLM pymupdf's best-effort cells, truncated. Don't send
    # the full raw cells when they're huge — we trust the image more.
    raw_preview: list[str] = []
    for row in candidate.raw_cells[:12]:
        raw_preview.append(
            " | ".join(str(c)[:80].replace("\n", " ⏎ ") for c in row)
        )
    raw_block = "\n".join(raw_preview) or "  (pymupdf extracted no cells)"

    user_text = (
        f"Table candidate from page {candidate.page} "
        f"(pymupdf saw {candidate.nrows}×{candidate.ncols} cells).\n\n"
        f"pymupdf raw-cell preview (may be wrong — trust the image):\n"
        f"{raw_block}\n\n"
        f"Caption candidates near this page:\n{cap_block}\n\n"
        "Return the JSON described in the system prompt."
    )

    _raise_if_ingest_cancelled(ctx, "ingest.table_parse.worker.before_image_read")
    image = VlmImage.from_path(candidate.image_path)
    _raise_if_ingest_cancelled(ctx, "ingest.table_parse.worker.after_image_read")
    _raise_if_ingest_cancelled(ctx, "ingest.table_parse.worker.before_vlm")
    result = vlm_call_json(
        settings=ctx.settings,
        model=ctx.settings.ingest_model,
        system=_TABLE_PARSE_PROMPT,
        user_text=user_text,
        images=[image],
        max_tokens=4096,
        timeout_s=_table_parse_timeout_s(ctx),
        max_retries=_table_parse_max_retries(ctx),
        cancellation_token=_ingest_cancellation_token(ctx),
    )
    _raise_if_ingest_cancelled(ctx, "ingest.table_parse.worker.after_vlm")

    local_idx = result.get("matched_idx")
    global_idx = None
    caption_text = str(result.get("title", "")).strip()
    if isinstance(local_idx, int) and local_idx in local_to_global:
        global_idx = local_to_global[local_idx]
        caption_text = str(all_tables[global_idx].get("caption", "")) or caption_text

    rows = result.get("rows") or []
    headers = result.get("headers") or []
    col_rule = result.get("col_highlight_rule") or []
    # Light sanitation: coerce to str, drop empty trailing rows.
    rows = [[str(c) if c is not None else "" for c in row] for row in rows]
    headers = [str(c) if c is not None else "" for c in headers]
    col_rule = [str(c) if c is not None else "" for c in col_rule]
    # Normalize rule length to match header count (pad with "" / truncate).
    if headers:
        n = len(headers)
        col_rule = col_rule[:n] + [""] * max(0, n - len(col_rule))
    while rows and not any(c.strip() for c in rows[-1]):
        rows.pop()

    _raise_if_ingest_cancelled(ctx, "ingest.table_parse.worker.before_return")
    return {
        "is_table": bool(result.get("is_table", False)),
        "headers": headers,
        "rows": rows,
        "col_highlight_rule": col_rule,
        "title": caption_text,
        "matched_idx": global_idx,
        "caption_text": caption_text,
        # v2.3 — ≤15-char label for deck/poster slots where the full
        # table caption won't fit. Clipped defensively.
        "short_title": str(result.get("short_title", "") or "")[:40],
        "reason": str(result.get("reason", ""))[:200],
    }


_TABLE_PARSE_TRANSPORT_ERROR_TERMS = (
    "unable to download image",
    "image transport",
    "transport failed",
    "curl transport failed",
    "vlm curl transport",
    "openai-compatible curl transport",
    "downstream_http",
    "upstream_vendor",
    "upstreamstatus",
    "status=0",
    "status=400",
    "status=500",
    "error code: 400",
    "error code: 500",
    "timed out",
    "timeout",
)


def _is_table_parse_transport_error(reason: Any) -> bool:
    text = str(reason or "").strip().lower()
    if not text:
        return False
    return any(term in text for term in _TABLE_PARSE_TRANSPORT_ERROR_TERMS)


def _captioned_table_parse_fallback_allowed(
    cand: PdfTableCandidate,
    info: Mapping[str, Any],
    manifest: dict[str, Any] | None = None,
) -> bool:
    meta = _source_group_metadata_for_candidate(
        kind="table",
        candidate=cand,
        manifest=manifest or {},
        caption=str(info.get("caption_text") or info.get("title") or ""),
        source="table_parse",
    )
    if not _candidate_has_captioned_source_group(cand, meta):
        return False
    if str(getattr(cand, "source_group_kind", "") or getattr(cand, "anchor_kind", "") or "table").lower() != "table":
        return False
    if not _table_candidate_crop_plausible(cand):
        return False
    reason = str(info.get("reason") or "").lower()
    if info and not reason and not info.get("is_table"):
        return False
    if _table_parse_semantic_rejection(reason):
        return False
    return True


def _candidate_has_captioned_source_group(candidate: Any, metadata: Mapping[str, Any] | None = None) -> bool:
    metadata = metadata or {}
    if bool(getattr(candidate, "captioned_source_group", False)):
        return True
    if not _candidate_has_source_group(candidate) and not str(metadata.get("source_group_id") or "").strip():
        return False
    return bool(
        str(metadata.get("source_group_label") or getattr(candidate, "source_group_label", "") or "").strip()
        or str(metadata.get("source_group_caption") or getattr(candidate, "source_group_caption", "") or "").strip()
        or str(getattr(candidate, "anchor_label", "") or "").strip()
    )


def _table_candidate_crop_plausible(cand: PdfTableCandidate) -> bool:
    return cand.width_px >= 120 and cand.height_px >= 60 and (cand.width_px * cand.height_px) >= 18_000


def _table_parse_semantic_rejection(reason: str) -> bool:
    text = str(reason or "").lower()
    if not text:
        return False
    failure_terms = (
        "parse failed",
        "unable to parse",
        "could not parse",
        "invalid json",
        "malformed",
        "schema",
        "timeout",
        "timed out",
        "transport",
        "download",
        "empty rows",
        "empty headers",
    )
    if _is_table_parse_transport_error(text) or any(term in text for term in failure_terms):
        return False
    reject_terms = (
        "not a table",
        "not a data table",
        "diagram",
        "flowchart",
        "equation",
        "formula",
        "figure",
        "layout artifact",
    )
    return any(term in text for term in reject_terms)


def _candidate_source_group_payload(
    cand: Any,
    *,
    kind: str,
    manifest: dict[str, Any] | None = None,
    caption: str = "",
    source: str = "",
) -> dict[str, Any]:
    metadata = _source_group_metadata_for_candidate(
        kind=kind,
        candidate=cand,
        manifest=manifest or {},
        caption=caption,
        source=source,
    )
    payload = {
        "captioned_source_group": _candidate_has_captioned_source_group(cand, metadata),
        "source_group_id": metadata.get("source_group_id") or "",
        "source_group_kind": metadata.get("source_group_kind") or "",
        "source_group_label": metadata.get("source_group_label") or "",
        "source_group_caption": metadata.get("source_group_caption") or "",
        "source_group_source": metadata.get("source_group_source") or "",
    }
    return {k: v for k, v in payload.items() if v not in ("", False)}


def _unparsed_source_table_info(
    cand: PdfTableCandidate,
    info: Mapping[str, Any],
    manifest: dict[str, Any],
) -> dict[str, Any]:
    source_group = _source_group_metadata_for_candidate(
        kind="table",
        candidate=cand,
        manifest=manifest,
        caption=str(info.get("caption_text") or info.get("title") or ""),
        source="table_parse_fallback",
    )
    label = _normalize_source_group_display_label(
        str(source_group.get("source_group_label") or "").strip(),
        kind="table",
    )
    caption = (
        str(info.get("caption_text") or info.get("title") or "").strip()
        or str(source_group.get("source_group_caption") or "").strip()
        or (f"Table {label}" if label else "Captioned source table")
    )
    short_title = (
        str(info.get("short_title") or "").strip()
        or (f"Table {label}" if label else "Source table")
    )[:40]
    reason = str(info.get("reason") or "table parse failed")[:500]
    return {
        **dict(info),
        "is_table": True,
        "headers": [],
        "rows": [],
        "col_highlight_rule": [],
        "title": caption,
        "caption_text": caption,
        "short_title": short_title,
        "table_parse_status": "unparsed_source_crop",
        "table_parse_error": reason,
    }


def _normalize_source_group_display_label(label: str, *, kind: str) -> str:
    text = str(label or "").strip()
    if not text:
        return ""
    prefix = "Table" if str(kind or "").lower() == "table" else "Figure"
    match = re.match(rf"(?i)^\s*{prefix}\s+(.+?)\s*$", text)
    return match.group(1).strip() if match else text


def _register_tables(
    *,
    candidates: list[PdfTableCandidate],
    parsed: dict[int, dict[str, Any]],
    ctx: ToolContext,
    pdf_path: Path,
    manifest: dict[str, Any],
) -> list[str]:
    """Register VLM-validated tables as `kind="table"` layers in
    `ctx.state["rendered_layers"]`. Rejects go to `reject_fake` logs.

    The PSD/SVG fallback PNG (a cleanly re-drawn table image) is
    generated lazily by the renderer via `util.table_png`, so we only
    need to persist the structured rows/headers here.
    """
    _raise_if_ingest_cancelled(ctx, "ingest.table_register.start")
    registered: list[str] = []
    next_idx = 1
    unparsed_source_tables = 0
    source_pdf_sha256 = sha256_file(pdf_path)
    _raise_if_ingest_cancelled(ctx, "ingest.table_register.after_source_hash")
    table_quality_doc = None
    try:
        table_quality_doc = fitz.open(pdf_path)
    except Exception as e:  # noqa: BLE001 - table quality flags are advisory
        log("ingest.pdf.table_crop_quality.open_fail", file=pdf_path.name, error=str(e)[:200])

    try:
        for cand_idx, cand in enumerate(candidates):
            _raise_if_ingest_cancelled(ctx, "ingest.table_register.before_candidate")
            info = parsed.get(cand_idx)
            if info is None:
                log("ingest.pdf.reject_table",
                    page=cand.page, reason="table_parse_missing")
                try:
                    _raise_if_ingest_cancelled(ctx, "ingest.table_register.before_missing_unlink")
                    cand.image_path.unlink()
                    _raise_if_ingest_cancelled(ctx, "ingest.table_register.after_missing_unlink")
                except OSError:
                    pass
                continue
            if not info.get("is_table", False):
                if _captioned_table_parse_fallback_allowed(cand, info, manifest):
                    info = _unparsed_source_table_info(cand, info, manifest)
                    unparsed_source_tables += 1
                else:
                    log("ingest.pdf.reject_table",
                        page=cand.page, reason=info.get("reason", ""))
                    try:
                        _raise_if_ingest_cancelled(ctx, "ingest.table_register.before_reject_unlink")
                        cand.image_path.unlink()
                        _raise_if_ingest_cancelled(ctx, "ingest.table_register.after_reject_unlink")
                    except OSError:
                        pass
                    continue

            rows = info.get("rows") or []
            headers = info.get("headers") or []
            col_rule = info.get("col_highlight_rule") or []
            if not rows and not headers and info.get("table_parse_status") != "unparsed_source_crop":
                empty_info = {**dict(info), "reason": str(info.get("reason") or "empty rows+headers")}
                if _captioned_table_parse_fallback_allowed(cand, empty_info, manifest):
                    info = _unparsed_source_table_info(cand, empty_info, manifest)
                    rows = []
                    headers = []
                    col_rule = []
                    unparsed_source_tables += 1
                else:
                    log("ingest.pdf.reject_table",
                        page=cand.page, reason="empty rows+headers")
                    try:
                        _raise_if_ingest_cancelled(ctx, "ingest.table_register.before_empty_unlink")
                        cand.image_path.unlink()
                        _raise_if_ingest_cancelled(ctx, "ingest.table_register.after_empty_unlink")
                    except OSError:
                        pass
                    continue

            layer_id = f"ingest_table_{next_idx:02d}"
            next_idx += 1
            table_caption = str(info.get("caption_text") or info.get("title") or "")
            crop_quality_flags: list[str] = []
            if (
                table_quality_doc is not None
                and cand.bbox_pt is not None
                and 1 <= cand.page <= len(table_quality_doc)
            ):
                crop_quality_flags = _table_crop_quality_flags(
                    table_quality_doc[cand.page - 1],
                    fitz.Rect(cand.bbox_pt),
                    table_caption,
                )

            # Keep the PDF author's table appearance as the source evidence.
            # The stable img_ingest_table_NN name remains for downstream
            # compatibility, but the pixels are the PyMuPDF bbox crop, not a
            # table reconstructed from parsed rows.
            final_png = ctx.layers_dir / f"img_{layer_id}.png"
            try:
                _raise_if_ingest_cancelled(ctx, "ingest.table_register.before_rename")
                if cand.image_path.resolve() != final_png.resolve():
                    shutil.move(str(cand.image_path), str(final_png))
                _raise_if_ingest_cancelled(ctx, "ingest.table_register.after_rename")
            except Exception as e:
                _raise_if_ingest_cancelled(ctx, "ingest.table_register.before_rename_failure_log")
                log("ingest.pdf.table_crop_preserve_fail",
                    layer_id=layer_id, error=str(e))
                _raise_if_ingest_cancelled(ctx, "ingest.table_register.after_rename_failure_log")
                final_png = cand.image_path

            _raise_if_ingest_cancelled(ctx, "ingest.table_register.before_state_merge")
            ctx.state["rendered_layers"][layer_id] = {
                "layer_id": layer_id,
                "name": f"table_{next_idx - 1}",
                "kind": "table",
                "z_index": 5,
                "bbox": None,
                "src_path": str(final_png),        # original PDF table crop
                "table_visual_source": "original_pdf_crop",
                "aspect_ratio": _aspect_from_dims(cand.width_px, cand.height_px),
                "image_size": f"{cand.width_px}x{cand.height_px}",
                "sha256": sha256_file(final_png),
                "source": "ingested_pdf",
                "source_file": str(pdf_path),
                "source_pdf": pdf_path.name,
                "source_pdf_sha256": source_pdf_sha256,
                "source_page": cand.page,
                "source_bbox_pdf_points": list(cand.bbox_pt),
                "caption": table_caption,
                # v2.3 — ≤15-char label for tight table slots (tables share the
                # "caption_short" key with figures so renderers can read one field).
                "caption_short": info.get("short_title", ""),
                # structured data:
                "rows": rows,
                "headers": headers,
                "col_highlight_rule": col_rule,
                "title": info.get("title", ""),
                "extract_strategy": "table",
                "table_parse_status": info.get("table_parse_status") or "parsed",
                "table_parse_error": info.get("table_parse_error") or "",
                "protected_anchor": bool(getattr(cand, "protected_anchor", False)),
                "anchor_kind": getattr(cand, "anchor_kind", "") or "",
                "anchor_label": getattr(cand, "anchor_label", "") or "",
                "anchor_reason": getattr(cand, "anchor_reason", "") or "",
                **_candidate_source_group_payload(
                    cand,
                    kind="table",
                    manifest=manifest,
                    caption=table_caption,
                    source="table_parse",
                ),
                "crop_quality_flags": crop_quality_flags,
            }
            _raise_if_ingest_cancelled(ctx, "ingest.table_register.after_state_merge")
            registered.append(layer_id)
    finally:
        if table_quality_doc is not None:
            table_quality_doc.close()

    _raise_if_ingest_cancelled(ctx, "ingest.table_register.before_done_log")
    log("ingest.pdf.register_tables",
        kept=len(registered),
        unparsed_source_tables=unparsed_source_tables,
        dropped=len(candidates) - len(registered))
    _raise_if_ingest_cancelled(ctx, "ingest.table_register.after_done_log")
    return registered


# ───────────────────────── Markdown branch ─────────────────────────────

def _ingest_markdown(fp: Path, ctx: ToolContext) -> dict[str, Any]:
    text = fp.read_text(encoding="utf-8")
    registered: list[str] = []
    skipped: list[str] = []

    for m in _MD_IMG_RE.finditer(text):
        alt_text = m.group(1)
        ref = m.group(2).strip()
        if ref.startswith(("http://", "https://", "data:")):
            skipped.append(ref[:80])
            continue
        src = Path(ref)
        if not src.is_absolute():
            src = (fp.parent / src).resolve()
        if not src.exists() or not src.is_file():
            skipped.append(str(src))
            continue
        try:
            layer_id = _register_image_file(src, ctx, name_hint=alt_text or src.stem)
        except RuntimeError:
            skipped.append(str(src))
            continue
        registered.append(layer_id)

    log("ingest.md.done", file=fp.name, chars=len(text),
        registered=len(registered), skipped=len(skipped))

    return {
        "file": str(fp), "type": "markdown",
        "raw_text": text,
        "n_chars": len(text),
        "registered_layer_ids": registered,
        "skipped_images": skipped,
        "summary": f"{fp.name} — {len(text)} chars, "
                   f"{len(registered)} image(s)"
                   + (f", {len(skipped)} skipped" if skipped else ""),
    }


# ─────────────────────────── .docx branch ─────────────────────────────

def _ingest_docx(fp: Path, ctx: ToolContext) -> dict[str, Any]:
    """Read a Word document into the same manifest shape as PDF.

    Docx has real structural metadata (heading styles, inline images,
    captions), so unlike PDF this branch does NOT need a VLM call —
    we read the docx tree directly, which is faster, free, and more
    faithful.
    """
    from docx import Document

    doc = Document(str(fp))

    body_paras: list[tuple[str, str]] = []  # (style_name, text)
    for para in doc.paragraphs:
        text = (para.text or "").strip()
        if not text:
            continue
        style = (para.style.name if para.style is not None else "") or ""
        body_paras.append((style, text))

    title = _docx_pick_title(body_paras)
    sections = _docx_build_sections(body_paras, title)

    registered_figure_ids: list[str] = []
    for rel_id, rel in doc.part.rels.items():
        if rel.reltype and "image" in rel.reltype:
            try:
                blob = rel.target_part.blob
                mime = getattr(rel.target_part, "content_type", "") or ""
                ext = _ext_for_image_mime(mime)
                layer_id = _register_image_blob(
                    blob, ext, ctx,
                    name_hint=f"{fp.stem}_{rel_id}",
                    source_file=fp,
                    source_ref=f"rel={rel_id}",
                )
                registered_figure_ids.append(layer_id)
            except (RuntimeError, OSError) as e:
                log("ingest.docx.image_skip", rel=rel_id, error=str(e))

    manifest = {
        "title": title or fp.stem,
        "authors": [],
        "venue": None,
        "abstract": sections[0]["summary"] if sections else "",
        "sections": sections,
        "figures": [],
        "tables": [],
        "key_quotes": [],
    }
    _normalize_manifest_lists(manifest)

    log("ingest.docx.done", file=fp.name,
        sections=len(sections), figures=len(registered_figure_ids))

    # v2.7 — verbatim body text for provenance validator (paired with the
    # PDF branch's `raw_text`). docx body is already in memory as
    # body_paras tuples; flatten them.
    docx_raw_text = "\n".join(t for _, t in body_paras)
    return {
        "file": str(fp), "type": "docx", "manifest": manifest,
        "registered_layer_ids": registered_figure_ids,
        "registered_figure_ids": registered_figure_ids,
        "registered_table_ids": [],
        "raw_text": docx_raw_text,
        "summary": f"{manifest['title']} — "
                   f"{len(registered_figure_ids)} figure(s), "
                   f"{len(sections)} section(s)",
    }


_DOCX_HEADING_PREFIXES = ("Heading", "Title")


def _docx_pick_title(paras: list[tuple[str, str]]) -> str:
    for style, text in paras:
        if style.startswith("Title"):
            return text
    for style, text in paras:
        if style.startswith("Heading 1") or style == "Heading 1":
            return text
    for _style, text in paras:
        return text
    return ""


def _docx_build_sections(
    paras: list[tuple[str, str]], title: str,
) -> list[dict[str, Any]]:
    """Group paragraphs under their nearest preceding heading. Non-
    heading paras become the section body; the first sentence (or 2-3
    sentences up to ~400 chars) becomes `summary`, and short bullet-like
    lines become `key_points`."""
    sections: list[dict[str, Any]] = []
    current_heading: str | None = None
    current_body: list[str] = []

    def flush() -> None:
        nonlocal current_heading, current_body
        if current_heading is None and not current_body:
            return
        heading = current_heading or "Body"
        body_text = "\n".join(current_body).strip()
        sections.append({
            "idx": len(sections) + 1,
            "heading": heading,
            "summary": _first_sentences(body_text, max_chars=400),
            "key_points": _pick_key_points(current_body),
        })
        current_heading = None
        current_body = []

    for style, text in paras:
        is_heading = any(style.startswith(p) for p in _DOCX_HEADING_PREFIXES)
        if is_heading and text == title:
            continue
        if is_heading:
            flush()
            current_heading = text
        else:
            current_body.append(text)
    flush()

    return sections


def _first_sentences(text: str, max_chars: int = 400) -> str:
    if not text:
        return ""
    pieces: list[str] = re.split(r"(?<=[.!?。！？])\s+", text)
    acc = ""
    for p in pieces:
        if not p:
            continue
        if len(acc) + len(p) + 1 > max_chars:
            break
        acc = (acc + " " + p).strip() if acc else p
    return acc or text[:max_chars]


def _pick_key_points(body: list[str]) -> list[str]:
    """Pull bullet-like lines (short, starts with punctuation/enumeration
    marker, OR simply short standalone paras ≤ 120 chars). Max 5."""
    out: list[str] = []
    for line in body:
        s = line.strip()
        if not s or len(s) > 160:
            continue
        if re.match(r"^[-\*•]|^\d+[\.\)]\s", s) or len(s) <= 120:
            out.append(s.lstrip("-*•").strip())
        if len(out) >= 5:
            break
    return out


def _ext_for_image_mime(mime: str) -> str:
    mime = (mime or "").lower()
    if "png" in mime:
        return ".png"
    if "jpeg" in mime or "jpg" in mime:
        return ".jpg"
    if "gif" in mime:
        return ".gif"
    if "webp" in mime:
        return ".webp"
    if "bmp" in mime:
        return ".bmp"
    if "tif" in mime:
        return ".tif"
    return ".png"


# ─────────────────────────── .pptx branch ─────────────────────────────

def _ingest_pptx(fp: Path, ctx: ToolContext) -> dict[str, Any]:
    """Read a PowerPoint file into the same manifest shape. Each slide
    becomes one section; title placeholder → section heading, body
    placeholders → section body; picture shapes → ingest_fig_NN layers.
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(str(fp))

    sections: list[dict[str, Any]] = []
    registered_figure_ids: list[str] = []
    deck_title: str | None = None

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_title: str | None = None
        body_parts: list[str] = []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    blob = shape.image.blob
                    mime = shape.image.content_type or ""
                    ext = _ext_for_image_mime(mime)
                    layer_id = _register_image_blob(
                        blob, ext, ctx,
                        name_hint=f"{fp.stem}_slide{slide_idx:02d}",
                        source_file=fp,
                        source_ref=f"slide={slide_idx}",
                    )
                    registered_figure_ids.append(layer_id)
                except (RuntimeError, OSError) as e:
                    log("ingest.pptx.image_skip",
                        slide=slide_idx, error=str(e))
                continue
            if not shape.has_text_frame:
                continue
            text = (shape.text_frame.text or "").strip()
            if not text:
                continue
            is_title = (
                hasattr(shape, "placeholder_format")
                and shape.placeholder_format is not None
                and shape.placeholder_format.idx == 0
            )
            if is_title and slide_title is None:
                slide_title = text
            else:
                body_parts.append(text)

        if slide_idx == 1 and slide_title:
            deck_title = slide_title

        heading = slide_title or f"Slide {slide_idx}"
        body_text = "\n".join(body_parts).strip()
        sections.append({
            "idx": slide_idx,
            "heading": heading,
            "summary": _first_sentences(body_text, max_chars=400),
            "key_points": _pick_key_points(body_parts),
        })

    manifest = {
        "title": deck_title or fp.stem,
        "authors": [],
        "venue": None,
        "abstract": sections[0]["summary"] if sections else "",
        "sections": sections,
        "figures": [],
        "tables": [],
        "key_quotes": [],
    }
    _normalize_manifest_lists(manifest)

    log("ingest.pptx.done", file=fp.name,
        slides=len(sections), figures=len(registered_figure_ids))

    return {
        "file": str(fp), "type": "pptx", "manifest": manifest,
        "registered_layer_ids": registered_figure_ids,
        "registered_figure_ids": registered_figure_ids,
        "registered_table_ids": [],
        "summary": f"{manifest['title']} — "
                   f"{len(sections)} slide(s), "
                   f"{len(registered_figure_ids)} figure(s)",
    }


def _register_image_blob(
    blob: bytes, ext: str, ctx: ToolContext, *,
    name_hint: str, source_file: Path, source_ref: str,
) -> str:
    """Register an in-memory image blob as an `ingest_fig_NN` layer.

    Unlike `_register_image_file` which expects an on-disk source and
    uses the sha-based layer_id shape (`ingest_img_<sha8>`), this helper
    allocates a sequential `ingest_fig_NN` id so .docx / .pptx images
    show up to the planner the same way PDF figures do — and the
    figure-cross-reference detector in composite picks them up too.
    """
    import hashlib
    from PIL import Image as _Image
    import io

    # Sequential id — peek at rendered_layers to find next free index.
    existing = [
        k for k in ctx.state["rendered_layers"]
        if k.startswith("ingest_fig_")
    ]
    next_idx = len(existing) + 1
    layer_id = f"ingest_fig_{next_idx:02d}"

    dest = ctx.layers_dir / f"img_{layer_id}{ext}"
    dest.write_bytes(blob)

    try:
        with _Image.open(io.BytesIO(blob)) as im:
            w, h = im.size
    except Exception as e:
        try:
            dest.unlink()
        except OSError:
            pass
        raise RuntimeError(f"blob not readable ({e})")

    source_page = None
    if source_ref.startswith("slide="):
        source_page = _safe_int(source_ref.split("=", 1)[1], default=0) or None

    sha = hashlib.sha256(blob).hexdigest()
    ctx.state["rendered_layers"][layer_id] = {
        "layer_id": layer_id,
        "name": _sanitize_name(name_hint) or layer_id,
        "kind": "image",
        "z_index": 5,
        "bbox": None,
        "src_path": str(dest),
        "aspect_ratio": _aspect_from_dims(w, h),
        "image_size": f"{w}x{h}",
        "sha256": sha,
        "source": "ingested_" + ("pptx" if source_ref.startswith("slide=") else "docx"),
        "source_file": str(source_file),
        "source_ref": source_ref,
        "source_page": source_page,
        "caption": "",
        "caption_short": _sanitize_name(name_hint) or layer_id,
        "extract_strategy": "embedded",
        "caption_confidence": 0.0,
    }
    return layer_id


# ─────────────────────────── Image branch ──────────────────────────────

def _ingest_image(fp: Path, ctx: ToolContext) -> dict[str, Any]:
    layer_id = _register_image_file(fp, ctx, name_hint=fp.stem)
    rec = ctx.state["rendered_layers"][layer_id]
    w_s, h_s = rec["image_size"].split("x")
    log("ingest.image.done", file=fp.name, layer_id=layer_id,
        w=int(w_s), h=int(h_s))
    return {
        "file": str(fp), "type": "image",
        "registered_layer_ids": [layer_id],
        "width": int(w_s), "height": int(h_s),
        "summary": f"{fp.name} — {w_s}×{h_s}",
    }


# ───────────────────────────── helpers ─────────────────────────────────

def _register_image_file(src: Path, ctx: ToolContext, *, name_hint: str) -> str:
    from PIL import Image as _Image

    sha = sha256_file(src)
    ext = src.suffix.lower() if src.suffix else ".png"
    if ext == ".jpeg":
        ext = ".jpg"
    layer_id = f"ingest_img_{sha[:8]}"
    dest = ctx.layers_dir / f"img_{layer_id}{ext}"
    if not dest.exists():
        shutil.copy2(src, dest)

    try:
        with _Image.open(dest) as im:
            w, h = im.size
    except Exception as e:
        raise RuntimeError(f"image not readable: {src} ({e})")

    ctx.state["rendered_layers"][layer_id] = {
        "layer_id": layer_id,
        "name": _sanitize_name(name_hint) or layer_id,
        "kind": "image",
        "z_index": 5,
        "bbox": None,
        "src_path": str(dest),
        "aspect_ratio": _aspect_from_dims(w, h),
        "image_size": f"{w}x{h}",
        "sha256": sha,
        "source": "ingested",
        "source_file": str(src),
    }
    return layer_id


def _curation_payload(rec: dict[str, Any]) -> dict[str, Any]:
    out: dict[str, Any] = {}
    for key in (
        "visual_role",
        "visual_score",
        "curation_reason",
        "curation_flags",
        "protected_anchor",
        "anchor_kind",
        "anchor_label",
        "anchor_reason",
        "captioned_source_group",
        "source_group_id",
        "source_group_kind",
        "source_group_label",
        "source_group_caption",
        "source_group_source",
        "table_parse_status",
        "crop_quality_flags",
        "placement_quality_flags",
        "caption_association_method",
        "visual_selection_tier",
        "eligibility_policy_version",
        "unmatched_caption",
        "designer_eligible",
        "planner_eligible",
        "planner_visible",
        "designer_reject_reasons",
        "planner_reject_reasons",
        "severe_crop_flags",
    ):
        if key in rec:
            out[key] = rec.get(key)
    return out


def _merged_visual_flags(*records: dict[str, Any] | None) -> list[str]:
    out: list[str] = []
    for rec in records:
        if not isinstance(rec, dict):
            continue
        for key in ("crop_quality_flags", "curation_flags", "severe_crop_flags"):
            for raw in list(rec.get(key) or []):
                flag = str(raw or "").strip()
                if flag and flag not in out:
                    out.append(flag)
        material_quality = rec.get("material_quality")
        if isinstance(material_quality, dict):
            for raw in list(material_quality.get("warnings") or []):
                flag = str(raw or "").strip()
                if flag and flag not in out:
                    out.append(flag)
    return out


def _severe_crop_flags_for_records(*records: dict[str, Any] | None) -> list[str]:
    return [flag for flag in _merged_visual_flags(*records) if flag in _SEVERE_CROP_CURATION_FLAGS]


def _explicit_false(key: str, *records: dict[str, Any] | None) -> bool:
    return any(isinstance(rec, dict) and rec.get(key) is False for rec in records)


def _unique_reason_strings(values: list[str]) -> list[str]:
    out: list[str] = []
    for value in values:
        item = str(value or "").strip()
        if item and item not in out:
            out.append(item)
    return out


def _source_visual_quality_reject_reasons(
    layer_id: str,
    *records: dict[str, Any] | None,
) -> list[str]:
    layer_id = str(layer_id or "").strip()
    if not layer_id.startswith(("ingest_fig_", "ingest_table_")):
        return []

    reasons: list[str] = []
    flags = set(_merged_visual_flags(*records))
    non_severe_blocking = flags & (_SELECTED_VISUAL_BLOCKING_FLAGS - _SEVERE_CROP_CURATION_FLAGS)
    for flag in sorted(non_severe_blocking):
        reasons.append(f"selected_blocking_flag:{flag}")

    score_values = [
        _safe_int((rec or {}).get("visual_score"), default=0)
        for rec in records
        if isinstance(rec, dict)
    ]
    score = max(score_values) if score_values else 0
    role = ""
    for rec in records:
        if isinstance(rec, dict) and str(rec.get("visual_role") or "").strip():
            role = str(rec.get("visual_role") or "").strip().lower()
            break
    has_captioned_group = any(
        isinstance(rec, dict)
        and bool(
            rec.get("captioned_source_group")
            or rec.get("source_group_id")
            or rec.get("source_group_label")
            or rec.get("source_group_caption")
        )
        for rec in records
    )
    caption_text = " ".join(
        str(rec.get(key) or "")
        for rec in records
        if isinstance(rec, dict)
        for key in ("caption", "caption_text", "caption_full", "caption_short", "source_group_caption", "title")
    ).strip()
    strong_captioned_anchor = bool(has_captioned_group and caption_text)
    if "low_caption_confidence" in flags and (not strong_captioned_anchor or score < 72):
        reasons.append("selected_weak_caption")
    if "no_caption" in flags and not strong_captioned_anchor:
        reasons.append("selected_missing_caption")
    if "source_page_unknown" in flags and not strong_captioned_anchor:
        reasons.append("selected_source_page_unknown")

    low_detail_is_weak = (
        score < 72
        or not strong_captioned_anchor
        or "no_caption" in flags
        or role not in {"method", "table", "evidence"}
    )
    if "low_detail_visual_content" in flags and low_detail_is_weak:
        reasons.append("selected_low_detail_weak_crop")

    mostly_white_is_weak = (
        score < 72
        or "no_caption" in flags
        or low_detail_is_weak
        or ("low_caption_confidence" in flags and not strong_captioned_anchor)
    )
    if "mostly_white_visual" in flags and mostly_white_is_weak:
        reasons.append("selected_mostly_white_weak_crop")

    edge_whitespace_is_weak = (
        score < 68
        or low_detail_is_weak
        or ("low_caption_confidence" in flags and not strong_captioned_anchor)
        or (role not in {"method", "table", "evidence"} and not (strong_captioned_anchor and score >= 72))
    )
    if "high_edge_whitespace" in flags and edge_whitespace_is_weak:
        reasons.append("selected_edge_whitespace_weak_crop")

    is_table = layer_id.startswith("ingest_table_") or any(
        isinstance(rec, dict) and str(rec.get("kind") or "").lower() == "table"
        for rec in records
    )
    if is_table:
        has_rows = any(isinstance(rec, dict) and bool(rec.get("rows")) for rec in records)
        has_headers = any(isinstance(rec, dict) and bool(rec.get("headers")) for rec in records)
        parse_status = next(
            (
                str(rec.get("table_parse_status") or "").strip()
                for rec in records
                if isinstance(rec, dict) and str(rec.get("table_parse_status") or "").strip()
            ),
            "",
        )
        if not has_rows and not has_headers and not strong_captioned_anchor:
            reasons.append("selected_table_missing_structure")
        unparsed_table_hard_flags = flags & {
            "body_text_leak",
            "caption_in_crop",
            "no_caption",
            "table_body_text_leak",
        }
        if "low_caption_confidence" in flags and not strong_captioned_anchor:
            unparsed_table_hard_flags.add("low_caption_confidence")
        if parse_status == "unparsed_source_crop" and unparsed_table_hard_flags:
            reasons.append("selected_unparsed_table_weak_crop")

    return _unique_reason_strings(reasons)


def _designer_reject_reasons(
    layer_id: str,
    rec: dict[str, Any] | None,
    *extra_records: dict[str, Any] | None,
) -> list[str]:
    if _is_source_visual_layer_id(layer_id):
        return list(
            classify_source_visual(layer_id, rec, *extra_records).get("designer_reject_reasons") or []
        )
    records = (rec or {}, *extra_records)
    reasons: list[str] = []
    if _explicit_false("designer_eligible", *records):
        reasons.append("designer_eligible=false")
    if _is_uncaptioned_object_level_crop(layer_id, rec or {}):
        reasons.append("uncaptioned_object_level_crop")
    if str((rec or {}).get("kind") or (rec or {}).get("extract_strategy") or "") == "source_table_crop_candidate":
        reasons.append("audit_only_source_table_crop_candidate")
    for flag in _severe_crop_flags_for_records(*records):
        reasons.append(f"severe_crop:{flag}")
    reasons.extend(_source_visual_quality_reject_reasons(layer_id, *records))
    return _unique_reason_strings(reasons)


def _planner_reject_reasons(
    layer_id: str,
    rec: dict[str, Any] | None,
    *extra_records: dict[str, Any] | None,
) -> list[str]:
    if _is_source_visual_layer_id(layer_id):
        return list(
            classify_source_visual(layer_id, rec, *extra_records).get("planner_reject_reasons") or []
        )
    rec = rec or {}
    records = (rec, *extra_records)
    reasons = _designer_reject_reasons(layer_id, rec, *extra_records)
    if _explicit_false("planner_visible", *records):
        reasons.append("planner_visible=false")
    if _explicit_false("planner_eligible", *records):
        reasons.append("planner_eligible=false")
    if _is_partial_pdf_crop_record(layer_id, rec):
        reasons.append("partial_pdf_crop")
    if str(rec.get("kind") or rec.get("extract_strategy") or "") == "source_table_crop_candidate":
        reasons.append("audit_only_source_table_crop_candidate")
    flags = set(_merged_visual_flags(*records))
    if "unlocated_raster_component" in flags:
        reasons.append("unlocated_raster_component")
    return _unique_reason_strings(reasons)


def _is_designer_eligible(
    layer_id: str,
    rec: dict[str, Any] | None,
    *extra_records: dict[str, Any] | None,
) -> bool:
    return not _designer_reject_reasons(layer_id, rec, *extra_records)


def _is_planner_visible(
    layer_id: str,
    rec: dict[str, Any] | None,
    *extra_records: dict[str, Any] | None,
) -> bool:
    return not _planner_reject_reasons(layer_id, rec, *extra_records)


def _visual_eligibility_payload(
    layer_id: str,
    rec: dict[str, Any] | None,
    *extra_records: dict[str, Any] | None,
) -> dict[str, Any]:
    if _is_source_visual_layer_id(layer_id):
        return classify_source_visual(layer_id, rec, *extra_records)
    designer_reasons = _designer_reject_reasons(layer_id, rec, *extra_records)
    planner_reasons = _planner_reject_reasons(layer_id, rec, *extra_records)
    return {
        "designer_eligible": not designer_reasons,
        "planner_eligible": not planner_reasons,
        "planner_visible": not planner_reasons,
        "designer_reject_reasons": designer_reasons,
        "planner_reject_reasons": planner_reasons,
        "severe_crop_flags": _severe_crop_flags_for_records(rec, *extra_records),
    }


def _is_source_visual_layer_id(layer_id: str) -> bool:
    return str(layer_id or "").startswith(("ingest_fig_", "ingest_table_"))


def _selected_visual_reject_reasons(
    layer_id: str,
    *records: dict[str, Any] | None,
    forbidden_ids: set[str] | None = None,
    allow_reserve: bool = False,
) -> list[str]:
    layer_id = str(layer_id or "").strip()
    reasons: list[str] = []
    if not layer_id:
        return ["missing_layer_id"]
    if not _is_source_visual_layer_id(layer_id):
        return []
    if forbidden_ids and layer_id in forbidden_ids:
        reasons.append("forbidden_source_visual")
    primary = next((rec for rec in records if isinstance(rec, dict) and rec), {})
    classification = classify_source_visual(layer_id, *records)
    if (
        classification.get("visual_selection_tier") == "reserve_unmatched"
        and not allow_reserve
    ):
        reasons.append("reserve_unmatched_shortfall_only")
    reasons.extend(_planner_reject_reasons(layer_id, primary, *records[1:]))
    return _unique_reason_strings(reasons)


def _selected_visual_score_by_id(visual_candidate_scores: list[dict[str, Any]]) -> dict[str, dict[str, Any]]:
    return {
        str(item.get("layer_id") or ""): item
        for item in visual_candidate_scores or []
        if isinstance(item, dict) and str(item.get("layer_id") or "").strip()
    }


def _selected_visual_record_sources(
    layer_id: str,
    rendered: dict[str, dict[str, Any]],
    score_by_id: dict[str, dict[str, Any]],
    provenance_assets: dict[str, dict[str, Any]] | None = None,
) -> tuple[dict[str, Any], dict[str, Any], dict[str, Any]]:
    return (
        rendered.get(layer_id) or {},
        score_by_id.get(layer_id) or {},
        (provenance_assets or {}).get(layer_id) or {},
    )


def _selected_visual_preference_score(
    layer_id: str,
    rendered: dict[str, dict[str, Any]],
    score_by_id: dict[str, dict[str, Any]],
    provenance_assets: dict[str, dict[str, Any]] | None = None,
) -> tuple[int, int, int, int, int]:
    rec, score_rec, asset = _selected_visual_record_sources(layer_id, rendered, score_by_id, provenance_assets)
    flags = set(_merged_visual_flags(rec, score_rec, asset))
    role = str(asset.get("visual_role") or rec.get("visual_role") or score_rec.get("visual_role") or "").lower()
    role_score = 2 if role in {"method", "table", "evidence"} else 1 if role == "qualitative" else 0
    visual_score = max(
        _safe_int(rec.get("visual_score"), default=0),
        _safe_int(score_rec.get("visual_score"), default=0),
        _safe_int(asset.get("visual_score"), default=0),
    )
    width, height = _image_size_tuple(rec.get("image_size"))
    width = width or _safe_int(asset.get("output_width_px"), default=0)
    height = height or _safe_int(asset.get("output_height_px"), default=0)
    area_score = min(width * height, 3_000_000) // 10_000
    clean_bonus = 20 - min(20, len(flags & (_SELECTED_VISUAL_BLOCKING_FLAGS | _SELECTED_VISUAL_CONDITIONAL_FLAGS)) * 5)
    anchor_bonus = 1 if any(
        bool(source.get("protected_anchor") or source.get("captioned_source_group") or source.get("source_group_id"))
        for source in (rec, score_rec, asset)
        if isinstance(source, dict)
    ) else 0
    return (clean_bonus, role_score, visual_score, anchor_bonus, area_score)


def _bbox_tuple_from_visual_records(*records: dict[str, Any]) -> tuple[float, float, float, float] | None:
    for rec in records:
        raw = rec.get("source_bbox_pdf_points") if isinstance(rec, dict) else None
        if not isinstance(raw, (list, tuple)) or len(raw) < 4:
            continue
        try:
            x0, y0, x1, y1 = [float(v) for v in raw[:4]]
        except (TypeError, ValueError):
            continue
        if x1 > x0 and y1 > y0:
            return x0, y0, x1, y1
    return None


def _bbox_overlap_stats(
    a: tuple[float, float, float, float],
    b: tuple[float, float, float, float],
) -> tuple[float, float]:
    ax0, ay0, ax1, ay1 = a
    bx0, by0, bx1, by1 = b
    inter_w = max(0.0, min(ax1, bx1) - max(ax0, bx0))
    inter_h = max(0.0, min(ay1, by1) - max(ay0, by0))
    inter = inter_w * inter_h
    area_a = max(1.0, (ax1 - ax0) * (ay1 - ay0))
    area_b = max(1.0, (bx1 - bx0) * (by1 - by0))
    union = max(area_a + area_b - inter, 1.0)
    return inter / union, inter / max(1.0, min(area_a, area_b))


def _visual_average_hash(layer_id: str, rec: dict[str, Any], cache: dict[str, int | None]) -> int | None:
    if layer_id in cache:
        return cache[layer_id]
    raw_path = str(rec.get("src_path") or rec.get("output_file") or "").strip()
    if not raw_path:
        cache[layer_id] = None
        return None
    path = Path(raw_path)
    if not path.is_absolute():
        cache[layer_id] = None
        return None
    try:
        with Image.open(path) as im:
            gray = im.convert("L").resize((8, 8))
            pixels = list(gray.getdata())
    except Exception:
        cache[layer_id] = None
        return None
    mean = sum(pixels) / max(1, len(pixels))
    value = 0
    for idx, pixel in enumerate(pixels):
        if pixel >= mean:
            value |= 1 << idx
    cache[layer_id] = value
    return value


def _hamming_distance(a: int, b: int) -> int:
    return int((a ^ b).bit_count())


def _selected_visuals_are_duplicate(
    a_id: str,
    b_id: str,
    rendered: dict[str, dict[str, Any]],
    score_by_id: dict[str, dict[str, Any]],
    provenance_assets: dict[str, dict[str, Any]] | None,
    hash_cache: dict[str, int | None],
) -> bool:
    if a_id == b_id:
        return True
    a_rec, a_score, a_asset = _selected_visual_record_sources(a_id, rendered, score_by_id, provenance_assets)
    b_rec, b_score, b_asset = _selected_visual_record_sources(b_id, rendered, score_by_id, provenance_assets)
    for key in ("source_group_id", "output_sha256", "sha256"):
        a_val = str(a_asset.get(key) or a_rec.get(key) or a_score.get(key) or "").strip()
        b_val = str(b_asset.get(key) or b_rec.get(key) or b_score.get(key) or "").strip()
        if a_val and b_val and a_val == b_val:
            return True
    a_page = _safe_int(a_asset.get("source_page") or a_rec.get("source_page") or a_score.get("source_page"), default=-1)
    b_page = _safe_int(b_asset.get("source_page") or b_rec.get("source_page") or b_score.get("source_page"), default=-2)
    if a_page >= 0 and a_page == b_page:
        a_xref = str(a_asset.get("source_image_xref") or a_rec.get("source_image_xref") or "").strip()
        b_xref = str(b_asset.get("source_image_xref") or b_rec.get("source_image_xref") or "").strip()
        if a_xref and b_xref and a_xref == b_xref:
            return True
        a_bbox = _bbox_tuple_from_visual_records(a_asset, a_rec, a_score)
        b_bbox = _bbox_tuple_from_visual_records(b_asset, b_rec, b_score)
        if a_bbox and b_bbox:
            iou, coverage = _bbox_overlap_stats(a_bbox, b_bbox)
            if iou >= 0.72 or coverage >= 0.88:
                return True
            a_flags = set(_merged_visual_flags(a_rec, a_score, a_asset))
            b_flags = set(_merged_visual_flags(b_rec, b_score, b_asset))
            if coverage >= 0.62 and (a_flags | b_flags) & {
                "high_edge_whitespace",
                "low_caption_confidence",
                "low_detail_visual_content",
                "low_information_visual",
                "neighbor_asset_leak",
                "partial_visual_crop",
            }:
                return True
    if a_page >= 0 and a_page == b_page:
        a_hash = _visual_average_hash(a_id, a_rec, hash_cache)
        b_hash = _visual_average_hash(b_id, b_rec, hash_cache)
        if a_hash is not None and b_hash is not None and _hamming_distance(a_hash, b_hash) <= 5:
            return True
    return False


def _sanitize_selected_visual_ids(
    values: list[str],
    *,
    rendered: dict[str, dict[str, Any]],
    visual_candidate_scores: list[dict[str, Any]],
    provenance_assets: dict[str, dict[str, Any]] | None = None,
    forbidden_visual_ids: list[str] | None = None,
    limit: int = 0,
    max_tables: int = 3,
    allow_reserve: bool = False,
) -> list[str]:
    score_by_id = _selected_visual_score_by_id(visual_candidate_scores)
    forbidden = {str(item or "").strip() for item in (forbidden_visual_ids or []) if str(item or "").strip()}
    out: list[str] = []
    hash_cache: dict[str, int | None] = {}
    for raw in values or []:
        layer_id = str(raw or "").strip()
        if not layer_id or layer_id not in rendered:
            continue
        rec, score_rec, asset = _selected_visual_record_sources(layer_id, rendered, score_by_id, provenance_assets)
        if _selected_visual_reject_reasons(
            layer_id,
            rec,
            score_rec,
            asset,
            forbidden_ids=forbidden,
            allow_reserve=allow_reserve,
        ):
            continue
        if layer_id.startswith("ingest_table_") and sum(1 for item in out if item.startswith("ingest_table_")) >= max_tables:
            continue
        duplicate_idx = None
        for idx, existing_id in enumerate(out):
            if _selected_visuals_are_duplicate(
                existing_id,
                layer_id,
                rendered,
                score_by_id,
                provenance_assets,
                hash_cache,
            ):
                duplicate_idx = idx
                break
        if duplicate_idx is not None:
            existing_id = out[duplicate_idx]
            if _selected_visual_preference_score(layer_id, rendered, score_by_id, provenance_assets) > _selected_visual_preference_score(
                existing_id,
                rendered,
                score_by_id,
                provenance_assets,
            ):
                out[duplicate_idx] = layer_id
            continue
        out.append(layer_id)
        if limit > 0 and len(out) >= limit:
            break
    return out


def _is_uncaptioned_object_level_crop(layer_id: str, rec: dict[str, Any]) -> bool:
    kind = str(rec.get("kind") or "").lower()
    if kind == "table" or layer_id.startswith("ingest_table_"):
        return False
    strategy = str(rec.get("extract_strategy") or "").lower()
    if strategy not in {"raster", "vector", "sub_panel", "embedded"}:
        return False
    if str(rec.get("caption") or rec.get("title") or rec.get("caption_short") or "").strip():
        return False
    if bool(rec.get("protected_anchor") or rec.get("captioned_source_group")):
        return False
    return not bool(
        str(rec.get("source_group_label") or "").strip()
        or str(rec.get("source_group_caption") or "").strip()
        or str(rec.get("anchor_label") or "").strip()
    )


def _annotate_visual_curation(
    layer_ids: list[str],
    rendered: dict[str, dict[str, Any]],
) -> list[dict[str, Any]]:
    """Attach deterministic planner-facing curation metadata to PDF visuals.

    This pass only sees registered figure/table ids, so VLM-rejected fake
    figures never enter `visual_candidate_scores`.
    """
    scored: list[dict[str, Any]] = []
    for layer_id in layer_ids:
        rec = rendered.get(layer_id)
        if not isinstance(rec, dict):
            continue
        if _is_partial_pdf_crop_record(str(layer_id), rec):
            continue
        rec.update(_score_visual_candidate(layer_id, rec))
        if _hide_from_planner_visual_catalog(str(layer_id), rec):
            continue
        scored.append(_visual_candidate_score_record(layer_id, rec))
    return sorted(scored, key=_visual_candidate_score_key, reverse=True)


def _score_visual_candidate(layer_id: str, rec: dict[str, Any]) -> dict[str, Any]:
    role = _primary_visual_role(layer_id, rec)
    flags = _visual_curation_flags(layer_id, rec)
    score = 20
    score += {
        "method": 34,
        "evidence": 32,
        "table": 32,
        "qualitative": 10,
        "fallback": 10,
    }.get(role, 10)

    caption = str(rec.get("caption") or rec.get("title") or "").strip()
    if caption:
        score += 12
    else:
        score -= 4
    if str(rec.get("caption_short") or "").strip():
        score += 3

    strategy = str(rec.get("extract_strategy") or "").lower()
    score += {
        "vector": 12,
        "sub_panel": 8,
        "raster": 5,
        "embedded": 4,
    }.get(strategy, 3)
    if str(rec.get("kind") or "") == "table" or layer_id.startswith("ingest_table_"):
        rows = len(rec.get("rows") or [])
        cols = len(rec.get("headers") or [])
        score += min(12, rows // 2 + cols)
    else:
        score += _image_size_score(rec.get("image_size"))

    page = _safe_int(rec.get("source_page"), default=999)
    if page <= 2:
        score += 10
    elif page <= 5:
        score += 8
    elif page <= 10:
        score += 6
    elif page <= 20:
        score += 3

    if layer_id.startswith("ingest_fig_"):
        confidence = _safe_float(rec.get("caption_confidence"), default=0.0)
        if confidence >= 0.75:
            score += 6
        elif confidence >= _CAPTION_MATCH_MIN_CONFIDENCE:
            score += 3
        else:
            score -= 6

    if "low_information_visual" in flags:
        score -= 18
    if "low_resolution" in flags:
        score -= 5
    if "no_caption" in flags:
        score -= 4
    if "table_too_dense" in flags:
        score -= 2
    if "caption_in_crop" in flags:
        score -= 12
    if "body_text_leak" in flags:
        score -= 24
    if "page_like_figure_crop" in flags:
        score -= 24
    if "low_value_example_crop" in flags:
        score -= 22
    severe_flags = [flag for flag in flags if flag in _SEVERE_CROP_CURATION_FLAGS]
    if severe_flags:
        score -= 40 + (len(severe_flags) - 1) * 8

    score = max(0, min(100, int(round(score))))
    scored_rec = {**rec, "visual_score": score, "curation_flags": flags}
    eligibility = _visual_eligibility_payload(layer_id, scored_rec)
    if not eligibility["designer_eligible"]:
        score = min(score, 20)
    return {
        "visual_role": role,
        "visual_score": score,
        "curation_reason": _visual_curation_reason(role, score, rec, flags),
        "curation_flags": flags,
        **eligibility,
    }


def _visual_curation_flags(layer_id: str, rec: dict[str, Any]) -> list[str]:
    flags: list[str] = []
    def add(flag: str) -> None:
        flag_s = str(flag or "").strip()
        if flag_s and flag_s not in flags:
            flags.append(flag_s)

    for flag in rec.get("crop_quality_flags") or []:
        add(str(flag or ""))
    if _is_partial_pdf_crop_record(layer_id, rec):
        add("partial_visual_crop")
    caption = str(rec.get("caption") or rec.get("title") or "").strip()
    if not caption:
        add("no_caption")
    if layer_id.startswith("ingest_fig_"):
        confidence = _safe_float(rec.get("caption_confidence"), default=0.0)
        if confidence < _CAPTION_MATCH_MIN_CONFIDENCE:
            add("low_caption_confidence")
    min_side = _min_image_side(rec.get("image_size"))
    if 0 < min_side < 220:
        add("low_resolution")
    if _safe_int(rec.get("source_page"), default=0) <= 0:
        add("source_page_unknown")
    strategy = str(rec.get("extract_strategy") or "").lower()
    source_bbox = rec.get("source_bbox_pdf_points")
    width_px, height_px = _image_size_tuple(rec.get("image_size"))
    max_side = max(width_px or 0, height_px or 0)
    if strategy == "raster" and not source_bbox and 0 < max_side < 900:
        add("unlocated_raster_component")
        add("low_information_visual")
    text = " ".join(str(rec.get(k) or "") for k in (
        "layer_id", "name", "caption", "caption_short", "title",
        "source_group_caption", "source_group_kind", "source_group_label",
    )).lower()
    if any(term in text for term in _LOW_INFORMATION_VISUAL_TERMS):
        add("low_information_visual")
    for flag in _image_content_curation_flags(layer_id, rec):
        add(flag)
    if _looks_like_low_value_example_visual(rec):
        add("low_value_example_crop")
    if str(rec.get("kind") or "") == "table" or layer_id.startswith("ingest_table_"):
        rows = len(rec.get("rows") or [])
        cols = len(rec.get("headers") or [])
        if rows > 24 or cols > 8:
            add("table_too_dense")
    return flags


def _image_content_curation_flags(layer_id: str, rec: dict[str, Any]) -> list[str]:
    """Flag visually weak source crops that metadata scoring cannot catch."""
    if not str(layer_id or "").startswith("ingest_fig_"):
        return []
    if str(rec.get("kind") or "").lower() == "table":
        return []
    raw_path = str(rec.get("src_path") or rec.get("output_file") or "").strip()
    if not raw_path:
        return []
    path = Path(raw_path)
    if not path.exists() or not path.is_file():
        return []
    try:
        from PIL import Image
    except Exception:
        return []
    try:
        with Image.open(path) as im:
            if im.width <= 0 or im.height <= 0:
                return []
            im = im.convert("RGB")
            scale = min(1.0, 128.0 / float(max(im.width, im.height)))
            sample_size = (
                max(1, int(round(im.width * scale))),
                max(1, int(round(im.height * scale))),
            )
            resampling = getattr(getattr(Image, "Resampling", Image), "LANCZOS", 1)
            sample = im.resize(sample_size, resampling)
            gray = sample.convert("L")
            pixels = list(gray.getdata())
            w, h = gray.size
            if w <= 1 and h <= 1:
                return []
            total_edges = 0
            strong_edges = 0
            diff_sum = 0
            for y in range(h):
                row = y * w
                for x in range(1, w):
                    diff = abs(int(pixels[row + x]) - int(pixels[row + x - 1]))
                    diff_sum += diff
                    total_edges += 1
                    if diff > 18:
                        strong_edges += 1
            for y in range(1, h):
                row = y * w
                prev = (y - 1) * w
                for x in range(w):
                    diff = abs(int(pixels[row + x]) - int(pixels[prev + x]))
                    diff_sum += diff
                    total_edges += 1
                    if diff > 18:
                        strong_edges += 1
            if total_edges <= 0:
                return []
            edge_density = strong_edges / float(total_edges)
            mean_edge_delta = diff_sum / float(total_edges)
    except Exception:
        return []

    if edge_density >= 0.07 or mean_edge_delta >= 3.5:
        return []

    flags = ["low_detail_visual_content"]
    caption = str(rec.get("caption") or rec.get("title") or rec.get("caption_short") or "").strip()
    caption_confidence = _safe_float(rec.get("caption_confidence"), default=0.0)
    has_anchor = bool(
        rec.get("protected_anchor")
        or rec.get("captioned_source_group")
        or str(rec.get("source_group_label") or "").strip()
        or str(rec.get("anchor_label") or "").strip()
    )
    if not caption or caption_confidence < _CAPTION_MATCH_MIN_CONFIDENCE or not has_anchor:
        flags.append("low_information_visual")
    return flags


def _hide_from_planner_visual_catalog(layer_id: str, rec: dict[str, Any]) -> bool:
    return bool(_planner_reject_reasons(layer_id, rec))


def _looks_like_low_value_example_visual(rec: dict[str, Any]) -> bool:
    if str(rec.get("kind") or "").lower() == "table":
        return False
    text = " ".join(str(rec.get(k) or "") for k in (
        "layer_id", "name", "caption", "caption_short", "title",
        "source_group_caption", "visual_role",
    )).lower()
    if not any(token in text for token in ("example", "sample", "demo", "case", "qualitative")):
        return False
    high_value_terms = (
        "method", "architecture", "pipeline", "framework", "benchmark",
        "result", "evaluation", "ablation", "comparison", "performance",
    )
    return not any(term in text for term in high_value_terms)


def _visual_curation_reason(
    role: str,
    score: int,
    rec: dict[str, Any],
    flags: list[str],
) -> str:
    caption = str(rec.get("caption") or rec.get("title") or "").strip()
    page = rec.get("source_page") or "?"
    strategy = rec.get("extract_strategy") or rec.get("kind") or "visual"
    parts = [
        f"role={role}",
        "captioned" if caption else "no caption",
        f"strategy={strategy}",
        f"page={page}",
        f"score={score}",
    ]
    if flags:
        parts.append("flags=" + ",".join(flags[:4]))
    return "; ".join(parts)


def _visual_candidate_score_record(layer_id: str, rec: dict[str, Any]) -> dict[str, Any]:
    kind = rec.get("kind") or (
        "table" if layer_id.startswith("ingest_table_") else "image"
    )
    return {
        "layer_id": layer_id,
        "kind": kind,
        "source_page": rec.get("source_page"),
        "caption_short": rec.get("caption_short"),
        "visual_role": rec.get("visual_role"),
        "visual_score": rec.get("visual_score"),
        "curation_reason": rec.get("curation_reason"),
        "curation_flags": list(rec.get("curation_flags") or []),
        "crop_quality_flags": list(rec.get("crop_quality_flags") or []),
        "placement_quality_flags": list(rec.get("placement_quality_flags") or []),
        "caption_association_method": rec.get("caption_association_method") or "unmatched",
        "visual_selection_tier": rec.get("visual_selection_tier"),
        "eligibility_policy_version": rec.get("eligibility_policy_version"),
        "unmatched_caption": bool(rec.get("unmatched_caption")),
        "designer_eligible": bool(rec.get("designer_eligible", True)),
        "planner_eligible": bool(rec.get("planner_eligible", True)),
        "planner_visible": bool(rec.get("planner_visible", True)),
        "designer_reject_reasons": list(rec.get("designer_reject_reasons") or []),
        "planner_reject_reasons": list(rec.get("planner_reject_reasons") or []),
        "severe_crop_flags": list(rec.get("severe_crop_flags") or []),
        "protected_anchor": bool(rec.get("protected_anchor")),
        "anchor_kind": rec.get("anchor_kind"),
        "anchor_label": rec.get("anchor_label"),
        "anchor_reason": rec.get("anchor_reason"),
        "captioned_source_group": bool(rec.get("captioned_source_group")),
        "source_group_id": rec.get("source_group_id"),
        "source_group_kind": rec.get("source_group_kind"),
        "source_group_label": rec.get("source_group_label"),
        "source_group_caption": rec.get("source_group_caption"),
        "source_group_source": rec.get("source_group_source"),
        "table_parse_status": rec.get("table_parse_status"),
    }


def _visual_candidate_score_key(item: dict[str, Any]) -> tuple[int, ...]:
    flags = set(item.get("curation_flags") or [])
    role = str(item.get("visual_role") or "")
    return (
        1 if _is_planner_visible(str(item.get("layer_id") or ""), {}, item) else 0,
        1 if role in {"method", "table", "evidence"} else 0,
        0 if "low_value_example_crop" in flags else 1,
        _safe_int(item.get("visual_score"), default=0),
        _VISUAL_ROLE_PRIORITY.get(role, 0),
        -_safe_int(item.get("source_page"), default=999),
    )


def _primary_visual_role(layer_id: str, rec: dict[str, Any]) -> str:
    roles = _infer_visual_roles(layer_id, rec)
    if "table" in roles:
        return "table"
    for role in ("method", "evidence", "qualitative"):
        if role in roles:
            return role
    return roles[0] if roles else "fallback"


def _infer_visual_roles(layer_id: str, rec: dict[str, Any]) -> list[str]:
    if str(rec.get("kind") or "") == "table" or layer_id.startswith("ingest_table_"):
        return ["table", "evidence"]
    text = " ".join(str(rec.get(k) or "") for k in (
        "layer_id", "name", "caption", "caption_short", "title",
        "anchor_kind", "anchor_label", "anchor_reason",
        "source_group_caption", "source_group_kind", "source_group_label",
    )).lower()
    roles: list[str] = []
    if any(k in text for k in (
        "method", "overview", "architecture", "pipeline", "framework",
        "model", "algorithm", "approach", "system", "workflow",
        "流程", "方法", "架构",
    )):
        roles.append("method")
    if any(k in text for k in (
        "result", "benchmark", "experiment", "evaluation", "ablation",
        "scaling", "performance", "metric", "comparison", "baseline",
        "sota", "state-of-the-art", "accuracy", "结果", "实验", "消融",
    )):
        roles.append("evidence")
    if any(k in text for k in (
        "qualitative", "visual", "example", "sample", "demo", "motion",
        "image", "video", "sequence", "case", "showcase", "示例", "案例",
    )):
        roles.append("qualitative")
    if not roles:
        roles.append("evidence")
    return roles


def _visual_rank_key(layer_id: str, rec: dict[str, Any]) -> tuple[int, ...]:
    caption = str(rec.get("caption") or rec.get("title") or "").strip()
    page = _safe_int(rec.get("source_page"), default=999)
    role = str(rec.get("visual_role") or "")
    flags = set(rec.get("curation_flags") or [])
    density_or_side = (
        _table_density(rec)
        if layer_id.startswith("ingest_table_")
        else _min_image_side(rec.get("image_size"))
    )
    return (
        1 if _is_planner_visible(layer_id, rec) else 0,
        1 if role in {"method", "table", "evidence"} else 0,
        0 if "low_value_example_crop" in flags else 1,
        _safe_int(rec.get("visual_score"), default=0),
        _VISUAL_ROLE_PRIORITY.get(str(rec.get("visual_role") or ""), 0),
        1 if caption else 0,
        density_or_side,
        -page,
    )


def _rank_figure_ids_for_planner(
    figure_ids: list[str], rendered: dict[str, dict[str, Any]],
) -> list[str]:
    """Order figure ids so the ones most useful to a planner come first.

    When PDF curation metadata is present, `visual_score` is the primary
    signal. Older non-PDF ingest records keep the legacy ranking:

    1. has a non-empty caption (carries explicit author intent),
    2. larger min-dimension (more printable at poster scale),
    3. source strategy = "vector" (composite diagrams > raster sub-panels),
    4. smaller page number (main-paper figures before appendix ones).
    """
    figure_ids = [
        str(fid)
        for fid in figure_ids
        if str(fid or "").strip()
        and not _hide_from_planner_visual_catalog(str(fid), rendered.get(str(fid)) or {})
    ]
    if any("visual_score" in (rendered.get(fid) or {}) for fid in figure_ids):
        return sorted(
            figure_ids,
            key=lambda fid: _visual_rank_key(fid, rendered.get(fid) or {}),
            reverse=True,
        )

    def key(fid: str) -> tuple[int, int, int, int]:
        rec = rendered.get(fid) or {}
        cap = (rec.get("caption") or "").strip()
        size = rec.get("image_size") or "0x0"
        try:
            w_s, h_s = size.split("x")
            w = int(w_s); h = int(h_s)
            side = min(w, h)
        except Exception:
            side = 0
        strat_rank = 1 if rec.get("extract_strategy") == "vector" else 0
        has_caption_rank = 1 if cap else 0
        try:
            page = int(rec.get("source_page") or 999)
        except (TypeError, ValueError):
            page = 999
        # Sort DESC by caption/strategy/size, ASC by page → negate page.
        return (has_caption_rank, strat_rank, side, -page)

    return sorted(figure_ids, key=key, reverse=True)


def _rank_visual_ids_for_planner(
    layer_ids: list[str], rendered: dict[str, dict[str, Any]],
) -> list[str]:
    """Rank figure + table ids for dense paper-poster planning."""
    layer_ids = [
        str(layer_id)
        for layer_id in layer_ids
        if str(layer_id or "").strip()
        and not _is_partial_pdf_crop_record(str(layer_id), rendered.get(str(layer_id)) or {})
        and not _hide_from_planner_visual_catalog(str(layer_id), rendered.get(str(layer_id)) or {})
    ]
    if any("visual_score" in (rendered.get(layer_id) or {}) for layer_id in layer_ids):
        return sorted(
            layer_ids,
            key=lambda layer_id: _visual_rank_key(layer_id, rendered.get(layer_id) or {}),
            reverse=True,
        )

    ranked_figs = _rank_figure_ids_for_planner(
        [fid for fid in layer_ids if str(fid).startswith("ingest_fig_")],
        rendered,
    )
    table_ids = [tid for tid in layer_ids if str(tid).startswith("ingest_table_")]

    def table_key(tid: str) -> tuple[int, int, int]:
        rec = rendered.get(tid) or {}
        rows = len(rec.get("rows") or [])
        cols = len(rec.get("headers") or [])
        try:
            page = int(rec.get("source_page") or 999)
        except (TypeError, ValueError):
            page = 999
        cap = 1 if (rec.get("caption") or rec.get("title")) else 0
        return (cap, min(rows, 24) + min(cols, 12), -page)

    ranked_tables = sorted(table_ids, key=table_key, reverse=True)
    # Interleave tables after the strongest early figures so evidence tables
    # are planner-visible without crowding out method/qualitative visuals.
    return ranked_figs[:4] + ranked_tables[:4] + ranked_figs[4:] + ranked_tables[4:]


def _image_size_score(raw: Any) -> int:
    side = _min_image_side(raw)
    if side >= 900:
        return 14
    if side >= 600:
        return 11
    if side >= 360:
        return 8
    if side >= 220:
        return 4
    return 0


def _min_image_side(raw: Any) -> int:
    try:
        w_s, h_s = str(raw or "0x0").split("x")
        return min(int(w_s), int(h_s))
    except Exception:
        return 0


def _table_density(rec: dict[str, Any]) -> int:
    return min(len(rec.get("rows") or []), 24) + min(len(rec.get("headers") or []), 12)


def _safe_int(value: Any, *, default: int) -> int:
    try:
        return int(value)
    except (TypeError, ValueError):
        return default


def _safe_float(value: Any, *, default: float) -> float:
    try:
        return float(value)
    except (TypeError, ValueError):
        return default


def _build_ingest_contact_sheet(
    layer_ids: list[str],
    rendered: dict[str, dict[str, Any]],
    ctx: ToolContext,
    pdf_path: Path,
) -> Path | None:
    if not layer_ids:
        return None
    try:
        layer_ids = [
            str(layer_id)
            for layer_id in layer_ids
            if str(layer_id or "").strip()
            and not _is_partial_pdf_crop_record(str(layer_id), rendered.get(str(layer_id)) or {})
            and not _hide_from_planner_visual_catalog(str(layer_id), rendered.get(str(layer_id)) or {})
        ]
        if not layer_ids:
            return None
        cols = 4
        cell_w, cell_h = 520, 390
        thumb_w, thumb_h = 480, 255
        rows = (len(layer_ids) + cols - 1) // cols
        sheet = Image.new("RGB", (cols * cell_w, rows * cell_h), "#f8f8f4")
        draw = ImageDraw.Draw(sheet)
        font = _contact_sheet_font(ctx, 17)
        small_font = _contact_sheet_font(ctx, 13)

        for idx, layer_id in enumerate(layer_ids):
            rec = rendered.get(layer_id) or {}
            col = idx % cols
            row = idx // cols
            x = col * cell_w
            y = row * cell_h
            draw.rectangle((x + 10, y + 10, x + cell_w - 10, y + cell_h - 10),
                           outline="#d6d3cc", width=2)
            src = rec.get("src_path")
            if src and Path(str(src)).exists():
                with Image.open(str(src)) as img:
                    img = img.convert("RGB")
                    img.thumbnail((thumb_w, thumb_h), Image.Resampling.LANCZOS)
                    px = x + 20 + (thumb_w - img.width) // 2
                    py = y + 20 + (thumb_h - img.height) // 2
                    sheet.paste(img, (px, py))
            else:
                draw.rectangle((x + 20, y + 20, x + 20 + thumb_w, y + 20 + thumb_h),
                               outline="#b8b5ad", fill="#ece9e2")
                draw.text((x + 34, y + 128), "structured table", fill="#5f5a50", font=font)

            label = f"{layer_id} · p.{rec.get('source_page', '?')} · {rec.get('image_size', '')}"
            draw.text((x + 20, y + 290), label[:58], fill="#1f2933", font=font)
            caption = (
                rec.get("caption_short")
                or rec.get("caption")
                or rec.get("title")
                or rec.get("name")
                or ""
            )
            for line_i, line in enumerate(_wrap_contact_text(str(caption), 62)[:3]):
                draw.text((x + 20, y + 316 + line_i * 19),
                          line, fill="#55514a", font=small_font)

        out_path = ctx.layers_dir / f"ingest_contact_sheet_{pdf_path.stem}.png"
        sheet.save(out_path)
        log("ingest.pdf.contact_sheet",
            path=str(out_path), n=len(layer_ids), pdf=pdf_path.name)
        return out_path
    except Exception as e:  # noqa: BLE001 - contact sheet is advisory
        log("ingest.pdf.contact_sheet_failed", pdf=pdf_path.name, error=str(e))
        return None


def _contact_sheet_font(ctx: ToolContext, size: int) -> ImageFont.ImageFont:
    candidates = [
        ctx.settings.fonts_dir / "NotoSansSC-Bold.otf",
        ctx.settings.fonts_dir / "NotoSerifSC-Bold.otf",
    ]
    for path in candidates:
        if path.exists():
            try:
                return ImageFont.truetype(str(path), size=size)
            except OSError:
                pass
    return ImageFont.load_default()


def _wrap_contact_text(text: str, width: int) -> list[str]:
    text = " ".join(str(text).replace("\n", " ").split())
    if not text:
        return []
    words = text.split(" ")
    lines: list[str] = []
    cur = ""
    for word in words:
        nxt = word if not cur else f"{cur} {word}"
        if len(nxt) <= width:
            cur = nxt
            continue
        if cur:
            lines.append(cur)
        cur = word[:width]
    if cur:
        lines.append(cur)
    return lines


def _recommend_paper_visuals(
    ranked_ids: list[str],
    rendered: dict[str, dict[str, Any]],
) -> dict[str, list[str]]:
    buckets: dict[str, list[str]] = {
        "method": [],
        "evidence": [],
        "qualitative": [],
        "table": [],
        "fallback": [],
    }
    for layer_id in ranked_ids:
        rec = rendered.get(layer_id) or {}
        if _is_partial_pdf_crop_record(str(layer_id), rec):
            continue
        if _hide_from_planner_visual_catalog(str(layer_id), rec):
            continue
        flags = set(rec.get("curation_flags") or [])
        if "low_information_visual" in flags and not (
            str(rec.get("kind") or "") == "table" or str(layer_id).startswith("ingest_table_")
        ):
            continue
        categories = _paper_visual_bucket_names(layer_id, rec)
        for category in categories:
            if category not in buckets:
                continue
            if len(buckets[category]) < 8 and layer_id not in buckets[category]:
                buckets[category].append(layer_id)
        if len(buckets["fallback"]) < 12 and "low_information_visual" not in flags:
            buckets["fallback"].append(layer_id)
    return buckets


def _paper_visual_bucket_names(layer_id: str, rec: dict[str, Any]) -> list[str]:
    buckets = _infer_visual_roles(layer_id, rec)
    role = str(rec.get("visual_role") or "")
    if role and role not in buckets:
        buckets.insert(0, role)
    return buckets


def _figure_catalog_summary(
    ranked_ids: list[str],
    rendered: dict[str, dict[str, Any]],
) -> dict[str, Any]:
    pages: list[int] = []
    figures = 0
    tables = 0
    role_counts: dict[str, int] = {}
    for layer_id in ranked_ids:
        rec = rendered.get(layer_id) or {}
        if str(rec.get("kind") or "") == "table" or layer_id.startswith("ingest_table_"):
            tables += 1
        else:
            figures += 1
        role = str(rec.get("visual_role") or _primary_visual_role(layer_id, rec))
        role_counts[role] = role_counts.get(role, 0) + 1
        try:
            pages.append(int(rec.get("source_page")))
        except (TypeError, ValueError):
            pass
    recommended = _recommend_paper_visuals(ranked_ids, rendered)
    top_scored_ids = [
        {
            "layer_id": layer_id,
            "visual_role": (rendered.get(layer_id) or {}).get("visual_role"),
            "visual_score": (rendered.get(layer_id) or {}).get("visual_score"),
            "designer_eligible": (rendered.get(layer_id) or {}).get("designer_eligible"),
            "planner_visible": (rendered.get(layer_id) or {}).get("planner_visible"),
            "planner_reject_reasons": list((rendered.get(layer_id) or {}).get("planner_reject_reasons") or []),
        }
        for layer_id in ranked_ids
    ]
    return {
        "total_visuals": len(ranked_ids),
        "n_figures": figures,
        "n_tables": tables,
        "role_counts": role_counts,
        "n_planner_visible": len(ranked_ids),
        "page_min": min(pages) if pages else None,
        "page_max": max(pages) if pages else None,
        "top_ids": ranked_ids,
        "top_scored_ids": top_scored_ids,
        "recommended_counts": {k: len(v) for k, v in recommended.items()},
    }


def _empty_text_unit_buckets() -> dict[str, list[dict[str, Any]]]:
    return {bucket: [] for bucket in _TEXT_UNIT_BUCKETS}


def _merge_text_unit_buckets(
    target: dict[str, list[dict[str, Any]]],
    source: dict[str, Any],
) -> None:
    for bucket in _TEXT_UNIT_BUCKETS:
        for item in list(source.get(bucket) or []):
            if not isinstance(item, dict):
                continue
            _append_text_unit(
                target,
                bucket,
                str(item.get("text") or ""),
                str(item.get("source") or "ingest"),
                claim_id=item.get("claim_id"),
                source_ids=item.get("source_ids"),
                intended_panel_role=item.get("intended_panel_role"),
            )


def _merge_recommended_visual_buckets(
    target: dict[str, list[str]],
    source: dict[str, Any],
) -> None:
    for bucket in ("method", "evidence", "qualitative", "table", "fallback"):
        values = source.get(bucket) or []
        if bucket not in target:
            target[bucket] = []
        for value in values:
            layer_id = str(value or "").strip()
            if layer_id and layer_id not in target[bucket]:
                target[bucket].append(layer_id)


def _filter_visual_candidate_scores_for_planner(
    scores: list[dict[str, Any]],
    rendered: dict[str, dict[str, Any]],
) -> list[dict[str, Any]]:
    out: list[dict[str, Any]] = []
    for item in list(scores or []):
        if not isinstance(item, dict):
            continue
        layer_id = str(item.get("layer_id") or "").strip()
        if not layer_id:
            continue
        rec = rendered.get(layer_id) or {}
        if _is_planner_visible(layer_id, rec, item):
            out.append(item)
    return out


def _filter_recommended_visual_buckets_for_planner(
    recommended: dict[str, list[str]],
    rendered: dict[str, dict[str, Any]],
    visual_candidate_scores: list[dict[str, Any]],
) -> dict[str, list[str]]:
    score_by_id = {
        str(item.get("layer_id") or ""): item
        for item in visual_candidate_scores or []
        if isinstance(item, dict) and str(item.get("layer_id") or "").strip()
    }
    out: dict[str, list[str]] = {}
    for bucket, values in (recommended or {}).items():
        clean: list[str] = []
        for value in values or []:
            layer_id = str(value or "").strip()
            if not layer_id:
                continue
            rec = rendered.get(layer_id) or {}
            if _is_planner_visible(layer_id, rec, score_by_id.get(layer_id) or {}):
                clean.append(layer_id)
        out[str(bucket)] = clean
    return out


def _cached_ingest_summary(
    fp: Path,
    ingested: list[Any],
) -> dict[str, Any] | None:
    """Return an existing ingest summary for the same absolute file path."""
    try:
        target = fp.resolve()
    except OSError:
        target = fp
    for item in ingested:
        if not isinstance(item, dict):
            continue
        raw_file = item.get("file")
        if not raw_file:
            continue
        try:
            candidate = Path(str(raw_file)).expanduser().resolve()
        except OSError:
            candidate = Path(str(raw_file)).expanduser()
        if candidate == target:
            return dict(item)
    return None


def _is_partial_pdf_crop_record(layer_id: str, rec: dict[str, Any] | None) -> bool:
    if not isinstance(rec, dict):
        return _is_partial_pdf_crop_id(layer_id)
    if str(rec.get("extract_strategy") or "").lower() == "sub_panel":
        return True
    if str(rec.get("parent_layer_id") or "").strip():
        return True
    if not _is_partial_pdf_crop_id(layer_id):
        return False
    return not rec or str(rec.get("source") or "ingested_pdf") == "ingested_pdf"


def _is_partial_pdf_crop_id(layer_id: str) -> bool:
    return bool(re.match(r"^ingest_fig_\d{2}_[A-Za-z0-9][A-Za-z0-9_-]*$", str(layer_id or "")))


def _sanitize_pdf_summary_visuals(
    summary: dict[str, Any],
    *,
    rendered: dict[str, dict[str, Any]] | None = None,
    partial_ids: set[str] | None = None,
) -> dict[str, Any]:
    rendered = rendered if isinstance(rendered, dict) else {}
    blocked = set(partial_ids or set())

    def is_blocked(raw: Any) -> bool:
        layer_id = str(raw or "").strip()
        if not layer_id:
            return True
        if layer_id in blocked:
            return True
        rec = rendered.get(layer_id)
        if _is_partial_pdf_crop_record(layer_id, rec):
            blocked.add(layer_id)
            return True
        if _hide_from_planner_visual_catalog(layer_id, rec or {}):
            blocked.add(layer_id)
            return True
        return False

    for key in ("registered_layer_ids", "registered_figure_ids"):
        values = [str(v) for v in list(summary.get(key) or []) if str(v or "").strip()]
        summary[key] = [v for v in values if not is_blocked(v)]
    values = [str(v) for v in list(summary.get("registered_table_ids") or []) if str(v or "").strip()]
    summary["registered_table_ids"] = [v for v in values if not is_blocked(v)]

    scores = []
    for item in list(summary.get("visual_candidate_scores") or []):
        if not isinstance(item, dict):
            continue
        layer_id = str(item.get("layer_id") or "").strip()
        if not is_blocked(layer_id) and _is_planner_visible(layer_id, rendered.get(layer_id) or {}, item):
            scores.append(item)
    if scores or isinstance(summary.get("visual_candidate_scores"), list):
        summary["visual_candidate_scores"] = scores

    recommended = summary.get("recommended_figures")
    if isinstance(recommended, dict):
        clean_recommended: dict[str, list[str]] = {}
        for bucket, values in recommended.items():
            clean_recommended[str(bucket)] = [
                str(v) for v in list(values or [])
                if str(v or "").strip() and not is_blocked(v)
            ]
        summary["recommended_figures"] = clean_recommended

    catalog = summary.get("figure_catalog_summary")
    if isinstance(catalog, dict):
        top_ids = [
            str(v) for v in list(catalog.get("top_ids") or [])
            if str(v or "").strip() and not is_blocked(v)
        ]
        catalog = dict(catalog)
        catalog["top_ids"] = top_ids
        catalog["total_visuals"] = len(top_ids)
        catalog["n_planner_visible"] = len(top_ids)
        catalog["top_scored_ids"] = [
            item for item in list(catalog.get("top_scored_ids") or [])
            if isinstance(item, dict) and not is_blocked(item.get("layer_id"))
        ]
        summary["figure_catalog_summary"] = catalog

    return summary


def _cached_pdf_ingest_summary(fp: Path, ctx: ToolContext) -> dict[str, Any] | None:
    key = _pdf_ingest_cache_key(fp, ctx)
    if not key or not pipeline_cache_enabled("ingest_pdf"):
        return None
    entry_dir = cache_entry_dir(ctx.settings, "ingest_pdf", key)
    payload_path = entry_dir / "payload.json"
    try:
        payload = json.loads(payload_path.read_text(encoding="utf-8"))
    except Exception:
        return None
    if not isinstance(payload, dict):
        return None
    summary = _json_clone(payload.get("summary"))
    rendered_layers = _json_clone(payload.get("rendered_layers"))
    if not isinstance(summary, dict) or not isinstance(rendered_layers, dict):
        return None
    asset_files = payload.get("asset_files") if isinstance(payload.get("asset_files"), dict) else {}
    assets_dir = entry_dir / "assets"
    restored: dict[str, dict[str, Any]] = {}
    partial_ids: set[str] = set()
    quality_doc: fitz.Document | None = None
    try:
        quality_doc = fitz.open(fp)
    except Exception as e:  # noqa: BLE001 - stale cache can still be used without reflagging
        log("ingest.pdf.cache_reflag.open_fail", file=fp.name, error=str(e)[:200])
    try:
        for layer_id, raw_rec in rendered_layers.items():
            if not isinstance(raw_rec, dict):
                continue
            rec = dict(raw_rec)
            layer_key = str(layer_id)
            if _is_partial_pdf_crop_record(layer_key, rec):
                partial_ids.add(layer_key)
                continue
            asset_name = str(asset_files.get(layer_id) or "").strip()
            if asset_name:
                src = assets_dir / asset_name
                if not src.exists():
                    return None
                dest = ctx.layers_dir / Path(asset_name).name
                try:
                    if src.resolve() != dest.resolve():
                        shutil.copy2(src, dest)
                except OSError:
                    return None
                rec["src_path"] = str(dest)
                try:
                    rec["sha256"] = sha256_file(dest)
                except Exception:
                    pass
            rec["source_file"] = str(fp)
            rec["source_pdf"] = fp.name
            try:
                rec["source_pdf_sha256"] = sha256_file(fp)
            except Exception:
                pass
            _refresh_pdf_record_curation_from_current_rules(
                layer_id=layer_key,
                rec=rec,
                pdf_doc=quality_doc,
            )
            restored[layer_key] = rec
    finally:
        if quality_doc is not None:
            quality_doc.close()
    summary = _sanitize_pdf_summary_visuals(
        summary,
        rendered=restored,
        partial_ids=partial_ids,
    )
    if restored:
        ranked_ids = _rank_visual_ids_for_planner(list(restored), restored)
        summary["registered_layer_ids"] = ranked_ids
        summary["registered_figure_ids"] = [layer_id for layer_id in ranked_ids if layer_id.startswith("ingest_fig_")]
        summary["registered_table_ids"] = [layer_id for layer_id in ranked_ids if layer_id.startswith("ingest_table_")]
        visual_scores = _annotate_visual_curation(ranked_ids, restored)
        summary["visual_candidate_scores"] = visual_scores
        summary["recommended_figures"] = _recommend_paper_visuals(ranked_ids, restored)
        summary["figure_catalog_summary"] = _figure_catalog_summary(ranked_ids, restored)
        summary = _sanitize_pdf_summary_visuals(
            summary,
            rendered=restored,
            partial_ids=partial_ids,
        )
    if restored:
        state_rendered = ctx.state.get("rendered_layers")
        if not isinstance(state_rendered, dict):
            state_rendered = {}
            ctx.state["rendered_layers"] = state_rendered
        state_rendered.update(restored)

    contact_asset = str(payload.get("contact_sheet_file") or "").strip()
    if contact_asset:
        src = assets_dir / contact_asset
        if src.exists():
            dest = ctx.layers_dir / Path(contact_asset).name
            try:
                if src.resolve() != dest.resolve():
                    shutil.copy2(src, dest)
                summary["contact_sheet_path"] = str(dest)
            except OSError:
                return None
    summary["file"] = str(fp)
    log(
        "ingest.pdf.cache_hit",
        file=fp.name,
        cache_key=key[:12],
        registered=len(summary.get("registered_layer_ids") or []),
    )
    return summary


def _ensure_pdf_summary_paper_memory(
    fp: Path,
    summary: dict[str, Any],
    ctx: ToolContext,
) -> dict[str, Any]:
    if summary.get("type") != "pdf" or isinstance(summary.get("paper_memory"), dict):
        return summary
    manifest = summary.get("manifest") if isinstance(summary.get("manifest"), dict) else {}
    body_window = (
        summary.get("page_window")
        if isinstance(summary.get("page_window"), dict)
        else manifest.get("_page_window") if isinstance(manifest.get("_page_window"), dict) else {}
    )
    rendered = ctx.state.get("rendered_layers") if isinstance(ctx.state, dict) else {}
    registered_layer_ids = list(summary.get("registered_layer_ids") or [])
    recommended_text_units = (
        summary.get("recommended_text_units")
        if isinstance(summary.get("recommended_text_units"), dict)
        else {}
    )
    memory_key = paper_memory_cache_key(
        pdf_path=fp,
        body_window=body_window,
        manifest=manifest,
        rendered_layers=rendered if isinstance(rendered, dict) else {},
        registered_layer_ids=registered_layer_ids,
        recommended_text_units=recommended_text_units,
    )
    paper_memory = read_paper_memory_cache(ctx.settings, memory_key)
    if paper_memory is None:
        raw_text = str(summary.get("raw_text") or "").strip()
        page_texts = [raw_text] if raw_text else _summary_manifest_page_texts(manifest)
        paper_memory = build_paper_memory(
            pdf_path=fp,
            page_texts=page_texts,
            manifest=manifest,
            body_window=body_window,
            rendered_layers=rendered if isinstance(rendered, dict) else {},
            registered_layer_ids=registered_layer_ids,
            recommended_text_units=recommended_text_units,
        )
        write_paper_memory_cache(ctx.settings, paper_memory)
        log(
            "paper_memory.backfilled",
            file=fp.name,
            cache_key=str(paper_memory.get("cache_key") or "")[:12],
            chunks=paper_memory.get("chunk_count"),
        )
    else:
        paper_memory = dict(paper_memory)
        paper_memory["source_file"] = str(fp)
    updated = dict(summary)
    updated["paper_memory"] = paper_memory
    _write_pdf_ingest_cache(fp, ctx, updated)
    return updated


def _ensure_paper_memory_dossier(
    *,
    ctx: ToolContext,
    paper_memory: dict[str, Any],
    paper_manifest: dict[str, Any],
    paper_visual_provenance: dict[str, Any],
    recommended_text_units: dict[str, Any],
    recommended_figures: dict[str, Any],
) -> dict[str, Any]:
    if not isinstance(paper_memory, dict) or paper_memory.get("kind") != "paper_memory":
        return {}
    model = str(getattr(ctx.settings, "paper_memory_model", "") or "")
    cached = read_paper_memory_dossier_cache(ctx.settings, paper_memory, model=model or None)
    if cached:
        log(
            "paper_memory_dossier.cache_hit",
            cache_key=str(paper_memory.get("cache_key") or "")[:12],
            sections=len(cached.get("sections") or []),
        )
        return cached
    if not getattr(ctx.settings, "enable_paper_memory_agent", True):
        log("paper_memory_agent.skipped", reason="disabled")
        return {}
    try:
        _raise_if_ingest_cancelled(ctx, "ingest.paper_memory_agent.before_build")
        dossier = PaperMemoryAgent(ctx.settings).build(
            memory=paper_memory,
            manifest=paper_manifest,
            visual_provenance=paper_visual_provenance,
            recommended_text_units=recommended_text_units,
            recommended_figures=recommended_figures,
            cancellation_token=_ingest_cancellation_token(ctx),
        )
        _raise_if_ingest_cancelled(ctx, "ingest.paper_memory_agent.after_build")
    except Exception as e:  # noqa: BLE001
        log("paper_memory_agent.failed", error=f"{type(e).__name__}: {e}")
        return {}
    if not dossier:
        return {}
    _raise_if_ingest_cancelled(ctx, "ingest.paper_memory_agent.before_cache_write")
    write_paper_memory_dossier_cache(ctx.settings, paper_memory, dossier)
    _raise_if_ingest_cancelled(ctx, "ingest.paper_memory_agent.after_cache_write")
    _raise_if_ingest_cancelled(ctx, "ingest.paper_memory_agent.before_created_log")
    log(
        "paper_memory_dossier.created",
        cache_key=str(paper_memory.get("cache_key") or "")[:12],
        sections=len(dossier.get("sections") or []),
    )
    _raise_if_ingest_cancelled(ctx, "ingest.paper_memory_agent.after_created_log")
    return dossier


def _summary_manifest_page_texts(manifest: dict[str, Any]) -> list[str]:
    parts = [
        str(manifest.get("title") or ""),
        str(manifest.get("abstract") or ""),
    ]
    for section in manifest.get("sections") or []:
        if not isinstance(section, dict):
            continue
        parts.append(str(section.get("title") or ""))
        parts.append(str(section.get("summary") or section.get("text") or ""))
    text = "\n\n".join(part.strip() for part in parts if part and part.strip())
    return [text] if text else []


def _write_pdf_ingest_cache(fp: Path, ctx: ToolContext, summary: dict[str, Any]) -> None:
    key = _pdf_ingest_cache_key(fp, ctx)
    if not key or not pipeline_cache_enabled("ingest_pdf"):
        return
    entry_dir = cache_entry_dir(ctx.settings, "ingest_pdf", key)
    assets_dir = entry_dir / "assets"
    try:
        assets_dir.mkdir(parents=True, exist_ok=True)
    except OSError:
        return
    rendered = ctx.state.get("rendered_layers") if isinstance(ctx.state, dict) else {}
    rendered_payload: dict[str, Any] = {}
    asset_files: dict[str, str] = {}
    registered_ids = [
        str(layer_id)
        for layer_id in (summary.get("registered_layer_ids") or [])
        if str(layer_id or "").strip()
    ]
    for layer_id in registered_ids:
        rec = rendered.get(layer_id) if isinstance(rendered, dict) else None
        if not isinstance(rec, dict):
            continue
        rec_payload = _json_clone(rec)
        if not isinstance(rec_payload, dict):
            continue
        src_path = Path(str(rec.get("src_path") or ""))
        if src_path.exists():
            asset_name = _cache_asset_name(layer_id, src_path)
            try:
                shutil.copy2(src_path, assets_dir / asset_name)
            except OSError:
                return
            asset_files[layer_id] = asset_name
            rec_payload["src_path"] = asset_name
        rendered_payload[layer_id] = rec_payload

    contact_sheet_file = ""
    contact_sheet = Path(str(summary.get("contact_sheet_path") or ""))
    if contact_sheet.exists():
        contact_sheet_file = _cache_asset_name("contact_sheet", contact_sheet)
        try:
            shutil.copy2(contact_sheet, assets_dir / contact_sheet_file)
        except OSError:
            contact_sheet_file = ""

    payload = {
        "version": 1,
        "cache_key": key,
        "pdf_sha256": _safe_sha256_file(fp),
        "summary": _json_clone(summary),
        "rendered_layers": rendered_payload,
        "asset_files": asset_files,
        "contact_sheet_file": contact_sheet_file,
    }
    try:
        atomic_write_json(entry_dir / "payload.json", payload)
    except Exception:
        return
    log(
        "ingest.pdf.cache_write",
        file=fp.name,
        cache_key=key[:12],
        registered=len(registered_ids),
        assets=len(asset_files),
    )


def _pdf_ingest_cache_key(fp: Path, ctx: ToolContext) -> str:
    try:
        pdf_sha = sha256_file(fp)
    except Exception:
        return ""
    try:
        pages = page_count(fp)
    except Exception:
        pages = None
    return stable_cache_key({
        "stage": "ingest_pdf",
        "version": 18,
        "crop_quality_heuristics": 4,
        "pdf_sha256": pdf_sha,
        "page_count": pages,
        "ingest_model": getattr(ctx.settings, "ingest_model", ""),
        "source_scope": _paper_source_scope(),
        "harness_mode": _poster_harness_mode(ctx),
        "claim_graph_sha256": _claim_graph_cache_fingerprint(ctx.state.get("claim_graph")),
        "canvas_plan_sha256": _json_cache_fingerprint(ctx.state.get("canvas_plan")),
        "prompts": {
            "structure": _hash_text(_INGEST_STRUCTURE_PROMPT),
            "caption": _hash_text(_CAPTION_MATCH_PROMPT),
            "table": _hash_text(_TABLE_PARSE_PROMPT),
            "ocr": _hash_text(_OCR_PROMPT),
        },
        "env": {
            name: os.getenv(name, "")
            for name in (
                "INGEST_VLM_PARALLELISM",
                "INGEST_CAPTION_PARALLELISM",
                "INGEST_VLM_MAX_RETRIES",
                "INGEST_STRUCTURE_COVER_IMAGE",
                "PAPER_SOURCE_SCOPE",
                "POSTER_HARNESS_MODE",
            )
        },
    })


def _claim_graph_cache_fingerprint(value: Any) -> str:
    if value is None:
        return "none"
    return _json_cache_fingerprint(value)


def _json_cache_fingerprint(value: Any) -> str:
    return _hash_text(json.dumps(_json_clone(value), ensure_ascii=False, sort_keys=True, default=str))


def _json_clone(value: Any) -> Any:
    try:
        if hasattr(value, "model_dump"):
            value = value.model_dump(mode="json")
        return json.loads(json.dumps(value, ensure_ascii=False, default=str))
    except Exception:
        return value


def _cache_asset_name(layer_id: str, path: Path) -> str:
    suffix = path.suffix or ".png"
    if path.name.startswith("img_ingest_"):
        return path.name
    stem = re.sub(r"[^A-Za-z0-9_.-]+", "_", str(layer_id or "asset")).strip("._")
    if not stem:
        stem = "asset"
    return f"{stem}{suffix}"


def _hash_text(text: str) -> str:
    import hashlib

    return hashlib.sha256(str(text or "").encode("utf-8")).hexdigest()


def _recommend_paper_text_units(
    manifest: dict[str, Any],
    rendered: dict[str, dict[str, Any]],
    ranked_visual_ids: list[str],
    claim_graph: Any | None,
) -> dict[str, list[dict[str, Any]]]:
    buckets = _empty_text_unit_buckets()

    if claim_graph is not None:
        thesis = str(getattr(claim_graph, "thesis", "") or "")
        _append_text_unit(buckets, "takeaways", thesis, "claim_graph.thesis")
        for node in list(getattr(claim_graph, "tensions", []) or []):
            _append_text_unit(
                buckets,
                "problem",
                f"{getattr(node, 'name', '')}: {getattr(node, 'description', '')}",
                "claim_graph.tensions",
                claim_id=getattr(node, "id", None),
            )
        for node in list(getattr(claim_graph, "mechanisms", []) or []):
            _append_text_unit(
                buckets,
                "method",
                f"{getattr(node, 'name', '')}: {getattr(node, 'description', '')}",
                "claim_graph.mechanisms",
                claim_id=getattr(node, "id", None),
            )
        for node in list(getattr(claim_graph, "evidence", []) or []):
            metric = str(getattr(node, "metric", "") or "")
            quote = str(getattr(node, "raw_quote", "") or "")
            _append_text_unit(
                buckets,
                "evidence",
                f"{metric}: {quote}",
                str(getattr(node, "source", "") or "claim_graph.evidence"),
                claim_id=getattr(node, "id", None),
            )
        for node in list(getattr(claim_graph, "implications", []) or []):
            _append_text_unit(
                buckets,
                "takeaways",
                str(getattr(node, "description", "") or ""),
                "claim_graph.implications",
                claim_id=getattr(node, "id", None),
            )

    abstract = str(manifest.get("abstract") or "")
    for idx, sentence in enumerate(_text_unit_sentences(abstract)[:3]):
        _append_text_unit(
            buckets,
            "takeaways" if idx == 0 else "problem",
            sentence,
            "manifest.abstract",
        )

    for quote in list(manifest.get("key_quotes") or [])[:6]:
        _append_text_unit(buckets, "takeaways", str(quote), "manifest.key_quotes")

    for section in list(manifest.get("sections") or []):
        if not isinstance(section, dict):
            continue
        heading = str(
            section.get("heading")
            or section.get("title")
            or section.get("name")
            or "section"
        )
        source = f"section:{heading[:60]}"
        candidates: list[str] = []
        candidates.extend(str(p) for p in list(section.get("key_points") or []))
        summary = str(section.get("summary") or "")
        candidates.extend(_text_unit_sentences(summary)[:2])
        for text in candidates:
            bucket = _text_unit_bucket_for(f"{heading} {text}", default="takeaways")
            _append_text_unit(buckets, bucket, text, source)

    for item in list(manifest.get("figures") or [])[:16]:
        if not isinstance(item, dict):
            continue
        caption = str(item.get("caption") or item.get("description") or "")
        bucket = _text_unit_bucket_for(caption, default="evidence")
        page = item.get("page")
        source = f"figure:p.{page}" if page else "figure"
        _append_text_unit(buckets, bucket, caption, source)

    for item in list(manifest.get("tables") or [])[:8]:
        if not isinstance(item, dict):
            continue
        caption = str(item.get("caption") or item.get("title") or "")
        page = item.get("page")
        source = f"table:p.{page}" if page else "table"
        _append_text_unit(buckets, "evidence", caption, source)

    for layer_id in ranked_visual_ids[:16]:
        rec = rendered.get(layer_id) or {}
        caption = str(rec.get("caption") or rec.get("title") or "")
        bucket = _text_unit_bucket_for(caption, default="evidence")
        _append_text_unit(buckets, bucket, caption, str(layer_id))

    return buckets


def _append_text_unit(
    buckets: dict[str, list[dict[str, Any]]],
    bucket: str,
    text: str,
    source: str,
    *,
    claim_id: Any | None = None,
    source_ids: Any | None = None,
    intended_panel_role: str | None = None,
) -> None:
    if bucket not in buckets:
        return
    clean = _clip_text_unit(text)
    if len(clean.split()) < 3:
        return
    existing = {str(item.get("text") or "").lower() for item in buckets[bucket]}
    if clean.lower() in existing or len(buckets[bucket]) >= 8:
        return
    normalized_source_ids = _normalize_text_unit_source_ids(source_ids, source)
    item: dict[str, Any] = {
        "bucket": bucket,
        "text": clean,
        "source": source[:120],
        "source_ids": normalized_source_ids,
        "intended_panel_role": intended_panel_role or _text_unit_intended_panel_role(bucket, clean),
    }
    if claim_id:
        item["claim_id"] = str(claim_id)
    buckets[bucket].append(item)


def _normalize_text_unit_source_ids(source_ids: Any, source: str) -> list[str]:
    values: list[str] = []
    if isinstance(source_ids, list):
        values.extend(str(item).strip() for item in source_ids if str(item or "").strip())
    elif source_ids:
        values.append(str(source_ids).strip())
    source_text = str(source or "").strip()
    if source_text.startswith(("ingest_fig_", "ingest_table_", "ingest_img_", "claim_graph.")):
        values.append(source_text)
    elif source_text.startswith(("figure:", "table:", "section:", "manifest.")):
        values.append(source_text)
    out: list[str] = []
    seen: set[str] = set()
    for value in values:
        if value and value not in seen:
            out.append(value[:120])
            seen.add(value)
    return out


def _text_unit_intended_panel_role(bucket: str, text: str) -> str:
    mapped = {
        "problem": "motivation",
        "method": "method_pipeline",
        "evidence": "results_table",
        "limitations": "limitations_future",
        "takeaways": "contribution",
    }
    role = mapped.get(bucket, bucket)
    lower = str(text or "").lower()
    if any(token in lower for token in ("ablation", "analysis", "tradeoff", "failure")):
        return "ablation_analysis"
    if any(token in lower for token in ("pipeline", "framework", "workflow", "encoder", "decoder", "training")):
        return "method_pipeline"
    if any(token in lower for token in ("model", "parameter", "architecture", "backbone")):
        return "model_card"
    if any(token in lower for token in ("table", "benchmark", "leaderboard", "accuracy", "score")):
        return "results_table"
    return role


def _clip_text_unit(text: str, limit: int = 180) -> str:
    clean = " ".join(str(text).replace("\n", " ").split())
    if len(clean) <= limit:
        return clean
    return clean[: limit - 1].rstrip() + "…"


def _text_unit_sentences(text: str) -> list[str]:
    clean = " ".join(str(text).replace("\n", " ").split())
    if not clean:
        return []
    pieces = re.split(r"(?<=[.!?])\s+", clean)
    return [piece.strip() for piece in pieces if len(piece.split()) >= 4]


def _text_unit_bucket_for(text: str, *, default: str) -> str:
    lower = str(text).lower()
    if any(k in lower for k in (
        "limitation", "limitations", "failure", "fails", "future work",
        "caveat", "cannot", "challenge remains",
    )):
        return "limitations"
    if any(k in lower for k in (
        "method", "architecture", "pipeline", "framework", "model",
        "algorithm", "encoder", "decoder", "conditioning", "module",
        "objective", "loss", "training", "inference", "system",
        "diagram", "overview",
    )):
        return "method"
    if any(k in lower for k in (
        "result", "benchmark", "experiment", "evaluation", "ablation",
        "performance", "metric", "comparison", "baseline", "score",
        "accuracy", "qualitative", "table", "figure", "fig.",
    )):
        return "evidence"
    if any(k in lower for k in (
        "problem", "challenge", "gap", "bottleneck", "motivation",
        "difficulty", "under", "conflict", "trade-off", "tradeoff",
    )):
        return "problem"
    if default in _TEXT_UNIT_BUCKETS:
        return default
    return "takeaways"


def _poster_visual_unit_target(
    *,
    target_source_visual_count: int,
    max_visuals: int,
) -> int:
    floor = 4
    target = max(floor, min(max_visuals or 12, target_source_visual_count or floor))
    return max(1, min(target, 8))


def _text_units_by_panel_role(
    recommended_text_units: dict[str, list[dict[str, Any]]],
) -> dict[str, list[dict[str, Any]]]:
    roles: dict[str, list[dict[str, Any]]] = {}
    for bucket, items in (recommended_text_units or {}).items():
        for item in items or []:
            if not isinstance(item, dict):
                continue
            text = str(item.get("text") or "").strip()
            if not text:
                continue
            role = str(item.get("intended_panel_role") or _text_unit_intended_panel_role(str(bucket), text))
            roles.setdefault(role, []).append({
                "text": _clip_text_unit(text, limit=150),
                "source": item.get("source"),
                "source_ids": item.get("source_ids") or [],
            })
    return roles


def _role_text_units_for_native_task(
    roles: dict[str, list[dict[str, Any]]],
    role_names: tuple[str, ...],
    *,
    fallback_roles: tuple[str, ...] = ("contribution", "motivation", "results_table"),
    limit: int = 3,
) -> list[dict[str, Any]]:
    out: list[dict[str, Any]] = []
    for role in role_names + fallback_roles:
        for item in roles.get(role, []) or []:
            text = str(item.get("text") or "").strip()
            if text and all(str(existing.get("text") or "") != text for existing in out):
                out.append(item)
            if len(out) >= limit:
                return out
    return out


def _supplemental_native_visual_tasks(
    *,
    recommended_text_units: dict[str, list[dict[str, Any]]],
    selected_source_visuals: list[str],
    poster_visual_unit_target: int,
) -> list[dict[str, Any]]:
    shortfall = max(0, poster_visual_unit_target - len(selected_source_visuals))
    if shortfall <= 0:
        return []
    roles = _text_units_by_panel_role(recommended_text_units)
    task_specs = [
        (
            "method_flow",
            "Native method flow",
            ("method_pipeline", "model_card"),
            "Draw a compact editable flow/pipeline from paper-grounded method text; use boxes/arrows/labels, not generated imagery.",
        ),
        (
            "benchmark_summary_table",
            "Native result summary table",
            ("results_table", "ablation_analysis"),
            "Build a small booktabs-style summary table or comparison board from verified paper result/evidence text only.",
        ),
        (
            "result_discussion",
            "Short result discussion",
            ("results_table", "contribution"),
            "Write a short result discussion or source-grounded bullets tied to nearby source evidence; use compact comparison table rows only when verified paper facts support them.",
        ),
        (
            "training_recipe",
            "Native model/training method notes",
            ("model_card", "method_pipeline"),
            "Summarize model components, data, training, losses, or setup as source-grounded method notes or compact comparison table rows.",
        ),
        (
            "ablation_or_analysis",
            "Native analysis/limitations panel",
            ("ablation_analysis", "limitations_future"),
            "Use concise paper-backed analysis, ablation, limitation, or future-work facts as a native visual/text hybrid panel.",
        ),
    ]
    tasks: list[dict[str, Any]] = []
    for task_id, title, role_names, instruction in task_specs:
        if len(tasks) >= shortfall:
            break
        source_units = _role_text_units_for_native_task(roles, role_names)
        if not source_units:
            continue
        tasks.append({
            "task_id": task_id,
            "kind": "native_visual_unit",
            "title": title,
            "instruction": instruction,
            "source_text_roles": list(role_names),
            "source_text_units": source_units,
            "requires_source_grounding": True,
            "policy": "native editable HTML/SVG/table only; do not use rejected source crops, remote images, generated scientific imagery, or unsupported numbers",
        })
    return tasks


def _build_poster_content_brief(
    *,
    summaries: list[dict[str, Any]],
    rendered: dict[str, dict[str, Any]],
    recommended_figures: dict[str, list[str]],
    recommended_text_units: dict[str, list[dict[str, Any]]],
    visual_candidate_scores: list[dict[str, Any]],
    canvas_plan: Any,
    paper_visual_provenance: dict[str, Any] | None = None,
    paper_visual_storyboard: dict[str, Any] | None = None,
    raw_brief: str = "",
    required_color_system: dict[str, Any] | None = None,
) -> dict[str, Any]:
    """Build a compact paper-poster planning brief from ingest outputs.

    This is deterministic scaffolding for the planner, not a separate design
    engine. It prevents the common failure where the model copies the abstract
    into the poster and then treats figures as decoration.
    """
    paper_summaries = [s for s in summaries if s.get("type") in {"pdf", "docx", "pptx"}]
    if not paper_summaries:
        return {}

    manifest = paper_summaries[0].get("manifest") or {}
    budget = canvas_plan.get("density_budget") if isinstance(canvas_plan, dict) else {}
    target_min = _safe_int((budget or {}).get("target_visuals_min"), default=6)
    target_max = _safe_int((budget or {}).get("target_visuals_max"), default=10)
    max_visuals = _safe_int((budget or {}).get("max_visuals"), default=12)
    if max_visuals <= 0:
        max_visuals = 12
    target_visual_count = min(
        max(target_min, 8),
        max_visuals,
        max(target_max, target_min, 8),
    )
    visual_candidate_scores = _filter_visual_candidate_scores_for_planner(
        visual_candidate_scores,
        rendered,
    )
    recommended_figures = _filter_recommended_visual_buckets_for_planner(
        recommended_figures,
        rendered,
        visual_candidate_scores,
    )
    ineligible_source_ids = [
        str(layer_id)
        for layer_id, rec in sorted(rendered.items())
        if str(layer_id).startswith(("ingest_fig_", "ingest_table_"))
        and isinstance(rec, dict)
        and not _is_planner_visible(str(layer_id), rec)
    ]

    storyboard = _compact_paper_visual_storyboard_for_planner(paper_visual_storyboard)
    storyboard_selected = [
        str(item.get("asset_id") or "").strip()
        for item in list(storyboard.get("selected_assets") or [])
        if isinstance(item, dict) and str(item.get("asset_id") or "").strip()
    ]
    storyboard_primary = [
        str(item.get("asset_id") or "").strip()
        for item in list(storyboard.get("primary_assets") or [])
        if isinstance(item, dict) and str(item.get("asset_id") or "").strip()
    ]
    storyboard_secondary = [
        str(item.get("asset_id") or "").strip()
        for item in list(storyboard.get("secondary_assets") or [])
        if isinstance(item, dict) and str(item.get("asset_id") or "").strip()
    ]
    storyboard_reserve = [
        str(item.get("asset_id") or "").strip()
        for item in list(storyboard.get("reserve_assets") or [])
        if isinstance(item, dict) and str(item.get("asset_id") or "").strip()
    ]
    storyboard_rejected = [
        str(item.get("asset_id") or "").strip()
        for item in list(storyboard.get("rejected_assets") or [])
        if isinstance(item, dict) and str(item.get("asset_id") or "").strip()
    ]
    storyboard_rejected.extend(
        str(asset_id or "").strip()
        for asset_id in storyboard.get("rejected_asset_ids") or []
        if str(asset_id or "").strip()
    )
    storyboard_rejected = _take_unique_values(storyboard_rejected + ineligible_source_ids)
    if storyboard.get("target_visual_count"):
        target_visual_count = min(
            target_visual_count,
            max(1, _safe_int(storyboard.get("target_visual_count"), default=target_visual_count)),
        )
    selected_visuals = _select_poster_visual_ids(
        recommended_figures=recommended_figures,
        rendered=rendered,
        visual_candidate_scores=visual_candidate_scores,
        target_visual_count=target_visual_count,
        preferred_visual_ids=storyboard_selected,
        forbidden_visual_ids=storyboard_rejected,
    )
    visual_records = [_poster_visual_record(layer_id, rendered.get(layer_id) or {}) for layer_id in selected_visuals]
    dense_targets = _clamp_dense_synthesis_targets(_load_dense_synthesis_targets())
    reference_metadata = _load_reference_metadata_targets()
    reference_profile = _reference_profile_target(reference_metadata)
    if not reference_profile:
        reference_profile = _reference_profile_from_brief(raw_brief)
    if not reference_profile and _canvas_plan_requests_editorial_flow(canvas_plan):
        reference_profile = _EDITORIAL_FLOW_PROFILE
    if reference_metadata:
        dense_targets = _clamp_dense_synthesis_targets(
            _dense_targets_with_reference_metadata(dense_targets, reference_metadata)
        )
    use_dense_synthesis = _should_use_dense_synthesis_profile(
        manifest=manifest,
        rendered=rendered,
        selected_visuals=selected_visuals,
        recommended_text_units=recommended_text_units,
        visual_candidate_scores=visual_candidate_scores,
        dense_targets=dense_targets,
        raw_brief=raw_brief,
    )
    if reference_profile in _DENSE_PAPER_POSTER_PROFILES:
        use_dense_synthesis = True
    if use_dense_synthesis:
        target_visual_count = min(max(target_visual_count, 8), max_visuals)
        selected_visuals = _dense_synthesis_support_visual_ids(
            recommended_figures=recommended_figures,
            selected_visuals=selected_visuals,
            limit=target_visual_count,
            forbidden_visual_ids=storyboard_rejected,
        )
        selected_visuals = _sanitize_selected_visual_ids(
            selected_visuals,
            rendered=rendered,
            visual_candidate_scores=visual_candidate_scores,
            forbidden_visual_ids=storyboard_rejected,
            limit=target_visual_count,
            max_tables=3,
        )
        visual_records = [_poster_visual_record(layer_id, rendered.get(layer_id) or {}) for layer_id in selected_visuals]
    active_dense_reference_profile = (
        (reference_profile or _EDITORIAL_FLOW_PROFILE)
        if use_dense_synthesis else reference_profile
    )
    if active_dense_reference_profile == _EDITORIAL_FLOW_PROFILE:
        editorial_primary_candidates = (
            _take_unique_values(storyboard_primary + selected_visuals)
            if storyboard_primary else selected_visuals
        )
        editorial_primary = [
            layer_id for layer_id in editorial_primary_candidates
            if layer_id in rendered
        ][:6]
        if editorial_primary:
            selected_visuals = _sanitize_selected_visual_ids(
                editorial_primary,
                rendered=rendered,
                visual_candidate_scores=visual_candidate_scores,
                forbidden_visual_ids=storyboard_rejected,
                limit=len(editorial_primary),
                max_tables=3,
            )
            visual_records = [
                _poster_visual_record(layer_id, rendered.get(layer_id) or {})
                for layer_id in selected_visuals
            ]
    provenance = paper_visual_provenance if isinstance(paper_visual_provenance, dict) else {}
    provenance_assets = {
        str(asset.get("asset_id") or ""): asset
        for asset in list(provenance.get("assets") or [])
        if isinstance(asset, dict) and str(asset.get("asset_id") or "").strip()
    }
    selected_visuals = _sanitize_selected_visual_ids(
        selected_visuals,
        rendered=rendered,
        visual_candidate_scores=visual_candidate_scores,
        provenance_assets=provenance_assets,
        forbidden_visual_ids=storyboard_rejected,
        limit=target_visual_count,
        max_tables=3,
    )
    storyboard_primary = _sanitize_selected_visual_ids(
        storyboard_primary,
        rendered=rendered,
        visual_candidate_scores=visual_candidate_scores,
        provenance_assets=provenance_assets,
        forbidden_visual_ids=storyboard_rejected,
        limit=len(storyboard_primary),
        max_tables=3,
    )
    storyboard_secondary = _sanitize_selected_visual_ids(
        storyboard_secondary,
        rendered=rendered,
        visual_candidate_scores=visual_candidate_scores,
        provenance_assets=provenance_assets,
        forbidden_visual_ids=storyboard_rejected,
        limit=len(storyboard_secondary),
        max_tables=3,
        allow_reserve=True,
    )
    storyboard_reserve = _sanitize_selected_visual_ids(
        storyboard_reserve,
        rendered=rendered,
        visual_candidate_scores=visual_candidate_scores,
        provenance_assets=provenance_assets,
        forbidden_visual_ids=storyboard_rejected,
        limit=len(storyboard_reserve),
        max_tables=3,
        allow_reserve=True,
    )
    storyboard_selected = _sanitize_selected_visual_ids(
        storyboard_selected,
        rendered=rendered,
        visual_candidate_scores=visual_candidate_scores,
        provenance_assets=provenance_assets,
        forbidden_visual_ids=storyboard_rejected,
        limit=len(storyboard_selected),
        max_tables=3,
        allow_reserve=True,
    )
    minimum_source_visual_count = _storyboard_minimum_source_visual_count(storyboard)
    constrained_optional_ids = constrain_optional_source_visual_ids(
        [*storyboard_secondary, *storyboard_selected, *storyboard_reserve],
        rendered,
        minimum_count=minimum_source_visual_count,
    )
    allowed_unmatched_reserve_ids = [
        layer_id for layer_id in constrained_optional_ids
        if classify_source_visual(layer_id, rendered.get(layer_id) or {}).get("visual_selection_tier")
        == "reserve_unmatched"
    ]
    storyboard_secondary = _take_unique_values([
        *[
            layer_id for layer_id in storyboard_secondary
            if classify_source_visual(layer_id, rendered.get(layer_id) or {}).get("visual_selection_tier") == "eligible"
        ],
        *allowed_unmatched_reserve_ids,
    ])
    storyboard_reserve = [
        layer_id for layer_id in storyboard_reserve
        if classify_source_visual(layer_id, rendered.get(layer_id) or {}).get("visual_selection_tier") == "eligible"
    ]
    storyboard_selected = _take_unique_values([
        *storyboard_primary,
        *[
            layer_id for layer_id in storyboard_selected
            if classify_source_visual(layer_id, rendered.get(layer_id) or {}).get("visual_selection_tier") == "eligible"
        ],
        *storyboard_secondary,
    ])
    visual_records = [_poster_visual_record(layer_id, rendered.get(layer_id) or {}) for layer_id in selected_visuals]
    high_priority_visual_ids = _high_priority_source_visual_ids(
        provenance_assets=provenance_assets,
        rendered=rendered,
        visual_candidate_scores=visual_candidate_scores,
        preferred_ids=storyboard_primary + storyboard_selected + selected_visuals,
        limit=8,
    )
    high_priority_visual_ids = [
        layer_id
        for layer_id in high_priority_visual_ids
        if layer_id not in storyboard_rejected
    ]
    if high_priority_visual_ids:
        merged_selected = _merge_visual_priority(
            high_priority_visual_ids,
            selected_visuals,
            limit=max(target_visual_count, len(high_priority_visual_ids)),
        )
        selected_visuals = [
            layer_id for layer_id in merged_selected
            if layer_id in rendered and layer_id not in storyboard_rejected
        ]
        selected_visuals = _sanitize_selected_visual_ids(
            _limit_source_table_visuals(selected_visuals, max_tables=3),
            rendered=rendered,
            visual_candidate_scores=visual_candidate_scores,
            provenance_assets=provenance_assets,
            forbidden_visual_ids=storyboard_rejected,
            limit=max(target_visual_count, len(high_priority_visual_ids)),
            max_tables=3,
        )
        target_visual_count = max(target_visual_count, min(max_visuals, len(selected_visuals)))
        visual_records = [
            _poster_visual_record(layer_id, rendered.get(layer_id) or {})
            for layer_id in selected_visuals
        ]

    poster_visual_unit_target = _poster_visual_unit_target(
        target_source_visual_count=len(selected_visuals),
        max_visuals=max_visuals,
    )
    source_visual_shortfall = max(0, poster_visual_unit_target - len(selected_visuals))
    supplemental_native_visual_tasks = _supplemental_native_visual_tasks(
        recommended_text_units=recommended_text_units,
        selected_source_visuals=selected_visuals,
        poster_visual_unit_target=poster_visual_unit_target,
    )

    default_sections = [
        _poster_section(
            "problem",
            "Problem",
            recommended_text_units,
            visual_ids=[],
            purpose="Why the paper exists; use as a compact left/top context block.",
        ),
        _poster_section(
            "method",
            "Method",
            recommended_text_units,
            visual_ids=_take_unique(recommended_figures.get("method") or [], selected_visuals, limit=2),
            purpose="Explain the system or algorithm through figures, not prose.",
        ),
        _poster_section(
            "key_contribution",
            "Key Contribution",
            recommended_text_units,
            buckets=("takeaways", "method"),
            visual_ids=[],
            purpose="Turn the abstract into 2-4 posterized claims.",
        ),
        _poster_section(
            "main_evidence",
            "Main Evidence",
            recommended_text_units,
            buckets=("evidence",),
            visual_ids=_take_unique(
                (recommended_figures.get("evidence") or [])
                + (recommended_figures.get("table") or []),
                selected_visuals,
                limit=4,
            ),
            purpose="Results, benchmark tables, plots, ablations, or comparisons.",
        ),
        _poster_section(
            "takeaway",
            "Takeaway",
            recommended_text_units,
            buckets=("takeaways", "evidence"),
            visual_ids=_take_unique(recommended_figures.get("qualitative") or [], selected_visuals, limit=2),
            purpose="Close the story with the strongest result and what it enables.",
        ),
        _poster_section(
            "limitation_future",
            "Limitations / Future",
            recommended_text_units,
            buckets=("limitations", "takeaways"),
            visual_ids=[],
            purpose="One honest caveat or future direction; never a long paragraph.",
        ),
    ]
    sections = (
        _editorial_flow_sections(
            recommended_text_units,
            recommended_figures=recommended_figures,
            selected_visuals=selected_visuals,
        )
        if use_dense_synthesis and active_dense_reference_profile == _EDITORIAL_FLOW_PROFILE
        else _dense_synthesis_sections(
            recommended_text_units,
            recommended_figures=recommended_figures,
            selected_visuals=selected_visuals,
        )
        if use_dense_synthesis
        else default_sections
    )

    default_panel_plan = [
        {
            "slot_id": "title_meta",
            "purpose": "Paper title, authors, and school/institution/company names only.",
            "text_budget": "three compact lines: title, authors, organization names",
            "visual_ids": [],
            "space_fill": "compact band; do not consume figure space",
        },
        {
            "slot_id": "problem_contribution",
            "purpose": "Problem + key contribution summary.",
            "text_budget": "2-4 bullets total, no abstract paragraph",
            "visual_ids": [],
            "space_fill": "use metrics, callout rules, or a small diagram label if empty",
        },
        {
            "slot_id": "method_visual",
            "purpose": "Hero method/system visual with nearby readout.",
            "text_budget": "1 short section label + local readout",
            "visual_ids": _take_unique(recommended_figures.get("method") or [], selected_visuals, limit=2),
            "space_fill": "large figure/table area; never leave top-half gutters empty",
        },
        {
            "slot_id": "evidence_grid",
            "purpose": "Dense visual evidence grid.",
            "text_budget": "short labels/readouts only",
            "visual_ids": _take_unique(
                (recommended_figures.get("evidence") or [])
                + (recommended_figures.get("qualitative") or [])
                + (recommended_figures.get("fallback") or []),
                selected_visuals,
                limit=6,
            ),
            "space_fill": "fill empty cells with cropped sub-panels, tables, or compact evidence bullets",
        },
        {
            "slot_id": "results_table",
            "purpose": "Benchmark table or quantitative result.",
            "text_budget": "local readout + 1 takeaway",
            "visual_ids": _take_unique(recommended_figures.get("table") or [], selected_visuals, limit=2),
            "space_fill": "prefer one legible table over several unreadable strips",
        },
        {
            "slot_id": "footer_takeaway",
            "purpose": "Takeaway, code/project/contact/citation.",
            "text_budget": "1-2 concise rows",
            "visual_ids": [],
            "space_fill": "use thin rules, QR/contact, or source metadata; no generated ambience",
        },
    ]
    panel_plan = (
        _editorial_flow_panel_plan(
            recommended_figures=recommended_figures,
            selected_visuals=selected_visuals,
        )
        if use_dense_synthesis and active_dense_reference_profile == _EDITORIAL_FLOW_PROFILE
        else _dense_synthesis_panel_plan(
            dense_targets,
            recommended_figures=recommended_figures,
            selected_visuals=selected_visuals,
        )
        if use_dense_synthesis
        else default_panel_plan
    )
    primary_visual_ids = [
        layer_id for layer_id in _merge_visual_priority(
            high_priority_visual_ids,
            (
                storyboard_primary
                or (selected_visuals if use_dense_synthesis else selected_visuals[:min(6, len(selected_visuals))])
            ),
            limit=max(
                len(high_priority_visual_ids),
                len(storyboard_primary)
                if storyboard_primary else len(selected_visuals if use_dense_synthesis else selected_visuals[:min(6, len(selected_visuals))]),
            ),
        )
        if layer_id in selected_visuals
    ]
    secondary_visual_ids = [
        layer_id for layer_id in storyboard_secondary
        if layer_id in rendered and layer_id not in set(primary_visual_ids)
    ]
    required = dict(required_color_system or {})
    recommended_color_system = required or select_academic_color_system(
        raw_brief=raw_brief,
        manifest=manifest,
        recommended_text_units=recommended_text_units,
    )
    color_system_options = _complete_academic_color_system_options(
        raw_brief=raw_brief,
        manifest=manifest,
        recommended_text_units=recommended_text_units,
    )
    color_roles = (
        recommended_color_system.get("roles")
        if isinstance(recommended_color_system.get("roles"), dict)
        else {}
    )
    affiliations = [
        str(item).strip()
        for item in (manifest.get("affiliations") or manifest.get("institutions") or [])
        if str(item or "").strip()
    ][:12]
    institution_color_signals = _institution_color_signals(
        affiliations,
        list(manifest.get("authors") or []),
    )
    allowed_fill = _take_unique_values([
        "#FFFFFF",
        "#FAFDFB",
        "#FBFBF7",
        color_roles.get("background"),
    ])
    aesthetic_contract = _paper_poster_aesthetic_contract()

    payload = {
        "kind": "paper_poster_content_brief",
        "title": manifest.get("title"),
        "authors": list(manifest.get("authors") or [])[:12],
        "affiliations": affiliations,
        "institutions": affiliations,
        "venue": manifest.get("venue"),
        "background_contract": {
            "default": "native white/cream academic canvas",
            "use_generated_background": False,
            "allowed_fill": allowed_fill,
            "structure": "mostly white surfaces, fixed white identity header with a single top accent rule, filled primary section heading bands, thin neutral separators",
        },
        "aesthetic_contract": aesthetic_contract,
        "recommended_color_system": recommended_color_system,
        "color_system": recommended_color_system,
        "color_system_options": color_system_options,
        "institution_color_signals": institution_color_signals,
        "typography_contract": {
            "font_family": '"Times New Roman", Times, Georgia, serif',
            "title_font_size_px": 56,
            "identity_rows_font_size_px": 28,
            "section_heading_font_size_px": 36,
            "body_font_size_px": 24,
            "readout_font_size_px": 24,
            "table_text_font_size_px": 24,
            "caption_label_font_size_px": 20,
            "font_size_tolerance_px": 0.5,
            "times_new_roman_family_ratio_required": 1.0,
            "section_heading_words_max": 8,
            "body_or_readout_words_max": 55,
            "section_bullets": "2-4 short bullets; prefer sourced text units",
        },
        "sections": sections,
        "visual_selection": {
            "target_visual_count": target_visual_count,
            "target_source_visual_count": target_visual_count,
            "selected_source_visual_count": len(selected_visuals),
            "poster_visual_unit_target": poster_visual_unit_target,
            "source_visual_shortfall": source_visual_shortfall,
            "supplemental_native_visual_task_count": len(supplemental_native_visual_tasks),
            "supplemental_native_visual_tasks": supplemental_native_visual_tasks,
            "max_visual_count": (
                len([
                    layer_id for layer_id in (storyboard_primary + storyboard_secondary)
                    if layer_id in rendered
                ])
                if active_dense_reference_profile == _EDITORIAL_FLOW_PROFILE and (storyboard_primary or storyboard_secondary)
                else target_visual_count
            ),
            "primary_visual_ids": primary_visual_ids,
            "high_priority_visual_ids": high_priority_visual_ids,
            "secondary_visual_ids": secondary_visual_ids,
            "reserve_visual_ids": storyboard_reserve,
            "forbidden_visual_ids": storyboard_rejected,
            "visual_records": visual_records,
            "storyboard_selected_asset_ids": storyboard_selected,
            "storyboard_primary_asset_ids": storyboard_primary,
            "storyboard_secondary_asset_ids": storyboard_secondary,
            "storyboard_reserve_asset_ids": storyboard_reserve,
            "storyboard_rejected_asset_ids": storyboard_rejected,
            "source_provenance_manifest": (
                "paper_visual_provenance.json"
                if provenance.get("assets")
                else None
            ),
            "source_asset_records": [
                _poster_source_asset_record(layer_id, provenance_assets.get(layer_id) or {})
                for layer_id in _merge_visual_priority(
                    high_priority_visual_ids,
                    [*selected_visuals, *secondary_visual_ids, *storyboard_reserve],
                    limit=max(
                        len(selected_visuals) + len(secondary_visual_ids) + len(storyboard_reserve),
                        len(high_priority_visual_ids),
                    ),
                )
                if provenance_assets.get(layer_id)
            ],
            "role_buckets": {
                key: list(value or [])[:8]
                for key, value in sorted((recommended_figures or {}).items())
            },
            "selection_rules": [
                "use method + evidence + qualitative/table diversity",
                "skip low-information logos/watermarks/decorative fragments unless user asks",
                "prefer high-score visuals and readable source dimensions",
                "use sub-panels when a composite figure is too dense for a small slot",
                "treat primary_visual_ids as mandatory source assets; secondary_visual_ids are optional capacity fillers; reserve_visual_ids are replacements only",
                "never use forbidden_visual_ids or rejected paper_visual_storyboard assets unless the user explicitly overrides curation",
                *(
                    [
                        "source visual shortfall is active: do not add rejected source crops back into selected_visuals",
                        "satisfy poster_visual_unit_target with supplemental_native_visual_tasks rendered as native HTML/SVG/table units",
                    ]
                    if source_visual_shortfall > 0 else []
                ),
                *(
                    [
                        "conference editorial-flow profile is active: selected source figures/tables are section subjects, not decoration",
                        "use native tables, formulas, concise visual interpretation, method notes, short result discussion, and source-grounded bullets around those source visuals",
                    ]
                    if active_dense_reference_profile == _EDITORIAL_FLOW_PROFILE else
                    [
                        "dense synthesis profile is active: selected source figures/tables are panel subjects, not decoration",
                        "use native model cards, pipelines, result tables, ablation or limitation notes, and takeaway sentences around those source visuals",
                    ]
                    if use_dense_synthesis else []
                ),
            ],
        },
        "panel_plan": panel_plan,
        "editorial_column_plan": (
            panel_plan
            if use_dense_synthesis and active_dense_reference_profile == _EDITORIAL_FLOW_PROFILE
            else None
        ),
        "text_synthesis_targets": _content_brief_text_synthesis_targets(reference_metadata, dense_targets, use_dense_synthesis=use_dense_synthesis),
        "content_fidelity_targets": _content_fidelity_targets(reference_metadata, use_dense_synthesis=use_dense_synthesis),
        "native_reconstruction_hints": _native_reconstruction_hints(
            rendered=rendered,
            recommended_figures=recommended_figures,
            recommended_text_units=recommended_text_units,
            selected_visuals=selected_visuals,
            use_dense_synthesis=use_dense_synthesis,
        ),
        "supplemental_native_visual_tasks": supplemental_native_visual_tasks,
        "visual_storyboard": storyboard or None,
        "space_fill_contract": {
            "empty_panel_repair": (
                "fill with compact comparison tables, method notes, ablation or limitation notes, source-grounded bullets, or takeaway sentences before adding source crops"
                if use_dense_synthesis else
                "fill with sourced figure/table/sub-panel, compact evidence bullets, or thin visual rules"
            ),
            "forbid": [
                "abstract paragraph pasted into a panel",
                "decorative generated background to hide whitespace",
                "orphan captions far from their visual",
                "large blank top-half gutters",
            ],
        },
        "source_asset_policy": {
            "all_images_must_use_declared_assets": bool(provenance.get("assets")),
            "external_images_allowed": False,
            "generated_imagery_allowed": False,
            "selected_source_visual_count": len(selected_visuals),
            "poster_visual_unit_target": poster_visual_unit_target,
            "source_visual_shortfall": source_visual_shortfall,
            "supplemental_native_visual_task_count": len(supplemental_native_visual_tasks),
            "supplemental_native_visual_tasks": supplemental_native_visual_tasks,
            "supplemental_tasks_are_source_assets": False,
            "primary_assets_mandatory": [
                layer_id for layer_id in _merge_visual_priority(
                    high_priority_visual_ids,
                    (storyboard_primary or selected_visuals[:min(6, len(selected_visuals))]),
                    limit=max(
                        len(high_priority_visual_ids),
                        len(storyboard_primary) if storyboard_primary else min(6, len(selected_visuals)),
                    ),
                )
                if layer_id in selected_visuals
            ],
            "secondary_assets_optional": secondary_visual_ids,
            "reserve_assets_replacements_only": storyboard_reserve,
            "forbidden_source_ids": storyboard_rejected,
            "required_manifest": (
                "paper_visual_provenance.json"
                if provenance.get("assets")
                else None
            ),
        },
    }
    if required:
        payload["required_color_system"] = required
    if reference_profile:
        payload["reference_profile"] = reference_profile
    if use_dense_synthesis:
        active_reference_profile = active_dense_reference_profile
        payload["reference_profile"] = active_reference_profile
        payload["native_reference_targets"] = (
            _compact_dense_synthesis_targets(dense_targets)
            if active_reference_profile == _RESEARCH_SYNTHESIS_PROFILE
            else _compact_editorial_flow_targets(dense_targets)
        )
        payload["reference_archetype_skeleton"] = (
            _dense_reference_archetype_skeleton(reference_metadata, dense_targets)
            if active_reference_profile == _RESEARCH_SYNTHESIS_PROFILE
            else _editorial_flow_reference_archetype_skeleton()
        )
        if reference_metadata:
            payload["reference_metadata"] = _compact_reference_metadata_for_brief(reference_metadata)
    return payload


def _load_dense_synthesis_targets() -> dict[str, Any]:
    try:
        payload = json.loads(_DENSE_SYNTHESIS_TARGETS_PATH.read_text(encoding="utf-8"))
    except Exception:
        return {}
    return payload if isinstance(payload, dict) else {}


def _load_reference_metadata_targets() -> dict[str, Any]:
    raw_json = os.getenv("POSTER_REFERENCE_METADATA_JSON", "").strip()
    if raw_json:
        try:
            payload = json.loads(raw_json)
            return payload if isinstance(payload, dict) else {}
        except Exception:
            return {}
    raw_path = os.getenv("POSTER_REFERENCE_METADATA_PATH", "").strip()
    if not raw_path:
        return {}
    path = Path(raw_path).expanduser()
    if not path.exists():
        return {}
    try:
        payload = json.loads(path.read_text(encoding="utf-8"))
    except Exception:
        return {}
    if isinstance(payload, dict):
        payload = dict(payload)
        payload.setdefault("path", str(path))
        return payload
    return {}


def _reference_profile_env() -> str:
    return (
        os.getenv("AUTODESIGN_POSTER_REFERENCE_PROFILE", "")
        or os.getenv("DESIGN_ANYTHING_POSTER_REFERENCE_PROFILE", "")
        or os.getenv("POSTER_REFERENCE_PROFILE", "")
        or os.getenv("POSTER_FORCE_REFERENCE_PROFILE", "")
    ).strip().lower()


def _reference_profile_target(reference_metadata: dict[str, Any]) -> str:
    profile = str(
        reference_metadata.get("reference_profile")
        or reference_metadata.get("profile")
        or ""
    ).strip().lower()
    return profile or _reference_profile_env()


def _reference_profile_from_brief(raw_brief: str) -> str:
    text = str(raw_brief or "").lower()
    if _RESEARCH_SYNTHESIS_PROFILE in text:
        return _RESEARCH_SYNTHESIS_PROFILE
    if _EDITORIAL_FLOW_PROFILE in text or "conference editorial flow" in text or "editorial-flow" in text:
        return _EDITORIAL_FLOW_PROFILE
    return ""


def _canvas_plan_requests_editorial_flow(canvas_plan: Any) -> bool:
    if not isinstance(canvas_plan, dict):
        return False
    body_grid = canvas_plan.get("body_grid") if isinstance(canvas_plan.get("body_grid"), dict) else {}
    values = (
        canvas_plan.get("grid_family"),
        canvas_plan.get("layout_mode"),
        body_grid.get("family"),
        body_grid.get("layout_mode"),
    )
    normalized = {str(value or "").strip().lower().replace("_", "-") for value in values}
    return bool(normalized & {"editorial-3col", "editorial-flow", "conference-editorial-flow"})


def _reference_template_requests_landscape(reference_metadata: dict[str, Any] | None = None) -> bool:
    metadata = reference_metadata if isinstance(reference_metadata, dict) else {}
    env_template = (
        os.getenv("AUTODESIGN_POSTER_TEMPLATE", "")
        or os.getenv("DESIGN_ANYTHING_POSTER_TEMPLATE", "")
        or os.getenv("POSTER_CANVAS_TEMPLATE", "")
        or os.getenv("POSTER_TEMPLATE", "")
    )
    text = " ".join(
        str(value or "")
        for value in (
            env_template,
            metadata.get("preferred_template"),
            metadata.get("layout_archetype"),
        )
    ).lower()
    return any(token in text for token in ("landscape", "wide", "2x1"))


def _normalize_landscape_dense_reference_value(value: Any) -> Any:
    if isinstance(value, str):
        text = value
        text = re.sub(
            r"\bbuild a portrait dense synthesis board\b",
            "build a landscape contract-grid dense synthesis board",
            text,
            flags=re.IGNORECASE,
        )
        text = re.sub(
            r"\buse a portrait dense board\b",
            "use a landscape contract-grid dense board",
            text,
            flags=re.IGNORECASE,
        )
        text = re.sub(
            r"\bacross\s+(?:10|12)\s+(?:numbered\s+)?panels\b",
            (
                "across the contract-defined 3x2, 4x2, or 3x3 main panels, with extra density carried as "
                "internal lanes, native tables, compact comparison tables, concise visual interpretation, and section headings"
            ),
            text,
            flags=re.IGNORECASE,
        )
        text = re.sub(
            r"\b(?:10|12)\s+(?:numbered\s+)?panels\b",
            "contract-defined 3x2, 4x2, or 3x3 main panels plus internal lanes/native units",
            text,
            flags=re.IGNORECASE,
        )
        return text
    if isinstance(value, list):
        return [_normalize_landscape_dense_reference_value(item) for item in value]
    if isinstance(value, dict):
        return {
            key: _normalize_landscape_dense_reference_value(item)
            for key, item in value.items()
        }
    return value


def _dense_targets_with_reference_metadata(
    dense_targets: dict[str, Any],
    reference_metadata: dict[str, Any],
) -> dict[str, Any]:
    if not isinstance(reference_metadata, dict) or not reference_metadata:
        return dense_targets
    out = dict(dense_targets or {})
    landscape_template = _reference_template_requests_landscape(reference_metadata)
    if str(reference_metadata.get("reference_profile") or reference_metadata.get("profile") or ""):
        out["profile"] = str(reference_metadata.get("reference_profile") or reference_metadata.get("profile"))
    hint = reference_metadata.get("reference_metrics_hint") if isinstance(reference_metadata.get("reference_metrics_hint"), dict) else {}
    mapping = {
        "target_panel_count": "target_panel_count",
        "min_panel_count": "min_panel_count",
        "target_native_information_units": "target_native_information_units",
        "min_native_information_units": "min_native_information_units",
        "target_visible_words": "target_visible_words",
        "min_visible_words": "min_visible_words",
    }
    for src, dst in mapping.items():
        if hint.get(src) is not None:
            out[dst] = hint.get(src)
    if landscape_template:
        out["target_panel_count"] = min(9, max(6, _safe_int(out.get("target_panel_count"), default=8)))
        out["min_panel_count"] = 6
    for key in ("required_units", "generation_prior"):
        if isinstance(reference_metadata.get(key), list) and reference_metadata.get(key):
            value = list(reference_metadata.get(key) or [])
            out[key] = (
                _normalize_landscape_dense_reference_value(value)
                if landscape_template else value
            )
    panel_jobs = list(out.get("panel_jobs") or [])
    for unit in reference_metadata.get("required_units") or []:
        value = str(unit or "").strip()
        if value and value not in panel_jobs:
            panel_jobs.append(value)
    if panel_jobs:
        out["panel_jobs"] = panel_jobs
    for key in ("text_synthesis_targets", "acceptance_focus", "negative_guidance"):
        if reference_metadata.get(key):
            value = reference_metadata.get(key)
            out[key] = (
                _normalize_landscape_dense_reference_value(value)
                if landscape_template else value
            )
    return out


def _clamp_dense_synthesis_targets(targets: dict[str, Any]) -> dict[str, Any]:
    if not isinstance(targets, dict):
        return {}
    out = dict(targets)
    def ratio(key: str, *, default: float, floor: float, ceiling: float = 1.0) -> float:
        return min(ceiling, max(floor, _safe_float(out.get(key), default=default)))

    min_native = min(20, max(14, _safe_int(out.get("min_native_information_units"), default=16)))
    target_native = min(30, max(
        22,
        min_native,
        _safe_int(
            out.get("target_native_information_units")
            or out.get("target_native_information_unit_count"),
            default=24,
        ),
    ))
    min_words = min(760, max(520, _safe_int(out.get("min_visible_words"), default=520)))
    target_words = min(950, max(
        700,
        min_words,
        _safe_int(
            out.get("target_visible_words")
            or out.get("target_visible_text_word_count"),
            default=850,
        ),
    ))
    out["min_native_information_units"] = min_native
    out["target_native_information_units"] = target_native
    out["target_native_information_unit_count"] = target_native
    out["min_visible_words"] = min_words
    out["target_visible_words"] = target_words
    out["target_visible_text_word_count"] = target_words
    target_panels = min(9, max(6, _safe_int(out.get("target_panel_count"), default=8)))
    out["target_panel_count"] = target_panels
    out["min_panel_count"] = min(target_panels, max(6, _safe_int(out.get("min_panel_count"), default=6)))
    min_fill = ratio("min_effective_content_fill_ratio", default=0.78, floor=0.72, ceiling=0.92)
    target_fill = ratio(
        "target_effective_content_fill_ratio",
        default=0.86,
        floor=max(0.82, min_fill),
        ceiling=0.96,
    )
    out["min_effective_content_fill_ratio"] = min_fill
    out["target_effective_content_fill_ratio"] = target_fill
    out["max_low_information_blank_ratio"] = ratio(
        "max_low_information_blank_ratio",
        default=0.14,
        floor=0.04,
        ceiling=0.22,
    )
    out["min_panel_fill_ratio"] = ratio("min_panel_fill_ratio", default=0.70, floor=0.64, ceiling=0.90)
    out["target_panel_fill_ratio"] = ratio("target_panel_fill_ratio", default=0.82, floor=0.74, ceiling=0.94)
    out["min_text_fragment_count"] = max(14, _safe_int(out.get("min_text_fragment_count"), default=14))
    out["target_text_fragment_count"] = min(28, max(22, _safe_int(out.get("target_text_fragment_count"), default=24)))
    out["min_mixed_text_evidence_panel_count"] = max(
        4,
        _safe_int(out.get("min_mixed_text_evidence_panel_count"), default=4),
    )
    out["min_readable_source_visual_area_ratio"] = ratio(
        "min_readable_source_visual_area_ratio",
        default=0.24,
        floor=0.16,
        ceiling=0.34,
    )
    out["min_readable_source_visual_unit_ratio"] = ratio(
        "min_readable_source_visual_unit_ratio",
        default=0.035,
        floor=0.025,
        ceiling=0.08,
    )
    return out


def _content_brief_text_synthesis_targets(
    reference_metadata: dict[str, Any],
    dense_targets: dict[str, Any],
    *,
    use_dense_synthesis: bool,
) -> dict[str, Any]:
    raw = reference_metadata.get("text_synthesis_targets") if isinstance(reference_metadata, dict) else None
    if isinstance(raw, dict) and raw:
        out = dict(raw)
        if _reference_template_requests_landscape(reference_metadata):
            out = _normalize_landscape_dense_reference_value(out)
        return out
    if not use_dense_synthesis:
        return {}
    return {
        "priority": "high",
        "density_goal": "Dense but edited poster copy: short panel theses, compact bullets, labels, and highlighted takeaways.",
        "correctness_goal": "Use source-backed paper claims and numbers; do not invent benchmark values or model facts.",
        "coherence_goal": "Connect panels from problem to method, evidence, analysis, limitations, and provenance.",
        "condensation_goal": "Compress the paper into native synthesis units instead of pasting abstract prose.",
        "preferred_text_units": [
            "panel thesis",
            "compact evidence bullet",
            "table/chart takeaway",
            "method-step label",
            "limitation/future bullet",
        ],
        "avoid_text_failures": [
            "abstract dumping",
            "unsupported generic claims",
            "benchmark hallucination",
            "disconnected panels",
        ],
        "source": dense_targets.get("source_set") or "dense_synthesis_targets",
    }


def _content_fidelity_targets(
    reference_metadata: dict[str, Any],
    *,
    use_dense_synthesis: bool,
) -> dict[str, Any]:
    if not use_dense_synthesis:
        return {}
    focus = reference_metadata.get("acceptance_focus") if isinstance(reference_metadata.get("acceptance_focus"), list) else []
    negative = reference_metadata.get("negative_guidance") if isinstance(reference_metadata.get("negative_guidance"), list) else []
    return {
        "priority": "high",
        "source_backed_text_required": True,
        "source_backed_core_panels_min": 5,
        "claim_correctness_policy": "Core claims, model facts, benchmark numbers, and limitations must be grounded in paper sections, claim_graph, figures, tables, or source metadata.",
        "abstract_dump_policy": "Do not paste the abstract as a panel; split it into sourced problem, contribution, method, evidence, and limitation units.",
        "acceptance_focus": [str(item) for item in focus if str(item or "").strip()][:10],
        "negative_guidance": [str(item) for item in negative if str(item or "").strip()][:10],
    }


def _native_reconstruction_hints(
    *,
    rendered: dict[str, dict[str, Any]],
    recommended_figures: dict[str, list[str]],
    recommended_text_units: dict[str, list[dict[str, Any]]],
    selected_visuals: list[str],
    use_dense_synthesis: bool,
) -> dict[str, Any]:
    direct_visual_limit = (
        min(10, max(6, len(selected_visuals)))
        if use_dense_synthesis else 8
    )
    direct_visuals = _take_unique(
        (recommended_figures.get("method") or [])
        + (recommended_figures.get("table") or [])
        + (recommended_figures.get("evidence") or [])
        + (recommended_figures.get("qualitative") or [])
        + (recommended_figures.get("fallback") or []),
        selected_visuals,
        limit=direct_visual_limit,
    )
    native_tables = _take_unique(
        (recommended_figures.get("table") or [])
        + [layer_id for layer_id in selected_visuals if str(layer_id).startswith("ingest_table_")],
        selected_visuals,
        limit=4,
    )
    text_roles: dict[str, list[dict[str, Any]]] = {}
    for bucket, items in recommended_text_units.items():
        for item in items or []:
            if not isinstance(item, dict):
                continue
            role = str(item.get("intended_panel_role") or _text_unit_intended_panel_role(bucket, str(item.get("text") or "")))
            text_roles.setdefault(role, []).append({
                "text": item.get("text"),
                "source": item.get("source"),
                "source_ids": item.get("source_ids") or [],
            })
    return {
        "direct_source_visual": [
            _poster_visual_record(layer_id, rendered.get(layer_id) or {})
            for layer_id in direct_visuals
        ],
        "native_table": native_tables,
        "native_pipeline": text_roles.get("method_pipeline", [])[:6],
        "native_model_card": text_roles.get("model_card", [])[:6],
        "native_chart_or_stat_card": text_roles.get("results_table", [])[:8],
        "text_synthesis_panel": {
            key: value[:6]
            for key, value in sorted(text_roles.items())
            if key not in {"method_pipeline", "model_card", "results_table"}
        },
        "policy": (
            "Use a source-asset range, typically 5-8 selected figures/tables with about 6 as a starting point; choose the final count from the available column height and carry text as compact local readouts, labels, and takeaways around native editable tables, cards, charts, and pipelines."
            if use_dense_synthesis
            else "Use source visuals/tables as primary evidence and fill gaps with compact sourced text units."
        ),
    }


def _compact_reference_metadata_for_brief(metadata: dict[str, Any]) -> dict[str, Any]:
    keys = (
        "case_id",
        "source_set",
        "paper_title",
        "profile",
        "reference_profile",
        "quality_tier",
        "preferred_template",
        "layout_archetype",
        "text_synthesis_targets",
        "acceptance_focus",
        "negative_guidance",
    )
    out = {key: metadata.get(key) for key in keys if metadata.get(key) is not None}
    if _reference_template_requests_landscape(metadata):
        out = _normalize_landscape_dense_reference_value(out)
    return out


def _should_use_dense_synthesis_profile(
    *,
    manifest: dict[str, Any],
    rendered: dict[str, dict[str, Any]],
    selected_visuals: list[str],
    recommended_text_units: dict[str, list[dict[str, Any]]],
    visual_candidate_scores: list[dict[str, Any]],
    dense_targets: dict[str, Any],
    raw_brief: str = "",
) -> bool:
    explicit_profile = _reference_profile_env()
    if explicit_profile and explicit_profile not in _DENSE_PAPER_POSTER_PROFILES:
        return False
    if dense_targets.get("profile") != _RESEARCH_SYNTHESIS_PROFILE:
        return False
    if _force_dense_synthesis_profile(raw_brief):
        return True
    source_visual_count = len([
        layer_id for layer_id, rec in rendered.items()
        if str(layer_id).startswith(("ingest_fig_", "ingest_table_"))
        and isinstance(rec, dict)
        and "logo" not in " ".join(str(rec.get(key) or "") for key in ("name", "caption", "title")).lower()
    ])

    haystack = _dense_synthesis_signal_text(
        manifest=manifest,
        rendered=rendered,
        recommended_text_units=recommended_text_units,
        visual_candidate_scores=visual_candidate_scores,
    )
    if _word_count_for_dense_signal(haystack) < 120:
        return False
    signal_families = {
        "model_card": ("model", "parameter", "architecture", "modality", "backbone", "tokenizer"),
        "method_pipeline": ("pipeline", "framework", "stage", "encoder", "decoder", "training", "inference", "system"),
        "benchmark_table": ("benchmark", "leaderboard", "table", "accuracy", "score", "baseline", "ablation", "result"),
        "analysis": ("analysis", "ablation", "limitation", "future", "tradeoff", "failure"),
    }
    family_hits = sum(
        1 for tokens in signal_families.values()
        if any(token in haystack for token in tokens)
    )
    structure_degraded = bool(manifest.get("_structure_degraded"))
    strong_selected_visuals = _strong_selected_visual_count(selected_visuals, visual_candidate_scores)
    if 0 < source_visual_count <= 3 and family_hits >= 2:
        return True
    if source_visual_count >= 4 and family_hits >= 2:
        return True
    if source_visual_count == 0 and selected_visuals and len(selected_visuals) <= 3 and family_hits >= 2:
        return True
    if structure_degraded and family_hits >= 2:
        return True
    if family_hits >= 3 and strong_selected_visuals <= 8:
        return True
    return family_hits >= 2 and source_visual_count <= 6


def _force_dense_synthesis_profile(raw_brief: str) -> bool:
    env_profile = _reference_profile_env()
    if env_profile:
        return env_profile in _DENSE_PAPER_POSTER_PROFILES
    text = str(raw_brief or "").lower()
    return (
        _RESEARCH_SYNTHESIS_PROFILE in text
        or _EDITORIAL_FLOW_PROFILE in text
        or "conference editorial flow" in text
        or "editorial-flow" in text
        or "native html dense-synthesis reference prior" in text
        or "dense-synthesis reference prior" in text
    )


def _strong_selected_visual_count(
    selected_visuals: list[str],
    visual_candidate_scores: list[dict[str, Any]],
) -> int:
    scores = {
        str(item.get("layer_id") or ""): item
        for item in visual_candidate_scores or []
        if isinstance(item, dict) and str(item.get("layer_id") or "").strip()
    }
    count = 0
    for layer_id in selected_visuals:
        rec = scores.get(str(layer_id) or "")
        if not rec:
            continue
        score = _safe_float(rec.get("visual_score"), default=0.0)
        role = str(rec.get("visual_role") or "").lower()
        if score >= 72 or role in {"method", "evidence", "table", "qualitative"}:
            count += 1
    return count


def _dense_synthesis_support_visual_ids(
    *,
    recommended_figures: dict[str, list[str]],
    selected_visuals: list[str],
    limit: int,
    forbidden_visual_ids: list[str] | None = None,
) -> list[str]:
    forbidden = {str(item or "").strip() for item in (forbidden_visual_ids or []) if str(item or "").strip()}
    ordered: list[str] = []
    _extend_unique(
        ordered,
        [item for item in selected_visuals if str(item or "").strip() not in forbidden],
        limit=max(1, limit),
    )
    for bucket in ("method", "table", "evidence", "qualitative", "fallback"):
        _extend_unique(
            ordered,
            [item for item in (recommended_figures.get(bucket) or []) if str(item or "").strip() not in forbidden],
            limit=max(1, limit),
        )
    return ordered[:max(1, limit)]


def _dense_synthesis_signal_text(
    *,
    manifest: dict[str, Any],
    rendered: dict[str, dict[str, Any]],
    recommended_text_units: dict[str, list[dict[str, Any]]],
    visual_candidate_scores: list[dict[str, Any]],
) -> str:
    parts: list[str] = [
        str(manifest.get("title") or ""),
        str(manifest.get("abstract") or ""),
        str(manifest.get("summary") or ""),
    ]
    for section in manifest.get("sections") or []:
        if not isinstance(section, dict):
            continue
        parts.extend(str(section.get(key) or "") for key in ("title", "heading", "summary", "text"))
        parts.extend(str(item or "") for item in (section.get("key_points") or []))
    for bucket in recommended_text_units.values():
        for item in bucket or []:
            if isinstance(item, dict):
                parts.append(str(item.get("text") or ""))
    for rec in rendered.values():
        if isinstance(rec, dict):
            parts.extend(str(rec.get(key) or "") for key in ("caption", "title", "caption_short", "visual_role"))
    for item in visual_candidate_scores or []:
        if isinstance(item, dict):
            parts.extend(str(item.get(key) or "") for key in ("visual_role", "reason", "caption", "title"))
    return re.sub(r"\s+", " ", " ".join(parts)).lower()


def _word_count_for_dense_signal(text: str) -> int:
    return len(re.findall(r"[^\W_]+", text, flags=re.UNICODE))


def _dense_synthesis_sections(
    recommended_text_units: dict[str, list[dict[str, Any]]],
    *,
    recommended_figures: dict[str, list[str]],
    selected_visuals: list[str],
) -> list[dict[str, Any]]:
    def visuals(*buckets: str, limit: int) -> list[str]:
        values: list[str] = []
        for bucket in buckets:
            values.extend(recommended_figures.get(bucket) or [])
        return _take_unique(values, selected_visuals, limit=limit)

    return [
        _poster_section(
            "identity_header",
            "Identity Header",
            recommended_text_units,
            buckets=("takeaways",),
            visual_ids=[],
            purpose=(
                "Identity header for exactly three text rows only: paper title, authors, "
                "and school/institution/company names; keep this as a compact title band rather than "
                "a separate content panel."
            ),
            bullet_limit=2,
        ),
        _poster_section(
            "problem_contribution",
            "Problem and Contributions",
            recommended_text_units,
            buckets=("problem", "takeaways"),
            visual_ids=[],
            purpose="One compressed context panel: why the problem matters, what bottleneck the paper targets, and the 2-4 claims that make the paper worth reading.",
            bullet_limit=6,
        ),
        _poster_section(
            "model_card",
            "Model Card",
            recommended_text_units,
            buckets=("method", "takeaways"),
            visual_ids=visuals("method", limit=1),
            purpose="Native card covering model, architecture, modalities, data, and capability fields.",
            bullet_limit=5,
        ),
        _poster_section(
            "method_pipeline",
            "Method Pipeline",
            recommended_text_units,
            buckets=("method",),
            visual_ids=visuals("method", "fallback", limit=2),
            purpose="Editable stages for training/inference/system pipeline.",
            bullet_limit=5,
        ),
        _poster_section(
            "results_table",
            "Benchmark Table",
            recommended_text_units,
            buckets=("evidence",),
            visual_ids=visuals("table", "evidence", limit=2),
            purpose="Native result table, leaderboard, compact comparison table, or short result discussion.",
            bullet_limit=5,
        ),
        _poster_section(
            "ablation_analysis",
            "Ablation / Analysis",
            recommended_text_units,
            buckets=("evidence", "takeaways"),
            visual_ids=visuals("evidence", "table", limit=2),
            purpose="Analysis of deltas, ablations, tradeoffs, or failure modes.",
            bullet_limit=5,
        ),
        _poster_section(
            "limitations_future",
            "Limitations / Future",
            recommended_text_units,
            buckets=("limitations", "takeaways"),
            visual_ids=[],
            purpose="Required caveats, failure modes, and future directions.",
            bullet_limit=5,
        ),
        _poster_section(
            "synthesis_takeaway",
            "Understanding Generation Unification",
            recommended_text_units,
            buckets=("takeaways", "evidence"),
            visual_ids=[],
            purpose="Understanding generation unification: close the research story with implications, strongest evidence, and what the result enables.",
            bullet_limit=5,
        ),
    ]


def _editorial_flow_sections(
    recommended_text_units: dict[str, list[dict[str, Any]]],
    *,
    recommended_figures: dict[str, list[str]],
    selected_visuals: list[str],
) -> list[dict[str, Any]]:
    def visuals(*buckets: str, limit: int) -> list[str]:
        values: list[str] = []
        for bucket in buckets:
            values.extend(recommended_figures.get(bucket) or [])
        return _take_unique(values, selected_visuals, limit=limit)

    return [
        _poster_section(
            "motivation",
            "The Motivation",
            recommended_text_units,
            buckets=("problem", "takeaways"),
            visual_ids=visuals("fallback", "qualitative", limit=1),
            purpose="Left-column setup: problem, bottleneck, and why the task matters. Keep prose short.",
            bullet_limit=3,
        ),
        _poster_section(
            "comparison",
            "Existing Methods",
            recommended_text_units,
            buckets=("evidence", "problem"),
            visual_ids=visuals("table", "evidence", limit=1),
            purpose="A compact comparison table, source-grounded bullet group, or visual contrast; source table/figure first.",
            bullet_limit=3,
        ),
        _poster_section(
            "method_overview",
            "The Method",
            recommended_text_units,
            buckets=("method",),
            visual_ids=visuals("method", "fallback", limit=2),
            purpose="Middle-column method story with one large source method figure and a short readout.",
            bullet_limit=4,
        ),
        _poster_section(
            "method_details",
            "How It Works",
            recommended_text_units,
            buckets=("method", "takeaways"),
            visual_ids=visuals("method", "qualitative", limit=1),
            purpose="Equations, pipeline stages, or native table/formula support around the main method visual.",
            bullet_limit=3,
        ),
        _poster_section(
            "results",
            "The Results",
            recommended_text_units,
            buckets=("evidence",),
            visual_ids=visuals("table", "evidence", limit=2),
            purpose="Right-column benchmark/result evidence. Prefer one readable source table/chart plus compact comparison table rows.",
            bullet_limit=4,
        ),
        _poster_section(
            "analysis",
            "Analysis",
            recommended_text_units,
            buckets=("evidence", "limitations"),
            visual_ids=visuals("evidence", "qualitative", limit=2),
            purpose="Ablations, qualitative evidence, limitations, or failure modes with local visual readouts.",
            bullet_limit=4,
        ),
        _poster_section(
            "takeaway",
            "Takeaway",
            recommended_text_units,
            buckets=("takeaways", "limitations"),
            visual_ids=[],
            purpose="One concise bottom readout: strongest result, caveat, and practical implication.",
            bullet_limit=3,
        ),
    ]


def _compact_dense_synthesis_targets(targets: dict[str, Any]) -> dict[str, Any]:
    keys = (
        "profile",
        "version",
        "source_set",
        "gold_regression_references",
        "reference_count",
        "target_panel_count",
        "min_panel_count",
        "target_native_information_units",
        "target_native_information_unit_count",
        "min_native_information_units",
        "target_visible_words",
        "target_visible_text_word_count",
        "min_visible_words",
        "min_effective_content_fill_ratio",
        "target_effective_content_fill_ratio",
        "max_low_information_blank_ratio",
        "min_panel_fill_ratio",
        "target_panel_fill_ratio",
        "min_text_fragment_count",
        "target_text_fragment_count",
        "min_mixed_text_evidence_panel_count",
        "min_readable_source_visual_area_ratio",
        "min_readable_source_visual_unit_ratio",
        "slot_fill_targets",
        "slot_native_information_unit_targets",
        "required_units",
        "panel_jobs",
        "generation_prior",
        "text_synthesis_targets",
        "acceptance_focus",
        "negative_guidance",
    )
    return {key: targets.get(key) for key in keys if targets.get(key) is not None}


def _compact_editorial_flow_targets(targets: dict[str, Any]) -> dict[str, Any]:
    return {
        "profile": _EDITORIAL_FLOW_PROFILE,
        "column_count": 3,
        "min_sections_total": 7,
        "target_sections_total": 9,
        "target_source_visual_count": min(6, max(5, _safe_int(targets.get("target_visual_count"), default=6))),
        "max_source_visual_count": 8,
        "min_native_information_units": 10,
        "target_native_information_units": 16,
        "min_visible_words": 360,
        "target_visible_words": 560,
        "source_flow_unit": "one source figure/table plus one local readout as direct siblings",
        "readout_target_words": "20-90",
        "forbid": [
            "visible figcaption",
            "local readouts starting with Fig./Figure/Table numbers",
            "panel_content_plan",
            "poster-grid six-card body",
            "child lane DOM",
        ],
    }


def _dense_reference_archetype_skeleton(
    reference_metadata: dict[str, Any],
    targets: dict[str, Any],
) -> dict[str, Any]:
    metadata = reference_metadata if isinstance(reference_metadata, dict) else {}
    style = metadata.get("visual_style") if isinstance(metadata.get("visual_style"), dict) else {}
    content_fill_contract = {
        "min_effective_content_fill_ratio": _safe_float(
            targets.get("min_effective_content_fill_ratio"),
            default=0.78,
        ),
        "target_effective_content_fill_ratio": _safe_float(
            targets.get("target_effective_content_fill_ratio"),
            default=0.86,
        ),
        "max_low_information_blank_ratio": _safe_float(
            targets.get("max_low_information_blank_ratio"),
            default=0.14,
        ),
        "min_panel_fill_ratio": _safe_float(targets.get("min_panel_fill_ratio"), default=0.70),
        "target_panel_fill_ratio": _safe_float(targets.get("target_panel_fill_ratio"), default=0.82),
        "min_text_fragment_count": _safe_int(targets.get("min_text_fragment_count"), default=14),
        "target_text_fragment_count": _safe_int(targets.get("target_text_fragment_count"), default=22),
        "min_mixed_text_evidence_panel_count": _safe_int(
            targets.get("min_mixed_text_evidence_panel_count"),
            default=4,
        ),
        "min_readable_source_visual_area_ratio": _safe_float(
            targets.get("min_readable_source_visual_area_ratio"),
            default=0.12,
        ),
        "min_readable_source_visual_unit_ratio": _safe_float(
            targets.get("min_readable_source_visual_unit_ratio"),
            default=0.035,
        ),
    }
    slot_fill_targets = (
        targets.get("slot_fill_targets")
        if isinstance(targets.get("slot_fill_targets"), dict)
        else {}
    )
    slot_native_targets = (
        targets.get("slot_native_information_unit_targets")
        if isinstance(targets.get("slot_native_information_unit_targets"), dict)
        else {}
    )
    gold_refs = (
        targets.get("gold_regression_references")
        if isinstance(targets.get("gold_regression_references"), dict)
        else {}
    )
    case_slug = str(metadata.get("case_id") or metadata.get("slug") or "").strip()
    gold_reference = (
        gold_refs.get(case_slug)
        if case_slug and isinstance(gold_refs.get(case_slug), dict)
        else {}
    )
    gold_density_floor = (
        gold_reference.get("visual_density_floor")
        if isinstance(gold_reference, dict) and isinstance(gold_reference.get("visual_density_floor"), dict)
        else {}
    )
    if gold_density_floor:
        content_fill_contract["gold_min_nonwhite_pixel_ratio"] = _safe_float(
            gold_density_floor.get("min_nonwhite_pixel_ratio"),
            default=0.0,
        )
        content_fill_contract["gold_max_longest_blank_vertical_run_ratio"] = _safe_float(
            gold_density_floor.get("max_longest_blank_vertical_run_ratio"),
            default=0.0,
        )
        content_fill_contract["gold_min_leaf_visible_words"] = _safe_int(
            gold_density_floor.get("min_leaf_visible_words"),
            default=0,
        )
        content_fill_contract["gold_min_source_figure_area_ratio"] = _safe_float(
            gold_density_floor.get("min_source_figure_area_ratio"),
            default=0.0,
        )
        content_fill_contract["gold_min_native_information_units"] = _safe_int(
            gold_density_floor.get("min_native_information_units"),
            default=0,
        )
    slot_geometry_contract = {
        "portrait": [
            {"slot_id": "title_meta", "band": "top", "canvas_area_ratio_max": 0.09, "content_fill_ratio_min": 0.65, "native_units_min": 1, "text_fragments_min": 2},
            {"slot_id": "method_pipeline", "band": "upper_middle", "canvas_area_ratio_min": 0.12, "content_fill_ratio_min": max(0.82, _safe_float(slot_fill_targets.get("method_pipeline"), default=0.82)), "native_units_min": _safe_int(slot_native_targets.get("method_pipeline"), default=4), "text_fragments_min": 3, "role": "hero_content"},
            {"slot_id": "results_table", "band": "middle", "canvas_area_ratio_min": 0.10, "content_fill_ratio_min": max(0.82, _safe_float(slot_fill_targets.get("results_table"), default=0.82)), "native_units_min": _safe_int(slot_native_targets.get("results_table"), default=4), "text_fragments_min": 3, "role": "hero_content"},
            {"slot_id": "benchmark_table", "band": "middle_lower", "canvas_area_ratio_min": 0.08, "content_fill_ratio_min": max(0.82, _safe_float(slot_fill_targets.get("benchmark_table"), default=0.82)), "native_units_min": _safe_int(slot_native_targets.get("benchmark_table"), default=4), "text_fragments_min": 3},
            {"slot_id": "ablation_analysis", "band": "lower_middle", "canvas_area_ratio_min": 0.06, "content_fill_ratio_min": max(0.78, _safe_float(slot_fill_targets.get("ablation_analysis"), default=0.78)), "native_units_min": _safe_int(slot_native_targets.get("ablation_analysis"), default=3), "text_fragments_min": 3},
            {"slot_id": "limitations_future", "band": "bottom", "canvas_area_ratio_min": 0.05, "content_fill_ratio_min": max(0.72, _safe_float(slot_fill_targets.get("limitations_future"), default=0.72)), "native_units_min": _safe_int(slot_native_targets.get("limitations_future"), default=2), "text_fragments_min": 3},
        ],
        "wide": [
            {"slot_id": "title_meta", "band": "top", "canvas_area_ratio_max": 0.08, "content_fill_ratio_min": 0.65, "native_units_min": 1, "text_fragments_min": 2},
            {"slot_id": "method_pipeline", "band": "left_or_center_hero", "canvas_area_ratio_min": 0.14, "content_fill_ratio_min": max(0.82, _safe_float(slot_fill_targets.get("method_pipeline"), default=0.82)), "native_units_min": _safe_int(slot_native_targets.get("method_pipeline"), default=4), "text_fragments_min": 3, "role": "hero_content"},
            {"slot_id": "results_table", "band": "right_or_center_hero", "canvas_area_ratio_min": 0.12, "content_fill_ratio_min": max(0.82, _safe_float(slot_fill_targets.get("results_table"), default=0.82)), "native_units_min": _safe_int(slot_native_targets.get("results_table"), default=4), "text_fragments_min": 3, "role": "hero_content"},
            {"slot_id": "qualitative_evidence", "band": "lower_strip", "canvas_area_ratio_min": 0.08, "content_fill_ratio_min": max(0.76, _safe_float(slot_fill_targets.get("qualitative_evidence"), default=0.76)), "native_units_min": 2, "text_fragments_min": 3},
            {"slot_id": "ablation_analysis", "band": "lower_grid", "canvas_area_ratio_min": 0.06, "content_fill_ratio_min": max(0.78, _safe_float(slot_fill_targets.get("ablation_analysis"), default=0.78)), "native_units_min": _safe_int(slot_native_targets.get("ablation_analysis"), default=3), "text_fragments_min": 3},
            {"slot_id": "synthesis_takeaway", "band": "footer", "canvas_area_ratio_max": 0.06, "content_fill_ratio_min": 0.68, "native_units_min": _safe_int(slot_native_targets.get("synthesis_takeaway"), default=2), "text_fragments_min": 2},
        ],
    }
    return {
        "profile": "research_synthesis_dense",
        "reference_archetype": str(
            metadata.get("layout_archetype")
            or "dense_native_synthesis_board"
        ),
        "hard_constraints": {
            "target_panel_count": _safe_int(targets.get("target_panel_count"), default=8),
            "min_panel_count": _safe_int(targets.get("min_panel_count"), default=6),
            "target_native_information_units": _safe_int(targets.get("target_native_information_units"), default=24),
            "min_native_information_units": _safe_int(targets.get("min_native_information_units"), default=16),
            "target_visible_words": _safe_int(targets.get("target_visible_words"), default=850),
            "min_visible_words": _safe_int(targets.get("min_visible_words"), default=520),
            "target_table_count": _safe_int(
                (metadata.get("reference_metrics_hint") or {}).get("target_table_count")
                if isinstance(metadata.get("reference_metrics_hint"), dict) else None,
                default=2,
            ),
        },
        "content_fill_contract": content_fill_contract,
        "gold_regression_reference": {
            "id": gold_reference.get("id"),
            "source_iteration": gold_reference.get("source_iteration"),
            "visual_density_floor": dict(gold_density_floor),
            "dense_skeleton": gold_reference.get("dense_skeleton"),
        } if gold_reference else {},
        "slot_fill_targets": dict(slot_fill_targets),
        "slot_native_information_unit_targets": dict(slot_native_targets),
        "slot_geometry_contract": slot_geometry_contract,
        "orientation_templates": {
            "portrait": "Compact three-row identity header, two-column upper context, content-filled method/results bands, and dense lower analysis grid.",
            "wide": "Compact three-row identity header, contract-defined 3x2/4x2/3x3 main panels with content-filled method/results/table slots, and optional lower qualitative/ablation/synthesis lane.",
        },
        "visual_style": {
            "palette": list(style.get("palette") or [])[:8],
            "typography": {
                "font_family": '"Times New Roman", Times, Georgia, serif',
                "title_font_size_px": 56,
                "identity_rows_font_size_px": 28,
                "section_heading_font_size_px": 36,
                "body_font_size_px": 24,
                "readout_font_size_px": 24,
                "table_text_font_size_px": 24,
                "caption_label_font_size_px": 20,
                "font_size_tolerance_px": 0.5,
                "times_new_roman_family_ratio_required": 1.0,
            },
            "density": style.get("density") or "high text and table density with narrow gutters",
            "identity_treatment": style.get("identity_treatment") or "compact three-line title band",
        },
        "layout_slots": [
            {"slot_id": "title_meta", "placement": "top identity band", "must_include": ["title", "authors", "school_institution_company_names"], "must_exclude": ["venue", "archive", "citation_contact_text", "project_code_resource_links", "logos_icons_qr_or_text_badges", "thesis", "right_side_summary", "tagline", "generic_source_backed_descriptor", "pipeline_process_descriptor", "authored_html_descriptor"]},
            {"slot_id": "problem_contribution", "placement": "top context panel", "must_include": ["problem", "bottleneck", "numbered claims"]},
            {"slot_id": "model_card", "placement": "top or upper-middle panel", "must_include": ["native fields", "method notes"]},
            {"slot_id": "method_pipeline", "placement": "large middle band", "must_include": ["flow boxes", "arrows", "stage labels"]},
            {"slot_id": "results_table", "placement": "middle or lower evidence panel", "must_include": ["native benchmark table", "highlighted deltas"]},
            {"slot_id": "ablation_analysis", "placement": "analysis panel", "must_include": ["mini table", "takeaway bullets"]},
            {"slot_id": "limitations_future", "placement": "lower synthesis panel", "must_include": ["limitations", "future work"]},
            {"slot_id": "synthesis_takeaway", "placement": "footer or lower synthesis band", "must_include": ["scientific takeaway", "limitation or result interpretation"]},
        ],
        "layout_rules": [
            "Use colored numbered section bars so the reading path is obvious.",
            "Keep gutters narrow and fill panel interiors with native evidence instead of whitespace.",
            "Bind every selected source visual/table to a local claim and explanatory caption in the same panel.",
            "Use source figures/tables plus native editable tables/cards/pipelines/formulas as the primary evidence mass.",
            "Optimize for readable source visual/table area, native units, and filled panel interiors before adding text.",
            "Keep source figures/tables readable when used; let native reconstruction carry only genuinely figure-light papers.",
            "Do not place detached screenshot walls or generic abstract cards.",
        ],
    }


def _editorial_flow_reference_archetype_skeleton() -> dict[str, Any]:
    return {
        "profile": _EDITORIAL_FLOW_PROFILE,
        "reference_archetype": "conference_editorial_flow_board",
        "layout": {
            "identity_header": "compact full-width three-row paper title/authors/organization-names band",
            "body": "exactly three poster columns, each a normal-flow stack of multiple poster-section blocks",
            "section_style": "dark section bars, compact subsection labels, thin separators",
        },
        "hard_constraints": {
            "column_count": 3,
            "min_sections_total": 7,
            "target_sections_total": 9,
            "source_assets_are_subjects": True,
            "visible_figcaption_allowed": False,
            "panel_content_plan_allowed": False,
        },
        "forbid": [
            "old poster-grid six-card body",
            "flow-panel wall",
            "child lane DOM",
            "visible Fig./Figure/Table-number caption rows",
            "one shared text flow for multiple figures",
        ],
        "source_flow_unit": (
            "Use one figure-flow-unit/source-flow-unit per source figure or table. "
            "The bound asset and short local readout are direct siblings so text can "
            "wrap around that asset with float/shape-outside when appropriate."
        ),
    }


def _editorial_flow_panel_plan(
    *,
    recommended_figures: dict[str, list[str]],
    selected_visuals: list[str],
) -> list[dict[str, Any]]:
    def visual_ids(*buckets: str, limit: int) -> list[str]:
        values: list[str] = []
        for bucket in buckets:
            values.extend(recommended_figures.get(bucket) or [])
        candidates = _take_unique(values, selected_visuals, limit=max(limit, len(selected_visuals)))
        candidates.extend(layer_id for layer_id in selected_visuals if layer_id not in candidates)
        return candidates[:limit]

    return [
        {
            "column_id": "left_story",
            "role": "motivation_and_context",
            "section_targets": ["The Motivation", "Existing Methods", "Contributions"],
            "visual_ids": visual_ids("fallback", "qualitative", "table", limit=2),
            "layout": "normal-flow poster sections with dark bars; no card grid",
            "text_budget": "short local readouts only; source visuals/tables carry the evidence",
        },
        {
            "column_id": "middle_method",
            "role": "method_flow",
            "section_targets": ["The Method", "How It Works", "Applications"],
            "visual_ids": visual_ids("method", "fallback", limit=3),
            "layout": "large method figures/tables with one flow unit per source asset",
            "text_budget": "readout paragraphs of 20-90 words; no visible Fig./Table caption rows",
        },
        {
            "column_id": "right_results",
            "role": "results_and_analysis",
            "section_targets": ["The Results", "Ablation / Analysis", "Takeaway"],
            "visual_ids": visual_ids("table", "evidence", "qualitative", limit=4),
            "layout": "source table/chart first, compact comparison table rows and concise visual interpretation nearby",
            "text_budget": "compact claims and takeaways; avoid paragraph-heavy summary blocks",
        },
    ]


def _dense_synthesis_panel_plan(
    targets: dict[str, Any],
    *,
    recommended_figures: dict[str, list[str]],
    selected_visuals: list[str],
) -> list[dict[str, Any]]:
    assigned_visuals: list[str] = []

    def visual_ids(*buckets: str, limit: int) -> list[str]:
        values: list[str] = []
        for bucket in buckets:
            values.extend(recommended_figures.get(bucket) or [])
        candidates = _take_unique(values, selected_visuals, limit=max(limit, len(selected_visuals)))
        candidates.extend(layer_id for layer_id in selected_visuals if layer_id not in candidates)
        fresh = [layer_id for layer_id in candidates if layer_id not in assigned_visuals]
        out = fresh[:limit]
        _extend_unique(assigned_visuals, out, limit=len(selected_visuals))
        return out

    target_panels = _safe_int(targets.get("target_panel_count"), default=8)
    return [
        {
            "slot_id": "title_meta",
            "purpose": "Identity header: exactly three text rows only for paper title, authors, and school/institution/company names.",
            "panel_job": "title_meta",
            "text_budget": "title/authors/school-institution-company names only; no venue, resource links, citation/contact text, logos, icons, QR codes, thesis, tagline, right-side summary, or generator/process labels",
            "visual_ids": [],
            "space_fill_policy": "compact text-only three-line title band; venue/project/provenance metadata is omitted from the generated header; never use visible labels like Paper poster, authored HTML, or no generated evidence imagery",
        },
        {
            "slot_id": "problem_contribution",
            "purpose": "Problem and contribution: why the paper exists and the 2-4 claims that make it matter.",
            "panel_job": "problem_contribution",
            "text_budget": "4-6 compact bullets total; no abstract paragraph",
            "visual_ids": visual_ids("method", "evidence", "qualitative", "fallback", limit=1),
            "space_fill_policy": "use one source figure/table, concise visual interpretation, and source-grounded bullets; do not split into separate oversized motivation/contribution panels",
        },
        {
            "slot_id": "model_card",
            "purpose": "Model/system card: architecture, modalities, parameters, data, or capabilities.",
            "panel_job": "model_card",
            "text_budget": "compact labeled fields, method notes, and source-grounded bullets",
            "visual_ids": visual_ids("method", limit=1),
            "space_fill_policy": "build a native model card when source figures are limited",
        },
        {
            "slot_id": "method_pipeline",
            "purpose": "Method or training/inference pipeline as stages.",
            "panel_job": "method_pipeline",
            "text_budget": "short stage labels and one caption",
            "visual_ids": visual_ids("method", "fallback", limit=2),
            "space_fill_policy": "use editable pipeline boxes/arrows and formulas instead of raw prose",
        },
        {
            "slot_id": "results_table",
            "purpose": "Benchmark/results table or leaderboard.",
            "panel_job": "results_table",
            "text_budget": "native table plus one takeaway",
            "visual_ids": visual_ids("table", "evidence", limit=2),
            "space_fill_policy": "prefer a legible native table, compact comparison table, or short result discussion",
        },
        {
            "slot_id": "ablation_analysis",
            "purpose": "Ablation, analysis, comparison, or result interpretation.",
            "panel_job": "ablation_analysis",
            "text_budget": "2-4 analysis bullets with highlighted deltas",
            "visual_ids": visual_ids("evidence", "table", limit=2),
            "space_fill_policy": "use ablation or limitation notes, compact comparison tables, and source-grounded bullets",
        },
        {
            "slot_id": "limitations_future",
            "purpose": "Limitations, failure modes, and future work.",
            "panel_job": "limitations_future",
            "text_budget": "2-4 honest caveats or next steps",
            "visual_ids": visual_ids("qualitative", "evidence", "fallback", limit=1),
            "space_fill_policy": "use a source error/qualitative/example visual when available plus compact caveat cards; do not omit this panel",
        },
        {
            "slot_id": "synthesis_takeaway",
            "purpose": "Final scientific takeaway, result interpretation, or limitation synthesis.",
            "panel_job": "synthesis_takeaway",
            "text_budget": "2-4 source-backed takeaway/limitation/result rows",
            "visual_ids": visual_ids("evidence", "table", "qualitative", "method", "fallback", limit=1),
            "space_fill_policy": "anchor the final synthesis with a source figure/table, short result discussion, or takeaway sentence when available; provenance remains metadata-only unless the user explicitly asks for visible citation/contact text",
        },
    ][:max(7, min(8, target_panels))]


def _compact_paper_visual_storyboard_for_planner(
    storyboard: dict[str, Any] | None,
) -> dict[str, Any]:
    if not isinstance(storyboard, dict) or not storyboard:
        return {}
    raw_selected = [
        item for item in list(storyboard.get("selected_assets") or [])
        if isinstance(item, dict) and str(item.get("asset_id") or "").strip()
        and _storyboard_item_planner_visible(item)
    ]
    raw_primary = [
        item for item in list(storyboard.get("primary_assets") or [])
        if isinstance(item, dict) and str(item.get("asset_id") or "").strip()
        and _storyboard_item_planner_visible(item)
    ]
    raw_secondary = [
        item for item in list(storyboard.get("secondary_assets") or [])
        if isinstance(item, dict) and str(item.get("asset_id") or "").strip()
        and _storyboard_item_planner_visible(item)
    ]
    reserve = [
        item for item in list(storyboard.get("reserve_assets") or [])
        if isinstance(item, dict) and str(item.get("asset_id") or "").strip()
        and _storyboard_item_planner_visible(item)
    ]
    records_by_id = {
        str(item.get("asset_id") or ""): item
        for item in [*raw_selected, *raw_primary, *raw_secondary, *reserve]
        if str(item.get("asset_id") or "").strip()
    }
    minimum_count = _storyboard_minimum_source_visual_count(storyboard)
    selected_ids = set(constrain_optional_source_visual_ids(
        [str(item.get("asset_id") or "") for item in raw_selected],
        records_by_id,
        minimum_count=minimum_count,
    ))
    secondary_ids = set(constrain_optional_source_visual_ids(
        [str(item.get("asset_id") or "") for item in [*raw_secondary, *raw_selected, *reserve]],
        records_by_id,
        minimum_count=minimum_count,
    ))
    selected = [item for item in raw_selected if str(item.get("asset_id") or "") in selected_ids]
    primary = [
        item for item in raw_primary
        if classify_source_visual(str(item.get("asset_id") or ""), item).get("visual_selection_tier") == "eligible"
    ]
    primary_ids = {str(item.get("asset_id") or "") for item in primary}
    secondary: list[dict[str, Any]] = []
    seen_secondary: set[str] = set()
    for item in [*raw_secondary, *raw_selected]:
        asset_id = str(item.get("asset_id") or "")
        if asset_id in secondary_ids and asset_id not in primary_ids and asset_id not in seen_secondary:
            secondary.append(item)
            seen_secondary.add(asset_id)
    panel_jobs = [
        job for job in list(storyboard.get("panel_jobs") or [])
        if isinstance(job, dict)
    ]
    rejected = [
        item for item in list(storyboard.get("rejected_assets") or [])
        if isinstance(item, dict)
    ]
    return {
        "kind": storyboard.get("kind") or "paper_visual_storyboard",
        "version": storyboard.get("version") or 1,
        "manifest_path": "paper_visual_storyboard.json",
        "central_thesis": storyboard.get("central_thesis"),
        "storyline": list(storyboard.get("storyline") or []),
        "target_visual_count": storyboard.get("target_visual_count"),
        "selected_assets": selected,
        "primary_assets": primary or [
            item for item in selected
            if classify_source_visual(str(item.get("asset_id") or ""), item).get("visual_selection_tier") == "eligible"
        ][:min(6, len(selected))],
        "secondary_assets": secondary,
        "reserve_assets": [
            item for item in reserve
            if classify_source_visual(str(item.get("asset_id") or ""), item).get("visual_selection_tier") == "eligible"
        ],
        "rejected_assets": [
            {
                "asset_id": item.get("asset_id"),
                "reason": item.get("reason"),
            }
            for item in rejected[:12]
            if str(item.get("asset_id") or "").strip()
        ],
        "rejected_asset_ids": _take_unique_values([
            str(item.get("asset_id") or "")
            for item in rejected
            if str(item.get("asset_id") or "").strip()
        ]),
        "panel_jobs": panel_jobs,
        "selection_policy": storyboard.get("selection_policy") or {},
        "metrics": storyboard.get("metrics") or {},
        "omitted_rejected_asset_count": len(rejected),
    }


def _storyboard_item_planner_visible(item: dict[str, Any]) -> bool:
    asset_id = str(item.get("asset_id") or item.get("layer_id") or "")
    return _is_planner_visible(asset_id, {}, item)


def _compact_paper_visual_provenance_for_planner(
    provenance: dict[str, Any] | None,
    *,
    paper_visual_storyboard: dict[str, Any] | None = None,
) -> dict[str, Any]:
    if not isinstance(provenance, dict) or not provenance:
        return {}
    assets = [
        asset for asset in list(provenance.get("assets") or [])
        if isinstance(asset, dict) and str(asset.get("asset_id") or "").strip()
        and not _is_audit_only_source_asset(asset)
    ]
    planner_assets = [
        asset for asset in assets
        if _is_planner_visible(str(asset.get("asset_id") or ""), {}, asset)
    ]
    selected_ids = {
        str(item.get("asset_id") or "")
        for item in list((paper_visual_storyboard or {}).get("selected_assets") or [])
        if isinstance(item, dict) and str(item.get("asset_id") or "").strip()
        and _storyboard_item_planner_visible(item)
    }
    compact_assets = [asset for asset in planner_assets if str(asset.get("asset_id") or "") in selected_ids]
    if not compact_assets:
        compact_assets = planner_assets[:12]
    metrics = dict(provenance.get("metrics") or {})
    metrics["full_asset_count"] = len(assets)
    metrics["planner_eligible_asset_count"] = len(planner_assets)
    metrics["planner_visible_asset_count"] = len(compact_assets)
    metrics["omitted_asset_count"] = max(0, len(assets) - len(compact_assets))
    return {
        "kind": provenance.get("kind") or "paper_visual_provenance",
        "version": provenance.get("version") or 1,
        "manifest_path": "paper_visual_provenance.json",
        "source_documents": list(provenance.get("source_documents") or []),
        "generation_policy": provenance.get("generation_policy") or {},
        "assets": compact_assets,
        "metrics": metrics,
    }


def _is_audit_only_source_asset(asset: dict[str, Any]) -> bool:
    return str(asset.get("kind") or asset.get("extract_strategy") or "") == "source_table_crop_candidate"


def _compact_paper_memory_dossier_for_planner(
    dossier: dict[str, Any] | None,
) -> dict[str, Any]:
    if not isinstance(dossier, dict) or dossier.get("kind") != "paper_memory_dossier":
        return {}
    return {
        "kind": "paper_memory_dossier",
        "version": dossier.get("version"),
        "manifest_path": "paper_memory_dossier.json",
        "markdown_path": "paper_memory_dossier.md",
        "source_memory_cache_key": dossier.get("source_memory_cache_key"),
        "model": dossier.get("model"),
        "section_count": len(dossier.get("sections") or []),
        "sections": [
            {
                "id": section.get("id"),
                "panel_role": section.get("panel_role"),
                "title": section.get("title"),
                "claim": section.get("claim"),
                "poster_copy_suggestion": section.get("poster_copy_suggestion"),
                "visual_ids": section.get("visual_ids") or [],
                "evidence_refs": [
                    {
                        "chunk_id": ref.get("chunk_id"),
                        "page": ref.get("page"),
                        "source_id": ref.get("source_id"),
                        "safe_to_quote": ref.get("safe_to_quote"),
                    }
                    for ref in (section.get("evidence_refs") or [])
                    if isinstance(ref, dict)
                ],
            }
            for section in (dossier.get("sections") or [])
            if isinstance(section, dict)
        ],
    }


def _compact_paper_resources_for_planner(resources: dict[str, Any] | None) -> dict[str, Any]:
    if not isinstance(resources, dict) or not resources:
        return {}
    compact_resources = [
        {
            "type": item.get("type"),
            "label": item.get("label"),
            "href": item.get("href"),
            "title": item.get("title"),
            "icon": item.get("icon"),
            "source": item.get("source"),
            "confidence": item.get("confidence"),
            "evidence": item.get("evidence"),
        }
        for item in list(resources.get("resources") or [])[:12]
        if isinstance(item, dict)
    ]
    return {
        "kind": resources.get("kind") or "paper_resource_manifest",
        "version": resources.get("version") or 1,
        "manifest_path": "paper_resource_manifest.json",
        "resource_count": resources.get("resource_count") or len(compact_resources),
        "resources": compact_resources,
        "resource_chips": list(resources.get("resource_chips") or [])[:12],
        "resource_recall_audit": resources.get("resource_recall_audit") or {},
        "recall_audit_path": "paper_resource_recall_audit.json",
        "warnings": list(resources.get("warnings") or [])[:8],
    }


def _select_poster_visual_ids(
    *,
    recommended_figures: dict[str, list[str]],
    rendered: dict[str, dict[str, Any]],
    visual_candidate_scores: list[dict[str, Any]],
    target_visual_count: int,
    preferred_visual_ids: list[str] | None = None,
    forbidden_visual_ids: list[str] | None = None,
) -> list[str]:
    forbidden = {str(item or "").strip() for item in (forbidden_visual_ids or []) if str(item or "").strip()}
    selected: list[str] = []
    pool_limit = max(
        target_visual_count * 4,
        target_visual_count + 12,
        len(visual_candidate_scores or []),
        24,
    )
    _extend_unique(
        selected,
        [
            item for item in (preferred_visual_ids or [])
            if str(item or "").strip() not in forbidden
        ],
        limit=pool_limit,
    )
    for bucket in ("method", "table", "evidence", "qualitative"):
        _extend_unique(
            selected,
            [
                item for item in (recommended_figures.get(bucket) or [])
                if str(item or "").strip() not in forbidden
            ],
            limit=pool_limit,
        )
    scored = [
        str(item.get("layer_id") or "")
        for item in sorted(
            visual_candidate_scores or [],
            key=lambda item: (
                1 if _is_planner_visible(str(item.get("layer_id") or ""), {}, item) else 0,
                1 if str(item.get("visual_role") or "") in {"method", "table", "evidence"} else 0,
                0 if "low_value_example_crop" in set(item.get("curation_flags") or []) else 1,
                _safe_int(item.get("visual_score"), default=0),
                _VISUAL_ROLE_PRIORITY.get(str(item.get("visual_role") or ""), 0),
                -_safe_int(item.get("source_page"), default=999),
            ),
            reverse=True,
        )
        if str(item.get("layer_id") or "")
        and str(item.get("layer_id") or "") not in forbidden
        and _is_planner_visible(str(item.get("layer_id") or ""), {}, item)
    ]
    _extend_unique(selected, scored, limit=pool_limit)
    _extend_unique(
        selected,
        [item for item in (recommended_figures.get("fallback") or []) if str(item or "").strip() not in forbidden],
        limit=pool_limit,
    )
    return _sanitize_selected_visual_ids(
        _limit_source_table_visuals(selected, max_tables=3),
        rendered=rendered,
        visual_candidate_scores=visual_candidate_scores,
        forbidden_visual_ids=list(forbidden),
        limit=target_visual_count,
        max_tables=3,
    )


def _limit_source_table_visuals(values: list[str], *, max_tables: int) -> list[str]:
    out: list[str] = []
    table_count = 0
    for value in values:
        layer_id = str(value or "").strip()
        if not layer_id:
            continue
        if layer_id.startswith("ingest_table_"):
            if table_count >= max_tables:
                continue
            table_count += 1
        out.append(layer_id)
    return out


def _extend_unique(target: list[str], values: list[str], *, limit: int) -> None:
    for value in values:
        item = str(value or "").strip()
        if not item or item in target:
            continue
        target.append(item)
        if len(target) >= limit:
            return


def _take_unique(values: list[str], allowed: list[str], *, limit: int) -> list[str]:
    allowed_set = set(allowed)
    out: list[str] = []
    for value in values:
        item = str(value or "").strip()
        if item and item in allowed_set and item not in out:
            out.append(item)
        if len(out) >= limit:
            break
    return out


def _take_unique_values(values: list[Any]) -> list[str]:
    out: list[str] = []
    for value in values:
        item = str(value or "").strip()
        if item and item not in out:
            out.append(item)
    return out


def _merge_visual_priority(
    priority_ids: list[str],
    existing_ids: list[str],
    *,
    limit: int,
) -> list[str]:
    out: list[str] = []
    for value in list(priority_ids or []) + list(existing_ids or []):
        item = str(value or "").strip()
        if item and item not in out:
            out.append(item)
        if limit > 0 and len(out) >= max(limit, len(priority_ids or [])):
            break
    return out


def _high_priority_source_visual_ids(
    *,
    provenance_assets: dict[str, dict[str, Any]],
    rendered: dict[str, dict[str, Any]],
    visual_candidate_scores: list[dict[str, Any]],
    preferred_ids: list[str],
    limit: int,
) -> list[str]:
    score_by_id = {
        str(item.get("layer_id") or ""): item
        for item in visual_candidate_scores or []
        if isinstance(item, dict) and str(item.get("layer_id") or "").strip()
    }
    preferred_set = {
        str(item or "").strip()
        for item in preferred_ids or []
        if str(item or "").strip()
    }

    def candidate_score(asset_id: str) -> tuple[int, ...]:
        asset = provenance_assets.get(asset_id) or {}
        rec = rendered.get(asset_id) or {}
        score_rec = score_by_id.get(asset_id) or {}
        role = str(asset.get("visual_role") or rec.get("visual_role") or score_rec.get("visual_role") or "").lower()
        role_boost = 1 if role in {"method", "table", "evidence"} else 0
        page = _safe_int(asset.get("source_page") or rec.get("source_page"), default=999)
        visual_score = _safe_int(asset.get("visual_score") or rec.get("visual_score") or score_rec.get("visual_score"), default=0)
        width_px = _safe_int(asset.get("output_width_px"), default=0)
        height_px = _safe_int(asset.get("output_height_px"), default=0)
        min_side = min(width_px, height_px)
        return (
            0 if _selected_visual_reject_reasons(asset_id, rec, asset, score_rec) else 1,
            role_boost,
            visual_score,
            min(min_side, 2400),
            -page,
        )

    ids: list[str] = []
    for asset_id, asset in provenance_assets.items():
        rec = rendered.get(asset_id) or {}
        if _hide_from_planner_visual_catalog(asset_id, rec):
            continue
        score_rec = score_by_id.get(asset_id) or {}
        if _selected_visual_reject_reasons(asset_id, rec, asset, score_rec):
            continue
        role = str(asset.get("visual_role") or rec.get("visual_role") or score_rec.get("visual_role") or "").lower()
        visual_score = _safe_int(
            asset.get("visual_score") or rec.get("visual_score") or score_rec.get("visual_score"),
            default=0,
        )
        if role not in {"method", "table", "evidence"}:
            continue
        if visual_score < 60 and asset_id not in preferred_set:
            continue
        if asset_id not in preferred_set:
            if not bool(asset.get("protected_anchor") or rec.get("protected_anchor") or score_rec.get("protected_anchor")):
                if visual_score < 72 and not bool(asset.get("captioned_source_group") or rec.get("captioned_source_group")):
                    continue
        ids.append(asset_id)
    preferred_order = {str(item or "").strip(): idx for idx, item in enumerate(preferred_ids or []) if str(item or "").strip()}
    ids = sorted(
        ids,
        key=lambda asset_id: (
            candidate_score(asset_id),
            -preferred_order.get(asset_id, 9999),
        ),
        reverse=True,
    )
    out: list[str] = []
    table_count = 0
    for asset_id in ids:
        if asset_id.startswith("ingest_table_"):
            if table_count >= 3:
                continue
            table_count += 1
        if asset_id in rendered and asset_id not in out:
            out.append(asset_id)
        if limit > 0 and len(out) >= limit:
            break
    return out


def _poster_visual_record(layer_id: str, rec: dict[str, Any]) -> dict[str, Any]:
    out = {
        "layer_id": layer_id,
        "kind": rec.get("kind") or ("table" if layer_id.startswith("ingest_table_") else "image"),
        "visual_role": rec.get("visual_role"),
        "visual_score": rec.get("visual_score"),
        "caption_short": rec.get("caption_short"),
        "source_page": rec.get("source_page"),
        "source_bbox_pdf_points": rec.get("source_bbox_pdf_points"),
        "source_pdf_sha256": rec.get("source_pdf_sha256"),
        "image_size": rec.get("image_size"),
        "curation_flags": list(rec.get("curation_flags") or []),
        "crop_quality_flags": list(rec.get("crop_quality_flags") or []),
        "designer_eligible": bool(rec.get("designer_eligible", True)),
        "planner_eligible": bool(rec.get("planner_eligible", True)),
        "planner_visible": bool(rec.get("planner_visible", True)),
        "designer_reject_reasons": list(rec.get("designer_reject_reasons") or []),
        "planner_reject_reasons": list(rec.get("planner_reject_reasons") or []),
        "severe_crop_flags": list(rec.get("severe_crop_flags") or []),
        "extract_strategy": rec.get("extract_strategy") or rec.get("kind"),
        "protected_anchor": bool(rec.get("protected_anchor")),
        "anchor_kind": rec.get("anchor_kind"),
        "anchor_label": rec.get("anchor_label"),
        "anchor_reason": rec.get("anchor_reason"),
        "captioned_source_group": bool(rec.get("captioned_source_group")),
        "source_group_id": rec.get("source_group_id"),
        "source_group_kind": rec.get("source_group_kind"),
        "source_group_label": rec.get("source_group_label"),
        "source_group_caption": rec.get("source_group_caption"),
        "source_group_source": rec.get("source_group_source"),
    }
    table_metadata = _table_metadata_payload(rec)
    if table_metadata:
        out["table_metadata"] = table_metadata
    return out


def _poster_source_asset_record(layer_id: str, asset: dict[str, Any]) -> dict[str, Any]:
    out = {
        "asset_id": layer_id,
        "kind": asset.get("kind"),
        "source_page": asset.get("source_page"),
        "source_bbox_pdf_points": asset.get("source_bbox_pdf_points"),
        "output_file": asset.get("output_file"),
        "output_sha256": asset.get("output_sha256"),
        "output_width_px": asset.get("output_width_px"),
        "output_height_px": asset.get("output_height_px"),
        "caption_short": asset.get("caption_short"),
        "caption_full": asset.get("caption_full"),
        "visual_role": asset.get("visual_role"),
        "visual_score": asset.get("visual_score"),
        "curation_flags": list(asset.get("curation_flags") or []),
        "crop_quality_flags": list(asset.get("crop_quality_flags") or []),
        "placement_quality_flags": list(asset.get("placement_quality_flags") or []),
        "caption_association_method": asset.get("caption_association_method") or "unmatched",
        "visual_selection_tier": asset.get("visual_selection_tier"),
        "eligibility_policy_version": asset.get("eligibility_policy_version"),
        "unmatched_caption": bool(asset.get("unmatched_caption")),
        "designer_eligible": bool(asset.get("designer_eligible", True)),
        "planner_eligible": bool(asset.get("planner_eligible", True)),
        "planner_visible": bool(asset.get("planner_visible", True)),
        "designer_reject_reasons": list(asset.get("designer_reject_reasons") or []),
        "planner_reject_reasons": list(asset.get("planner_reject_reasons") or []),
        "severe_crop_flags": list(asset.get("severe_crop_flags") or []),
        "extract_strategy": asset.get("extract_strategy"),
        "protected_anchor": bool(asset.get("protected_anchor")),
        "anchor_kind": asset.get("anchor_kind"),
        "anchor_label": asset.get("anchor_label"),
        "anchor_reason": asset.get("anchor_reason"),
        "captioned_source_group": bool(asset.get("captioned_source_group")),
        "source_group_id": asset.get("source_group_id"),
        "source_group_kind": asset.get("source_group_kind"),
        "source_group_label": asset.get("source_group_label"),
        "source_group_caption": asset.get("source_group_caption"),
        "source_group_source": asset.get("source_group_source"),
    }
    if isinstance(asset.get("table_metadata"), dict):
        out["table_metadata"] = _json_clone(asset.get("table_metadata"))
    return out


def _table_metadata_payload(rec: dict[str, Any]) -> dict[str, Any]:
    if str(rec.get("kind") or "").lower() != "table":
        return {}
    return {
        "headers": _json_clone(rec.get("headers") or []),
        "rows": _json_clone(rec.get("rows") or []),
        "col_highlight_rule": _json_clone(rec.get("col_highlight_rule") or []),
        "title": rec.get("title") or "",
        "table_visual_source": rec.get("table_visual_source") or "original_pdf_crop",
        "table_parse_status": rec.get("table_parse_status") or "parsed",
        "table_parse_error": rec.get("table_parse_error") or "",
    }


def _build_paper_visual_provenance(
    *,
    layer_ids: list[str],
    rendered: dict[str, dict[str, Any]],
    run_dir: Path,
    summaries: list[dict[str, Any]],
) -> dict[str, Any]:
    """Build a Codex-style source-backed manifest for paper visual assets."""
    seen: set[str] = set()
    assets: list[dict[str, Any]] = []
    for raw_id in layer_ids:
        layer_id = str(raw_id or "").strip()
        if not layer_id or layer_id in seen:
            continue
        seen.add(layer_id)
        rec = rendered.get(layer_id) or {}
        if not isinstance(rec, dict) or str(rec.get("source") or "") != "ingested_pdf":
            continue
        src_path = Path(str(rec.get("src_path") or ""))
        output_sha = None
        output_w = None
        output_h = None
        if src_path.exists():
            try:
                output_sha = sha256_file(src_path)
            except OSError:
                output_sha = rec.get("sha256")
            try:
                with Image.open(src_path) as im:
                    output_w, output_h = im.size
            except Exception:
                output_w, output_h = _image_size_tuple(rec.get("image_size"))
        else:
            output_sha = rec.get("sha256")
            output_w, output_h = _image_size_tuple(rec.get("image_size"))
        if output_w is None or output_h is None:
            output_w, output_h = _image_size_tuple(rec.get("image_size"))
        material_quality = _paper_visual_material_quality(src_path)
        material_flags = list(rec.get("curation_flags") or [])
        for flag in material_quality.get("warnings") or []:
            if flag not in material_flags:
                material_flags.append(flag)
        asset: dict[str, Any] = {
            "asset_id": layer_id,
            "kind": rec.get("kind") or ("table" if layer_id.startswith("ingest_table_") else "image"),
            "source_pdf": rec.get("source_pdf") or Path(str(rec.get("source_file") or "")).name,
            "source_pdf_sha256": rec.get("source_pdf_sha256"),
            "source_page": rec.get("source_page"),
            "source_bbox_pdf_points": rec.get("source_bbox_pdf_points"),
            "caption_short": rec.get("caption_short"),
            "caption_full": rec.get("caption") or rec.get("title"),
            "output_file": _path_relative_to(src_path, run_dir) if src_path else None,
            "output_sha256": output_sha,
            "output_width_px": output_w,
            "output_height_px": output_h,
            "visual_role": rec.get("visual_role"),
            "visual_score": rec.get("visual_score"),
            "curation_reason": rec.get("curation_reason"),
            "curation_flags": material_flags,
            "crop_quality_flags": list(rec.get("crop_quality_flags") or []),
            "placement_quality_flags": list(rec.get("placement_quality_flags") or []),
            "caption_association_method": rec.get("caption_association_method") or "unmatched",
            "material_quality": material_quality,
            "extract_strategy": rec.get("extract_strategy") or rec.get("kind"),
            "protected_anchor": bool(rec.get("protected_anchor")),
            "anchor_kind": rec.get("anchor_kind"),
            "anchor_label": rec.get("anchor_label"),
            "anchor_reason": rec.get("anchor_reason"),
            "captioned_source_group": bool(rec.get("captioned_source_group")),
            "source_group_id": rec.get("source_group_id"),
            "source_group_kind": rec.get("source_group_kind"),
            "source_group_label": rec.get("source_group_label"),
            "source_group_caption": rec.get("source_group_caption"),
            "source_group_source": rec.get("source_group_source"),
            "table_parse_status": rec.get("table_parse_status"),
            "extraction": _provenance_extraction_label(rec),
        }
        if rec.get("source_image_xref") is not None:
            asset["source_image_xref"] = rec.get("source_image_xref")
        if rec.get("parent_layer_id"):
            asset["parent_asset_id"] = rec.get("parent_layer_id")
        table_metadata = _table_metadata_payload(rec)
        if table_metadata:
            asset["table_metadata"] = table_metadata
        asset.update(_visual_eligibility_payload(layer_id, rec, asset))
        assets.append(asset)

    source_documents: list[dict[str, Any]] = []
    seen_docs: set[str] = set()
    for summary in summaries:
        if not isinstance(summary, dict) or summary.get("type") != "pdf":
            continue
        pdf_path = Path(str(summary.get("file") or ""))
        key = str(pdf_path)
        if not key or key in seen_docs:
            continue
        seen_docs.add(key)
        source_documents.append({
            "source_pdf": pdf_path.name,
            "source_pdf_sha256": _safe_sha256_file(pdf_path),
            "source_page_count": (summary.get("manifest") or {}).get("page_count"),
            "title": (summary.get("manifest") or {}).get("title"),
        })

    return {
        "kind": "paper_visual_provenance",
        "version": 1,
        "source_documents": source_documents,
        "generation_policy": {
            "used_ai_generated_imagery": False,
            "used_external_images": False,
            "all_raster_assets_derived_from_source_pdf": True,
            "text_and_layout_authored_later": True,
        },
        "assets": assets,
        "metrics": {
            "asset_count": len(assets),
            "source_backed_asset_count": sum(1 for a in assets if a.get("source_bbox_pdf_points")),
            "captioned_asset_count": sum(1 for a in assets if a.get("caption_full") or a.get("caption_short")),
            "missing_output_count": sum(1 for a in assets if not a.get("output_sha256")),
            "designer_eligible_asset_count": sum(1 for a in assets if a.get("designer_eligible")),
            "planner_visible_asset_count": sum(1 for a in assets if a.get("planner_visible")),
            "severe_crop_rejected_asset_count": sum(1 for a in assets if a.get("severe_crop_flags")),
            "high_edge_whitespace_count": sum(
                1 for a in assets
                if "high_edge_whitespace" in ((a.get("material_quality") or {}).get("warnings") or [])
            ),
            "mostly_white_visual_count": sum(
                1 for a in assets
                if "mostly_white_visual" in ((a.get("material_quality") or {}).get("warnings") or [])
            ),
        },
    }


def _paper_visual_material_quality(path: Path) -> dict[str, Any]:
    quality: dict[str, Any] = {
        "material_score": None,
        "white_ratio": None,
        "edge_white_ratio": None,
        "warnings": [],
    }
    if not path or not path.exists() or not path.is_file():
        quality["warnings"].append("image_payload_unavailable")
        return quality
    try:
        with Image.open(path) as im:
            rgb = im.convert("RGB")
            quality["width_px"] = rgb.width
            quality["height_px"] = rgb.height
            white_ratio, edge_white_ratio = _paper_visual_white_ratios(rgb)
            quality["white_ratio"] = round(white_ratio, 4)
            quality["edge_white_ratio"] = round(edge_white_ratio, 4)
            if min(rgb.size) < 180:
                quality["warnings"].append("low_resolution_visual")
            if edge_white_ratio >= 0.82:
                quality["warnings"].append("high_edge_whitespace")
            if white_ratio >= 0.82:
                quality["warnings"].append("mostly_white_visual")
            score = 1.0
            score -= max(0.0, edge_white_ratio - 0.35) * 0.75
            score -= max(0.0, white_ratio - 0.65) * 0.5
            if min(rgb.size) < 180:
                score -= 0.25
            quality["material_score"] = round(max(0.0, min(1.0, score)), 4)
    except Exception as exc:  # noqa: BLE001 - provenance should survive bad images
        quality["warnings"].append(f"image_decode_failed:{type(exc).__name__}")
    return quality


def _paper_visual_white_ratios(img: Image.Image) -> tuple[float, float]:
    thumb = img.copy()
    thumb.thumbnail((320, 320))
    pixels = list(thumb.getdata())
    if not pixels:
        return 0.0, 0.0
    white = sum(1 for pixel in pixels if min(pixel) >= 245)
    edge = max(2, int(min(thumb.size) * 0.08))
    edge_pixels = []
    width, height = thumb.size
    for y in range(height):
        for x in range(width):
            if x < edge or y < edge or x >= width - edge or y >= height - edge:
                edge_pixels.append(thumb.getpixel((x, y)))
    edge_white = sum(1 for pixel in edge_pixels if min(pixel) >= 245)
    return white / len(pixels), edge_white / max(1, len(edge_pixels))


def _safe_sha256_file(path: Path) -> str | None:
    try:
        if path.exists():
            return sha256_file(path)
    except OSError:
        return None
    return None


def _image_size_tuple(raw: Any) -> tuple[int | None, int | None]:
    try:
        w_s, h_s = str(raw or "").split("x", 1)
        return int(w_s), int(h_s)
    except Exception:
        return None, None


def _path_relative_to(path: Path, root: Path) -> str:
    try:
        return str(path.resolve().relative_to(root.resolve()))
    except Exception:
        return str(path)


def _provenance_extraction_label(rec: dict[str, Any]) -> str:
    strategy = str(rec.get("extract_strategy") or rec.get("kind") or "asset")
    if strategy == "sub_panel":
        return "Pillow crop from a PyMuPDF source-PDF parent crop; no generated imagery"
    if strategy == "table":
        return "Original PDF table crop localized from source PDF; parsed rows are metadata for native summaries only"
    return f"PyMuPDF {strategy} extraction from source PDF; no generated imagery"


def _poster_section(
    section_id: str,
    title: str,
    recommended_text_units: dict[str, list[dict[str, Any]]],
    *,
    buckets: tuple[str, ...] | None = None,
    visual_ids: list[str],
    purpose: str,
    bullet_limit: int = 4,
) -> dict[str, Any]:
    buckets = buckets or (section_id,)
    limit = max(1, min(6, int(bullet_limit or 4)))
    bullets: list[dict[str, Any]] = []
    for bucket in buckets:
        for item in list(recommended_text_units.get(bucket) or []):
            if not isinstance(item, dict):
                continue
            text = str(item.get("text") or "").strip()
            if not text:
                continue
            bullets.append({
                "text": text,
                "source": item.get("source"),
                "bucket": item.get("bucket") or bucket,
                "source_ids": list(item.get("source_ids") or []),
                "intended_panel_role": item.get("intended_panel_role") or section_id,
                **({"claim_id": item.get("claim_id")} if item.get("claim_id") else {}),
            })
            if len(bullets) >= limit:
                break
        if len(bullets) >= limit:
            break
    return {
        "section_id": section_id,
        "title": title,
        "purpose": purpose,
        "bullet_budget": "2-4" if limit <= 4 else f"3-{limit}",
        "bullets": bullets[:limit],
        "visual_ids": visual_ids,
    }

def _aspect_from_dims(w: int, h: int) -> str:
    if h <= 0 or w <= 0:
        return "1:1"
    from math import gcd
    g = gcd(w, h)
    return f"{w // g}:{h // g}" if max(w // g, h // g) <= 32 else (
        "16:9" if w > h else "3:4" if h > w else "1:1"
    )


def _sanitize_name(s: str) -> str:
    s = (s or "").strip()
    return re.sub(r"[^\w\- ]", "", s)[:60]
