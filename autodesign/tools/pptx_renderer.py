"""PPTX renderer for deck artifacts (v1.0 #7).

One top-level LayerNode with `kind="slide"` per slide; each slide's `children`
hold `kind in {"text","image","background"}` elements positioned by bbox in
slide-canvas pixel coords. We emit native PowerPoint shapes with native
TextFrames — no Pillow rasterization of text — so the .pptx opens
type-editable in PowerPoint / Keynote / Google Slides.

Per-slide simplified PNGs are also written (for chat preview + critic).

Font embedding is intentionally NOT performed: .pptx delegates font rendering
to the consuming app's font engine. CJK output therefore depends on compatible
fonts being installed in the consuming PowerPoint environment. This mirrors
Paper2Any's approach.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Pt

from ..schema import SafeZone
from ._contract import ToolContext


# 1 pixel at 96 DPI ≈ 9525 EMU. python-pptx accepts Emu(int).
PX_TO_EMU = 9525

DEFAULT_SLIDE_W = 1920
DEFAULT_SLIDE_H = 1080

_ALIGN_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
}


# Child-name substrings that mark a text layer as the slide title. Used by
# `_with_section_prefix` to decide which text node should receive the section
# number prefix at render time.
_TITLE_NAME_HINTS = ("title", "headline", "section_title")


# Coordinates retained only to render persisted layer-graph deck specs that
# predate the HTML-first deck contract. They are not a file-backed template.
_LEGACY_SLOT_BBOX: dict[str, dict[str, tuple[int, int, int, int]]] = {
    "cover": {
        "title": (96, 280, 880, 280),
        "authors": (96, 580, 880, 60),
        "badge": (1660, 80, 180, 40),
        "image_slot": (1000, 0, 920, 1080),
    },
    "section_divider": {
        "section_number": (200, 320, 1520, 50),
        "title": (200, 380, 1520, 200),
        "subtitle": (200, 600, 1520, 60),
    },
    "content": {
        "section_label": (96, 80, 1728, 30),
        "title": (96, 120, 1728, 80),
        "footer": (96, 1020, 1200, 30),
        "slide_number": (1700, 1020, 124, 30),
        "body": (96, 260, 1728, 740),
    },
    "content_with_figure": {
        "section_label": (96, 80, 1728, 30),
        "title": (96, 120, 1728, 80),
        "footer": (96, 1020, 1200, 30),
        "slide_number": (1700, 1020, 124, 30),
        "body": (96, 260, 920, 740),
        "image_slot": (1056, 260, 768, 740),
    },
    "content_with_table": {
        "section_label": (96, 80, 1728, 30),
        "title": (96, 120, 1728, 80),
        "footer": (96, 1020, 1200, 30),
        "slide_number": (1700, 1020, 124, 30),
        "body": (96, 260, 800, 740),
        "table_anchor": (920, 260, 904, 740),
    },
    "closing": {
        "title": (200, 360, 1520, 160),
        "subtitle": (200, 540, 1520, 50),
        "links": (200, 620, 1520, 30),
    },
}


def _legacy_slot_bbox(
    role: str | None,
    slot: str | None,
) -> tuple[int, int, int, int] | None:
    if not role or not slot:
        return None
    return _LEGACY_SLOT_BBOX.get(role, {}).get(slot)


def _with_legacy_bbox(node: Any, role: str | None) -> Any:
    """Supply deterministic geometry only for pre-HTML-first saved specs."""
    if getattr(node, "bbox", None) is not None:
        return node
    slot_bbox = _legacy_slot_bbox(role, getattr(node, "template_slot", None))
    if slot_bbox is None:
        return node
    x, y, w, h = slot_bbox
    try:
        return node.model_copy(update={"bbox": SafeZone(x=x, y=y, w=w, h=h)})
    except Exception:
        return node


def _is_title_child(child: Any) -> bool:
    name = (getattr(child, "name", None) or "").lower()
    return getattr(child, "role", None) == "title" or any(
        h in name for h in _TITLE_NAME_HINTS
    )


def _with_section_prefix(child: Any, section_number: str | None) -> Any:
    """Return a per-render copy of ``child`` with ``section_number``
    prepended to its text (e.g. "§2.2 · Vision tokenizer").

    Pure: never mutates the input child. Returns the original child
    unchanged when there is no section_number, no text, or the prefix
    already appears at the start of the text (idempotent — safe to
    re-render).
    """
    if not section_number:
        return child
    text = (getattr(child, "text", None) or "").strip()
    if not text:
        return child
    if text.startswith(section_number):
        return child
    new_text = f"{section_number} · {text}"
    try:
        return child.model_copy(update={"text": new_text})
    except Exception:
        # Fallback for non-pydantic proxies — shouldn't happen for
        # LayerNode but defensive.
        return child


def write_pptx(spec: Any, pptx_path: Path, ctx: ToolContext) -> int:
    """Walk the slide tree and emit a bbox-positioned .pptx file.

    Deck generation is HTML-first. This renderer remains only for explicit
    compatibility exports and never loads a repository PPTX template.
    """

    canvas = spec.canvas or {}
    slide_w = int(canvas.get("w_px") or DEFAULT_SLIDE_W)
    slide_h = int(canvas.get("h_px") or DEFAULT_SLIDE_H)

    prs = Presentation()
    prs.slide_width = Emu(slide_w * PX_TO_EMU)
    prs.slide_height = Emu(slide_h * PX_TO_EMU)

    blank_layout = prs.slide_layouts[6]  # "Blank" layout — we position everything.

    slide_count = 0
    for node in (spec.layer_graph or []):
        if getattr(node, "kind", None) != "slide":
            continue
        slide = prs.slides.add_slide(blank_layout)
        _render_slide(slide, node, slide_w, slide_h, ctx)
        slide_count += 1

    prs.save(str(pptx_path))
    return slide_count


def write_pptx_hybrid(
    spec: Any,
    pptx_path: Path,
    ctx: ToolContext,
    *,
    slide_pngs: list[Path],
    export_mode: str = "hybrid",
) -> int:
    """Write a frame-first PPTX.

    ``visual`` mode places one full-slide PNG per slide. ``hybrid`` mode uses
    the supplied PNGs as the browser-rendered visual base and overlays editable
    PowerPoint text/table primitives on top.
    """
    canvas = spec.canvas or {}
    slide_w = int(canvas.get("w_px") or DEFAULT_SLIDE_W)
    slide_h = int(canvas.get("h_px") or DEFAULT_SLIDE_H)
    slides = [
        n for n in (getattr(spec, "layer_graph", None) or [])
        if getattr(n, "kind", None) == "slide"
    ]

    prs = Presentation()
    prs.slide_width = Emu(slide_w * PX_TO_EMU)
    prs.slide_height = Emu(slide_h * PX_TO_EMU)
    blank_layout = prs.slide_layouts[6]

    mode = (export_mode or "hybrid").lower()

    for idx, slide_node in enumerate(slides):
        slide = prs.slides.add_slide(blank_layout)
        if idx < len(slide_pngs) and Path(slide_pngs[idx]).exists():
            slide.shapes.add_picture(
                str(slide_pngs[idx]),
                Emu(0), Emu(0),
                width=Emu(slide_w * PX_TO_EMU),
                height=Emu(slide_h * PX_TO_EMU),
            )
        if mode == "hybrid":
            _render_hybrid_overlays(slide, slide_node, slide_w, slide_h, ctx)
        else:
            notes = getattr(slide_node, "speaker_notes", None)
            if notes:
                slide.notes_slide.notes_text_frame.text = notes

    prs.save(str(pptx_path))
    return len(slides)


def _render_hybrid_overlays(
    slide: Any,
    slide_node: Any,
    slide_w: int,
    slide_h: int,
    ctx: ToolContext,
) -> None:
    role = getattr(slide_node, "role", None) or "content"
    children = sorted(
        list(getattr(slide_node, "children", None) or []),
        key=lambda c: int(getattr(c, "z_index", 0) or 0),
    )

    section_number = getattr(slide_node, "section_number", None)
    title_seen = False
    for child in children:
        effective = _with_legacy_bbox(child, role)
        kind = getattr(effective, "kind", None)
        if kind == "text" and not title_seen and _is_title_child(effective):
            effective = _with_section_prefix(effective, section_number)
            title_seen = True
        if kind == "text":
            if getattr(effective, "bbox", None) is not None:
                _add_text_frame(slide, effective, slide_w, slide_h)
        elif kind == "table":
            if getattr(effective, "bbox", None) is not None:
                _add_table(slide, effective, slide_w, slide_h)

    notes = getattr(slide_node, "speaker_notes", None)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def _default_render_slide(slide: Any, slide_node: Any, slide_w: int,
                          slide_h: int, ctx: ToolContext) -> None:
    """Original inline per-slide layout (pre-v2.8.1).

    Walks the slide's children in z_index order, renders every text /
    image / background / table child via the matching `_add_*` helper,
    decorates the first title-bearing text child with v2.7.2's
    `section_number` prefix, and writes `speaker_notes` to the notes
    pane.

    This is the byte-identical fallback the v2.8.1 dispatcher reaches
    for whenever an archetype renderer is unavailable (Phase 2/3
    placeholder OR `archetype="evidence_snapshot"` on a slide without
    a big-number child).
    """
    children = list(getattr(slide_node, "children", None) or [])
    # Sort by z_index so higher z draws on top (pptx respects insertion order).
    children.sort(key=lambda c: int(getattr(c, "z_index", 0) or 0))

    # v2.7.2 — prepend slide.section_number to the first title-bearing
    # text child. Only the first match is decorated so multi-text slides
    # don't double-stamp.
    role = getattr(slide_node, "role", None) or "content"
    section_number = getattr(slide_node, "section_number", None)
    title_seen = False

    for child in children:
        effective = _with_legacy_bbox(child, role)
        kind = getattr(effective, "kind", None)
        if kind == "text" and not title_seen and _is_title_child(effective):
            effective = _with_section_prefix(effective, section_number)
            title_seen = True
        if kind == "background":
            _add_background(slide, effective, slide_w, slide_h)
        elif kind == "image":
            _add_picture(slide, effective, slide_w, slide_h)
        elif kind == "text":
            _add_text_frame(slide, effective, slide_w, slide_h)
        elif kind == "table":
            _add_table(slide, effective, slide_w, slide_h)
        # silently skip unknown kinds; planner enforces vocab

    # v2.3 — populate PowerPoint's notes pane from slide.speaker_notes.
    # `notes_slide` is auto-created on first access by python-pptx; we only
    # write when the planner provided actual text, so slides without notes
    # keep an empty (but valid) notes_slide underneath.
    notes = getattr(slide_node, "speaker_notes", None)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def _render_slide(slide: Any, slide_node: Any, slide_w: int, slide_h: int,
                  ctx: ToolContext) -> None:
    """v2.8.1 dispatcher.

    Layered approach so backgrounds / images / tables stay on the
    default render path even when an archetype claims the text layout:

    1. Render every non-text child via `_add_background` / `_add_picture`
       / `_add_table` (z_index sorted) — archetype renderers only emit
       text shapes, so this preserves figures + cover photos.
    2. Look up the archetype renderer in `archetypes.get_renderer`. When
       a renderer exists, call it (it handles text + speaker_notes +
       section_number prefix). When it returns None — Phase 2/3
       placeholders, missing archetype, OR evidence_snapshot without a
       big-number child — fall through to `_default_render_slide`.

    The default render path remains byte-identical for every pre-v2.8.1
    deck because the schema default `archetype="evidence_snapshot"`
    routes through the placeholder branch when no big-number is
    present.
    """
    # Lazy import to avoid a circular `tools.archetypes._common` →
    # `tools.pptx_renderer` cycle at module load time.
    from .archetypes import get_renderer
    from .archetypes.evidence_snapshot import has_big_number

    archetype = getattr(slide_node, "archetype", None)
    renderer = get_renderer(archetype)
    if renderer is not None and archetype == "evidence_snapshot" \
            and not has_big_number(slide_node):
        renderer = None

    if renderer is None:
        _default_render_slide(slide, slide_node, slide_w, slide_h, ctx)
        return

    # Archetype path — first place non-text children (figures /
    # backgrounds / tables) so the archetype's text shapes layer on
    # top. The archetype itself owns title decoration + speaker_notes.
    role = getattr(slide_node, "role", None) or "content"
    children = list(getattr(slide_node, "children", None) or [])
    children.sort(key=lambda c: int(getattr(c, "z_index", 0) or 0))
    for child in children:
        effective = _with_legacy_bbox(child, role)
        kind = getattr(effective, "kind", None)
        if kind == "background":
            _add_background(slide, effective, slide_w, slide_h)
        elif kind == "image":
            _add_picture(slide, effective, slide_w, slide_h)
        elif kind == "table":
            _add_table(slide, effective, slide_w, slide_h)
        # text + unknown kinds → archetype's responsibility (or skipped)
    renderer(slide_node, slide, slide_w, slide_h, ctx)


def _bbox_to_emu(bbox: Any, slide_w: int, slide_h: int,
                 default_bbox: tuple[int, int, int, int] | None = None
                 ) -> tuple[Emu, Emu, Emu, Emu]:
    """Resolve a bbox (or fallback to slide-sized) to EMU left/top/width/height."""
    if bbox is not None:
        x = int(getattr(bbox, "x", 0) or 0)
        y = int(getattr(bbox, "y", 0) or 0)
        w = int(getattr(bbox, "w", slide_w) or slide_w)
        h = int(getattr(bbox, "h", slide_h) or slide_h)
    elif default_bbox is not None:
        x, y, w, h = default_bbox
    else:
        x, y, w, h = 0, 0, slide_w, slide_h
    # clamp to slide bounds
    x = max(0, min(x, slide_w - 1))
    y = max(0, min(y, slide_h - 1))
    w = max(1, min(w, slide_w - x))
    h = max(1, min(h, slide_h - y))
    return Emu(x * PX_TO_EMU), Emu(y * PX_TO_EMU), Emu(w * PX_TO_EMU), Emu(h * PX_TO_EMU)


def _add_background(slide: Any, node: Any, slide_w: int, slide_h: int) -> None:
    """Full-slide picture background. Planner can pass bbox or leave None to
    cover the whole slide."""
    src = getattr(node, "src_path", None)
    if not src:
        return  # no background image available; PowerPoint default is white
    if not Path(src).exists():
        return
    left, top, width, height = _bbox_to_emu(
        getattr(node, "bbox", None), slide_w, slide_h,
        default_bbox=(0, 0, slide_w, slide_h),
    )
    slide.shapes.add_picture(src, left, top, width=width, height=height)


def _add_callout(
    slide: Any, node: Any,
    placed_anchors: dict[str, tuple[int, int, int, int]],
    slide_w: int, slide_h: int,
) -> None:
    """v2.6 — overlay an annotation shape on top of a sibling picture/table.

    Resolves `node.callout_region` (in slide-pixel coordinates, top-left
    origin) against optional `node.anchor_layer_id` (looked up in
    `placed_anchors`). Renders one of three shape styles:

    - "highlight" → MSO_SHAPE.RECTANGLE with no fill, oxblood outline 2px
    - "circle"    → MSO_SHAPE.OVAL,      no fill, oxblood outline 2px
    - "label"     → text box w/ thin border + cream fill + Inter 12pt text;
                    optional thin connector from label center to region
                    center if `arrow=True`

    All three target the same EMU bbox computed from callout_region. If
    region is None and anchor exists, uses the whole anchor bbox. If
    nothing resolves, silent no-op (defensive — callout shouldn't crash
    a slide that's otherwise fine).
    """
    from pptx.dml.color import RGBColor as _RGBColor
    from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE

    style = (getattr(node, "callout_style", None) or "highlight").lower()
    region = getattr(node, "callout_region", None)
    anchor_id = getattr(node, "anchor_layer_id", None)
    anchor_bbox = placed_anchors.get(anchor_id) if anchor_id else None

    # Compute target EMU bbox.
    if region is not None:
        # callout_region is in slide-pixel coords (top-left origin, same
        # as every other LayerNode bbox in the codebase).
        cx_emu = Emu(int(region.x) * PX_TO_EMU)
        cy_emu = Emu(int(region.y) * PX_TO_EMU)
        cw_emu = Emu(int(region.w) * PX_TO_EMU)
        ch_emu = Emu(int(region.h) * PX_TO_EMU)
    elif anchor_bbox is not None:
        cx_emu, cy_emu, cw_emu, ch_emu = anchor_bbox
    else:
        return  # nothing to render

    accent = _RGBColor(0x7F, 0x1D, 0x1D)
    cream = _RGBColor(0xFA, 0xF7, 0xF0)
    ink = _RGBColor(0x0F, 0x17, 0x2A)

    if style == "highlight":
        rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, cx_emu, cy_emu, cw_emu, ch_emu,
        )
        rect.fill.background()
        rect.line.color.rgb = accent
        rect.line.width = Emu(2 * PX_TO_EMU)
        rect.name = getattr(node, "layer_id", None) or "callout_highlight"

    elif style == "circle":
        oval = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, cx_emu, cy_emu, cw_emu, ch_emu,
        )
        oval.fill.background()
        oval.line.color.rgb = accent
        oval.line.width = Emu(2 * PX_TO_EMU)
        oval.name = getattr(node, "layer_id", None) or "callout_circle"

    elif style == "label":
        text = (getattr(node, "callout_text", None) or "").strip()
        # Label dims — heuristic, ~12pt Inter at ~9px per char + padding.
        label_w_px = max(80, len(text) * 11 + 24)
        label_h_px = 36
        # Try to place to the right of the region; fall back below.
        slide_right_emu = Emu(slide_w * PX_TO_EMU)
        right_emu = cx_emu + cw_emu + Emu(8 * PX_TO_EMU)
        label_w_emu = Emu(label_w_px * PX_TO_EMU)
        label_h_emu = Emu(label_h_px * PX_TO_EMU)
        if right_emu + label_w_emu <= slide_right_emu:
            lx, ly = right_emu, cy_emu
        else:
            lx = cx_emu
            ly = cy_emu + ch_emu + Emu(8 * PX_TO_EMU)
        tb = slide.shapes.add_textbox(lx, ly, label_w_emu, label_h_emu)
        tb.fill.solid()
        tb.fill.fore_color.rgb = cream
        tb.line.color.rgb = accent
        tb.line.width = Emu(1 * PX_TO_EMU)
        tb.name = getattr(node, "layer_id", None) or "callout_label"
        tf = tb.text_frame
        tf.margin_left = Emu(4 * PX_TO_EMU)
        tf.margin_right = Emu(4 * PX_TO_EMU)
        tf.margin_top = Emu(4 * PX_TO_EMU)
        tf.margin_bottom = Emu(4 * PX_TO_EMU)
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.name = "Inter"
        run.font.size = Pt(12)
        run.font.color.rgb = ink

        # Optional arrow: thin connector from region center to label center.
        if getattr(node, "arrow", False):
            region_cx = cx_emu + cw_emu // 2
            region_cy = cy_emu + ch_emu // 2
            label_cx = lx + label_w_emu // 2
            label_cy = ly + label_h_emu // 2
            try:
                conn = slide.shapes.add_connector(
                    MSO_CONNECTOR_TYPE.STRAIGHT,
                    region_cx, region_cy, label_cx, label_cy,
                )
                conn.line.color.rgb = accent
                conn.line.width = Emu(1 * PX_TO_EMU)
            except Exception:
                # Connector not supported on this python-pptx version — skip.
                pass


def _add_picture(slide: Any, node: Any, slide_w: int, slide_h: int) -> None:
    """Place an image into the slide letterbox-fit inside the planner's bbox.

    `python-pptx`'s `add_picture(left, top, width=W, height=H)` force-stretches
    the source to W×H. For paper figures pulled from `ingest_document` the
    source aspect ratio rarely matches the planner's slot bbox, so the
    stretch makes captions / axis labels / equation glyphs unreadable
    (2026-04-25 dogfood feedback). We mirror v1.2.3 poster behavior here:
    compute contain-fit dimensions from the source's real pixel aspect,
    center inside bbox, leave letterbox bands transparent so the slide
    background shows through.
    """
    src = getattr(node, "src_path", None)
    if not src or not Path(src).exists():
        return
    left, top, width, height = _bbox_to_emu(
        getattr(node, "bbox", None), slide_w, slide_h,
    )
    fit_left, fit_top, fit_w, fit_h = _aspect_fit_emu(src, left, top, width, height)
    slide.shapes.add_picture(src, fit_left, fit_top, width=fit_w, height=fit_h)


def _aspect_fit_emu(
    src_path: str,
    bbox_left: int,
    bbox_top: int,
    bbox_width: int,
    bbox_height: int,
) -> tuple[int, int, int, int]:
    """Letterbox-fit `src_path` into the EMU bbox; return (left, top, w, h).

    Falls back to the original bbox if source dims are unreadable so a
    stretched render is still better than a missing image.
    """
    if bbox_width <= 0 or bbox_height <= 0:
        return bbox_left, bbox_top, bbox_width, bbox_height
    try:
        with Image.open(src_path) as im:
            sw, sh = im.size
    except Exception:
        return bbox_left, bbox_top, bbox_width, bbox_height
    if sw <= 0 or sh <= 0:
        return bbox_left, bbox_top, bbox_width, bbox_height
    src_ratio = sw / sh
    bbox_ratio = bbox_width / bbox_height
    if src_ratio > bbox_ratio:
        new_w = bbox_width
        new_h = int(round(bbox_width / src_ratio))
    else:
        new_h = bbox_height
        new_w = int(round(bbox_height * src_ratio))
    new_left = bbox_left + (bbox_width - new_w) // 2
    new_top = bbox_top + (bbox_height - new_h) // 2
    return new_left, new_top, new_w, new_h


def _add_table(slide: Any, node: Any, slide_w: int, slide_h: int) -> None:
    """Render a `kind="table"` layer as a native PowerPoint table.

    Expects `node.rows: list[list[str]]` and optional `node.headers:
    list[str]`. When `headers` is empty, the first row of `rows` is
    promoted to the header. When `rows` is empty we bail out silently
    (planner is expected to not reference empty tables, but defensive).

    Sizing:
    - bbox width/height set the table's outer frame.
    - Row heights are even; header row is slightly taller.
    - Column widths are proportional to the max string length seen in
      that column — rough but prevents one wide column from collapsing
      everything else.
    - Font size auto-shrinks based on row count to keep cells legible
      at deck scale (floor 10pt, ceiling 18pt).

    Optional: `node.col_highlight_rule: list[str]` — per-column "max"
    / "min" / "". When set, the winning row per column is bolded.
    """
    rows = getattr(node, "rows", None) or []
    headers = list(getattr(node, "headers", None) or [])
    col_rule = list(getattr(node, "col_highlight_rule", None) or [])
    if not rows and not headers:
        return

    # Promote first row if headers empty.
    if not headers and rows:
        headers = [str(c) for c in rows[0]]
        rows = rows[1:]

    # Normalize all rows to header width (pad / truncate).
    n_cols = max(len(headers), max((len(r) for r in rows), default=0))
    if n_cols == 0:
        return
    headers = [str(h) for h in headers] + [""] * (n_cols - len(headers))
    headers = headers[:n_cols]
    rows = [
        [str(c) for c in r] + [""] * (n_cols - len(r))
        for r in rows
    ]
    rows = [r[:n_cols] for r in rows]

    # v2.7 — wide-table safety net. Planner.md rule #7 asks the planner
    # to subset wide tables down to 4-6 cols (deck cells go illegible
    # past ~8). Two consecutive 2026-04-25 dogfoods produced 12-15 col
    # tables anyway. Cap rendering at 8 cols, keep first 6, append a
    # marker so the audience knows columns were dropped. Loud log so
    # reviews can flag the designer.
    _WIDE_CAP = 8
    _WIDE_KEEP = 6
    if n_cols > _WIDE_CAP:
        from ..util.logging import log as _wlog
        original_cols = n_cols
        _wlog("pptx.table.truncate",
              layer_id=getattr(node, "layer_id", "?"),
              original_cols=original_cols, kept_cols=_WIDE_KEEP)
        headers = headers[:_WIDE_KEEP]
        rows = [r[:_WIDE_KEEP] for r in rows]
        n_cols = _WIDE_KEEP
        # Stuff a marker into the caption so the slide carries evidence.
        marker = (f" [Truncated: showing {_WIDE_KEEP}/{original_cols} "
                  f"cols — see paper for full table]")
        try:
            cur_cap = (getattr(node, "caption", None) or "").strip()
            node.caption = (cur_cap + marker).strip() if cur_cap else marker.strip()
        except (AttributeError, TypeError):
            pass  # frozen / proxy — best-effort

    # Normalize rule list length.
    if col_rule:
        col_rule = col_rule[:n_cols] + [""] * max(0, n_cols - len(col_rule))

    # Winner map for bold highlighting. Import here to avoid cycles.
    from ..util.table_png import _compute_winner_rows
    winner_rows = _compute_winner_rows(rows, col_rule)

    n_rows = len(rows) + 1  # +1 for header
    left, top, width, height = _bbox_to_emu(
        getattr(node, "bbox", None), slide_w, slide_h,
    )

    shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = shape.table

    # Column widths proportional to max string length in column.
    total_w = sum(col.width for col in table.columns)
    col_weights = []
    for c in range(n_cols):
        cells = [headers[c]] + [row[c] for row in rows]
        max_len = max((len(str(v)) for v in cells), default=1)
        col_weights.append(max(1, min(max_len, 30)))
    total_weight = sum(col_weights)
    for c, col in enumerate(table.columns):
        col.width = Emu(int(total_w * col_weights[c] / total_weight))

    # Font-size autoscale: smaller when the table has many rows.
    body_pt = 18 if n_rows <= 6 else 14 if n_rows <= 12 else 11
    header_pt = body_pt + 1

    _fill_table_row(table, 0, headers, font_pt=header_pt,
                    is_header=True, winner_cols=set())
    for r_idx, row in enumerate(rows, start=1):
        # Data row idx in the rows list is r_idx - 1.
        winning_cols = {c for c, win_r in winner_rows.items()
                        if win_r == r_idx - 1}
        _fill_table_row(table, r_idx, row, font_pt=body_pt,
                        is_header=False, winner_cols=winning_cols)


def _fill_table_row(table: Any, r: int, values: list[str],
                    *, font_pt: int, is_header: bool,
                    winner_cols: set[int] | None = None) -> None:
    winner_cols = winner_cols or set()
    for c, val in enumerate(values):
        cell = table.cell(r, c)
        cell.text = ""  # clear default
        tf = cell.text_frame
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER if is_header else PP_ALIGN.LEFT
        run = para.add_run()
        run.text = val
        font = run.font
        font.size = Pt(font_pt)
        # Bold for header row OR for the winning data cell per column.
        font.bold = is_header or (c in winner_cols)
        if is_header:
            font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            # Fill header cell with a dark accent (python-pptx solid-fill API).
            from pptx.oxml.ns import qn
            tcPr = cell._tc.get_or_add_tcPr()
            for existing in tcPr.findall(qn("a:solidFill")):
                tcPr.remove(existing)
            from lxml import etree
            fill = etree.SubElement(tcPr, qn("a:solidFill"))
            etree.SubElement(fill, qn("a:srgbClr"), val="1F2A44")


def _add_text_frame(slide: Any, node: Any, slide_w: int, slide_h: int) -> None:
    text = (getattr(node, "text", None) or "").strip()
    if not text:
        return
    left, top, width, height = _bbox_to_emu(
        getattr(node, "bbox", None), slide_w, slide_h,
    )
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True

    # First paragraph holds the initial run; subsequent lines become new paras.
    lines = text.splitlines() or [text]
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.text = ""  # clear any default run
        align = getattr(node, "align", None)
        if align in _ALIGN_MAP:
            para.alignment = _ALIGN_MAP[align]
        run = para.add_run()
        run.text = line
        font = run.font
        # python-pptx wants pts; we get px. pt ≈ px * 72/96 = px * 0.75.
        size_px = int(getattr(node, "font_size_px", None) or 36)
        font.size = Pt(max(6, round(size_px * 0.75)))
        family = getattr(node, "font_family", None)
        if family:
            font.name = family
        effects = getattr(node, "effects", None)
        fill_hex = None
        if effects is not None:
            fill_hex = getattr(effects, "fill", None)
        if fill_hex and isinstance(fill_hex, str) and fill_hex.startswith("#"):
            rgb = _hex_to_rgb(fill_hex)
            if rgb is not None:
                font.color.rgb = RGBColor(*rgb)


def _hex_to_rgb(hx: str) -> tuple[int, int, int] | None:
    s = hx.lstrip("#")
    if len(s) == 3:
        s = "".join(c * 2 for c in s)
    if len(s) != 6:
        return None
    try:
        return int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16)
    except ValueError:
        return None


def render_slide_preview_png(slide_node: Any, slide_w: int, slide_h: int,
                             out_path: Path, ctx: ToolContext,
                             *, hide_editable: bool = False) -> None:
    """Pillow-render a simplified preview of one slide.

    Not pixel-accurate with PowerPoint's renderer — it's an at-a-glance thumb
    for chat UX + for stitching into the grid preview. Shows the bg image (if
    any), image shapes (as resized thumbs), and text (approximated font).
    ``hide_editable`` skips text/table layers so the result can be used as a
    hybrid-PPTX visual base under native editable overlays.
    """
    # Downscale to a chat-friendly size while preserving aspect.
    max_w = 960
    scale = min(1.0, max_w / slide_w)
    w = max(1, int(slide_w * scale))
    h = max(1, int(slide_h * scale))

    img = Image.new("RGB", (w, h), (255, 255, 255))

    children = sorted(
        list(getattr(slide_node, "children", None) or []),
        key=lambda c: int(getattr(c, "z_index", 0) or 0),
    )
    role = getattr(slide_node, "role", None) or "content"
    for child in children:
        effective = _with_legacy_bbox(child, role)
        kind = getattr(effective, "kind", None)
        if kind in ("background", "image") or (kind == "table" and not hide_editable):
            # Tables use their pre-rendered src_path (PIL-drawn PNG) —
            # ingest_document baked it. PPTX itself holds a live table
            # shape; this preview just needs a raster for the thumbnail.
            _paste_image(img, effective, slide_w, slide_h, scale)
        elif kind == "text" and not hide_editable:
            _draw_text(img, effective, slide_w, slide_h, scale, ctx)

    img.save(out_path, format="PNG", optimize=True)


def _scaled_bbox(bbox: Any, slide_w: int, slide_h: int,
                 scale: float,
                 default_full: bool = False,
                 ) -> tuple[int, int, int, int]:
    if bbox is not None:
        x = int(getattr(bbox, "x", 0) or 0)
        y = int(getattr(bbox, "y", 0) or 0)
        w = int(getattr(bbox, "w", slide_w) or slide_w)
        h = int(getattr(bbox, "h", slide_h) or slide_h)
    elif default_full:
        x, y, w, h = 0, 0, slide_w, slide_h
    else:
        x, y, w, h = 0, 0, slide_w // 2, slide_h // 4
    return (
        max(0, int(x * scale)),
        max(0, int(y * scale)),
        max(1, int(w * scale)),
        max(1, int(h * scale)),
    )


def _paste_image(canvas: Image.Image, node: Any, slide_w: int, slide_h: int,
                 scale: float) -> None:
    src = getattr(node, "src_path", None)
    if not src or not Path(src).exists():
        return
    try:
        tile = Image.open(src).convert("RGBA")
    except Exception:
        return
    kind = getattr(node, "kind", None)
    default_full = kind == "background"
    sx, sy, sw, sh = _scaled_bbox(
        getattr(node, "bbox", None), slide_w, slide_h, scale,
        default_full=default_full,
    )
    if kind == "image" and tile.size != (sw, sh):
        # Letterbox-fit content figures so the preview matches the PPTX
        # render path (same v1.2.3-style aspect-preserve as poster/SVG).
        # Backgrounds keep cover-fit (force-resize) since the cover is
        # always full-bleed by design and any minor seedream aspect drift
        # is better cropped than letterboxed with white bars on a slide.
        src_w, src_h = tile.size
        if src_w > 0 and src_h > 0 and sw > 0 and sh > 0:
            src_ratio = src_w / src_h
            bbox_ratio = sw / sh
            if src_ratio > bbox_ratio:
                new_w = sw
                new_h = max(1, int(round(sw / src_ratio)))
            else:
                new_h = sh
                new_w = max(1, int(round(sh * src_ratio)))
            tile = tile.resize((new_w, new_h), Image.LANCZOS)
            sx += (sw - new_w) // 2
            sy += (sh - new_h) // 2
    else:
        tile = tile.resize((sw, sh), Image.LANCZOS)
    canvas.alpha_composite(tile, dest=(sx, sy)) if canvas.mode == "RGBA" else canvas.paste(tile, (sx, sy), tile)


def _draw_text(canvas: Image.Image, node: Any, slide_w: int, slide_h: int,
               scale: float, ctx: ToolContext) -> None:
    text = (getattr(node, "text", None) or "").strip()
    if not text:
        return
    sx, sy, sw, sh = _scaled_bbox(
        getattr(node, "bbox", None), slide_w, slide_h, scale,
        default_full=False,
    )

    # Pick a reasonable approximated font size from the node metadata.
    size_px = int(getattr(node, "font_size_px", None) or 36)
    approx = max(10, min(120, int(size_px * scale)))
    family = getattr(node, "font_family", None) or ctx.settings.default_text_font
    if _contains_cjk(text) and family in {"PlayfairDisplay", "Inter"}:
        family = "NotoSerifSC-Bold" if _is_title_child(node) else "NotoSansSC-Bold"

    fonts = ctx.settings.fonts
    fname = fonts.get(family) or fonts[ctx.settings.default_text_font]
    try:
        font = ImageFont.truetype(str(ctx.settings.fonts_dir / fname), size=approx)
    except Exception:
        font = ImageFont.load_default()

    effects = getattr(node, "effects", None)
    fill_hex = getattr(effects, "fill", None) if effects is not None else None
    rgb = _hex_to_rgb(fill_hex) if isinstance(fill_hex, str) else None
    fill = rgb or (15, 23, 42)

    draw = ImageDraw.Draw(canvas)
    # Word-wrap to bbox: simple char-by-char wrap (works for CJK; latin uses spaces).
    lines = _wrap_for_width(text, font, sw, draw)
    line_h = approx + 6
    y = sy
    for line in lines:
        if y + line_h > sy + sh:
            break
        draw.text((sx, y), line, fill=fill, font=font)
        y += line_h


def _contains_cjk(text: str) -> bool:
    return any("\u4e00" <= ch <= "\u9fff" for ch in text)


def _wrap_for_width(text: str, font: Any, max_w: int, draw: Any) -> list[str]:
    """Simple word/char wrap to fit a pixel width."""
    out: list[str] = []
    for paragraph in text.splitlines() or [text]:
        if " " in paragraph:
            # space-delimited: greedy word-wrap
            words = paragraph.split()
            line = ""
            for w in words:
                probe = (line + " " + w).strip()
                if _measure(probe, font, draw) <= max_w:
                    line = probe
                else:
                    if line:
                        out.append(line)
                    line = w
            if line:
                out.append(line)
        else:
            # CJK-style: char-by-char
            line = ""
            for ch in paragraph:
                probe = line + ch
                if _measure(probe, font, draw) <= max_w:
                    line = probe
                else:
                    if line:
                        out.append(line)
                    line = ch
            if line:
                out.append(line)
    return out


def _measure(s: str, font: Any, draw: Any) -> int:
    try:
        bbox = draw.textbbox((0, 0), s, font=font)
        return int(bbox[2] - bbox[0])
    except Exception:
        return len(s) * 10
