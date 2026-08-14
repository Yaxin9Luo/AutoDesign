#!/usr/bin/env python3
"""No-API regression checks for reference-poster style extraction."""

from __future__ import annotations

import json
import shutil
import sys
import tempfile
from pathlib import Path
from types import SimpleNamespace

import fitz
from bs4 import BeautifulSoup
from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.util import Inches, Pt

from autodesign.agents.external_designer_author import (
    ExternalDesignerAuthor,
    _active_color_system,
    _aesthetic_contract_lines,
    _author_visible_brief,
    _author_must_read_first,
    _best_available_artifact_fallback_acceptance,
    _best_candidate_fallback_acceptance,
    _build_repair_context,
    _color_system_brief_lines,
    _formula_authoring_contract,
    _format_repair_prompt_block,
    _identity_header_authoring_contract,
    _reference_layout_authoring_contract,
    _reference_style_contract,
    _required_color_system,
    _reference_style_brief_lines,
    _stage_reference_style_inputs,
    _synchronize_staged_color_system,
    _table_authoring_contract,
    _typography_authoring_contract,
)
from autodesign.agents.reference_style_agent import (
    ReferenceStyleAgentError,
    _compile_reference_style_contract,
    _normalized_header,
    _normalized_section_heading,
    _normalized_typography,
    _reference_style_repair_prompt,
    _validate_raw_reference_style_blueprint,
    _validate_rendered_style_tokens,
    prepare_reference_style_contract,
)
from autodesign.tools import ToolContext
from autodesign.tools.propose_paper_poster_html import (
    _active_paper_typography_contract,
    _authored_palette_diagnostics,
    _authored_reference_style_diagnostics,
    _ensure_dom_block_ids,
    _measure_dom_bboxes,
    _paper_poster_typography_contract_error,
    _reference_style_contract_error,
    _typography_family_matches,
)
from autodesign.util.reference_poster import normalize_reference_poster


def main() -> int:
    with tempfile.TemporaryDirectory(prefix="reference-poster-smoke-") as raw_tmp:
        root = Path(raw_tmp)
        fixtures = _write_fixtures(root)
        for name, path in fixtures.items():
            output = root / f"normalized_{name}"
            metadata = normalize_reference_poster(path, output)
            assert (output / "reference.png").exists(), name
            assert metadata["preview_width_px"] > 0 and metadata["preview_height_px"] > 0

        run_dir = root / "run"
        ctx = ToolContext(
            settings=SimpleNamespace(
                harness_api_key=None,
                poster_harness_mode="dogfood",
                repo_root=Path.cwd(),
                skills_dir=Path.cwd() / "skills",
            ),
            run_dir=run_dir,
            layers_dir=run_dir / "layers",
            run_id="reference-style-smoke",
        )
        run_dir.mkdir(parents=True)
        ctx.layers_dir.mkdir(parents=True)
        reference_dir = run_dir / "reference_poster"
        reference_dir.mkdir(parents=True)
        (reference_dir / "reference_style_agent_review.json").write_text(
            json.dumps({
                "status": "ok",
                "rendered_blueprint_inspected": True,
                "header_matches_reference": True,
                "body_region_geometry_matches_reference": True,
                "chrome_avoids_content": True,
                "blueprint_sha256": "stale",
            }),
            encoding="utf-8",
        )
        ctx.state["attachments"] = [str(root / "paper.pdf")]
        harness = _write_fake_harness(root)
        contract = prepare_reference_style_contract(
            ctx,
            fixtures["png"],
            command=f"{sys.executable} {harness}",
            harness="custom",
            timeout_s=30,
        )
        style_id = contract["style_reference_id"]
        assert style_id.startswith("reference_")
        assert contract["content_transfer_forbidden"] is True
        assert contract["color_system"]["palette_id"] == style_id
        assert contract["version"] == 4
        assert contract["body_region_schema_version"] == 1
        assert contract["sanitizer_version"] == 4
        assert contract["extraction_skill"] == "reference_style_agent_skill.md"
        assert len(contract["extraction_skill_sha256"]) == 64
        assert len(contract["extraction_prompt_schema_sha256"]) == 64
        assert contract["extraction_attempt_count"] == 2
        assert (reference_dir / "reference_style_attempts" / "attempt_01" / "failure.json").exists()
        assert (reference_dir / "reference_style_agent_skill.md").exists()
        audit = json.loads((run_dir / "reference_style_audit.json").read_text(encoding="utf-8"))
        assert audit["status"] == "pass"
        assert audit["checks"]["skill_bundle_binding"] is True
        assert audit["checks"]["skill_resources_binding"] is True
        assert audit["checks"]["all_top_level_sections_owned"] is True
        assert audit["checks"]["forbidden_pipeline_artifacts_absent"] is True
        process_path = reference_dir / "reference_style_agent_process.json"
        process_mtime = process_path.stat().st_mtime_ns
        cached_contract = prepare_reference_style_contract(
            ctx,
            fixtures["png"],
            command=f"{sys.executable} {harness}",
            harness="custom",
            timeout_s=30,
        )
        assert cached_contract["blueprint"]["sha256"] == contract["blueprint"]["sha256"]
        assert process_path.stat().st_mtime_ns == process_mtime
        assert _normalized_header({
            "mode": "tinted_open", "background_role": "secondary",
            "rule_placement": "bottom", "rule_color_role": "primary", "rule_width_px": 8,
        })["rule_placement"] == "bottom"
        outlined = _normalized_section_heading({
            "mode": "outlined_band", "fill_role": "secondary",
            "border_role": "primary", "border_width_px": 2, "corner_style": "capsule",
        })
        assert outlined["mode"] == "outlined_band" and outlined["corner_style"] == "capsule"
        mixed_type = _normalized_typography({
            "display_family_category": "serif", "body_family_category": "sans_serif"
        })
        assert mixed_type["display_family_category"] == "serif"
        assert mixed_type["body_family_category"] == "sans_serif"
        multicolor_analysis = json.loads(
            (reference_dir / "reference_style_analysis.json").read_text(encoding="utf-8")
        )
        multicolor_analysis["palette"]["additional_roles"] = {
            "panel_accent_green": "#A8C61E",
        }
        multicolor_analysis["palette"]["panel_accent_pink"] = "#FF6EB4"
        multicolor_contract = _compile_reference_style_contract(
            multicolor_analysis,
            json.loads(
                (reference_dir / "reference_source_metadata.json").read_text(encoding="utf-8")
            ),
        )
        assert multicolor_contract["color_system"]["roles"]["panel_accent_green"] == "#A8C61E"
        assert multicolor_contract["color_system"]["roles"]["panel_accent_pink"] == "#FF6EB4"
        assert "#A8C61E" in multicolor_contract["color_system"]["allowed_hexes"]
        assert "#FF6EB4" in multicolor_contract["color_system"]["allowed_hexes"]
        assert (
            multicolor_contract["color_system"]["css_variables"]["--poster-reference-panel-accent-green"]
            == "#A8C61E"
        )
        repair_root = root / "repair_prompt"
        repair_runtime_skill = {
            "resources": [{
                "id": "output_contract_v4",
                "path": "references/output_contract_v4.md",
            }],
        }
        prior_archive = repair_root / "reference_style_attempts" / "attempt_02"
        prior_archive.mkdir(parents=True)
        (prior_archive / "reference_style_analysis.json").write_text("{}", encoding="utf-8")
        (prior_archive / "reference_style_blueprint.html").write_text("", encoding="utf-8")
        repair_prompt = _reference_style_repair_prompt(
            repair_root,
            attempt_index=3,
            failure="rendered outlined section heading is missing its full border",
            failures=[
                "raw reference style blueprint uses colors outside its analysis palette: ['#A8C61E']",
                "rendered outlined section heading is missing its full border",
            ],
            runtime_skill=repair_runtime_skill,
        )
        assert "All deterministic failures seen so far" in repair_prompt
        assert "outside its analysis palette" in repair_prompt
        assert "missing its full border" in repair_prompt
        assert "Preserve the existing body-region decomposition" in repair_prompt
        assert "Do not regress earlier fixes" in repair_prompt
        assert "Begin from the archived prior files" in repair_prompt
        missing_baseline_prompt = _reference_style_repair_prompt(
            root / "missing_repair_prompt",
            attempt_index=2,
            failure="reference style agent did not write a valid reference_style_analysis.json",
            failures=[],
            runtime_skill=repair_runtime_skill,
        )
        assert "did not produce a complete patch baseline" in missing_baseline_prompt
        assert "Reconstruct only the missing files" in missing_baseline_prompt
        assert "Do not regenerate either file" not in missing_baseline_prompt
        chrome_repair_prompt = _reference_style_repair_prompt(
            Path("/tmp/reference-style-repair"),
            attempt_index=4,
            failure="reference style extraction audit failed: chrome_presence",
            failures=["reference style extraction audit failed: chrome_presence"],
            runtime_skill=repair_runtime_skill,
        )
        assert "root-level `chrome-layer`" in chrome_repair_prompt
        assert "inside `identity-header`" in chrome_repair_prompt
        assert "Do not change the body-region map" in chrome_repair_prompt
        computed_snapshot = {
            "header": {
                "background_color": "rgb(255, 255, 255)",
                "border_top_width_px": 0,
                "border_bottom_width_px": 0,
            },
            "title": {"font_family": "Arial, sans-serif"},
            "heading": {
                "background_color": "rgb(255, 255, 255)",
                "font_family": "Arial, sans-serif",
                "border_top_width_px": 0,
                "border_right_width_px": 0,
                "border_bottom_width_px": 3,
                "border_left_width_px": 0,
                "border_radius_px": 0,
                "height_px": 50,
            },
            "body": {"font_family": "Arial, sans-serif"},
        }
        _validate_rendered_style_tokens(computed_snapshot, contract)
        bad_computed = json.loads(json.dumps(computed_snapshot))
        bad_computed["heading"]["border_bottom_width_px"] = 0
        try:
            _validate_rendered_style_tokens(bad_computed, contract)
        except ReferenceStyleAgentError:
            pass
        else:
            raise AssertionError("computed section-heading mismatch must fail")
        try:
            _compile_reference_style_contract(
                {"version": 3, "column_structure": {"major_sections_per_column": [1, 1, 1]}},
                {"source_sha256": "a" * 64},
            )
        except ReferenceStyleAgentError:
            pass
        else:
            raise AssertionError("legacy column-only reference analysis must be rejected")
        invalid_v4 = json.loads(
            (reference_dir / "reference_style_analysis.json").read_text(encoding="utf-8")
        )
        invalid_v4["body_region_structure"]["region_count"] += 1
        try:
            _compile_reference_style_contract(
                invalid_v4,
                json.loads(
                    (reference_dir / "reference_source_metadata.json").read_text(encoding="utf-8")
                ),
            )
        except ReferenceStyleAgentError:
            pass
        else:
            raise AssertionError("inconsistent version-4 region counts must be rejected")
        oversized_v4 = json.loads(
            (reference_dir / "reference_style_analysis.json").read_text(encoding="utf-8")
        )
        oversized_v4["body_region_structure"]["regions"] = [
            {
                "region_id": f"region_{index}",
                "region_role": "column",
                "section_count": 1,
                "reading_order": index,
            }
            for index in range(1, 8)
        ]
        oversized_v4["body_region_structure"]["region_count"] = 7
        oversized_v4["body_region_structure"]["major_section_count"] = 7
        oversized_v4["body_region_structure"]["major_sections_per_region"] = [1] * 7
        try:
            _compile_reference_style_contract(
                oversized_v4,
                json.loads(
                    (reference_dir / "reference_source_metadata.json").read_text(encoding="utf-8")
                ),
            )
        except ReferenceStyleAgentError as exc:
            assert "two to six" in str(exc)
        else:
            raise AssertionError("reference region count beyond downstream contract must be rejected")
        assert contract["typography_contract"]["subsection_heading_font_size_px"] == 24
        assert contract["transfer_mode"] == "reference_first_reconstruction"
        assert contract["style_tokens"]["header_treatment"]["alignment"] == "left"
        assert contract["style_tokens"]["header_treatment"]["top_rule"] == "none"
        assert contract["style_tokens"]["lead_band"]["present"] is True
        region_structure = contract["style_tokens"]["body_region_structure"]
        assert region_structure["layout_mode"] == "freeform_regions"
        assert region_structure["major_sections_per_region"] == [1, 1, 1, 1, 1]
        assert [region["region_role"] for region in region_structure["regions"]] == [
            "column",
            "column",
            "column",
            "side_callout",
            "footer_band",
        ]
        assert contract["style_tokens"]["column_structure"]["major_sections_per_column"] == [1, 1, 1, 1, 1]
        assert contract["style_tokens"]["layout_rhythm"]["region_proportions"] == [0.7, 1.1, 1.4, 0.8, 1.8]
        assert all(
            box["x_pct"] + box["w_pct"] <= 100
            and box["y_pct"] + box["h_pct"] <= 100
            for box in contract["style_tokens"]["layout_rhythm"]["region_boxes"]
        )
        assert contract["style_tokens"]["chrome_treatment"]["present"] is True
        sanitized_blueprint = (run_dir / "reference_style_blueprint.html").read_text(
            encoding="utf-8"
        )
        assert ".reference-rail::before" in sanitized_blueprint
        assert "content:''" in sanitized_blueprint
        assert "COPIED CHROME LABEL" not in sanitized_blueprint
        assert contract["blueprint_review"]["chrome_avoids_content"] is True
        assert contract["blueprint_review"]["blueprint_sha256"] != "stale"
        assert (run_dir / "reference_style_blueprint_preview.png").exists()
        measured_boxes = contract["style_tokens"]["layout_rhythm"]["region_boxes"]
        assert measured_boxes[-1]["x_pct"] < 90
        assert measured_boxes[-1]["w_pct"] > 15
        assert contract["style_tokens"]["table_treatment"]["rule_style"] == "none"
        assert contract["style_tokens"]["formula_treatment"]["frame"] == "none"
        assert contract["style_tokens"]["section_structure"]["inter_section_dividers"] == "none"
        assert contract["typography_contract"]["family_category"] == "sans_serif"
        assert contract["color_system"]["roles"]["header_text"] == "#7A1F5C"
        blueprint = (run_dir / "reference_style_blueprint.html").read_text(encoding="utf-8")
        assert "{{PAPER_TITLE}}" in blueprint
        assert "COPIED REFERENCE TITLE" not in blueprint
        assert "<img" not in blueprint.lower()
        assert "javascript:" not in blueprint.lower()
        assert "justify-content:flex-start" in blueprint.replace(" ", "")
        duplicate_identity = root / "duplicate_identity_blueprint.html"
        duplicate_identity.write_text(
            blueprint.replace("{{AUTHORS}}", "{{AUTHORS}}{{AUTHORS}}", 1),
            encoding="utf-8",
        )
        try:
            _validate_raw_reference_style_blueprint(duplicate_identity, contract)
        except ReferenceStyleAgentError as exc:
            assert "exactly once" in str(exc)
        else:
            raise AssertionError("duplicate identity placeholders must be rejected")
        inline_color = root / "inline_color_blueprint.html"
        inline_color.write_text(
            blueprint.replace(
                'data-style-role="identity-header"',
                'data-style-role="identity-header" style="color:#123456"',
                1,
            ),
            encoding="utf-8",
        )
        try:
            _validate_raw_reference_style_blueprint(inline_color, contract)
        except ReferenceStyleAgentError as exc:
            assert "outside its analysis palette" in str(exc)
        else:
            raise AssertionError("off-palette inline style must be rejected")
        functional_color = root / "functional_color_blueprint.html"
        functional_color.write_text(
            blueprint.replace(
                "</style>",
                ".bad{--bad:rgb(1,2,3);color:var(--bad)}</style>",
                1,
            ),
            encoding="utf-8",
        )
        try:
            _validate_raw_reference_style_blueprint(functional_color, contract)
        except ReferenceStyleAgentError as exc:
            assert "non-hex color" in str(exc)
        else:
            raise AssertionError("functional off-palette CSS color must be rejected")
        blueprint_soup = BeautifulSoup(blueprint, "html.parser")
        assert blueprint_soup.select_one(".paper-poster.poster-section") is None
        assert len(blueprint_soup.select(".poster-column")) == 5
        assert [
            len(column.find_all(class_="poster-section", recursive=False))
            for column in blueprint_soup.select(".poster-column")
        ] == [1, 1, 1, 1, 1]
        assert blueprint_soup.select_one(
            '.poster-column[data-region-role="footer_band"]'
        ) is not None
        assert all(
            section.find_parent(attrs={"data-style-role": "body-region"}) is not None
            for section in blueprint_soup.select('[data-style-role="section"]')
        )
        assert blueprint_soup.select_one(
            ".paper-poster > .reference-chrome[data-style-role='chrome-layer']"
        ) is not None
        assert blueprint_soup.select_one(".poster-header[data-panel-role='identity_header']") is not None
        assert blueprint_soup.select_one(
            ".poster-body-regions[data-layout-region='reference_body_regions']"
        ) is not None
        assert blueprint_soup.select_one(
            f'.paper-poster[data-reference-style-id="{style_id}"][data-major-section-count="5"]'
        ) is not None
        assert blueprint_soup.select_one(
            '.paper-poster[data-reference-layout-mode="freeform_regions"][data-reference-region-count="5"]'
        ) is not None
        assert "border-top:12px" not in blueprint.replace(" ", "")
        compact_blueprint = blueprint.replace(" ", "")
        assert ".table-slot{border-top" not in compact_blueprint
        assert ".formula-slot{border-top" not in compact_blueprint
        assert ".poster-section::before" not in compact_blueprint
        assert "section.chrome::after" not in compact_blueprint
        assert '[data-style-role="section-heading"]{border-bottom:3px' in compact_blueprint
        assert ctx.state["attachments"] == [str(root / "paper.pdf")]
        assert ctx.state["rendered_layers"] == {}
        author_prompt = ExternalDesignerAuthor(ctx.settings, "")._build_prompt(
            ctx,
            brief="Generate a paper poster from the target paper.",
            attempt_dir=run_dir,
        )
        assert "with three target-paper columns" not in author_prompt

        active = _active_color_system(ctx, {}, {}, "Generate a paper poster")
        required = _required_color_system(ctx, {}, {}, "Generate a paper poster")
        assert active["palette_id"] == style_id
        assert required["palette_id"] == style_id
        ctx.state["raw_user_brief"] = "Use Teal Coral for this poster."
        explicit = _active_color_system(ctx, {}, {}, ctx.state["raw_user_brief"])
        assert explicit["palette_id"] == "teal_coral"
        ctx.state.pop("raw_user_brief")
        ctx.state["raw_user_brief"] = "Use the reference layout from a blank canvas."
        assert _author_visible_brief(ctx, "Times New Roman default skill text") == ctx.state["raw_user_brief"]
        ctx.state.pop("raw_user_brief")

        attributes = contract["required_root_attributes"]
        body = (
            '<main class="paper-poster" '
            + " ".join(f'{key}="{value}"' for key, value in attributes.items())
            + f' data-palette-id="{style_id}"></main>'
        )
        css = ".paper-poster{" + ";".join(
            f"{key}:{value}" for key, value in contract["color_system"]["css_variables"].items()
        ) + "}"
        assert _authored_palette_diagnostics(body, css, contract["color_system"]) == []
        body = body.replace(
            "</main>",
            '<div class="reference-lead-band" data-style-role="lead-band"></div></main>',
        )
        assert _authored_reference_style_diagnostics(body, css, contract) == []
        active_typography = _active_paper_typography_contract(ctx)
        assert active_typography["primary_font_family"] == "Arial"
        assert _typography_family_matches("Arial, sans-serif", "sans_serif") is True
        assert _typography_family_matches('"Times New Roman", serif', "sans_serif") is False
        leakage = _authored_reference_style_diagnostics(
            body,
            '.paper-poster{font-family:"Times New Roman"}.poster-section{border-top:2px solid #7A1F5C}.lead-key{border-left:3px solid #7A1F5C}',
            contract,
        )
        leakage_ids = {item["issue_id"] for item in leakage}
        assert "paper_poster_html_reference_default_typography_leakage" in leakage_ids
        assert "paper_poster_html_reference_section_divider_leakage" in leakage_ids
        assert "paper_poster_html_reference_vertical_rule_leakage" in leakage_ids
        reference_feedback = _check_reference_style_hard_gate(ctx, contract)
        _check_reference_hard_issue_blocks_accepted_fallback(ctx, root, reference_feedback)
        reference_feedback["summary"]["local_repair_hint"] = "Use one to three sections per column."
        repair_dir = root / "reference_repair"
        repair_dir.mkdir()
        repair_context = _build_repair_context(ctx, repair_dir, reference_feedback)
        assert repair_context["classification"] == "reference_style_failure"
        assert repair_context["reference_style"]["style_reference_id"] == style_id
        assert "exact major-section counts [1, 1, 1, 1, 1]" in repair_context["hint"]
        assert "one to three sections" not in repair_context["hint"].lower()
        assert "structured issues and repair_scope" in repair_context["local_repair_hint"]
        reference_feedback["summary"]["repair_context"] = repair_context
        repair_prompt = _format_repair_prompt_block(
            reference_feedback,
            attempt_index=2,
            max_attempts=3,
            reference_style=contract,
        )
        assert "reference_style_failure" in repair_prompt
        assert "Do not restore the normal AutoDesign skin" in repair_prompt
        assert "exact major-section counts [1, 1, 1, 1, 1]" in repair_prompt
        assert "one to three sections per column" not in repair_prompt
        for classification in (
            "row_allocation_failure",
            "section_content_overflow",
            "command_no_output",
        ):
            scoped_prompt = _format_repair_prompt_block(
                {
                    "summary": {
                        "issue_id": "paper_poster_html_validation_failed",
                        "repair_context": {"classification": classification},
                    }
                },
                attempt_index=2,
                max_attempts=3,
                reference_style=contract,
            )
            assert "three-column" not in scoped_prompt
            assert "all three column" not in scoped_prompt

        brief_lines = _reference_style_brief_lines(contract)
        assert any("REFERENCE-FIRST RECONSTRUCTION MODE" in line for line in brief_lines)
        assert any("absence" in line for line in brief_lines)
        identity_prompt = " ".join(_identity_header_authoring_contract(ctx)).lower()
        assert "left" in identity_prompt
        assert "centered text rows only" not in identity_prompt
        assert "no top rule" in identity_prompt
        typography_prompt = _typography_authoring_contract(ctx)
        assert "Arial" in typography_prompt
        assert "replaces the normal default typography contract" in typography_prompt
        assert "[1, 1, 1, 1, 1]" in _reference_layout_authoring_contract(ctx)
        assert "freeform_regions" in _reference_layout_authoring_contract(ctx)
        visual_treatment = "\n".join(_aesthetic_contract_lines({}, contract))
        assert "wrapper borders transparent" in visual_treatment
        assert "booktabs" not in _table_authoring_contract(ctx).lower()
        assert "no top/bottom separator rules" in _formula_authoring_contract(ctx)
        color_prompt = " ".join(_color_system_brief_lines(
            contract["color_system"],
            reference_style=contract,
        )).lower()
        assert "single top accent rule only" not in color_prompt
        assert "filled section heading bands" not in color_prompt

        attempt_dir = root / "attempt"
        attempt_dir.mkdir()
        shutil.copy2(run_dir / "reference_style_contract.json", attempt_dir / "reference_style_contract.json")
        staged = _stage_reference_style_inputs(ctx, attempt_dir)
        assert "reference_style_contract.json" in staged
        assert "reference_style_blueprint.html" in staged
        assert "reference_poster/reference.png" in staged
        assert not list((attempt_dir / "reference_poster").glob("reference_source.*"))
        assert not (attempt_dir / "reference_poster/reference_style_analysis.json").exists()
        assert not (attempt_dir / "reference_poster/reference_source_metadata.json").exists()
        must_read = _author_must_read_first([], staged + ["poster_content_brief.json"])
        assert must_read[:3] == [
            "reference_style_contract.json",
            "reference_style_blueprint.html",
            "reference_poster/reference.png",
        ]
        repair_must_read = _author_must_read_first(
            ["repair_context.json", "visual_repair_packet.md"],
            staged + ["author_quick_brief.md"],
        )
        assert repair_must_read[:2] == ["repair_context.json", "visual_repair_packet.md"]
        content_path = attempt_dir / "poster_content_brief.json"
        contract_path = attempt_dir / "poster_plan_contract.json"
        content_path.write_text(
            json.dumps({
                "typography_contract": {"font_family": '"Times New Roman", serif'},
                "sections": [{"section_id": "motivation", "title": "Motivation"}],
                "panel_plan": [{"column_id": "left_story", "section_targets": ["Motivation"]}],
                "editorial_column_plan": [{"column_id": "left_story", "section_targets": ["Motivation"]}],
                "background_contract": {"structure": "single top accent rule and filled section bands"},
                "reference_archetype_skeleton": {
                    "hard_constraints": {"min_sections_total": 7, "target_sections_total": 9},
                },
                "native_reference_targets": {"min_sections_total": 7, "target_sections_total": 9},
            }),
            encoding="utf-8",
        )
        contract_path.write_text(
            json.dumps({
                "typography_targets": {"primary_font_family": "Times New Roman"},
                "required_sections": [{"section_id": "motivation", "title": "Motivation"}],
                "background_policy": {"structure": "single top accent rule and filled section bands"},
                "reference_archetype_skeleton": {
                    "hard_constraints": {"min_sections_total": 7, "target_sections_total": 9},
                },
                "authored_html_skeleton": {
                    "html": "<div>nine default sections</div>",
                    "css": ".poster-header{border-top:10px solid red}",
                },
                "layout_storyboard_targets": {
                    "column_flow_contract": {
                        "min_sections_total": 7,
                        "target_sections_total": 9,
                        "section_bar_required": True,
                    }
                },
                "editorial_flow_contract": {
                    "column_capacity_contract": {"max_sections_per_column": 3},
                    "hard_rules": [
                        "Each column is normal document flow with one to three .poster-section blocks.",
                        "Use dark section bars and compact subsection headings like a real conference poster.",
                    ],
                },
            }),
            encoding="utf-8",
        )
        synced = _synchronize_staged_color_system(ctx, attempt_dir, "Generate a paper poster")
        assert synced["typography_contract"]["family_category"] == "sans_serif"
        synced_content = json.loads((attempt_dir / "poster_content_brief.json").read_text(encoding="utf-8"))
        synced_contract = json.loads((attempt_dir / "poster_plan_contract.json").read_text(encoding="utf-8"))
        assert synced_content["typography_contract"]["primary_font_family"] == "Arial"
        assert synced_contract["typography_targets"]["primary_font_family"] == "Arial"
        assert "times_new_roman_family_ratio_required" not in synced_contract["typography_targets"]
        assert synced_content["sections"][0]["reference_hierarchy_role"] == "nested_subsection_topic"
        assert synced_content["panel_plan"][0]["major_section_count"] == 1
        assert synced_contract["required_sections"][0]["top_level_major_section"] is False
        assert synced_contract["editorial_flow_contract"]["column_capacity_contract"]["max_sections_per_column"] == 1
        assert "dark section bars" not in " ".join(synced_contract["editorial_flow_contract"]["hard_rules"]).lower()
        assert len(synced_content["color_system_options"]) == 1
        assert "single top accent rule" not in synced_content["background_contract"]["structure"]
        assert synced_content["reference_archetype_skeleton"]["hard_constraints"]["target_sections_total"] == 5
        assert synced_content["native_reference_targets"]["target_sections_total"] == 5
        assert synced_contract["reference_archetype_skeleton"]["hard_constraints"]["min_sections_total"] == 5
        assert synced_contract["layout_storyboard_targets"]["column_flow_contract"]["target_sections_total"] == 5
        assert synced_contract["layout_storyboard_targets"]["column_flow_contract"]["section_bar_required"] is False
        authored_skeleton = synced_contract["authored_html_skeleton"]
        assert authored_skeleton["source"] == "reference_style_blueprint.html"
        assert authored_skeleton["html"].count("poster-section") == 5
        assert "border-top:10px solid red" not in authored_skeleton["css"]
        ctx.state["raw_user_brief"] = "Use Teal Coral for the reference-styled poster."
        _synchronize_staged_color_system(ctx, attempt_dir, ctx.state["raw_user_brief"])
        explicit_reference_content = json.loads(content_path.read_text(encoding="utf-8"))
        explicit_reference_contract = json.loads(contract_path.read_text(encoding="utf-8"))
        serialized_color_contracts = json.dumps(
            {
                "content": explicit_reference_content.get("color_system"),
                "required": explicit_reference_content.get("required_color_system"),
                "options": explicit_reference_content.get("color_system_options"),
                "contract": explicit_reference_contract.get("color_system"),
            },
            ensure_ascii=False,
        ).lower()
        assert "single top accent rule" not in serialized_color_contracts
        assert "use palette color elsewhere for compact filled section" not in serialized_color_contracts
        assert len(explicit_reference_content["color_system_options"]) == 2
        assert "teal_coral" in explicit_reference_content["aesthetic_contract"]["palette_usage_policy"]
        ctx.state.pop("raw_user_brief", None)
        ctx.state["poster_plan_contract"]["reference_profile"] = "conference_editorial_flow"
        _check_reference_typography_validator(ctx)

    print("reference poster style no-API checks passed")
    return 0


def _check_reference_typography_validator(ctx: ToolContext) -> None:
    soup = BeautifulSoup(
        """<main class="paper-poster"><header data-panel-role="identity_header">
        <h1 data-block-id="title">Target Paper Title</h1>
        <p class="authors" data-block-id="authors">A. Author and B. Author</p>
        <p class="institutions" data-block-id="institutions">Target University</p></header>
        <div class="reference-lead-band" data-style-role="lead-band" data-block-id="lead">Target paper summary</div>
        <div class="poster-columns"><section class="poster-column"><article class="poster-section">
        <h2 data-block-id="heading">Target Section</h2>
        <h3 class="inline_colored_label" data-block-id="subsection">Ablations</h3>
        <p data-block-id="body">This is target paper body copy with enough words.</p>
        <p class="caption" data-block-id="caption">Target paper figure caption</p>
        </article></section></div></main>""",
        "html.parser",
    )
    roles = {
        "title": (72, 700),
        "authors": (26, 500),
        "institutions": (26, 500),
        "lead": (38, 700),
        "heading": (34, 700),
        "subsection": (24, 700),
        "body": (24, 400),
        "caption": (20, 400),
    }
    bboxes = {
        block_id: {
            "x": 0,
            "y": 0,
            "w": 500,
            "h": 100,
            "_computed_style": {
                "font_family": "Arial, sans-serif",
                "font_size_px": size,
                "font_weight": weight,
                "font_style": "normal",
                "line_height": 1.18,
            },
        }
        for block_id, (size, weight) in roles.items()
    }
    assert _paper_poster_typography_contract_error(soup, bboxes, ctx) is None
    bboxes["body"]["_computed_style"]["font_family"] = "Times New Roman"
    error = _paper_poster_typography_contract_error(soup, bboxes, ctx)
    assert error is not None
    assert error.payload["issues"][0]["failure_kind"] == "font_family_not_active_contract"


def _check_reference_style_hard_gate(
    ctx: ToolContext,
    contract: dict[str, object],
) -> dict[str, object]:
    ctx.state["reference_style_contract"] = contract
    body = """<header class="poster-header identity-header" data-panel-role="identity_header" data-block-id="identity">
    <h1>Target title</h1><p>Target authors</p><p>Target institution</p></header>
    <div class="reference-lead-band" data-style-role="lead-band">Target paper summary</div>
    <div class="reference-chrome" data-style-role="chrome-layer" aria-hidden="true"><span class="chrome-group"><span class="reference-rail" style="display:block;width:8px;height:900px;background:#137A78"></span></span></div>
    <div class="poster-columns">
      <div class="poster-column" data-block-id="c1"><section class="poster-section" data-block-id="s1"><h2>One</h2><div class="formula" data-block-id="formula">\\[x=1\\]</div></section></div>
      <div class="poster-column" data-block-id="c2"><section class="poster-section" data-block-id="s2"><h2>Two</h2><table class="booktabs" data-block-id="table"><tr><th>A</th></tr><tr><td>B</td></tr></table></section></div>
      <div class="poster-column" data-block-id="c3"><section class="poster-section" data-block-id="s3"><h2>Three</h2></section><section class="poster-section" data-block-id="s4"><h2>Extra</h2></section></div>
      <div class="poster-column" data-block-id="c4"><section class="poster-section" data-block-id="s5"><h2>Four</h2></section></div>
      <div class="poster-column" data-block-id="c5"><section class="poster-section" data-block-id="s6"><h2>Five</h2></section></div>
    </div>"""
    css = """
    .poster-header{border-top:12px solid #7A1F5C}
    .poster-section{border-bottom:2px solid #7A1F5C}
    .formula{border-top:1px solid #171717;border-bottom:1px solid #171717}
    .booktabs{border-top:3px solid #171717;border-bottom:3px solid #171717}
    .poster-section::before{content:'';position:absolute;border:12px solid #7A1F5C}
    """
    def border_style(**overrides: object) -> dict[str, object]:
        style: dict[str, object] = {}
        for side in ("top", "right", "bottom", "left"):
            style[f"border_{side}_width_px"] = 0
            style[f"border_{side}_style"] = "none"
            style[f"border_{side}_color"] = "rgba(0, 0, 0, 0)"
        style.update(overrides)
        return {"_computed_style": style}

    bboxes = {
        "__paper_poster_root__": {**border_style(), "x": 0, "y": 0, "w": 3072, "h": 1536},
        "c1": {**border_style(), "x": 0, "y": 260, "w": 614, "h": 1276},
        "c2": {**border_style(), "x": 614, "y": 260, "w": 614, "h": 1276},
        "c3": {**border_style(), "x": 1228, "y": 260, "w": 614, "h": 1276},
        "c4": {**border_style(), "x": 1842, "y": 260, "w": 614, "h": 1276},
        "c5": {**border_style(), "x": 2456, "y": 260, "w": 616, "h": 1276},
        "identity": border_style(border_top_width_px=12, border_top_style="solid", border_top_color="rgb(122, 31, 92)"),
        "s1": border_style(border_bottom_width_px=2, border_bottom_style="solid", border_bottom_color="rgb(122, 31, 92)"),
        "s2": border_style(border_bottom_width_px=2, border_bottom_style="solid", border_bottom_color="rgb(122, 31, 92)"),
        "s3": border_style(border_bottom_width_px=2, border_bottom_style="solid", border_bottom_color="rgb(122, 31, 92)"),
        "s4": border_style(border_bottom_width_px=2, border_bottom_style="solid", border_bottom_color="rgb(122, 31, 92)"),
        "s5": border_style(border_bottom_width_px=2, border_bottom_style="solid", border_bottom_color="rgb(122, 31, 92)"),
        "s6": border_style(border_bottom_width_px=2, border_bottom_style="solid", border_bottom_color="rgb(122, 31, 92)"),
        "formula": border_style(
            border_top_width_px=1, border_top_style="solid", border_top_color="rgb(23, 23, 23)",
            border_bottom_width_px=1, border_bottom_style="solid", border_bottom_color="rgb(23, 23, 23)",
        ),
        "table": border_style(
            border_top_width_px=3, border_top_style="solid", border_top_color="rgb(23, 23, 23)",
            border_bottom_width_px=3, border_bottom_style="solid", border_bottom_color="rgb(23, 23, 23)",
        ),
    }
    error = _reference_style_contract_error(
        BeautifulSoup(body, "html.parser"), css, ctx, bboxes=bboxes
    )
    assert error is not None
    kinds = {item["failure_kind"] for item in error.payload["issues"]}
    assert "reference_major_section_count_mismatch" in kinds
    assert "reference_header_top_rule_leakage" in kinds
    assert "reference_section_divider_leakage" in kinds
    assert "reference_formula_frame_leakage" in kinds
    assert "reference_chrome_attached_to_content_regions" in kinds
    assert "reference_region_geometry_mismatch" in kinds
    assert error.payload["blocks_soft_accept"] is True

    clean = body.replace(
        '<section class="poster-section" data-block-id="s3"><h2>Three</h2></section><section class="poster-section" data-block-id="s4"><h2>Extra</h2></section>',
        '<section class="poster-section" data-block-id="s3"><h2>Three</h2><h3>Extra</h3></section>',
    )
    assert _reference_style_contract_error(BeautifulSoup(clean, "html.parser"), "", ctx) is None
    harmless_pseudo = ".poster-section::before{content:none;border:20px solid red}"
    assert _reference_style_contract_error(
        BeautifulSoup(clean, "html.parser"), harmless_pseudo, ctx
    ) is None
    transparent_bboxes = {
        key: border_style(
            border_top_width_px=2 if key == "identity" else 0,
            border_top_style="solid" if key == "identity" else "none",
            border_top_color="rgb(0 0 0 / 0)" if key == "identity" else "rgba(0, 0, 0, 0)",
        )
        for key in ("__paper_poster_root__", "identity", "s1", "s2", "s3", "s5", "s6", "formula", "table")
    }
    transparent_bboxes["__paper_poster_root__"].update({"x": 0, "y": 0, "w": 3072, "h": 1536})
    for index, expected_box in enumerate(
        contract["style_tokens"]["layout_rhythm"]["region_boxes"],
        start=1,
    ):
        transparent_bboxes[f"c{index}"] = {
            **border_style(),
            "x": expected_box["x_pct"] * 30.72,
            "y": expected_box["y_pct"] * 15.36,
            "w": expected_box["w_pct"] * 30.72,
            "h": expected_box["h_pct"] * 15.36,
        }
    transparent_error = _reference_style_contract_error(
        BeautifulSoup(clean, "html.parser"), "", ctx, bboxes=transparent_bboxes
    )
    assert transparent_error is None, transparent_error.payload if transparent_error else None
    clean = _with_reference_region_geometry(clean, contract)
    canvas = {"w_px": 3072, "h_px": 1536, "dpi": 150, "aspect_ratio": "2:1", "color_mode": "RGB"}
    crossing_chrome_soup = BeautifulSoup(clean, "html.parser")
    _ensure_dom_block_ids(crossing_chrome_soup, ctx, panel_flow_mode=True)
    crossing_chrome_css = (
        ".reference-chrome{position:absolute;inset:0}"
        ".reference-rail{position:absolute!important;left:0!important;top:0!important;"
        "width:100%!important;height:100%!important;background:#137A78!important}"
    )
    crossing_chrome_bboxes = _measure_dom_bboxes(
        str(crossing_chrome_soup),
        crossing_chrome_css,
        canvas=canvas,
        ctx=ctx,
        stage="reference_crossing_chrome_smoke",
    )
    crossing_chrome_error = _reference_style_contract_error(
        crossing_chrome_soup,
        crossing_chrome_css,
        ctx,
        bboxes=crossing_chrome_bboxes,
    )
    assert crossing_chrome_error is not None
    assert "reference_chrome_crosses_content" in {
        item["failure_kind"] for item in crossing_chrome_error.payload["issues"]
    }
    pseudo_chrome_soup = BeautifulSoup(clean, "html.parser")
    pseudo_rail = pseudo_chrome_soup.select_one(".reference-rail")
    pseudo_rail["style"] = "position:absolute;left:0;top:0;width:8px;height:900px"
    _ensure_dom_block_ids(pseudo_chrome_soup, ctx, panel_flow_mode=True)
    pseudo_chrome_css = (
        ".reference-chrome{position:absolute;inset:0}"
        ".reference-rail::before{content:'';display:block;width:100%;height:100%;background:#137A78}"
    )
    pseudo_chrome_bboxes = _measure_dom_bboxes(
        str(pseudo_chrome_soup),
        pseudo_chrome_css,
        canvas=canvas,
        ctx=ctx,
        stage="reference_pseudo_chrome_smoke",
    )
    pseudo_chrome_error = _reference_style_contract_error(
        pseudo_chrome_soup,
        pseudo_chrome_css,
        ctx,
        bboxes=pseudo_chrome_bboxes,
    )
    pseudo_chrome_kinds = {
        item["failure_kind"]
        for item in (pseudo_chrome_error.payload["issues"] if pseudo_chrome_error else [])
    }
    assert "reference_chrome_layer_not_visible" not in pseudo_chrome_kinds
    for suffix, override in (
        ("display_none", "display:none"),
        ("opacity_zero", "opacity:0"),
    ):
        cascaded_pseudo_css = (
            ".poster-section::before{content:'';border:2px solid red}"
            f".poster-section::before{{{override}}}"
        )
        cascaded_pseudo_soup = BeautifulSoup(clean, "html.parser")
        _ensure_dom_block_ids(cascaded_pseudo_soup, ctx, panel_flow_mode=True)
        cascaded_pseudo_bboxes = _measure_dom_bboxes(
            str(cascaded_pseudo_soup),
            cascaded_pseudo_css,
            canvas=canvas,
            ctx=ctx,
            stage=f"reference_cascaded_pseudo_{suffix}_smoke",
        )
        assert _reference_style_contract_error(
            cascaded_pseudo_soup,
            cascaded_pseudo_css,
            ctx,
            bboxes=cascaded_pseudo_bboxes,
        ) is None
    hidden_chrome_soup = BeautifulSoup(clean, "html.parser")
    hidden_chrome_soup.select_one(".reference-rail")["style"] = "display:none"
    _ensure_dom_block_ids(hidden_chrome_soup, ctx, panel_flow_mode=True)
    hidden_chrome_bboxes = _measure_dom_bboxes(
        str(hidden_chrome_soup),
        "",
        canvas=canvas,
        ctx=ctx,
        stage="reference_hidden_chrome_smoke",
    )
    hidden_chrome_error = _reference_style_contract_error(
        hidden_chrome_soup,
        "",
        ctx,
        bboxes=hidden_chrome_bboxes,
    )
    assert hidden_chrome_error is not None
    assert "reference_chrome_layer_not_visible" in {
        item["failure_kind"] for item in hidden_chrome_error.payload["issues"]
    }
    panel_flow_soup = BeautifulSoup(clean, "html.parser")
    formula_slot = panel_flow_soup.new_tag("div")
    formula_slot["class"] = "formula-slot"
    formula_slot["data-block-id"] = "formula_slot"
    formula_slot.string = "Nested equation"
    panel_flow_soup.select_one(".poster-section").append(formula_slot)
    _ensure_dom_block_ids(panel_flow_soup, ctx, panel_flow_mode=True)
    formula_block = panel_flow_soup.select_one(".formula")
    assert formula_block is not None
    assert formula_block.get("data-block-id") == "formula"
    assert panel_flow_soup.select_one(".formula-slot").get("data-block-id") == "formula_slot"
    panel_flow_bboxes = _measure_dom_bboxes(
        str(panel_flow_soup),
        "",
        canvas=canvas,
        ctx=ctx,
        stage="reference_formula_measurement_smoke",
    )
    assert "formula" in panel_flow_bboxes
    assert "formula_slot" in panel_flow_bboxes
    filled_header_css = (
        ".poster-header{background:rgb(0,85,170);text-align:center;border-bottom:12px solid rgb(0,85,170)}"
        ".poster-header h1{text-align:center}"
    )
    filled_header_bboxes = _measure_dom_bboxes(
        str(panel_flow_soup),
        filled_header_css,
        canvas=canvas,
        ctx=ctx,
        stage="reference_filled_header_smoke",
    )
    filled_header_error = _reference_style_contract_error(
        panel_flow_soup, filled_header_css, ctx, bboxes=filled_header_bboxes
    )
    assert filled_header_error is not None
    filled_header_kinds = {item["failure_kind"] for item in filled_header_error.payload["issues"]}
    assert "reference_header_top_rule_leakage" in filled_header_kinds
    assert "reference_header_background_leakage" in filled_header_kinds
    assert "reference_header_alignment_mismatch" in filled_header_kinds
    chrome_header_css = (
        "body{background:rgb(0,85,170)}.paper-poster{background:transparent}"
        ".poster-header{background:transparent;display:grid;justify-content:center;box-shadow:0 0 0 4px red}"
        ".poster-header h1{text-align:left}"
        ".poster-header::after{content:'';display:block;border-bottom:5px solid red}"
    )
    chrome_header_bboxes = _measure_dom_bboxes(
        str(panel_flow_soup),
        chrome_header_css,
        canvas=canvas,
        ctx=ctx,
        stage="reference_header_chrome_smoke",
    )
    chrome_header_error = _reference_style_contract_error(
        panel_flow_soup, chrome_header_css, ctx, bboxes=chrome_header_bboxes
    )
    assert chrome_header_error is not None
    chrome_header_kinds = {item["failure_kind"] for item in chrome_header_error.payload["issues"]}
    assert "reference_header_background_leakage" in chrome_header_kinds
    assert "reference_header_layout_alignment_mismatch" in chrome_header_kinds
    assert "reference_header_chrome_leakage" in chrome_header_kinds
    assert "reference_header_pseudo_chrome_leakage" in chrome_header_kinds

    nested_soup = BeautifulSoup(clean, "html.parser")
    nested_section = nested_soup.new_tag("section")
    nested_section["class"] = "poster-section"
    nested_section.string = "Nested panel"
    nested_soup.select_one(".poster-section").append(nested_section)
    nested_error = _reference_style_contract_error(nested_soup, "", ctx)
    assert nested_error is not None
    assert "reference_nested_major_section_leakage" in {
        item["failure_kind"] for item in nested_error.payload["issues"]
    }

    internal_rule_soup = BeautifulSoup(clean, "html.parser")
    internal_rule_soup.select_one(".formula").append(internal_rule_soup.new_tag("hr"))
    role_formula = internal_rule_soup.new_tag("div")
    role_formula["data-role"] = "formula-main"
    role_formula["data-block-id"] = "formula_role"
    role_formula.string = "Role equation"
    internal_rule_soup.select_one(".poster-section").append(role_formula)
    internal_rule_css = (
        ".booktabs tbody tr{border-bottom:1px solid #171717}"
        ".formula::after{content:'';display:block;border-bottom:2px solid #171717}"
        "[data-role='formula-main']::before{content:'';display:block;border-top:2px solid #171717}"
    )
    _ensure_dom_block_ids(internal_rule_soup, ctx, panel_flow_mode=True)
    internal_rule_bboxes = _measure_dom_bboxes(
        str(internal_rule_soup),
        internal_rule_css,
        canvas=canvas,
        ctx=ctx,
        stage="reference_internal_rule_smoke",
    )
    internal_rule_error = _reference_style_contract_error(
        internal_rule_soup, internal_rule_css, ctx, bboxes=internal_rule_bboxes
    )
    assert internal_rule_error is not None
    internal_rule_kinds = {item["failure_kind"] for item in internal_rule_error.payload["issues"]}
    assert "reference_formula_internal_rule_leakage" in internal_rule_kinds
    assert "reference_formula_pseudo_rule_leakage" in internal_rule_kinds
    assert "reference_table_row_rule_leakage" in internal_rule_kinds
    assert any(
        item.get("id") == "formula_role"
        for item in ctx.state["paper_poster_html_reference_rule_measurements"]["formulaPseudos"]
    )
    overridden_rule_soup = BeautifulSoup(clean, "html.parser")
    hidden_hr = overridden_rule_soup.new_tag("hr")
    hidden_hr["style"] = "display:none"
    overridden_rule_soup.select_one(".formula").append(hidden_hr)
    _ensure_dom_block_ids(overridden_rule_soup, ctx, panel_flow_mode=True)
    overridden_rule_css = (
        internal_rule_css
        + ".booktabs tbody tr{border-bottom:0}"
        + ".formula::after{display:none;border-bottom:0}"
    )
    overridden_rule_bboxes = _measure_dom_bboxes(
        str(overridden_rule_soup),
        overridden_rule_css,
        canvas=canvas,
        ctx=ctx,
        stage="reference_overridden_rule_smoke",
    )
    assert _reference_style_contract_error(
        overridden_rule_soup, overridden_rule_css, ctx, bboxes=overridden_rule_bboxes
    ) is None
    browser_visible_css = ".poster-header{border-top:12px solid rgb(255, 165, 0)}"
    browser_visible_bboxes = _measure_dom_bboxes(
        clean,
        browser_visible_css,
        canvas=canvas,
        ctx=ctx,
        stage="reference_visible_border_smoke",
    )
    browser_visible_error = _reference_style_contract_error(
        BeautifulSoup(clean, "html.parser"), browser_visible_css, ctx, bboxes=browser_visible_bboxes
    )
    assert browser_visible_error is not None
    assert "reference_header_top_rule_leakage" in {
        item["failure_kind"] for item in browser_visible_error.payload["issues"]
    }
    browser_transparent_css = (
        ".poster-header{border-top:12px solid rgb(255, 165, 0)!important}"
        "header.poster-header{border-top-color:rgb(0 0 0 / 0)!important}"
    )
    browser_transparent_bboxes = _measure_dom_bboxes(
        clean,
        browser_transparent_css,
        canvas=canvas,
        ctx=ctx,
        stage="reference_transparent_border_smoke",
    )
    assert _reference_style_contract_error(
        BeautifulSoup(clean, "html.parser"),
        browser_transparent_css,
        ctx,
        bboxes=browser_transparent_bboxes,
    ) is None
    hidden_soup = BeautifulSoup(clean, "html.parser")
    hidden_formula = hidden_soup.new_tag("div")
    hidden_formula["class"] = "formula"
    hidden_formula["data-block-id"] = "hidden_formula"
    hidden_formula["style"] = "display:none;border:4px solid red"
    hidden_formula.string = "Hidden equation"
    hidden_soup.select_one(".poster-section").append(hidden_formula)
    _ensure_dom_block_ids(hidden_soup, ctx, panel_flow_mode=True)
    hidden_bboxes = _measure_dom_bboxes(
        str(hidden_soup),
        "",
        canvas=canvas,
        ctx=ctx,
        stage="reference_hidden_formula_smoke",
    )
    assert "hidden_formula" not in hidden_bboxes
    assert _reference_style_contract_error(hidden_soup, "", ctx, bboxes=hidden_bboxes) is None
    ancestor_hidden_soup = BeautifulSoup(clean, "html.parser")
    hidden_wrapper = ancestor_hidden_soup.new_tag("div")
    hidden_wrapper["style"] = "display:none"
    ancestor_hidden_formula = ancestor_hidden_soup.new_tag("div")
    ancestor_hidden_formula["class"] = "formula"
    ancestor_hidden_formula["data-block-id"] = "ancestor_hidden_formula"
    ancestor_hidden_formula["style"] = "display:block;border:4px solid red"
    ancestor_hidden_formula.string = "Ancestor-hidden equation"
    hidden_wrapper.append(ancestor_hidden_formula)
    ancestor_hidden_soup.select_one(".poster-section").append(hidden_wrapper)
    _ensure_dom_block_ids(ancestor_hidden_soup, ctx, panel_flow_mode=True)
    ancestor_hidden_bboxes = _measure_dom_bboxes(
        str(ancestor_hidden_soup),
        "",
        canvas=canvas,
        ctx=ctx,
        stage="reference_ancestor_hidden_formula_smoke",
    )
    assert "ancestor_hidden_formula" not in ancestor_hidden_bboxes
    assert _reference_style_contract_error(
        ancestor_hidden_soup, "", ctx, bboxes=ancestor_hidden_bboxes
    ) is None
    ctx.state.pop("paper_poster_html_computed_style_measurements", None)
    missing_computed = {
        key: {"_computed_style": {"font_size_px": 24}}
        for key in transparent_bboxes
    }
    missing_error = _reference_style_contract_error(
        BeautifulSoup(clean, "html.parser"), "", ctx, bboxes=missing_computed
    )
    assert missing_error is not None
    assert "reference_style_computed_measurement_unavailable" in {
        item["failure_kind"] for item in missing_error.payload["issues"]
    }
    ctx.state["reference_style_contract"] = contract
    assert _reference_style_contract(ctx)["version"] == 4
    ctx.state["reference_style_contract"] = contract
    return {
        "summary": {
            "issue_id": error.payload["issue_id"],
            "repair_route": error.payload["repair_route"],
            "issues": error.payload["issues"],
            "hint": error.payload["hint"],
        },
        "payload": dict(error.payload),
    }


def _with_reference_region_geometry(html: str, contract: dict[str, object]) -> str:
    soup = BeautifulSoup(html, "html.parser")
    columns = soup.select_one(".poster-columns")
    if columns is None:
        return html
    columns["style"] = "position:absolute;inset:0"
    boxes = contract["style_tokens"]["layout_rhythm"]["region_boxes"]
    for column, box in zip(soup.select(".poster-column"), boxes, strict=True):
        column["style"] = (
            f"position:absolute;left:{box['x_pct']}%;top:{box['y_pct']}%;"
            f"width:{box['w_pct']}%;height:{box['h_pct']}%"
        )
    return str(soup)


def _check_reference_hard_issue_blocks_accepted_fallback(
    ctx: ToolContext,
    root: Path,
    feedback: dict[str, object],
) -> None:
    candidate_dir = root / "accepted_candidate"
    candidate_dir.mkdir()
    html = "<!doctype html><html><body><main class='paper-poster'>Target poster</main></body></html>"
    paths = {
        "measure_html": candidate_dir / "measure.html",
        "body_html": candidate_dir / "body.html",
        "style_css": candidate_dir / "style.css",
        "preview_png": candidate_dir / "preview.png",
        "measurement_json": candidate_dir / "measurement.json",
    }
    paths["measure_html"].write_text(html, encoding="utf-8")
    paths["body_html"].write_text(html, encoding="utf-8")
    paths["style_css"].write_text(".paper-poster{width:3072px;height:1536px}", encoding="utf-8")
    paths["preview_png"].write_bytes(b"png")
    paths["measurement_json"].write_text("{}", encoding="utf-8")
    candidate = {
        "status": "accepted",
        "candidate_score": 1200,
        **{f"_{key}_abs": str(path) for key, path in paths.items()},
    }
    acceptance = _best_candidate_fallback_acceptance(ctx, candidate, feedback)
    assert acceptance["accepted"] is False
    assert acceptance["reason"] == "current_reference_style_hard_issue"
    available_acceptance = _best_available_artifact_fallback_acceptance(ctx, candidate, feedback)
    assert available_acceptance["accepted"] is False
    assert available_acceptance["reason"] == "current_reference_style_hard_issue"


def _write_fixtures(root: Path) -> dict[str, Path]:
    image_path = root / "reference.png"
    image = Image.new("RGB", (1200, 600), "white")
    draw = ImageDraw.Draw(image)
    draw.rectangle((0, 0, 1200, 54), fill="#7A1F5C")
    draw.rectangle((35, 110, 380, 155), fill="#7A1F5C")
    draw.rectangle((425, 110, 770, 155), fill="#7A1F5C")
    draw.rectangle((815, 110, 1160, 155), fill="#7A1F5C")
    image.save(image_path)

    html_path = root / "reference.html"
    html_path.write_text(
        """<!doctype html><html><head><style>
        html,body{margin:0;background:white}.paper-poster{width:1200px;height:600px;color:#161616;
        font-family:serif;border-top:12px solid #7A1F5C}.cols{display:grid;grid-template-columns:repeat(3,1fr);gap:24px;padding:30px}
        h2{background:#7A1F5C;color:white;margin:0;padding:8px}section{min-height:350px}
        </style></head><body><main class="paper-poster"><div class="cols"><section><h2>One</h2></section><section><h2>Two</h2></section><section><h2>Three</h2></section></div></main></body></html>""",
        encoding="utf-8",
    )

    pdf_path = root / "reference.pdf"
    pdf = fitz.open()
    page = pdf.new_page(width=1200, height=600)
    page.draw_rect(fitz.Rect(0, 0, 1200, 12), color=(0.48, 0.12, 0.36), fill=(0.48, 0.12, 0.36))
    page.insert_text((40, 90), "Reference poster", fontsize=36)
    pdf.save(pdf_path)
    pdf.close()

    pptx_path = root / "reference.pptx"
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(6.666)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.0), Inches(1.0))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = "Reference poster"
    run.font.size = Pt(36)
    presentation.save(pptx_path)
    return {"png": image_path, "html": html_path, "pdf": pdf_path, "pptx": pptx_path}


def _write_fake_harness(root: Path) -> Path:
    harness = root / "fake_reference_style_agent.py"
    harness.write_text(
        """import json
from pathlib import Path

repair = Path('reference_style_attempts/attempt_01').exists()
analysis = {
  'version': 4,
  'transfer_mode': 'reference_first_reconstruction',
  'summary': 'Asymmetric academic poster with left identity, four body regions, and gutter chrome.',
  'palette': {
    'background': '#FFFFFF', 'ink': '#171717', 'primary': '#7A1F5C',
    'secondary': '#F1E4EC', 'accent': '#137A78',
    'header_text': '#7A1F5C', 'section_heading_text': '#171717'
  },
  'header_treatment': {'mode': 'open_white', 'alignment': 'left', 'composition': 'left_identity_cluster', 'background_role': 'background', 'title_color_role': 'primary', 'rule_placement': 'none', 'rule_color_role': 'primary', 'rule_width_px': 0},
  'lead_band': {'present': True, 'placement': 'below_identity', 'background_role': 'primary', 'text_color_role': 'on_primary', 'alignment': 'center', 'height_px': 64, 'text_size_px': 38},
  'section_heading_treatment': {'mode': 'underline', 'text_color_role': 'ink', 'fill_role': 'background', 'border_role': 'ink', 'border_width_px': 0, 'corner_style': 'square', 'rule_color_role': 'ink', 'rule_width_px': 3},
  'section_structure': {'inter_section_dividers': 'none', 'outer_border': 'none', 'vertical_accent_rules': 'none'},
  'body_region_structure': {
    'layout_mode': 'freeform_regions', 'region_count': 5, 'major_section_count': 5,
    'major_sections_per_region': [1, 1, 1, 1, 1],
    'regions': [
      {'region_id':'region_1','region_role':'column','section_count':1,'reading_order':1},
      {'region_id':'region_2','region_role':'column','section_count':1,'reading_order':2},
      {'region_id':'region_3','region_role':'column','section_count':1,'reading_order':3},
      {'region_id':'region_4','region_role':'side_callout','section_count':1,'reading_order':4},
      {'region_id':'region_5','region_role':'footer_band','section_count':1,'reading_order':5}
    ],
    'subsection_treatment': 'inline_colored_label'
  },
  'surfaces': {'panel_fill': 'white', 'border_style': 'none', 'corner_style': 'square', 'shadow_style': 'none'},
  'spacing': {'outer_margin_px': 36, 'column_gap_px': 28, 'section_gap_px': 14, 'panel_padding_px': 10},
  'layout_rhythm': {
    'region_proportions': [0.7, 1.1, 1.4, 0.8, 1.8], 'density': 'dense',
    'region_boxes': [
      {'region_id':'region_1','x_pct':1,'y_pct':18,'w_pct':18,'h_pct':62},
      {'region_id':'region_2','x_pct':20,'y_pct':18,'w_pct':23,'h_pct':62},
      {'region_id':'region_3','x_pct':44,'y_pct':18,'w_pct':34,'h_pct':62},
      {'region_id':'region_4','x_pct':79,'y_pct':18,'w_pct':20,'h_pct':38},
      {'region_id':'region_5','x_pct':79 if repair else 70,'y_pct':58 if repair else 30,'w_pct':20,'h_pct':22 if repair else 38}
    ]
  },
  'chrome_treatment': {'present': True, 'placement': 'gutters', 'density': 'sparse', 'crossing_policy': 'never_cross_content'},
  'typography_style': {'display_family_category': 'sans_serif', 'body_family_category': 'sans_serif', 'family_category': 'sans_serif', 'title_weight': 700, 'identity_weight': 500, 'section_heading_weight': 700, 'body_weight': 400, 'title_size_px': 72, 'identity_size_px': 26, 'section_heading_size_px': 34, 'body_size_px': 24, 'caption_size_px': 20},
  'table_treatment': {'observed': False, 'rule_style': 'none', 'header_fill': 'none'},
  'formula_treatment': {'frame': 'none', 'background': 'none'},
  'figure_treatment': {'frame': 'none', 'caption_alignment': 'left'}
}
Path('reference_style_analysis.json').write_text(json.dumps(analysis), encoding='utf-8')
Path('reference_style_blueprint.html').write_text('''<!doctype html><html><head><style>
.reference-style-blueprint{width:4800px;height:2400px;background:#FFFFFF;color:#171717;font-family:Arial,sans-serif}
[data-style-role="identity-header"]{display:flex;justify-content:flex-start;text-align:left}[data-style-role="lead-band"]{background:#7A1F5C;color:#FFFFFF}
[data-style-role="body-regions"]{position:relative;height:1880px}
[data-style-role="chrome-layer"]{position:absolute;inset:0;z-index:0;pointer-events:none}.reference-rail{position:absolute;left:.2%;top:20px;width:8px;height:120px}.reference-rail::before{content:'';display:block;width:100%;height:100%;background:#137A78}
[data-style-role="body-region"]{position:absolute;z-index:1}[data-region-id="region_1"]{left:1%;top:18%;width:18%;height:62%}[data-region-id="region_2"]{left:20%;top:18%;width:23%;height:62%}[data-region-id="region_3"]{left:44%;top:18%;width:34%;height:62%}[data-region-id="region_4"]{left:79%;top:18%;width:20%;height:38%}[data-region-id="region_5"]{left:__R5_LEFT__%;top:__R5_TOP__%;width:20%;height:__R5_HEIGHT__%}
[data-style-role="section-heading"]{border-bottom:3px solid #171717}
</style></head><body><main class="reference-style-blueprint"><header data-style-role="identity-header">{{PAPER_TITLE}}<div>{{AUTHORS}}</div><div>{{INSTITUTIONS}}</div></header><div data-style-role="lead-band">{{TARGET_PAPER_SUMMARY}}</div><div data-style-role="chrome-layer"><span class="reference-rail"></span></div><div data-style-role="body-regions"><div data-style-role="body-region" data-region-id="region_1" data-region-role="column"><section data-style-role="section"><h2 data-style-role="section-heading">{{SECTION_TITLE}}</h2><div class="formula-slot">{{TARGET_PAPER_CONTENT}}</div></section></div><div data-style-role="body-region" data-region-id="region_2" data-region-role="column"><section data-style-role="section"><h2 data-style-role="section-heading">{{SECTION_TITLE}}</h2>{{TARGET_PAPER_FIGURE}}</section></div><div data-style-role="body-region" data-region-id="region_3" data-region-role="column"><section data-style-role="section"><h2 data-style-role="section-heading">{{SECTION_TITLE}}</h2><div class="table-slot">{{TARGET_PAPER_TABLE}}</div></section></div><div data-style-role="body-region" data-region-id="region_4" data-region-role="side_callout"><section data-style-role="section"><h2 data-style-role="section-heading">{{SECTION_TITLE}}</h2>{{TARGET_PAPER_CONTENT}}</section></div><div data-style-role="body-region" data-region-id="region_5" data-region-role="footer_band"><section data-style-role="section"><h2 data-style-role="section-heading">{{SECTION_TITLE}}</h2>{{TARGET_PAPER_CONTENT}}</section></div></div></main></body></html>'''.replace('__R5_LEFT__', '79' if repair else '70').replace('__R5_TOP__', '58' if repair else '30').replace('__R5_HEIGHT__', '22' if repair else '38'), encoding='utf-8')
import hashlib
import time
time.sleep(4)
blueprint_sha256 = hashlib.sha256(Path('reference_style_blueprint.html').read_bytes()).hexdigest()
Path('reference_style_agent_review.json').write_text(json.dumps({
  'status':'ok', 'rendered_blueprint_inspected':True, 'header_matches_reference':True,
  'body_region_geometry_matches_reference':True, 'chrome_avoids_content':True,
  'blueprint_sha256': blueprint_sha256
}), encoding='utf-8')
Path('reference_style_agent_done.json').write_text('{\"status\":\"ok\"}', encoding='utf-8')
""",
        encoding="utf-8",
    )
    return harness


if __name__ == "__main__":
    raise SystemExit(main())
