"""HTTP transport and job wiring for the AutoDesign Workbench.

Lifecycle:
  1. uvicorn boots this module — FastAPI server stays up, "agent on standby"
  2. browser reserves a run, streams declared inputs, then starts generation
  3. RunSupervisor owns one cancellable OS worker for the pipeline
  4. when the run completes, locates the artifact under out/runs/<run_id>/final/
     and returns a JSON response pointing the browser at the file URL

Run:
    uv run uvicorn scripts.web_server:app --reload --port 8000

Key env: OPENROUTER_API_KEY (read by autodesign.config.load_settings).
Mocking has been removed in this revision — every /api/generate call costs
real money and takes 30s-5min depending on artifact type.
"""

from __future__ import annotations

import asyncio
import base64
import hashlib
import html
import json
import mimetypes
import os
import re
import secrets
import shlex
import shutil
import signal
import subprocess
import sys
import tempfile
import threading
import time
import uuid
import warnings
from contextlib import AsyncExitStack, asynccontextmanager
from copy import copy
from dataclasses import asdict, dataclass, replace
from datetime import datetime, timedelta, timezone
from pathlib import Path
from typing import Any, AsyncIterator, Callable, Literal
from urllib.parse import unquote, urlparse

from fastapi import FastAPI, File, Form, HTTPException, Query, Request, UploadFile
from fastapi.responses import FileResponse, JSONResponse, StreamingResponse
from fastapi.staticfiles import StaticFiles
from starlette.datastructures import MutableHeaders
from PIL import Image, UnidentifiedImageError
from pydantic import BaseModel, Field, ValidationError
from bs4 import BeautifulSoup

from autodesign.config import (
    Settings,
    authoring_max_attempts_for,
    clear_harness_login_marker,
    coding_agent_smoke_command_for_harness,
    code_editor_command_for_harness,
    designer_author_command_for_harness,
    harness_auth_dir,
    harness_auth_read_dir,
    harness_login_present,
    harness_subprocess_env,
    load_settings,
    mark_harness_login,
    normalize_model_id,
    poster_canvas_preset_catalog,
    resolve_codex_runtime,
    resolve_deepseek_harness_runtime,
    resolve_harness_binary,
    resolve_template,
)
from autodesign.harness_matrix import (
    CODING_HARNESSES,
    HarnessMatrixCellSpec,
    build_coding_harness_capabilities,
    load_harness_matrix,
    run_harness_matrix,
)
from autodesign.runner import _load_resume_state
from autodesign.util.canvas_planner import CanvasIntentError, parse_canvas_intent
from autodesign.run_control import (
    CancellationToken,
    InvalidRunTransition,
    RunCancelled,
    RunControlError,
    RunControlStore,
    durable_replace_json,
    validate_run_id,
)
from autodesign.run_file_access import (
    OpenedRunFile as _OpenedPublicRunFile,
    RunFileAccessError,
    canonical_run_file_parts,
    open_run_file,
)
from autodesign.run_supervisor import (
    RunSupervisor,
    TerminalReconciliation,
    WorkerExitDiagnostic,
    WorkerOutcome,
)
from autodesign.run_worker_protocol import (
    ArtifactEditWorkerRequest,
    AttemptForkWorkerRequest,
    CandidatePublishWorkerRequest,
    EditableVideoRenderWorkerRequest,
    PipelineWorkerRequest,
    PosterCodeEditWorkerRequest,
    PptxExportWorkerRequest,
    ProtocolError as WorkerProtocolError,
    RunWorkerRequest,
    VideoExportRetryWorkerRequest,
    decode_worker_result,
    format_worker_error_message,
    parse_worker_result_json,
)
from autodesign.candidate_publish import (
    deliver_video_candidate_draft as _deliver_video_candidate_draft_core,
    reconcile_video_delivery_context_promotion,
    validate_video_delivery_context_promotion,
    validate_candidate_draft as _validate_candidate_draft_core,
)
from autodesign.web_run_services import (
    CancelResult,
    InputSlot,
    InvalidInputSlot,
    InvalidReservation,
    ReservationConflict,
    ReservationNotFound,
    RunNotReady,
    UploadAuthorizationError,
    UploadCancelled,
    UploadConflict,
    UploadIntegrityError,
    WebRunServices,
)
from autodesign.paper_bundle_jobs import (
    ChildStateSnapshot,
    PaperBundleBarrierClosed,
    PaperBundleChildDescriptor,
    PaperBundleConflict,
    PaperBundleError,
    PaperBundleInputSlot,
    PaperBundleJobRecord,
    PaperBundleJobStore,
    PaperBundleNotFound,
)
from autodesign.poster_code_edit import (
    run_poster_code_edit_sync as _run_poster_code_edit_core,
)
from autodesign.agents.atomic_artifact_promotion import (
    reconcile_artifact_promotion,
)
from autodesign.schema import ApplyEditsResult, RunResult
from autodesign.attempt_candidates import (
    candidate_summary,
    is_browser_preview_resource_path,
    load_attempt_candidate,
    load_attempt_candidates,
    load_candidate_index,
    load_selection_journal,
)
from autodesign.attempt_selection import (
    complete_source_run_with_candidate_fork,
    promote_pending_selection,
    request_attempt_selection,
)
from autodesign.agents.openresearch_gui_submitter import (
    AGENT_PROMPT_FILE as OPENRESEARCH_AGENT_PROMPT_FILE,
    PROCESS_FILE as OPENRESEARCH_GUI_PROCESS_FILE,
    submit_openresearch_gui,
)
from autodesign.util.design_events import (
    append_design_event,
    attachment_event_data,
    infer_follow_up_sentiment,
    layer_edit_events,
)
from autodesign.util.ids import new_run_id
from autodesign.util.academic_palette import (
    AcademicPaletteCatalogError,
    academic_palette_catalog_payload,
    require_academic_color_system,
)
from autodesign.util.layer_parse import (
    parse_deck_html_as_layer_mode,
    parse_html_layers,
)
from autodesign.util.math_typesetting import ensure_poster_katex_document
from autodesign.util.logging import append_jsonl_event, log, run_context
from autodesign.util.openresearch_api import OpenResearchApiClient
from autodesign.util.browser_render import (
    export_html_pdf,
    screenshot_html,
)
from autodesign.util.editable_html import ensure_editable_html_contract
from autodesign.util.io import atomic_write_json, sha256_file
from autodesign.video_delivery_validation import (
    CurrentVideoDeliveryValidation,
    VideoDeliverySnapshot,
    revalidate_current_video_delivery_snapshots,
    validate_current_video_delivery as _validate_current_video_delivery,
)
from autodesign.tools._contract import ToolContext
from autodesign.video_runtime import (
    video_environment_profile as _video_environment_profile,
)

_REPO_ROOT = Path(__file__).resolve().parents[1]


def _environment_profile(
    *,
    repo_root: Path | None = None,
    path_env: str | None = None,
) -> dict[str, Any]:
    codex = resolve_codex_runtime(required=("--ephemeral",))
    return {
        "video": _video_environment_profile(
            repo_root=repo_root,
            path_env=path_env,
        ),
        "coding_agent": {
            "harness": "codex",
            "ready": bool(codex["available"]),
            "binary": codex["binary"],
            "binary_source": codex["source"],
            "version": codex["version"],
            "capabilities": codex["capabilities"],
            "missing": codex["missing"],
            "rejected_candidates": codex["rejected_candidates"],
        },
    }


def _require_artifact_runtime(
    artifact_type: str,
    *,
    environment: dict[str, Any] | None = None,
) -> None:
    if artifact_type != "video":
        return
    profile = environment or {"video": _video_environment_profile()}
    video = profile.get("video") if isinstance(profile, dict) else None
    if isinstance(video, dict) and video.get("ready") is True:
        return
    missing = (
        list(video.get("missing") or [])
        if isinstance(video, dict)
        else ["video runtime diagnostics"]
    )
    repair = (
        str(video.get("repair") or "")
        if isinstance(video, dict)
        else "Run `autodesign doctor` and `autodesign setup`."
    )
    raise HTTPException(
        412,
        detail={
            "code": "video_runtime_unavailable",
            "message": (
                "Video generation is unavailable because required local "
                f"runtime components are missing: {', '.join(missing)}."
            ),
            "missing": missing,
            "repair": repair,
        },
    )

# ---------- Wire shapes (mirror web/src/lib/types.ts) ----------

ArtifactType = Literal["poster", "landing", "deck", "video"]
NativeFormat = Literal["svg", "html", "pptx", "mp4", "png"]
ViewFormat = Literal["svg", "html", "mp4", "png"]


def _validated_authoring_max_attempts(
    value: int | None,
    artifact_type: ArtifactType,
    settings: Settings,
) -> int:
    if value is None or not isinstance(value, int):
        return authoring_max_attempts_for(settings, artifact_type)
    if value < 1 or value > 12:
        raise HTTPException(
            422,
            detail={
                "code": "invalid_authoring_max_attempts",
                "message": "Authoring attempts must be an integer from 1 through 12.",
            },
        )
    return value


def _settings_with_authoring_max_attempts(
    settings: Settings,
    artifact_type: ArtifactType,
    value: int | None,
) -> Settings:
    resolved = _validated_authoring_max_attempts(value, artifact_type, settings)
    try:
        return replace(settings, authoring_max_attempts_override=resolved)
    except TypeError:
        copied_settings = copy(settings)
        setattr(copied_settings, "authoring_max_attempts_override", resolved)
        return copied_settings


_TRUTHY = frozenset({"1", "true", "yes", "on"})
_WEB_REFERENCE_POSTER_IMAGE_FORMATS = {
    ".png": "PNG",
    ".jpg": "JPEG",
    ".jpeg": "JPEG",
    ".webp": "WEBP",
}


def _validated_web_reference_poster_name(filename: str) -> str:
    name = Path(str(filename or "")).name
    suffix = Path(name).suffix.lower()
    if suffix not in _WEB_REFERENCE_POSTER_IMAGE_FORMATS:
        raise HTTPException(
            400,
            detail={
                "code": "unsupported_reference_poster_image",
                "message": "Reference style accepts PNG, JPEG, or WebP poster images only.",
            },
        )
    return name


def _validate_web_reference_poster_file(path: Path) -> None:
    name = _validated_web_reference_poster_name(path.name)
    expected_format = _WEB_REFERENCE_POSTER_IMAGE_FORMATS[Path(name).suffix.lower()]
    try:
        with warnings.catch_warnings():
            warnings.simplefilter("error", Image.DecompressionBombWarning)
            with Image.open(path) as image:
                detected_format = str(image.format or "").upper()
                image.verify()
            with Image.open(path) as image:
                image.load()
    except (
        Image.DecompressionBombError,
        Image.DecompressionBombWarning,
        OSError,
        UnidentifiedImageError,
        ValueError,
    ) as exc:
        raise HTTPException(
            400,
            detail={
                "code": "invalid_reference_poster_image",
                "message": "The reference poster is not a valid PNG, JPEG, or WebP image.",
            },
        ) from exc
    if detected_format != expected_format:
        raise HTTPException(
            400,
            detail={
                "code": "reference_poster_image_format_mismatch",
                "message": "The reference poster contents do not match its file extension.",
            },
        )


def _first_env_value(*names: str) -> str:
    for name in names:
        value = (os.getenv(name, "") or "").strip()
        if value:
            return value
    return ""


def _truthy_env(*names: str) -> bool:
    return _first_env_value(*names).lower() in _TRUTHY


def _int_env(names: str | tuple[str, ...], default: int, *, minimum: int = 0) -> int:
    env_names = (names,) if isinstance(names, str) else names
    raw = _first_env_value(*env_names)
    try:
        value = int(raw)
    except ValueError:
        return default
    return max(minimum, value)


_DEMO_MODE = _truthy_env(
    "AUTODESIGN_DEMO_MODE",
    "DESIGN_ANYTHING_DEMO_MODE",
    "DEMO_MODE",
)
_PUBLIC_USER_ISOLATION = _truthy_env(
    "AUTODESIGN_PUBLIC_USER_ISOLATION",
    "DESIGN_ANYTHING_PUBLIC_USER_ISOLATION",
    "PUBLIC_USER_ISOLATION",
)
_RUN_ACCESS_CONTROL = _DEMO_MODE or _PUBLIC_USER_ISOLATION
_DEMO_DAILY_LIMIT = _int_env(
    ("AUTODESIGN_DEMO_DAILY_LIMIT", "DESIGN_ANYTHING_DEMO_DAILY_LIMIT", "DEMO_DAILY_LIMIT"),
    3,
    minimum=1,
)
_DEMO_CONCURRENCY = _int_env(
    ("AUTODESIGN_DEMO_CONCURRENCY", "DESIGN_ANYTHING_DEMO_CONCURRENCY", "DEMO_CONCURRENCY"),
    1,
    minimum=1,
)
_DEMO_QUEUE_MAX = _int_env(
    ("AUTODESIGN_DEMO_QUEUE_MAX", "DESIGN_ANYTHING_DEMO_QUEUE_MAX", "DEMO_QUEUE_MAX"),
    8,
    minimum=1,
)
_DEMO_RUN_TTL_HOURS = _int_env(
    ("AUTODESIGN_DEMO_RUN_TTL_HOURS", "DESIGN_ANYTHING_DEMO_RUN_TTL_HOURS", "DEMO_RUN_TTL_HOURS"),
    24,
    minimum=1,
)
_DEMO_MAX_PDF_BYTES = _int_env(
    ("AUTODESIGN_DEMO_MAX_PDF_BYTES", "DESIGN_ANYTHING_DEMO_MAX_PDF_BYTES", "DEMO_MAX_PDF_BYTES"),
    25 * 1024 * 1024,
    minimum=1,
)
_DEMO_USER_HEADER = "x-demo-user"
_DEMO_FIXED_TEMPLATE = "cvpr-landscape"


class Canvas(BaseModel):
    w: int
    h: int
    background: str | None = None


class Artifact(BaseModel):
    artifact_id: str
    name: str
    artifact_type: ArtifactType
    canvas: Canvas
    canvas_plan: dict[str, Any] = Field(default_factory=dict)
    deck_plan: dict[str, Any] = Field(default_factory=dict)
    # Path 1 — direct render of agent's actual file. Frontend embeds this
    # via <img>/iframe rather than iterating over `layers`.
    native_file_url: str | None = None
    native_format: NativeFormat | None = None
    view_file_url: str | None = None
    view_format: ViewFormat | None = None
    download_url: str | None = None
    pdf_url: str | None = None
    downloads: dict[str, str] = Field(default_factory=dict)
    # PNG render of the artifact (vision-critic preview). Frontend uses
    # this for the chat thumbnail when present, falling back to an iframe.
    preview_url: str | None = None
    # Compact viewport thumbnail. Full-page QA/critic imagery remains in
    # preview_url and must not be cropped into a misleading card preview.
    card_preview_url: str | None = None
    quality_status: Literal["ready", "ready_with_warnings"] | None = None
    quality_diagnostics: list[str] = Field(default_factory=list)
    layers: list[dict[str, Any]] = []
    video_project: dict[str, Any] | None = None
    openresearch: dict[str, Any] | None = None
    parent_artifact_id: str | None = None
    candidate_draft: bool = False
    attempt_lineage: dict[str, Any] | None = None


class Failure(BaseModel):
    """Structured failure metadata so the frontend can render a real
    Failure Card instead of regex-parsing the assistant text. None of
    these fields are required — when artifact is None we always set at
    least `status` so the UI can branch."""

    status: str
    """Terminal run status, e.g. `max_turns`,
    `cancelled`, `error`. Falls back to `error` for exception cases."""

    phase: str | None = None
    """Best-effort coarse phase the run reached, e.g. `ingest`,
    `planning`, `rendering`. Inferred from on-disk artifacts."""

    error_code: str | None = None
    """Stable user-facing failure category, e.g. `provider_rate_limit`."""

    error_message: str | None = None
    """Concise, safe explanation of the interruption and the next action."""

    error_detail: str | None = None
    """Redacted coding-harness output retained for expandable diagnostics."""

    resume_available: bool = False
    """Whether the run has a validated external-author checkpoint."""

    resume_from_attempt: int | None = None
    """Attempt whose authored output will be used as the repair baseline."""

    next_attempt: int | None = None
    """Attempt number that a resumed run will create next."""

    retry_route: Literal[
        "full_authoring",
        "export_only",
        "setup_required",
        "none",
    ] | None = None
    """Machine-readable retry boundary for failed derived operations."""

    parent_run_id: str | None = None
    """Source run to use when a derived operation needs a bounded fallback."""

    agent_last_note: str | None = None
    """The designer's last natural-language emission, if any."""

    pointer_cleanup_warnings: list[str] = Field(default_factory=list)
    """Validated diagnostics retained when final-pointer cleanup was incomplete."""

    produced_files: list[str] = []
    """Files that did land on disk under the run dir, relative to it.
    Same enumerator as `_list_produced_artifacts`."""

    suggested_designer: str | None = None
    """If we know a model that's likely to do better next time
    (e.g., switch off Kimi K2.6 for paper-poster), surface it so the
    Retry CTA can name it. None when no concrete suggestion."""

    suggested_planner: str | None = None
    """Deprecated compatibility alias for `suggested_designer`."""

    elapsed_ms: int | None = None
    """Wall-clock spent inside the run, in milliseconds. Helps the UX
    say "stalled after 4m 12s" instead of just "stalled"."""

    critic_verdict: str | None = None
    """Vision critic's verdict — `pass` / `revise` / `fail`. Populated
    for degraded runs (artifact published but agent self-graded as
    sub-pass) so the quality-warning banner can show context."""

    critic_score: float | None = None
    """Vision critic's numeric score in [0, 1]. Pairs with verdict for
    the warning banner."""


class Message(BaseModel):
    id: str
    role: Literal["user", "assistant"]
    text: str
    ts: int
    run_id: str | None = None
    artifact_id: str | None = None
    status: Literal["streaming", "done", "error"] | None = None
    failure: Failure | None = None
    download_url: str | None = None
    download_filename: str | None = None
    download_mime_type: str | None = None
    task_type: str | None = None
    task_payload: dict[str, Any] | None = None
    source_artifact_id: str | None = None


class GenerateResponse(BaseModel):
    message: Message
    artifact: Artifact | None = None


class GenerateAck(BaseModel):
    """Returned immediately from /api/generate. The actual run continues
    in the background; the frontend should open an EventSource on
    /api/runs/{run_id}/events and then GET /api/runs/{run_id}/artifact
    once a `done` SSE event arrives."""
    run_id: str
    placeholder_message: Message
    progress_mode: str | None = None
    reference_poster_handle: str | None = None
    start_token: str | None = None


class RunInputSlotRequest(BaseModel):
    name: str
    role: Literal["attachment", "reference_poster"] = "attachment"
    sha256: str
    size: int = Field(ge=0)


class RunReserveRequest(BaseModel):
    brief: str
    artifact_type: str | None = None
    palette_id: str | None = None
    baseline_artifact: str | None = None
    conversation_history: str | None = None
    prior_artifacts: str | None = None
    conversation_id: str | None = None
    template: str | None = None
    canvas_preset_id: str | None = None
    authoring_max_attempts: int | None = None
    input_slots: list[RunInputSlotRequest] = Field(default_factory=list)


class RunReserveResponse(BaseModel):
    run_id: str
    upload_token: str
    input_slots: list[RunInputSlotRequest]
    request_digest: str
    run_state: str
    expires_at: float
    reused: bool = False


class RunUploadResponse(BaseModel):
    run_id: str
    slot: str
    sha256: str
    size: int
    run_state: str
    idempotent: bool = False


class RunStatusResponse(BaseModel):
    run_id: str
    run_state: str
    revision: int
    publishable: bool
    cancellation_pending: str | None = None
    worker_pid: int | None = None
    terminal_event: str | None = None


class RunCancelResponse(BaseModel):
    run_id: str
    status: Literal[
        "cancelled", "already_cancelled", "already_terminal", "cancellation_pending"
    ]
    run_state: str
    confirmed: bool
    terminated_pids: list[int] = Field(default_factory=list)
    surviving_pids: list[int] = Field(default_factory=list)


class PaperBundleCreateRequest(BaseModel):
    job_id: str
    conversation_id: str
    source_name: str
    prompt_version: str
    children: dict[str, RunReserveRequest]


class DesignEventRequest(BaseModel):
    conversation_id: str
    event: str
    run_id: str | None = None
    artifact_id: str | None = None
    data: dict[str, Any] = Field(default_factory=dict)


class EditorAssetUploadResponse(BaseModel):
    url: str
    filename: str
    content_type: str
    size: int


class ArtifactAsset(BaseModel):
    asset_id: str
    name: str
    kind: Literal["figure", "table", "image"]
    url: str
    filename: str
    run_id: str
    source: str
    size: int


class VideoRenderRequest(BaseModel):
    artifact: dict[str, Any]
    conversation_id: str | None = None


class PosterCodeEditRequest(BaseModel):
    artifact: dict[str, Any]
    instruction: str
    palette_id: str
    conversation_id: str | None = None
    conversation_history: list[dict[str, Any]] = Field(default_factory=list)
    source_run_id: str | None = None
    selection_context: dict[str, Any] | None = None


class ArtifactExportRequest(BaseModel):
    artifact: dict[str, Any]
    format: Literal["pdf", "pptx", "original_html", "standalone_html"]


class ArtifactPptxExportRequest(BaseModel):
    artifact: dict[str, Any]
    conversation_id: str | None = None


class CodingAgentSmokeRequest(BaseModel):
    timeout_s: int = Field(default=60, ge=1, le=90)


class AttemptSelectionRequest(BaseModel):
    idempotency_key: str = Field(min_length=1, max_length=128)
    expected_candidate_sha256: str = Field(min_length=64, max_length=64)


class AttemptForkRequest(BaseModel):
    conversation_id: str | None = None


class CandidateDraftPublishRequest(BaseModel):
    conversation_id: str | None = None


class DirectAttemptPublishRequest(BaseModel):
    conversation_id: str | None = None
    expected_candidate_sha256: str = Field(min_length=64, max_length=64)
    idempotency_key: str = Field(min_length=1, max_length=128)


class ArtifactExportResponse(BaseModel):
    url: str
    filename: str
    format: str
    mime_type: str


class OpenResearchProjectRequest(BaseModel):
    artifact: dict[str, Any] | None = None
    artifact_id: str | None = None
    source_run_id: str | None = None
    conversation_id: str | None = None
    org_id: str | None = None
    paper_id: str | None = None
    paper_url: str | None = None
    repo_full_name: str | None = None
    agent_prompt: str | None = None


class OpenResearchProjectAck(BaseModel):
    job_id: str
    status: Literal["running"]


class HarnessMatrixHarnessRequest(BaseModel):
    id: Literal["codex", "claude", "deepseek", "opencode", "kimi", "mimo", "pi", "zcode"]
    model: str | None = None


class HarnessMatrixRequest(BaseModel):
    paper_path: str
    prompt: str
    template: str = "cvpr-landscape"
    harnesses: list[HarnessMatrixHarnessRequest] = Field(default_factory=list)
    attempts: int = Field(default=12, ge=1, le=24)
    timeout_s: int = Field(default=3600, ge=60, le=14400)
    concurrency: Literal["by_harness"] = "by_harness"
    reuse_ingest_run: str | None = None


class HarnessMatrixAck(BaseModel):
    matrix_id: str
    status: Literal["running", "completed", "cancelled", "error"]
    matrix: dict[str, Any]


class OpenResearchProjectResult(BaseModel):
    job_id: str
    source_run_id: str
    artifact_id: str
    status: Literal["running", "submitted", "error"]
    project_id: str | None = None
    project_url: str | None = None
    org_id: str | None = None
    paper_id: str | None = None
    repo_full_name: str | None = None
    gui_submitter_status: str | None = None
    gui_submitter_reason: str | None = None
    gui_submitter_error: str | None = None
    gui_submitter_session_url: str | None = None
    agent_prompt_url: str | None = None
    submitter_log_url: str | None = None
    latest_report_id: str | None = None
    latest_report_url: str | None = None
    latest_report_markdown: str | None = None
    result_url: str | None = None
    api_log_url: str | None = None
    error: str | None = None
    details: dict[str, Any] = Field(default_factory=dict)


class HistoryConversationRequest(BaseModel):
    conversation: dict[str, Any]


# ---------- App ----------

app = FastAPI(title="AutoDesign web shim", version="0.3.0")

# Boot tolerantly: if the host has no .env credentials, we still come up
# (so the user can drop keys into the Settings UI and POST). The actual
# enforcement is per-request via _settings_for_request.
SETTINGS: "Settings | None"  # type: ignore[name-defined]
try:
    SETTINGS = load_settings()
    log("web.boot", out_dir=str(SETTINGS.out_dir),
        designer_model=SETTINGS.designer_model,
        image_model=SETTINGS.image_model)
except RuntimeError as e:
    SETTINGS = None
    log("web.boot.no_credential", reason=str(e)[:200])

# We still need a writable runs dir for the static mount and per-request
# Settings construction. Resolve it from a heuristic when there's no
# bootstrap Settings (rare path — only first-time users without .env).
_BOOT_OUT_DIR = SETTINGS.out_dir if SETTINGS else (Path.cwd() / "out")
RUNS_DIR = _BOOT_OUT_DIR / "runs"
RUNS_DIR.mkdir(parents=True, exist_ok=True)

EDITOR_ASSETS_DIR = _BOOT_OUT_DIR / "editor_assets"
EDITOR_ASSETS_DIR.mkdir(parents=True, exist_ok=True)
app.mount(
    "/api/files/editor-assets",
    StaticFiles(directory=EDITOR_ASSETS_DIR),
    name="editor_assets",
)

# Per-request upload staging area. Each request gets its own subdir so we
# can clean up a single failed run without touching others.
UPLOADS_DIR = _BOOT_OUT_DIR / "uploads"
UPLOADS_DIR.mkdir(parents=True, exist_ok=True)
UPLOADS_INDEX_PATH = UPLOADS_DIR / "conversation_attachments.json"
_UPLOADS_INDEX_LOCK = threading.RLock()


def _immutable_pointer_cleanup_warnings(raw: Any) -> tuple[str, ...]:
    if type(raw) not in {list, tuple} or not all(
        type(warning) is str for warning in raw
    ):
        return ()
    return tuple(raw)


def _validated_json_pointer_cleanup_warnings(raw: Any) -> tuple[str, ...]:
    if type(raw) is not list or not all(
        type(warning) is str for warning in raw
    ):
        return ()
    return tuple(raw)


@dataclass
class _WebRunRuntime:
    runs_dir: Path
    control_store: RunControlStore
    supervisor: RunSupervisor
    services: WebRunServices


_WEB_RUN_RUNTIME: _WebRunRuntime | None = None
_PAPER_BUNDLE_STORE: PaperBundleJobStore | None = None
_PAPER_BUNDLE_STORE_DIR: Path | None = None
_WEB_SERVER_LOCK_HANDLE: Any | None = None
_WEB_RUN_MONITOR_QUIESCE_TIMEOUT_S = 5.0


def _acquire_web_server_singleton_lock() -> None:
    global _WEB_SERVER_LOCK_HANDLE
    if _WEB_SERVER_LOCK_HANDLE is not None:
        raise RuntimeError("AutoDesign Web server singleton lock is already held")
    _BOOT_OUT_DIR.mkdir(parents=True, exist_ok=True)
    lock_path = _BOOT_OUT_DIR / ".autodesign-web-server.lock"
    handle = lock_path.open("a+b")
    try:
        if os.name == "posix":
            import fcntl

            fcntl.flock(handle.fileno(), fcntl.LOCK_EX | fcntl.LOCK_NB)
        else:
            import msvcrt

            handle.seek(0, os.SEEK_END)
            if handle.tell() == 0:
                handle.write(b"\0")
                handle.flush()
            handle.seek(0)
            msvcrt.locking(handle.fileno(), msvcrt.LK_NBLCK, 1)
    except (BlockingIOError, OSError) as exc:
        handle.close()
        raise RuntimeError(
            f"another AutoDesign Web server already owns {_BOOT_OUT_DIR}"
        ) from exc
    _WEB_SERVER_LOCK_HANDLE = handle


def _release_web_server_singleton_lock() -> None:
    global _WEB_SERVER_LOCK_HANDLE
    handle = _WEB_SERVER_LOCK_HANDLE
    if handle is None:
        return
    try:
        if os.name == "posix":
            import fcntl

            fcntl.flock(handle.fileno(), fcntl.LOCK_UN)
        else:
            import msvcrt

            handle.seek(0)
            msvcrt.locking(handle.fileno(), msvcrt.LK_UNLCK, 1)
    finally:
        handle.close()
        _WEB_SERVER_LOCK_HANDLE = None


def _reconcile_run_terminal_artifact(request: TerminalReconciliation) -> None:
    """Commit or roll back adapter-local final-directory promotion."""
    if request.phase != "commit":
        _reconcile_candidate_publish_terminal(request)
        return
    artifact_name = {
        "poster": "poster",
        "deck": "slides",
        "slides": "slides",
        "landing": "landing",
        "video": "video",
    }.get(request.record.artifact_type)
    if artifact_name is None:
        return
    run_dir = RUNS_DIR / request.run_id
    accept = request.decision == "accept"
    if artifact_name == "video" and accept:
        validate_video_delivery_context_promotion(run_dir)
    reconcile_artifact_promotion(
        run_dir / "final",
        artifact_name=artifact_name,
        accept=accept,
    )
    if artifact_name == "video":
        reconcile_video_delivery_context_promotion(run_dir, accept=accept)
    _reconcile_candidate_publish_terminal(request)


def _candidate_publish_source_disposition(
    source_run_id: str,
) -> Literal["mutable", "immutable_failed", "cancelled"]:
    source_dir = RUNS_DIR / source_run_id
    control_path = source_dir / "run_control.json"
    if control_path.is_file():
        record = RunControlStore(RUNS_DIR).read(source_run_id)
        if record.state in {"cancelling", "cancelled"} or record.writes_frozen:
            return "cancelled"
        if record.state == "failed":
            return "immutable_failed"
        return "mutable"

    terminal_event, _owner_pid, started = _disk_run_lifecycle(source_dir)
    if terminal_event == "run.cancelled":
        return "cancelled"
    if terminal_event == "run.error":
        return "immutable_failed"
    completion = _history_run_completion_from_disk(source_dir)
    terminal_status = str(completion.get("terminal_status") or "").strip().lower()
    if terminal_status in {"fail", "abort", "max_turns"} and (
        terminal_event == "run.done" or not started
    ):
        return "immutable_failed"
    return "mutable"


def _validate_canvas_candidate_publish_source(
    direct_publish: dict[str, Any],
    artifact_type: str,
    *,
    source_draft_run_id: str | None = None,
) -> str:
    source_draft_run_id = str(
        source_draft_run_id
        or direct_publish.get("source_draft_run_id")
        or ""
    )
    if not source_draft_run_id:
        raise ValueError("canvas candidate publish source draft is missing")
    try:
        validate_run_id(source_draft_run_id)
    except RunControlError as exc:
        raise ValueError(
            "canvas candidate publish source draft identity is unsafe"
        ) from exc
    draft_lineage = _read_json_file(
        RUNS_DIR / source_draft_run_id / "candidate_draft_lineage.json"
    )
    if (
        not isinstance(draft_lineage, dict)
        or draft_lineage.get("status") != "draft"
        or draft_lineage.get("artifact_type") != artifact_type
        or draft_lineage.get("source_run_id") != direct_publish["source_run_id"]
        or draft_lineage.get("source_attempt") != direct_publish["source_attempt"]
        or draft_lineage.get("source_candidate_id")
        != direct_publish["source_candidate_id"]
        or draft_lineage.get("source_candidate_sha256")
        != direct_publish["source_candidate_sha256"]
    ):
        raise ValueError("canvas candidate publish draft lineage is inconsistent")
    draft_descriptor = _read_derived_job_descriptor(source_draft_run_id)
    if (
        draft_descriptor is None
        or draft_descriptor["run_id"] != source_draft_run_id
        or draft_descriptor["artifact_type"] != artifact_type
        or draft_descriptor["job_kind"]
        not in {
            "artifact_edit",
            "attempt_fork",
            "editable_video_render",
            "poster_code_edit",
        }
    ):
        raise ValueError("canvas candidate publish draft descriptor is inconsistent")
    try:
        draft_control = RunControlStore(RUNS_DIR).read(source_draft_run_id)
    except RunControlError as exc:
        raise ValueError("canvas candidate publish draft control is unavailable") from exc
    if (
        draft_control.run_id != source_draft_run_id
        or draft_control.artifact_type != artifact_type
        or draft_control.parent_job_id != draft_descriptor["parent_run_id"]
        or draft_control.state != "completed"
        or not draft_control.publishable
        or draft_control.writes_frozen
    ):
        raise ValueError("canvas candidate publish draft control is inconsistent")
    return source_draft_run_id


def _reconcile_candidate_publish_terminal(
    request: TerminalReconciliation,
) -> None:
    descriptor_path = RUNS_DIR / request.run_id / "derived_job.json"
    raw_descriptor = _read_json_file(descriptor_path)
    if not isinstance(raw_descriptor, dict) or raw_descriptor.get("job_kind") != "candidate_publish":
        return
    if request.decision != "accept" or request.terminal_state != "completed":
        return
    descriptor = _read_derived_job_descriptor(request.run_id)
    if descriptor is None:
        raise ValueError("candidate publish descriptor is missing")
    lineage_path = RUNS_DIR / request.run_id / "candidate_draft_lineage.json"
    lineage = _read_json_file(lineage_path)
    if not isinstance(lineage, dict) or lineage.get("status") not in {
        "validated",
        "published",
    }:
        raise ValueError("candidate publish lineage is not validated")
    source_run_id = str(lineage.get("source_run_id") or "")
    source_attempt = lineage.get("source_attempt")
    source_sha256 = str(lineage.get("source_candidate_sha256") or "")
    if (
        not source_run_id
        or not isinstance(source_attempt, int)
        or source_attempt <= 0
        or len(source_sha256) != 64
    ):
        raise ValueError("candidate publish source lineage is incomplete")
    artifact_id = f"art_{request.run_id}"
    direct_publish = _read_direct_candidate_publish_descriptor(request.run_id)
    if direct_publish is not None and direct_publish["version"] in {2, 3}:
        bundle_artifact_type = str(
            direct_publish["paper_bundle_artifact_type"]
        )
        source_draft_run_id: str | None = None
        immediate_parent_run_id = source_run_id
        if direct_publish["version"] == 3:
            source_draft_run_id = _validate_canvas_candidate_publish_source(
                direct_publish,
                bundle_artifact_type,
            )
            immediate_parent_run_id = source_draft_run_id
        if (
            descriptor["job_kind"] != "candidate_publish"
            or descriptor["parent_run_id"] != immediate_parent_run_id
            or descriptor["artifact_type"] != bundle_artifact_type
            or request.record.run_id != request.run_id
            or request.record.artifact_type != bundle_artifact_type
            or request.record.parent_job_id != immediate_parent_run_id
            or direct_publish["run_id"] != request.run_id
            or direct_publish["source_run_id"] != source_run_id
            or direct_publish["source_attempt"] != source_attempt
            or direct_publish["source_candidate_id"]
            != str(lineage.get("source_candidate_id") or "")
            or direct_publish["source_candidate_sha256"] != source_sha256
            or str(lineage.get("artifact_type") or "") != bundle_artifact_type
        ):
            raise ValueError("bundle candidate publish lineage is inconsistent")
        if _RUN_ACCESS_CONTROL:
            owner_id = str(direct_publish["paper_bundle_owner_id"])
            owned_run_ids = [source_run_id, request.run_id]
            if source_draft_run_id is not None:
                owned_run_ids.insert(1, source_draft_run_id)
            for owned_run_id in owned_run_ids:
                access = _demo_run_access(owned_run_id)
                if (
                    access is None
                    or str(access.get("owner") or "") != owner_id
                ):
                    raise ValueError(
                        "bundle candidate publish owner is inconsistent"
                    )
        candidate = load_attempt_candidate(
            RUNS_DIR / source_run_id,
            source_attempt,
        )
        if (
            candidate.artifact_type.value != bundle_artifact_type
            or candidate.source_sha256 != source_sha256
            or candidate.candidate_id != direct_publish["source_candidate_id"]
        ):
            raise ValueError("bundle candidate publish source identity changed")
        if (
            direct_publish["version"] == 2
            and candidate.safety_state == "blocked"
        ):
            raise ValueError("blocked attempt candidate cannot be published")
        if request.phase != "commit":
            return
        store = _paper_bundle_store()
        store.read_owned(
            str(direct_publish["paper_bundle_job_id"]),
            str(direct_publish["paper_bundle_owner_id"]),
            child_status_provider=_paper_bundle_child_snapshot,
        )
        artifact = _build_artifact_response(
            RUNS_DIR / request.run_id,
            request.run_id,
            bundle_artifact_type,  # type: ignore[arg-type]
            baseline_artifact_json=str(
                descriptor.get("baseline_artifact_json") or ""
            ),
        )
        if (
            artifact is None
            or artifact.artifact_id != artifact_id
            or artifact.artifact_type != bundle_artifact_type
        ):
            raise ValueError("bundle candidate publish final artifact is invalid")
        try:
            commit = store.commit_child_publication(
                str(direct_publish["paper_bundle_job_id"]),
                str(direct_publish["paper_bundle_owner_id"]),
                bundle_artifact_type,
                source_run_id,
                publication_run_id=request.run_id,
                artifact_id=artifact_id,
                source_attempt=source_attempt,
                source_candidate_id=candidate.candidate_id,
                source_candidate_sha256=candidate.source_sha256,
                generation=int(direct_publish["publication_generation"]),
            )
            reconciliation_status = commit.status
        except PaperBundleBarrierClosed:
            reconciliation_status = "bundle_cancelled_not_attached"
        if lineage.get("status") != "published":
            lineage = {
                **lineage,
                "status": "published",
                "published_at": datetime.now(timezone.utc).isoformat(),
            }
            atomic_write_json(lineage_path, lineage)
        atomic_write_json(
            RUNS_DIR / request.run_id / "candidate_publish_reconciliation.json",
            {
                "run_id": request.run_id,
                "source_run_id": source_run_id,
                "artifact_id": artifact_id,
                "paper_bundle_job_id": direct_publish["paper_bundle_job_id"],
                "paper_bundle_artifact_type": bundle_artifact_type,
                "publication_generation": direct_publish[
                    "publication_generation"
                ],
                "status": reconciliation_status,
            },
        )
        return
    canvas_source_draft_run_id: str | None = None
    if (
        direct_publish is not None
        and direct_publish["version"] == 1
        and descriptor["parent_run_id"] != source_run_id
    ):
        canvas_source_draft_run_id = _validate_canvas_candidate_publish_source(
            direct_publish,
            descriptor["artifact_type"],
            source_draft_run_id=descriptor["parent_run_id"],
        )
        if (
            request.record.run_id != request.run_id
            or request.record.artifact_type != descriptor["artifact_type"]
            or request.record.parent_job_id != canvas_source_draft_run_id
            or direct_publish["run_id"] != request.run_id
            or direct_publish["source_run_id"] != source_run_id
            or direct_publish["source_attempt"] != source_attempt
            or direct_publish["source_candidate_id"]
            != str(lineage.get("source_candidate_id") or "")
            or direct_publish["source_candidate_sha256"] != source_sha256
            or str(lineage.get("artifact_type") or "")
            != descriptor["artifact_type"]
        ):
            raise ValueError("canvas candidate publish lineage is inconsistent")
    source_disposition = _candidate_publish_source_disposition(source_run_id)
    source_cancellation_frozen = source_disposition == "cancelled"
    if source_disposition == "immutable_failed":
        candidate = load_attempt_candidate(
            RUNS_DIR / source_run_id,
            source_attempt,
        )
        if (
            candidate.source_sha256 != source_sha256
            or candidate.candidate_id
            != str(lineage.get("source_candidate_id") or "")
        ):
            raise ValueError("candidate publish source identity changed")
        if (
            canvas_source_draft_run_id is None
            and candidate.safety_state == "blocked"
        ):
            raise ValueError("blocked attempt candidate cannot be published")
        if request.phase != "commit":
            return
        if lineage.get("status") != "published":
            lineage = {
                **lineage,
                "status": "published",
                "published_at": datetime.now(timezone.utc).isoformat(),
            }
            atomic_write_json(lineage_path, lineage)
        atomic_write_json(
            RUNS_DIR / request.run_id / "candidate_publish_reconciliation.json",
            {
                "run_id": request.run_id,
                "source_run_id": source_run_id,
                "artifact_id": artifact_id,
                "status": "source_immutable_failed",
            },
        )
        return
    committed_selection = load_selection_journal(RUNS_DIR / source_run_id)
    selection_matches = bool(
        committed_selection is not None
        and committed_selection.state == "complete"
        and committed_selection.artifact_id == artifact_id
        and committed_selection.source_attempt == source_attempt
        and committed_selection.candidate_sha256 == source_sha256
    )
    if source_cancellation_frozen:
        if lineage.get("status") != "published" and not selection_matches:
            raise ValueError("candidate publish source is cancelling")
        if request.phase != "commit":
            return
        if lineage.get("status") != "published":
            lineage = {
                **lineage,
                "status": "published",
                "published_at": datetime.now(timezone.utc).isoformat(),
            }
            atomic_write_json(lineage_path, lineage)
        atomic_write_json(
            RUNS_DIR / request.run_id / "candidate_publish_reconciliation.json",
            {
                "run_id": request.run_id,
                "source_run_id": source_run_id,
                "artifact_id": artifact_id,
                "status": "source_selection_already_committed",
            },
        )
        return
    if request.phase != "commit":
        return
    completion = complete_source_run_with_candidate_fork(
        run_dir=RUNS_DIR / source_run_id,
        run_id=source_run_id,
        attempt=source_attempt,
        expected_candidate_sha256=source_sha256,
        artifact_id=artifact_id,
    )
    if lineage.get("status") != "published":
        lineage = {
            **lineage,
            "status": "published",
            "published_at": datetime.now(timezone.utc).isoformat(),
        }
        atomic_write_json(lineage_path, lineage)
    atomic_write_json(
        RUNS_DIR / request.run_id / "candidate_publish_reconciliation.json",
        {
            "run_id": request.run_id,
            "source_run_id": source_run_id,
            "artifact_id": artifact_id,
            "status": completion,
        },
    )


async def _quiesce_web_completion_monitor(
    run_id: str,
    *,
    timeout_s: float = _WEB_RUN_MONITOR_QUIESCE_TIMEOUT_S,
) -> bool:
    """Join every Web-side writer owned by a run before cancellation is terminal."""
    current = asyncio.current_task()
    deadline = asyncio.get_running_loop().time() + max(0.0, timeout_s)
    while True:
        async with _RUNS_LOCK:
            state = _RUNS.get(run_id)
            completion_task = state.task if state is not None else None
        with _WEB_RUN_START_GATE_LOCK:
            operation_tasks = tuple(
                task
                for task, operation_ids in _WEB_RUN_OPERATION_TASKS.items()
                if f"run:{run_id}:" in operation_ids
            )
        selection_task = _ATTEMPT_SELECTION_TASKS.get(run_id)
        pending = {
            task
            for task in (completion_task, selection_task, *operation_tasks)
            if task is not None and task is not current and not task.done()
        }
        if not pending:
            return True
        remaining = deadline - asyncio.get_running_loop().time()
        if remaining <= 0:
            return False
        try:
            done, still_pending = await asyncio.wait(
                pending,
                timeout=remaining,
                return_when=asyncio.ALL_COMPLETED,
            )
        except asyncio.CancelledError:
            if current is not None and current.cancelling():
                raise
            return False
        if still_pending:
            return False
        for task in done:
            try:
                task.result()
            except (asyncio.CancelledError, Exception):
                pass


def _web_run_runtime() -> _WebRunRuntime:
    global _WEB_RUN_RUNTIME
    resolved_runs_dir = RUNS_DIR.resolve()
    if (
        _WEB_RUN_RUNTIME is not None
        and _WEB_RUN_RUNTIME.runs_dir == resolved_runs_dir
    ):
        return _WEB_RUN_RUNTIME
    control_store = RunControlStore(resolved_runs_dir)
    supervisor = RunSupervisor(
        resolved_runs_dir,
        control_store=control_store,
        terminal_reconciler=_reconcile_run_terminal_artifact,
        cancellation_quiescer=_quiesce_web_completion_monitor,
    )
    _WEB_RUN_RUNTIME = _WebRunRuntime(
        runs_dir=resolved_runs_dir,
        control_store=control_store,
        supervisor=supervisor,
        services=WebRunServices(
            resolved_runs_dir,
            control_store=control_store,
            supervisor=supervisor,
        ),
    )
    return _WEB_RUN_RUNTIME


def _reset_web_run_runtime_for_tests() -> None:
    global _WEB_RUN_RUNTIME
    _WEB_RUN_RUNTIME = None


def _paper_bundle_store() -> PaperBundleJobStore:
    global _PAPER_BUNDLE_STORE, _PAPER_BUNDLE_STORE_DIR
    jobs_dir = (RUNS_DIR.parent / "paper-bundles").absolute()
    if _PAPER_BUNDLE_STORE is None or _PAPER_BUNDLE_STORE_DIR != jobs_dir:
        _PAPER_BUNDLE_STORE = PaperBundleJobStore(jobs_dir)
        _PAPER_BUNDLE_STORE_DIR = jobs_dir
    return _PAPER_BUNDLE_STORE


def _reset_paper_bundle_store_for_tests() -> None:
    global _PAPER_BUNDLE_STORE, _PAPER_BUNDLE_STORE_DIR
    _PAPER_BUNDLE_STORE = None
    _PAPER_BUNDLE_STORE_DIR = None


@app.on_event("startup")
async def _startup_web_server_singleton_lock() -> None:
    _acquire_web_server_singleton_lock()
    try:
        _open_web_run_start_gate()
    except BaseException:
        _release_web_server_singleton_lock()
        raise


WEB_HISTORY_PATH = _BOOT_OUT_DIR / "web_history.json"
_WEB_HISTORY_LOCK = threading.RLock()

DEMO_ACCESS_PATH = _BOOT_OUT_DIR / "demo_run_access.json"
DEMO_USAGE_PATH = _BOOT_OUT_DIR / "demo_usage.json"
_DEMO_ACCESS_LOCK = threading.RLock()
_DEMO_USAGE_LOCK = threading.RLock()
_RUN_FILE_TOKEN_COOKIE_PREFIX = "autodesign_run_token_"
_LEGACY_RUN_FILE_TOKEN_COOKIE_PREFIX = "designanything_run_token_"
@app.api_route("/api/files/runs/{rel_path:path}", methods=["GET", "HEAD"])
async def run_file(rel_path: str, request: Request) -> FileResponse:
    if _RUN_ACCESS_CONTROL:
        return _demo_run_file_response(rel_path, request)
    return _OpenedPublicRunFileResponse(_open_public_run_file(rel_path))


# Lock for the env-mutation path in _settings_for_request — see comment
# there. Held only across load_settings() (sub-100ms), not across the
# actual run. RLock so the auto-switch / retry paths can wrap an outer
# lock-held block around `_settings_for_request` without self-deadlock.
_SETTINGS_LOCK = threading.RLock()


# Per-request HTTP headers the frontend can set to override .env values.
# Each entry maps a (lowercased) header name to the env var name that
# `autodesign.config.load_settings` already reads. Adding a new field
# is one line here — no other backend code needs to change.
#
# Naming convention: `X-<thing>-Key` for credentials, `X-<thing>-Base`
# for endpoint URLs, `X-Model-<agent>` for per-agent model picks,
# `X-Provider-<agent>` for the rare case where a user wants to FORCE
# a specific backend (anthropic | openai_compat | gemini) when multiple
# credentials are configured.
_HEADER_TO_ENV: dict[str, str] = {
    # ---- Provider credentials ----
    "x-openrouter-key": "OPENROUTER_API_KEY",
    "x-anthropic-key": "ANTHROPIC_API_KEY",
    "x-openai-key": "OPENAI_COMPAT_API_KEY",
    "x-gemini-key": "GEMINI_API_KEY",
    # Custom OpenAI-compatible endpoints (Together / Fireworks / Moonshot
    # direct / DeepSeek direct / vLLM / etc.). The user supplies both the
    # key and the base URL.
    "x-custom-openai-base": "OPENAI_COMPAT_BASE_URL",
    "x-anthropic-base": "ANTHROPIC_BASE_URL",
    # ---- Per-agent model overrides ----
    "x-model-designer": "DESIGNER_MODEL",
    "x-model-planner": "DESIGNER_MODEL",
    "x-model-enhancer": "ENHANCER_MODEL",
    "x-model-claim-graph": "CLAIM_GRAPH_MODEL",
    "x-model-deck-outline": "DECK_OUTLINE_MODEL",
    "x-model-paper-memory": "PAPER_MEMORY_MODEL",
    "x-designer-author-model": "AUTODESIGN_DESIGNER_AUTHOR_MODEL",
    "x-planner-author-model": "AUTODESIGN_DESIGNER_AUTHOR_MODEL",
    "x-code-editor-model": "AUTODESIGN_CODE_EDITOR_MODEL",
    "x-model-critic": "CRITIC_MODEL",
    "x-model-composer": "COMPOSER_MODEL",
    "x-model-ingest": "INGEST_MODEL",
    "x-model-image": "IMAGE_MODEL",
    "x-model-image-fallback": "IMAGE_FALLBACK_MODEL",
    # ---- Local coding-agent harness overrides ----
    # Explicit API key for the harness CLI subprocess (claude → ANTHROPIC_API_KEY,
    # codex → OPENAI_API_KEY). NOT a pipeline provider credential, so it is
    # deliberately absent from _KEY_HEADERS and never satisfies the "has key" gate.
    "x-designer-author-key": "AUTODESIGN_DESIGNER_AUTHOR_API_KEY",
    "x-designer-author-harness": "AUTODESIGN_DESIGNER_AUTHOR_HARNESS",
    "x-code-editor-harness": "AUTODESIGN_CODE_EDITOR_HARNESS",
    # ---- OpenResearch submission ----
    "x-openresearch-submitter": "AUTODESIGN_OPENRESEARCH_SUBMITTER",
    "x-openresearch-org-id": "AUTODESIGN_OPENRESEARCH_ORG_ID",
    "x-openresearch-repo": "AUTODESIGN_OPENRESEARCH_REPO",
    "x-openresearch-token": "AUTODESIGN_OPENRESEARCH_TOKEN",
    "x-openresearch-submitter-timeout": "AUTODESIGN_OPENRESEARCH_SUBMITTER_TIMEOUT_SECONDS",
    # ---- Per-agent provider overrides (advanced) ----
    "x-provider-designer": "DESIGNER_PROVIDER",
    "x-provider-planner": "DESIGNER_PROVIDER",
    "x-provider-enhancer": "ENHANCER_PROVIDER",
    "x-provider-claim-graph": "CLAIM_GRAPH_PROVIDER",
    "x-provider-deck-outline": "DECK_OUTLINE_PROVIDER",
    "x-provider-paper-memory": "PAPER_MEMORY_PROVIDER",
    "x-provider-critic": "CRITIC_PROVIDER",
    "x-provider-composer": "COMPOSER_PROVIDER",
}

_MODEL_ENV_NAMES: frozenset[str] = frozenset({
    "DESIGNER_MODEL",
    "ENHANCER_MODEL",
    "CLAIM_GRAPH_MODEL",
    "DECK_OUTLINE_MODEL",
    "PAPER_MEMORY_MODEL",
    "AUTODESIGN_DESIGNER_AUTHOR_MODEL",
    "AUTODESIGN_CODE_EDITOR_MODEL",
    "CRITIC_MODEL",
    "COMPOSER_MODEL",
    "INGEST_MODEL",
    "IMAGE_MODEL",
    "IMAGE_FALLBACK_MODEL",
})

# Which header keys actually count as "credentials" (vs. base URLs or
# model picks). Used to decide whether the request can authenticate at
# all — model overrides without any key are still a setup error.
_KEY_HEADERS: frozenset[str] = frozenset({
    "x-openrouter-key", "x-anthropic-key", "x-openai-key", "x-gemini-key",
})


def _request_env_overrides(request: Request) -> tuple[dict[str, str], bool]:
    overrides: dict[str, str] = {}
    has_key = False
    for header_name, env_name in _HEADER_TO_ENV.items():
        v = (request.headers.get(header_name, "") or "").strip()
        if not v:
            continue
        if (
            header_name in {"x-model-planner", "x-provider-planner"}
            and env_name in overrides
        ):
            continue
        if env_name in _MODEL_ENV_NAMES:
            v = normalize_model_id(v)
        if _RUN_ACCESS_CONTROL and header_name in {
            "x-openresearch-submitter",
            "x-openresearch-submitter-timeout",
        }:
            continue
        if _RUN_ACCESS_CONTROL and header_name in {
            "x-custom-openai-base",
            "x-anthropic-base",
        }:
            raise HTTPException(
                400,
                detail={
                    "code": "custom_provider_url_disabled",
                    "message": (
                        "Custom provider endpoints are available only in "
                        "the loopback local installation."
                    ),
                },
            )
        overrides[env_name] = v
        if header_name in _KEY_HEADERS:
            has_key = True
    return overrides, has_key


def _apply_request_network_policy(settings: "Settings") -> "Settings":  # type: ignore[name-defined]
    if not _RUN_ACCESS_CONTROL:
        return settings
    return replace(
        settings,
        allow_private_network=False,
        allow_remote_image_urls=False,
    )


def _settings_for_request(request: Request) -> "Settings":  # type: ignore[name-defined]
    """Resolve a per-request `Settings` object.

    Resolution order: client-supplied headers override boot-time `.env`
    values. If no credential header is set AND there's no `.env`-derived
    SETTINGS, we fail with 412 + `code: "no_api_key"` so the frontend
    can pop the Settings drawer.

    The full mapping table is `_HEADER_TO_ENV` above — it covers
    credentials, custom base URLs, per-agent model overrides, and per-
    agent backend overrides.
    """
    if _DEMO_MODE:
        if SETTINGS is None:
            raise HTTPException(
                503,
                detail={
                    "code": "demo_not_configured",
                    "message": "Demo mode requires server-side provider credentials in .env.",
                },
            )
        return _apply_request_network_policy(_demo_settings(SETTINGS))

    overrides, has_key = _request_env_overrides(request)

    if not overrides:
        if SETTINGS is None:
            raise HTTPException(
                412,
                detail={
                    "code": "no_api_key",
                    "message": "API key required. Open Settings and paste "
                               "your OpenRouter key (https://openrouter.ai/keys).",
                },
            )
        return _apply_request_network_policy(SETTINGS)

    # If only model/provider/base-URL overrides came in (no key), the
    # request can still succeed iff the boot-time SETTINGS has creds.
    # Otherwise it's the same setup-required path.
    if not has_key and SETTINGS is None:
        raise HTTPException(
            412,
            detail={
                "code": "no_api_key",
                "message": "API key required — Settings has model overrides "
                           "but no credential. Paste a key under Providers.",
            },
        )

    # `load_settings()` requires OPENROUTER_API_KEY *or* ANTHROPIC_API_KEY
    # to bootstrap the Anthropic SDK construct. A user who configured
    # ONLY Custom OpenAI-compatible (Together / Fireworks / vLLM / …)
    # would fail that check even though every agent in their plan can be
    # routed through `openai_compat`. Inject a sentinel that satisfies
    # the bootstrap; the SDK never gets called when each agent's
    # `*_PROVIDER` is `openai_compat`.
    if (
        "OPENAI_COMPAT_API_KEY" in overrides
        and "OPENROUTER_API_KEY" not in overrides
        and "ANTHROPIC_API_KEY" not in overrides
        and not (os.environ.get("OPENROUTER_API_KEY") or os.environ.get("ANTHROPIC_API_KEY"))
    ):
        overrides["OPENROUTER_API_KEY"] = "sk-or-v1-bootstrap-placeholder"

    # Use load_settings()' native env-driven resolution so all the cascading
    # defaults (base URLs and model fallbacks) stay in lock-
    # step with the .env path. We mutate process env briefly under a lock,
    # then restore. Safe because the lock is held for sub-100ms and the
    # generated Settings object is immutable from then on.
    saved = {k: os.environ.get(k) for k in overrides}
    with _SETTINGS_LOCK:
        try:
            os.environ.update(overrides)
            return _apply_request_network_policy(load_settings())
        except RuntimeError as e:
            raise HTTPException(
                412,
                detail={"code": "no_api_key", "message": str(e)},
            ) from e
        finally:
            for k, v in saved.items():
                if v is None:
                    os.environ.pop(k, None)
                else:
                    os.environ[k] = v


def _settings_for_openresearch_request(request: Request) -> "Settings":  # type: ignore[name-defined]
    try:
        return _web_openresearch_settings(_settings_for_request(request))
    except HTTPException as exc:
        detail = exc.detail if isinstance(exc.detail, dict) else {}
        if exc.status_code != 412 or detail.get("code") != "no_api_key":
            raise
    overrides, _has_key = _request_env_overrides(request)
    if (
        "OPENROUTER_API_KEY" not in overrides
        and "ANTHROPIC_API_KEY" not in overrides
        and not (os.environ.get("OPENROUTER_API_KEY") or os.environ.get("ANTHROPIC_API_KEY"))
    ):
        overrides["OPENROUTER_API_KEY"] = "sk-or-v1-bootstrap-placeholder"
    saved = {k: os.environ.get(k) for k in overrides}
    with _SETTINGS_LOCK:
        try:
            os.environ.update(overrides)
            return _web_openresearch_settings(
                _apply_request_network_policy(load_settings())
            )
        finally:
            for k, v in saved.items():
                if v is None:
                    os.environ.pop(k, None)
                else:
                    os.environ[k] = v


def _settings_for_code_editor_request(request: Request) -> "Settings":  # type: ignore[name-defined]
    try:
        return _settings_for_request(request)
    except HTTPException as exc:
        detail = exc.detail if isinstance(exc.detail, dict) else {}
        if exc.status_code != 412 or detail.get("code") != "no_api_key":
            raise
    overrides, _has_key = _request_env_overrides(request)
    if (
        "OPENROUTER_API_KEY" not in overrides
        and "ANTHROPIC_API_KEY" not in overrides
        and not (os.environ.get("OPENROUTER_API_KEY") or os.environ.get("ANTHROPIC_API_KEY"))
    ):
        overrides["OPENROUTER_API_KEY"] = "sk-or-v1-bootstrap-placeholder"
    saved = {k: os.environ.get(k) for k in overrides}
    with _SETTINGS_LOCK:
        try:
            os.environ.update(overrides)
            return load_settings()
        finally:
            for k, v in saved.items():
                if v is None:
                    os.environ.pop(k, None)
                else:
                    os.environ[k] = v


# In-memory map of in-flight runs. Each entry holds the asyncio Task so
# /api/runs/{id}/events can know when the run finished, plus the artifact
# type the frontend asked for so /api/runs/{id}/artifact can resolve the
# right file from final/. Cleared after the artifact is fetched (or after
# 1 hour, whichever comes first — keeps the dict bounded).
class _RunState:
    __slots__ = ("artifact_type", "task", "created_at", "result_message",
                 "result_artifact", "error", "cancelled", "designer_model",
                 "has_pdf", "brief", "attach_paths", "reference_poster_path", "baseline_artifact_json",
                 "conversation_id", "template", "canvas_preset_id", "palette_id", "authoring_max_attempts",
                 "demo_user_id", "queued", "reservation_token", "input_slot_roles",
                 "reference_poster_handle")

    def __init__(
        self,
        artifact_type: ArtifactType,
        *,
        designer_model: str = "",
        has_pdf: bool = False,
        brief: str = "",
        attach_paths: list[Path] | None = None,
        reference_poster_path: Path | None = None,
        baseline_artifact_json: str | None = None,
        conversation_id: str = "",
        template: str | None = None,
        canvas_preset_id: str | None = None,
        palette_id: str | None = None,
        authoring_max_attempts: int | None = None,
        input_slot_roles: dict[str, str] | None = None,
    ) -> None:
        self.artifact_type: ArtifactType = artifact_type
        self.task: asyncio.Task | None = None
        self.created_at: float = time.time()
        self.result_message: Message | None = None
        self.result_artifact: Artifact | None = None
        self.error: str | None = None
        self.cancelled: bool = False
        # Captured at submit so the failure builder has the right model
        # name even if Settings was mutated later. Also lets the retry
        # suggestion know whether we were on Kimi.
        self.designer_model: str = designer_model
        self.has_pdf: bool = has_pdf
        # The full effective brief (post-prologue, post-conversation-context)
        # plus attachment paths and the baseline artifact JSON. These let
        # /api/runs/{run_id}/retry re-fire the run with a single click —
        # the user doesn't have to re-upload files or re-type the brief.
        self.brief: str = brief
        self.attach_paths: list[Path] = attach_paths or []
        self.reference_poster_path: Path | None = reference_poster_path
        self.baseline_artifact_json: str | None = baseline_artifact_json
        self.conversation_id: str = conversation_id
        self.template: str | None = template
        self.canvas_preset_id: str | None = canvas_preset_id
        self.palette_id: str | None = palette_id
        self.authoring_max_attempts: int | None = authoring_max_attempts
        self.demo_user_id: str = ""
        self.queued: bool = False
        self.reservation_token: str = ""
        self.input_slot_roles: dict[str, str] = dict(input_slot_roles or {})
        self.reference_poster_handle: str | None = None


_RUNS: dict[str, _RunState] = {}
_RUNS_LOCK = asyncio.Lock()
_PptxExportKey = tuple[str, str, str]
_PPTX_EXPORT_RUNS: dict[_PptxExportKey, str] = {}
_PPTX_EXPORT_LOCK = asyncio.Lock()


def _persisted_run_log(event: str, run_id: str, **data: Any) -> None:
    with run_context(run_id, RUNS_DIR / run_id):
        log(event, run_id=run_id, **data)


class _DemoQueuedRun:
    __slots__ = (
        "run_id", "brief", "attach_paths", "template", "a_type",
        "baseline_artifact_json", "state", "settings",
    )

    def __init__(
        self,
        *,
        run_id: str,
        brief: str,
        attach_paths: list[Path],
        template: str | None,
        a_type: ArtifactType,
        baseline_artifact_json: str | None,
        state: _RunState,
        settings: Settings,
    ) -> None:
        self.run_id = run_id
        self.brief = brief
        self.attach_paths = attach_paths
        self.template = template
        self.a_type = a_type
        self.baseline_artifact_json = baseline_artifact_json
        self.state = state
        self.settings = settings


_DEMO_RUN_QUEUE: asyncio.Queue[_DemoQueuedRun] | None = None
_DEMO_WORKERS: list[asyncio.Task] = []
_WEB_RUN_START_GATE_LOCK = threading.RLock()
_WEB_RUN_STARTS_BLOCKED = False
_WEB_RUN_START_TASKS: dict[asyncio.Task[Any], tuple[str, ...]] = {}
_WEB_RUN_OPERATION_TASKS: dict[asyncio.Task[Any], tuple[str, ...]] = {}
_RUN_TREE_LOCKS: dict[str, tuple[asyncio.AbstractEventLoop, asyncio.Lock]] = {}


def _run_tree_lock(run_id: str) -> asyncio.Lock:
    loop = asyncio.get_running_loop()
    with _WEB_RUN_START_GATE_LOCK:
        existing = _RUN_TREE_LOCKS.get(run_id)
        if existing is not None and existing[0] is loop:
            return existing[1]
        lock = asyncio.Lock()
        _RUN_TREE_LOCKS[run_id] = (loop, lock)
        return lock


def _derived_ancestor_chain(parent_run_id: str) -> tuple[str, ...]:
    """Return the persisted derived ancestry in root-to-parent order."""
    chain: list[str] = []
    seen: set[str] = set()
    current_run_id = parent_run_id
    store = _web_run_runtime().control_store
    while current_run_id:
        if current_run_id in seen:
            raise RunNotReady("derived run ancestry contains a cycle")
        seen.add(current_run_id)
        chain.append(current_run_id)
        try:
            descriptor = _read_derived_job_descriptor(current_run_id)
        except ValueError as exc:
            raise RunNotReady("derived run ancestry is unreadable") from exc
        if descriptor is None:
            break
        try:
            record = store.read(current_run_id)
        except RunControlError as exc:
            raise RunNotReady("derived run ancestry control is unreadable") from exc
        ancestor_run_id = str(descriptor["parent_run_id"])
        if record.parent_job_id != ancestor_run_id:
            raise RunNotReady("derived run ancestry does not match run control")
        current_run_id = ancestor_run_id
    return tuple(reversed(chain))


@asynccontextmanager
async def _derived_tree_locks(parent_run_id: str):
    async with AsyncExitStack() as stack:
        for ancestor_run_id in _derived_ancestor_chain(parent_run_id):
            await stack.enter_async_context(_run_tree_lock(ancestor_run_id))
        yield


def _open_web_run_start_gate() -> None:
    global _WEB_RUN_STARTS_BLOCKED
    with _WEB_RUN_START_GATE_LOCK:
        for task in tuple(_WEB_RUN_START_TASKS):
            if task.done():
                _WEB_RUN_START_TASKS.pop(task, None)
        for task in tuple(_WEB_RUN_OPERATION_TASKS):
            if task.done():
                _WEB_RUN_OPERATION_TASKS.pop(task, None)
        if _WEB_RUN_START_TASKS or _WEB_RUN_OPERATION_TASKS:
            raise RuntimeError("cannot reopen Web run starts while prior starts are active")
        _WEB_RUN_STARTS_BLOCKED = False


def _close_web_run_start_gate() -> None:
    global _WEB_RUN_STARTS_BLOCKED
    with _WEB_RUN_START_GATE_LOCK:
        _WEB_RUN_STARTS_BLOCKED = True


def _inflight_web_run_starts() -> set[str]:
    with _WEB_RUN_START_GATE_LOCK:
        return {
            run_id
            for task, run_ids in _WEB_RUN_START_TASKS.items()
            if not task.done()
            for run_id in run_ids
        }


def _inflight_web_run_operations() -> set[str]:
    with _WEB_RUN_START_GATE_LOCK:
        return {
            operation_id
            for task, operation_ids in _WEB_RUN_OPERATION_TASKS.items()
            if not task.done()
            for operation_id in operation_ids
        }


def _task_has_web_run_admission(task: asyncio.Task[Any]) -> bool:
    return task in _WEB_RUN_START_TASKS or task in _WEB_RUN_OPERATION_TASKS


@asynccontextmanager
async def _web_run_start_guard(run_id: str) -> AsyncIterator[None]:
    task = asyncio.current_task()
    if task is None:
        raise RunNotReady("run start requires an asyncio task")
    with _WEB_RUN_START_GATE_LOCK:
        if _WEB_RUN_STARTS_BLOCKED and not _task_has_web_run_admission(task):
            raise RunNotReady("Web server shutdown has blocked new run starts")
        _WEB_RUN_START_TASKS[task] = (*_WEB_RUN_START_TASKS.get(task, ()), run_id)
    try:
        yield
    finally:
        with _WEB_RUN_START_GATE_LOCK:
            run_ids = list(_WEB_RUN_START_TASKS.get(task, ()))
            if run_id in run_ids:
                run_ids.remove(run_id)
            if run_ids:
                _WEB_RUN_START_TASKS[task] = tuple(run_ids)
            else:
                _WEB_RUN_START_TASKS.pop(task, None)


@asynccontextmanager
async def _web_run_operation_guard(operation_id: str) -> AsyncIterator[None]:
    task = asyncio.current_task()
    if task is None:
        raise RunNotReady("Web operation requires an asyncio task")
    with _WEB_RUN_START_GATE_LOCK:
        if _WEB_RUN_STARTS_BLOCKED and not _task_has_web_run_admission(task):
            raise RunNotReady("Web server shutdown has blocked new run starts")
        _WEB_RUN_OPERATION_TASKS[task] = (
            *_WEB_RUN_OPERATION_TASKS.get(task, ()),
            operation_id,
        )
    try:
        yield
    finally:
        with _WEB_RUN_START_GATE_LOCK:
            operation_ids = list(_WEB_RUN_OPERATION_TASKS.get(task, ()))
            if operation_id in operation_ids:
                operation_ids.remove(operation_id)
            if operation_ids:
                _WEB_RUN_OPERATION_TASKS[task] = tuple(operation_ids)
            else:
                _WEB_RUN_OPERATION_TASKS.pop(task, None)


class _OpenResearchJobState:
    __slots__ = (
        "job_id", "source_run_id", "artifact_id", "conversation_id",
        "request", "task", "created_at", "result", "error",
    )

    def __init__(
        self,
        *,
        job_id: str,
        source_run_id: str,
        artifact_id: str,
        conversation_id: str,
        request: BaseModel,
    ) -> None:
        self.job_id = job_id
        self.source_run_id = source_run_id
        self.artifact_id = artifact_id
        self.conversation_id = conversation_id
        self.request = request
        self.task: asyncio.Task | None = None
        self.created_at = time.time()
        self.result: BaseModel | None = None
        self.error: str | None = None


_OPENRESEARCH_JOBS: dict[str, _OpenResearchJobState] = {}
_OPENRESEARCH_JOBS_LOCK = asyncio.Lock()


class _HarnessMatrixJobState:
    __slots__ = ("matrix_id", "matrix_dir", "task", "created_at", "cancel_event", "error")

    def __init__(self, *, matrix_id: str, matrix_dir: Path) -> None:
        self.matrix_id = matrix_id
        self.matrix_dir = matrix_dir
        self.task: asyncio.Task | None = None
        self.created_at = time.time()
        self.cancel_event = threading.Event()
        self.error: str | None = None


_HARNESS_MATRIX_JOBS: dict[str, _HarnessMatrixJobState] = {}
_HARNESS_MATRIX_JOBS_LOCK = asyncio.Lock()
_ATTEMPT_SELECTION_TASKS: dict[str, asyncio.Task[None]] = {}
_RUN_LIFECYCLE_TASK: asyncio.Task[None] | None = None
_WEB_RUN_SHUTDOWN_TIMEOUT_S = 10.0
_WEB_RUN_SHUTDOWN_POLL_S = 0.05


@app.on_event("startup")
async def _startup_demo_workers() -> None:
    if not _RUN_ACCESS_CONTROL:
        await asyncio.to_thread(_load_web_history_summaries)
    if SETTINGS is not None:
        for journal_path in RUNS_DIR.glob("*/attempt_candidates/selection.json"):
            run_dir = journal_path.parent.parent
            try:
                journal = load_selection_journal(run_dir)
            except (OSError, ValueError, ValidationError) as exc:
                log(
                    "attempt_selection.recovery_skipped",
                    run_id=run_dir.name,
                    error=f"{type(exc).__name__}: {exc}",
                )
                continue
            if journal is not None and journal.state in {
                "requested",
                "terminating",
                "promoting",
                "delivering",
            }:
                _schedule_attempt_selection_recovery(run_dir.name, SETTINGS)
    if not _DEMO_MODE:
        return
    _demo_cleanup_expired_runs()
    global _DEMO_RUN_QUEUE
    if _DEMO_RUN_QUEUE is None:
        _DEMO_RUN_QUEUE = asyncio.Queue(maxsize=_DEMO_QUEUE_MAX)
    if _DEMO_WORKERS:
        return
    for idx in range(_DEMO_CONCURRENCY):
        _DEMO_WORKERS.append(asyncio.create_task(_demo_queue_worker(idx + 1)))
    log(
        "demo.queue.started",
        concurrency=_DEMO_CONCURRENCY,
        max_queue=_DEMO_QUEUE_MAX,
        daily_limit=_DEMO_DAILY_LIMIT,
        run_ttl_hours=_DEMO_RUN_TTL_HOURS,
    )


_DERIVED_DESCRIPTOR_FIELDS = frozenset({
    "version",
    "job_kind",
    "run_id",
    "parent_run_id",
    "artifact_type",
    "conversation_id",
    "baseline_artifact_json",
    "source_artifact_id",
    "artifact_name",
    "source_relative_path",
})

_DIRECT_CANDIDATE_PUBLISH_FIELDS = frozenset({
    "version",
    "run_id",
    "source_run_id",
    "source_attempt",
    "source_candidate_id",
    "source_candidate_sha256",
    "idempotency_key_digest",
    "request_digest",
})
_BUNDLE_CANDIDATE_PUBLISH_FIELDS = frozenset({
    *_DIRECT_CANDIDATE_PUBLISH_FIELDS,
    "paper_bundle_job_id",
    "paper_bundle_owner_id",
    "paper_bundle_artifact_type",
    "publication_generation",
})
_CANVAS_BUNDLE_CANDIDATE_PUBLISH_FIELDS = frozenset({
    *_BUNDLE_CANDIDATE_PUBLISH_FIELDS,
    "source_draft_run_id",
})


def _read_direct_candidate_publish_descriptor(
    run_id: str,
) -> dict[str, Any] | None:
    path = RUNS_DIR / run_id / "candidate_publish_request.json"
    if not path.exists():
        return None
    if path.is_symlink():
        raise ValueError("candidate publish request must not be a symlink")
    try:
        payload = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError) as exc:
        raise ValueError("candidate publish request is unreadable") from exc
    if not isinstance(payload, dict):
        raise ValueError("candidate publish request fields are invalid")
    version = payload.get("version")
    expected_fields = (
        _DIRECT_CANDIDATE_PUBLISH_FIELDS
        if version == 1
        else _BUNDLE_CANDIDATE_PUBLISH_FIELDS
        if version == 2
        else _CANVAS_BUNDLE_CANDIDATE_PUBLISH_FIELDS
        if version == 3
        else None
    )
    if expected_fields is None or set(payload) != expected_fields:
        raise ValueError("candidate publish request fields are invalid")
    if payload.get("run_id") != run_id:
        raise ValueError("candidate publish request identity is invalid")
    if (
        type(payload.get("source_attempt")) is not int
        or payload["source_attempt"] <= 0
    ):
        raise ValueError("candidate publish source attempt is invalid")
    for field_name in _DIRECT_CANDIDATE_PUBLISH_FIELDS - {
        "version",
        "source_attempt",
    }:
        if not isinstance(payload.get(field_name), str) or not payload[field_name]:
            raise ValueError(
                f"candidate publish request {field_name} must be a non-empty string"
            )
    for field_name in (
        "source_candidate_sha256",
        "idempotency_key_digest",
        "request_digest",
    ):
        if re.fullmatch(r"[0-9a-f]{64}", payload[field_name]) is None:
            raise ValueError(
                f"candidate publish request {field_name} must be a SHA-256 digest"
            )
    for field_name in ("run_id", "source_run_id", "source_candidate_id"):
        try:
            validate_run_id(payload[field_name])
        except RunControlError as exc:
            raise ValueError(
                f"candidate publish request {field_name} is unsafe"
            ) from exc
    if version in {2, 3}:
        try:
            validate_run_id(payload["paper_bundle_job_id"])
        except RunControlError as exc:
            raise ValueError("candidate publish request bundle identity is unsafe") from exc
        owner_id = payload.get("paper_bundle_owner_id")
        if (
            not isinstance(owner_id, str)
            or not owner_id.strip()
            or len(owner_id) > 512
            or "\x00" in owner_id
        ):
            raise ValueError("candidate publish request bundle owner is invalid")
        if payload.get("paper_bundle_artifact_type") not in _PAPER_BUNDLE_ARTIFACT_TYPES:
            raise ValueError("candidate publish request bundle artifact type is invalid")
        generation = payload.get("publication_generation")
        if type(generation) is not int or generation <= 0:
            raise ValueError("candidate publish request publication generation is invalid")
    if version == 3:
        source_draft_run_id = payload.get("source_draft_run_id")
        try:
            validate_run_id(source_draft_run_id)
        except RunControlError as exc:
            raise ValueError(
                "candidate publish request source draft identity is unsafe"
            ) from exc
        if source_draft_run_id == payload["source_run_id"]:
            raise ValueError(
                "candidate publish request source draft must differ from the bundle child"
            )
    return payload


def _read_derived_job_descriptor(run_id: str) -> dict[str, Any] | None:
    path = RUNS_DIR / run_id / "derived_job.json"
    if not path.exists():
        return None
    if path.is_symlink():
        raise ValueError("derived job descriptor must not be a symlink")
    try:
        payload = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError) as exc:
        raise ValueError("derived job descriptor is unreadable") from exc
    if not isinstance(payload, dict) or set(payload) != _DERIVED_DESCRIPTOR_FIELDS:
        raise ValueError("derived job descriptor fields are invalid")
    if payload.get("version") != _DERIVED_JOB_VERSION:
        raise ValueError("derived job descriptor version is unsupported")
    for field_name in _DERIVED_DESCRIPTOR_FIELDS - {"version"}:
        if not isinstance(payload.get(field_name), str):
            raise ValueError(f"derived job descriptor {field_name} must be a string")
    if payload["run_id"] != run_id:
        raise ValueError("derived job descriptor run identity does not match")
    if payload["job_kind"] not in _DERIVED_JOB_KINDS:
        raise ValueError("derived job descriptor kind is unsupported")
    if not payload["parent_run_id"] or payload["parent_run_id"] == run_id:
        raise ValueError("derived job descriptor parent identity is invalid")
    relative = Path(payload["source_relative_path"])
    if relative.is_absolute() or ".." in relative.parts:
        raise ValueError("derived job descriptor source path is invalid")
    return payload


def _candidate_publish_reservation_digest(
    descriptor: dict[str, Any],
) -> str:
    reservation_descriptor = {
        "job_kind": descriptor["job_kind"],
        "source_artifact_id": descriptor["source_artifact_id"],
        "artifact_name": descriptor["artifact_name"],
        "source_relative_path": descriptor["source_relative_path"],
    }
    return hashlib.sha256(
        json.dumps(
            reservation_descriptor,
            ensure_ascii=False,
            sort_keys=True,
            separators=(",", ":"),
        ).encode("utf-8")
    ).hexdigest()


def _candidate_publish_recovery_request(
    run_id: str,
    descriptor: dict[str, Any],
    direct_publish: dict[str, Any],
    settings: Settings,
) -> tuple[CandidatePublishWorkerRequest, _RunState]:
    if (
        descriptor["job_kind"] != "candidate_publish"
        or direct_publish["run_id"] != run_id
        or descriptor["artifact_type"] not in _ARTIFACT_TYPES
    ):
        raise ValueError("candidate publish recovery identity is inconsistent")
    source_run_id = str(direct_publish["source_run_id"])
    source_attempt = int(direct_publish["source_attempt"])
    candidate = load_attempt_candidate(
        RUNS_DIR / source_run_id,
        source_attempt,
    )
    if (
        candidate.artifact_type.value != descriptor["artifact_type"]
        or candidate.candidate_id != direct_publish["source_candidate_id"]
        or candidate.source_sha256
        != direct_publish["source_candidate_sha256"]
    ):
        raise ValueError("candidate publish recovery source identity changed")

    direct_attempt = descriptor["parent_run_id"] == source_run_id
    if direct_publish["version"] == 2:
        if (
            not direct_attempt
            or direct_publish["paper_bundle_artifact_type"]
            != descriptor["artifact_type"]
        ):
            raise ValueError("bundle candidate publish recovery is inconsistent")
    elif direct_publish["version"] == 3:
        source_draft_run_id = _validate_canvas_candidate_publish_source(
            direct_publish,
            descriptor["artifact_type"],
        )
        if (
            descriptor["parent_run_id"] != source_draft_run_id
            or direct_publish["paper_bundle_artifact_type"]
            != descriptor["artifact_type"]
        ):
            raise ValueError("bundle canvas publish recovery is inconsistent")
        direct_attempt = False
    elif not direct_attempt:
        _validate_canvas_candidate_publish_source(
            direct_publish,
            descriptor["artifact_type"],
            source_draft_run_id=descriptor["parent_run_id"],
        )

    if direct_attempt and candidate.safety_state == "blocked":
        raise ValueError("blocked attempt candidate cannot be published")
    if direct_publish["version"] in {2, 3} and _RUN_ACCESS_CONTROL:
        owner_id = str(direct_publish["paper_bundle_owner_id"])
        owned_run_ids = [source_run_id, run_id]
        if direct_publish["version"] == 3:
            owned_run_ids.insert(1, str(direct_publish["source_draft_run_id"]))
        for owned_run_id in owned_run_ids:
            access = _demo_run_access(owned_run_id)
            if access is None or str(access.get("owner") or "") != owner_id:
                raise ValueError(
                    "bundle candidate publish recovery owner is inconsistent"
                )

    request = CandidatePublishWorkerRequest(
        job_kind="candidate_publish",
        run_id=run_id,
        parent_run_id=str(descriptor["parent_run_id"]),
        conversation_id=str(descriptor["conversation_id"]),
        settings=settings,
        source_attempt=source_attempt if direct_attempt else None,
        expected_candidate_sha256=(
            candidate.source_sha256 if direct_attempt else None
        ),
    )
    state = _RunState(
        artifact_type=descriptor["artifact_type"],
        brief=f"publish attempt {source_attempt}",
        conversation_id=str(descriptor["conversation_id"]),
        baseline_artifact_json=str(
            descriptor.get("baseline_artifact_json") or ""
        ),
    )
    return request, state


async def _recover_queued_candidate_publish_context(
    run_id: str,
    settings: Settings,
) -> tuple[_RunState, str, dict[str, Any]]:
    runtime = _web_run_runtime()
    record = runtime.control_store.read(run_id)
    if record.state != "queued":
        raise RunNotReady("candidate publish recovery requires a queued run")
    descriptor = _read_derived_job_descriptor(run_id)
    direct_publish = _read_direct_candidate_publish_descriptor(run_id)
    if descriptor is None or direct_publish is None:
        raise RunNotReady("candidate publish recovery metadata is unavailable")
    worker_request, state = _candidate_publish_recovery_request(
        run_id,
        descriptor,
        direct_publish,
        settings,
    )
    descriptor_path = RUNS_DIR / run_id / "derived_job.json"
    descriptor_sha256 = hashlib.sha256(descriptor_path.read_bytes()).hexdigest()
    reservation = await runtime.services.recover_queued_derived_reservation(
        run_id=run_id,
        artifact_type=state.artifact_type,
        parent_job_id=worker_request.parent_run_id,
        idempotency_key=f"derived:{run_id}",
        request_digest=_candidate_publish_reservation_digest(descriptor),
        descriptor_sha256=descriptor_sha256,
        settings=settings,
        payload={"request": worker_request},
    )
    state.reservation_token = reservation.upload_token
    if _RUN_ACCESS_CONTROL:
        access = _demo_run_access(run_id)
        state.demo_user_id = str(access.get("owner") or "") if access else ""
    async with _RUNS_LOCK:
        existing = _RUNS.get(run_id)
        if existing is not None and existing.task is not None and not existing.task.done():
            raise RunNotReady("candidate publish recovery conflicts with an active run")
        _RUNS[run_id] = state
    return state, reservation.upload_token, descriptor


def _decoded_worker_result_file(
    run_id: str,
    expected_job_kind: str,
) -> dict[str, Any] | None:
    path = RUNS_DIR / run_id / "worker_result.json"
    if path.is_symlink():
        return None
    try:
        envelope = parse_worker_result_json(path.read_text(encoding="utf-8"))
        return decode_worker_result(
            envelope,
            expected_run_id=run_id,
            expected_job_kind=expected_job_kind,
        )
    except (OSError, WorkerProtocolError):
        return None


def _decoded_pointer_cleanup_warnings(
    decoded: dict[str, Any],
) -> tuple[str, ...]:
    value = decoded.get("result") if decoded.get("ok") is True else decoded.get("error")
    if not isinstance(value, dict):
        return ()
    return _validated_json_pointer_cleanup_warnings(
        value.get("pointer_cleanup_warnings")
    )


def _worker_result_pointer_cleanup_warnings(
    run_id: str,
    expected_job_kind: str,
) -> tuple[str, ...]:
    decoded = _decoded_worker_result_file(run_id, expected_job_kind)
    return _decoded_pointer_cleanup_warnings(decoded) if decoded is not None else ()


def _derived_worker_pointer_cleanup_warnings(run_id: str) -> tuple[str, ...]:
    try:
        descriptor = _read_derived_job_descriptor(run_id)
    except (OSError, ValueError):
        return ()
    if descriptor is None:
        return ()
    return _worker_result_pointer_cleanup_warnings(
        run_id,
        str(descriptor["job_kind"]),
    )


def _recovered_worker_outcome(
    run_id: str,
    expected_job_kind: str = "pipeline",
) -> WorkerOutcome | None:
    decoded = _decoded_worker_result_file(run_id, expected_job_kind)
    if decoded is None:
        return None
    pointer_cleanup_warnings = _decoded_pointer_cleanup_warnings(decoded)
    if decoded.get("ok") is True:
        result = decoded.get("result")
        return WorkerOutcome(
            run_id=run_id, job_kind=expected_job_kind, returncode=0, ok=True,
            result=result if isinstance(result, dict) else None,
            error=None, relayed_events=0,
            pointer_cleanup_warnings=pointer_cleanup_warnings,
        )
    error = decoded.get("error")
    return WorkerOutcome(
        run_id=run_id, job_kind=expected_job_kind, returncode=1, ok=False,
        result=None,
        error=(
            format_worker_error_message(
                str(error.get("message") or error.get("type") or "worker failed"),
                pointer_cleanup_warnings,
            )
            if isinstance(error, dict)
            else "worker failed"
        ),
        relayed_events=0,
        failure_phase=(
            str(error.get("phase"))
            if isinstance(error, dict)
            and isinstance(error.get("phase"), str)
            and error.get("phase")
            else None
        ),
        pointer_cleanup_warnings=pointer_cleanup_warnings,
    )


async def _recover_web_run_controls() -> None:
    runtime = _web_run_runtime()
    bundle_store = _paper_bundle_store()
    cancellation_barriers = await asyncio.to_thread(
        bundle_store.recover_cancellation_barriers_after_restart,
    )
    for barrier in cancellation_barriers:
        for child_run_id in barrier.child_run_ids:
            try:
                runtime.control_store.request_cancel(child_run_id)
            except RunControlError:
                continue
    cancelled_roots: list[str] = []
    for control_path in sorted(RUNS_DIR.glob("*/run_control.json")):
        try:
            record = runtime.control_store.read(control_path.parent.name)
        except RunControlError:
            continue
        if record.state in {"cancelling", "cancelled"}:
            cancelled_roots.append(record.run_id)
    for root_run_id in cancelled_roots:
        for descendant_run_id in _controlled_derived_descendants(root_run_id):
            try:
                runtime.control_store.request_cancel(descendant_run_id)
            except RunControlError:
                continue
    for control_path in sorted(RUNS_DIR.glob("*/run_control.json")):
        run_id = control_path.parent.name
        try:
            record = runtime.control_store.read(run_id)
        except RunControlError:
            continue
        if record.state == "queued":
            raw_descriptor = _read_json_file(
                RUNS_DIR / run_id / "derived_job.json"
            )
            if (
                isinstance(raw_descriptor, dict)
                and raw_descriptor.get("job_kind") == "candidate_publish"
            ):
                try:
                    await _recover_queued_candidate_publish_context(
                        run_id,
                        SETTINGS or _runtime_only_settings(),
                    )
                except Exception as exc:  # noqa: BLE001
                    log(
                        "candidate_publish.recovery_failed",
                        run_id=run_id,
                        error=f"{type(exc).__name__}: {exc}",
                    )
                else:
                    continue
        await runtime.supervisor.recover(run_id)
        record = runtime.control_store.read(run_id)
        if record.state != "completing" or record.artifact_type not in _ARTIFACT_TYPES:
            continue
        try:
            descriptor = _read_derived_job_descriptor(run_id)
        except ValueError as exc:
            digest = hashlib.sha256(str(exc).encode("utf-8")).hexdigest()
            try:
                await runtime.supervisor.accept_completion(
                    run_id,
                    terminal_state="failed",
                    publishable=False,
                    result_digest=digest,
                )
            except Exception:
                pass
            continue
        if descriptor is not None and (
            descriptor["artifact_type"] != record.artifact_type
            or descriptor["parent_run_id"] != (record.parent_job_id or "")
        ):
            digest = hashlib.sha256(
                b"derived descriptor does not match run control"
            ).hexdigest()
            try:
                await runtime.supervisor.accept_completion(
                    run_id,
                    terminal_state="failed",
                    publishable=False,
                    result_digest=digest,
                )
            except Exception:
                pass
            continue
        expected_job_kind = (
            str(descriptor["job_kind"]) if descriptor is not None else "pipeline"
        )
        outcome = _recovered_worker_outcome(run_id, expected_job_kind)
        if outcome is None:
            continue
        state = _RunState(
            artifact_type=record.artifact_type,  # type: ignore[arg-type]
            brief="recovered supervised run",
            conversation_id=(
                str(descriptor["conversation_id"])
                if descriptor is not None
                else _event_conversation_id(None, run_id)
            ),
            baseline_artifact_json=(
                str(descriptor["baseline_artifact_json"])
                if descriptor is not None
                else None
            ),
        )
        async with _RUNS_LOCK:
            _RUNS.setdefault(run_id, state)
        if descriptor is not None:
            state.task = asyncio.create_task(
                _monitor_supervised_derived_job(
                    run_id=run_id,
                    state=state,
                    job_kind=expected_job_kind,
                    parent_run_id=str(descriptor["parent_run_id"]),
                    descriptor=descriptor,
                    recovered_outcome=outcome,
                )
            )
        else:
            state.task = asyncio.create_task(
                _monitor_supervised_pipeline(
                    run_id=run_id,
                    state=state,
                    recovered_outcome=outcome,
                )
            )
    for barrier in cancellation_barriers:
        if not barrier.pending_creation:
            continue
        recovered = await bundle_store.cancel_pending_creation(
            barrier.job_id,
            barrier.owner_id,
            cleanup_child=_cleanup_paper_bundle_child,
        )
        if recovered != "cancelled":
            raise PaperBundleError(
                f"paper bundle creation cancellation did not recover: {barrier.job_id}"
            )


async def _run_lifecycle_reaper() -> None:
    while True:
        try:
            runtime = _web_run_runtime()
            expired = await runtime.services.reconcile_expired_reservations()
            for run_id in expired:
                await _reconcile_paper_bundle_for_run(run_id)
            for control_path in sorted(RUNS_DIR.glob("*/run_control.json")):
                run_id = control_path.parent.name
                try:
                    record = runtime.control_store.read(run_id)
                except RunControlError:
                    continue
                if record.state == "cancelling":
                    await runtime.supervisor.recover(run_id)
                    await _reconcile_paper_bundle_for_run(run_id)
            await _reconcile_all_paper_bundles()
            await asyncio.sleep(1.0)
        except asyncio.CancelledError:
            raise
        except Exception as exc:  # noqa: BLE001
            log("web.run_lifecycle_reaper.error", error=type(exc).__name__)
            await asyncio.sleep(1.0)


@app.on_event("startup")
async def _startup_run_supervision() -> None:
    global _RUN_LIFECYCLE_TASK
    try:
        await _recover_web_run_controls()
        await _reconcile_all_paper_bundles()
        if _RUN_LIFECYCLE_TASK is None or _RUN_LIFECYCLE_TASK.done():
            _RUN_LIFECYCLE_TASK = asyncio.create_task(_run_lifecycle_reaper())
    except BaseException:
        _release_web_server_singleton_lock()
        raise


def _durable_nonterminal_run_ids(runtime: _WebRunRuntime) -> set[str]:
    run_ids: set[str] = set()
    for control_path in sorted(RUNS_DIR.glob("*/run_control.json")):
        run_id = control_path.parent.name
        try:
            record = runtime.control_store.read(run_id)
        except RunControlError:
            continue
        if record.state in {
            "reserved",
            "uploading",
            "queued",
            "running",
            "completing",
            "cancelling",
        }:
            run_ids.add(run_id)
    return run_ids


async def _stop_demo_run_queue() -> None:
    global _DEMO_RUN_QUEUE
    current = asyncio.current_task()
    workers = tuple(task for task in _DEMO_WORKERS if task is not current)
    for task in workers:
        task.cancel()
    if workers:
        await asyncio.gather(*workers, return_exceptions=True)
    _DEMO_WORKERS[:] = [
        task
        for task in _DEMO_WORKERS
        if task is current or not task.done()
    ]

    queue = _DEMO_RUN_QUEUE
    if queue is None:
        return
    while True:
        try:
            job = queue.get_nowait()
        except asyncio.QueueEmpty:
            break
        try:
            state = getattr(job, "state", None)
            if state is not None:
                state.queued = False
        finally:
            queue.task_done()
    _DEMO_RUN_QUEUE = None


async def _shutdown_supervised_runs(
    *,
    timeout_s: float | None = None,
    poll_s: float | None = None,
) -> None:
    """Block starts and quiesce every durable Web run before exit."""
    _close_web_run_start_gate()
    runtime = _web_run_runtime()
    if timeout_s is None:
        timeout_s = _WEB_RUN_SHUTDOWN_TIMEOUT_S
    if poll_s is None:
        poll_s = _WEB_RUN_SHUTDOWN_POLL_S
    deadline = time.monotonic() + max(0.0, timeout_s)

    async def within_deadline(awaitable: Any) -> Any:
        remaining = max(0.001, deadline - time.monotonic())
        return await asyncio.wait_for(awaitable, timeout=remaining)

    targets: set[str] = set()
    cancel_dispatched: set[str] = set()
    surviving_pids: dict[str, tuple[int, ...]] = {}
    demo_stopped = False
    stable_scans = 0
    while True:
        inflight_starts = _inflight_web_run_starts()
        inflight_operations = _inflight_web_run_operations()
        discovered = set(runtime.supervisor.active_run_ids())
        discovered.update(_durable_nonterminal_run_ids(runtime))
        discovered.update(inflight_starts)
        targets.update(discovered)

        new_targets = sorted(discovered - cancel_dispatched)
        if new_targets:
            initial_results = await asyncio.gather(
                *(
                    within_deadline(
                        _cancel_controlled_run(run_id, "web_server_shutdown")
                    )
                    for run_id in new_targets
                ),
                return_exceptions=True,
            )
            cancel_dispatched.update(new_targets)
            for run_id, result in zip(new_targets, initial_results):
                if isinstance(result, BaseException):
                    log(
                        "web.run.shutdown_cancel_retrying",
                        run_id=run_id,
                        error=type(result).__name__,
                    )

        if not demo_stopped:
            try:
                await within_deadline(_stop_demo_run_queue())
            except asyncio.CancelledError:
                raise
            except Exception as exc:
                log(
                    "web.run.shutdown_demo_retrying",
                    error=type(exc).__name__,
                )
            demo_stopped = not _DEMO_WORKERS and (
                _DEMO_RUN_QUEUE is None or _DEMO_RUN_QUEUE.empty()
            )

        inflight_starts = _inflight_web_run_starts()
        inflight_operations = _inflight_web_run_operations()
        active = set(runtime.supervisor.active_run_ids())
        durable_nonterminal = _durable_nonterminal_run_ids(runtime)
        targets.update(active)
        targets.update(durable_nonterminal)
        targets.update(inflight_starts)
        unresolved: set[str] = set()
        retryable: set[str] = set()
        for run_id in tuple(targets):
            try:
                record = runtime.control_store.read(run_id)
            except RunControlError:
                if run_id in active or run_id in inflight_starts:
                    unresolved.add(run_id)
                else:
                    targets.discard(run_id)
                continue
            if (
                record.state not in {"completed", "failed", "cancelled"}
                or run_id in active
                or run_id in inflight_starts
            ):
                unresolved.add(run_id)
                retryable.add(run_id)
        if (
            not unresolved
            and not active
            and not durable_nonterminal
            and not inflight_starts
            and not inflight_operations
            and demo_stopped
        ):
            stable_scans += 1
            if stable_scans >= 2:
                return
        else:
            stable_scans = 0

        if time.monotonic() >= deadline:
            pending: dict[str, str] = {}
            for run_id in sorted(unresolved):
                try:
                    record = runtime.control_store.read(run_id)
                    if record.state not in {"completed", "failed", "cancelled", "cancelling"}:
                        record = runtime.control_store.request_cancel(run_id)
                    if record.state == "cancelling" and not record.cancellation_pending:
                        record = runtime.control_store.finalize_cancel(
                            run_id,
                            {
                                "termination_verified": False,
                                "cancellation_pending": "web_server_shutdown_timeout",
                                "reason": "web_server_shutdown",
                            },
                        )
                    pending[run_id] = (
                        record.cancellation_pending or f"state:{record.state}"
                    )
                except RunControlError:
                    pending[run_id] = "start_task_inflight_without_run_control"
            if not demo_stopped:
                pending["demo-background"] = "demo_worker_or_queue_not_quiesced"
            for operation_id in sorted(inflight_operations):
                pending[f"operation:{operation_id}"] = "web_operation_not_quiesced"
            log(
                "web.run.shutdown_cancel_unresolved",
                severity="critical",
                pending=pending,
                surviving_pids={
                    run_id: list(surviving_pids.get(run_id, ()))
                    for run_id in sorted(unresolved)
                },
            )
            raise RuntimeError(
                "shutdown could not verify cancellation for: "
                + ", ".join(sorted(pending))
            )

        if retryable:
            retry_ids = sorted(retryable)
            retry_results = await asyncio.gather(
                *(
                    within_deadline(
                        runtime.services.cancel(run_id, "web_server_shutdown")
                    )
                    for run_id in retry_ids
                ),
                return_exceptions=True,
            )
            for run_id, result in zip(retry_ids, retry_results):
                if isinstance(result, BaseException):
                    continue
                surviving_pids[run_id] = tuple(result.surviving_pids)
                if result.cancel_request_event_required:
                    append_jsonl_event(
                        RUNS_DIR / run_id / "run_events.jsonl",
                        {
                            "run_id": run_id,
                            "event": "run.cancel_requested",
                            "reason": "web_server_shutdown",
                        },
                        event_id=f"{run_id}:cancel-requested",
                    )
            await asyncio.gather(
                *(
                    within_deadline(runtime.supervisor.recover(run_id))
                    for run_id in retry_ids
                ),
                return_exceptions=True,
            )
        await asyncio.sleep(
            min(max(0.001, poll_s), max(0.001, deadline - time.monotonic()))
        )


@app.on_event("shutdown")
async def _shutdown_run_supervision() -> None:
    global _RUN_LIFECYCLE_TASK
    await _shutdown_supervised_runs()
    if _RUN_LIFECYCLE_TASK is not None:
        _RUN_LIFECYCLE_TASK.cancel()
        await asyncio.gather(_RUN_LIFECYCLE_TASK, return_exceptions=True)
        _RUN_LIFECYCLE_TASK = None
    _release_web_server_singleton_lock()


async def _demo_queue_worker(worker_id: int) -> None:
    assert _DEMO_RUN_QUEUE is not None
    while True:
        job = await _DEMO_RUN_QUEUE.get()
        try:
            if job.state.cancelled:
                job.state.queued = False
                log("demo.queue.skip_cancelled", run_id=job.run_id, worker_id=worker_id)
                continue
            log("demo.queue.start", run_id=job.run_id, worker_id=worker_id)
            try:
                await _start_legacy_pipeline_worker(
                    run_id=job.run_id,
                    brief=job.brief,
                    attach_paths=job.attach_paths,
                    reference_poster_path=job.state.reference_poster_path,
                    template=job.template,
                    state=job.state,
                    settings=job.settings,
                    resume_run=None,
                )
            except (RunNotReady, InvalidRunTransition):
                cancellation_observed = job.state.cancelled
                if not cancellation_observed:
                    try:
                        record = _web_run_runtime().control_store.read(job.run_id)
                    except RunControlError:
                        record = None
                    cancellation_observed = (
                        record is not None
                        and record.state in {"cancelling", "cancelled"}
                    )
                if not cancellation_observed:
                    raise
                job.state.cancelled = True
                job.state.queued = False
                log(
                    "demo.queue.cancelled_before_start",
                    run_id=job.run_id,
                    worker_id=worker_id,
                )
                continue
            job.state.task = asyncio.create_task(
                _monitor_supervised_pipeline(run_id=job.run_id, state=job.state),
            )
            job.state.queued = False
            try:
                await job.state.task
            except asyncio.CancelledError:
                log("demo.queue.cancelled", run_id=job.run_id, worker_id=worker_id)
                current = asyncio.current_task()
                if current is not None and current.cancelling():
                    raise
            finally:
                log("demo.queue.finish", run_id=job.run_id, worker_id=worker_id)
        finally:
            if job.state.task is None:
                job.state.queued = False
            _DEMO_RUN_QUEUE.task_done()


async def _latest_conversation_attach_paths(
    conversation_id: str | None,
) -> list[Path]:
    """Return the newest still-readable uploads for this web conversation.

    Upload `File` handles are one-shot on the browser side, but the server
    keeps the staged files under out/uploads. A follow-up like "重新生成一个"
    in the same conversation should still be able to ingest the same paper.
    """
    clean = str(conversation_id or "").strip()
    if not clean:
        return []
    async with _RUNS_LOCK:
        candidates = [
            s for s in _RUNS.values()
            if s.conversation_id == clean and s.attach_paths
        ]
    if candidates:
        latest = max(candidates, key=lambda s: s.created_at)
        live = [p for p in latest.attach_paths if p.exists() and p.is_file()]
        if live:
            return live
    return _latest_persisted_conversation_attach_paths(clean)


def _load_uploads_index() -> dict[str, Any]:
    if not UPLOADS_INDEX_PATH.exists():
        return {"v": 1, "conversations": {}}
    try:
        data = json.loads(UPLOADS_INDEX_PATH.read_text(encoding="utf-8"))
    except Exception:  # noqa: BLE001
        return {"v": 1, "conversations": {}}
    if not isinstance(data, dict):
        return {"v": 1, "conversations": {}}
    conversations = data.get("conversations")
    if not isinstance(conversations, dict):
        data["conversations"] = {}
    return data


def _write_uploads_index(data: dict[str, Any]) -> None:
    tmp = UPLOADS_INDEX_PATH.with_suffix(".tmp")
    tmp.write_text(json.dumps(data, indent=2, ensure_ascii=False), encoding="utf-8")
    os.replace(tmp, UPLOADS_INDEX_PATH)


def _record_conversation_attach_paths(
    conversation_id: str | None,
    paths: list[Path],
) -> None:
    clean = str(conversation_id or "").strip()
    live = [str(p) for p in paths if p.exists() and p.is_file()]
    if not clean or not live:
        return
    with _UPLOADS_INDEX_LOCK:
        data = _load_uploads_index()
        conversations = data.setdefault("conversations", {})
        conversations[clean] = {
            "updated_at": time.time(),
            "paths": live,
        }
        _write_uploads_index(data)


def _latest_persisted_conversation_attach_paths(conversation_id: str) -> list[Path]:
    clean = str(conversation_id or "").strip()
    if not clean:
        return []
    with _UPLOADS_INDEX_LOCK:
        data = _load_uploads_index()
    entry = data.get("conversations", {}).get(clean)
    if not isinstance(entry, dict):
        return []
    raw_paths = entry.get("paths")
    if not isinstance(raw_paths, list):
        return []
    out: list[Path] = []
    for raw in raw_paths:
        if not isinstance(raw, str):
            continue
        p = Path(raw)
        if p.exists() and p.is_file():
            out.append(p)
    return out


def _reference_upload_scope(owner_id: str) -> str:
    """Keep reference uploads per user when run access control is enabled."""
    return owner_id if _RUN_ACCESS_CONTROL else "local"


def _record_conversation_reference_poster(
    conversation_id: str | None,
    reference_path: Path | None,
    *,
    owner_id: str = "",
    reference_handle: str | None = None,
) -> str | None:
    clean = str(conversation_id or "").strip()
    if not clean or reference_path is None:
        return None
    try:
        resolved = reference_path.resolve()
        resolved.relative_to(UPLOADS_DIR.resolve())
    except (OSError, ValueError):
        return None
    if not resolved.is_file() or resolved.parent.name != "reference_poster":
        return None
    handle = str(reference_handle or "").strip() or f"ref_{uuid.uuid4().hex}"
    with _UPLOADS_INDEX_LOCK:
        data = _load_uploads_index()
        references = data.setdefault("reference_posters", {})
        if not isinstance(references, dict):
            references = {}
            data["reference_posters"] = references
        per_conversation = references.setdefault(clean, {})
        if not isinstance(per_conversation, dict):
            per_conversation = {}
            references[clean] = per_conversation
        scope = _reference_upload_scope(owner_id)
        entry = per_conversation.setdefault(scope, {})
        if not isinstance(entry, dict):
            entry = {}
            per_conversation[scope] = entry
        updated_at = time.time()
        entry.update({"updated_at": updated_at, "path": str(resolved)})
        by_handle = entry.setdefault("by_handle", {})
        if not isinstance(by_handle, dict):
            by_handle = {}
            entry["by_handle"] = by_handle
        by_handle[handle] = {
            "updated_at": updated_at,
            "path": str(resolved),
        }
        _write_uploads_index(data)
    return handle


def _persisted_conversation_reference_poster_by_handle(
    conversation_id: str | None,
    reference_handle: str,
    *,
    owner_id: str = "",
) -> Path | None:
    clean = str(conversation_id or "").strip()
    handle = str(reference_handle or "").strip()
    if not clean or not handle:
        return None
    with _UPLOADS_INDEX_LOCK:
        data = _load_uploads_index()
    references = data.get("reference_posters")
    if not isinstance(references, dict):
        return None
    per_conversation = references.get(clean)
    if not isinstance(per_conversation, dict):
        return None
    scope_entry = per_conversation.get(_reference_upload_scope(owner_id))
    if not isinstance(scope_entry, dict):
        return None
    by_handle = scope_entry.get("by_handle")
    if not isinstance(by_handle, dict):
        return None
    entry = by_handle.get(handle)
    if not isinstance(entry, dict):
        return None
    raw_path = entry.get("path")
    if not isinstance(raw_path, str) or not raw_path.strip():
        return None
    try:
        reference_path = Path(raw_path).resolve()
        reference_path.relative_to(UPLOADS_DIR.resolve())
    except (OSError, ValueError):
        return None
    if reference_path.is_file() and reference_path.parent.name == "reference_poster":
        return reference_path
    return None


def _promote_completed_run_reference_poster(
    run_id: str,
    state: _RunState,
) -> None:
    source = state.reference_poster_path
    handle = state.reference_poster_handle
    if source is None or not handle:
        return
    run_uploads = (RUNS_DIR / run_id / "uploads").resolve()
    if source.is_symlink():
        return
    try:
        resolved = source.resolve(strict=True)
        resolved.relative_to(run_uploads)
    except (OSError, ValueError):
        return
    if not resolved.is_file():
        return
    destination_dir = UPLOADS_DIR / run_id / "reference_poster"
    destination_dir.mkdir(parents=True, exist_ok=True)
    destination = destination_dir / resolved.name
    temporary = destination.with_name(
        f".{destination.stem}.{uuid.uuid4().hex}.partial{destination.suffix}"
    )
    try:
        with resolved.open("rb") as reader, temporary.open("xb") as writer:
            shutil.copyfileobj(reader, writer)
            writer.flush()
            os.fsync(writer.fileno())
        _validate_web_reference_poster_file(temporary)
        os.replace(temporary, destination)
        _record_conversation_reference_poster(
            state.conversation_id,
            destination,
            owner_id=state.demo_user_id,
            reference_handle=handle,
        )
    finally:
        temporary.unlink(missing_ok=True)


def _reference_poster_handle(refs_json: str | None) -> str | None:
    for item in _safe_load_json_list(refs_json):
        if not isinstance(item, dict):
            continue
        handle = str(item.get("reference_handle") or "").strip()
        if handle:
            return handle
    return None


def _reference_poster_ref_requires_exact_lookup(refs_json: str | None) -> bool:
    return any(
        isinstance(item, dict) and bool(str(item.get("id") or "").strip())
        for item in _safe_load_json_list(refs_json)
    )


def _latest_persisted_conversation_reference_poster(
    conversation_id: str | None,
    *,
    owner_id: str = "",
) -> Path | None:
    clean = str(conversation_id or "").strip()
    if not clean:
        return None
    with _UPLOADS_INDEX_LOCK:
        data = _load_uploads_index()
    references = data.get("reference_posters")
    if not isinstance(references, dict):
        return None
    per_conversation = references.get(clean)
    if not isinstance(per_conversation, dict):
        return None
    entry = per_conversation.get(_reference_upload_scope(owner_id))
    if not isinstance(entry, dict):
        return None
    raw_path = entry.get("path")
    if not isinstance(raw_path, str) or not raw_path.strip():
        return None
    try:
        reference_path = Path(raw_path).resolve()
        reference_path.relative_to(UPLOADS_DIR.resolve())
    except (OSError, ValueError):
        return None
    if reference_path.is_file() and reference_path.parent.name == "reference_poster":
        return reference_path
    return None


async def _latest_conversation_reference_poster(
    conversation_id: str | None,
    *,
    owner_id: str = "",
) -> Path | None:
    clean = str(conversation_id or "").strip()
    if not clean:
        return None
    async with _RUNS_LOCK:
        candidates = [
            state
            for state in _RUNS.values()
            if state.conversation_id == clean
            and state.reference_poster_path is not None
            and (not _RUN_ACCESS_CONTROL or state.demo_user_id == owner_id)
        ]
    if candidates:
        latest = max(candidates, key=lambda state: state.created_at)
        reference_path = latest.reference_poster_path
        if (
            reference_path is not None
            and reference_path.exists()
            and reference_path.is_file()
            and reference_path.parent.name == "reference_poster"
        ):
            return reference_path
    return _latest_persisted_conversation_reference_poster(clean, owner_id=owner_id)


def _find_uploads_by_attachment_refs(refs_json: str | None) -> list[Path]:
    """Best-effort recovery for old frontend conversations.

    Before the upload index existed, a backend restart lost the in-memory
    conversation -> upload mapping. The browser still persists the prior
    user message's attachment name and byte size, so Resume can pass those
    refs and we can find the newest matching staged upload.
    """
    refs = _safe_load_json_list(refs_json)
    if not refs:
        return []
    matches: list[Path] = []
    for ref in refs:
        if not isinstance(ref, dict):
            continue
        name = Path(str(ref.get("name") or "")).name
        if not name:
            continue
        try:
            size = int(ref.get("size"))
        except (TypeError, ValueError):
            size = -1
        candidates: list[Path] = []
        for p in UPLOADS_DIR.rglob(name):
            if not p.is_file():
                continue
            if size >= 0 and p.stat().st_size != size:
                continue
            candidates.append(p)
        if candidates:
            matches.append(max(candidates, key=lambda p: p.stat().st_mtime))
    return matches


# Maps the per-agent X-Model-* header → models[role] field in the
# health response. Lets the footer display reflect a user's overrides
# *before* a real run hits load_settings().
_HEALTH_MODEL_HEADERS: dict[str, str] = {
    "x-model-designer": "designer",
    "x-model-planner": "designer",
    "x-model-enhancer": "enhancer",
    "x-model-claim-graph": "claim_graph",
    "x-model-deck-outline": "deck_outline",
    "x-model-paper-memory": "paper_memory",
    "x-designer-author-model": "designer_author",
    "x-planner-author-model": "designer_author",
    "x-code-editor-model": "code_editor",
    "x-model-critic": "critic",
    "x-model-composer": "composer",
    "x-model-ingest": "ingest",
    "x-model-image": "image",
    "x-model-image-fallback": "image_fallback",
}

_DEFAULT_MODELS: dict[str, str] = {
    "designer": "gpt-5.5",
    "planner": "gpt-5.5",
    "enhancer": "gpt-5.4-nano",
    "claim_graph": "gpt-5.4-nano",
    "deck_outline": "gpt-5.4-nano",
    "paper_memory": "gpt-5.4-nano",
    "designer_author": "codex",
    "code_editor": "codex",
    "critic": "gpt-5.4-nano",
    "composer": "gpt-5.4-nano",
    "ingest": "gpt-5.4-nano",
    "image": "google/gemini-2.5-flash-image",
    "image_fallback": "openai/gpt-5-image-mini",
}

_PAPER_POSTER_WEB_TEMPLATE = "cvpr-landscape"
_PAPER_POSTER_WEB_DESIGNER_AUTHOR_ARGS = (
    "--search",
    "exec",
    "--dangerously-bypass-approvals-and-sandbox",
    "-",
)
_PAPER_POSTER_WEB_DESIGNER_AUTHOR_MISSING_MESSAGE = (
    "Paper poster generation needs a local coding-agent CLI reachable from the "
    "AutoDesign backend. Install the CLI in the same shell/runtime, put "
    "`codex`/`claude`/`dsh`/`opencode` on that PATH, or set "
    "AUTODESIGN_DESIGNER_AUTHOR_CMD in .env and restart the backend."
)
_CODE_EDITOR_MISSING_MESSAGE = (
    "Poster revision needs a local coding-agent CLI reachable from the AutoDesign "
    "backend. Install the CLI in the same shell/runtime, put "
    "`codex`/`claude`/`dsh`/`opencode` on that PATH, or set AUTODESIGN_CODE_EDITOR_CMD "
    "in .env and restart the backend."
)
_OPENRESEARCH_WEB_SUBMITTER_ARGS = (
    "exec",
    "--skip-git-repo-check",
    "--dangerously-bypass-approvals-and-sandbox",
    "--cd",
    ".",
    "--model",
    "gpt-5.5",
    "-",
)
_OPENRESEARCH_SUBMITTER_MISSING_MESSAGE = (
    "OpenResearch submission needs a local Codex submitter command. Install "
    "the ChatGPT/Codex app, put `codex` on PATH, or configure the external submitter "
    "command in Settings."
)


def _codex_stdin_args_with_model(model: str | None) -> list[str]:
    args = list(_PAPER_POSTER_WEB_DESIGNER_AUTHOR_ARGS)
    model = str(model or "").strip()
    if model:
        insert_at = max(0, len(args) - 1)
        args[insert_at:insert_at] = ["--model", model]
    return args


def _compatible_codex_runtime_for_args(args: list[str] | tuple[str, ...]) -> dict[str, Any]:
    required = ["--ephemeral"]
    for flag in (
        "--search",
        "--skip-git-repo-check",
        "--sandbox",
        "--model",
        "--dangerously-bypass-approvals-and-sandbox",
    ):
        if flag in args and flag not in required:
            required.append(flag)
    return resolve_codex_runtime(required=tuple(required))


def _codex_runtime_command(
    args: list[str] | tuple[str, ...],
) -> dict[str, Any]:
    runtime = _compatible_codex_runtime_for_args(args)
    if not runtime["available"]:
        return {
            "available": False,
            "source": "incompatible",
            "cmd": "",
            "message": (
                "Installed Codex CLI does not support the required capabilities: "
                f"{', '.join(runtime['missing'])}."
            ),
            "binary": runtime["binary"],
            "binary_version": runtime["version"],
            "capabilities": runtime["capabilities"],
            "rejected_candidates": runtime["rejected_candidates"],
        }
    return {
        "available": True,
        "source": runtime["source"],
        "cmd": shlex.join([str(runtime["binary"]), *args]),
        "message": "",
        "binary": runtime["binary"],
        "binary_version": runtime["version"],
        "capabilities": runtime["capabilities"],
        "rejected_candidates": runtime["rejected_candidates"],
    }


def _configured_codex_runtime_command(command: str) -> dict[str, Any] | None:
    try:
        parts = shlex.split(command)
    except ValueError:
        return None
    if not parts:
        return None
    executable_name = Path(parts[0]).name.lower()
    if executable_name != "codex":
        return None
    return _codex_runtime_command(parts[1:])


def _resolve_bare_harness_command(command: str, harness: str) -> str | None:
    """Replace an unlaunchable harness executable with the resolved binary.

    Desktop app bundles often expose a CLI to the user's interactive shell but
    not to the LaunchAgent PATH used by the Web backend. Saved commands may
    also retain the legacy Codex.app path after the CLI moves to ChatGPT.app.
    """

    try:
        parts = shlex.split(command)
    except ValueError:
        return command
    if not parts:
        return command
    executable = parts[0]
    known_executables = {harness}
    if harness == "codex":
        known_executables.update({
            "/Applications/ChatGPT.app/Contents/Resources/codex",
            "/Applications/Codex.app/Contents/Resources/codex",
        })
    if executable not in known_executables:
        return command
    if executable != harness and Path(executable).exists() and os.access(executable, os.X_OK):
        return command
    binary = resolve_harness_binary(harness)
    if not binary:
        return None
    parts[0] = binary
    return shlex.join(parts)


def _harness_binary_source(harness: str, binary: str) -> str:
    if harness == "codex" and binary.startswith("/Applications/"):
        return "app_bundle"
    return "path"


_HARNESS_BINARY_NAMES: dict[str, str] = {
    "codex": "codex",
    "claude": "claude",
    "deepseek": "dsh",
    "opencode": "opencode",
    "kimi": "kimi",
    "mimo": "mimo",
    "zcode": "zcode",
}


def _known_harness_binary_status(
    harness: str,
    *,
    env_keys: tuple[str, ...] = (),
    fallback_path: Path | None = None,
    prefer_fallback: bool = False,
) -> dict[str, Any]:
    explicit = _first_env_value(*env_keys)
    if explicit:
        return {"available": True, "source": "configured", "binary": explicit}

    if (
        prefer_fallback
        and fallback_path is not None
        and fallback_path.exists()
        and os.access(fallback_path, os.X_OK)
    ):
        return {"available": True, "source": "codex_app", "binary": str(fallback_path)}

    binary_name = _HARNESS_BINARY_NAMES.get(harness, harness)
    resolved = resolve_harness_binary(harness)
    if resolved:
        source = "app_bundle" if harness == "codex" and resolved.startswith("/Applications/") else "path"
        return {"available": True, "source": source, "binary": resolved}

    if (
        fallback_path is not None
        and fallback_path.exists()
        and os.access(fallback_path, os.X_OK)
    ):
        return {"available": True, "source": "fallback_path", "binary": str(fallback_path)}

    return {"available": False, "source": "missing", "binary": binary_name}


def _deepseek_binary_status(*, env_keys: tuple[str, ...]) -> dict[str, Any]:
    runtime = resolve_deepseek_harness_runtime(configured_env_keys=env_keys)
    if runtime["available"]:
        message = ""
    elif runtime["source"] == "incompatible":
        version = f" {runtime['version']}" if runtime.get("version") else ""
        message = (
            f"Installed DeepSeek Harness{version} does not expose the released "
            "headless profile. Upgrade with: npm install -g @deepseek-ai/dsh@latest"
        )
    else:
        message = (
            "DeepSeek Harness CLI was not found. Install it with: "
            "npm install -g @deepseek-ai/dsh@latest"
        )
    return {
        "available": bool(runtime["available"]),
        "source": runtime["source"],
        "binary": runtime["binary"],
        "binary_version": runtime.get("version", ""),
        "capabilities": runtime.get("capabilities", {}),
        "missing": runtime.get("missing", []),
        "rejected_candidates": runtime.get("rejected_candidates", []),
        "message": message,
    }


def _designer_author_binary_status(harness: str) -> dict[str, Any]:
    if harness == "codex":
        return _known_harness_binary_status(
            harness,
            env_keys=(
                "AUTODESIGN_DESIGNER_AUTHOR_CODEX_BIN",
                "DESIGN_ANYTHING_DESIGNER_AUTHOR_CODEX_BIN",
                "DESIGN_ANYTHING_PLANNER_AUTHOR_CODEX_BIN",
            ),
            prefer_fallback=True,
        )
    if harness == "claude":
        return _known_harness_binary_status(
            harness,
            env_keys=(
                "AUTODESIGN_DESIGNER_AUTHOR_CLAUDE_BIN",
                "DESIGN_ANYTHING_DESIGNER_AUTHOR_CLAUDE_BIN",
                "DESIGN_ANYTHING_PLANNER_AUTHOR_CLAUDE_BIN",
            ),
        )
    if harness == "deepseek":
        return _deepseek_binary_status(env_keys=(
            "AUTODESIGN_DESIGNER_AUTHOR_DEEPSEEK_BIN",
            "DESIGN_ANYTHING_DESIGNER_AUTHOR_DEEPSEEK_BIN",
            "DESIGN_ANYTHING_PLANNER_AUTHOR_DEEPSEEK_BIN",
        ))
    if harness == "opencode":
        return _known_harness_binary_status(
            harness,
            env_keys=(
                "AUTODESIGN_DESIGNER_AUTHOR_OPENCODE_BIN",
                "DESIGN_ANYTHING_DESIGNER_AUTHOR_OPENCODE_BIN",
                "DESIGN_ANYTHING_PLANNER_AUTHOR_OPENCODE_BIN",
            ),
        )
    if harness == "kimi":
        return _known_harness_binary_status(
            harness,
            env_keys=(
                "AUTODESIGN_DESIGNER_AUTHOR_KIMI_BIN",
                "DESIGN_ANYTHING_DESIGNER_AUTHOR_KIMI_BIN",
                "DESIGN_ANYTHING_PLANNER_AUTHOR_KIMI_BIN",
            ),
        )
    if harness == "pi":
        return _known_harness_binary_status(
            harness,
            env_keys=(
                "AUTODESIGN_DESIGNER_AUTHOR_PI_BIN",
                "DESIGN_ANYTHING_DESIGNER_AUTHOR_PI_BIN",
                "DESIGN_ANYTHING_PLANNER_AUTHOR_PI_BIN",
            ),
        )
    if harness == "zcode":
        return _known_harness_binary_status(
            harness,
            env_keys=(
                "AUTODESIGN_DESIGNER_AUTHOR_ZCODE_BIN",
                "DESIGN_ANYTHING_DESIGNER_AUTHOR_ZCODE_BIN",
                "DESIGN_ANYTHING_PLANNER_AUTHOR_ZCODE_BIN",
            ),
        )
    return _known_harness_binary_status(harness)


def _code_editor_binary_status(harness: str) -> dict[str, Any]:
    if harness == "codex":
        return _known_harness_binary_status(
            harness,
            env_keys=("AUTODESIGN_CODE_EDITOR_CODEX_BIN", "DESIGN_ANYTHING_CODE_EDITOR_CODEX_BIN"),
            prefer_fallback=True,
        )
    if harness == "claude":
        return _known_harness_binary_status(
            harness,
            env_keys=("AUTODESIGN_CODE_EDITOR_CLAUDE_BIN", "DESIGN_ANYTHING_CODE_EDITOR_CLAUDE_BIN"),
        )
    if harness == "deepseek":
        return _deepseek_binary_status(env_keys=(
            "AUTODESIGN_CODE_EDITOR_DEEPSEEK_BIN",
            "DESIGN_ANYTHING_CODE_EDITOR_DEEPSEEK_BIN",
        ))
    if harness == "opencode":
        return _known_harness_binary_status(
            harness,
            env_keys=("AUTODESIGN_CODE_EDITOR_OPENCODE_BIN", "DESIGN_ANYTHING_CODE_EDITOR_OPENCODE_BIN"),
        )
    if harness == "kimi":
        return _known_harness_binary_status(
            harness,
            env_keys=("AUTODESIGN_CODE_EDITOR_KIMI_BIN", "DESIGN_ANYTHING_CODE_EDITOR_KIMI_BIN"),
        )
    if harness == "pi":
        return _known_harness_binary_status(
            harness,
            env_keys=("AUTODESIGN_CODE_EDITOR_PI_BIN", "DESIGN_ANYTHING_CODE_EDITOR_PI_BIN"),
        )
    if harness == "zcode":
        return _known_harness_binary_status(
            harness,
            env_keys=("AUTODESIGN_CODE_EDITOR_ZCODE_BIN", "DESIGN_ANYTHING_CODE_EDITOR_ZCODE_BIN"),
        )
    return _known_harness_binary_status(harness)


def _paper_poster_author_cmd_resolution(settings: Settings | None) -> dict[str, Any]:
    """Resolve the Web paper-poster external author command.

    Product default is deliberately host-friendly: use the bundled macOS app
    binary when available, fall back to PATH, then finally honor an explicit
    advanced .env/Settings command before reporting setup-required.
    """
    configured = ""
    harness = "codex"
    if settings is not None:
        harness = str(getattr(settings, "designer_author_harness", "codex") or "codex").strip()
        configured = str(getattr(settings, "designer_author_cmd", "") or "").strip()
    model = getattr(settings, "designer_author_model", None) if settings is not None else None
    configured = configured or _first_env_value(
        "AUTODESIGN_DESIGNER_AUTHOR_CMD",
        "DESIGN_ANYTHING_DESIGNER_AUTHOR_CMD",
        "DESIGN_ANYTHING_PLANNER_AUTHOR_CMD",
    )
    if harness == "custom" and not configured:
        harness = "codex"
    if harness not in {"custom", "codex"}:
        binary = _designer_author_binary_status(harness)
        if not binary["available"]:
            return {
                **binary,
                "cmd": "",
                "message": binary.get("message") or _PAPER_POSTER_WEB_DESIGNER_AUTHOR_MISSING_MESSAGE,
            }
        cmd = designer_author_command_for_harness(
            harness,
            getattr(settings, "designer_author_model", None) if settings is not None else None,
            explicit_cmd=configured,
        )
        if cmd:
            return {
                **binary,
                "available": True,
                "source": "configured" if configured else binary["source"],
                "cmd": cmd,
                "message": "",
            }
        return {
            "available": False,
            "source": "missing",
            "cmd": "",
            "message": _PAPER_POSTER_WEB_DESIGNER_AUTHOR_MISSING_MESSAGE,
        }

    codex_command = _codex_runtime_command(_codex_stdin_args_with_model(model))
    if codex_command["available"]:
        return codex_command

    if configured:
        resolved_command = _resolve_bare_harness_command(configured, "codex")
        if resolved_command is None:
            return {
                "available": False,
                "source": "missing",
                "cmd": "",
                "message": _PAPER_POSTER_WEB_DESIGNER_AUTHOR_MISSING_MESSAGE,
            }
        return {
            "available": True,
            "source": "configured",
            "cmd": resolved_command,
            "message": "",
        }

    return {
        **codex_command,
        "message": codex_command["message"]
        or _PAPER_POSTER_WEB_DESIGNER_AUTHOR_MISSING_MESSAGE,
    }


def _code_editor_cmd_resolution(settings: Settings | None) -> dict[str, Any]:
    """Resolve the Web multi-turn poster revision command."""
    configured = ""
    harness = "codex"
    if settings is not None:
        harness = str(getattr(settings, "code_editor_harness", "codex") or "codex").strip()
        configured = str(getattr(settings, "code_editor_cmd", "") or "").strip()
    model = getattr(settings, "code_editor_model", None) if settings is not None else None
    configured = configured or _first_env_value(
        "AUTODESIGN_CODE_EDITOR_CMD",
        "DESIGN_ANYTHING_CODE_EDITOR_CMD",
    )
    if harness != "codex":
        binary = _code_editor_binary_status(harness)
        if not configured and not binary["available"]:
            return {
                **binary,
                "cmd": "",
                "message": binary.get("message") or _CODE_EDITOR_MISSING_MESSAGE,
            }
        cmd = code_editor_command_for_harness(
            harness,
            getattr(settings, "code_editor_model", None) if settings is not None else None,
            explicit_cmd=configured,
        )
        if cmd:
            return {
                **binary,
                "available": True,
                "source": "configured" if configured else binary["source"],
                "cmd": cmd,
                "message": "",
            }
        return {
            "available": False,
            "source": "missing",
            "cmd": "",
            "message": _CODE_EDITOR_MISSING_MESSAGE,
        }

    if configured:
        configured_runtime = _configured_codex_runtime_command(configured)
        if configured_runtime is not None:
            if configured_runtime["available"]:
                return configured_runtime
            return {
                **configured_runtime,
                "message": configured_runtime["message"] or _CODE_EDITOR_MISSING_MESSAGE,
            }
        resolved_command = _resolve_bare_harness_command(configured, harness)
        if resolved_command is None:
            return {
                "available": False,
                "source": "missing",
                "cmd": "",
                "message": _CODE_EDITOR_MISSING_MESSAGE,
            }
        return {
            "available": True,
            "source": "configured",
            "cmd": resolved_command,
            "message": "",
        }

    codex_command = _codex_runtime_command(_codex_stdin_args_with_model(model))
    if codex_command["available"]:
        return codex_command

    return {
        **codex_command,
        "message": codex_command["message"] or _CODE_EDITOR_MISSING_MESSAGE,
    }


def _web_paper_poster_settings(settings: Settings) -> Settings:
    """Return the Web UI's dogfood paper-poster profile for PDF poster runs."""
    author_cmd = str(_paper_poster_author_cmd_resolution(settings)["cmd"])
    author_harness = str(getattr(settings, "designer_author_harness", "codex") or "codex").strip()
    if author_harness == "custom" and not str(getattr(settings, "designer_author_cmd", "") or "").strip():
        author_harness = "codex"
    return replace(
        settings,
        designer_author_mode="external",
        designer_author_harness=author_harness,  # type: ignore[arg-type]
        designer_author_cmd=author_cmd,
        designer_author_timeout_s=3600,
        designer_author_max_attempts=12,
    )


def _paper_poster_profile(settings: Settings | None) -> dict[str, Any]:
    author_cmd = _paper_poster_author_cmd_resolution(settings)
    return {
        "template": _PAPER_POSTER_WEB_TEMPLATE,
        "designer_author": getattr(settings, "designer_author_mode", "external"),
        "designer_author_harness": getattr(settings, "designer_author_harness", "custom"),
        "designer_author_model": getattr(settings, "designer_author_model", None),
        "designer_author_cmd": author_cmd["cmd"],
        "designer_author_cmd_available": author_cmd["available"],
        "designer_author_cmd_source": author_cmd["source"],
        "designer_author_cmd_message": author_cmd["message"],
        "designer_author_timeout_s": getattr(settings, "designer_author_timeout_s", 3600),
        "designer_author_max_attempts": getattr(settings, "designer_author_max_attempts", 12),
    }


def _code_editor_profile(settings: Settings | None) -> dict[str, Any]:
    cmd = _code_editor_cmd_resolution(settings)
    auth_status = "not_verified" if cmd["available"] else "unavailable"
    return {
        "available": cmd["available"],
        "harness": getattr(settings, "code_editor_harness", "codex"),
        "model": getattr(settings, "code_editor_model", None),
        "cmd": cmd["cmd"],
        "cmd_source": cmd["source"],
        "command_detected": bool(cmd["available"]),
        "auth_status": auth_status,
        "auth_message": (
            "Command detected; auth not verified. Run a coding-agent smoke test."
            if cmd["available"]
            else cmd["message"]
        ),
        "message": cmd["message"],
        "timeout_s": getattr(settings, "code_editor_timeout_s", 600),
        "max_attempts": getattr(settings, "code_editor_max_attempts", 2),
    }


def _openresearch_submitter_cmd_resolution(settings: Settings | None) -> dict[str, Any]:
    configured = ""
    if settings is not None:
        configured = str(getattr(settings, "openresearch_submitter_cmd", "") or "").strip()
    configured = configured or _first_env_value(
        "AUTODESIGN_OPENRESEARCH_SUBMITTER_CMD",
        "DESIGN_ANYTHING_OPENRESEARCH_SUBMITTER_CMD",
        "OPEN_DESIGN_OPENRESEARCH_SUBMITTER_CMD",
    )
    if configured:
        resolved_command = _resolve_bare_harness_command(configured, "codex")
        if resolved_command is None:
            return {
                "available": False,
                "source": "missing",
                "cmd": "",
                "message": _OPENRESEARCH_SUBMITTER_MISSING_MESSAGE,
            }
        return {
            "available": True,
            "source": "configured",
            "cmd": resolved_command,
            "message": "",
        }

    codex_binary = resolve_harness_binary("codex")
    if codex_binary:
        return {
            "available": True,
            "source": _harness_binary_source("codex", codex_binary),
            "cmd": shlex.join([
                codex_binary,
                *_OPENRESEARCH_WEB_SUBMITTER_ARGS,
            ]),
            "message": "",
        }

    return {
        "available": False,
        "source": "missing",
        "cmd": "",
        "message": _OPENRESEARCH_SUBMITTER_MISSING_MESSAGE,
    }


def _coding_agent_smoke_cmd_resolution(settings: Settings | None) -> dict[str, Any]:
    harness = str((getattr(settings, "code_editor_harness", "codex") if settings else "codex") or "codex").strip()
    model = str((getattr(settings, "code_editor_model", None) if settings else None) or "").strip()
    settings_cmd = ""
    if settings is not None:
        settings_cmd = str(getattr(settings, "code_editor_cmd", "") or "").strip()
    env_cmd = _first_env_value("AUTODESIGN_CODE_EDITOR_CMD", "DESIGN_ANYTHING_CODE_EDITOR_CMD")
    configured = env_cmd or (settings_cmd if harness == "custom" else "")

    if harness == "deepseek" and not configured:
        binary = _code_editor_binary_status(harness)
        if not binary["available"]:
            return {
                **binary,
                "cmd": "",
                "message": binary.get("message") or _CODE_EDITOR_MISSING_MESSAGE,
            }
        cmd = coding_agent_smoke_command_for_harness(harness, model)
        return {
            **binary,
            "available": True,
            "cmd": cmd,
            "message": "",
        }

    if harness == "claude" and not configured:
        binary = _code_editor_binary_status(harness)
        if not binary["available"]:
            return {
                "available": False,
                "source": "missing",
                "cmd": "",
                "message": _CODE_EDITOR_MISSING_MESSAGE,
            }
        cmd = [
            str(binary["binary"]),
            "--debug",
            "--debug-file",
            ".coding_agent_smoke.claude-debug.log",
            "--print",
            "--permission-mode",
            "bypassPermissions",
            "--dangerously-skip-permissions",
            "--no-session-persistence",
        ]
        if model:
            cmd.extend(["--model", model])
        cmd.append(
            "Read the full smoke-test prompt from coding_agent_smoke_prompt.md in the current directory, "
            "write coding_agent_smoke_output.json exactly as requested, and exit."
        )
        return {
            "available": True,
            "source": binary["source"],
            "cmd": shlex.join(cmd),
            "message": "",
        }

    if harness == "opencode" and not configured:
        binary = _code_editor_binary_status(harness)
        if not binary["available"]:
            return {
                "available": False,
                "source": "missing",
                "cmd": "",
                "message": _CODE_EDITOR_MISSING_MESSAGE,
            }
        cmd = [
            str(binary["binary"]),
            "run",
            "--dir",
            ".",
        ]
        if model:
            cmd.extend(["--model", model])
        if (_first_env_value("AUTODESIGN_CODE_EDITOR_OPENCODE_SKIP_PERMISSIONS", "DESIGN_ANYTHING_CODE_EDITOR_OPENCODE_SKIP_PERMISSIONS") or "1").lower() not in {"0", "false", "no", "off"}:
            cmd.append("--dangerously-skip-permissions")
        cmd.append(
            "Read coding_agent_smoke_prompt.md in the current directory and follow it exactly. "
            "Work only in this directory, write coding_agent_smoke_output.json, and exit."
        )
        return {
            "available": True,
            "source": binary["source"],
            "cmd": shlex.join(cmd),
            "message": "",
        }

    resolved = _code_editor_cmd_resolution(settings)
    if harness != "codex" or not resolved.get("available"):
        return resolved
    codex_runtime = resolve_codex_runtime(required=("--ephemeral",))
    if not codex_runtime["available"]:
        return {
            **resolved,
            "available": False,
            "source": "incompatible",
            "cmd": "",
            "message": (
                "Installed Codex CLI does not support the required smoke-test "
                f"capabilities: {', '.join(codex_runtime['missing'])}."
            ),
            "binary": codex_runtime["binary"],
            "binary_version": codex_runtime["version"],
            "capabilities": codex_runtime["capabilities"],
            "rejected_candidates": codex_runtime["rejected_candidates"],
        }
    try:
        parts = shlex.split(str(resolved.get("cmd") or ""))
    except ValueError:
        return resolved
    if "exec" not in parts:
        return resolved
    parts[0] = str(codex_runtime["binary"])
    parts = [part for part in parts if part != "--search"]
    exec_index = parts.index("exec")
    if "--ephemeral" not in parts:
        parts.insert(exec_index + 1, "--ephemeral")
    return {
        **resolved,
        "source": codex_runtime["source"],
        "cmd": shlex.join(parts),
        "binary": codex_runtime["binary"],
        "binary_version": codex_runtime["version"],
        "capabilities": codex_runtime["capabilities"],
        "rejected_candidates": codex_runtime["rejected_candidates"],
    }


def _run_coding_agent_smoke(settings: Settings | None, *, timeout_s: int = 60) -> dict[str, Any]:
    resolved = _coding_agent_smoke_cmd_resolution(settings)
    harness = str((getattr(settings, "code_editor_harness", "codex") if settings else "codex") or "codex").strip()
    model = str((getattr(settings, "code_editor_model", None) if settings else None) or "").strip()
    base = {
        "harness": harness,
        "model": model or None,
        "command": resolved.get("cmd", ""),
        "command_source": resolved.get("source", "missing"),
        "command_detected": bool(resolved.get("available")),
        "binary": resolved.get("binary", ""),
        "binary_version": resolved.get("binary_version", ""),
        "capabilities": resolved.get("capabilities", {}),
        "rejected_candidates": resolved.get("rejected_candidates", []),
        "timeout_s": timeout_s,
    }
    if not resolved.get("available") or not str(resolved.get("cmd") or "").strip():
        return {
            **base,
            "ok": False,
            "status": "missing_command",
            "auth_status": "unavailable",
            "reason": resolved.get("message") or _CODE_EDITOR_MISSING_MESSAGE,
            "elapsed_s": 0,
            "stdout_excerpt": "",
            "stderr_excerpt": "",
        }

    prompt = (
        "This is an AutoDesign coding-agent smoke test.\n"
        "In the current working directory, write exactly one file named "
        "coding_agent_smoke_output.json with this JSON object:\n"
        "{\"ok\": true, \"message\": \"coding-agent-smoke-passed\"}\n"
        "Do not edit any other files. Exit immediately after writing the file.\n"
    )

    try:
        cmd = shlex.split(str(resolved["cmd"]))
    except ValueError as exc:
        return {
            **base,
            "ok": False,
            "status": "failed",
            "auth_status": "not_verified",
            "reason": f"command_parse_error: {exc}",
            "elapsed_s": 0,
            "stdout_excerpt": "",
            "stderr_excerpt": "",
        }
    if not cmd:
        return {
            **base,
            "ok": False,
            "status": "failed",
            "auth_status": "not_verified",
            "reason": "empty_command",
            "elapsed_s": 0,
            "stdout_excerpt": "",
            "stderr_excerpt": "",
        }

    with tempfile.TemporaryDirectory(prefix="autodesign-coding-agent-smoke-") as tmp:
        smoke_dir = Path(tmp)
        marker = smoke_dir / "coding_agent_smoke_output.json"
        stdout_path = smoke_dir / ".coding_agent_smoke.stdout.tmp"
        stderr_path = smoke_dir / ".coding_agent_smoke.stderr.tmp"
        debug_path = smoke_dir / ".coding_agent_smoke.claude-debug.log"
        (smoke_dir / "coding_agent_smoke_prompt.md").write_text(prompt, encoding="utf-8")
        start = time.monotonic()
        returncode: int | None = None
        timed_out = False
        reason = "process_exit"
        try:
            env = harness_subprocess_env(
                os.environ,
                harness=harness,
                api_key=getattr(settings, "harness_api_key", None) if settings else None,
            )
            author_python = _first_env_value(
                "AUTODESIGN_AUTHOR_PYTHON",
                "DESIGN_ANYTHING_AUTHOR_PYTHON",
            ) or sys.executable
            env["AUTODESIGN_AUTHOR_PYTHON"] = author_python
            env.setdefault("DESIGN_ANYTHING_AUTHOR_PYTHON", author_python)
            with stdout_path.open("w", encoding="utf-8") as stdout_f, stderr_path.open("w", encoding="utf-8") as stderr_f:
                proc = subprocess.Popen(
                    cmd,
                    stdin=subprocess.PIPE,
                    stdout=stdout_f,
                    stderr=stderr_f,
                    text=True,
                    cwd=str(smoke_dir),
                    env=env,
                    start_new_session=(os.name == "posix"),
                )
                try:
                    if proc.stdin is not None:
                        proc.stdin.write(prompt)
                        proc.stdin.close()
                except BrokenPipeError:
                    pass
                deadline = start + timeout_s
                while True:
                    returncode = proc.poll()
                    if marker.exists():
                        returncode = _terminate_subprocess_group(proc)
                        reason = "marker_written"
                        break
                    if returncode is not None:
                        break
                    if time.monotonic() >= deadline:
                        timed_out = True
                        returncode = _terminate_subprocess_group(proc)
                        reason = "timeout"
                        break
                    time.sleep(0.1)
        except OSError as exc:
            return {
                **base,
                "ok": False,
                "status": "failed",
                "auth_status": "not_verified",
                "reason": f"command_start_error: {exc}",
                "returncode": returncode,
                "timed_out": timed_out,
                "elapsed_s": round(time.monotonic() - start, 3),
                "stdout_excerpt": _coding_agent_smoke_log_excerpt(stdout_path, limit=1400),
                "stderr_excerpt": _coding_agent_smoke_stderr_excerpt(stderr_path, debug_path, limit=900),
            }

        marker_ok = False
        if marker.exists():
            try:
                payload = json.loads(marker.read_text(encoding="utf-8"))
                marker_ok = bool(isinstance(payload, dict) and payload.get("ok") is True)
                if not marker_ok:
                    reason = "marker_json_not_ok"
            except Exception as exc:  # noqa: BLE001
                reason = f"invalid_marker_json: {exc}"
        elif timed_out:
            reason = "timeout"
        else:
            reason = "missing_smoke_output"

        ok = bool(marker_ok and not timed_out)
        return {
            **base,
            "ok": ok,
            "status": "passed" if ok else ("timeout" if timed_out else "failed"),
            "auth_status": "verified" if ok else "not_verified",
            "reason": reason,
            "returncode": returncode,
            "timed_out": timed_out,
            "elapsed_s": round(time.monotonic() - start, 3),
            "stdout_excerpt": _coding_agent_smoke_log_excerpt(stdout_path, limit=1400),
            "stderr_excerpt": _coding_agent_smoke_stderr_excerpt(stderr_path, debug_path, limit=900),
        }


def _web_openresearch_settings(settings: Settings) -> Settings:
    if str(getattr(settings, "openresearch_submitter_mode", "off") or "off") != "custom":
        return settings
    if str(getattr(settings, "openresearch_submitter_cmd", "") or "").strip():
        return settings
    resolved = _openresearch_submitter_cmd_resolution(settings)
    if not resolved["available"]:
        return settings
    return replace(settings, openresearch_submitter_cmd=str(resolved["cmd"]))


def _demo_designer_author_cmd() -> str:
    codex_binary = resolve_harness_binary("codex")
    if codex_binary:
        return shlex.join([
            codex_binary,
            *_PAPER_POSTER_WEB_DESIGNER_AUTHOR_ARGS,
        ])
    return ""


def _demo_settings(settings: Settings) -> Settings:
    return replace(
        settings,
        designer_author_mode="external",
        designer_author_harness="custom",
        designer_author_cmd=_demo_designer_author_cmd(),
        designer_author_timeout_s=3600,
        designer_author_max_attempts=12,
        openresearch_submitter_mode="off",
        openresearch_submitter_cmd="",
        poster_harness_mode="dogfood",
    )


def _openresearch_profile(settings: Settings | None) -> dict[str, Any]:
    if _DEMO_MODE:
        return {
            "submitter": "off",
            "submitter_cmd": "",
            "submitter_cmd_available": False,
            "submitter_cmd_source": "disabled",
            "submitter_cmd_message": "OpenResearch submission is disabled in demo mode.",
            "submitter_timeout_s": 0,
            "org_id_configured": False,
            "repo_configured": False,
            "api_url": "",
            "api_token_configured": False,
        }
    resolved = _openresearch_submitter_cmd_resolution(settings)
    return {
        "submitter": getattr(settings, "openresearch_submitter_mode", "off"),
        "submitter_cmd": resolved["cmd"],
        "submitter_cmd_available": resolved["available"],
        "submitter_cmd_source": resolved["source"],
        "submitter_cmd_message": resolved["message"],
        "submitter_timeout_s": getattr(settings, "openresearch_submitter_timeout_s", 300),
        "org_id_configured": bool(str(getattr(settings, "openresearch_org_id", "") or "").strip()),
        "repo_configured": bool(str(getattr(settings, "openresearch_default_repo_full_name", "") or "").strip()),
        "api_url": getattr(settings, "openresearch_api_url", "https://api.openresearch.sh"),
        "api_token_configured": bool(str(getattr(settings, "openresearch_token", "") or "").strip()),
    }


@app.get("/api/palettes")
def palettes(artifact_type: str = Query(...)) -> dict[str, Any]:
    if artifact_type != "poster":
        raise HTTPException(
            400,
            detail={
                "code": "unsupported_palette_artifact_type",
                "message": "Palette catalog is available for Poster only.",
            },
        )
    try:
        return academic_palette_catalog_payload()
    except AcademicPaletteCatalogError as exc:
        raise HTTPException(
            503,
            detail={
                "code": "palette_catalog_unavailable",
                "message": str(exc),
            },
        ) from exc


@app.get("/api/canvas-presets")
def canvas_presets(artifact_type: str = Query(...)) -> dict[str, Any]:
    if artifact_type != "poster":
        raise HTTPException(
            400,
            detail={
                "code": "unsupported_canvas_preset_artifact_type",
                "message": "Canvas presets are available for Poster only.",
            },
        )
    return poster_canvas_preset_catalog()


def _validated_web_canvas_selection(
    artifact_type: str,
    template: str | None,
    canvas_preset_id: str | None,
) -> tuple[str | None, str | None]:
    requested_template = (
        template.strip().lower().replace("_", "-")
        if isinstance(template, str) and template.strip()
        else None
    )
    requested_selection = (
        canvas_preset_id.strip()
        if isinstance(canvas_preset_id, str) and canvas_preset_id.strip()
        else None
    )
    if artifact_type != "poster":
        if requested_template is not None or requested_selection is not None:
            raise HTTPException(
                422,
                detail={
                    "code": "canvas_preset_not_supported_for_artifact",
                    "message": "Canvas presets are supported for Poster only.",
                },
            )
        return None, None
    if requested_template is not None and resolve_template(requested_template) is None:
        raise HTTPException(
            422,
            detail={
                "code": "unknown_canvas_preset",
                "message": f"Unknown Poster canvas preset: {requested_template}",
            },
        )
    if requested_selection is None:
        return requested_template, None

    catalog_ids = {
        str(item.get("id") or "")
        for item in poster_canvas_preset_catalog()["presets"]
        if isinstance(item, dict)
    }
    if requested_selection not in catalog_ids:
        raise HTTPException(
            422,
            detail={
                "code": "unknown_canvas_preset",
                "message": f"Unknown Poster canvas preset: {requested_selection}",
            },
        )
    expected_template = None if requested_selection == "auto" else requested_selection
    if requested_template not in {None, expected_template} or (
        requested_selection == "auto" and requested_template is not None
    ):
        raise HTTPException(
            422,
            detail={
                "code": "canvas_preset_mismatch",
                "message": "Canvas selection does not match the requested template.",
            },
        )
    return expected_template, requested_selection


def _validated_canvas_prompt(brief: str, artifact_type: str) -> None:
    try:
        intent = parse_canvas_intent(brief)
        if (
            artifact_type != "poster"
            and intent is not None
            and intent.get("template_id")
        ):
            raise CanvasIntentError(
                "conflicting_canvas_directives",
                "A Poster template cannot be combined with a different artifact type.",
            )
    except CanvasIntentError as exc:
        raise HTTPException(
            422,
            detail={"code": exc.code, "message": str(exc)},
        ) from exc


def _validated_web_palette_id(
    artifact_type: str,
    palette_id: str | None,
) -> str | None:
    if artifact_type != "poster":
        if palette_id is not None:
            raise HTTPException(
                422,
                detail={
                    "code": "palette_not_supported_for_artifact",
                    "message": "Palettes are supported for Poster only.",
                },
            )
        return None

    requested_id = str(palette_id or "").strip()
    if not requested_id:
        raise HTTPException(
            422,
            detail={
                "code": "poster_palette_required",
                "message": "Select a Poster palette before generating or revising.",
            },
        )
    try:
        catalog = academic_palette_catalog_payload()
        canonical_ids = {
            str(item.get("id") or "").strip()
            for item in catalog["palettes"]
            if isinstance(item, dict)
        }
        if requested_id not in canonical_ids:
            raise ValueError(f"unknown academic palette: {requested_id}")
        color_system = require_academic_color_system(requested_id)
    except AcademicPaletteCatalogError as exc:
        raise HTTPException(
            503,
            detail={
                "code": "palette_catalog_unavailable",
                "message": str(exc),
            },
        ) from exc
    except ValueError as exc:
        raise HTTPException(
            422,
            detail={
                "code": "unknown_poster_palette",
                "message": f"Unknown Poster palette: {requested_id}",
            },
        ) from exc
    normalized_id = str(color_system.get("palette_id") or "").strip()
    if not normalized_id:
        raise HTTPException(
            503,
            detail={
                "code": "palette_catalog_unavailable",
                "message": "The Poster palette catalog returned an invalid palette.",
            },
        )
    return normalized_id


@app.get("/api/health")
def health(request: Request) -> dict[str, Any]:
    """Reports every agent role's resolved model id, so the frontend
    footer can show all roles, not just designer. `designer_model` and
    `image_model` are kept at the top level for v0.2 wire-shape backward
    compat — new callers should read `models.<role>` instead.

    Model resolution: starts from SETTINGS (or hardcoded defaults if the
    backend booted credentials-less), then overlays any `X-Model-*`
    headers the client sent. This makes the footer reflect Settings-
    drawer picks the moment the user clicks Save, with no need for the
    user to also paste a key just to preview the binding.
    """
    effective_settings = _demo_settings(SETTINGS) if (_DEMO_MODE and SETTINGS is not None) else SETTINGS
    if not _DEMO_MODE:
        overrides, _has_key = _request_env_overrides(request)
        if overrides:
            try:
                effective_settings = _settings_for_request(request)
            except HTTPException:
                effective_settings = SETTINGS
    if effective_settings is None:
        try:
            effective_settings = _settings_for_request(request)
        except HTTPException:
            effective_settings = None

    if effective_settings is None:
        models = dict(_DEFAULT_MODELS)
        needs_setup = True
    else:
        models = {
            "designer": effective_settings.designer_model,
            "planner": effective_settings.designer_model,
            "enhancer": effective_settings.enhancer_model,
            "claim_graph": effective_settings.claim_graph_model,
            "deck_outline": effective_settings.deck_outline_model,
            "paper_memory": effective_settings.paper_memory_model,
            "designer_author": effective_settings.designer_author_model or effective_settings.designer_author_harness,
            "code_editor": effective_settings.code_editor_model or effective_settings.code_editor_harness,
            "critic": effective_settings.critic_model,
            "composer": effective_settings.composer_model,
            "ingest": effective_settings.ingest_model,
            "image": effective_settings.image_model,
            "image_fallback": effective_settings.image_fallback_model,
        }
        needs_setup = False

    if not _DEMO_MODE:
        # Header overlay — empty headers are ignored.
        for header_name, role in _HEALTH_MODEL_HEADERS.items():
            v = (request.headers.get(header_name, "") or "").strip()
            if v:
                models[role] = normalize_model_id(v)

    paper_poster_settings = (
        _web_paper_poster_settings(effective_settings)
        if effective_settings is not None
        else None
    )
    try:
        openresearch_settings = _settings_for_openresearch_request(request)
    except HTTPException:
        openresearch_settings = (
            _web_openresearch_settings(effective_settings)
            if effective_settings is not None
            else None
        )

    return {
        "status": "needs_setup" if needs_setup else "ok",
        "mode": "demo" if _DEMO_MODE else "real",
        "demo_mode": _DEMO_MODE,
        "public_user_isolation": _PUBLIC_USER_ISOLATION,
        "user_isolation": _RUN_ACCESS_CONTROL,
        "needs_setup": needs_setup,
        "designer_model": models["designer"],
        "planner_model": models["designer"],
        "image_model": models["image"],
        "models": models,
        "demo": {
            "artifact_type": "poster",
            "template": _DEMO_FIXED_TEMPLATE,
            "daily_limit": _DEMO_DAILY_LIMIT,
            "concurrency": _DEMO_CONCURRENCY,
            "queue_max": _DEMO_QUEUE_MAX,
            "run_ttl_hours": _DEMO_RUN_TTL_HOURS,
            "max_pdf_bytes": _DEMO_MAX_PDF_BYTES,
            "settings_locked": _DEMO_MODE,
            "openresearch_enabled": False if _DEMO_MODE else True,
            "requires_low_privilege_user": _DEMO_MODE,
        },
        "backend_profile": {
            "paper_poster": _paper_poster_profile(paper_poster_settings),
            "code_editor": _code_editor_profile(effective_settings),
            "harness_capabilities": build_coding_harness_capabilities(effective_settings),
            "openresearch": _openresearch_profile(openresearch_settings),
            "environment": _environment_profile(),
        },
    }


@app.post("/api/coding-agent/smoke")
def coding_agent_smoke(req: CodingAgentSmokeRequest, request: Request) -> dict[str, Any]:
    if _DEMO_MODE:
        return {
            "ok": False,
            "status": "disabled",
            "auth_status": "unavailable",
            "command_detected": False,
            "command_source": "disabled",
            "command": "",
            "harness": "demo",
            "model": None,
            "timeout_s": req.timeout_s,
            "elapsed_s": 0,
            "reason": "Coding-agent smoke tests are disabled in demo mode.",
            "stdout_excerpt": "",
            "stderr_excerpt": "",
        }
    settings = _settings_for_code_editor_request(request)
    result = _run_coding_agent_smoke(settings, timeout_s=req.timeout_s)
    log(
        "web.coding_agent.smoke",
        ok=result.get("ok"),
        status=result.get("status"),
        harness=result.get("harness"),
        command_source=result.get("command_source"),
        elapsed_s=result.get("elapsed_s"),
    )
    return result


@app.get("/api/history")
def history(
    request: Request,
    limit: int = Query(25, ge=1, le=200),
    include_design_sessions: bool = Query(False),
) -> dict[str, Any]:
    """Server-backed history for local browser clients.

    The frontend still keeps a localStorage cache for snappy rendering, but
    this endpoint is the shared source that lets Chrome, the Codex in-app
    browser, and a freshly-cleared browser rediscover real generated runs.

    The response contains only list-safe conversation summaries. Full message
    and editable artifact payloads are loaded from the detail endpoint after a
    user selects a conversation.

    By default this imports recent valid disk runs and skips the older
    design-session scan that can touch hundreds of stale runs in dogfood
    worktrees. Call with `include_design_sessions=true` when a full
    event-backed reconstruction is explicitly needed.

    Sources that may be merged:
      1. out/web_history.json — optional full conversation snapshots saved by
         clients.
      2. out/design_sessions/*.jsonl + out/runs/* — authoritative real run
         events written by the backend during generation/editing.
    """
    if _RUN_ACCESS_CONTROL:
        imported = _import_history_from_server_events(
            limit=None,
            include_design_sessions=False,
            demo_user_id=_demo_user_id(request),
            compact=True,
            summary_limit=limit,
        )
        conversations = _limit_history_conversations(imported, limit)
        return {
            "conversations": _history_conversation_summaries(conversations),
            "imported_runs": _count_history_artifacts(imported),
            "user_isolated": True,
        }

    stored = _load_web_history_summaries()
    imported = _import_history_from_server_events(
        limit=None,
        include_design_sessions=include_design_sessions,
        demo_user_id=None,
        compact=True,
        summary_limit=None if include_design_sessions else limit,
        summary_min_updated_at=(
            None
            if include_design_sessions
            else _history_summary_cutoff(stored, limit)
        ),
        existing_conversation_ids=set(stored),
    )
    merged = _merge_history_conversations(stored, imported)
    conversations = _limit_history_conversations(merged, limit)
    return {
        "conversations": _history_conversation_summaries(conversations),
        "imported_runs": _count_history_artifacts(imported),
        "user_isolated": False,
    }


@app.get("/api/history/conversations/{conversation_id}")
def history_conversation_detail(
    conversation_id: str,
    request: Request,
    include_design_sessions: bool = Query(False),
) -> dict[str, Any]:
    """Return the full saved/imported conversation after sidebar selection."""
    stored: dict[str, Any] = {}
    if not _RUN_ACCESS_CONTROL:
        data = _load_web_history()
        conversations = data.get("conversations", {})
        raw = conversations.get(conversation_id) if isinstance(conversations, dict) else None
        if isinstance(raw, dict):
            stored[conversation_id] = raw

    imported = _history_conversation_from_server_sources(
        conversation_id,
        include_design_sessions=include_design_sessions and not _RUN_ACCESS_CONTROL,
        demo_user_id=_demo_user_id(request) if _RUN_ACCESS_CONTROL else None,
    )
    merged = _merge_history_conversations(
        stored,
        {conversation_id: imported} if imported is not None else {},
    )
    conversation = merged.get(conversation_id)
    if conversation is None:
        raise HTTPException(404, detail=f"conversation not found: {conversation_id}")
    return {
        "conversation": _history_detail_conversation(conversation),
        "user_isolated": _RUN_ACCESS_CONTROL,
    }


@app.post("/api/history/conversation")
def save_history_conversation(req: HistoryConversationRequest) -> dict[str, Any]:
    """Persist one browser conversation snapshot server-side.

    This is intentionally small and permissive: it stores only JSON fields
    that are already part of the web Conversation shape, and strips pending
    run flags so reloading another browser never revives stale spinners.
    """
    if _RUN_ACCESS_CONTROL:
        return {"ok": True, "conversation_id": str(req.conversation.get("id") or "")}
    conv = _sanitize_history_conversation(req.conversation)
    if conv is None:
        raise HTTPException(400, detail="invalid conversation")
    with _WEB_HISTORY_LOCK:
        data = _load_web_history()
        conversations = data.setdefault("conversations", {})
        conversations[conv["id"]] = conv
        _write_web_history(data)
    return {"ok": True, "conversation_id": conv["id"]}


@app.delete("/api/history/conversations/{conversation_id}")
def delete_history_conversation(conversation_id: str) -> dict[str, Any]:
    if _RUN_ACCESS_CONTROL:
        return {"ok": True, "conversation_id": conversation_id}
    with _WEB_HISTORY_LOCK:
        data = _load_web_history()
        conversations = data.setdefault("conversations", {})
        conversations.pop(conversation_id, None)
        _write_web_history(data)
    return {"ok": True, "conversation_id": conversation_id}


def _safe_matrix_id(matrix_id: str) -> str:
    clean = re.sub(r"[^A-Za-z0-9_.-]", "", matrix_id.strip())
    if not clean or clean in {".", ".."}:
        raise HTTPException(400, detail="invalid matrix id")
    return clean


def _matrix_dir_for_id(matrix_id: str) -> Path:
    clean = _safe_matrix_id(matrix_id)
    return _BOOT_OUT_DIR / "harness_matrix" / clean


async def _lookup_matrix_dir(matrix_id: str) -> Path:
    clean = _safe_matrix_id(matrix_id)
    async with _HARNESS_MATRIX_JOBS_LOCK:
        state = _HARNESS_MATRIX_JOBS.get(clean)
    if state is not None:
        return state.matrix_dir
    return _matrix_dir_for_id(clean)


def _initial_harness_matrix_snapshot(
    matrix_id: str,
    req: HarnessMatrixRequest,
    *,
    matrix_dir: Path,
) -> dict[str, Any]:
    harnesses = req.harnesses or [HarnessMatrixHarnessRequest(id=h) for h in CODING_HARNESSES]
    return {
        "matrix_id": matrix_id,
        "status": "running",
        "paper_path": req.paper_path,
        "template": req.template,
        "attempts": req.attempts,
        "timeout_s": req.timeout_s,
        "concurrency": req.concurrency,
        "matrix_dir": str(matrix_dir),
        "report_path": str(matrix_dir / "report.md"),
        "strict_success": False,
        "hard_failure_count": 0,
        "summary": {
            "total_cells": len(harnesses),
            "terminal_cells": 0,
            "strict_success_count": 0,
            "usable_count": 0,
            "hard_failure_count": 0,
            "outcome_counts": {},
        },
        "rows": [
            {
                "matrix_id": matrix_id,
                "harness": h.id,
                "requested_model": h.model or "",
                "model_selection_mode": (
                    "locked_config"
                    if h.id == "zcode"
                    else "config_overlay"
                    if h.id == "deepseek"
                    else "cli_flag"
                ),
                "status": "pending",
                "outcome_class": "",
                "terminal_status": "",
                "primary_blocker": "",
            }
            for h in harnesses
        ],
    }


async def _run_harness_matrix_in_background(
    *,
    state: _HarnessMatrixJobState,
    req: HarnessMatrixRequest,
    env_overrides: dict[str, str],
    out_dir: Path,
) -> None:
    try:
        await asyncio.to_thread(
            run_harness_matrix,
            paper_path=req.paper_path,
            prompt=req.prompt,
            template=req.template,
            harnesses=[
                HarnessMatrixCellSpec(harness=h.id, model=h.model)
                for h in (req.harnesses or [HarnessMatrixHarnessRequest(id=h) for h in CODING_HARNESSES])
            ],
            attempts=req.attempts,
            timeout_s=req.timeout_s,
            concurrency=req.concurrency,
            reuse_ingest_run=req.reuse_ingest_run,
            env_overrides=env_overrides,
            out_dir=out_dir,
            matrix_id=state.matrix_id,
            cancel_event=state.cancel_event,
        )
    except Exception as exc:  # noqa: BLE001
        state.error = str(exc)
        snapshot = _initial_harness_matrix_snapshot(state.matrix_id, req, matrix_dir=state.matrix_dir)
        fallback = {
            **snapshot,
            "matrix_id": state.matrix_id,
            "status": "error",
            "error": str(exc),
            "matrix_dir": str(state.matrix_dir),
            "report_path": str(state.matrix_dir / "report.md"),
            "updated_at": datetime.now().isoformat(timespec="seconds"),
        }
        state.matrix_dir.mkdir(parents=True, exist_ok=True)
        atomic_write_json(state.matrix_dir / "matrix.json", fallback)


def _prepare_harness_matrix_response(matrix: dict[str, Any]) -> dict[str, Any]:
    """Patch run-file URLs at the API boundary so demo tokens stay centralized."""

    for row in matrix.get("rows", []):
        if not isinstance(row, dict):
            continue
        run_id = str(row.get("run_id") or "").strip()
        if not run_id:
            run_dir = str(row.get("run_dir") or "").strip()
            if run_dir:
                run_id = Path(run_dir).name
        if not run_id:
            continue
        if row.get("final_html"):
            row["final_html_url"] = _run_file_url(run_id, "final/poster.html")
        if row.get("preview_png"):
            row["preview_url"] = _run_file_url(run_id, "final/preview.png")
    return matrix


def _cancel_persisted_matrix_processes(matrix: dict[str, Any]) -> int:
    killed = 0
    for row in matrix.get("rows", []):
        if not isinstance(row, dict):
            continue
        if str(row.get("status") or "") not in {"pending", "running"}:
            continue
        raw_pgid = row.get("process_group_id") or row.get("process_id")
        try:
            pgid = int(raw_pgid)
        except (TypeError, ValueError):
            row["status"] = "cancelled"
            row["primary_blocker"] = row.get("primary_blocker") or "cancelled without persisted process id"
            continue
        try:
            os.killpg(pgid, signal.SIGTERM)
            killed += 1
        except ProcessLookupError:
            pass
        except Exception as exc:  # noqa: BLE001
            row["primary_blocker"] = f"cancel failed for process group {pgid}: {exc}"
            continue
        row["status"] = "cancelled"
        row["primary_blocker"] = row.get("primary_blocker") or "cancel requested"
    return killed


@app.post("/api/harness-matrix", response_model=HarnessMatrixAck)
async def start_harness_matrix(req: HarnessMatrixRequest, request: Request) -> HarnessMatrixAck:
    if _DEMO_MODE:
        raise HTTPException(
            403,
            detail={
                "code": "demo_matrix_disabled",
                "message": "Harness matrix runs are disabled in demo mode.",
            },
        )
    settings = _settings_for_request(request)
    env_overrides, _has_key = _request_env_overrides(request)
    paper = Path(req.paper_path).expanduser()
    if not paper.is_absolute():
        paper = (REPO_ROOT / paper).resolve()
    if not paper.exists() or not paper.is_file():
        raise HTTPException(400, detail=f"paper not found: {paper}")
    if not req.prompt.strip():
        raise HTTPException(400, detail="prompt is required")
    if resolve_template(req.template) is None:
        raise HTTPException(400, detail=f"unknown template: {req.template}")
    requested_harnesses = req.harnesses or [HarnessMatrixHarnessRequest(id=h) for h in CODING_HARNESSES]
    requested_ids = [h.id for h in requested_harnesses]
    duplicates = sorted({h for h in requested_ids if requested_ids.count(h) > 1})
    if duplicates:
        raise HTTPException(400, detail=f"duplicate harness in matrix request: {', '.join(duplicates)}")
    matrix_id = "matrix-" + datetime.now().strftime("%Y%m%d-%H%M%S") + "-" + secrets.token_hex(3)
    matrix_dir = settings.out_dir / "harness_matrix" / matrix_id
    state = _HarnessMatrixJobState(matrix_id=matrix_id, matrix_dir=matrix_dir)
    snapshot = _initial_harness_matrix_snapshot(matrix_id, req, matrix_dir=matrix_dir)
    matrix_dir.mkdir(parents=True, exist_ok=True)
    atomic_write_json(matrix_dir / "matrix.json", snapshot)
    async with _HARNESS_MATRIX_JOBS_LOCK:
        _HARNESS_MATRIX_JOBS[matrix_id] = state
    state.task = asyncio.create_task(
        _run_harness_matrix_in_background(
            state=state,
            req=req.model_copy(update={"paper_path": str(paper)}),
            env_overrides=env_overrides,
            out_dir=settings.out_dir,
        )
    )
    return HarnessMatrixAck(matrix_id=matrix_id, status="running", matrix=snapshot)


@app.get("/api/harness-matrix/{matrix_id}")
async def get_harness_matrix(matrix_id: str) -> dict[str, Any]:
    matrix_dir = await _lookup_matrix_dir(matrix_id)
    try:
        return _prepare_harness_matrix_response(load_harness_matrix(matrix_dir))
    except FileNotFoundError:
        raise HTTPException(404, detail=f"matrix not found: {matrix_id}") from None


@app.get("/api/harness-matrix/{matrix_id}/report")
async def get_harness_matrix_report(matrix_id: str) -> FileResponse:
    matrix_dir = await _lookup_matrix_dir(matrix_id)
    report_path = matrix_dir / "report.md"
    if not report_path.exists():
        raise HTTPException(404, detail=f"matrix report not found: {matrix_id}")
    return FileResponse(report_path, media_type="text/markdown", filename=f"{matrix_id}-report.md")


@app.get("/api/harness-matrix/{matrix_id}/events")
async def harness_matrix_events(matrix_id: str, request: Request) -> StreamingResponse:
    matrix_id = _safe_matrix_id(matrix_id)
    matrix_dir = await _lookup_matrix_dir(matrix_id)
    matrix_path = matrix_dir / "matrix.json"

    async def event_generator():
        yield ":hello\n\n"
        last_payload = ""
        deadline = time.monotonic() + _SSE_DEADLINE_S
        while time.monotonic() < deadline:
            if await request.is_disconnected():
                break
            if matrix_path.exists():
                raw_payload = matrix_path.read_text(encoding="utf-8")
                try:
                    data = _prepare_harness_matrix_response(json.loads(raw_payload))
                except json.JSONDecodeError:
                    await asyncio.sleep(1.0)
                    continue
                payload = json.dumps(data, ensure_ascii=False, separators=(",", ":"))
                if payload != last_payload:
                    last_payload = payload
                    yield f"data: {payload}\n\n"
                    if data.get("status") in {"completed", "completed_with_failures", "cancelled", "error"}:
                        break
            await asyncio.sleep(1.0)
            async with _HARNESS_MATRIX_JOBS_LOCK:
                state = _HARNESS_MATRIX_JOBS.get(matrix_id)
            if state is not None and state.task is not None and state.task.done() and not matrix_path.exists():
                yield f"data: {json.dumps({'matrix_id': matrix_id, 'status': 'error', 'error': state.error or 'matrix failed before ledger was written'}, ensure_ascii=False)}\n\n"
                break
    return StreamingResponse(
        event_generator(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache, no-transform",
            "X-Accel-Buffering": "no",
            "Connection": "keep-alive",
        },
    )


@app.post("/api/harness-matrix/{matrix_id}/cancel")
async def cancel_harness_matrix(matrix_id: str) -> dict[str, Any]:
    matrix_id = _safe_matrix_id(matrix_id)
    async with _HARNESS_MATRIX_JOBS_LOCK:
        state = _HARNESS_MATRIX_JOBS.get(matrix_id)
    if state is None:
        matrix_dir = await _lookup_matrix_dir(matrix_id)
        try:
            matrix = load_harness_matrix(matrix_dir)
        except FileNotFoundError:
            raise HTTPException(404, detail=f"matrix not found: {matrix_id}") from None
        killed = _cancel_persisted_matrix_processes(matrix)
        matrix["status"] = "cancelled"
        matrix["cancelled_at"] = datetime.now().isoformat(timespec="seconds")
        atomic_write_json(matrix_dir / "matrix.json", matrix)
        return {"ok": True, "matrix_id": matrix_id, "killed": killed, "matrix": _prepare_harness_matrix_response(matrix)}
    state.cancel_event.set()
    matrix_path = state.matrix_dir / "matrix.json"
    matrix = load_harness_matrix(state.matrix_dir) if matrix_path.exists() else {"matrix_id": matrix_id, "status": "cancelling"}
    return {"ok": True, "matrix_id": matrix_id, "matrix": _prepare_harness_matrix_response(matrix)}


# ---------------------------------------------------------------------------
# Coding-agent harness account login (local single-user only).
#
# Lets the Web UI drive `claude auth login` / `codex login` so the user can
# connect a personal subscription without dropping to a terminal. Credentials
# land in the isolated `harness_auth_dir`, keeping the daily `~/.claude` setup
# untouched. Disabled in demo mode (server-side credentials only).
# ---------------------------------------------------------------------------

_HARNESS_LOGIN_HARNESSES: frozenset[str] = frozenset({"claude", "codex"})
_HARNESS_LOGIN_DEADLINE_S = 600
_OAUTH_URL_RE = re.compile(r"https?://\S+")


class _HarnessLoginState:
    __slots__ = ("login_id", "harness", "status", "url", "message", "lines", "returncode", "proc", "created_at")

    def __init__(self, *, login_id: str, harness: str) -> None:
        self.login_id = login_id
        self.harness = harness
        self.status = "starting"  # starting | awaiting_user | success | failed | cancelled
        self.url = ""
        self.message = ""
        self.lines: list[str] = []
        self.returncode: int | None = None
        self.proc: subprocess.Popen | None = None
        self.created_at = time.time()


_HARNESS_LOGIN_JOBS: dict[str, _HarnessLoginState] = {}
_HARNESS_LOGIN_JOBS_LOCK = asyncio.Lock()


def _harness_login_snapshot(state: _HarnessLoginState) -> dict[str, Any]:
    return {
        "login_id": state.login_id,
        "harness": state.harness,
        "status": state.status,
        "url": state.url,
        "message": state.message,
        "lines": state.lines[-12:],
        "returncode": state.returncode,
    }


def _drive_harness_login_blocking(
    state: _HarnessLoginState, *, binary: str, args: list[str], env: dict[str, str]
) -> None:
    try:
        proc = subprocess.Popen(
            [binary, *args],
            stdin=subprocess.PIPE,
            stdout=subprocess.PIPE,
            stderr=subprocess.STDOUT,
            text=True,
            env=env,
            start_new_session=(os.name == "posix"),
        )
    except OSError as exc:
        state.status = "failed"
        state.message = f"failed to start {state.harness} login: {exc}"
        return
    state.proc = proc
    # Some CLIs prompt "press Enter to open the browser" — nudge it once.
    try:
        if proc.stdin is not None:
            proc.stdin.write("\n")
            proc.stdin.flush()
    except (BrokenPipeError, OSError):
        pass
    # Hard wall-clock guard so a hung login can't leak the reader thread.
    watchdog = threading.Timer(_HARNESS_LOGIN_DEADLINE_S, lambda: _terminate_subprocess_group(proc))
    watchdog.daemon = True
    watchdog.start()
    try:
        if proc.stdout is not None:
            for raw in proc.stdout:
                line = raw.rstrip("\n")
                if not line.strip():
                    continue
                state.lines.append(line)
                if len(state.lines) > 200:
                    state.lines = state.lines[-200:]
                if not state.url:
                    match = _OAUTH_URL_RE.search(line)
                    if match:
                        state.url = match.group(0).rstrip(".,)]}'\"")
                        if state.status not in {"cancelled", "success", "failed"}:
                            state.status = "awaiting_user"
                            state.message = "Open the login URL in your browser to finish signing in."
        returncode = proc.wait()
    finally:
        watchdog.cancel()
    state.returncode = returncode
    if state.status == "cancelled":
        return
    if returncode == 0:
        state.status = "success"
        state.message = "Login complete."
        mark_harness_login(state.harness)
    else:
        state.status = "failed"
        state.message = state.message or f"Login process exited with code {returncode}."


def _run_auth_status_blocking(cmd: list[str], env: dict[str, str]) -> dict[str, Any]:
    try:
        proc = subprocess.run(
            cmd, capture_output=True, text=True, timeout=20, env=env, stdin=subprocess.DEVNULL
        )
    except subprocess.TimeoutExpired:
        return {"logged_in": False, "account": None, "reason": "auth status timed out", "returncode": None}
    except OSError as exc:
        return {"logged_in": False, "account": None, "reason": str(exc), "returncode": None}
    raw = (proc.stdout or "").strip()
    parsed: Any = None
    try:
        parsed = json.loads(raw) if raw else None
    except json.JSONDecodeError:
        parsed = None
    logged_in = False
    account: str | None = None
    if isinstance(parsed, dict):
        for key in ("loggedIn", "authenticated", "isAuthenticated", "signedIn"):
            if isinstance(parsed.get(key), bool):
                logged_in = parsed[key]
                break
        acct = parsed.get("account")
        if isinstance(acct, dict):
            account = acct.get("emailAddress") or acct.get("email") or acct.get("organizationName")
        account = account or parsed.get("email") or parsed.get("emailAddress")
    # Fall back to a filesystem credential-presence check when the CLI gave us
    # no parseable verdict (older CLIs, unexpected JSON, non-zero exit).
    if not logged_in and proc.returncode == 0 and parsed is None:
        logged_in = True
    return {
        "logged_in": bool(logged_in),
        "account": account,
        "returncode": proc.returncode,
        "raw": (raw[:600] if raw else (proc.stderr or "").strip()[:600]),
    }


class HarnessLoginRequest(BaseModel):
    harness: str = "claude"


@app.post("/api/harness/login")
async def start_harness_login(req: HarnessLoginRequest, request: Request) -> dict[str, Any]:
    if _DEMO_MODE:
        raise HTTPException(
            403,
            detail={"code": "demo_login_disabled", "message": "Harness account login is disabled in demo mode."},
        )
    harness = (req.harness or "claude").strip().lower()
    if harness not in _HARNESS_LOGIN_HARNESSES:
        raise HTTPException(400, detail=f"unsupported harness for login: {harness}")
    binary = resolve_harness_binary(harness)
    if not binary:
        raise HTTPException(
            400,
            detail={"code": "harness_cli_missing", "message": f"{harness} CLI was not found on PATH."},
        )
    auth_dir = harness_auth_dir(harness)
    auth_dir.mkdir(parents=True, exist_ok=True)
    env = harness_subprocess_env(os.environ, harness=harness, config_dir=str(auth_dir))
    args = ["auth", "login", "--claudeai"] if harness == "claude" else ["login"]
    login_id = "login-" + datetime.now().strftime("%Y%m%d-%H%M%S") + "-" + secrets.token_hex(3)
    state = _HarnessLoginState(login_id=login_id, harness=harness)
    async with _HARNESS_LOGIN_JOBS_LOCK:
        _HARNESS_LOGIN_JOBS[login_id] = state
    asyncio.create_task(
        asyncio.to_thread(_drive_harness_login_blocking, state, binary=binary, args=args, env=env)
    )
    log("harness.login.start", harness=harness, login_id=login_id, config_dir=str(auth_dir))
    return _harness_login_snapshot(state)


@app.get("/api/harness/login/{login_id}/events")
async def harness_login_events(login_id: str, request: Request) -> StreamingResponse:
    async with _HARNESS_LOGIN_JOBS_LOCK:
        state = _HARNESS_LOGIN_JOBS.get(login_id)
    if state is None:
        raise HTTPException(404, detail=f"login session not found: {login_id}")

    async def event_generator():
        yield ":hello\n\n"
        last_payload = ""
        deadline = time.monotonic() + _HARNESS_LOGIN_DEADLINE_S
        while time.monotonic() < deadline:
            if await request.is_disconnected():
                break
            payload = json.dumps(_harness_login_snapshot(state), ensure_ascii=False, separators=(",", ":"))
            if payload != last_payload:
                last_payload = payload
                yield f"data: {payload}\n\n"
            if state.status in {"success", "failed", "cancelled"}:
                break
            await asyncio.sleep(0.5)

    return StreamingResponse(
        event_generator(),
        media_type="text/event-stream",
        headers={"Cache-Control": "no-cache, no-transform", "X-Accel-Buffering": "no", "Connection": "keep-alive"},
    )


@app.post("/api/harness/login/{login_id}/cancel")
async def cancel_harness_login(login_id: str) -> dict[str, Any]:
    async with _HARNESS_LOGIN_JOBS_LOCK:
        state = _HARNESS_LOGIN_JOBS.get(login_id)
    if state is None:
        raise HTTPException(404, detail=f"login session not found: {login_id}")
    state.status = "cancelled"
    state.message = "Login cancelled."
    if state.proc is not None:
        _terminate_subprocess_group(state.proc)
    return _harness_login_snapshot(state)


@app.get("/api/harness/auth-status")
async def harness_auth_status(harness: str = Query("claude")) -> dict[str, Any]:
    harness = (harness or "claude").strip().lower()
    if harness not in _HARNESS_LOGIN_HARNESSES:
        raise HTTPException(400, detail=f"unsupported harness: {harness}")
    binary = resolve_harness_binary(harness)
    auth_dir = harness_auth_read_dir(harness)
    managed_login = harness_login_present(harness)
    if not binary:
        return {
            "harness": harness,
            "available": False,
            "logged_in": harness_login_present(harness),
            "account": None,
            "config_dir": str(auth_dir),
            "message": f"{harness} CLI was not found.",
        }
    env = harness_subprocess_env(
        os.environ,
        harness=harness,
        config_dir=str(auth_dir) if managed_login else None,
    )
    cmd = [binary, "auth", "status", "--json"] if harness == "claude" else [binary, "login", "status"]
    result = await asyncio.to_thread(_run_auth_status_blocking, cmd, env)
    cli_logged_in = bool(result.get("logged_in"))
    # Keep the on-disk login marker in sync with the CLI's own verdict so the
    # author/smoke path (which trusts `harness_login_present`) reflects reality —
    # this also self-heals logins made before the marker existed.
    if managed_login:
        if cli_logged_in:
            mark_harness_login(harness, config_dir=auth_dir)
        elif result.get("returncode") is not None:
            clear_harness_login_marker(harness)
    return {
        "harness": harness,
        "available": True,
        "logged_in": cli_logged_in or harness_login_present(harness),
        "account": result.get("account"),
        "config_dir": str(auth_dir),
        "returncode": result.get("returncode"),
    }


@app.post("/api/design-events")
async def design_events(req: DesignEventRequest, request: Request) -> dict[str, Any]:
    allowed = {
        "artifact.opened",
        "artifact.downloaded",
        "openresearch.project_requested",
        "openresearch.project_ready",
        "openresearch.project_failed",
    }
    if req.event not in allowed:
        raise HTTPException(400, detail=f"unsupported design event: {req.event}")
    if req.run_id:
        _demo_require_run_owner(req.run_id, request)
    settings = SETTINGS if SETTINGS is not None else load_settings()
    path = append_design_event(
        _settings_or_boot(),
        req.conversation_id,
        req.event,
        run_id=req.run_id,
        artifact_id=req.artifact_id,
        data=req.data,
    )
    # Download is the strongest user-visible acceptance signal we have.
    # Write a companion artifact.accepted event so the memory extractor
    # can use a single event type for "user approved this artifact" without
    # needing to know about the download-as-proxy convention.
    if req.event == "artifact.downloaded":
        append_design_event(
            settings,
            req.conversation_id,
            "artifact.accepted",
            run_id=req.run_id,
            artifact_id=req.artifact_id,
            data={
                **req.data,
                "acceptance_signal": "download",
            },
        )
    return {"ok": True, "path": str(path)}


_EDITOR_ASSET_MAX_BYTES = 10 * 1024 * 1024
_EDITOR_ASSET_TYPES: dict[str, str] = {
    "image/png": ".png",
    "image/jpeg": ".jpg",
    "image/webp": ".webp",
    "image/gif": ".gif",
}


@app.post("/api/assets/upload", response_model=EditorAssetUploadResponse)
async def upload_editor_asset(
    file: UploadFile = File(...),
) -> EditorAssetUploadResponse:
    content_type = (file.content_type or "").split(";")[0].strip().lower()
    ext = _EDITOR_ASSET_TYPES.get(content_type)
    if not ext:
        raise HTTPException(
            400,
            detail="unsupported image type; upload PNG, JPEG, WebP, or GIF",
        )

    data = await file.read(_EDITOR_ASSET_MAX_BYTES + 1)
    if len(data) > _EDITOR_ASSET_MAX_BYTES:
        raise HTTPException(400, detail="image must be 10MB or smaller")
    if not data:
        raise HTTPException(400, detail="image file is empty")

    safe_stem = Path(file.filename or "image").stem[:48] or "image"
    safe_stem = "".join(ch if ch.isalnum() or ch in {"-", "_"} else "-" for ch in safe_stem)
    filename = f"{uuid.uuid4().hex[:12]}-{safe_stem}{ext}"
    out_path = EDITOR_ASSETS_DIR / filename
    out_path.write_bytes(data)
    return EditorAssetUploadResponse(
        url=f"/api/files/editor-assets/{filename}",
        filename=filename,
        content_type=content_type,
        size=len(data),
    )


@app.get("/api/artifacts/{artifact_id}/assets", response_model=list[ArtifactAsset])
async def artifact_assets(artifact_id: str, request: Request) -> list[ArtifactAsset]:
    run_id = _run_id_from_artifact_id(artifact_id)
    if not run_id:
        raise HTTPException(400, detail="invalid artifact_id")
    _require_run_owner_before_lookup(run_id, request)
    _assert_controlled_run_source_usable(run_id, mode="artifact")
    run_dir = RUNS_DIR / run_id
    if not run_dir.exists():
        raise HTTPException(404, detail=f"artifact not found: {artifact_id}")

    def authorize_context_run(context_run_id: str) -> None:
        _require_run_owner_before_lookup(context_run_id, request)
        _assert_controlled_run_source_usable(context_run_id, mode="artifact")

    run_dirs = _poster_revision_context_run_dirs(
        run_id,
        {"artifact_id": artifact_id},
        authorize_run=authorize_context_run,
    ) or [run_dir]
    return _collect_artifact_assets(run_dirs)


# ---------- /api/video/render ----------


def _validate_editable_video_run_assets(
    artifact: dict[str, Any],
    *,
    source_run_id: str,
) -> None:
    project = artifact.get("video_project")
    scenes = project.get("scenes") if isinstance(project, dict) else None
    if not isinstance(scenes, list):
        return
    prefix = "/api/files/runs/"
    layer_groups: list[Any] = [artifact.get("layers")]
    for scene in scenes:
        layer_groups.append(scene.get("layers") if isinstance(scene, dict) else None)
    for layers in layer_groups:
        if not isinstance(layers, list):
            continue
        for layer in layers:
            src = layer.get("src") if isinstance(layer, dict) else None
            if not isinstance(src, str):
                continue
            clean = src.split("?", 1)[0]
            if not clean.startswith(prefix):
                continue
            relative = clean[len(prefix):].lstrip("/")
            try:
                with _open_public_run_file(
                    relative,
                    expected_run_id=source_run_id,
                ):
                    pass
            except HTTPException as exc:
                raise HTTPException(
                    400,
                    detail="editable video run asset is unavailable or unauthorized",
                ) from exc


@app.post("/api/video/render", response_model=GenerateAck)
async def video_render(req: VideoRenderRequest, request: Request) -> GenerateAck:
    """Render a layer-mode editable video artifact to MP4.

    This is deliberately no-agent/no-LLM: the browser sends the current
    editable scene manifest, we write a deterministic HyperFrames
    project, run the HyperFrames CLI, and expose the MP4 through the same
    run artifact endpoint as normal video generations.
    """
    if _DEMO_MODE:
        raise HTTPException(
            403,
            detail={
                "code": "demo_poster_only",
                "message": "Demo mode only supports paper poster generation.",
            },
        )
    artifact = req.artifact
    if artifact.get("artifact_type") != "video":
        raise HTTPException(400, detail="artifact_type must be video")
    _require_artifact_runtime("video")
    project = artifact.get("video_project")
    scenes = project.get("scenes") if isinstance(project, dict) else None
    if not isinstance(scenes, list) or not scenes:
        raise HTTPException(400, detail="video_project.scenes is required")

    source_run_id = _run_id_from_artifact_id(str(artifact.get("artifact_id") or "")) or ""
    if not source_run_id:
        raise HTTPException(400, detail="video render requires a source artifact run")
    _demo_require_run_owner(
        source_run_id,
        request,
        detail="video source not found",
    )
    try:
        _require_derived_source_ready(source_run_id, "editable_video_render")
    except RunNotReady as exc:
        raise HTTPException(409, detail=str(exc)) from exc
    _validate_editable_video_run_assets(
        artifact,
        source_run_id=source_run_id,
    )
    settings = SETTINGS if SETTINGS is not None else load_settings()
    run_id = new_run_id()
    event_conversation_id = _event_conversation_id(req.conversation_id, run_id)
    baseline_json = json.dumps(
        {"artifact_id": artifact.get("artifact_id")},
        ensure_ascii=False,
    )
    access_user_id = _demo_register_derived_run_access(
        run_id,
        request,
        parent_run_id=source_run_id or None,
        missing_detail="video source not found",
    )
    state = _RunState(
        artifact_type="video",
        brief="editable video render",
        baseline_artifact_json=baseline_json,
        conversation_id=event_conversation_id,
    )
    state.demo_user_id = access_user_id
    async with _RUNS_LOCK:
        _RUNS[run_id] = state
    worker_request = EditableVideoRenderWorkerRequest(
        job_kind="editable_video_render",
        run_id=run_id,
        parent_run_id=source_run_id,
        artifact=artifact,
        conversation_id=event_conversation_id,
        baseline_artifact_json=baseline_json,
        settings=settings,
    )
    reserve_only = _request_reserve_only(request)
    try:
        start_token = await _start_supervised_derived_job(
            request=worker_request,
            state=state,
            descriptor={
                "job_kind": "editable_video_render",
                "source_artifact_id": str(artifact.get("artifact_id") or ""),
                "artifact_name": str(artifact.get("name") or "Video"),
                "source_relative_path": ".",
            },
            start_immediately=not reserve_only,
        )
    except Exception as exc:
        raise _web_run_service_error(exc) from exc
    _append_event(
        _settings_or_boot(),
        event_conversation_id,
        "artifact.render_started",
        run_id=run_id,
        artifact_id=str(artifact.get("artifact_id") or ""),
        data={"artifact_type": "video", "source": "editable_video_demo"},
    )

    return GenerateAck(
        run_id=run_id,
        progress_mode="video_render",
        start_token=start_token if reserve_only else None,
        placeholder_message=Message(
            id=f"msg_{run_id}",
            role="assistant",
            text="",
            ts=int(time.time() * 1000),
            run_id=run_id,
            status="streaming",
        ),
    )


# ---------- durable run reservation / upload / start ----------


def _run_owner_id(request: Request) -> str:
    return _demo_user_id(request) if _RUN_ACCESS_CONTROL else "local"


def _require_run_owner_before_lookup(run_id: str, request: Request) -> None:
    if _RUN_ACCESS_CONTROL and not _demo_user_owns_run(
        run_id,
        _demo_user_id(request),
    ):
        raise HTTPException(404, detail=f"run not found: {run_id}")


def _reservation_digest(owner_id: str, payload: RunReserveRequest) -> str:
    nonsecret = {
        "owner": owner_id,
        "request": payload.model_dump(mode="json"),
    }
    encoded = json.dumps(
        nonsecret,
        ensure_ascii=False,
        sort_keys=True,
        separators=(",", ":"),
    ).encode("utf-8")
    return hashlib.sha256(encoded).hexdigest()


def _reservation_idempotency_key(owner_id: str, value: str) -> str:
    clean = value.strip()
    if not clean:
        raise HTTPException(400, detail="Idempotency-Key header is required")
    owner_namespace = hashlib.sha256(owner_id.encode("utf-8")).hexdigest()[:24]
    return f"{owner_namespace}:{clean}"


def _prepare_reservation(
    request: Request,
    payload: RunReserveRequest,
) -> tuple[ArtifactType, Settings, dict[str, Any], _RunState]:
    settings = _settings_for_request(request)
    if _DEMO_MODE:
        _demo_require_host_safety()
        _demo_cleanup_expired_runs()
        if payload.artifact_type and payload.artifact_type != "poster":
            raise HTTPException(
                403,
                detail={
                    "code": "demo_poster_only",
                    "message": "Demo mode only supports paper poster generation.",
                },
            )
        attachments = [
            slot for slot in payload.input_slots if slot.role == "attachment"
        ]
        if len(attachments) != 1 or Path(attachments[0].name).suffix.lower() != ".pdf":
            raise HTTPException(
                400,
                detail={
                    "code": "demo_pdf_required",
                    "message": "Demo mode accepts one PDF paper per poster run.",
                },
            )
        if attachments[0].size > _DEMO_MAX_PDF_BYTES:
            raise HTTPException(
                413,
                detail={
                    "code": "demo_pdf_too_large",
                    "message": (
                        f"Demo PDF limit is {_DEMO_MAX_PDF_BYTES // (1024 * 1024)} MB."
                    ),
                },
            )
        if any(slot.role == "reference_poster" for slot in payload.input_slots):
            raise HTTPException(
                400,
                detail={
                    "code": "demo_reference_poster_unsupported",
                    "message": "Demo mode does not accept a reference poster.",
                },
            )
        a_type: ArtifactType = "poster"
        requested_template = _DEMO_FIXED_TEMPLATE
        requested_canvas_preset_id = _DEMO_FIXED_TEMPLATE
    else:
        a_type = _coerce_artifact_type(payload.artifact_type, brief=payload.brief)
        requested_template, requested_canvas_preset_id = _validated_web_canvas_selection(
            a_type,
            payload.template,
            payload.canvas_preset_id,
        )
    _validated_canvas_prompt(payload.brief, a_type)
    attempts = _validated_authoring_max_attempts(
        payload.authoring_max_attempts,
        a_type,
        settings,
    )
    settings = _settings_with_authoring_max_attempts(settings, a_type, attempts)
    _require_artifact_runtime(a_type)
    palette_id = _validated_web_palette_id(a_type, payload.palette_id)
    names = [slot.name for slot in payload.input_slots if slot.role == "attachment"]
    has_pdf = any(Path(name).suffix.lower() == ".pdf" for name in names)
    reference_slots = [slot for slot in payload.input_slots if slot.role == "reference_poster"]
    if len(reference_slots) > 1:
        raise HTTPException(422, detail="only one reference poster input is supported")
    if reference_slots and a_type != "poster":
        raise HTTPException(400, detail="reference poster style is only supported for poster generation")
    if reference_slots:
        _validated_web_reference_poster_name(reference_slots[0].name)
        if not has_pdf:
            raise HTTPException(400, detail="attach a paper PDF when using a reference poster style")
    effective_template = requested_template
    if a_type == "poster" and has_pdf:
        settings = _web_paper_poster_settings(settings)
        author_cmd = _paper_poster_author_cmd_resolution(settings)
        if not author_cmd["available"]:
            raise HTTPException(
                412,
                detail={
                    "code": "missing_external_author_command",
                    "message": author_cmd["message"],
                },
            )
        user_picked_designer = bool(
            (request.headers.get("x-model-designer", "") or "").strip()
            or (request.headers.get("x-model-planner", "") or "").strip()
        )
        if (
            not user_picked_designer
            and _KIMI_PATTERN in settings.designer_model.lower()
        ):
            log(
                "designer.auto_switched",
                from_model=settings.designer_model,
                to_model=_OPUS_FALLBACK,
                reason="paper_poster_kimi_stall",
            )
            settings = replace(settings, designer_model=_OPUS_FALLBACK)
    effective_brief = _apply_type_prologue(payload.brief, a_type)
    effective_brief = _apply_conversation_prologue(
        effective_brief,
        history_json=payload.conversation_history,
        artifacts_json=payload.prior_artifacts,
    )
    run_payload = {
        "brief": effective_brief,
        "slot_roles": {slot.name: slot.role for slot in payload.input_slots},
        "template": effective_template,
        "canvas_preset_id": requested_canvas_preset_id,
        "palette_id": palette_id,
        "resume_run": None,
        "baseline_artifact": payload.baseline_artifact,
        "conversation_id": payload.conversation_id,
    }
    state = _RunState(
        artifact_type=a_type,
        designer_model=settings.designer_model,
        has_pdf=has_pdf,
        brief=effective_brief,
        baseline_artifact_json=payload.baseline_artifact,
        conversation_id="",
        template=effective_template,
        canvas_preset_id=requested_canvas_preset_id,
        palette_id=palette_id,
        authoring_max_attempts=attempts,
        input_slot_roles={slot.name: slot.role for slot in payload.input_slots},
    )
    return a_type, settings, run_payload, state


def _pipeline_request_factory(
    run_id: str,
    settings: Settings,
    payload: Any,
    completed_slots: dict[str, Path],
) -> PipelineWorkerRequest:
    if not isinstance(payload, dict):
        raise RunNotReady("pipeline reservation payload is unavailable")
    roles = payload.get("slot_roles")
    if not isinstance(roles, dict):
        roles = {}
    attachments = tuple(
        str(path)
        for name, path in sorted(completed_slots.items())
        if roles.get(name, "attachment") == "attachment"
    ) or tuple(str(path) for path in payload.get("direct_attachments", ()))
    references = [
        path for name, path in completed_slots.items()
        if roles.get(name) == "reference_poster"
    ]
    direct_reference = payload.get("direct_reference")
    if not references and direct_reference:
        references = [Path(str(direct_reference))]
    if references:
        _validate_web_reference_poster_file(references[0])
    state = _RUNS.get(run_id)
    if state is not None:
        state.attach_paths = [Path(value) for value in attachments]
        state.reference_poster_path = references[0] if references else None
        if references and not state.reference_poster_handle:
            state.reference_poster_handle = f"ref_{uuid.uuid4().hex}"
    return PipelineWorkerRequest(
        job_kind="pipeline",
        run_id=run_id,
        brief=str(payload.get("brief") or ""),
        attachments=attachments,
        template=payload.get("template"),
        palette_id=payload.get("palette_id"),
        resume_run=payload.get("resume_run"),
        reference_poster=str(references[0]) if references else None,
        settings=settings,
        canvas_preset_id=payload.get("canvas_preset_id"),
    )


def _legacy_pipeline_payload(
    *,
    brief: str,
    attach_paths: list[Path],
    reference_poster_path: Path | None,
    template: str | None,
    state: _RunState,
    resume_run: str | None,
) -> dict[str, Any]:
    return {
        "brief": brief,
        "slot_roles": {},
        "direct_attachments": [str(path) for path in attach_paths],
        "direct_reference": str(reference_poster_path) if reference_poster_path else None,
        "template": template,
        "canvas_preset_id": state.canvas_preset_id,
        "palette_id": state.palette_id,
        "resume_run": resume_run,
        "baseline_artifact": state.baseline_artifact_json,
        "conversation_id": state.conversation_id,
    }


async def _reserve_legacy_pipeline_worker(
    *,
    run_id: str,
    brief: str,
    attach_paths: list[Path],
    reference_poster_path: Path | None,
    template: str | None,
    state: _RunState,
    settings: Settings,
    resume_run: str | None,
) -> None:
    async with _web_run_start_guard(run_id):
        reservation = await _web_run_runtime().services.reserve(
            run_id=run_id,
            artifact_type=state.artifact_type,
            idempotency_key=f"legacy:{run_id}",
            request_digest=hashlib.sha256(run_id.encode("utf-8")).hexdigest(),
            settings=settings,
            payload=_legacy_pipeline_payload(
                brief=brief,
                attach_paths=attach_paths,
                reference_poster_path=reference_poster_path,
                template=template,
                state=state,
                resume_run=resume_run,
            ),
            input_slots=(),
        )
        state.reservation_token = reservation.upload_token


async def _start_legacy_pipeline_worker(
    *,
    run_id: str,
    brief: str,
    attach_paths: list[Path],
    reference_poster_path: Path | None,
    template: str | None,
    state: _RunState,
    settings: Settings | None,
    resume_run: str | None,
) -> None:
    async with _web_run_start_guard(run_id):
        runtime = _web_run_runtime()
        if not (RUNS_DIR / run_id / "run_control.json").is_file():
            if settings is None:
                raise RunNotReady("run reservation context is unavailable")
            await _reserve_legacy_pipeline_worker(
                run_id=run_id,
                brief=brief,
                attach_paths=attach_paths,
                reference_poster_path=reference_poster_path,
                template=template,
                state=state,
                settings=settings,
                resume_run=resume_run,
            )
        if not state.reservation_token:
            raise RunNotReady("reserved run no longer retains its start token")
        await runtime.services.start(
            run_id,
            state.reservation_token,
            _pipeline_request_factory,
        )
        state.reservation_token = ""


_DERIVED_JOB_VERSION = 1
_DERIVED_JOB_KINDS = frozenset({
    "artifact_edit",
    "attempt_fork",
    "candidate_publish",
    "editable_video_render",
    "poster_code_edit",
    "pptx_export",
    "video_export_retry",
})


def _require_derived_source_ready(
    parent_run_id: str,
    job_kind: str,
    request: RunWorkerRequest | None = None,
) -> None:
    if job_kind not in _DERIVED_JOB_KINDS:
        raise RunNotReady(f"unsupported derived job kind: {job_kind}")
    parent_dir = RUNS_DIR / parent_run_id
    if not parent_dir.is_dir():
        raise RunNotReady(f"derived source run is missing: {parent_run_id}")
    control_path = parent_dir / "run_control.json"
    if not control_path.is_file():
        source_state = _RUNS.get(parent_run_id)
        source_task = getattr(source_state, "task", None)
        if source_task is not None and not source_task.done():
            raise RunNotReady("legacy derived source is still active")
        return
    try:
        record = RunControlStore(RUNS_DIR).read(parent_run_id)
    except RunControlError as exc:
        raise RunNotReady("derived source control is unreadable") from exc
    if job_kind == "attempt_fork":
        if record.state in {"running", "completing", "failed"} or (
            record.state == "completed" and record.publishable
        ):
            return
        raise RunNotReady(
            f"attempt fork source is unavailable: {record.state}"
        )
    if (
        job_kind == "candidate_publish"
        and isinstance(request, CandidatePublishWorkerRequest)
        and request.source_attempt is not None
        and request.expected_candidate_sha256 is not None
    ):
        if record.state in {"running", "completing", "failed"} or (
            record.state == "completed" and record.publishable
        ):
            return
        raise RunNotReady(
            f"candidate publish source is unavailable: {record.state}"
        )
    if job_kind == "artifact_edit":
        if record.state in {"running", "completing"} or (
            record.state == "completed" and record.publishable
        ):
            return
        raise RunNotReady(
            f"artifact edit source is unavailable: {record.state}"
        )
    if job_kind == "video_export_retry":
        if record.state == "failed" or (
            record.state == "completed" and record.publishable
        ):
            return
        raise RunNotReady(
            f"video export retry source is not quiescent: {record.state}"
        )
    if record.state != "completed" or not record.publishable:
        raise RunNotReady(
            f"derived source is not publishable: {record.state}"
        )


def _require_derived_ancestors_not_cancelled(parent_run_id: str) -> None:
    store = _web_run_runtime().control_store
    for ancestor_run_id in _derived_ancestor_chain(parent_run_id):
        control_path = RUNS_DIR / ancestor_run_id / "run_control.json"
        if not control_path.is_file():
            if _read_derived_job_descriptor(ancestor_run_id) is not None:
                raise RunNotReady("derived run ancestor control is missing")
            continue
        try:
            record = store.read(ancestor_run_id)
        except RunControlError as exc:
            raise RunNotReady("derived run ancestor control is unreadable") from exc
        if record.state in {"cancelling", "cancelled"} or record.writes_frozen:
            raise RunNotReady(
                f"derived run ancestor is cancelling: {ancestor_run_id}"
            )


def _derived_request_factory(
    run_id: str,
    settings: Settings,
    payload: Any,
    _completed: Any,
) -> RunWorkerRequest:
    if not isinstance(payload, dict):
        raise RunNotReady("derived start payload is unavailable")
    request = payload.get("request")
    if not isinstance(
        request,
        (
            ArtifactEditWorkerRequest,
            EditableVideoRenderWorkerRequest,
            AttemptForkWorkerRequest,
            CandidatePublishWorkerRequest,
            PosterCodeEditWorkerRequest,
            PptxExportWorkerRequest,
            VideoExportRetryWorkerRequest,
        ),
    ):
        raise RunNotReady("derived start request is invalid")
    if request.run_id != run_id:
        raise RunNotReady("derived start request lost its reservation identity")
    if hasattr(request, "settings") and request.settings is not settings:
        raise RunNotReady("derived start request lost its Settings identity")
    if (
        request.job_kind == "video_export_retry"
        and Path(request.runs_dir).resolve() != RUNS_DIR.resolve()
    ):
        raise RunNotReady("video export retry lost its runs directory identity")
    _require_derived_source_ready(
        request.parent_run_id,
        request.job_kind,
        request,
    )
    return request


def _runtime_only_settings() -> Settings:
    """Provide non-secret process context for derived jobs that do not call models."""
    return Settings(
        anthropic_api_key="",
        anthropic_base_url=None,
        gemini_api_key="",
        designer_model="runtime-only",
        critic_model="runtime-only",
        repo_root=Path(__file__).resolve().parents[1],
        out_dir=RUNS_DIR.parent,
    )


def _request_reserve_only(request: Request) -> bool:
    return request.headers.get("x-autodesign-reserve-only", "").strip().lower() in {
        "1",
        "true",
        "yes",
    }


def _durable_derived_job_payload(
    request: RunWorkerRequest,
    state: _RunState,
    descriptor: dict[str, Any],
) -> dict[str, Any]:
    return {
        "version": _DERIVED_JOB_VERSION,
        **descriptor,
        "run_id": request.run_id,
        "parent_run_id": request.parent_run_id,
        "artifact_type": state.artifact_type,
        "conversation_id": state.conversation_id,
        "baseline_artifact_json": state.baseline_artifact_json or "",
    }


async def _start_supervised_derived_job(
    *,
    request: RunWorkerRequest,
    state: _RunState,
    descriptor: dict[str, Any],
    start_immediately: bool = True,
    settings: Settings | None = None,
    descriptor_prepared: bool = False,
) -> str:
    run_id = request.run_id
    async with _derived_tree_locks(request.parent_run_id), _web_run_start_guard(run_id):
        if request.job_kind not in _DERIVED_JOB_KINDS:
            raise RunNotReady("derived job kind is not supported")
        if descriptor.get("job_kind") != request.job_kind:
            raise RunNotReady("derived descriptor does not match its request")
        _require_derived_ancestors_not_cancelled(request.parent_run_id)
        _require_derived_source_ready(
            request.parent_run_id,
            request.job_kind,
            request,
        )
        durable_descriptor = _durable_derived_job_payload(
            request,
            state,
            descriptor,
        )
        if descriptor_prepared:
            try:
                prepared_descriptor = _read_derived_job_descriptor(run_id)
            except ValueError as exc:
                raise RunNotReady("prepared derived descriptor is unreadable") from exc
            if prepared_descriptor != durable_descriptor:
                raise RunNotReady("prepared derived descriptor changed")
        runtime = _web_run_runtime()
        digest = hashlib.sha256(
            json.dumps(
                descriptor,
                ensure_ascii=False,
                sort_keys=True,
                separators=(",", ":"),
            ).encode("utf-8")
        ).hexdigest()
        request_settings = settings or getattr(request, "settings", None) or _runtime_only_settings()
        reservation = await runtime.services.reserve(
            run_id=run_id,
            artifact_type=state.artifact_type,
            idempotency_key=f"derived:{run_id}",
            request_digest=digest,
            settings=request_settings,
            payload={"request": request},
            input_slots=(),
            parent_job_id=request.parent_run_id,
        )
        state.reservation_token = reservation.upload_token
        if not descriptor_prepared:
            durable_replace_json(
                RUNS_DIR / run_id / "derived_job.json",
                durable_descriptor,
            )
            try:
                _require_derived_ancestors_not_cancelled(request.parent_run_id)
                _require_derived_source_ready(
                    request.parent_run_id,
                    request.job_kind,
                    request,
                )
            except BaseException:
                await runtime.services.cancel(run_id, "derived_source_unavailable")
                raise
        if start_immediately:
            await _start_reserved_derived_job(
                run_id=run_id,
                token=reservation.upload_token,
                state=state,
                descriptor=durable_descriptor,
            )
        return reservation.upload_token


async def _start_reserved_derived_job(
    *,
    run_id: str,
    token: str,
    state: _RunState,
    descriptor: dict[str, Any],
) -> None:
    try:
        await _web_run_runtime().services.start(
            run_id,
            token,
            _derived_request_factory,
        )
        async with _RUNS_LOCK:
            state.reservation_token = ""
            if state.task is None:
                state.task = asyncio.create_task(
                    _monitor_supervised_derived_job(
                        run_id=run_id,
                        state=state,
                        job_kind=str(descriptor["job_kind"]),
                        parent_run_id=str(descriptor["parent_run_id"]),
                        descriptor=descriptor,
                    )
                )
    except (UploadAuthorizationError, ReservationNotFound):
        raise
    except BaseException as start_error:
        cancel_task = asyncio.create_task(
            _web_run_runtime().services.cancel(run_id, "derived_start_failed")
        )
        await _join_owned_task(cancel_task)
        try:
            cancellation = cancel_task.result()
            if cancellation.cancel_request_event_required:
                _append_cancel_request_event(run_id, "derived_start_failed")
        except BaseException as cancel_error:
            log(
                "derived_start.cleanup_deferred",
                run_id=run_id,
                start_error=type(start_error).__name__,
                cancel_error=type(cancel_error).__name__,
            )
        raise


def _web_run_service_error(exc: Exception) -> HTTPException:
    if isinstance(exc, (ReservationConflict, UploadConflict, RunNotReady)):
        return HTTPException(409, detail=str(exc))
    if isinstance(exc, (UploadAuthorizationError, ReservationNotFound)):
        return HTTPException(404, detail="run reservation not found")
    if isinstance(exc, UploadCancelled):
        return HTTPException(409, detail={"code": "run_cancelling", "message": str(exc)})
    if isinstance(exc, (InvalidReservation, InvalidInputSlot, UploadIntegrityError)):
        return HTTPException(422, detail=str(exc))
    return HTTPException(500, detail="run lifecycle operation failed")


def _append_reservation_history(
    *,
    settings: Settings,
    state: _RunState,
    payload: RunReserveRequest,
    run_id: str,
) -> None:
    submitted: dict[str, Any] = {
        "brief": payload.brief.strip(),
        "artifact_type": state.artifact_type,
        "palette_id": state.palette_id,
        "template": state.template,
        "canvas_preset_id": state.canvas_preset_id,
        "has_baseline": payload.baseline_artifact is not None,
        "authoring_max_attempts": state.authoring_max_attempts,
        "attachment_count": sum(
            slot.role == "attachment" for slot in payload.input_slots
        ),
    }
    reference = next(
        (slot for slot in payload.input_slots if slot.role == "reference_poster"),
        None,
    )
    if reference is not None:
        submitted["has_reference_poster"] = True
        submitted["reference_poster_name"] = reference.name
    if payload.baseline_artifact is not None:
        submitted["follow_up_sentiment"] = infer_follow_up_sentiment(payload.brief)
    _append_event(
        settings,
        state.conversation_id,
        "message.user_submitted",
        run_id=run_id,
        data=submitted,
    )


_PAPER_BUNDLE_ARTIFACT_TYPES = frozenset({"poster", "deck", "landing", "video"})


def _paper_bundle_request_digest(
    owner_id: str,
    payload: PaperBundleCreateRequest,
) -> str:
    encoded = json.dumps(
        {"owner": owner_id, "request": payload.model_dump(mode="json")},
        ensure_ascii=False,
        sort_keys=True,
        separators=(",", ":"),
    ).encode("utf-8")
    return hashlib.sha256(encoded).hexdigest()


def _paper_bundle_child_snapshot(run_id: str) -> ChildStateSnapshot:
    runtime = _web_run_runtime()
    try:
        record = runtime.control_store.read(run_id)
    except RunControlError:
        return ChildStateSnapshot(
            state="failed",
            terminal=True,
            process_free=False,
            diagnostic="child_control_missing",
        )
    terminal = record.state in {"completed", "cancelled", "failed"}
    # PaperBundleJobStore invokes this synchronous provider from a worker thread.
    process_free = terminal and runtime.supervisor.is_durably_quiescent(run_id)
    return ChildStateSnapshot(
        state=record.state,
        terminal=terminal,
        process_free=process_free,
        diagnostic=record.cancellation_pending,
    )


def _paper_bundle_parent_for_run(run_id: str) -> str | None:
    try:
        record = _web_run_runtime().control_store.read(run_id)
    except RunControlError:
        return None
    return str(record.parent_job_id or "").strip() or None


def _reserve_candidate_publish_bundle_binding(
    source_run_id: str,
    artifact_type: str,
    owner_id: str,
) -> dict[str, Any] | None:
    try:
        source_control = _web_run_runtime().control_store.read(source_run_id)
    except RunControlError:
        return None
    parent_job_id = str(source_control.parent_job_id or "").strip()
    if not parent_job_id:
        return None
    try:
        derived_descriptor = _read_derived_job_descriptor(source_run_id)
    except ValueError as exc:
        raise PaperBundleConflict(
            "candidate publication source descriptor is unreadable"
        ) from exc
    if derived_descriptor is not None:
        return None
    store = _paper_bundle_store()
    parent = store.read_owned(
        parent_job_id,
        owner_id,
        child_status_provider=_paper_bundle_child_snapshot,
    )
    child = parent.children.get(artifact_type)
    if child is None or child.run_id != source_run_id:
        raise PaperBundleConflict(
            "candidate publication source does not match the bundle child"
        )
    generation = store.reserve_child_publication(
        parent_job_id,
        owner_id,
        artifact_type,
        source_run_id,
    )
    return {
        "paper_bundle_job_id": parent_job_id,
        "paper_bundle_owner_id": owner_id,
        "paper_bundle_artifact_type": artifact_type,
        "publication_generation": generation,
    }


def _validate_replayed_candidate_publish_binding(
    *,
    source_run_id: str,
    artifact_type: str,
    candidate_id: str,
    owner_id: str,
    write_ahead_request: dict[str, Any],
) -> None:
    if write_ahead_request["source_candidate_id"] != candidate_id:
        raise PaperBundleConflict(
            "candidate publication source identity changed"
        )
    version = int(write_ahead_request["version"])
    try:
        source_control = RunControlStore(RUNS_DIR).read(source_run_id)
    except RunControlError as exc:
        if version != 2:
            return
        raise PaperBundleConflict(
            "candidate publication source control is unavailable"
        ) from exc
    parent_job_id = str(source_control.parent_job_id or "").strip()
    source_descriptor: dict[str, Any] | None = None
    if parent_job_id or version == 2:
        try:
            source_descriptor = _read_derived_job_descriptor(source_run_id)
        except ValueError as exc:
            raise PaperBundleConflict(
                "candidate publication source descriptor is unreadable"
            ) from exc
    if version != 2:
        if parent_job_id and source_descriptor is None:
            raise PaperBundleConflict(
                "candidate publication bundle binding is missing"
            )
        return
    if (
        write_ahead_request["paper_bundle_owner_id"] != owner_id
        or write_ahead_request["paper_bundle_artifact_type"] != artifact_type
    ):
        raise PaperBundleConflict(
            "candidate publication bundle binding changed"
        )
    if (
        source_control.artifact_type != artifact_type
        or parent_job_id != write_ahead_request["paper_bundle_job_id"]
    ):
        raise PaperBundleConflict(
            "candidate publication bundle child changed"
        )
    if source_descriptor is not None:
        raise PaperBundleConflict(
            "candidate publication source is no longer a bundle child"
        )
    parent = _paper_bundle_store().read_owned(parent_job_id, owner_id)
    child = parent.children.get(artifact_type)
    generation = int(write_ahead_request["publication_generation"])
    if (
        child is None
        or child.run_id != source_run_id
        or parent.cancel_requested
        or parent.state in {"cancelling", "cancelled"}
        or parent.publication_generations.get(artifact_type) != generation
    ):
        raise PaperBundleConflict(
            "candidate publication bundle reservation changed"
        )
    committed = parent.publications.get(artifact_type)
    if committed is not None and committed.generation >= generation:
        raise PaperBundleConflict(
            "candidate publication bundle generation is already committed"
        )


async def _validate_replayed_candidate_publish_binding_async(
    **kwargs: Any,
) -> None:
    task = asyncio.create_task(
        asyncio.to_thread(_validate_replayed_candidate_publish_binding, **kwargs)
    )
    caller_cancelled = await _join_owned_task(task)
    if caller_cancelled:
        try:
            task.result()
        except BaseException as exc:
            log(
                "candidate_publish.binding_validation_finished_after_cancel",
                source_run_id=str(kwargs.get("source_run_id") or ""),
                error=type(exc).__name__,
            )
        raise asyncio.CancelledError
    task.result()


async def _reserve_candidate_publish_bundle_binding_async(
    source_run_id: str,
    artifact_type: str,
    owner_id: str,
) -> dict[str, Any] | None:
    async with _run_tree_lock(source_run_id):
        task = asyncio.create_task(
            asyncio.to_thread(
                _reserve_candidate_publish_bundle_binding,
                source_run_id,
                artifact_type,
                owner_id,
            )
        )
        caller_cancelled = await _join_owned_task(task)
        if caller_cancelled:
            try:
                task.result()
            except BaseException as exc:
                log(
                    "candidate_publish.binding_finished_after_cancel",
                    source_run_id=source_run_id,
                    error=type(exc).__name__,
                )
            raise asyncio.CancelledError
        return task.result()


async def _quiesce_bundle_candidate_publication_source(
    publication_run_id: str,
) -> None:
    """Quiesce an active source after validation and before publication commit.

    The caller must hold the source's derived-tree lock. Bundle publication
    reservations take the same lock, so the generation check stays current
    until the source is terminal and process-free.
    """
    direct_publish = _read_direct_candidate_publish_descriptor(
        publication_run_id
    )
    if direct_publish is None or direct_publish["version"] not in {2, 3}:
        return
    store = _paper_bundle_store()
    job_id = str(direct_publish["paper_bundle_job_id"])
    owner_id = str(direct_publish["paper_bundle_owner_id"])
    artifact_type = str(direct_publish["paper_bundle_artifact_type"])
    source_run_id = str(direct_publish["source_run_id"])
    generation = int(direct_publish["publication_generation"])
    parent = await asyncio.to_thread(
        store.read_owned,
        job_id,
        owner_id,
        child_status_provider=_paper_bundle_child_snapshot,
    )
    allocated_generation = parent.publication_generations[artifact_type]
    if generation < allocated_generation:
        raise PaperBundleConflict(
            "candidate publication was superseded by a newer request"
        )
    if generation != allocated_generation:
        raise PaperBundleConflict("publication generation was not reserved")
    if parent.cancel_requested or parent.state in {"cancelling", "cancelled"}:
        raise PaperBundleBarrierClosed(
            "paper bundle cancellation won the publication race"
        )
    child = parent.children[artifact_type]
    if child.run_id != source_run_id:
        raise PaperBundleConflict(
            "publication source does not match the bundle child"
        )
    if child.state in {"running", "completing"}:
        selection = await asyncio.to_thread(
            request_attempt_selection,
            run_dir=RUNS_DIR / source_run_id,
            run_id=source_run_id,
            attempt=int(direct_publish["source_attempt"]),
            expected_candidate_sha256=str(
                direct_publish["source_candidate_sha256"]
            ),
            idempotency_key=f"candidate-publication:{publication_run_id}",
            writable_guard=lambda: RunControlStore(RUNS_DIR).assert_writable(
                source_run_id
            ),
        )
        if selection.status not in {"selection_accepted", "already_selected"}:
            raise PaperBundleConflict(
                f"source attempt could not be selected: {selection.status}"
            )
        _schedule_attempt_selection_recovery(
            source_run_id,
            SETTINGS or _runtime_only_settings(),
        )
    elif child.state not in {"completed", "failed"}:
        raise PaperBundleConflict(
            "publication source child cannot be safely quiesced"
        )

    joined = await _quiesce_web_completion_monitor(
        source_run_id,
        timeout_s=120.0,
    )
    snapshot = _paper_bundle_child_snapshot(source_run_id)
    if (
        not joined
        or snapshot.state not in {"completed", "failed"}
        or not snapshot.terminal
        or not snapshot.process_free
    ):
        raise PaperBundleConflict(
            "publication source child did not become quiescent"
        )
    await asyncio.to_thread(
        store.read_owned,
        job_id,
        owner_id,
        child_status_provider=_paper_bundle_child_snapshot,
    )


async def _join_owned_task(task: asyncio.Task[Any]) -> bool:
    """Join owned work without letting repeated caller cancellation orphan it."""
    caller_cancelled = False
    while not task.done():
        try:
            await asyncio.shield(task)
        except asyncio.CancelledError:
            caller_cancelled = True
        except BaseException:
            break
    return caller_cancelled


def _assert_paper_bundle_child_work_allowed(
    run_id: str,
    owner_id: str,
) -> str | None:
    parent_job_id = _paper_bundle_parent_for_run(run_id)
    if parent_job_id is None:
        return None
    _paper_bundle_store().assert_child_may_upload_or_start(
        parent_job_id,
        run_id,
        owner_id,
    )
    return parent_job_id


async def _start_run_with_paper_bundle_barrier(
    run_id: str,
    token: str,
    owner_id: str,
):
    async with _web_run_start_guard(run_id):
        return await _start_run_with_paper_bundle_barrier_unchecked(
            run_id,
            token,
            owner_id,
        )


async def _start_run_with_paper_bundle_barrier_unchecked(
    run_id: str,
    token: str,
    owner_id: str,
):
    runtime = _web_run_runtime()
    parent_job_id = _paper_bundle_parent_for_run(run_id)
    if parent_job_id is None:
        return await runtime.services.start(
            run_id,
            token,
            _pipeline_request_factory,
        )
    store = _paper_bundle_store()
    intent_id = f"start-{run_id}"
    intent = await asyncio.to_thread(
        store.claim_child_start,
        parent_job_id,
        run_id,
        owner_id,
        intent_id=intent_id,
    )
    if intent.state == "registered":
        await asyncio.to_thread(
            store.assert_child_may_upload_or_start,
            parent_job_id,
            run_id,
            owner_id,
        )
        return await runtime.services.start(
            run_id,
            token,
            _pipeline_request_factory,
        )
    await asyncio.to_thread(
        store.commit_child_start,
        parent_job_id,
        run_id,
        owner_id,
        intent_id,
    )
    try:
        supervised = await runtime.services.start(
            run_id,
            token,
            _pipeline_request_factory,
        )
    except BaseException:
        try:
            await asyncio.to_thread(
                store.resolve_child_start,
                parent_job_id,
                run_id,
                owner_id,
                intent_id,
                "aborted",
            )
        except PaperBundleError:
            pass
        raise
    try:
        await asyncio.to_thread(
            store.resolve_child_start,
            parent_job_id,
            run_id,
            owner_id,
            intent_id,
            "registered",
        )
    except PaperBundleError:
        await _cancel_controlled_run(
            run_id,
            f"paper_bundle_start_cancel_race:{parent_job_id}",
        )
        try:
            await asyncio.to_thread(
                store.resolve_child_start,
                parent_job_id,
                run_id,
                owner_id,
                intent_id,
                "aborted",
            )
        except PaperBundleError:
            pass
        raise
    return supervised


def _paper_bundle_error(exc: PaperBundleError) -> HTTPException:
    if isinstance(exc, PaperBundleNotFound):
        return HTTPException(404, detail="paper bundle not found")
    if isinstance(exc, (PaperBundleBarrierClosed, PaperBundleConflict)):
        return HTTPException(409, detail=str(exc))
    return HTTPException(422, detail=str(exc))


async def _cleanup_paper_bundle_child(run_id: str) -> None:
    if not (RUNS_DIR / run_id / "run_control.json").is_file():
        return
    try:
        result = await _cancel_controlled_run(run_id, "paper_bundle_creation_cleanup")
    except RunControlError:
        return
    if not result.confirmed:
        raise PaperBundleError(
            f"paper bundle child cancellation is still pending: {run_id}"
        )


async def _reconcile_paper_bundle_for_run(
    run_id: str,
    *,
    owner_id: str | None = None,
) -> PaperBundleJobRecord | None:
    parent_job_id = _paper_bundle_parent_for_run(run_id)
    if parent_job_id is None:
        return None
    if owner_id is None:
        async with _RUNS_LOCK:
            state = _RUNS.get(run_id)
        owner_id = (
            state.demo_user_id
            if state is not None and state.demo_user_id
            else "local"
        )
    try:
        return await asyncio.to_thread(
            _paper_bundle_store().read_owned,
            parent_job_id,
            owner_id,
            child_status_provider=_paper_bundle_child_snapshot,
        )
    except PaperBundleError:
        return None


async def _reconcile_all_paper_bundles() -> tuple[PaperBundleJobRecord, ...]:
    """Recover every durable parent without relying on in-memory owner state."""
    return await asyncio.to_thread(
        _paper_bundle_store().reconcile_all,
        _paper_bundle_child_snapshot,
    )


@app.post("/api/paper-bundles")
async def create_paper_bundle(
    payload: PaperBundleCreateRequest,
    request: Request,
) -> dict[str, Any]:
    if set(payload.children) != _PAPER_BUNDLE_ARTIFACT_TYPES:
        raise HTTPException(
            422,
            detail="paper bundle requires poster, deck, landing, and video children",
        )
    for artifact_type, child in payload.children.items():
        if child.artifact_type != artifact_type:
            raise HTTPException(
                422,
                detail=f"paper bundle child {artifact_type} has the wrong artifact_type",
            )
        if not child.input_slots:
            raise HTTPException(
                422,
                detail=f"paper bundle child {artifact_type} requires an input slot",
            )
    owner_id = _run_owner_id(request)
    idempotency_key = request.headers.get("idempotency-key", "").strip()
    if not idempotency_key:
        raise HTTPException(400, detail="Idempotency-Key header is required")
    request_digest = _paper_bundle_request_digest(owner_id, payload)
    runtime = _web_run_runtime()

    async def reserve_child(
        artifact_type: str,
        parent_job_id: str,
        run_id: str,
    ) -> PaperBundleChildDescriptor:
        async with _web_run_start_guard(run_id):
            child_payload = payload.children[artifact_type]
            a_type, settings, run_payload, state = _prepare_reservation(
                request,
                child_payload,
            )
            if a_type != artifact_type:
                raise HTTPException(
                    422,
                    detail=f"paper bundle child {artifact_type} resolved as {a_type}",
                )
            state.conversation_id = _event_conversation_id(
                child_payload.conversation_id,
                run_id,
            )
            state.demo_user_id = "" if owner_id == "local" else owner_id
            slots = tuple(
                InputSlot(
                    name=slot.name,
                    expected_sha256=slot.sha256,
                    expected_size=slot.size,
                )
                for slot in child_payload.input_slots
            )
            child_digest = _reservation_digest(owner_id, child_payload)
            reservation = await runtime.services.reserve(
                run_id=run_id,
                artifact_type=artifact_type,
                idempotency_key=f"paper-bundle:{parent_job_id}:{artifact_type}",
                request_digest=child_digest,
                settings=settings,
                payload=run_payload,
                input_slots=slots,
                parent_job_id=parent_job_id,
            )
            state.reservation_token = reservation.upload_token
            async with _RUNS_LOCK:
                _RUNS[run_id] = state
            if _RUN_ACCESS_CONTROL:
                _demo_register_run(run_id, owner_id)
            _append_reservation_history(
                settings=settings,
                state=state,
                payload=child_payload,
                run_id=run_id,
            )
            return PaperBundleChildDescriptor(
                run_id=run_id,
                artifact_type=artifact_type,
                conversation_id=state.conversation_id,
                input_slots=tuple(
                    PaperBundleInputSlot(
                        name=slot.name,
                        expected_sha256=slot.expected_sha256,
                        expected_size=slot.expected_size,
                    )
                    for slot in reservation.input_slots
                ),
                upload_token=reservation.upload_token,
                request_digest=child_digest,
                expires_at=reservation.expires_at,
            )

    try:
        operation_id = f"paper-bundle-create:{payload.job_id or request_digest[:12]}"
        async with _web_run_operation_guard(operation_id):
            creation = await _paper_bundle_store().create_with_factory(
                owner_id=owner_id,
                conversation_id=payload.conversation_id,
                source_name=payload.source_name,
                prompt_version=payload.prompt_version,
                idempotency_key=idempotency_key,
                request_digest=request_digest,
                child_reservation_factory=reserve_child,
                cleanup_child=_cleanup_paper_bundle_child,
                job_id=payload.job_id,
            )
            return creation.to_payload()
    except RunNotReady as exc:
        raise _web_run_service_error(exc) from exc
    except HTTPException:
        raise
    except PaperBundleError as exc:
        raise _paper_bundle_error(exc) from exc


@app.get("/api/paper-bundles")
async def list_paper_bundles(request: Request) -> list[dict[str, Any]]:
    owner_id = _run_owner_id(request)
    try:
        records = await asyncio.to_thread(
            _paper_bundle_store().list_owned,
            owner_id,
            child_status_provider=_paper_bundle_child_snapshot,
        )
    except PaperBundleError as exc:
        raise _paper_bundle_error(exc) from exc
    return [record.to_payload() for record in records]


@app.get("/api/paper-bundles/{job_id}")
async def get_paper_bundle(job_id: str, request: Request) -> dict[str, Any]:
    owner_id = _run_owner_id(request)
    try:
        record = await asyncio.to_thread(
            _paper_bundle_store().read_owned,
            job_id,
            owner_id,
            child_status_provider=_paper_bundle_child_snapshot,
        )
    except PaperBundleError as exc:
        raise _paper_bundle_error(exc) from exc
    return record.to_payload()


@app.post("/api/paper-bundles/{job_id}/cancel")
async def cancel_paper_bundle(job_id: str, request: Request) -> JSONResponse:
    owner_id = _run_owner_id(request)
    store = _paper_bundle_store()
    try:
        record = await asyncio.to_thread(store.request_cancel, job_id, owner_id)
    except PaperBundleNotFound:
        try:
            pending = await store.cancel_pending_creation(
                job_id,
                owner_id,
                cleanup_child=_cleanup_paper_bundle_child,
            )
            deadline = time.monotonic() + 5.0
            while pending == "pending" and time.monotonic() < deadline:
                await asyncio.sleep(0.02)
                pending = await store.cancel_pending_creation(
                    job_id,
                    owner_id,
                    cleanup_child=_cleanup_paper_bundle_child,
                )
        except PaperBundleError as exc:
            raise _paper_bundle_error(exc) from exc
        if pending == "not_found":
            raise HTTPException(404, detail="paper bundle not found")
        if pending == "published":
            try:
                record = await asyncio.to_thread(
                    store.request_cancel,
                    job_id,
                    owner_id,
                )
            except PaperBundleError as exc:
                raise _paper_bundle_error(exc) from exc
        else:
            response = {
                "job_id": job_id,
                "owner_id": owner_id,
                "state": "cancelled" if pending == "cancelled" else "cancelling",
                "confirmed": pending == "cancelled",
                "status": (
                    "cancelled" if pending == "cancelled" else "cancellation_pending"
                ),
                "pending_creation": True,
                "factory_quiesced": pending == "cancelled",
                "children": {},
            }
            return JSONResponse(
                status_code=200 if response["confirmed"] else 202,
                content=response,
            )
    except PaperBundleError as exc:
        raise _paper_bundle_error(exc) from exc

    if record.terminal:
        payload = record.to_payload()
        payload["confirmed"] = True
        payload["status"] = (
            "already_cancelled" if record.state == "cancelled" else "already_terminal"
        )
        return JSONResponse(status_code=200, content=payload)

    child_results = await asyncio.gather(
        *(
            _cancel_controlled_run(
                child.run_id,
                f"paper_bundle_cancel:{job_id}",
            )
            for child in record.children.values()
        ),
        return_exceptions=True,
    )
    try:
        reconciled = await asyncio.to_thread(
            store.reconcile,
            job_id,
            owner_id,
            _paper_bundle_child_snapshot,
        )
    except PaperBundleError as exc:
        raise _paper_bundle_error(exc) from exc
    confirmed = reconciled.state == "cancelled" and all(
        not isinstance(result, BaseException) and result.confirmed
        for result in child_results
    )
    payload = reconciled.to_payload()
    payload["confirmed"] = confirmed
    payload["status"] = "cancelled" if confirmed else "cancellation_pending"
    return JSONResponse(
        status_code=200 if confirmed else 202,
        content=payload,
    )


@app.post("/api/runs/reserve", response_model=RunReserveResponse)
async def reserve_run(
    payload: RunReserveRequest,
    request: Request,
) -> RunReserveResponse:
    owner_id = _run_owner_id(request)
    idempotency_key = _reservation_idempotency_key(
        owner_id,
        request.headers.get("idempotency-key", ""),
    )
    a_type, settings, run_payload, state = _prepare_reservation(request, payload)
    run_id = new_run_id()
    state.conversation_id = _event_conversation_id(payload.conversation_id, run_id)
    slots = tuple(
        InputSlot(
            name=slot.name,
            expected_sha256=slot.sha256,
            expected_size=slot.size,
        )
        for slot in payload.input_slots
    )
    request_digest = _reservation_digest(owner_id, payload)
    try:
        async with _web_run_start_guard(run_id):
            try:
                reservation = await _web_run_runtime().services.reserve(
                    run_id=run_id,
                    artifact_type=a_type,
                    idempotency_key=idempotency_key,
                    request_digest=request_digest,
                    settings=settings,
                    payload=run_payload,
                    input_slots=slots,
                )
            except Exception as exc:
                raise _web_run_service_error(exc) from exc
            actual_run_id = reservation.run_id
            if _RUN_ACCESS_CONTROL:
                if reservation.reused:
                    _require_run_owner_before_lookup(actual_run_id, request)
                else:
                    _demo_register_run(actual_run_id, owner_id)
            if not reservation.reused:
                state.conversation_id = _event_conversation_id(
                    payload.conversation_id,
                    actual_run_id,
                )
                state.demo_user_id = "" if owner_id == "local" else owner_id
                state.reservation_token = reservation.upload_token
                async with _RUNS_LOCK:
                    _RUNS[actual_run_id] = state
                _append_reservation_history(
                    settings=settings,
                    state=state,
                    payload=payload,
                    run_id=actual_run_id,
                )
            return RunReserveResponse(
                run_id=actual_run_id,
                upload_token=reservation.upload_token,
                input_slots=payload.input_slots,
                request_digest=request_digest,
                run_state=reservation.state,
                expires_at=reservation.expires_at,
                reused=reservation.reused,
            )
    except RunNotReady as exc:
        raise _web_run_service_error(exc) from exc


@app.put("/api/runs/{run_id}/inputs/{slot_name}", response_model=RunUploadResponse)
async def upload_run_input(
    run_id: str,
    slot_name: str,
    request: Request,
) -> RunUploadResponse:
    _require_run_owner_before_lookup(run_id, request)
    owner_id = _run_owner_id(request)
    try:
        _assert_paper_bundle_child_work_allowed(run_id, owner_id)
    except PaperBundleError as exc:
        raise _paper_bundle_error(exc) from exc
    token = request.headers.get("x-autodesign-upload-token", "")
    try:
        result = await _web_run_runtime().services.upload(
            run_id,
            token,
            slot_name,
            request.stream(),
        )
    except Exception as exc:
        raise _web_run_service_error(exc) from exc
    if not result.idempotent:
        async with _RUNS_LOCK:
            state = _RUNS.get(run_id)
        if state is not None:
            event_data = attachment_event_data(result.path)
            if state.input_slot_roles.get(slot_name) == "reference_poster":
                event_data["role"] = "style_reference"
            _append_event(
                _settings_or_boot(),
                state.conversation_id,
                "attachment.added",
                run_id=run_id,
                data=event_data,
            )
    return RunUploadResponse(
        run_id=result.run_id,
        slot=result.slot,
        sha256=result.sha256,
        size=result.size,
        run_state=result.state,
        idempotent=result.idempotent,
    )


@app.post("/api/runs/{run_id}/start", response_model=GenerateAck)
async def start_run(run_id: str, request: Request) -> GenerateAck:
    _require_run_owner_before_lookup(run_id, request)
    owner_id = _run_owner_id(request)
    token = request.headers.get("x-autodesign-upload-token", "")
    async with _RUNS_LOCK:
        state = _RUNS.get(run_id)
    if state is None:
        raise HTTPException(409, detail="run start context expired; reserve a new run")
    derived_descriptor = _read_json_file(RUNS_DIR / run_id / "derived_job.json")
    if isinstance(derived_descriptor, dict):
        job_kind = str(derived_descriptor.get("job_kind") or "")
        if job_kind not in _DERIVED_JOB_KINDS:
            raise HTTPException(409, detail="derived run descriptor is invalid")
        try:
            await _start_reserved_derived_job(
                run_id=run_id,
                token=token,
                state=state,
                descriptor=derived_descriptor,
            )
        except Exception as exc:
            raise _web_run_service_error(exc) from exc
        return GenerateAck(
            run_id=run_id,
            progress_mode={
                "editable_video_render": "video_render",
                "attempt_fork": "attempt_fork",
                "candidate_publish": "attempt_publish",
                "poster_code_edit": "poster_code_edit",
                "pptx_export": "artifact_export",
                "video_export_retry": "video_export",
            }[job_kind],
            placeholder_message=Message(
                id=f"msg_{run_id}",
                role="assistant",
                text="",
                ts=int(time.time() * 1000),
                run_id=run_id,
                status="streaming",
            ),
        )
    if _DEMO_MODE:
        if not state.reservation_token or not secrets.compare_digest(
            token,
            state.reservation_token,
        ):
            raise HTTPException(404, detail="run reservation not found")
        await _admit_demo_run(request, _DemoQueuedRun(
            run_id=run_id,
            brief=state.brief,
            attach_paths=state.attach_paths,
            template=state.template,
            a_type=state.artifact_type,
            baseline_artifact_json=state.baseline_artifact_json,
            state=state,
            settings=None,
        ))
    else:
        try:
            await _start_run_with_paper_bundle_barrier(
                run_id,
                token,
                owner_id,
            )
        except PaperBundleError as exc:
            raise _paper_bundle_error(exc) from exc
        except Exception as exc:
            raise _web_run_service_error(exc) from exc
        async with _RUNS_LOCK:
            if state.task is None:
                state.task = asyncio.create_task(
                    _monitor_supervised_pipeline(run_id=run_id, state=state)
                )
    return GenerateAck(
        run_id=run_id,
        progress_mode="generate",
        reference_poster_handle=state.reference_poster_handle,
        placeholder_message=Message(
            id=f"msg_{run_id}",
            role="assistant",
            text="",
            ts=int(time.time() * 1000),
            run_id=run_id,
            status="streaming",
        ),
    )


@app.get("/api/runs/{run_id}/status", response_model=RunStatusResponse)
async def run_status(run_id: str, request: Request) -> RunStatusResponse:
    _require_run_owner_before_lookup(run_id, request)
    try:
        record = _web_run_runtime().control_store.read(run_id)
    except RunControlError as exc:
        raise HTTPException(404, detail=f"run not found: {run_id}") from exc
    if record.state == "cancelling":
        await _web_run_runtime().supervisor.recover(run_id)
        record = _web_run_runtime().control_store.read(run_id)
    await _reconcile_paper_bundle_for_run(
        run_id,
        owner_id=_run_owner_id(request),
    )
    return RunStatusResponse(
        run_id=run_id,
        run_state=record.state,
        revision=record.revision,
        publishable=record.publishable,
        cancellation_pending=record.cancellation_pending,
        worker_pid=record.worker_pid,
        terminal_event=record.terminal_event,
    )


# ---------- /api/generate ----------

# We accept multipart/form-data so PDF attachments can ride along. Form
# fields are individually parsed (no JSON-in-form) so the contract is
# explicit and shows up cleanly in OpenAPI/Swagger.
#
# v0.3 wire shape: returns immediately with a run_id + a streaming
# placeholder message. The actual run continues in a background asyncio
# task; the client opens SSE on /api/runs/{id}/events for phase progress
# and then GETs /api/runs/{id}/artifact once a `done` event lands.
@app.post("/api/generate", response_model=GenerateAck)
async def generate(
    request: Request,
    brief: str = Form(...),
    artifact_type: str | None = Form(None),
    palette_id: str | None = Form(None),
    baseline_artifact: str | None = Form(None),
    # Memory: the frontend sends the last N turns of *this* conversation
    # plus a compact summary of artifacts already produced. Both are
    # stitched into the brief as a preamble so the designer sees the
    # continuity without us bolting a new tool onto the runner. JSON
    # strings (form fields don't support nested objects directly).
    conversation_history: str | None = Form(None),
    prior_artifacts: str | None = Form(None),
    attachment_refs: str | None = Form(None),
    reference_poster_ref: str | None = Form(None),
    conversation_id: str | None = Form(None),
    template: str | None = Form(None),
    canvas_preset_id: str | None = Form(None),
    authoring_max_attempts: int | None = Form(None),
    files: list[UploadFile] = File(default=[]),
    reference_poster: UploadFile | None = File(default=None),
) -> GenerateAck:
    # Resolves credentials before we touch the filesystem so a no-key
    # request fails fast with a 412 the frontend can pop a Settings dialog
    # on, instead of leaking a half-staged upload dir.
    req_settings = _settings_for_request(request)
    access_user_id = _demo_user_id(request) if _RUN_ACCESS_CONTROL else ""
    if _DEMO_MODE:
        _demo_require_host_safety()
        _demo_cleanup_expired_runs()
        if artifact_type and artifact_type != "poster":
            raise HTTPException(
                403,
                detail={
                    "code": "demo_poster_only",
                    "message": "Demo mode only supports paper poster generation.",
                },
            )
        if len(files) > 1:
            raise HTTPException(
                400,
                detail={
                    "code": "demo_single_pdf_only",
                    "message": "Demo mode accepts one PDF per poster run.",
                },
            )
        if (reference_poster is not None and reference_poster.filename) or reference_poster_ref:
            raise HTTPException(
                400,
                detail={
                    "code": "demo_reference_poster_unsupported",
                    "message": "Demo mode does not accept a reference poster.",
                },
            )
        if _demo_queue_full():
            raise HTTPException(
                429,
                detail={
                    "code": "demo_queue_full",
                    "message": "Demo queue is full. Try again later.",
                },
            )
        a_type: ArtifactType = "poster"
        requested_template = _DEMO_FIXED_TEMPLATE
        requested_canvas_preset_id = _DEMO_FIXED_TEMPLATE
    else:
        a_type = _coerce_artifact_type(artifact_type, brief=brief)
        requested_template, requested_canvas_preset_id = _validated_web_canvas_selection(
            a_type,
            template,
            canvas_preset_id,
        )
    _validated_canvas_prompt(brief, a_type)
    resolved_authoring_max_attempts = _validated_authoring_max_attempts(
        authoring_max_attempts,
        a_type,
        req_settings,
    )
    req_settings = _settings_with_authoring_max_attempts(
        req_settings,
        a_type,
        resolved_authoring_max_attempts,
    )
    _require_artifact_runtime(a_type)
    normalized_palette_id = _validated_web_palette_id(a_type, palette_id)
    if artifact_type not in _ARTIFACT_TYPES and a_type != "poster":
        log(
            "web.artifact_type.inferred",
            brief=brief.strip()[:120],
            inferred=a_type,
        )

    reference_name: str | None = None
    if reference_poster is not None and reference_poster.filename:
        reference_name = _validated_web_reference_poster_name(reference_poster.filename)

    upload_dir = UPLOADS_DIR / uuid.uuid4().hex[:12]
    upload_dir.mkdir(parents=True, exist_ok=True)
    attach_paths: list[Path] = []
    for f in files:
        if not f.filename:
            continue
        if _DEMO_MODE and Path(f.filename).suffix.lower() != ".pdf":
            raise HTTPException(
                400,
                detail={
                    "code": "demo_pdf_required",
                    "message": "Demo mode accepts PDF paper uploads only.",
                },
            )
        out_path = upload_dir / Path(f.filename).name
        with out_path.open("wb") as w:
            shutil.copyfileobj(f.file, w)
        if _DEMO_MODE and out_path.stat().st_size > _DEMO_MAX_PDF_BYTES:
            out_path.unlink(missing_ok=True)
            raise HTTPException(
                413,
                detail={
                    "code": "demo_pdf_too_large",
                    "message": f"Demo PDF limit is {_DEMO_MAX_PDF_BYTES // (1024 * 1024)} MB.",
                },
            )
        attach_paths.append(out_path)
    reference_poster_path: Path | None = None
    reference_handle = _reference_poster_handle(reference_poster_ref)
    if reference_poster is not None and reference_name is not None:
        reference_handle = None
        reference_upload_dir = upload_dir / "reference_poster"
        reference_upload_dir.mkdir(parents=True, exist_ok=True)
        reference_poster_path = reference_upload_dir / reference_name
        with reference_poster_path.open("wb") as writer:
            shutil.copyfileobj(reference_poster.file, writer)
        try:
            _validate_web_reference_poster_file(reference_poster_path)
        except HTTPException:
            reference_poster_path.unlink(missing_ok=True)
            raise
    if reference_poster_path is None and reference_poster_ref:
        if reference_handle:
            reference_poster_path = _persisted_conversation_reference_poster_by_handle(
                conversation_id,
                reference_handle,
                owner_id=access_user_id,
            )
        elif _reference_poster_ref_requires_exact_lookup(reference_poster_ref):
            raise HTTPException(
                422,
                detail={
                    "code": "reference_poster_handle_required",
                    "message": "The exact saved reference is unavailable. Reattach it to continue.",
                },
            )
        else:
            reference_poster_path = await _latest_conversation_reference_poster(
                conversation_id,
                owner_id=access_user_id,
            )
        if reference_poster_path is None:
            raise HTTPException(
                422,
                detail={
                    "code": "reference_poster_not_found",
                    "message": "The saved reference poster is unavailable. Reattach it to continue.",
                },
            )
        _validate_web_reference_poster_file(reference_poster_path)
    reused_prior_attachments = False
    if not attach_paths:
        prior_attach_paths = await _latest_conversation_attach_paths(conversation_id)
        if prior_attach_paths:
            attach_paths = prior_attach_paths
            reused_prior_attachments = True
    if not attach_paths:
        ref_attach_paths = _find_uploads_by_attachment_refs(attachment_refs)
        if ref_attach_paths:
            attach_paths = ref_attach_paths
            reused_prior_attachments = True
            if conversation_id:
                _record_conversation_attach_paths(conversation_id, attach_paths)
    elif conversation_id:
        _record_conversation_attach_paths(conversation_id, attach_paths)

    if _DEMO_MODE and not any(p.suffix.lower() == ".pdf" for p in attach_paths):
        raise HTTPException(
            400,
            detail={
                "code": "demo_pdf_required",
                "message": "Attach a PDF paper to generate a demo poster.",
            },
        )
    if reference_poster_path is not None and a_type != "poster":
        raise HTTPException(400, detail="reference poster style is only supported for poster generation")
    if reference_poster_path is not None and not any(p.suffix.lower() == ".pdf" for p in attach_paths):
        raise HTTPException(400, detail="attach a paper PDF when using a reference poster style")
    if reference_poster_path is not None:
        reference_handle = _record_conversation_reference_poster(
            conversation_id,
            reference_poster_path,
            owner_id=access_user_id,
            reference_handle=reference_handle,
        )

    if a_type == "poster" and any(path.suffix.lower() == ".pdf" for path in attach_paths):
        from autodesign.util.paper_source_sanity import (
            PaperSourceInputError,
            PaperSourceSanityError,
            assert_valid_paper_source_pdf,
        )

        for paper_path in (path for path in attach_paths if path.suffix.lower() == ".pdf"):
            try:
                assert_valid_paper_source_pdf(paper_path)
            except (PaperSourceSanityError, PaperSourceInputError) as exc:
                raise HTTPException(
                    422,
                    detail={
                        "code": exc.issue_id,
                        "message": str(exc),
                        "repair_route": exc.repair_route,
                        "paper_source_sanity": exc.report,
                    },
                ) from exc

    # Pass 2 — paper-poster + Kimi designer can stall on bbox geometry
    # for long PDF inputs.
    # Auto-switch to Claude Opus 4.7 for the duration of THIS run when
    # the user hasn't explicitly overridden the designer via a header.
    # The override path (`X-Model-Designer`, or legacy `X-Model-Planner`)
    # wins — power users choosing
    # Kimi on purpose are respected.
    has_pdf_attachment = any(p.suffix.lower() == ".pdf" for p in attach_paths)
    effective_template = requested_template
    web_paper_poster_profile_applied = False
    if a_type == "poster" and has_pdf_attachment:
        req_settings = _web_paper_poster_settings(req_settings)
        web_paper_poster_profile_applied = True
        author_cmd = _paper_poster_author_cmd_resolution(req_settings)
        if not author_cmd["available"]:
            raise HTTPException(
                412,
                detail={
                    "code": "missing_external_author_command",
                    "message": author_cmd["message"],
                },
            )
        log(
            "web.paper_poster_profile.applied",
            template=effective_template,
            designer_author=req_settings.designer_author_mode,
            designer_author_harness=req_settings.designer_author_harness,
            designer_author_cmd_source=author_cmd["source"],
        )
    user_picked_designer = bool(
        (request.headers.get("x-model-designer", "") or "").strip()
        or (request.headers.get("x-model-planner", "") or "").strip()
    )
    if (
        a_type == "poster"
        and has_pdf_attachment
        and not user_picked_designer
        and _KIMI_PATTERN in req_settings.designer_model.lower()
    ):
        log("designer.auto_switched",
            from_model=req_settings.designer_model,
            to_model=_OPUS_FALLBACK,
            reason="paper_poster_kimi_stall")
        # Re-resolve settings under the lock with DESIGNER_MODEL pinned.
        # MUST save+restore inside the lock — `_settings_for_request`
        # already returned and restored its own overrides; if we
        # mutated env here without restoring, the override would leak
        # across subsequent requests.
        with _SETTINGS_LOCK:
            saved = os.environ.get("DESIGNER_MODEL")
            try:
                os.environ["DESIGNER_MODEL"] = _OPUS_FALLBACK
                req_settings = load_settings()
                if web_paper_poster_profile_applied:
                    req_settings = _web_paper_poster_settings(req_settings)
                req_settings = _settings_with_authoring_max_attempts(
                    req_settings,
                    a_type,
                    resolved_authoring_max_attempts,
                )
            except RuntimeError:
                log("designer.auto_switch.failed",
                    reason="load_settings raised — keeping Kimi designer")
            finally:
                if saved is None:
                    os.environ.pop("DESIGNER_MODEL", None)
                else:
                    os.environ["DESIGNER_MODEL"] = saved

    # Build the layered prologue: type → conversation context → user brief.
    # Type prologue tells planner WHAT to make; conversation prologue tells
    # it what came before in this thread (so "now make slides" knows what
    # the slides should be ABOUT). Both are pure brief-augmentation — no
    # new tool, no runner change.
    effective_brief = _apply_type_prologue(brief, a_type)
    effective_brief = _apply_conversation_prologue(
        effective_brief,
        history_json=conversation_history,
        artifacts_json=prior_artifacts,
    )

    # Generate run_id up front so the client can subscribe to events before
    # the supervisor starts the pre-addressed worker (no race window).
    run_id = new_run_id()
    event_conversation_id = _event_conversation_id(conversation_id, run_id)
    if _RUN_ACCESS_CONTROL:
        _demo_register_run(run_id, access_user_id)

    log("web.run.start", run_id=run_id, artifact_type=a_type,
        brief_chars=len(effective_brief), attachments=len(attach_paths),
        reused_attachments=reused_prior_attachments,
        authoring_max_attempts=resolved_authoring_max_attempts,
        palette_id=normalized_palette_id or "",
        template=effective_template or "",
        canvas_preset_id=requested_canvas_preset_id or "",
        history_turns=_count_history_turns(conversation_history),
        prior_artifact_count=_count_prior_artifacts(prior_artifacts))
    _has_baseline = baseline_artifact is not None
    _submitted_data: dict[str, Any] = {
        "brief": brief.strip(),
        "artifact_type": a_type,
        "palette_id": normalized_palette_id,
        "template": effective_template,
        "canvas_preset_id": requested_canvas_preset_id,
        "has_baseline": _has_baseline,
        "authoring_max_attempts": resolved_authoring_max_attempts,
    }
    if reused_prior_attachments:
        _submitted_data["reused_prior_attachments"] = True
        _submitted_data["attachment_count"] = len(attach_paths)
    if reference_poster_path is not None:
        _submitted_data["has_reference_poster"] = True
        _submitted_data["reference_poster_name"] = reference_poster_path.name
    if _has_baseline:
        # Weak signal: classify whether this follow-up extends, corrects, or
        # rejects the prior artifact. Used by the memory extractor to decide
        # whether the prior version was implicitly accepted by the user before
        # they requested a change.
        _submitted_data["follow_up_sentiment"] = infer_follow_up_sentiment(brief)
    _append_event(
        req_settings, event_conversation_id, "message.user_submitted",
        run_id=run_id,
        data=_submitted_data,
    )
    for p in attach_paths:
        _append_event(
            req_settings, event_conversation_id, "attachment.added",
            run_id=run_id,
            data=attachment_event_data(p),
        )
    if reference_poster_path is not None:
        reference_event_data = {
            **attachment_event_data(reference_poster_path),
            "role": "style_reference",
        }
        if reference_handle:
            reference_event_data["reference_handle"] = reference_handle
        _append_event(
            req_settings,
            event_conversation_id,
            "attachment.added",
            run_id=run_id,
            data=reference_event_data,
        )

    state = _RunState(
        artifact_type=a_type,
        designer_model=req_settings.designer_model,
        has_pdf=any(p.suffix.lower() == ".pdf" for p in attach_paths),
        brief=effective_brief,
        attach_paths=attach_paths,
        reference_poster_path=reference_poster_path,
        baseline_artifact_json=baseline_artifact,
        conversation_id=event_conversation_id,
        template=effective_template,
        canvas_preset_id=requested_canvas_preset_id,
        palette_id=normalized_palette_id,
        authoring_max_attempts=resolved_authoring_max_attempts,
    )
    state.demo_user_id = access_user_id
    async with _RUNS_LOCK:
        _RUNS[run_id] = state

    reserve_only = _request_reserve_only(request)
    if _DEMO_MODE:
        await _reserve_legacy_pipeline_worker(
            run_id=run_id,
            brief=effective_brief,
            attach_paths=attach_paths,
            reference_poster_path=reference_poster_path,
            template=effective_template,
            state=state,
            settings=req_settings,
            resume_run=None,
        )
        if not reserve_only:
            await _admit_demo_run(request, _DemoQueuedRun(
                run_id=run_id,
                brief=effective_brief,
                attach_paths=attach_paths,
                template=effective_template,
                a_type=a_type,
                baseline_artifact_json=baseline_artifact,
                state=state,
                settings=req_settings,
            ))
    else:
        if reserve_only:
            await _reserve_legacy_pipeline_worker(
                run_id=run_id,
                brief=effective_brief,
                attach_paths=attach_paths,
                reference_poster_path=reference_poster_path,
                template=effective_template,
                state=state,
                settings=req_settings,
                resume_run=None,
            )
        else:
            await _start_legacy_pipeline_worker(
                run_id=run_id,
                brief=effective_brief,
                attach_paths=attach_paths,
                reference_poster_path=reference_poster_path,
                template=effective_template,
                state=state,
                settings=req_settings,
                resume_run=None,
            )
            state.task = asyncio.create_task(
                _monitor_supervised_pipeline(run_id=run_id, state=state),
            )

    return GenerateAck(
        run_id=run_id,
        progress_mode="generate",
        reference_poster_handle=reference_handle,
        start_token=state.reservation_token if reserve_only else None,
        placeholder_message=Message(
            id=f"msg_{run_id}",
            role="assistant",
            text="",
            ts=int(time.time() * 1000),
            run_id=run_id,
            status="streaming",
        ),
    )


# ---------- /api/code-edit/poster ----------


@app.post("/api/code-edit/poster", response_model=GenerateAck)
async def poster_code_edit(
    req: PosterCodeEditRequest,
    request: Request,
) -> GenerateAck:
    if _DEMO_MODE:
        raise HTTPException(
            403,
            detail={
                "code": "demo_code_edit_disabled",
                "message": "External poster code editing is disabled in demo mode.",
            },
        )
    artifact = req.artifact if isinstance(req.artifact, dict) else {}
    if artifact.get("artifact_type") != "poster":
        raise HTTPException(400, detail="artifact_type must be poster")
    if artifact.get("native_format") not in (None, "html"):
        raise HTTPException(400, detail="poster code edit requires an HTML artifact")
    source_run_id = (
        req.source_run_id
        or _artifact_declared_run_id(artifact)
        or ""
    ).strip()
    if not source_run_id:
        raise HTTPException(400, detail="missing source_run_id or poster artifact_id")
    _demo_require_run_owner(
        source_run_id,
        request,
        detail=f"source run not found: {source_run_id}",
    )
    _assert_controlled_run_source_usable(source_run_id, mode="artifact")
    try:
        _require_derived_source_ready(source_run_id, "poster_code_edit")
    except RunNotReady as exc:
        raise HTTPException(409, detail=str(exc)) from exc

    normalized_palette_id = _validated_web_palette_id("poster", req.palette_id)
    try:
        required_color_system = require_academic_color_system(normalized_palette_id or "")
    except AcademicPaletteCatalogError as exc:
        raise HTTPException(
            503,
            detail={
                "code": "palette_catalog_unavailable",
                "message": str(exc),
            },
        ) from exc
    settings = _settings_for_code_editor_request(request)
    command = _code_editor_cmd_resolution(settings)
    if not command["available"]:
        raise HTTPException(
            412,
            detail={
                "code": "missing_code_editor_command",
                "message": command["message"],
            },
        )
    settings = replace(settings, code_editor_cmd=str(command["cmd"]))

    instruction = req.instruction.strip()
    if not instruction:
        raise HTTPException(400, detail="instruction is required")
    selection_context = req.selection_context
    if selection_context is not None:
        try:
            selection_context_size = len(json.dumps(selection_context, ensure_ascii=False))
        except TypeError as exc:
            raise HTTPException(400, detail="selection_context must be JSON serializable") from exc
        if selection_context_size > 60_000:
            raise HTTPException(413, detail="selection_context is too large")

    source_poster_path = _source_poster_html_path(source_run_id)
    if source_poster_path is None:
        raise HTTPException(404, detail=f"source poster HTML not found for run {source_run_id}")

    run_id = new_run_id()
    access_user_id = _demo_register_derived_run_access(
        run_id,
        request,
        parent_run_id=source_run_id,
        missing_detail=f"source run not found: {source_run_id}",
    )
    event_conversation_id = _event_conversation_id(req.conversation_id, run_id)
    baseline_json = json.dumps(
        {"artifact_id": artifact.get("artifact_id") or f"art_{source_run_id}"},
        ensure_ascii=False,
    )
    state = _RunState(
        artifact_type="poster",
        designer_model=str(getattr(settings, "code_editor_model", None) or getattr(settings, "code_editor_harness", "codex")),
        has_pdf=False,
        brief=instruction,
        baseline_artifact_json=baseline_json,
        conversation_id=event_conversation_id,
        palette_id=normalized_palette_id,
    )
    state.demo_user_id = access_user_id
    async with _RUNS_LOCK:
        _RUNS[run_id] = state

    worker_request = PosterCodeEditWorkerRequest(
        job_kind="poster_code_edit",
        run_id=run_id,
        parent_run_id=source_run_id,
        source_poster=str(source_poster_path),
        artifact=artifact,
        instruction=instruction,
        conversation_history=tuple(req.conversation_history),
        selection_context=selection_context,
        palette_id=normalized_palette_id,
        required_color_system=required_color_system,
        conversation_id=event_conversation_id,
        baseline_artifact_json=baseline_json,
        settings=settings,
    )
    reserve_only = _request_reserve_only(request)
    try:
        start_token = await _start_supervised_derived_job(
            request=worker_request,
            state=state,
            descriptor={
                "job_kind": "poster_code_edit",
                "source_artifact_id": str(
                    artifact.get("artifact_id") or f"art_{source_run_id}"
                ),
                "artifact_name": str(artifact.get("name") or "Poster"),
                "source_relative_path": _run_relative_path(source_poster_path),
            },
            start_immediately=not reserve_only,
        )
    except Exception as exc:
        raise _web_run_service_error(exc) from exc

    _append_event(
        settings,
        event_conversation_id,
        "message.user_submitted",
        run_id=run_id,
        data={
            "brief": instruction,
            "artifact_type": "poster",
            "palette_id": normalized_palette_id,
            "required_color_system": required_color_system,
            "has_baseline": True,
            "revision_mode": "external_code_editor",
            "revision_scope": "area" if selection_context else "poster",
            "source_run_id": source_run_id,
            "has_selection_context": bool(selection_context),
            "follow_up_sentiment": infer_follow_up_sentiment(instruction),
        },
    )
    log(
        "code_editor.requested",
        run_id=run_id,
        source_run_id=source_run_id,
        palette_id=normalized_palette_id,
        command_source=command["source"],
    )
    return GenerateAck(
        run_id=run_id,
        progress_mode="poster_code_edit",
        start_token=start_token if reserve_only else None,
        placeholder_message=Message(
            id=f"msg_{run_id}",
            role="assistant",
            text="",
            ts=int(time.time() * 1000),
            run_id=run_id,
            status="streaming",
        ),
    )


@app.post("/api/artifacts/export", response_model=ArtifactExportResponse)
async def export_artifact(
    req: ArtifactExportRequest,
    request: Request,
) -> ArtifactExportResponse:
    settings = _settings_for_code_editor_request(request) if req.format == "pptx" else None
    access_user_id = _demo_user_id(request) if _RUN_ACCESS_CONTROL else None
    return await asyncio.to_thread(_export_artifact_sync, req, settings, access_user_id)


def _export_artifact_sync(
    req: ArtifactExportRequest,
    settings: Settings | None = None,
    access_user_id: str | None = None,
) -> ArtifactExportResponse:
    artifact = req.artifact if isinstance(req.artifact, dict) else {}
    referenced_run_ids = _artifact_referenced_run_ids(artifact)
    if access_user_id is not None:
        for referenced_run_id in referenced_run_ids:
            if not _demo_user_owns_run(referenced_run_id, access_user_id):
                raise HTTPException(404, detail="HTML source not found for this artifact")
    if len(referenced_run_ids) > 1:
        raise HTTPException(409, detail="artifact references multiple runs")
    run_id = referenced_run_ids[0] if referenced_run_ids else None
    if run_id:
        _assert_controlled_run_source_usable(run_id, mode="artifact")
    source = _html_export_source_path(artifact)
    if source is None:
        raise HTTPException(404, detail="HTML source not found for this artifact")
    source_run_id = _run_id_for_path(source)
    if not source_run_id:
        raise HTTPException(400, detail="could not resolve artifact run_id")
    if run_id and run_id != source_run_id:
        raise HTTPException(409, detail="artifact source does not match its declared run")
    run_id = source_run_id
    if access_user_id is not None and not _demo_user_owns_run(run_id, access_user_id):
        raise HTTPException(404, detail="HTML source not found for this artifact")
    _assert_controlled_run_source_usable(run_id, mode="artifact")

    if req.format == "original_html":
        return ArtifactExportResponse(
            url=_run_file_url(run_id, _run_relative_path(source)),
            filename=source.name,
            format="original_html",
            mime_type="text/html",
        )

    export_dir = source.parent / "exports"
    export_dir.mkdir(parents=True, exist_ok=True)
    stem = _safe_export_stem(str(artifact.get("name") or source.stem or "artifact"))
    if req.format == "standalone_html":
        out = export_dir / f"{stem}-standalone.html"
        _write_standalone_html(source, out)
        return ArtifactExportResponse(
            url=_run_file_url(run_id, _run_relative_path(out)),
            filename=out.name,
            format="standalone_html",
            mime_type="text/html",
        )

    canvas = artifact.get("canvas") if isinstance(artifact.get("canvas"), dict) else {}
    width, height = _html_export_canvas_size(source, artifact, canvas)
    if req.format == "pdf":
        out = export_dir / f"{stem}.pdf"
        fallback = source.parent / "preview.png"
        result = export_html_pdf(
            source,
            out,
            viewport_width=width,
            viewport_height=height,
            page_width=f"{width}px",
            page_height=f"{height}px",
            fallback_pngs=[fallback] if fallback.exists() else None,
            enforce_single_page=str(artifact.get("artifact_type") or "") == "poster",
            canvas_selector="body > .paper-poster, body > .canvas, body > [data-autodesign-artifact-root], body > svg",
            canvas_width_px=width,
            canvas_height_px=height,
            timeout_ms=30_000,
        )
        if not out.exists():
            raise HTTPException(500, detail="PDF export failed: " + "; ".join(result.warnings))
        return ArtifactExportResponse(
            url=_run_file_url(run_id, _run_relative_path(out)),
            filename=out.name,
            format="pdf",
            mime_type="application/pdf",
        )

    if req.format == "pptx":
        out = export_dir / f"{stem}.pptx"
        _write_agent_pptx_export(
            source,
            out,
            width=width,
            height=height,
            artifact=artifact,
            settings=settings,
            run_id=run_id,
        )
        return ArtifactExportResponse(
            url=_run_file_url(run_id, _run_relative_path(out)),
            filename=out.name,
            format="pptx",
            mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

    raise HTTPException(400, detail=f"unsupported export format: {req.format}")


@app.post("/api/artifacts/export/pptx-run", response_model=GenerateAck)
async def export_artifact_pptx_run(
    req: ArtifactPptxExportRequest,
    request: Request,
) -> GenerateAck:
    if _DEMO_MODE:
        raise HTTPException(
            403,
            detail={
                "code": "demo_export_disabled",
                "message": "Agent-assisted PowerPoint export is disabled in demo mode.",
            },
        )
    artifact = req.artifact if isinstance(req.artifact, dict) else {}
    referenced_run_ids = _artifact_referenced_run_ids(artifact)
    access_user_id = ""
    for referenced_run_id in referenced_run_ids:
        access_user_id = _demo_require_run_owner(
            referenced_run_id,
            request,
            detail="HTML source not found for this artifact",
        )
    if len(referenced_run_ids) > 1:
        raise HTTPException(409, detail="artifact references multiple runs")
    source_run_id = referenced_run_ids[0] if referenced_run_ids else None
    if source_run_id:
        _assert_controlled_run_source_usable(source_run_id, mode="artifact")
    source = _html_export_source_path(artifact)
    if source is None:
        raise HTTPException(404, detail="HTML source not found for this artifact")
    actual_source_run_id = _run_id_for_path(source)
    if not actual_source_run_id:
        raise HTTPException(400, detail="could not resolve artifact run_id")
    if source_run_id and source_run_id != actual_source_run_id:
        raise HTTPException(409, detail="artifact source does not match its declared run")
    source_run_id = actual_source_run_id
    access_user_id = _demo_require_run_owner(
        source_run_id,
        request,
        detail="HTML source not found for this artifact",
    )
    _assert_controlled_run_source_usable(source_run_id, mode="artifact")
    try:
        _require_derived_source_ready(source_run_id, "pptx_export")
    except RunNotReady as exc:
        raise HTTPException(409, detail=str(exc)) from exc
    export_key = _pptx_export_key(access_user_id, source_run_id, source)
    async with _RUNS_LOCK:
        active_ack = _active_pptx_export_ack_locked(export_key)
        if active_ack is not None:
            return active_ack

    settings = _settings_for_code_editor_request(request)
    command = _code_editor_cmd_resolution(settings)
    if not command["available"]:
        raise HTTPException(
            412,
            detail={
                "code": "missing_code_editor_command",
                "message": command["message"],
            },
        )
    settings = replace(settings, code_editor_cmd=str(command["cmd"]))

    artifact_name = str(artifact.get("name") or source.stem or "artifact").strip()
    async with _PPTX_EXPORT_LOCK:
        async with _RUNS_LOCK:
            active_ack = _active_pptx_export_ack_locked(export_key)
            if active_ack is not None:
                return active_ack

            run_id = new_run_id()
            access_user_id = _demo_register_derived_run_access(
                run_id,
                request,
                parent_run_id=source_run_id,
                missing_detail="HTML source not found for this artifact",
            )
            event_conversation_id = _event_conversation_id(req.conversation_id, run_id)
            source_artifact_id = str(
                artifact.get("artifact_id") or f"art_{source_run_id}"
            )
            state = _RunState(
                artifact_type="poster",
                designer_model=str(
                    getattr(settings, "code_editor_model", None)
                    or getattr(settings, "code_editor_harness", "codex")
                ),
                has_pdf=False,
                brief=f"Export this design as an editable PPTX: {artifact_name}",
                baseline_artifact_json=json.dumps(
                    {"artifact_id": source_artifact_id},
                    ensure_ascii=False,
                ),
                conversation_id=event_conversation_id,
            )
            state.demo_user_id = access_user_id
            _RUNS[run_id] = state
            _PPTX_EXPORT_RUNS[export_key] = run_id

        worker_request = PptxExportWorkerRequest(
            job_kind="pptx_export",
            run_id=run_id,
            parent_run_id=source_run_id,
            source_html=str(source),
            artifact=artifact,
            artifact_name=artifact_name,
            conversation_id=event_conversation_id,
            settings=settings,
        )
        reserve_only = _request_reserve_only(request)
        try:
            start_token = await _start_supervised_derived_job(
                request=worker_request,
                state=state,
                descriptor={
                    "job_kind": "pptx_export",
                    "source_artifact_id": source_artifact_id,
                    "artifact_name": artifact_name,
                    "source_relative_path": _run_relative_path(source),
                },
                start_immediately=not reserve_only,
            )
        except Exception as exc:
            await _clear_pptx_export_registration(export_key, run_id)
            raise _web_run_service_error(exc) from exc
        log(
            "artifact_export.requested",
            run_id=run_id,
            source_run_id=source_run_id,
            format="pptx",
            command_source=command["source"],
        )
        _append_event(
            settings,
            event_conversation_id,
            "artifact.export_requested",
            run_id=run_id,
            data={
                "format": "pptx",
                "artifact_id": artifact.get("artifact_id"),
                "artifact_name": artifact_name,
                "source_run_id": source_run_id,
            },
        )
        return _pptx_export_ack(
            run_id,
            state,
            start_token=start_token if reserve_only else None,
        )


def _pptx_export_key(
    access_user_id: str,
    source_run_id: str,
    source: Path,
) -> _PptxExportKey:
    return access_user_id, source_run_id, str(source.resolve())


def _pptx_export_ack(
    run_id: str,
    state: _RunState,
    *,
    start_token: str | None = None,
) -> GenerateAck:
    return GenerateAck(
        run_id=run_id,
        progress_mode="artifact_export",
        start_token=start_token,
        placeholder_message=Message(
            id=f"msg_{run_id}",
            role="assistant",
            text="",
            ts=int(state.created_at * 1000),
            run_id=run_id,
            status="streaming",
        ),
    )


def _active_pptx_export_ack_locked(export_key: _PptxExportKey) -> GenerateAck | None:
    run_id = _PPTX_EXPORT_RUNS.get(export_key)
    if run_id is None:
        return None
    state = _RUNS.get(run_id)
    if state is None or (
        not state.reservation_token
        and (state.task is None or state.task.done())
    ):
        if _PPTX_EXPORT_RUNS.get(export_key) == run_id:
            _PPTX_EXPORT_RUNS.pop(export_key, None)
        return None
    return _pptx_export_ack(
        run_id,
        state,
        start_token=state.reservation_token or None,
    )


def _artifact_declared_run_id(artifact: dict[str, Any]) -> str | None:
    run_id = _run_id_from_artifact_id(str(artifact.get("artifact_id") or ""))
    if run_id:
        return run_id
    for key in ("native_file_url", "view_file_url", "download_url"):
        run_id = _run_id_from_run_file_url(artifact.get(key))
        if run_id:
            return run_id
    return None


def _artifact_referenced_run_ids(artifact: dict[str, Any]) -> tuple[str, ...]:
    run_ids: list[str] = []
    artifact_run_id = _run_id_from_artifact_id(
        str(artifact.get("artifact_id") or "")
    )
    if artifact_run_id:
        run_ids.append(artifact_run_id)
    for key in ("native_file_url", "view_file_url", "download_url"):
        run_id = _run_id_from_run_file_url(artifact.get(key))
        if run_id and run_id not in run_ids:
            run_ids.append(run_id)
    return tuple(run_ids)


def _run_id_from_run_file_url(raw: Any) -> str | None:
    if not isinstance(raw, str) or not raw.strip():
        return None
    parsed = urlparse(raw.strip())
    path = unquote(parsed.path or raw.strip()).split("#", 1)[0].split("?", 1)[0]
    prefix = "/api/files/runs/"
    if not path.startswith(prefix):
        return None
    rel = path[len(prefix):].lstrip("/")
    parts = rel.split("/")
    if not rel or any(part in {"", ".", ".."} for part in parts):
        return None
    return parts[0]


def _html_export_source_path(artifact: dict[str, Any]) -> Path | None:
    for key in ("native_file_url", "view_file_url", "download_url"):
        path = _path_from_run_file_url(artifact.get(key))
        if path and path.suffix.lower() == ".html" and path.exists():
            return path
    run_id = _run_id_from_artifact_id(str(artifact.get("artifact_id") or ""))
    if not run_id:
        return None
    final = RUNS_DIR / run_id / "final"
    a_type = artifact.get("artifact_type")
    preferred = {
        "poster": "poster.html",
        "landing": "index.html",
        "deck": "deck.html",
    }.get(str(a_type or ""))
    candidates: list[Path] = []
    if preferred:
        candidates.append(final / preferred)
    candidates.extend(sorted(final.glob("*.html")))
    for path in candidates:
        if path.exists():
            return path
    return None


def _path_from_run_file_url(raw: Any) -> Path | None:
    if not isinstance(raw, str) or not raw.strip():
        return None
    parsed = urlparse(raw.strip())
    path = unquote(parsed.path or raw.strip()).split("#", 1)[0].split("?", 1)[0]
    prefix = "/api/files/runs/"
    if not path.startswith(prefix):
        return None
    rel = path[len(prefix):].lstrip("/")
    if not rel or any(part in {"", ".", ".."} for part in rel.split("/")):
        return None
    candidate = (RUNS_DIR / rel).resolve()
    try:
        candidate.relative_to(RUNS_DIR.resolve())
    except ValueError:
        return None
    return candidate


def _run_id_for_path(path: Path) -> str | None:
    try:
        rel = path.resolve().relative_to(RUNS_DIR.resolve())
    except ValueError:
        return None
    return rel.parts[0] if rel.parts else None


def _run_relative_path(path: Path) -> str:
    try:
        rel = path.resolve().relative_to((RUNS_DIR / (_run_id_for_path(path) or "")).resolve())
    except ValueError as exc:
        raise HTTPException(400, detail="export path is outside run directory") from exc
    return str(rel).replace(os.sep, "/")


def _safe_export_stem(raw: str) -> str:
    clean = re.sub(r"[^a-zA-Z0-9._-]+", "-", raw.strip()).strip("-._")
    return (clean or "artifact")[:80]


def _html_export_canvas_size(
    source: Path,
    artifact: dict[str, Any],
    canvas: dict[str, Any],
) -> tuple[int, int]:
    if artifact.get("artifact_type") == "poster" and _is_authored_paper_poster_html(source):
        return _authored_paper_poster_size(source)
    width = _positive_int(canvas.get("w")) or _num_from_html_attr(source, "data-w") or 1440
    height = _positive_int(canvas.get("h")) or _num_from_html_attr(source, "data-h") or 900
    return max(1, width), max(1, height)


def _positive_int(raw: Any) -> int | None:
    try:
        value = int(float(str(raw)))
    except (TypeError, ValueError):
        return None
    return value if value > 0 else None


def _num_from_html_attr(path: Path, attr: str) -> int | None:
    try:
        text = path.read_text(encoding="utf-8", errors="ignore")
    except OSError:
        return None
    match = re.search(rf"\b{re.escape(attr)}=[\"']([0-9.]+)[\"']", text)
    return _positive_int(match.group(1)) if match else None


def _write_standalone_html(source: Path, out: Path) -> None:
    text = source.read_text(encoding="utf-8", errors="ignore")
    soup = BeautifulSoup(text, "html.parser")
    for tag in soup.find_all(src=True):
        inlined = _html_ref_data_url(tag.get("src"), source.parent)
        if inlined:
            tag["src"] = inlined
    for tag in soup.find_all(srcset=True):
        tag["srcset"] = _inline_srcset(str(tag.get("srcset") or ""), source.parent)
    for tag in soup.find_all(style=True):
        tag["style"] = _inline_css_urls(str(tag.get("style") or ""), source.parent)
    for style in soup.find_all("style"):
        if style.string:
            style.string.replace_with(_inline_css_urls(style.string, source.parent))
    out.write_text(str(soup), encoding="utf-8")


def _inline_srcset(raw: str, base_dir: Path) -> str:
    parts: list[str] = []
    for item in raw.split(","):
        bits = item.strip().split()
        if not bits:
            continue
        inlined = _html_ref_data_url(bits[0], base_dir) or bits[0]
        parts.append(" ".join([inlined, *bits[1:]]))
    return ", ".join(parts)


def _inline_css_urls(raw: str, base_dir: Path) -> str:
    def repl(match: re.Match[str]) -> str:
        ref = match.group(1).strip().strip("'\"")
        inlined = _html_ref_data_url(ref, base_dir)
        return f"url({inlined})" if inlined else match.group(0)

    return re.sub(r"url\(([^)]+)\)", repl, raw)


def _html_ref_data_url(raw: Any, base_dir: Path) -> str | None:
    path = _resolve_html_ref_path(raw, base_dir)
    if not path or not path.exists() or not path.is_file():
        return None
    mime = mimetypes.guess_type(path.name)[0] or "application/octet-stream"
    try:
        data = base64.b64encode(path.read_bytes()).decode("ascii")
    except OSError:
        return None
    return f"data:{mime};base64,{data}"


def _resolve_html_ref_path(raw: Any, base_dir: Path) -> Path | None:
    if not isinstance(raw, str) or not raw.strip():
        return None
    value = raw.strip()
    lowered = value.lower()
    if lowered.startswith(("data:", "http://", "https://", "javascript:", "mailto:")):
        return None
    parsed = urlparse(value)
    if parsed.scheme == "file":
        return Path(unquote(parsed.path))
    if parsed.path.startswith("/api/files/runs/"):
        return _path_from_run_file_url(parsed.path)
    if parsed.scheme or value.startswith("/"):
        return None
    clean = unquote(parsed.path).split("#", 1)[0].split("?", 1)[0]
    if not clean or any(part in {"", ".", ".."} for part in clean.split("/")):
        return None
    return (base_dir / clean).resolve()


def _write_agent_pptx_export(
    source: Path,
    out: Path,
    *,
    width: int,
    height: int,
    artifact: dict[str, Any],
    settings: Settings | None,
    run_id: str,
) -> None:
    command = _code_editor_cmd_resolution(settings)
    if not command.get("available") or not command.get("cmd"):
        raise HTTPException(
            412,
            detail={
                "code": "missing_code_editor_command",
                "message": command.get("message") or _CODE_EDITOR_MISSING_MESSAGE,
            },
        )
    timeout_s = max(30, int(getattr(settings, "code_editor_timeout_s", 600) if settings else 600))
    max_attempts = max(1, int(getattr(settings, "code_editor_max_attempts", 2) if settings else 2))
    export_root = source.parent / "exports" / "_agent_pptx" / f"{int(time.time())}-{secrets.token_hex(4)}"
    export_root.mkdir(parents=True, exist_ok=True)

    repair_feedback: dict[str, Any] | None = None
    attempts: list[dict[str, Any]] = []
    for attempt_index in range(1, max_attempts + 1):
        attempt_dir = export_root / f"attempt_{attempt_index:02d}"
        if attempt_dir.exists():
            shutil.rmtree(attempt_dir)
        attempt_dir.mkdir(parents=True, exist_ok=True)
        log(
            "artifact_export.attempt_start",
            run_id=run_id,
            attempt=attempt_index,
            max_attempts=max_attempts,
            format="pptx",
            repair=repair_feedback is not None,
        )
        _stage_agent_pptx_export_inputs(source, attempt_dir)
        if repair_feedback is not None:
            atomic_write_json(attempt_dir / "validation_feedback.json", repair_feedback)
        prompt = _build_agent_pptx_export_prompt(
            source_name=source.name,
            artifact=artifact,
            width=width,
            height=height,
            repair_feedback=repair_feedback,
        )
        (attempt_dir / "export_prompt.md").write_text(prompt, encoding="utf-8")
        invocation = _invoke_agent_export_command(
            str(command["cmd"]),
            prompt=prompt,
            attempt_dir=attempt_dir,
            timeout_s=timeout_s,
        )
        if invocation.get("stdout_excerpt") or invocation.get("stderr_excerpt"):
            log(
                "artifact_export.agent_output",
                run_id=run_id,
                attempt=attempt_index,
                max_attempts=max_attempts,
                status=invocation.get("status"),
                reason=invocation.get("reason"),
                elapsed_s=invocation.get("elapsed_s"),
                stdout_excerpt=invocation.get("stdout_excerpt") or "",
                stderr_excerpt=invocation.get("stderr_excerpt") or "",
            )
        validation = _validate_agent_pptx_export(attempt_dir / "export.pptx")
        attempt_record = {
            "attempt": attempt_index,
            "attempt_dir": str(attempt_dir),
            "invocation": invocation,
            "validation": validation,
        }
        attempts.append(attempt_record)
        atomic_write_json(attempt_dir / "export_attempt_result.json", attempt_record)
        log(
            "artifact_export.pptx.agent_attempt",
            run_id=run_id,
            attempt=attempt_index,
            status=invocation.get("status"),
            reason=invocation.get("reason"),
            ok=validation.get("ok"),
            errors=validation.get("errors") or [],
        )
        if invocation.get("status") == "ok" and validation.get("ok") is True:
            log("artifact_export.attempt_ok", run_id=run_id, attempt=attempt_index, format="pptx")
            out.parent.mkdir(parents=True, exist_ok=True)
            shutil.copy2(attempt_dir / "export.pptx", out)
            atomic_write_json(
                out.with_suffix(".agent-export.json"),
                {
                    "render_mode": "external_code_editor_pptx_export",
                    "source_html": str(source),
                    "run_id": run_id,
                    "code_editor_harness": getattr(settings, "code_editor_harness", "codex") if settings else "codex",
                    "command_source": command.get("source"),
                    "attempts": attempts,
                    "canvas": {"w": width, "h": height},
                },
            )
            return
        log(
            "artifact_export.attempt_rejected",
            run_id=run_id,
            attempt=attempt_index,
            format="pptx",
            reason=invocation.get("reason") if invocation.get("status") != "ok" else "validation_failed",
            errors=validation.get("errors") or [],
        )
        repair_feedback = {
            "reason": invocation.get("reason") if invocation.get("status") != "ok" else "validation_failed",
            "invocation": invocation,
            "validation": validation,
        }

    raise HTTPException(
        500,
        detail={
            "code": "pptx_agent_export_failed",
            "message": "PowerPoint export agent did not produce a valid .pptx",
            "attempts": attempts[-2:],
        },
    )


def _stage_agent_pptx_export_inputs(source: Path, attempt_dir: Path) -> None:
    shutil.copy2(source, attempt_dir / "current.html")
    if source.name != "current.html":
        shutil.copy2(source, attempt_dir / source.name)
    for child in source.parent.iterdir():
        if child.name == "exports":
            continue
        target = attempt_dir / child.name
        try:
            if child.resolve() == source.resolve():
                continue
        except OSError:
            pass
        if child.is_dir():
            shutil.copytree(
                child,
                target,
                dirs_exist_ok=True,
                ignore=shutil.ignore_patterns("__pycache__", ".DS_Store"),
            )
        elif child.is_file():
            try:
                if child.stat().st_size > 50 * 1024 * 1024:
                    continue
                shutil.copy2(child, target)
            except OSError:
                continue


async def _run_pptx_export_in_background(
    *,
    run_id: str,
    source_run_id: str,
    source: Path,
    artifact: dict[str, Any],
    settings: Settings,
    export_key: _PptxExportKey,
) -> None:
    async with _RUNS_LOCK:
        state = _RUNS.get(run_id)
    if state is None:
        await _clear_pptx_export_registration(export_key, run_id)
        return
    started_at = time.time()
    try:
        canvas = artifact.get("canvas") if isinstance(artifact.get("canvas"), dict) else {}
        width, height = _html_export_canvas_size(source, artifact, canvas)
        export_dir = source.parent / "exports"
        export_dir.mkdir(parents=True, exist_ok=True)
        stem = _safe_export_stem(str(artifact.get("name") or source.stem or "artifact"))
        out = export_dir / f"{stem}.pptx"
        log("artifact_export.prepare", run_id=run_id, source_run_id=source_run_id, format="pptx")
        await asyncio.to_thread(
            _write_agent_pptx_export,
            source,
            out,
            width=width,
            height=height,
            artifact=artifact,
            settings=settings,
            run_id=run_id,
        )
        url = _run_file_url(source_run_id, _run_relative_path(out))
        filename = out.name
        source_artifact_id = str(artifact.get("artifact_id") or f"art_{source_run_id}")
        task_payload = {
            "source_artifact_id": source_artifact_id,
            "export_format": "pptx",
        }
        state.result_message = Message(
            id=f"msg_{run_id}",
            role="assistant",
            text=(
                "Exported editable PowerPoint from the current design. "
                f"The file is downloading now: {filename}"
            ),
            ts=int(time.time() * 1000),
            run_id=run_id,
            status="done",
            download_url=url,
            download_filename=filename,
            download_mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            task_type="artifact_export_pptx",
            task_payload=task_payload,
            source_artifact_id=source_artifact_id,
        )
        state.result_artifact = None
        _append_event(
            settings,
            state.conversation_id,
            "artifact.exported",
            run_id=run_id,
            data={
                "format": "pptx",
                "artifact_id": artifact.get("artifact_id"),
                "source_run_id": source_run_id,
                "download_url": url,
                "filename": filename,
                "elapsed_ms": int((time.time() - started_at) * 1000),
            },
        )
        log("artifact_export.done", run_id=run_id, format="pptx", filename=filename)
        _persisted_run_log("run.done", run_id)
    except asyncio.CancelledError:
        log("artifact_export.cancelled", run_id=run_id)
        state.error = state.error or "cancelled by user"
        source_artifact_id = str(artifact.get("artifact_id") or f"art_{source_run_id}")
        failure = Failure(
            status="cancelled",
            phase="artifact_export",
            agent_last_note="PowerPoint export cancelled.",
            produced_files=_list_produced_artifacts(run_id),
            elapsed_ms=int((time.time() - started_at) * 1000),
        )
        state.result_message = Message(
            id=f"msg_{run_id}",
            role="assistant",
            text="PowerPoint export cancelled.",
            ts=int(time.time() * 1000),
            run_id=run_id,
            status="error",
            failure=failure,
            task_type="artifact_export_pptx",
            task_payload={
                "source_artifact_id": source_artifact_id,
                "export_format": "pptx",
            },
            source_artifact_id=source_artifact_id,
        )
    except Exception as e:  # noqa: BLE001
        state.error = f"{type(e).__name__}: {e}"
        log("artifact_export.error", run_id=run_id, error=type(e).__name__, msg=str(e)[:300])
        _persisted_run_log("run.error", run_id, msg=state.error[:200])
        source_artifact_id = str(artifact.get("artifact_id") or f"art_{source_run_id}")
        failure = Failure(
            status="error",
            phase="artifact_export",
            agent_last_note=state.error,
            produced_files=_list_produced_artifacts(run_id),
            elapsed_ms=int((time.time() - started_at) * 1000),
        )
        state.result_message = Message(
            id=f"msg_{run_id}",
            role="assistant",
            text=f"PowerPoint export failed: {state.error}",
            ts=int(time.time() * 1000),
            run_id=run_id,
            status="error",
            failure=failure,
            task_type="artifact_export_pptx",
            task_payload={
                "source_artifact_id": source_artifact_id,
                "export_format": "pptx",
            },
            source_artifact_id=source_artifact_id,
        )
    finally:
        cleanup_task = asyncio.create_task(
            _clear_pptx_export_registration(export_key, run_id)
        )
        try:
            await asyncio.shield(cleanup_task)
        except asyncio.CancelledError:
            await cleanup_task


async def _clear_pptx_export_registration(
    export_key: _PptxExportKey,
    run_id: str,
) -> None:
    async with _RUNS_LOCK:
        if _PPTX_EXPORT_RUNS.get(export_key) == run_id:
            _PPTX_EXPORT_RUNS.pop(export_key, None)


def _build_agent_pptx_export_prompt(
    *,
    source_name: str,
    artifact: dict[str, Any],
    width: int,
    height: int,
    repair_feedback: dict[str, Any] | None,
) -> str:
    name = str(artifact.get("name") or source_name or "AutoDesign artifact").strip()
    repair_block = ""
    if repair_feedback is not None:
        repair_block = (
            "\nThis is a repair attempt. Read `validation_feedback.json` first, then fix the export. "
            "Do not repeat a conversion strategy that created invalid or unreadable PPTX output.\n"
        )
    return f"""Export the current AutoDesign HTML artifact to an editable PowerPoint file.

Artifact name: {name}
Source HTML files:
- current.html
- {source_name}

Canvas size: {width} x {height} CSS pixels.

Output contract:
- Write `export.pptx`.
- Write `export_done.json` with a concise JSON summary.
- Keep all work inside this directory.
- Use this Python interpreter exactly: {sys.executable}
- Do not install packages, create virtual environments, or run setup commands.

Quality contract:
- Match the browser-rendered poster/design as closely as practical.
- Set the PowerPoint slide size to the same aspect ratio and logical size as the HTML canvas.
- Avoid naive text-node extraction that creates overlapping text boxes.
- Prefer section-level and table-cell-level extraction: one text box per visual paragraph/list/table cell is better than one per DOM text node.
- Use native editable PowerPoint text for titles, body copy, labels, section bars, metric chips, and tables wherever practical.
- Use images only for real source figures, charts, logos, or diagrams that are already images in the HTML.
- It is acceptable to rasterize complex SVG/canvas/figure artwork, but do not rasterize ordinary body text if it can remain editable.
- Preserve the visual hierarchy, colors, section header bands, panel boundaries, and spacing.
- The final PPTX should open as a one-slide poster/presentation without text collisions, tiny top-left rendering, or giant blank unused page space.
- Do not use remote, file:, script, iframe, event-handler, or unsafe URLs.

Suggested implementation:
1. Use Playwright to render `current.html` locally and inspect computed boxes at {width}x{height}.
2. Use python-pptx to create a single blank slide at the matching dimensions.
3. Rebuild the poster from semantic blocks/sections/tables rather than dumping every DOM text range.
4. Add source images with their rendered bbox and add native text/table shapes around them.
5. Save as `export.pptx`, then reopen it with python-pptx to verify it has at least one slide and useful shapes.
{repair_block}
"""


def _invoke_agent_export_command(
    command: str,
    *,
    prompt: str,
    attempt_dir: Path,
    timeout_s: int,
) -> dict[str, Any]:
    pptx_path = attempt_dir / "export.pptx"
    done_marker = attempt_dir / "export_done.json"
    stdout_path = attempt_dir / ".export_agent.stdout.tmp"
    stderr_path = attempt_dir / ".export_agent.stderr.tmp"
    for path in (pptx_path, done_marker):
        try:
            path.unlink()
        except OSError:
            pass
    try:
        cmd = shlex.split(command)
    except ValueError as exc:
        return {"status": "error", "reason": f"command_parse_error: {exc}"}
    if not cmd:
        return {"status": "error", "reason": "empty_command"}

    start = time.monotonic()
    timed_out = False
    returncode: int | None = None
    reason = "process_exit"
    try:
        env = os.environ.copy()
        author_python = _first_env_value(
            "AUTODESIGN_AUTHOR_PYTHON",
            "DESIGN_ANYTHING_AUTHOR_PYTHON",
        ) or sys.executable
        env["AUTODESIGN_AUTHOR_PYTHON"] = author_python
        env.setdefault("DESIGN_ANYTHING_AUTHOR_PYTHON", author_python)
        with stdout_path.open("w", encoding="utf-8") as stdout_f, stderr_path.open("w", encoding="utf-8") as stderr_f:
            proc = subprocess.Popen(
                cmd,
                stdin=subprocess.PIPE,
                stdout=stdout_f,
                stderr=stderr_f,
                text=True,
                cwd=str(attempt_dir),
                env=env,
                start_new_session=(os.name == "posix"),
            )
            try:
                if proc.stdin is not None:
                    proc.stdin.write(prompt)
                    proc.stdin.close()
            except BrokenPipeError:
                pass
            deadline = start + timeout_s
            while True:
                returncode = proc.poll()
                if done_marker.exists() and pptx_path.exists():
                    returncode = _terminate_subprocess_group(proc)
                    reason = "done_marker"
                    break
                if returncode is not None:
                    break
                if time.monotonic() >= deadline:
                    timed_out = True
                    returncode = _terminate_subprocess_group(proc)
                    reason = "timeout"
                    break
                time.sleep(0.25)
    except OSError as exc:
        return {
            "status": "error",
            "reason": f"command_start_error: {exc}",
            "elapsed_s": round(time.monotonic() - start, 3),
        }

    status = "ok" if pptx_path.exists() and done_marker.exists() and not timed_out else "error"
    if status != "ok" and pptx_path.exists() and not done_marker.exists():
        reason = "missing_done_marker"
    if status != "ok" and not pptx_path.exists():
        reason = "missing_export_pptx"
    return {
        "status": status,
        "reason": reason,
        "returncode": returncode,
        "timed_out": timed_out,
        "elapsed_s": round(time.monotonic() - start, 3),
        "stdout_excerpt": _tail_file_text(stdout_path, limit=1400),
        "stderr_excerpt": _tail_file_text(stderr_path, limit=900),
    }


def _validate_agent_pptx_export(path: Path) -> dict[str, Any]:
    errors: list[str] = []
    warnings: list[str] = []
    if not path.exists():
        return {"ok": False, "errors": ["export.pptx missing"], "warnings": warnings}
    try:
        size = path.stat().st_size
    except OSError as exc:
        return {"ok": False, "errors": [f"export.pptx unreadable: {exc}"], "warnings": warnings}
    if size < 2048:
        errors.append("export.pptx is too small")
    try:
        from pptx import Presentation

        prs = Presentation(str(path))
        slide_count = len(prs.slides)
        shape_count = sum(len(slide.shapes) for slide in prs.slides)
        if slide_count < 1:
            errors.append("export.pptx has no slides")
        if shape_count < 3:
            warnings.append("export.pptx has very few shapes")
    except Exception as exc:  # noqa: BLE001
        errors.append(f"export.pptx could not be opened: {exc}")
        slide_count = 0
        shape_count = 0
    return {
        "ok": not errors,
        "errors": errors,
        "warnings": warnings,
        "bytes": size,
        "slide_count": slide_count,
        "shape_count": shape_count,
    }


def _terminate_subprocess_group(proc: subprocess.Popen[str]) -> int | None:
    if proc.poll() is not None:
        return proc.returncode
    try:
        if os.name == "posix":
            os.killpg(proc.pid, signal.SIGTERM)
        else:
            proc.terminate()
        return proc.wait(timeout=5)
    except Exception:  # noqa: BLE001
        try:
            if os.name == "posix":
                os.killpg(proc.pid, signal.SIGKILL)
            else:
                proc.kill()
        except Exception:
            pass
    try:
        return proc.wait(timeout=1)
    except Exception:  # noqa: BLE001
        return proc.poll()


def _tail_file_text(path: Path, *, limit: int) -> str:
    try:
        text = path.read_text(encoding="utf-8", errors="replace").strip()
    except OSError:
        return ""
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text if len(text) <= limit else "..." + text[-limit:]


def _coding_agent_smoke_log_excerpt(path: Path, *, limit: int) -> str:
    try:
        text = path.read_text(encoding="utf-8", errors="replace").strip()
    except OSError:
        return ""
    if not text:
        return ""
    lines = []
    for line in text.splitlines():
        if (
            "codex_core_skills::loader: ignoring interface.icon_" in line
            or "codex_core_plugins::manifest: ignoring interface.defaultPrompt" in line
        ):
            continue
        lines.append(line)
    cleaned = "\n".join(lines).strip()
    cleaned = re.sub("\x1b\\[[0-9;]*m", "", cleaned)
    cleaned = re.sub(r"\n{3,}", "\n\n", cleaned)
    return cleaned if len(cleaned) <= limit else "..." + cleaned[-limit:]


def _coding_agent_smoke_stderr_excerpt(stderr_path: Path, debug_path: Path, *, limit: int) -> str:
    stderr_text = _coding_agent_smoke_log_excerpt(stderr_path, limit=limit)
    debug_text = _coding_agent_smoke_log_excerpt(debug_path, limit=limit)
    if stderr_text and debug_text:
        combined = f"{stderr_text}\n\n{debug_text}"
        return combined if len(combined) <= limit else "..." + combined[-limit:]
    return stderr_text or debug_text


def _write_html_pptx_export(source: Path, out: Path, *, width: int, height: int) -> None:
    dom = _extract_html_export_dom(source, width=width, height=height)
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Emu, Pt
    except Exception as exc:
        raise HTTPException(500, detail=f"PowerPoint export unavailable: {exc}") from exc

    px_to_emu = 9525
    prs = Presentation()
    prs.slide_width = Emu(max(1, int(round(width * px_to_emu))))
    prs.slide_height = Emu(max(1, int(round(height * px_to_emu))))
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    root_bg = _css_color_to_rgb(dom.get("background"))
    if root_bg:
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(*root_bg)

    def emu(v: Any) -> Emu:
        try:
            n = float(v)
        except (TypeError, ValueError):
            n = 0.0
        return Emu(int(round(n * px_to_emu)))

    for item in (dom.get("backgrounds") or [])[:180]:
        rect = item.get("rect") if isinstance(item, dict) else None
        color = _css_color_to_rgb(item.get("background") if isinstance(item, dict) else None)
        if not rect or not color:
            continue
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            emu(rect.get("x")),
            emu(rect.get("y")),
            emu(rect.get("w")),
            emu(rect.get("h")),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*color)
        line_color = _css_color_to_rgb(item.get("border") if isinstance(item, dict) else None)
        if line_color:
            shape.line.color.rgb = RGBColor(*line_color)
            shape.line.width = Pt(0.5)
        else:
            shape.line.fill.background()

    with tempfile.TemporaryDirectory(prefix="opendesign-pptx-export-") as tmp:
        tmp_dir = Path(tmp)
        for idx, item in enumerate((dom.get("images") or [])[:80]):
            rect = item.get("rect") if isinstance(item, dict) else None
            src = item.get("src") if isinstance(item, dict) else None
            if not rect or not src:
                continue
            img_path = _pptx_image_source(src, source.parent, tmp_dir / f"image_{idx:03d}")
            if not img_path or not img_path.exists():
                continue
            try:
                slide.shapes.add_picture(
                    str(img_path),
                    emu(rect.get("x")),
                    emu(rect.get("y")),
                    width=emu(rect.get("w")),
                    height=emu(rect.get("h")),
                )
            except Exception:
                continue

        for item in (dom.get("texts") or [])[:650]:
            if not isinstance(item, dict):
                continue
            rect = item.get("rect")
            text = str(item.get("text") or "").strip()
            if not rect or not text:
                continue
            box = slide.shapes.add_textbox(
                emu(rect.get("x")),
                emu(rect.get("y")),
                emu(max(2, float(rect.get("w") or 2))),
                emu(max(2, float(rect.get("h") or 2))),
            )
            tf = box.text_frame
            tf.margin_left = 0
            tf.margin_right = 0
            tf.margin_top = 0
            tf.margin_bottom = 0
            tf.word_wrap = True
            paragraph = tf.paragraphs[0]
            align = str(item.get("align") or "").lower()
            if align == "center":
                paragraph.alignment = PP_ALIGN.CENTER
            elif align == "right":
                paragraph.alignment = PP_ALIGN.RIGHT
            run = paragraph.add_run()
            run.text = text
            font = run.font
            family = _clean_font_family(str(item.get("fontFamily") or ""))
            if family:
                font.name = family
            size_px = _positive_float(item.get("fontSize")) or 12.0
            font.size = Pt(max(4.0, min(96.0, size_px * 0.75)))
            color = _css_color_to_rgb(item.get("color"))
            if color:
                font.color.rgb = RGBColor(*color)
            weight = str(item.get("fontWeight") or "").lower()
            font.bold = weight == "bold" or (_positive_float(weight) or 0) >= 600
            style = str(item.get("fontStyle") or "").lower()
            font.italic = style == "italic"

        out.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(out))


def _extract_html_export_dom(source: Path, *, width: int, height: int) -> dict[str, Any]:
    try:
        from playwright.sync_api import sync_playwright
    except Exception as exc:
        raise HTTPException(500, detail=f"PowerPoint export needs Playwright: {exc}") from exc

    script = r"""
() => {
  const root = document.querySelector('.paper-poster') || document.querySelector('.canvas') || document.body;
  const rootRect = root.getBoundingClientRect();
  const clean = (s) => String(s || '').replace(/\s+/g, ' ').trim();
  const transparent = (c) => !c || c === 'transparent' || /rgba\([^)]*,\s*0\)/i.test(c);
  const rectOf = (r) => ({
    x: Math.max(0, r.left - rootRect.left),
    y: Math.max(0, r.top - rootRect.top),
    w: Math.max(0, r.width),
    h: Math.max(0, r.height),
  });
  const visible = (el) => {
    const cs = getComputedStyle(el);
    const r = el.getBoundingClientRect();
    return cs.display !== 'none' && cs.visibility !== 'hidden' && Number(cs.opacity || 1) > 0.02 && r.width > 1 && r.height > 1;
  };
  const rootStyle = getComputedStyle(root);
  const backgrounds = [];
  const bgSelector = [
    '.paper-poster', '.poster-section', '.poster-header', '.poster-title',
    '.section-title', '.section-heading', '.metric', '.metric-card', '.badge',
    '.callout', '.figure-card', '.table-card', 'section', 'header', 'table', 'th', 'td'
  ].join(',');
  Array.from(root.querySelectorAll(bgSelector)).forEach((el) => {
    if (!visible(el)) return;
    const cs = getComputedStyle(el);
    if (transparent(cs.backgroundColor)) return;
    const rect = rectOf(el.getBoundingClientRect());
    if (rect.w * rect.h < 120) return;
    backgrounds.push({
      rect,
      background: cs.backgroundColor,
      border: transparent(cs.borderTopColor) || parseFloat(cs.borderTopWidth || '0') <= 0 ? '' : cs.borderTopColor,
    });
  });

  const images = Array.from(root.querySelectorAll('img')).filter(visible).map((el) => ({
    src: el.currentSrc || el.src || el.getAttribute('src') || '',
    alt: el.alt || '',
    rect: rectOf(el.getBoundingClientRect()),
  }));

  const texts = [];
  const walker = document.createTreeWalker(root, NodeFilter.SHOW_TEXT, {
    acceptNode(node) {
      const text = clean(node.textContent);
      if (!text) return NodeFilter.FILTER_REJECT;
      const parent = node.parentElement;
      if (!parent || !visible(parent)) return NodeFilter.FILTER_REJECT;
      if (parent.closest('.ld-toolbar,.ld-modal-backdrop,.od-flow-handle,.od-area-overlay')) {
        return NodeFilter.FILTER_REJECT;
      }
      return NodeFilter.FILTER_ACCEPT;
    }
  });
  while (texts.length < 650) {
    const node = walker.nextNode();
    if (!node) break;
    const parent = node.parentElement;
    const range = document.createRange();
    range.selectNodeContents(node);
    const rects = Array.from(range.getClientRects()).filter((r) => r.width > 1 && r.height > 1);
    range.detach();
    if (!rects.length) continue;
    const cs = getComputedStyle(parent);
    const text = clean(node.textContent);
    rects.slice(0, 3).forEach((r) => {
      const rect = rectOf(r);
      if (rect.w * rect.h < 4) return;
      texts.push({
        text,
        rect,
        color: cs.color,
        fontSize: parseFloat(cs.fontSize || '12') || 12,
        fontFamily: cs.fontFamily,
        fontWeight: cs.fontWeight,
        fontStyle: cs.fontStyle,
        align: cs.textAlign,
      });
    });
  }
  return {
    background: rootStyle.backgroundColor || getComputedStyle(document.body).backgroundColor,
    backgrounds,
    images,
    texts,
  };
}
"""
    try:
        with sync_playwright() as p:
            browser = p.chromium.launch(args=["--no-sandbox"])
            page = browser.new_page(
                viewport={"width": max(1, int(width)), "height": max(1, int(height))},
                device_scale_factor=1,
            )
            page.set_default_timeout(30_000)
            page.goto(source.resolve().as_uri(), wait_until="load", timeout=30_000)
            try:
                page.wait_for_load_state("networkidle", timeout=2000)
            except Exception:
                pass
            data = page.evaluate(script)
            browser.close()
    except Exception as exc:
        raise HTTPException(500, detail=f"PowerPoint export failed: {exc}") from exc
    return data if isinstance(data, dict) else {}


def _pptx_image_source(src: str, base_dir: Path, target_prefix: Path) -> Path | None:
    if src.startswith("data:"):
        match = re.match(r"data:([^;,]+)?;base64,(.*)$", src, flags=re.IGNORECASE | re.DOTALL)
        if not match:
            return None
        mime = match.group(1) or "image/png"
        ext = mimetypes.guess_extension(mime) or ".png"
        path = target_prefix.with_suffix(ext)
        try:
            path.write_bytes(base64.b64decode(match.group(2)))
        except Exception:
            return None
        return path
    return _resolve_html_ref_path(src, base_dir)


def _positive_float(raw: Any) -> float | None:
    try:
        value = float(str(raw).strip())
    except (TypeError, ValueError):
        return None
    return value if value > 0 else None


def _css_color_to_rgb(raw: Any) -> tuple[int, int, int] | None:
    if not isinstance(raw, str):
        return None
    value = raw.strip()
    if not value or value.lower() == "transparent":
        return None
    hex_match = re.match(r"^#([0-9a-f]{3}|[0-9a-f]{6})$", value, flags=re.IGNORECASE)
    if hex_match:
        h = hex_match.group(1)
        if len(h) == 3:
            h = "".join(ch * 2 for ch in h)
        return int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    rgb_match = re.match(
        r"^rgba?\(\s*([0-9.]+)[,\s]+([0-9.]+)[,\s]+([0-9.]+)(?:[,\s/]+([0-9.]+%?))?\s*\)$",
        value,
        flags=re.IGNORECASE,
    )
    if not rgb_match:
        return None
    alpha = rgb_match.group(4)
    if alpha is not None:
        try:
            opacity = float(alpha.strip("%")) / (100.0 if alpha.endswith("%") else 1.0)
        except ValueError:
            opacity = 1.0
        if opacity <= 0.05:
            return None
    return tuple(max(0, min(255, int(float(rgb_match.group(i))))) for i in (1, 2, 3))  # type: ignore[return-value]


def _clean_font_family(raw: str) -> str | None:
    if not raw.strip():
        return None
    first = raw.split(",", 1)[0].strip().strip("'\"")
    return first or None


def _derived_result_file(
    run_id: str,
    raw: Any,
    *,
    required_root: str | None = None,
) -> Path:
    run_dir = (RUNS_DIR / run_id).resolve()
    requested = Path(str(raw or ""))
    if not requested.is_absolute():
        requested = run_dir / requested
    if requested.is_symlink():
        raise ValueError("derived result file must not be a symlink")
    resolved = requested.resolve(strict=True)
    try:
        relative = resolved.relative_to(run_dir)
    except ValueError as exc:
        raise ValueError("derived result file is outside its run") from exc
    if required_root is not None and (
        not relative.parts or relative.parts[0] != required_root
    ):
        raise ValueError(
            f"derived result file is outside {required_root}/"
        )
    if not resolved.is_file():
        raise ValueError("derived result path is not a regular file")
    return resolved


def _prepare_derived_completion(
    *,
    run_id: str,
    state: _RunState,
    job_kind: str,
    parent_run_id: str,
    descriptor: dict[str, Any],
    outcome: WorkerOutcome,
) -> tuple[bool, Artifact | None, Message, dict[str, Any]]:
    if not outcome.ok or not isinstance(outcome.result, dict):
        raise ValueError(outcome.error or "derived worker failed")
    result = outcome.result
    artifact: Artifact | None = None
    event_data: dict[str, Any]
    if job_kind == "editable_video_render":
        mp4_path = _derived_result_file(run_id, result.get("mp4_path"))
        artifact = _build_video_artifact_from_path(
            RUNS_DIR / run_id,
            run_id,
            mp4_path,
            baseline_artifact_json=state.baseline_artifact_json,
        )
        if artifact is None:
            raise ValueError("editable video render has no publishable artifact")
        text = "Rendered editable video MP4."
        event_data = {
            "artifact_type": "video",
            "name": artifact.name,
            "native_format": artifact.native_format,
            "source": "editable_video_demo",
            "parent_run_id": parent_run_id,
        }
    elif job_kind == "artifact_edit":
        source = _derived_result_file(
            run_id,
            result.get("source_path"),
            required_root="final",
        )
        artifact_type = str(result.get("artifact_type") or "")
        if artifact_type not in {"poster", "deck", "landing", "video"}:
            raise ValueError("artifact edit result type is invalid")
        if artifact_type != state.artifact_type:
            raise ValueError("artifact edit result type changed during the run")
        candidate_lineage = result.get("candidate_lineage")
        if not isinstance(candidate_lineage, dict):
            raise ValueError("artifact edit candidate lineage is invalid")
        if candidate_lineage:
            lineage_path = RUNS_DIR / run_id / "candidate_draft_lineage.json"
            if _read_json_file(lineage_path) != candidate_lineage:
                raise ValueError("artifact edit lineage does not match its durable file")
            artifact = _candidate_draft_artifact_from_lineage(
                RUNS_DIR / run_id,
                run_id,
                artifact_type,
                candidate_lineage,
                source=source,
            )
        else:
            artifact = _build_artifact_response(
                RUNS_DIR / run_id,
                run_id,
                artifact_type,
                baseline_artifact_json=state.baseline_artifact_json,
            )
        if artifact is None:
            raise ValueError("artifact edit has no publishable artifact")
        text = (
            f"Applied {len(result.get('restored_layer_ids') or [])} layer edit"
            f"{'s' if len(result.get('restored_layer_ids') or []) != 1 else ''}."
        )
        event_data = {
            "artifact_type": artifact_type,
            "name": artifact.name,
            "native_format": artifact.native_format,
            "parent_artifact_id": artifact.parent_artifact_id,
            "source": "artifact_edit",
            "parent_run_id": parent_run_id,
            "restored_layer_ids": result.get("restored_layer_ids") or [],
            "skipped": result.get("skipped") or [],
        }
    elif job_kind == "attempt_fork":
        source = _derived_result_file(
            run_id,
            result.get("source_path"),
            required_root="final",
        )
        lineage = result.get("lineage")
        if not isinstance(lineage, dict):
            raise ValueError("attempt fork lineage is missing")
        lineage_path = RUNS_DIR / run_id / "candidate_draft_lineage.json"
        if _read_json_file(lineage_path) != lineage:
            raise ValueError("attempt fork lineage does not match its durable file")
        artifact_type = str(result.get("artifact_type") or "")
        if artifact_type not in {"poster", "deck", "landing", "video"}:
            raise ValueError("attempt fork artifact type is invalid")
        artifact = _candidate_draft_artifact_from_lineage(
            RUNS_DIR / run_id,
            run_id,
            artifact_type,
            lineage,
            source=source,
        )
        text = f"Prepared Attempt {lineage.get('source_attempt') or '?'} for editing."
        event_data = {
            "artifact_type": artifact_type,
            "name": artifact.name,
            "source": "attempt_fork",
            "parent_run_id": parent_run_id,
            "source_attempt": lineage.get("source_attempt"),
        }
    elif job_kind == "candidate_publish":
        _derived_result_file(
            run_id,
            result.get("source_path"),
            required_root="final",
        )
        lineage = result.get("lineage")
        if not isinstance(lineage, dict) or lineage.get("status") not in {
            "validated",
            "published",
        }:
            raise ValueError("candidate publish lineage is missing or unvalidated")
        lineage_path = RUNS_DIR / run_id / "candidate_draft_lineage.json"
        if _read_json_file(lineage_path) != lineage:
            raise ValueError("candidate publish lineage does not match its durable file")
        artifact_type = str(result.get("artifact_type") or "")
        if artifact_type not in {"poster", "deck", "landing", "video"}:
            raise ValueError("candidate publish artifact type is invalid")
        artifact = _build_artifact_response(
            RUNS_DIR / run_id,
            run_id,
            artifact_type,
            baseline_artifact_json=state.baseline_artifact_json,
        )
        if artifact is None:
            raise ValueError("candidate publish has no publishable artifact")
        artifact = artifact.model_copy(
            update={
                "candidate_draft": False,
                "attempt_lineage": {**lineage, "status": "published"},
            }
        )
        text = f"Published Attempt {lineage.get('source_attempt') or '?'} as the new final."
        event_data = {
            "artifact_type": artifact_type,
            "name": artifact.name,
            "source": "candidate_publish",
            "parent_run_id": parent_run_id,
            "source_run_id": lineage.get("source_run_id"),
            "source_attempt": lineage.get("source_attempt"),
        }
    elif job_kind == "poster_code_edit":
        run_dir = Path(str(result.get("run_dir") or "")).resolve()
        if run_dir != (RUNS_DIR / run_id).resolve():
            raise ValueError("poster edit result belongs to a different run")
        poster_path = _derived_result_file(run_id, result.get("poster_path"))
        preview_path = _derived_result_file(run_id, result.get("preview_path"))
        expected_final = (RUNS_DIR / run_id / "final").resolve()
        if poster_path.parent != expected_final or preview_path.parent != expected_final:
            raise ValueError("poster edit final files are outside final/")
        artifact = _build_artifact_response(
            RUNS_DIR / run_id,
            run_id,
            "poster",
            baseline_artifact_json=state.baseline_artifact_json,
        )
        if artifact is None:
            raise ValueError("poster edit has no publishable artifact")
        text = "Applied the poster revision with the external code editor."
        event_data = {
            "artifact_type": "poster",
            "name": artifact.name,
            "native_format": artifact.native_format,
            "parent_artifact_id": artifact.parent_artifact_id,
            "render_mode": "external_code_editor_revision",
            "source_run_id": parent_run_id,
            "palette_id": state.palette_id,
            "attempts": result.get("attempts") or [],
            "validation_summary": result.get("validation_summary") or {},
            "selection_context_summary": result.get("selection_context_summary") or {},
        }
    elif job_kind == "pptx_export":
        pptx_path = _derived_result_file(
            run_id,
            result.get("pptx_path"),
            required_root="exports",
        )
        source_artifact_id = str(
            descriptor.get("source_artifact_id") or f"art_{parent_run_id}"
        )
        download_url = _run_file_url(
            run_id,
            _run_relative_path(pptx_path),
        )
        message = Message(
            id=f"msg_{run_id}",
            role="assistant",
            text=(
                "Exported editable PowerPoint from the current design. "
                f"The file is downloading now: {pptx_path.name}"
            ),
            ts=int(time.time() * 1000),
            run_id=run_id,
            status="done",
            download_url=download_url,
            download_filename=pptx_path.name,
            download_mime_type=(
                "application/vnd.openxmlformats-officedocument."
                "presentationml.presentation"
            ),
            task_type="artifact_export_pptx",
            task_payload={
                "source_artifact_id": source_artifact_id,
                "export_format": "pptx",
            },
            source_artifact_id=source_artifact_id,
        )
        return True, None, message, {
            "format": "pptx",
            "artifact_id": source_artifact_id,
            "source_run_id": parent_run_id,
            "download_url": download_url,
            "filename": pptx_path.name,
        }
    elif job_kind == "video_export_retry":
        _derived_result_file(run_id, result.get("manifest_path"))
        _derived_result_file(run_id, result.get("media_probe_path"))
        _derived_result_file(run_id, result.get("mp4_path"))
        artifact = _build_artifact_response(
            RUNS_DIR / run_id,
            run_id,
            "video",
            baseline_artifact_json=state.baseline_artifact_json,
        )
        if artifact is None:
            raise ValueError("video export retry has no publishable artifact")
        text = "Retried Video export from the existing authored project."
        event_data = {
            "artifact_type": "video",
            "name": artifact.name,
            "native_format": artifact.native_format,
            "source": "video_export_retry",
            "parent_run_id": parent_run_id,
            "pointer_cleanup_warnings": list(
                result["pointer_cleanup_warnings"]
            ),
        }
    else:
        raise ValueError(f"unsupported derived result kind: {job_kind}")
    message = Message(
        id=f"msg_{run_id}",
        role="assistant",
        text=text,
        ts=int(time.time() * 1000),
        run_id=run_id,
        artifact_id=artifact.artifact_id,
        status="done",
    )
    return True, artifact, message, event_data


async def _monitor_supervised_derived_job(
    *,
    run_id: str,
    state: _RunState,
    job_kind: str,
    parent_run_id: str,
    descriptor: dict[str, Any],
    recovered_outcome: WorkerOutcome | None = None,
) -> None:
    runtime = _web_run_runtime()
    if recovered_outcome is None:
        try:
            outcome = await runtime.supervisor.wait(run_id)
        except asyncio.CancelledError:
            raise
        except Exception as exc:  # noqa: BLE001
            outcome = WorkerOutcome(
                run_id=run_id,
                job_kind=job_kind,
                returncode=1,
                ok=False,
                result=None,
                error=f"{type(exc).__name__}: {exc}",
                relayed_events=0,
            )
    else:
        outcome = recovered_outcome
    outcome_warnings = _immutable_pointer_cleanup_warnings(
        getattr(outcome, "pointer_cleanup_warnings", ())
    )
    record = runtime.control_store.read(run_id)
    if record.state in {"cancelling", "cancelled"}:
        existing_message = state.result_message
        if existing_message is None:
            state.result_message = Message(
                id=f"msg_{run_id}",
                role="assistant",
                text="Run cancelled.",
                ts=int(time.time() * 1000),
                run_id=run_id,
                status="error",
                failure=Failure(
                    status="cancelled",
                    pointer_cleanup_warnings=list(outcome_warnings),
                ),
            )
        else:
            failure = existing_message.failure or Failure(status="cancelled")
            failure = failure.model_copy(update={
                "status": "cancelled",
                "pointer_cleanup_warnings": list(outcome_warnings),
            })
            state.result_message = existing_message.model_copy(update={
                "artifact_id": None,
                "status": "error",
                "failure": failure,
            })
        state.result_artifact = None
        return
    if record.state != "completing":
        return

    success = False
    artifact: Artifact | None = None
    event_data: dict[str, Any] = {}
    direct_worker_failure = not outcome.ok or not isinstance(outcome.result, dict)
    if direct_worker_failure:
        state.error = outcome.error or f"worker exited with status {outcome.returncode}"
    else:
        try:
            success, artifact, message, event_data = _prepare_derived_completion(
                run_id=run_id,
                state=state,
                job_kind=job_kind,
                parent_run_id=parent_run_id,
                descriptor=descriptor,
                outcome=outcome,
            )
        except Exception as exc:  # noqa: BLE001
            state.error = f"{type(exc).__name__}: {exc}"

    def prepare_failure_message(*, direct_failure: bool) -> Message:
        state.error = state.error or "derived worker failed"
        exit_failure = _worker_exit_failure_fields(outcome)
        failure_phase = (
            outcome.failure_phase
            or exit_failure.get("phase")
            or job_kind
        )
        retry_route = (
            _video_export_retry_route(
                phase=failure_phase,
                error=outcome.error or state.error,
            )
            if job_kind == "video_export_retry"
            else None
        )
        return Message(
            id=f"msg_{run_id}",
            role="assistant",
            text=(
                state.error
                if direct_failure
                else f"Derived task failed: {state.error}"
            ),
            ts=int(time.time() * 1000),
            run_id=run_id,
            status="error",
            failure=Failure(
                status="error",
                phase=failure_phase,
                error_code=exit_failure.get("error_code"),
                error_message=exit_failure.get("error_message"),
                error_detail=exit_failure.get("error_detail"),
                retry_route=retry_route,
                parent_run_id=parent_run_id,
                agent_last_note=state.error,
                pointer_cleanup_warnings=list(outcome_warnings),
                produced_files=_list_produced_artifacts(run_id),
                artifact_type=state.artifact_type,
            ),
        )

    if not success:
        message = prepare_failure_message(
            direct_failure=direct_worker_failure,
        )

    state.result_artifact = artifact
    state.result_message = message
    digest_value: Any = outcome.result if outcome.result is not None else {
        "returncode": outcome.returncode,
        "error": outcome.error,
    }
    result_digest = hashlib.sha256(
        json.dumps(
            digest_value,
            ensure_ascii=False,
            sort_keys=True,
            separators=(",", ":"),
        ).encode("utf-8")
    ).hexdigest()
    try:
        if job_kind in {"artifact_edit", "candidate_publish"}:
            async with _derived_tree_locks(parent_run_id):
                if success:
                    try:
                        readiness_request: RunWorkerRequest | None = None
                        if job_kind == "candidate_publish":
                            direct_publish = (
                                _read_direct_candidate_publish_descriptor(run_id)
                            )
                            if direct_publish is not None:
                                recovery_settings = SETTINGS
                                if recovery_settings is None:
                                    raise ValueError(
                                        "candidate publish settings are unavailable"
                                    )
                                readiness_request, _recovered_state = (
                                    _candidate_publish_recovery_request(
                                        run_id,
                                        descriptor,
                                        direct_publish,
                                        recovery_settings,
                                    )
                                )
                        _require_derived_ancestors_not_cancelled(parent_run_id)
                        _require_derived_source_ready(
                            parent_run_id,
                            job_kind,
                            readiness_request,
                        )
                        if job_kind == "candidate_publish":
                            await _quiesce_bundle_candidate_publication_source(
                                run_id
                            )
                    except Exception as exc:  # fail closed before publication
                        current = runtime.control_store.read(run_id)
                        if current.state in {"cancelling", "cancelled"}:
                            state.result_artifact = None
                            state.result_message = None
                            return
                        success = False
                        artifact = None
                        state.error = _redacted_error_detail(
                            f"{type(exc).__name__}: {exc}"
                        ) or "candidate publication validation failed"
                        message = prepare_failure_message(direct_failure=False)
                        state.result_artifact = None
                        state.result_message = message
                accepted = await runtime.supervisor.accept_completion(
                    run_id,
                    terminal_state="completed" if success else "failed",
                    publishable=success,
                    result_digest=result_digest,
                )
        else:
            accepted = await runtime.supervisor.accept_completion(
                run_id,
                terminal_state="completed" if success else "failed",
                publishable=success,
                result_digest=result_digest,
            )
    except Exception as exc:  # completion may lose to cancellation
        current = runtime.control_store.read(run_id)
        if current.state in {"cancelling", "cancelled"}:
            state.result_artifact = None
            state.result_message = None
            return
        expected_terminal = "completed" if success else "failed"
        if current.state == expected_terminal:
            accepted = current
        elif current.state == "completing":
            artifact = None
            if success:
                state.error = _redacted_error_detail(
                    f"{type(exc).__name__}: {exc}"
                ) or "candidate publication validation failed"
                message = prepare_failure_message(direct_failure=False)
            success = False
            state.result_artifact = None
            state.result_message = message
            try:
                accepted = await runtime.supervisor.accept_completion(
                    run_id,
                    terminal_state="failed",
                    publishable=False,
                    result_digest=result_digest,
                )
            except Exception:
                current = runtime.control_store.read(run_id)
                if current.state in {"cancelling", "cancelled"}:
                    state.result_artifact = None
                    state.result_message = None
                    return
                if current.state != "failed":
                    return
                accepted = current
        else:
            state.result_artifact = None
            state.result_message = None
            return
    expected_terminal = "completed" if success else "failed"
    if accepted.state != expected_terminal:
        state.result_artifact = None
        state.result_message = None
        return
    if not success:
        _append_event(
            _settings_or_boot(),
            state.conversation_id,
            "artifact.generation_failed",
            run_id=run_id,
            data={
                "status": "error",
                "artifact_type": state.artifact_type,
                "error": state.error,
                "produced_files": _list_produced_artifacts(run_id),
                "failure": _dump_model(message.failure),
                "canvas_plan": _read_canvas_plan(RUNS_DIR / run_id),
                "deck_plan": _read_deck_plan(RUNS_DIR / run_id),
            },
        )
        return
    event_name = "artifact.exported" if job_kind == "pptx_export" else "artifact.generated"
    if job_kind == "artifact_edit":
        return
    _append_event(
        _settings_or_boot(),
        state.conversation_id,
        event_name,
        run_id=run_id,
        artifact_id=artifact.artifact_id if artifact is not None else None,
        data=event_data,
    )


async def _monitor_supervised_pipeline(
    *,
    run_id: str,
    state: _RunState,
    recovered_outcome: WorkerOutcome | None = None,
) -> None:
    started_at = time.time()
    runtime = _web_run_runtime()
    if recovered_outcome is None:
        try:
            outcome = await runtime.supervisor.wait(run_id)
        except asyncio.CancelledError:
            raise
        except Exception as exc:  # noqa: BLE001
            outcome = WorkerOutcome(
                run_id=run_id, job_kind="pipeline", returncode=1, ok=False,
                result=None, error=f"{type(exc).__name__}: {exc}", relayed_events=0,
            )
    else:
        outcome = recovered_outcome

    record = runtime.control_store.read(run_id)
    if record.state in {"cancelling", "cancelled"} or record.state != "completing":
        return

    result: RunResult | None = None
    artifact: Artifact | None = None
    if outcome.ok and outcome.result is not None:
        try:
            result = RunResult.model_validate(outcome.result)
        except ValidationError as exc:
            outcome = replace(outcome, ok=False, result=None, error=f"invalid RunResult: {exc}")
    if result is not None:
        run_dir = Path(result.run_dir).resolve()
        if result.run_id != run_id or run_dir != (RUNS_DIR / run_id).resolve():
            outcome = replace(
                outcome, ok=False, result=None,
                error="worker RunResult does not belong to the reserved run directory",
            )
            result = None
        else:
            result_type = _coerce_result_artifact_type(result, fallback=state.artifact_type)
            candidate = _build_artifact_response(
                run_dir, run_id, result_type,
                baseline_artifact_json=state.baseline_artifact_json,
            )
            artifact = candidate if _should_publish_artifact(result, candidate) else None

    if result is None:
        state.error = outcome.error or f"worker exited with status {outcome.returncode}"
        artifact = _recover_video_artifact_after_exception(
            run_id=run_id,
            a_type=state.artifact_type,
            baseline_artifact_json=state.baseline_artifact_json,
            error=state.error,
        )

    elapsed_ms = int((time.time() - started_at) * 1000)
    if artifact is not None:
        prepared_message = Message(
            id=f"msg_{run_id}", role="assistant",
            text=(
                _build_assistant_text(result, artifact, state.attach_paths)
                if result is not None
                else f"Recovered rendered video after a post-render error: {state.error}"
            ),
            ts=int(time.time() * 1000), run_id=run_id,
            artifact_id=artifact.artifact_id, status="done",
        )
    else:
        failure = (
            _failure_for_no_artifact(
                result=result,
                a_type=state.artifact_type,
                designer_model=state.designer_model,
                has_pdf=state.has_pdf,
                elapsed_ms=elapsed_ms,
            )
            if result is not None
            else _failure_from_disk(
                run_id=run_id,
                a_type=state.artifact_type,
                status="error",
                designer_model=state.designer_model,
                has_pdf=state.has_pdf,
                elapsed_ms=elapsed_ms,
            )
        )
        exit_failure = _worker_exit_failure_fields(outcome)
        if exit_failure and not failure.error_code:
            failure = failure.model_copy(update={
                "phase": exit_failure.get("phase") or failure.phase,
                "error_code": exit_failure.get("error_code"),
                "error_message": exit_failure.get("error_message"),
                "error_detail": exit_failure.get("error_detail"),
            })
        prepared_message = Message(
            id=f"msg_{run_id}", role="assistant",
            text=(
                _build_assistant_text(result, None, state.attach_paths)
                if result is not None
                else f"Run errored: {state.error or 'Worker did not produce a result.'}"
            ),
            ts=int(time.time() * 1000), run_id=run_id,
            status="error", failure=failure,
        )

    # The accepted terminal event can wake an SSE client immediately. Make the
    # matching artifact response observable first; quarantine still prevents
    # publication until the durable control CAS accepts completion.
    state.result_artifact = artifact
    state.result_message = prepared_message

    digest_payload: Any = outcome.result if outcome.result is not None else {
        "returncode": outcome.returncode,
        "error": outcome.error,
    }
    result_digest = hashlib.sha256(
        json.dumps(
            digest_payload, ensure_ascii=False, sort_keys=True,
            separators=(",", ":"),
        ).encode("utf-8")
    ).hexdigest()
    try:
        accepted = await runtime.supervisor.accept_completion(
            run_id,
            terminal_state="completed" if artifact is not None else "failed",
            publishable=artifact is not None,
            result_digest=result_digest,
        )
    except Exception as exc:  # completion may lose the race to cancellation
        current = runtime.control_store.read(run_id)
        state.result_artifact = None
        state.result_message = None
        if current.state in {"cancelling", "cancelled"}:
            return
        state.error = f"{type(exc).__name__}: {exc}"
        return

    expected_terminal = "completed" if artifact is not None else "failed"
    if accepted.state != expected_terminal:
        state.result_artifact = None
        state.result_message = None
        return
    state.result_artifact = artifact
    if artifact is not None:
        _promote_completed_run_reference_poster(run_id, state)
        _append_event(
            _settings_or_boot(), state.conversation_id, "artifact.generated",
            run_id=run_id, artifact_id=artifact.artifact_id,
            data={
                "artifact_type": artifact.artifact_type,
                "name": artifact.name,
                "native_format": artifact.native_format,
                "parent_artifact_id": artifact.parent_artifact_id,
                "terminal_status": result.terminal_status if result is not None else "pass",
                "critic_verdict": result.critic_verdict if result is not None else None,
                "critic_score": result.critic_score if result is not None else None,
                "canvas_plan": result.canvas_plan if result is not None else artifact.canvas_plan,
                "deck_plan": result.deck_plan if result is not None else artifact.deck_plan,
            },
        )
        if result is not None and result.style_snapshot:
            _append_event(
                _settings_or_boot(), state.conversation_id, "artifact.style_snapshot",
                run_id=run_id, artifact_id=artifact.artifact_id,
                data=result.style_snapshot,
            )
        await _reconcile_paper_bundle_for_run(run_id)
        return
    _append_event(
        _settings_or_boot(), state.conversation_id, "artifact.generation_failed",
        run_id=run_id,
        data={
            "status": result.terminal_status if result is not None else "error",
            "artifact_type": state.artifact_type,
            "error": (state.error or "")[:500],
            "produced_files": _list_produced_artifacts(run_id),
            "failure": _dump_model(failure),
            "canvas_plan": _read_canvas_plan(RUNS_DIR / run_id),
            "deck_plan": _read_deck_plan(RUNS_DIR / run_id),
        },
    )
    await _reconcile_paper_bundle_for_run(run_id)


async def _run_poster_code_edit_in_background(
    *,
    run_id: str,
    source_run_id: str,
    source_run_dir: Path,
    source_poster_path: Path,
    artifact: dict[str, Any],
    instruction: str,
    conversation_history: list[dict[str, Any]],
    selection_context: dict[str, Any] | None,
    required_color_system: dict[str, Any],
    baseline_artifact_json: str,
    state: _RunState,
    settings: Settings,
) -> None:
    started_at = time.time()
    try:
        result = await asyncio.to_thread(
            _run_poster_code_edit_sync,
            run_id=run_id,
            source_run_id=source_run_id,
            source_run_dir=source_run_dir,
            source_poster_path=source_poster_path,
            artifact=artifact,
            instruction=instruction,
            conversation_history=conversation_history,
            selection_context=selection_context,
            required_color_system=required_color_system,
            settings=settings,
        )
    except asyncio.CancelledError:
        log("web.code_editor.cancelled", run_id=run_id)
        state.error = state.error or "cancelled by user"
        failure = _failure_from_disk(
            run_id=run_id,
            a_type="poster",
            status="cancelled",
            designer_model=state.designer_model,
            has_pdf=False,
            elapsed_ms=int((time.time() - started_at) * 1000),
        )
        state.result_message = Message(
            id=f"msg_{run_id}",
            role="assistant",
            text="Poster revision cancelled.",
            ts=int(time.time() * 1000),
            run_id=run_id,
            status="error",
            failure=failure,
        )
        _append_event(
            settings,
            state.conversation_id,
            "artifact.generation_failed",
            run_id=run_id,
            data={
                "status": "cancelled",
                "artifact_type": "poster",
                "palette_id": state.palette_id,
                "required_color_system": required_color_system,
                "failure": _dump_model(failure),
            },
        )
        return
    except Exception as e:  # noqa: BLE001
        state.error = f"{type(e).__name__}: {e}"
        log("web.code_editor.error", run_id=run_id, error=type(e).__name__, msg=str(e)[:300])
        _persisted_run_log("run.error", run_id, msg=state.error[:200])
        failure = _failure_from_disk(
            run_id=run_id,
            a_type="poster",
            status="error",
            designer_model=state.designer_model,
            has_pdf=False,
            elapsed_ms=int((time.time() - started_at) * 1000),
        )
        state.result_message = Message(
            id=f"msg_{run_id}",
            role="assistant",
            text=f"Poster revision failed: {state.error}",
            ts=int(time.time() * 1000),
            run_id=run_id,
            status="error",
            failure=failure,
        )
        _append_event(
            settings,
            state.conversation_id,
            "artifact.generation_failed",
            run_id=run_id,
            data={
                "status": "error",
                "artifact_type": "poster",
                "palette_id": state.palette_id,
                "required_color_system": required_color_system,
                "error": state.error[:500],
                "failure": _dump_model(failure),
            },
        )
        return

    run_dir = Path(result["run_dir"])
    artifact_response = _build_artifact_response(
        run_dir,
        run_id,
        "poster",
        baseline_artifact_json=baseline_artifact_json,
    )
    state.result_artifact = artifact_response
    state.result_message = Message(
        id=f"msg_{run_id}",
        role="assistant",
        text="Applied the poster revision with the external code editor.",
        ts=int(time.time() * 1000),
        run_id=run_id,
        artifact_id=artifact_response.artifact_id if artifact_response else None,
        status="done" if artifact_response else "error",
    )
    if artifact_response is None:
        state.error = "poster revision completed without a publishable artifact"
        _persisted_run_log("run.error", run_id, msg=state.error)
        return

    _append_event(
        settings,
        state.conversation_id,
        "artifact.generated",
        run_id=run_id,
        artifact_id=artifact_response.artifact_id,
        data={
            "artifact_type": "poster",
            "name": artifact_response.name,
            "native_format": artifact_response.native_format,
            "parent_artifact_id": artifact_response.parent_artifact_id,
            "render_mode": "external_code_editor_revision",
            "source_run_id": source_run_id,
            "palette_id": state.palette_id,
            "required_color_system": required_color_system,
            "attempts": result.get("attempts") or [],
            "validation_summary": result.get("validation_summary") or {},
            "selection_context_summary": result.get("selection_context_summary"),
            "canvas_plan": artifact_response.canvas_plan,
            "deck_plan": artifact_response.deck_plan,
        },
    )
    _persisted_run_log("run.done", run_id)


def _run_poster_code_edit_sync(
    *,
    run_id: str,
    source_run_id: str,
    source_run_dir: Path,
    source_poster_path: Path,
    artifact: dict[str, Any],
    instruction: str,
    conversation_history: list[dict[str, Any]],
    selection_context: dict[str, Any] | None,
    required_color_system: dict[str, Any],
    settings: Settings,
    cancellation_token: Any = None,
) -> dict[str, Any]:
    return _run_poster_code_edit_core(
        run_id=run_id,
        runs_dir=settings.out_dir / "runs",
        source_run_id=source_run_id,
        source_run_dir=source_run_dir,
        source_poster_path=source_poster_path,
        artifact=artifact,
        instruction=instruction,
        conversation_history=conversation_history,
        selection_context=selection_context,
        required_color_system=required_color_system,
        settings=settings,
        cancellation_token=cancellation_token,
    )


def _openresearch_file_url(source_run_id: str, job_id: str, name: str) -> str:
    return _run_file_url(source_run_id, f"openresearch/{job_id}/{name}")


@app.post("/api/openresearch/projects", response_model=OpenResearchProjectAck)
async def openresearch_projects(
    req: OpenResearchProjectRequest,
    request: Request,
) -> OpenResearchProjectAck:
    if _DEMO_MODE:
        raise HTTPException(
            403,
            detail={
                "code": "demo_openresearch_disabled",
                "message": "OpenResearch submission is disabled in demo mode.",
            },
        )
    artifact = req.artifact if isinstance(req.artifact, dict) else {}
    if artifact:
        artifact_type = str(artifact.get("artifact_type") or "")
        if artifact_type != "poster":
            raise HTTPException(400, detail="OpenResearch projects are only supported for poster artifacts")
    artifact_id = (req.artifact_id or str(artifact.get("artifact_id") or "")).strip()
    source_run_id = (req.source_run_id or _run_id_from_artifact_id(artifact_id) or "").strip()
    if not source_run_id:
        raise HTTPException(400, detail="missing source_run_id or poster artifact_id")
    artifact_id = artifact_id or f"art_{source_run_id}"
    _demo_require_run_owner(
        source_run_id,
        request,
        detail=f"source run not found: {source_run_id}",
    )
    _assert_controlled_run_source_usable(source_run_id, mode="artifact")
    run_dir = RUNS_DIR / source_run_id
    if not run_dir.exists():
        raise HTTPException(404, detail=f"source run not found: {source_run_id}")

    settings = _settings_for_openresearch_request(request)
    conversation_id = str(req.conversation_id or "").strip()
    job_id = new_run_id()
    state = _OpenResearchJobState(
        job_id=job_id,
        source_run_id=source_run_id,
        artifact_id=artifact_id,
        conversation_id=conversation_id,
        request=req,
    )
    async with _OPENRESEARCH_JOBS_LOCK:
        _OPENRESEARCH_JOBS[job_id] = state
    state.task = asyncio.create_task(
        _submit_openresearch_project_in_background(state=state, settings=settings)
    )
    _append_event(
        settings,
        conversation_id,
        "openresearch.project_requested",
        run_id=source_run_id,
        artifact_id=artifact_id,
        data={"job_id": job_id, "status": "running"},
    )
    log("openresearch.project_start", run_id=job_id, source_run_id=source_run_id)
    return OpenResearchProjectAck(job_id=job_id, status="running")


@app.get("/api/openresearch/projects/{job_id}", response_model=OpenResearchProjectResult)
async def openresearch_project(job_id: str, request: Request) -> OpenResearchProjectResult:
    if _DEMO_MODE:
        raise HTTPException(
            403,
            detail={
                "code": "demo_openresearch_disabled",
                "message": "OpenResearch submission is disabled in demo mode.",
            },
        )
    async with _OPENRESEARCH_JOBS_LOCK:
        state = _OPENRESEARCH_JOBS.get(job_id)
    if state is None:
        recovered = _read_openresearch_project_result_by_job(job_id)
        if recovered is None:
            raise HTTPException(404, detail=f"OpenResearch project job not found: {job_id}")
        _demo_require_run_owner(recovered.source_run_id, request, detail=f"OpenResearch project job not found: {job_id}")
        settings = _settings_for_openresearch_request(request)
        return await asyncio.to_thread(_refresh_openresearch_project_result, recovered, settings)
    _demo_require_run_owner(state.source_run_id, request, detail=f"OpenResearch project job not found: {job_id}")
    if state.task and not state.task.done():
        try:
            await asyncio.wait_for(asyncio.shield(state.task), timeout=120)
        except asyncio.TimeoutError:
            return OpenResearchProjectResult(
                job_id=job_id,
                source_run_id=state.source_run_id,
                artifact_id=state.artifact_id,
                status="running",
            )
    if isinstance(state.result, OpenResearchProjectResult):
        settings = _settings_for_openresearch_request(request)
        return await asyncio.to_thread(_refresh_openresearch_project_result, state.result, settings)
    if state.result is not None:
        raise HTTPException(500, detail="OpenResearch job type mismatch")
    raise HTTPException(500, detail=state.error or "OpenResearch project completed without a result")


async def _submit_openresearch_project_in_background(
    *,
    state: _OpenResearchJobState,
    settings: Settings,
) -> None:
    try:
        state.result = await asyncio.to_thread(_submit_openresearch_project_sync, state, settings)
        if isinstance(state.result, OpenResearchProjectResult):
            event_name = (
                "openresearch.project_failed"
                if state.result.status == "error"
                else "openresearch.project_ready"
            )
            _append_event(
                settings,
                state.conversation_id,
                event_name,
                run_id=state.source_run_id,
                artifact_id=state.artifact_id,
                data=state.result.model_dump(mode="json"),
            )
            terminal_event = "openresearch.error" if state.result.status == "error" else "openresearch.done"
            log(terminal_event, run_id=state.job_id, source_run_id=state.source_run_id, status=state.result.status)
    except Exception as exc:  # noqa: BLE001
        state.error = f"{type(exc).__name__}: {exc}"
        result = OpenResearchProjectResult(
            job_id=state.job_id,
            source_run_id=state.source_run_id,
            artifact_id=state.artifact_id,
            status="error",
            error=state.error,
            result_url=_openresearch_file_url(state.source_run_id, state.job_id, "openresearch_project_result.json"),
            api_log_url=_openresearch_file_url(state.source_run_id, state.job_id, "openresearch_api.jsonl"),
            agent_prompt_url=_openresearch_file_url(state.source_run_id, state.job_id, OPENRESEARCH_AGENT_PROMPT_FILE),
            submitter_log_url=_openresearch_file_url(state.source_run_id, state.job_id, OPENRESEARCH_GUI_PROCESS_FILE),
        )
        state.result = result
        _write_openresearch_project_result(state.source_run_id, state.job_id, result)
        _append_event(
            settings,
            state.conversation_id,
            "openresearch.project_failed",
            run_id=state.source_run_id,
            artifact_id=state.artifact_id,
            data=result.model_dump(mode="json"),
        )
        log("openresearch.error", run_id=state.job_id, error=state.error[:300])


def _submit_openresearch_project_sync(
    state: _OpenResearchJobState,
    settings: Settings,
) -> OpenResearchProjectResult:
    req = state.request
    if not isinstance(req, OpenResearchProjectRequest):
        raise TypeError("OpenResearch project job received the wrong request type")
    run_dir = RUNS_DIR / state.source_run_id
    job_dir = run_dir / "openresearch" / state.job_id
    job_dir.mkdir(parents=True, exist_ok=True)
    result_url = _openresearch_file_url(state.source_run_id, state.job_id, "openresearch_project_result.json")
    artifact = req.artifact if isinstance(req.artifact, dict) else {}
    paper_id = _normalize_openresearch_paper_id(
        req.paper_id
        or req.paper_url
        or _infer_openresearch_paper_id(run_dir, artifact)
    )
    repo_full_name = str(
        req.repo_full_name
        or getattr(settings, "openresearch_default_repo_full_name", "")
        or ""
    ).strip() or None
    base = OpenResearchProjectResult(
        job_id=state.job_id,
        source_run_id=state.source_run_id,
        artifact_id=state.artifact_id,
        status="running",
        paper_id=paper_id,
        repo_full_name=repo_full_name,
        result_url=result_url,
        agent_prompt_url=_openresearch_file_url(state.source_run_id, state.job_id, OPENRESEARCH_AGENT_PROMPT_FILE),
        submitter_log_url=_openresearch_file_url(state.source_run_id, state.job_id, OPENRESEARCH_GUI_PROCESS_FILE),
        details={"kind": "openresearch_gui_agent_submission"},
    )

    org_id = str(req.org_id or getattr(settings, "openresearch_org_id", "") or "").strip()
    base.org_id = org_id or None
    if not paper_id:
        base.status = "error"
        base.error = "missing OpenResearch paper id or arXiv URL"
        _write_openresearch_project_result(state.source_run_id, state.job_id, base)
        log("openresearch.project_error", run_id=state.job_id, error=base.error)
        return base
    request_payload = {
        "org_id": org_id,
        "openresearch_url": _openresearch_entry_url(org_id),
        "paper_id": paper_id,
        "paper_url": req.paper_url,
        "repo_full_name": repo_full_name,
        "user_prompt": req.agent_prompt,
    }
    atomic_write_json(job_dir / "openresearch_project_request.json", request_payload)

    agent_prompt = str(req.agent_prompt or "").strip() or _default_openresearch_agent_prompt(
        paper_id=paper_id,
        repo_full_name=repo_full_name,
    )
    submitter = submit_openresearch_gui(
        settings=settings,
        job_dir=job_dir,
        project_url=_openresearch_entry_url(org_id),
        agent_prompt=agent_prompt,
        project_id=None,
        org_id=base.org_id,
        paper_id=paper_id,
        paper_url=req.paper_url,
        repo_full_name=repo_full_name,
        source_run_id=state.source_run_id,
        artifact_id=state.artifact_id,
    )
    base.details["gui_submitter"] = submitter.to_dict()
    base.gui_submitter_status = submitter.status
    base.gui_submitter_reason = submitter.reason
    base.gui_submitter_error = submitter.error
    base.gui_submitter_session_url = submitter.session_url
    base.project_url = submitter.project_url or _openresearch_project_url_from_text(submitter.session_url)
    base.project_id = _openresearch_project_id_from_url(base.project_url)
    if submitter.status != "submitted":
        base.status = "error"
        base.error = submitter.error or submitter.reason or "failed to submit prompt to OpenResearch Auto Research"
        _write_openresearch_project_result(state.source_run_id, state.job_id, base)
        log("openresearch.project_error", run_id=state.job_id, error=base.error[:300])
        return base

    base.status = "submitted"
    if base.project_id:
        base = _refresh_openresearch_project_result(base, settings)
    _write_openresearch_project_result(state.source_run_id, state.job_id, base)
    log("openresearch.project_submitted", run_id=state.job_id, project_id=base.project_id)
    return base


def _refresh_openresearch_project_result(
    result: OpenResearchProjectResult,
    settings: Settings,
    *,
    client: OpenResearchApiClient | None = None,
) -> OpenResearchProjectResult:
    if result.status == "error" or not result.project_id:
        return result
    job_dir = RUNS_DIR / result.source_run_id / "openresearch" / result.job_id
    job_dir.mkdir(parents=True, exist_ok=True)
    client = client or OpenResearchApiClient(
        api_url=getattr(settings, "openresearch_api_url", "https://api.openresearch.sh"),
        token=getattr(settings, "openresearch_token", ""),
        timeout_s=getattr(settings, "openresearch_timeout_s", 120),
        request_log_path=job_dir / "openresearch_api.jsonl",
        allow_private_network=bool(
            getattr(settings, "allow_private_network", True)
        ),
    )
    project = client.get_project(result.project_id)
    result.details["project_status"] = project.to_dict()
    project_data = _openresearch_project_data(project.data)
    if project.ok and project_data:
        result.paper_id = str(project_data.get("paperId") or result.paper_id or "").strip() or None
        result.repo_full_name = _openresearch_repo_full_name(project_data) or result.repo_full_name

    reports = client.list_reports(result.project_id)
    result.details["reports"] = reports.to_dict()
    if reports.ok:
        latest = _latest_openresearch_report(reports.data)
        if latest:
            report_id = str(latest.get("id") or latest.get("reportId") or "").strip()
            result.latest_report_id = report_id or result.latest_report_id
            if result.project_url and result.latest_report_id:
                result.latest_report_url = f"{result.project_url}/reports/{result.latest_report_id}"
            if result.latest_report_id:
                report = client.get_report(result.project_id, result.latest_report_id)
                result.details["latest_report"] = report.to_dict()
                markdown = _openresearch_report_markdown(report.data)
                if markdown:
                    result.latest_report_markdown = markdown
    _write_openresearch_project_result(result.source_run_id, result.job_id, result)
    return result


def _write_openresearch_project_result(
    source_run_id: str,
    job_id: str,
    result: OpenResearchProjectResult,
) -> None:
    path = RUNS_DIR / source_run_id / "openresearch" / job_id / "openresearch_project_result.json"
    atomic_write_json(path, result.model_dump(mode="json"))


def _read_openresearch_project_result_by_job(job_id: str) -> OpenResearchProjectResult | None:
    if not RUNS_DIR.exists():
        return None
    for path in RUNS_DIR.glob(f"*/openresearch/{job_id}/openresearch_project_result.json"):
        payload = _read_json_file(path)
        if not isinstance(payload, dict):
            continue
        source_run_id = str(payload.get("source_run_id") or path.parents[2].name)
        payload.setdefault("job_id", job_id)
        payload.setdefault("source_run_id", source_run_id)
        payload.setdefault("artifact_id", f"art_{source_run_id}")
        payload.setdefault("status", "submitted")
        payload.setdefault("result_url", _openresearch_file_url(source_run_id, job_id, "openresearch_project_result.json"))
        payload.setdefault("api_log_url", _openresearch_file_url(source_run_id, job_id, "openresearch_api.jsonl"))
        payload.setdefault("agent_prompt_url", _openresearch_file_url(source_run_id, job_id, OPENRESEARCH_AGENT_PROMPT_FILE))
        payload.setdefault("submitter_log_url", _openresearch_file_url(source_run_id, job_id, OPENRESEARCH_GUI_PROCESS_FILE))
        try:
            return OpenResearchProjectResult.model_validate(payload)
        except ValidationError:
            return None
    return None


_OPENRESEARCH_ARXIV_RE = re.compile(r"(?<!\d)(\d{4}\.\d{4,5})(?:v\d+)?")


def _normalize_openresearch_paper_id(value: str | None) -> str | None:
    text = str(value or "").strip()
    if not text:
        return None
    arxiv_id = _extract_openresearch_arxiv_id(text)
    if arxiv_id:
        return arxiv_id
    return text


def _infer_openresearch_paper_id(run_dir: Path, artifact: dict[str, Any]) -> str | None:
    return _infer_openresearch_paper_id_from_run(run_dir, artifact, seen=set())


def _infer_openresearch_paper_id_from_run(
    run_dir: Path,
    artifact: dict[str, Any],
    *,
    seen: set[str],
) -> str | None:
    run_key = str(run_dir.resolve())
    if run_key in seen:
        return None
    seen.add(run_key)
    for text in _iter_openresearch_strings(artifact):
        paper_id = _extract_openresearch_arxiv_id(text)
        if paper_id:
            return paper_id
    for name in (
        "paper_resource_manifest.json",
        "poster_content_brief.json",
        "paper_memory.json",
        "paper_memory.md",
        "paper_memory_dossier.md",
    ):
        path = run_dir / name
        if not path.exists():
            continue
        try:
            text = path.read_text(encoding="utf-8", errors="replace")[:200000]
        except OSError:
            continue
        paper_id = _extract_openresearch_arxiv_id(text)
        if paper_id:
            return paper_id
    prior_paper_id = _openresearch_paper_id_from_prior_jobs(run_dir)
    if prior_paper_id:
        return prior_paper_id
    parent_run_id = _openresearch_parent_run_id(run_dir, artifact)
    if parent_run_id:
        parent_dir = RUNS_DIR / parent_run_id
        if parent_dir.exists():
            return _infer_openresearch_paper_id_from_run(
                parent_dir,
                {},
                seen=seen,
            )
    return None


def _openresearch_paper_id_from_prior_jobs(run_dir: Path) -> str | None:
    root = run_dir / "openresearch"
    if not root.exists():
        return None
    paths = sorted(
        [
            *root.glob("*/openresearch_project_result.json"),
            *root.glob("*/openresearch_project_request.json"),
        ],
        key=lambda p: p.stat().st_mtime if p.exists() else 0,
        reverse=True,
    )
    for path in paths:
        payload = _read_json_file(path)
        if not isinstance(payload, dict):
            continue
        paper_id = _normalize_openresearch_paper_id(
            payload.get("paper_id") or payload.get("paper_url")
        )
        if paper_id:
            return paper_id
    return None


def _openresearch_parent_run_id(run_dir: Path, artifact: dict[str, Any]) -> str | None:
    for raw in (artifact.get("parent_run_id"), artifact.get("parent_artifact_id")):
        if isinstance(raw, str) and raw.strip():
            value = raw.strip()
            return value[len("art_"):] if value.startswith("art_") else value
    for name in ("code_editor_revision_manifest.json", "authored_poster_edit_manifest.json"):
        manifest = _read_json_file(run_dir / "final" / name)
        if isinstance(manifest, dict):
            raw = manifest.get("parent_run_id") or manifest.get("parent_artifact_id")
            parent = _run_id_from_maybe_artifact_ref(raw)
            if parent:
                return parent
    return None


def _extract_openresearch_arxiv_id(text: str) -> str | None:
    match = _OPENRESEARCH_ARXIV_RE.search(str(text or ""))
    return match.group(1) if match else None


def _iter_openresearch_strings(value: Any) -> list[str]:
    out: list[str] = []
    stack = [value]
    while stack and len(out) < 200:
        item = stack.pop()
        if isinstance(item, str):
            out.append(item)
        elif isinstance(item, dict):
            stack.extend(item.values())
        elif isinstance(item, list):
            stack.extend(item)
    return out


def _default_openresearch_agent_prompt(
    *,
    paper_id: str | None,
    repo_full_name: str | None,
) -> str:
    parts = [
        "Reproduce the code and experiments for this paper, then write a formal reproduction report.",
    ]
    if paper_id:
        parts.append(f"Paper: {paper_id}.")
    if repo_full_name:
        parts.append(f"Starting repository: {repo_full_name}.")
    return "\n".join(parts)


def _openresearch_entry_url(org_id: str | None) -> str:
    if org_id:
        return f"https://openresearch.sh/orgs/{org_id}/projects"
    return "https://openresearch.sh/"


def _openresearch_project_url(org_id: str | None, project_id: str | None) -> str | None:
    if not org_id or not project_id:
        return None
    return f"https://openresearch.sh/orgs/{org_id}/projects/{project_id}"


def _openresearch_project_url_from_text(text: str | None) -> str | None:
    value = str(text or "").strip()
    if "/projects/" not in value:
        return None
    match = re.search(r"https?://openresearch\.sh/orgs/[^/]+/projects/[0-9a-f-]+", value)
    return match.group(0) if match else value


def _openresearch_project_id_from_url(url: str | None) -> str | None:
    value = str(url or "")
    match = re.search(r"/projects/([0-9a-f-]+)", value)
    return match.group(1) if match else None


def _openresearch_items(value: Any, key: str) -> list[dict[str, Any]]:
    if isinstance(value, list):
        return [item for item in value if isinstance(item, dict)]
    if isinstance(value, dict):
        nested = value.get(key)
        if isinstance(nested, list):
            return [item for item in nested if isinstance(item, dict)]
        for fallback in ("items", "data", "results"):
            nested = value.get(fallback)
            if isinstance(nested, list):
                return [item for item in nested if isinstance(item, dict)]
    return []


def _openresearch_project_data(value: Any) -> dict[str, Any]:
    if isinstance(value, dict):
        nested = value.get("project")
        if isinstance(nested, dict):
            return nested
        return value
    return {}


def _openresearch_repo_full_name(project: dict[str, Any]) -> str | None:
    direct = str(project.get("repoFullName") or "").strip()
    if direct:
        return direct
    owner = str(project.get("githubOwner") or "").strip()
    repo = str(project.get("githubRepo") or "").strip()
    if owner and repo:
        return f"{owner}/{repo}"
    return None


def _latest_openresearch_report(value: Any) -> dict[str, Any] | None:
    items = _openresearch_items(value, "reports")
    if not items:
        return None
    return sorted(
        items,
        key=lambda item: str(item.get("updatedAt") or item.get("createdAt") or item.get("id") or ""),
        reverse=True,
    )[0]


def _openresearch_report_markdown(value: Any) -> str | None:
    if isinstance(value, str):
        return value
    if isinstance(value, dict):
        for key in ("markdown", "content", "body", "report"):
            candidate = value.get(key)
            if isinstance(candidate, str) and candidate.strip():
                return candidate
    return None


def _attempt_candidate_payload(run_id: str, attempt: int) -> dict[str, Any]:
    candidate = load_attempt_candidate(RUNS_DIR / run_id, attempt)
    payload = candidate_summary(candidate)
    source_relative = str(payload.pop("source_relative_path"))
    preview_relatives = [
        str(value) for value in payload.pop("preview_relative_paths", [])
    ]
    payload.pop("dependency_relative_paths", None)
    payload.pop("browser_resource_relative_paths", None)
    payload.pop("validation_summary_relative_path", None)
    payload["source_url"] = _run_file_url(run_id, source_relative)
    payload["preview_urls"] = [
        _run_file_url(run_id, relative) for relative in preview_relatives
    ]
    return payload


def _attempt_state_payload(run_id: str) -> dict[str, Any]:
    run_dir = RUNS_DIR / run_id
    index = load_candidate_index(run_dir)
    candidates: list[dict[str, Any]] = []
    if index is not None:
        for manifest_relative in index.manifest_relative_paths:
            manifest_path = run_dir / manifest_relative
            try:
                raw = json.loads(manifest_path.read_text(encoding="utf-8"))
                attempt = int(raw.get("attempt"))
                candidates.append(_attempt_candidate_payload(run_id, attempt))
            except (OSError, TypeError, ValueError, json.JSONDecodeError):
                continue
    journal = load_selection_journal(run_dir)
    return {
        "run_id": run_id,
        "candidates": sorted(candidates, key=lambda item: int(item["attempt"])),
        "selection": journal.model_dump(mode="json") if journal else None,
    }


def _complete_attempt_selection_sync(run_id: str, settings: Settings) -> None:
    run_dir = RUNS_DIR / run_id
    cancellation_token = CancellationToken.for_run(
        RunControlStore(RUNS_DIR),
        run_id,
    )
    cancellation_token.raise_if_cancelled(
        "attempt_selection.web.before_recovery",
    )
    ctx = ToolContext(
        settings=settings,
        run_dir=run_dir,
        layers_dir=run_dir / "layers",
        run_id=run_id,
        cancellation_token=cancellation_token,
    )
    promote_pending_selection(ctx)
    cancellation_token.raise_if_cancelled(
        "attempt_selection.web.after_recovery",
    )


async def _complete_attempt_selection_in_background(
    run_id: str,
    settings: Settings,
) -> None:
    try:
        await asyncio.to_thread(
            _complete_attempt_selection_sync,
            run_id,
            settings,
        )
    except RunCancelled:
        return
    finally:
        current = asyncio.current_task()
        if _ATTEMPT_SELECTION_TASKS.get(run_id) is current:
            _ATTEMPT_SELECTION_TASKS.pop(run_id, None)


def _schedule_attempt_selection_recovery(
    run_id: str,
    settings: Settings,
) -> None:
    existing = _ATTEMPT_SELECTION_TASKS.get(run_id)
    if existing is not None and not existing.done():
        return
    _ATTEMPT_SELECTION_TASKS[run_id] = asyncio.create_task(
        _complete_attempt_selection_in_background(run_id, settings)
    )


def _failed_attempt_selection_error(
    run_id: str,
    attempt: int,
    expected_candidate_sha256: str,
) -> HTTPException | None:
    try:
        record = RunControlStore(RUNS_DIR).read(run_id)
    except RunControlError:
        return None
    if record.state != "failed":
        return None
    try:
        candidate = load_attempt_candidate(RUNS_DIR / run_id, attempt)
    except ValueError:
        return None
    if candidate.source_sha256 != expected_candidate_sha256:
        return HTTPException(409, detail={"code": "candidate_changed"})
    if candidate.safety_state == "blocked":
        return HTTPException(422, detail={"code": "candidate_blocked"})
    return HTTPException(409, detail={"code": "run_not_selectable"})


@app.get("/api/runs/{run_id}/attempts")
async def run_attempts(run_id: str, request: Request) -> dict[str, Any]:
    _demo_require_run_owner(run_id, request)
    run_dir = RUNS_DIR / run_id
    if not run_dir.is_dir():
        raise HTTPException(404, detail=f"run not found: {run_id}")
    try:
        journal = await asyncio.to_thread(load_selection_journal, run_dir)
    except (OSError, ValueError, ValidationError):
        journal = None
    if (
        journal is not None
        and not _controlled_run_is_cancellation_frozen(run_id)
        and journal.state in {
            "requested",
            "terminating",
            "promoting",
            "delivering",
        }
    ):
        try:
            recovery_settings = _settings_for_request(request)
        except HTTPException:
            recovery_settings = None
        if recovery_settings is not None:
            _schedule_attempt_selection_recovery(run_id, recovery_settings)
    return await asyncio.to_thread(_attempt_state_payload, run_id)


@app.post("/api/runs/{run_id}/attempts/{attempt}/select", status_code=202)
async def select_run_attempt(
    run_id: str,
    attempt: int,
    req: AttemptSelectionRequest,
    request: Request,
) -> dict[str, Any]:
    async with _web_run_operation_guard(f"run:{run_id}:attempt-selection"):
        try:
            validate_run_id(run_id)
        except RunControlError as exc:
            raise HTTPException(404, detail="run not found") from exc
        _demo_require_run_owner(run_id, request)
        run_dir = RUNS_DIR / run_id
        if not run_dir.is_dir():
            raise HTTPException(404, detail=f"run not found: {run_id}")
        try:
            _assert_controlled_run_source_usable(run_id, mode="mutation")
        except HTTPException as exc:
            failed_error = _failed_attempt_selection_error(
                run_id,
                attempt,
                req.expected_candidate_sha256,
            )
            if failed_error is not None:
                raise failed_error from exc
            raise
        control_store = RunControlStore(RUNS_DIR)
        result = await asyncio.to_thread(
            request_attempt_selection,
            run_dir=run_dir,
            run_id=run_id,
            attempt=attempt,
            expected_candidate_sha256=req.expected_candidate_sha256,
            idempotency_key=req.idempotency_key,
            writable_guard=lambda: control_store.assert_writable(run_id),
        )
        if result.status == "candidate_blocked":
            raise HTTPException(422, detail={"code": result.status})
        if result.status == "candidate_changed":
            raise HTTPException(409, detail={"code": result.status})
        if result.status == "run_not_selectable":
            raise HTTPException(409, detail={"code": result.status})
        settings = _settings_for_request(request)
        _schedule_attempt_selection_recovery(run_id, settings)
        state = await asyncio.to_thread(_attempt_state_payload, run_id)
        return {
            "status": result.status,
            "candidate_id": result.candidate_id,
            **state,
        }


def _direct_candidate_publish_ack(
    run_id: str,
    *,
    start_token: str | None,
) -> GenerateAck:
    return GenerateAck(
        run_id=run_id,
        progress_mode="attempt_publish",
        start_token=start_token,
        placeholder_message=Message(
            id=f"msg_{run_id}",
            role="assistant",
            text="",
            ts=int(time.time() * 1000),
            run_id=run_id,
            status="streaming",
        ),
    )


async def _cleanup_failed_candidate_publish_setup(
    *,
    run_id: str,
    state: _RunState | None,
    access_lease: _DemoDerivedRunAccessLease,
    durable_handoff: bool,
    reservation_claimed: bool,
    discard_write_ahead: bool,
) -> None:
    task = asyncio.create_task(
        _cleanup_failed_candidate_publish_setup_owned(
            run_id=run_id,
            state=state,
            access_lease=access_lease,
            durable_handoff=durable_handoff,
            reservation_claimed=reservation_claimed,
            discard_write_ahead=discard_write_ahead,
        )
    )
    await _join_owned_task(task)
    task.result()


async def _cleanup_failed_candidate_publish_setup_owned(
    *,
    run_id: str,
    state: _RunState | None,
    access_lease: _DemoDerivedRunAccessLease,
    durable_handoff: bool,
    reservation_claimed: bool,
    discard_write_ahead: bool,
) -> None:
    control_exists = (RUNS_DIR / run_id / "run_control.json").is_file()
    cancellation_confirmed = not control_exists or not reservation_claimed
    if control_exists and reservation_claimed:
        try:
            cancel_task = asyncio.create_task(
                _web_run_runtime().services.cancel(
                    run_id,
                    "candidate_publish_setup_failed",
                )
            )
            await _join_owned_task(cancel_task)
            cancellation = cancel_task.result()
            if getattr(cancellation, "cancel_request_event_required", False):
                _append_cancel_request_event(
                    run_id,
                    "candidate_publish_setup_failed",
                )
            cancellation_confirmed = cancellation.confirmed
        except BaseException as exc:  # preserve ownership until durable recovery
            log(
                "candidate_publish.setup_cleanup_deferred",
                run_id=run_id,
                error=type(exc).__name__,
            )
            return
    if not cancellation_confirmed:
        return
    if state is not None:
        async with _RUNS_LOCK:
            if _RUNS.get(run_id) is state:
                _RUNS.pop(run_id, None)
    if not durable_handoff:
        _demo_release_derived_run_access(access_lease)
        if not control_exists and discard_write_ahead:
            run_dir = RUNS_DIR / run_id
            for name in (
                "candidate_publish_request.json",
                "derived_job.json",
            ):
                (run_dir / name).unlink(missing_ok=True)
            try:
                run_dir.rmdir()
            except OSError:
                pass


@app.post(
    "/api/runs/{run_id}/attempts/{attempt}/publish",
    response_model=GenerateAck,
)
async def publish_run_attempt(
    run_id: str,
    attempt: int,
    req: DirectAttemptPublishRequest,
    request: Request,
) -> GenerateAck:
    _demo_require_run_owner(run_id, request)
    if re.fullmatch(r"[0-9a-f]{64}", req.expected_candidate_sha256) is None:
        raise HTTPException(
            422,
            detail="expected_candidate_sha256 must be a lowercase SHA-256 digest",
        )
    owner_id = _run_owner_id(request)
    idempotency_key = req.idempotency_key.strip()
    if not idempotency_key:
        raise HTTPException(422, detail="idempotency_key must not be blank")
    idempotency_key_digest = hashlib.sha256(
        json.dumps(
            [owner_id, idempotency_key],
            ensure_ascii=False,
            separators=(",", ":"),
        ).encode("utf-8")
    ).hexdigest()
    request_digest = hashlib.sha256(
        json.dumps(
            {
                "owner": owner_id,
                "source_run_id": run_id,
                "source_attempt": attempt,
                "source_candidate_sha256": req.expected_candidate_sha256,
                "conversation_id": str(req.conversation_id or ""),
            },
            ensure_ascii=False,
            sort_keys=True,
            separators=(",", ":"),
        ).encode("utf-8")
    ).hexdigest()
    published_run_id = f"candidate-publish-{idempotency_key_digest[:32]}"
    reserve_only = _request_reserve_only(request)

    async with _run_tree_lock(f"candidate-publish:{published_run_id}"):
        write_ahead_request: dict[str, Any] | None = None
        try:
            existing = _read_direct_candidate_publish_descriptor(published_run_id)
        except ValueError as exc:
            raise HTTPException(409, detail=str(exc)) from exc
        if existing is not None:
            if (
                existing["version"] not in {1, 2}
                or existing["idempotency_key_digest"] != idempotency_key_digest
                or existing["request_digest"] != request_digest
                or existing["source_run_id"] != run_id
                or existing["source_attempt"] != attempt
                or existing["source_candidate_sha256"]
                != req.expected_candidate_sha256
            ):
                raise HTTPException(
                    409,
                    detail="idempotency key was already used for another attempt publication",
                )
            try:
                child_control = RunControlStore(RUNS_DIR).read(published_run_id)
            except (RunControlError, OSError, json.JSONDecodeError) as exc:
                if (RUNS_DIR / published_run_id / "run_control.json").exists():
                    raise HTTPException(
                        409,
                        detail="attempt publication lifecycle is unavailable",
                    ) from exc
                write_ahead_request = existing
                child_control = None
            start_token: str | None = None
            if child_control is not None and child_control.state == "queued":
                async with _RUNS_LOCK:
                    existing_state = _RUNS.get(published_run_id)
                durable_descriptor: dict[str, Any] | None = None
                if (
                    existing_state is None
                    or not existing_state.reservation_token
                ):
                    try:
                        (
                            existing_state,
                            recovered_token,
                            durable_descriptor,
                        ) = await _recover_queued_candidate_publish_context(
                            published_run_id,
                            _settings_for_request(request),
                        )
                    except Exception as exc:  # noqa: BLE001
                        raise HTTPException(
                            409,
                            detail="attempt publication recovery is unavailable",
                        ) from exc
                    start_token = recovered_token
                else:
                    start_token = existing_state.reservation_token
                if not reserve_only:
                    if durable_descriptor is None:
                        durable_descriptor = _read_derived_job_descriptor(
                            published_run_id
                        )
                    if durable_descriptor is None:
                        raise HTTPException(
                            409,
                            detail="attempt publication descriptor is unavailable",
                        )
                    try:
                        await _start_reserved_derived_job(
                            run_id=published_run_id,
                            token=start_token,
                            state=existing_state,
                            descriptor=durable_descriptor,
                        )
                    except Exception as exc:
                        raise _web_run_service_error(exc) from exc
                    start_token = None
            if child_control is not None:
                return _direct_candidate_publish_ack(
                    published_run_id,
                    start_token=start_token if reserve_only else None,
                )

        _assert_controlled_run_source_usable(run_id, mode="snapshot")
        if _candidate_publish_source_disposition(run_id) == "cancelled":
            raise HTTPException(410, detail="cancelled run output is diagnostic only")
        source_run_dir = RUNS_DIR / run_id
        if not source_run_dir.is_dir():
            raise HTTPException(404, detail=f"run not found: {run_id}")
        try:
            candidate = load_attempt_candidate(source_run_dir, attempt)
        except ValueError as exc:
            raise HTTPException(404, detail=str(exc)) from exc
        if candidate.source_sha256 != req.expected_candidate_sha256:
            raise HTTPException(409, detail={"code": "candidate_changed"})
        if candidate.safety_state == "blocked":
            raise HTTPException(422, detail={"code": "candidate_blocked"})
        if write_ahead_request is not None:
            try:
                await _validate_replayed_candidate_publish_binding_async(
                    source_run_id=run_id,
                    artifact_type=candidate.artifact_type.value,
                    candidate_id=candidate.candidate_id,
                    owner_id=owner_id,
                    write_ahead_request=write_ahead_request,
                )
            except PaperBundleError as exc:
                raise HTTPException(409, detail=str(exc)) from exc

        access_lease = _demo_acquire_derived_run_access_lease(
            published_run_id,
            request,
            parent_run_id=run_id,
        )
        state: _RunState | None = None
        durable_handoff = False
        reservation_claimed = False
        try:
            if write_ahead_request is None:
                try:
                    bundle_binding = await _reserve_candidate_publish_bundle_binding_async(
                        run_id,
                        candidate.artifact_type.value,
                        owner_id,
                    )
                except PaperBundleError as exc:
                    raise _paper_bundle_error(exc) from exc
            elif write_ahead_request["version"] == 2:
                bundle_binding = {
                    field_name: write_ahead_request[field_name]
                    for field_name in (
                        "paper_bundle_job_id",
                        "paper_bundle_owner_id",
                        "paper_bundle_artifact_type",
                        "publication_generation",
                    )
                }
            else:
                bundle_binding = None

            conversation_id = _event_conversation_id(
                str(req.conversation_id or ""),
                published_run_id,
            )
            state = _RunState(
                artifact_type=candidate.artifact_type.value,
                brief=f"publish attempt {attempt}",
                conversation_id=conversation_id,
                baseline_artifact_json=json.dumps({"artifact_id": f"art_{run_id}"}),
            )
            state.demo_user_id = access_lease.user_id
            async with _RUNS_LOCK:
                _RUNS[published_run_id] = state
            settings = _settings_for_request(request)
            worker_request = CandidatePublishWorkerRequest(
                job_kind="candidate_publish",
                run_id=published_run_id,
                parent_run_id=run_id,
                conversation_id=conversation_id,
                settings=settings,
                source_attempt=attempt,
                expected_candidate_sha256=candidate.source_sha256,
            )
            candidate_publish_request = write_ahead_request or {
                "version": 2 if bundle_binding is not None else 1,
                "run_id": published_run_id,
                "source_run_id": run_id,
                "source_attempt": attempt,
                "source_candidate_id": candidate.candidate_id,
                "source_candidate_sha256": candidate.source_sha256,
                "idempotency_key_digest": idempotency_key_digest,
                "request_digest": request_digest,
                **(bundle_binding or {}),
            }
            if write_ahead_request is None:
                durable_replace_json(
                    RUNS_DIR
                    / published_run_id
                    / "candidate_publish_request.json",
                    candidate_publish_request,
                )
            derived_descriptor = {
                "job_kind": "candidate_publish",
                "source_artifact_id": f"art_{run_id}",
                "artifact_name": f"Published Attempt {attempt}",
                "source_relative_path": candidate.source_relative_path,
            }
            durable_derived_descriptor = _durable_derived_job_payload(
                worker_request,
                state,
                derived_descriptor,
            )
            try:
                prepared_descriptor = _read_derived_job_descriptor(
                    published_run_id
                )
            except ValueError as exc:
                raise HTTPException(
                    409,
                    detail="attempt publication descriptor is unreadable",
                ) from exc
            if prepared_descriptor is None:
                durable_replace_json(
                    RUNS_DIR / published_run_id / "derived_job.json",
                    durable_derived_descriptor,
                )
            elif prepared_descriptor != durable_derived_descriptor:
                raise HTTPException(
                    409,
                    detail="attempt publication descriptor changed",
                )
            reservation_task = asyncio.create_task(
                _start_supervised_derived_job(
                    request=worker_request,
                    state=state,
                    descriptor=derived_descriptor,
                    start_immediately=False,
                    settings=settings,
                    descriptor_prepared=True,
                )
            )
            caller_cancelled = await _join_owned_task(reservation_task)
            try:
                start_token = reservation_task.result()
            except Exception as exc:
                raise _web_run_service_error(exc) from exc
            reservation_claimed = True
            durable_handoff = True
            if caller_cancelled:
                raise asyncio.CancelledError
            if not reserve_only:
                durable_descriptor = _read_derived_job_descriptor(published_run_id)
                if durable_descriptor is None:
                    raise HTTPException(
                        500,
                        detail="candidate publish descriptor is missing",
                    )
                try:
                    await _start_reserved_derived_job(
                        run_id=published_run_id,
                        token=start_token,
                        state=state,
                        descriptor=durable_descriptor,
                    )
                except Exception as exc:
                    raise _web_run_service_error(exc) from exc
        except BaseException:
            await _cleanup_failed_candidate_publish_setup(
                run_id=published_run_id,
                state=state,
                access_lease=access_lease,
                durable_handoff=durable_handoff,
                reservation_claimed=reservation_claimed,
                discard_write_ahead=write_ahead_request is None,
            )
            raise
    _append_event(
        _settings_or_boot(),
        conversation_id,
        "artifact.attempt_publish_started",
        run_id=published_run_id,
        data={
            "artifact_type": candidate.artifact_type.value,
            "source_run_id": run_id,
            "source_attempt": attempt,
        },
    )
    return _direct_candidate_publish_ack(
        published_run_id,
        start_token=start_token if reserve_only else None,
    )


def _copy_poster_validation_context(
    source_run_dir: Path,
    target_run_dir: Path,
) -> None:
    for name in (
        "paper_visual_provenance.json",
        "paper_memory.json",
        "paper_memory_dossier.json",
        "paper_visual_storyboard.json",
        "poster_content_brief.json",
        "poster_plan_contract.json",
        "poster_contract_preflight.json",
        "canvas_plan.json",
        "paper_memory.md",
        "paper_memory_dossier.md",
    ):
        source_state = source_run_dir / name
        if source_state.is_file():
            shutil.copy2(source_state, target_run_dir / name)
    for name in ("layers", "paper_evidence_packs"):
        source_state = source_run_dir / name
        if source_state.is_dir():
            shutil.copytree(
                source_state,
                target_run_dir / name,
                dirs_exist_ok=True,
            )


def _poster_palette_id_for_run(run_dir: Path) -> str:
    run_brief = _read_json_file(run_dir / "run_brief.json")
    palette_id = (
        str(run_brief.get("palette_id") or "").strip()
        if isinstance(run_brief, dict)
        else ""
    )
    if palette_id:
        return palette_id
    for manifest_name in (
        "code_editor_revision_manifest.json",
        "authored_poster_edit_manifest.json",
        "apply_edits_palette_manifest.json",
    ):
        manifest = _read_json_file(run_dir / "final" / manifest_name)
        if isinstance(manifest, dict):
            palette_id = str(manifest.get("palette_id") or "").strip()
        if palette_id:
            return palette_id
    lineage = _read_json_file(run_dir / "candidate_draft_lineage.json")
    if isinstance(lineage, dict):
        return str(lineage.get("poster_palette_id") or "").strip()
    return ""


def _candidate_draft_artifact_from_lineage(
    run_dir: Path,
    run_id: str,
    artifact_type: str,
    lineage: dict[str, Any],
    *,
    source: Path | None = None,
) -> Artifact:
    run_dir = run_dir.resolve()
    if source is None:
        filename = {
            "poster": "poster.html",
            "landing": "index.html",
            "deck": "deck.html",
            "video": "deck.html",
        }.get(artifact_type)
        if filename is None:
            raise ValueError(f"unsupported candidate draft type: {artifact_type}")
        source = run_dir / "final" / filename
    source = source.resolve()
    if not source.is_file():
        raise ValueError(f"candidate draft HTML is missing: {source}")
    file_url = _run_file_url(
        run_id,
        source.relative_to(run_dir).as_posix(),
    )
    attempt = lineage.get("source_attempt") or "?"
    if artifact_type == "poster":
        canvas_width, canvas_height = _authored_paper_poster_size(source)
    elif artifact_type in {"deck", "video"}:
        canvas_width, canvas_height = (1920, 1080)
    else:
        canvas_width, canvas_height = (1440, 1200)
    return Artifact(
        artifact_id=f"art_{run_id}",
        name=f"Attempt {attempt} draft",
        artifact_type=artifact_type,
        canvas=Canvas(
            w=canvas_width,
            h=canvas_height,
        ),
        canvas_plan=_read_canvas_plan(run_dir),
        deck_plan=_read_deck_plan(run_dir),
        native_file_url=file_url,
        native_format="html",
        view_file_url=file_url,
        view_format="html",
        download_url=file_url,
        preview_url=(
            _run_file_url(run_id, "final/preview.png")
            if (run_dir / "final" / "preview.png").is_file()
            else None
        ),
        layers=parse_html_layers(source),
        candidate_draft=True,
        attempt_lineage=lineage,
    )


@app.post("/api/runs/{run_id}/attempts/{attempt}/fork", response_model=GenerateAck)
async def fork_run_attempt(
    run_id: str,
    attempt: int,
    req: AttemptForkRequest,
    request: Request,
) -> GenerateAck:
    _demo_require_run_owner(run_id, request)
    _assert_controlled_run_source_usable(run_id, mode="snapshot")
    source_run_dir = RUNS_DIR / run_id
    if not source_run_dir.is_dir():
        raise HTTPException(404, detail=f"run not found: {run_id}")
    try:
        candidate = load_attempt_candidate(source_run_dir, attempt)
    except ValueError as exc:
        raise HTTPException(404, detail=str(exc)) from exc
    draft_run_id = new_run_id()
    event_conversation_id = _event_conversation_id(
        str(req.conversation_id or ""),
        draft_run_id,
    )
    access_user_id = _demo_register_derived_run_access(
        draft_run_id,
        request,
        parent_run_id=run_id,
    )
    state = _RunState(
        artifact_type=candidate.artifact_type.value,
        brief=f"materialize attempt {attempt} for editing",
        conversation_id=event_conversation_id,
    )
    state.demo_user_id = access_user_id
    async with _RUNS_LOCK:
        _RUNS[draft_run_id] = state
    worker_request = AttemptForkWorkerRequest(
        job_kind="attempt_fork",
        run_id=draft_run_id,
        parent_run_id=run_id,
        attempt=attempt,
        expected_candidate_sha256=candidate.source_sha256,
        conversation_id=event_conversation_id,
        settings=_runtime_only_settings(),
    )
    reserve_only = _request_reserve_only(request)
    try:
        start_token = await _start_supervised_derived_job(
            request=worker_request,
            state=state,
            descriptor={
                "job_kind": "attempt_fork",
                "source_artifact_id": f"art_{run_id}",
                "artifact_name": f"Attempt {attempt} draft",
                "source_relative_path": candidate.source_relative_path,
            },
            start_immediately=not reserve_only,
            settings=worker_request.settings,
        )
    except Exception as exc:
        raise _web_run_service_error(exc) from exc
    _append_event(
        _settings_or_boot(),
        event_conversation_id,
        "artifact.attempt_fork_started",
        run_id=draft_run_id,
        data={
            "artifact_type": candidate.artifact_type.value,
            "source_run_id": run_id,
            "source_attempt": attempt,
        },
    )
    return GenerateAck(
        run_id=draft_run_id,
        progress_mode="attempt_fork",
        start_token=start_token if reserve_only else None,
        placeholder_message=Message(
            id=f"msg_{draft_run_id}",
            role="assistant",
            text="",
            ts=int(time.time() * 1000),
            run_id=draft_run_id,
            status="streaming",
        ),
    )


def _validate_candidate_draft(
    run_dir: Path,
    artifact_type: str,
    settings: Settings,
) -> list[dict[str, str]]:
    return _validate_candidate_draft_core(run_dir, artifact_type, settings)


def _deliver_video_candidate_draft(
    run_dir: Path,
    settings: Settings,
) -> list[dict[str, str]]:
    return _deliver_video_candidate_draft_core(run_dir, settings)


@app.post(
    "/api/artifacts/{artifact_id}/publish-candidate-draft",
    response_model=GenerateAck,
)
async def publish_candidate_draft(
    artifact_id: str,
    req: CandidateDraftPublishRequest,
    request: Request,
) -> GenerateAck:
    parent_run_id = _run_id_from_artifact_id(artifact_id)
    if not parent_run_id:
        raise HTTPException(400, detail="invalid artifact_id")
    _demo_require_run_owner(
        parent_run_id,
        request,
        detail=f"artifact not found: {artifact_id}",
    )
    _assert_controlled_run_source_usable(parent_run_id, mode="mutation")
    run_dir = RUNS_DIR / parent_run_id
    lineage_path = run_dir / "candidate_draft_lineage.json"
    lineage = _read_json_file(lineage_path)
    if not isinstance(lineage, dict) or lineage.get("status") != "draft":
        raise HTTPException(409, detail={"code": "candidate_draft_not_publishable"})
    artifact_type = str(
        lineage.get("artifact_type")
        or _detect_artifact_type_for_run(parent_run_id)
        or "poster"
    )
    if artifact_type not in {"poster", "deck", "landing", "video"}:
        raise HTTPException(
            409,
            detail={"code": "candidate_draft_artifact_type_invalid"},
        )
    source_run_id = str(lineage.get("source_run_id") or "")
    source_attempt = lineage.get("source_attempt")
    source_candidate_id = str(lineage.get("source_candidate_id") or "")
    source_candidate_sha256 = str(
        lineage.get("source_candidate_sha256") or ""
    )
    if (
        not source_run_id
        or type(source_attempt) is not int
        or source_attempt <= 0
        or not source_candidate_id
        or re.fullmatch(r"[0-9a-f]{64}", source_candidate_sha256) is None
    ):
        raise HTTPException(
            409,
            detail={"code": "candidate_draft_lineage_invalid"},
        )
    try:
        validate_run_id(source_run_id)
    except RunControlError as exc:
        raise HTTPException(
            409,
            detail={"code": "candidate_draft_lineage_invalid"},
        ) from exc
    _demo_require_run_owner(
        source_run_id,
        request,
        detail=f"source run not found: {source_run_id}",
    )
    _assert_controlled_run_source_usable(source_run_id, mode="snapshot")
    try:
        source_candidate = load_attempt_candidate(
            RUNS_DIR / source_run_id,
            source_attempt,
        )
    except ValueError as exc:
        raise HTTPException(
            409,
            detail={"code": "candidate_draft_source_invalid"},
        ) from exc
    if (
        source_candidate.artifact_type.value != artifact_type
        or source_candidate.candidate_id != source_candidate_id
        or source_candidate.source_sha256 != source_candidate_sha256
    ):
        raise HTTPException(
            409,
            detail={"code": "candidate_draft_source_changed"},
    )
    owner_id = _run_owner_id(request)
    published_run_id = new_run_id()
    access_lease = _demo_acquire_derived_run_access_lease(
        published_run_id,
        request,
        parent_run_id=parent_run_id,
        missing_detail=f"artifact not found: {artifact_id}",
    )
    state: _RunState | None = None
    durable_handoff = False
    reservation_claimed = False
    try:
        try:
            bundle_binding = await _reserve_candidate_publish_bundle_binding_async(
                source_run_id,
                artifact_type,
                owner_id,
            )
        except PaperBundleError as exc:
            raise _paper_bundle_error(exc) from exc
        conversation_id = _event_conversation_id(
            str(req.conversation_id or lineage.get("conversation_id") or ""),
            published_run_id,
        )
        state = _RunState(
            artifact_type=artifact_type,
            brief="validate and publish an edited attempt",
            conversation_id=conversation_id,
            baseline_artifact_json=json.dumps({"artifact_id": artifact_id}),
        )
        state.demo_user_id = access_lease.user_id
        async with _RUNS_LOCK:
            _RUNS[published_run_id] = state
        settings = _settings_for_request(request)
        worker_request = CandidatePublishWorkerRequest(
            job_kind="candidate_publish",
            run_id=published_run_id,
            parent_run_id=parent_run_id,
            conversation_id=conversation_id,
            settings=settings,
        )
        reserve_only = _request_reserve_only(request)
        idempotency_key_digest = hashlib.sha256(
            json.dumps(
                [owner_id, "candidate-draft", parent_run_id, published_run_id],
                ensure_ascii=False,
                separators=(",", ":"),
            ).encode("utf-8")
        ).hexdigest()
        request_digest = hashlib.sha256(
            json.dumps(
                {
                    "owner": owner_id,
                    "source_run_id": source_run_id,
                    "source_draft_run_id": parent_run_id,
                    "source_attempt": source_attempt,
                    "source_candidate_id": source_candidate_id,
                    "source_candidate_sha256": source_candidate_sha256,
                    "conversation_id": conversation_id,
                },
                ensure_ascii=False,
                sort_keys=True,
                separators=(",", ":"),
            ).encode("utf-8")
        ).hexdigest()
        durable_replace_json(
            RUNS_DIR / published_run_id / "candidate_publish_request.json",
            {
                "version": 3 if bundle_binding is not None else 1,
                "run_id": published_run_id,
                "source_run_id": source_run_id,
                "source_attempt": source_attempt,
                "source_candidate_id": source_candidate_id,
                "source_candidate_sha256": source_candidate_sha256,
                "idempotency_key_digest": idempotency_key_digest,
                "request_digest": request_digest,
                **(
                    {
                        **bundle_binding,
                        "source_draft_run_id": parent_run_id,
                    }
                    if bundle_binding is not None
                    else {}
                ),
            },
        )
        derived_descriptor = {
            "job_kind": "candidate_publish",
            "source_artifact_id": artifact_id,
            "artifact_name": f"Published Attempt {lineage.get('source_attempt') or '?'}",
            "source_relative_path": "candidate_draft_lineage.json",
        }
        durable_replace_json(
            RUNS_DIR / published_run_id / "derived_job.json",
            _durable_derived_job_payload(
                worker_request,
                state,
                derived_descriptor,
            ),
        )
        reservation_task = asyncio.create_task(
            _start_supervised_derived_job(
                request=worker_request,
                state=state,
                descriptor=derived_descriptor,
                start_immediately=False,
                settings=settings,
                descriptor_prepared=True,
            )
        )
        caller_cancelled = await _join_owned_task(reservation_task)
        try:
            start_token = reservation_task.result()
        except Exception as exc:
            raise _web_run_service_error(exc) from exc
        reservation_claimed = True
        durable_handoff = True
        if caller_cancelled:
            raise asyncio.CancelledError
        if not reserve_only:
            durable_descriptor = _read_derived_job_descriptor(published_run_id)
            if durable_descriptor is None:
                raise HTTPException(
                    500,
                    detail="candidate publish descriptor is missing",
                )
            try:
                await _start_reserved_derived_job(
                    run_id=published_run_id,
                    token=start_token,
                    state=state,
                    descriptor=durable_descriptor,
                )
            except Exception as exc:
                raise _web_run_service_error(exc) from exc
    except BaseException:
        await _cleanup_failed_candidate_publish_setup(
            run_id=published_run_id,
            state=state,
            access_lease=access_lease,
            durable_handoff=durable_handoff,
            reservation_claimed=reservation_claimed,
            discard_write_ahead=True,
        )
        raise
    _append_event(
        _settings_or_boot(),
        conversation_id,
        "artifact.attempt_publish_started",
        run_id=published_run_id,
        data={
            "artifact_type": artifact_type,
            "source_draft_run_id": parent_run_id,
            "source_run_id": source_run_id,
            "source_attempt": lineage.get("source_attempt"),
        },
    )
    return GenerateAck(
        run_id=published_run_id,
        progress_mode="attempt_publish",
        start_token=start_token if reserve_only else None,
        placeholder_message=Message(
            id=f"msg_{published_run_id}",
            role="assistant",
            text="",
            ts=int(time.time() * 1000),
            run_id=published_run_id,
            status="streaming",
        ),
    )


# ---------- SSE: /api/runs/{run_id}/events ----------

# Phase-bracket events the frontend cares about. Anything else gets
# streamed verbatim (the UI can ignore unknown event names).
_TERMINAL_EVENTS: frozenset[str] = frozenset({
    "run.done", "run.error", "apply.done", "run.cancelled",
    "openresearch.done", "openresearch.error",
})


# Longest a run can hold an SSE channel open before we force-disconnect.
# Tuned to cover the worst case observed in v2.8: paper2video with a
# heavy Composer turn + HyperFrames render is ~5–35 min, but cumulative
# planner retries can push past that. 60 min gives a comfortable margin
# without hoarding a connection forever.
_SSE_DEADLINE_S = 60 * 60
_DISK_RUN_RECOVERY_WINDOW_S = 2 * _SSE_DEADLINE_S


def _read_disk_run_events(
    path: Path,
    offset: int,
) -> tuple[int, list[dict[str, Any]]]:
    """Read complete JSONL events appended after ``offset``."""
    try:
        size = path.stat().st_size
        if size < offset:
            offset = 0
        events: list[dict[str, Any]] = []
        next_offset = offset
        with path.open("r", encoding="utf-8") as handle:
            handle.seek(offset)
            while True:
                line_offset = handle.tell()
                line = handle.readline()
                if not line:
                    break
                if not line.endswith("\n"):
                    next_offset = line_offset
                    break
                next_offset = handle.tell()
                try:
                    event = json.loads(line)
                except json.JSONDecodeError:
                    continue
                if isinstance(event, dict):
                    events.append(event)
        return next_offset, events
    except OSError:
        return offset, []


def _disk_run_lifecycle(run_dir: Path) -> tuple[str | None, int | None, bool]:
    control_path = run_dir / "run_control.json"
    if control_path.is_file():
        try:
            record = RunControlStore(run_dir.parent).read(run_dir.name)
        except RunControlError:
            pass
        else:
            terminal = record.terminal_event if record.state in {
                "completed", "failed", "cancelled"
            } else None
            started = record.state not in {"reserved", "uploading", "queued"}
            return terminal, record.worker_pid, started
    _offset, events = _read_disk_run_events(run_dir / "run_events.jsonl", 0)
    terminal: str | None = None
    owner_pid: int | None = None
    started = False
    for event in events:
        event_name = str(event.get("event") or "")
        if event_name == "run.start":
            started = True
            terminal = None
            try:
                candidate_pid = int(event.get("pid"))
            except (TypeError, ValueError):
                candidate_pid = 0
            owner_pid = candidate_pid if candidate_pid > 1 else None
        elif event_name in _TERMINAL_EVENTS:
            terminal = event_name
    return terminal, owner_pid, started


def _disk_run_terminal_event(run_dir: Path) -> str | None:
    terminal, _owner_pid, _started = _disk_run_lifecycle(run_dir)
    return terminal


def _in_memory_run_is_active(run_id: str) -> bool:
    state = _RUNS.get(run_id)
    if state is None:
        return False
    if bool(getattr(state, "queued", False)):
        return True
    task = getattr(state, "task", None)
    return task is not None and not task.done()


def _process_is_alive(pid: int) -> bool:
    try:
        os.kill(pid, 0)
    except ProcessLookupError:
        return False
    except PermissionError:
        return True
    except OSError:
        return False
    return True


def _disk_run_is_recoverable(run_dir: Path) -> bool:
    events_path = run_dir / "run_events.jsonl"
    try:
        age_s = max(0.0, time.time() - events_path.stat().st_mtime)
    except OSError:
        return False
    if age_s > _DISK_RUN_RECOVERY_WINDOW_S:
        return False
    terminal, owner_pid, started = _disk_run_lifecycle(run_dir)
    if terminal is not None or not started:
        return False
    if _in_memory_run_is_active(run_dir.name):
        return True
    if owner_pid is None or owner_pid == os.getpid():
        return False
    return _process_is_alive(owner_pid)


@app.get("/api/runs/{run_id}/events")
async def run_events(run_id: str, request: Request) -> StreamingResponse:
    """text/event-stream of log() events tagged with this run_id.
    Closes when a terminal event (`run.done` / `run.error` /
    `apply.done` / `run.cancelled`) arrives or after the absolute
    deadline."""
    _require_run_owner_before_lookup(run_id, request)
    run_events_path = RUNS_DIR / run_id / "run_events.jsonl"
    control_path = RUNS_DIR / run_id / "run_control.json"
    if not run_events_path.exists() and not control_path.exists():
        raise HTTPException(404, detail=f"run not found: {run_id}")
    if control_path.is_file():
        try:
            record = _web_run_runtime().control_store.read(run_id)
        except RunControlError:
            record = None
        if record is not None and record.state in {"cancelling", "completed", "failed", "cancelled"}:
            await _web_run_runtime().supervisor.recover(run_id)

    last_event_id = (request.headers.get("last-event-id", "") or "").strip()

    def sse_frame(event: dict[str, Any]) -> str:
        event_id = str(event.get("event_id") or "")
        lines = []
        if event_id:
            lines.append(f"id: {event_id}")
        lines.append(f"data: {json.dumps(event, ensure_ascii=False, default=str)}")
        return "\n".join(lines) + "\n\n"

    async def event_generator():
        disk_offset = 0
        replay_started = not last_event_id
        first_scan = True
        yield ":hello\n\n"
        deadline = time.monotonic() + _SSE_DEADLINE_S
        last_keepalive = time.monotonic()
        while time.monotonic() < deadline:
            if await request.is_disconnected():
                break
            disk_offset, disk_events = await asyncio.to_thread(
                _read_disk_run_events,
                run_events_path,
                disk_offset,
            )
            if first_scan and not replay_started:
                replay_started = not any(
                    str(event.get("event_id") or "") == last_event_id
                    for event in disk_events
                )
            first_scan = False
            terminal_found = False
            for disk_event in disk_events:
                event_id = str(disk_event.get("event_id") or "")
                if not replay_started:
                    if event_id == last_event_id:
                        replay_started = True
                    continue
                if disk_event.get("event") == "worker.exit":
                    continue
                yield sse_frame(disk_event)
                if disk_event.get("event") in _TERMINAL_EVENTS:
                    terminal_found = True
                    break
            if terminal_found:
                break
            now = time.monotonic()
            if now - last_keepalive >= 1.0:
                yield ": keepalive\n\n"
                last_keepalive = now
            await asyncio.sleep(0.05)

    return StreamingResponse(
        event_generator(),
        media_type="text/event-stream",
        headers={
            # Disable buffering at intermediaries (Vite proxy, nginx).
            "Cache-Control": "no-cache, no-transform",
            "X-Accel-Buffering": "no",
            "Connection": "keep-alive",
        },
    )


# ---------- POST /api/runs/{run_id}/retry ----------

def _video_export_retry_route(
    *,
    phase: str,
    error: str,
) -> Literal["full_authoring", "export_only", "setup_required", "none"]:
    normalized_phase = phase.strip().lower()
    normalized_error = error.strip().lower()
    setup_markers = (
        "cli is missing",
        "pinned hyperframes",
        "command not found",
        "executable missing",
        "executable not found",
        "ffmpeg not found",
        "ffprobe not found",
        "failed to start",
        "permission denied",
        "operation not permitted",
        "provider unavailable",
        "provider error",
    )
    if any(marker in normalized_error for marker in setup_markers):
        return "setup_required"
    authoring_markers = (
        "narration must be shortened",
        "required speech speed",
        "fitted speech duration",
        "speech coverage",
        "spoken wpm",
        "does not match the authored timeline",
    )
    if any(marker in normalized_error for marker in authoring_markers):
        return "full_authoring"
    if normalized_phase in {
        "validation",
        "authoring_lint",
        "subtitles",
        "delivery",
    }:
        return "full_authoring"
    if normalized_phase in {"tts", "lint", "render"}:
        return "export_only"
    return "none"


@app.post("/api/runs/{run_id}/retry-video-export", response_model=GenerateAck)
async def retry_video_export(
    run_id: str,
    request: Request,
    conversation_id: str | None = Form(None),
) -> GenerateAck:
    """Retry only lint, narration, render, probe, and final delivery."""

    if _DEMO_MODE:
        raise HTTPException(403, detail="Video export retry is unavailable in demo mode.")
    _require_run_owner_before_lookup(run_id, request)
    _assert_controlled_run_source_usable(run_id, mode="snapshot")
    try:
        _require_derived_source_ready(run_id, "video_export_retry")
    except RunNotReady as exc:
        raise HTTPException(409, detail=str(exc)) from exc
    _require_artifact_runtime("video")
    source_run_dir = RUNS_DIR / run_id
    if not source_run_dir.is_dir():
        raise HTTPException(404, detail=f"run not found: {run_id}")
    projects = sorted(
        (
            project
            for project in source_run_dir.glob("hyperframes-*")
            if (project / "index.html").is_file()
            and (project / "video_delivery_contract.json").is_file()
        ),
        key=lambda project: project.stat().st_mtime,
        reverse=True,
    )
    if not projects or not (source_run_dir / "design_spec.json").is_file():
        raise HTTPException(
            422,
            detail={
                "code": "video_export_retry_unavailable",
                "message": (
                    "This run does not contain a complete authored Video project. "
                    "Retry the full Video task instead."
                ),
            },
        )

    retry_run_id = new_run_id()
    event_conversation_id = _event_conversation_id(
        str(conversation_id or ""),
        retry_run_id,
    )
    baseline_artifact_json = json.dumps(
        {"artifact_id": f"art_{run_id}"},
        ensure_ascii=False,
    )
    access_user_id = _demo_register_derived_run_access(
        retry_run_id,
        request,
        parent_run_id=run_id,
    )
    state = _RunState(
        artifact_type="video",
        brief="retry existing video export",
        baseline_artifact_json=baseline_artifact_json,
        conversation_id=event_conversation_id,
    )
    state.demo_user_id = access_user_id
    async with _RUNS_LOCK:
        _RUNS[retry_run_id] = state
    worker_request = VideoExportRetryWorkerRequest(
        job_kind="video_export_retry",
        run_id=retry_run_id,
        parent_run_id=run_id,
        source_project=str(projects[0]),
        conversation_id=event_conversation_id,
        baseline_artifact_json=baseline_artifact_json,
        runs_dir=str(RUNS_DIR.resolve()),
    )
    reserve_only = _request_reserve_only(request)
    try:
        start_token = await _start_supervised_derived_job(
            request=worker_request,
            state=state,
            descriptor={
                "job_kind": "video_export_retry",
                "source_artifact_id": f"art_{run_id}",
                "artifact_name": "Video",
                "source_relative_path": _run_relative_path(projects[0]),
            },
            start_immediately=not reserve_only,
        )
    except Exception as exc:
        raise _web_run_service_error(exc) from exc
    _append_event(
        _settings_or_boot(),
        event_conversation_id,
        "artifact.retry_started",
        run_id=retry_run_id,
        data={
            "artifact_type": "video",
            "source": "video_export_retry",
            "parent_run_id": run_id,
        },
    )
    return GenerateAck(
        run_id=retry_run_id,
        progress_mode="video_export",
        start_token=start_token if reserve_only else None,
        placeholder_message=Message(
            id=f"msg_{retry_run_id}",
            role="assistant",
            text="",
            ts=int(time.time() * 1000),
            run_id=retry_run_id,
            status="streaming",
        ),
    )


async def _retry_video_export_in_background(
    *,
    run_id: str,
    source_run_id: str,
    source_project: Path,
    state: _RunState,
    conversation_id: str,
    baseline_artifact_json: str,
) -> None:
    started_at = time.time()
    run_dir = RUNS_DIR / run_id
    project_dir = run_dir / source_project.name
    settings = _settings_or_boot()
    failure_phase = "video_export"
    retry_route: Literal[
        "full_authoring",
        "export_only",
        "setup_required",
        "none",
    ] = "none"
    try:
        run_dir.mkdir(parents=True, exist_ok=True)
        shutil.copy2(
            RUNS_DIR / source_run_id / "design_spec.json",
            run_dir / "design_spec.json",
        )
        for optional in ("run_brief.json", "resume_state.json"):
            source = RUNS_DIR / source_run_id / optional
            if source.is_file():
                shutil.copy2(source, run_dir / optional)
        shutil.copytree(source_project, project_dir)
        _persisted_run_log(
            "run.start",
            run_id,
            pid=os.getpid(),
            artifact_type="video",
            source="video_export_retry",
            parent_run_id=source_run_id,
        )
        _append_event(
            settings,
            conversation_id,
            "artifact.retry_started",
            run_id=run_id,
            data={
                "artifact_type": "video",
                "source": "video_export_retry",
                "parent_run_id": source_run_id,
            },
        )
        from autodesign.tools.export_video import retry_video_export_project

        result = await asyncio.to_thread(
            retry_video_export_project,
            run_dir,
            project_dir,
        )
        if not result.get("ok"):
            failure_phase = str(result.get("phase") or "video_export")
            failure_error = str(
                result.get("error") or "Video export retry failed"
            )
            retry_route = _video_export_retry_route(
                phase=failure_phase,
                error=failure_error,
            )
            raise RuntimeError(failure_error)
        artifact = _build_artifact_response(
            run_dir,
            run_id,
            "video",
            baseline_artifact_json=baseline_artifact_json,
        )
        if artifact is None:
            raise RuntimeError("Video export retry completed without a publishable artifact")
        state.result_artifact = artifact
        state.result_message = Message(
            id=f"msg_{run_id}",
            role="assistant",
            text="Retried Video export from the existing authored project.",
            ts=int(time.time() * 1000),
            run_id=run_id,
            artifact_id=artifact.artifact_id,
            status="done",
        )
        _append_event(
            settings,
            conversation_id,
            "artifact.generated",
            run_id=run_id,
            artifact_id=artifact.artifact_id,
            data={
                "artifact_type": "video",
                "name": artifact.name,
                "native_format": artifact.native_format,
                "source": "video_export_retry",
                "parent_run_id": source_run_id,
            },
        )
        _persisted_run_log(
            "run.done",
            run_id,
            terminal_status="pass",
            parent_run_id=source_run_id,
        )
    except Exception as exc:  # noqa: BLE001
        state.error = f"{type(exc).__name__}: {exc}"
        state.result_message = Message(
            id=f"msg_{run_id}",
            role="assistant",
            text=f"Video export retry failed: {state.error}",
            ts=int(time.time() * 1000),
            run_id=run_id,
            status="error",
            failure=Failure(
                status="error",
                phase=failure_phase,
                retry_route=retry_route,
                parent_run_id=source_run_id,
                agent_last_note=state.error,
                produced_files=_list_produced_artifacts(run_id),
                elapsed_ms=int((time.time() - started_at) * 1000),
                artifact_type="video",
            ),
        )
        _append_event(
            settings,
            conversation_id,
            "artifact.generation_failed",
            run_id=run_id,
            data={
                "status": "error",
                "artifact_type": "video",
                "source": "video_export_retry",
                "parent_run_id": source_run_id,
                "phase": failure_phase,
                "retry_route": retry_route,
                "error": state.error[:500],
                "failure": _dump_model(state.result_message.failure),
            },
        )
        _persisted_run_log(
            "run.error",
            run_id,
            msg=state.error[:200],
            parent_run_id=source_run_id,
        )


@app.post("/api/runs/{run_id}/retry", response_model=GenerateAck)
async def run_retry(
    run_id: str,
    request: Request,
    designer_override: str | None = Form(None),
    planner_override: str | None = Form(None),
) -> GenerateAck:
    """Continue a failed run from its last validated author checkpoint.

    Runs without a checkpoint and explicit model overrides fall back to a fresh
    run using the original brief and attachments. Cancelled source trees are
    diagnostic-only and cannot be retried in place. The legacy
    `planner_override` field remains an alias for `designer_override`.
    """
    designer_override = designer_override or planner_override
    _require_run_owner_before_lookup(run_id, request)
    _assert_controlled_run_source_usable(run_id, mode="snapshot")
    async with _RUNS_LOCK:
        original = _RUNS.get(run_id)
    if original is None:
        raise HTTPException(404, detail=f"run not found: {run_id}")
    if original.task is not None and not original.task.done():
        raise HTTPException(
            409,
            detail={
                "code": "run_still_active",
                "message": "This run is still active and cannot be resumed yet.",
            },
        )
    if not original.brief:
        # Edge case: a state from before we recorded brief. Can't retry
        # without the original prompt; tell the user to retry from chat.
        raise HTTPException(
            422, detail="this run pre-dates retry support; resend your brief from chat",
        )
    access_user_id = ""
    if _RUN_ACCESS_CONTROL:
        access_user_id = _demo_user_id(request)
        if not _demo_user_owns_run(run_id, access_user_id):
            raise HTTPException(404, detail=f"run not found: {run_id}")
    if _DEMO_MODE:
        _demo_require_host_safety()
        if original.artifact_type != "poster":
            raise HTTPException(
                403,
                detail={
                    "code": "demo_poster_only",
                    "message": "Demo mode only supports paper poster generation.",
                },
            )
        if _demo_queue_full():
            raise HTTPException(
                429,
                detail={
                    "code": "demo_queue_full",
                    "message": "Demo queue is full. Try again later.",
                },
            )
        designer_override = None
    retry_palette_id = _validated_web_palette_id(
        original.artifact_type,
        original.palette_id,
    )

    # Re-resolve settings *with* the optional designer override applied.
    # Save + restore DESIGNER_MODEL under the lock so the override scope
    # is bounded to this request — without this, an Opus retry would
    # leak into subsequent runs that should have used the Settings
    # default.
    if designer_override:
        with _SETTINGS_LOCK:
            saved = os.environ.get("DESIGNER_MODEL")
            try:
                os.environ["DESIGNER_MODEL"] = designer_override.strip()
                req_settings = _settings_for_request(request)
            finally:
                if saved is None:
                    os.environ.pop("DESIGNER_MODEL", None)
                else:
                    os.environ["DESIGNER_MODEL"] = saved
    else:
        req_settings = _settings_for_request(request)
    retry_authoring_max_attempts = _validated_authoring_max_attempts(
        original.authoring_max_attempts,
        original.artifact_type,
        req_settings,
    )
    req_settings = _settings_with_authoring_max_attempts(
        req_settings,
        original.artifact_type,
        retry_authoring_max_attempts,
    )

    checkpoint = (
        _resume_checkpoint_from_disk(_settings_or_boot() / "runs" / run_id)
        if not _DEMO_MODE and not original.cancelled and not designer_override
        else None
    )
    # Terminal and cancelled runs are immutable diagnostic records. Resume is
    # always copy-based into a fresh run directory.
    retry_run_id = new_run_id()
    if _RUN_ACCESS_CONTROL and retry_run_id != run_id:
        _demo_register_run(retry_run_id, access_user_id)
    log(
        "web.run.resume.start" if checkpoint is not None else "web.run.retry.start",
        run_id=retry_run_id,
        parent_run_id=run_id,
        artifact_type=original.artifact_type,
        palette_id=retry_palette_id or "",
        authoring_max_attempts=retry_authoring_max_attempts,
        designer_override=designer_override or "",
        attachments=len(original.attach_paths),
        template=original.template or "",
        resume_from_attempt=(
            checkpoint["resume_from_attempt"] if checkpoint is not None else None
        ),
        next_attempt=checkpoint["next_attempt"] if checkpoint is not None else None,
    )
    _append_event(
        req_settings,
        original.conversation_id,
        "artifact.resume_started" if checkpoint is not None else "artifact.retry_started",
        run_id=retry_run_id,
        data={
            "parent_run_id": run_id,
            "artifact_type": original.artifact_type,
            "palette_id": retry_palette_id,
            "authoring_max_attempts": retry_authoring_max_attempts,
            "designer_override": designer_override or "",
            "planner_override": designer_override or "",
            "resume_from_attempt": (
                checkpoint["resume_from_attempt"] if checkpoint is not None else None
            ),
            "next_attempt": (
                checkpoint["next_attempt"] if checkpoint is not None else None
            ),
        },
    )

    new_state = _RunState(
        artifact_type=original.artifact_type,
        designer_model=req_settings.designer_model,
        has_pdf=original.has_pdf,
        brief=original.brief,
        attach_paths=original.attach_paths,
        reference_poster_path=original.reference_poster_path,
        baseline_artifact_json=original.baseline_artifact_json,
        conversation_id=original.conversation_id,
        template=original.template,
        canvas_preset_id=original.canvas_preset_id,
        palette_id=retry_palette_id,
        authoring_max_attempts=retry_authoring_max_attempts,
    )
    new_state.demo_user_id = access_user_id
    async with _RUNS_LOCK:
        _RUNS[retry_run_id] = new_state

    reserve_only = _request_reserve_only(request)
    if _DEMO_MODE:
        await _reserve_legacy_pipeline_worker(
            run_id=retry_run_id,
            brief=original.brief,
            attach_paths=original.attach_paths,
            reference_poster_path=original.reference_poster_path,
            template=original.template,
            state=new_state,
            settings=req_settings,
            resume_run=run_id if checkpoint is not None else None,
        )
        if not reserve_only:
            await _admit_demo_run(request, _DemoQueuedRun(
                run_id=retry_run_id,
                brief=original.brief,
                attach_paths=original.attach_paths,
                template=original.template,
                a_type=original.artifact_type,
                baseline_artifact_json=original.baseline_artifact_json,
                state=new_state,
                settings=req_settings,
            ))
    else:
        if reserve_only:
            await _reserve_legacy_pipeline_worker(
                run_id=retry_run_id,
                brief=original.brief,
                attach_paths=original.attach_paths,
                reference_poster_path=original.reference_poster_path,
                template=original.template,
                state=new_state,
                settings=req_settings,
                resume_run=run_id if checkpoint is not None else None,
            )
        else:
            await _start_legacy_pipeline_worker(
                run_id=retry_run_id,
                brief=original.brief,
                attach_paths=original.attach_paths,
                reference_poster_path=original.reference_poster_path,
                template=original.template,
                state=new_state,
                settings=req_settings,
                resume_run=run_id if checkpoint is not None else None,
            )
            new_state.task = asyncio.create_task(
                _monitor_supervised_pipeline(run_id=retry_run_id, state=new_state),
            )

    return GenerateAck(
        run_id=retry_run_id,
        start_token=new_state.reservation_token if reserve_only else None,
        placeholder_message=Message(
            id=f"msg_{retry_run_id}",
            role="assistant",
            text="",
            ts=int(time.time() * 1000),
            run_id=retry_run_id,
            status="streaming",
        ),
    )


# ---------- POST /api/runs/{run_id}/cancel ----------


def _controlled_derived_descendants(run_id: str) -> tuple[str, ...]:
    children_by_parent: dict[str, list[str]] = {}
    store = _web_run_runtime().control_store
    for control_path in sorted(RUNS_DIR.glob("*/run_control.json")):
        if control_path.parent.is_symlink():
            continue
        child_run_id = control_path.parent.name
        try:
            record = store.read(child_run_id)
        except (OSError, RunControlError, ValueError):
            continue
        if not record.parent_job_id:
            continue
        # The control record is the authoritative ownership edge. A corrupt or
        # missing derived descriptor must not let a live child escape parent
        # cancellation; descriptor validation remains a recovery concern.
        children_by_parent.setdefault(record.parent_job_id, []).append(child_run_id)

    descendants: list[str] = []
    pending = list(children_by_parent.get(run_id, ()))
    seen = {run_id}
    while pending:
        child_run_id = pending.pop(0)
        if child_run_id in seen:
            continue
        seen.add(child_run_id)
        descendants.append(child_run_id)
        pending.extend(children_by_parent.get(child_run_id, ()))
    return tuple(descendants)


def _append_cancel_request_event(run_id: str, reason: str) -> None:
    append_jsonl_event(
        RUNS_DIR / run_id / "run_events.jsonl",
        {"run_id": run_id, "event": "run.cancel_requested", "reason": reason},
        event_id=f"{run_id}:cancel-requested",
    )


async def _cancel_controlled_run(run_id: str, reason: str) -> CancelResult:
    runtime = _web_run_runtime()
    transitioned: list[str] = []
    async with _run_tree_lock(run_id):
        descendant_run_ids = _controlled_derived_descendants(run_id)
        for controlled_run_id in (run_id, *descendant_run_ids):
            before = runtime.control_store.read(controlled_run_id)
            requested = runtime.control_store.request_cancel(controlled_run_id)
            if (
                requested.state == "cancelling"
                and before.state not in {"cancelling", "cancelled", "completed", "failed"}
                and requested.revision == before.revision + 1
            ):
                transitioned.append(controlled_run_id)
    for transitioned_run_id in transitioned:
        _append_cancel_request_event(transitioned_run_id, reason)

    descendant_results: list[CancelResult | BaseException] = []
    for descendant_run_id in reversed(descendant_run_ids):
        try:
            child_result = await runtime.services.cancel(descendant_run_id, reason)
        except BaseException as exc:
            descendant_results.append(exc)
            continue
        descendant_results.append(child_result)
        if child_result.cancel_request_event_required:
            _append_cancel_request_event(descendant_run_id, reason)
        if child_result.state == "cancelled":
            async with _RUNS_LOCK:
                child_state = _RUNS.get(descendant_run_id)
                if child_state is not None:
                    child_state.cancelled = True
                    child_state.queued = False
                    child_state.error = reason

    result = await runtime.services.cancel(run_id, reason)
    if result.cancel_request_event_required:
        _append_cancel_request_event(run_id, reason)
    if result.state == "cancelled":
        async with _RUNS_LOCK:
            state = _RUNS.get(run_id)
            if state is not None:
                state.cancelled = True
                state.queued = False
                state.error = reason
    descendants_confirmed = all(
        isinstance(child_result, CancelResult) and child_result.confirmed
        for child_result in descendant_results
    )
    if descendants_confirmed:
        return result
    terminated_pids = {
        *result.terminated_pids,
        *(
            pid
            for child_result in descendant_results
            if isinstance(child_result, CancelResult)
            for pid in child_result.terminated_pids
        ),
    }
    surviving_pids = {
        *result.surviving_pids,
        *(
            pid
            for child_result in descendant_results
            if isinstance(child_result, CancelResult)
            for pid in child_result.surviving_pids
        ),
    }
    return replace(
        result,
        confirmed=False,
        terminated_pids=tuple(sorted(terminated_pids)),
        surviving_pids=tuple(sorted(surviving_pids)),
    )


@app.post("/api/runs/{run_id}/cancel", response_model=RunCancelResponse)
async def run_cancel(
    run_id: str,
    request: Request,
) -> RunCancelResponse | JSONResponse:
    """Confirm cancellation only after durable process and writer quiescence."""
    _require_run_owner_before_lookup(run_id, request)
    control_path = RUNS_DIR / run_id / "run_control.json"
    if control_path.is_file():
        try:
            result = await _cancel_controlled_run(run_id, "cancelled_by_user")
        except RunControlError as exc:
            raise HTTPException(404, detail=f"run not found: {run_id}") from exc
        await _reconcile_paper_bundle_for_run(
            run_id,
            owner_id=_run_owner_id(request),
        )
        if not result.confirmed:
            status_value = "cancellation_pending"
        elif result.state == "cancelled":
            status_value = "already_cancelled" if result.already_terminal else "cancelled"
        elif result.state in {"completed", "failed"}:
            status_value = "already_terminal"
        else:
            status_value = "cancellation_pending"
        payload = RunCancelResponse(
            run_id=run_id,
            status=status_value,
            run_state=result.state,
            confirmed=result.confirmed,
            terminated_pids=list(result.terminated_pids),
            surviving_pids=list(result.surviving_pids),
        )
        if not result.confirmed:
            return JSONResponse(status_code=202, content=payload.model_dump(mode="json"))
        return payload

    # Derived jobs are migrated separately. Preserve their legacy cancellation
    # behavior until they receive their own worker protocol variants.
    async with _RUNS_LOCK:
        state = _RUNS.get(run_id)
    if state is None:
        raise HTTPException(404, detail=f"run not found: {run_id}")
    if state.queued and state.task is None:
        state.cancelled = True
        state.error = "cancelled by user"
        failure = _failure_from_disk(
            run_id=run_id, a_type=state.artifact_type, status="cancelled",
            designer_model=state.designer_model, has_pdf=state.has_pdf,
            elapsed_ms=None,
        )
        state.result_message = Message(
            id=f"msg_{run_id}", role="assistant",
            text="Run cancelled.",
            ts=int(time.time() * 1000),
            run_id=run_id,
            status="error",
            failure=failure,
        )
        log("run.cancelled", run_id=run_id)
        _append_event(
            _settings_or_boot(), state.conversation_id, "artifact.generation_failed",
            run_id=run_id,
            data={
                "status": "cancelled",
                "artifact_type": state.artifact_type,
                "failure": _dump_model(failure),
                "canvas_plan": _read_canvas_plan(RUNS_DIR / run_id),
                "deck_plan": _read_deck_plan(RUNS_DIR / run_id),
            },
        )
        return RunCancelResponse(
            run_id=run_id,
            status="cancelled",
            run_state="cancelled",
            confirmed=True,
        )
    if state.task is None or state.task.done():
        return RunCancelResponse(
            run_id=run_id,
            status="already_terminal",
            run_state="completed" if state.result_artifact is not None else "failed",
            confirmed=True,
        )
    state.cancelled = True
    state.error = "cancelled by user"
    state.task.cancel()
    # Push a synthetic terminal event so live SSE listeners disconnect.
    log("run.cancelled", run_id=run_id)
    failure = _failure_from_disk(
        run_id=run_id, a_type=state.artifact_type, status="cancelled",
        designer_model=state.designer_model, has_pdf=state.has_pdf,
        elapsed_ms=None,
    )
    state.result_message = Message(
        id=f"msg_{run_id}", role="assistant",
        text="Run cancelled.",
        ts=int(time.time() * 1000),
        run_id=run_id,
        status="error",
        failure=failure,
    )
    _append_event(
        _settings_or_boot(), state.conversation_id, "artifact.generation_failed",
        run_id=run_id,
        data={
            "status": "cancelled",
            "artifact_type": state.artifact_type,
            "failure": _dump_model(failure),
            "canvas_plan": _read_canvas_plan(RUNS_DIR / run_id),
            "deck_plan": _read_deck_plan(RUNS_DIR / run_id),
        },
    )
    return RunCancelResponse(
        run_id=run_id,
        status="cancelled",
        run_state="cancelled",
        confirmed=True,
    )


# ---------- GET /api/runs/{run_id}/artifact ----------

def _controlled_terminal_diagnostic_response(
    run_id: str,
    state: _RunState | None,
) -> GenerateResponse | None:
    control_path = RUNS_DIR / run_id / "run_control.json"
    if not control_path.is_file():
        return None
    try:
        record = RunControlStore(RUNS_DIR).read(run_id)
    except RunControlError:
        return None
    if record.state not in {"failed", "cancelling", "cancelled"}:
        return None

    cancellation_diagnostic = record.state in {"cancelling", "cancelled"}
    status = "cancelled" if cancellation_diagnostic else "error"
    fallback_warnings = (
        _derived_worker_pointer_cleanup_warnings(run_id)
        if cancellation_diagnostic
        else ()
    )
    if state is not None and state.result_message is not None:
        message = state.result_message.model_copy(deep=True)
        failure = message.failure
        if failure is None:
            failure = Failure(status=status)
        existing_warnings = _immutable_pointer_cleanup_warnings(
            failure.pointer_cleanup_warnings
        )
        failure = failure.model_copy(update={
            "status": status,
            "pointer_cleanup_warnings": list(
                existing_warnings or fallback_warnings
            ),
            "produced_files": [],
        })
        message = message.model_copy(update={
            "artifact_id": None,
            "download_url": None,
            "download_filename": None,
            "download_mime_type": None,
            "task_payload": None,
            "status": "error",
            "failure": failure,
        })
    else:
        persisted_diagnostics = (
            _failure_diagnostics_from_disk(RUNS_DIR / run_id)
            if not cancellation_diagnostic
            else {}
        )
        message = Message(
            id=f"msg_{run_id}",
            role="assistant",
            text=(
                "Run cancelled."
                if cancellation_diagnostic
                else "Run ended without a publishable artifact."
            ),
            ts=int((record.terminal_at or record.updated_at) * 1000),
            run_id=run_id,
            status="error",
            failure=Failure(
                status=status,
                error_code=(
                    "run_cancellation_pending"
                    if record.state == "cancelling"
                    else "run_cancelled"
                    if record.state == "cancelled"
                    else persisted_diagnostics.get("error_code") or "run_failed"
                ),
                error_message=(
                    "Cancellation is still being verified."
                    if record.state == "cancelling"
                    else "The run was cancelled and its frozen output is diagnostic only."
                    if record.state == "cancelled"
                    else persisted_diagnostics.get("error_message")
                    or "The run failed before a publishable artifact was accepted."
                ),
                error_detail=persisted_diagnostics.get("error_detail"),
                phase=persisted_diagnostics.get("phase"),
                pointer_cleanup_warnings=list(fallback_warnings),
            ),
        )
    return GenerateResponse(message=message, artifact=None)


def _published_candidate_fork_artifact(
    source_run_id: str,
    *,
    authorize_run: Callable[[str], None] | None = None,
) -> Artifact | None:
    source_run_dir = RUNS_DIR / source_run_id
    try:
        journal = load_selection_journal(source_run_dir)
    except (OSError, ValueError) as exc:
        log(
            "web.attempt_selection_journal_invalid",
            run_id=source_run_id,
            error=f"{type(exc).__name__}: {exc}",
        )
        return None
    artifact_id = (
        journal.artifact_id
        if journal is not None and journal.state == "complete"
        else None
    )
    derived_run_id = _run_id_from_artifact_id(artifact_id or "")
    if not derived_run_id or derived_run_id == source_run_id:
        return None
    if authorize_run is not None:
        authorize_run(derived_run_id)
    derived_run_dir = RUNS_DIR / derived_run_id
    lineage = _read_json_file(
        derived_run_dir / "candidate_draft_lineage.json"
    )
    if (
        not isinstance(lineage, dict)
        or lineage.get("status") != "published"
        or lineage.get("source_run_id") != source_run_id
    ):
        return None
    artifact_type = str(lineage.get("artifact_type") or "")
    if artifact_type not in {"poster", "deck", "landing", "video"}:
        return None
    return _build_artifact_response(
        derived_run_dir,
        derived_run_id,
        artifact_type,  # type: ignore[arg-type]
        baseline_artifact_json=json.dumps(
            {"artifact_id": lineage.get("published_artifact_id_at_fork")}
        ),
    )


@app.get("/api/runs/{run_id}/artifact", response_model=GenerateResponse)
async def run_artifact(run_id: str, request: Request) -> GenerateResponse:
    """Final result of a run. Frontend calls this once the SSE stream
    delivers a `run.done` (or `apply.done`) event."""
    _require_run_owner_before_lookup(run_id, request)
    async with _RUNS_LOCK:
        state = _RUNS.get(run_id)
    terminal_diagnostic = _controlled_terminal_diagnostic_response(run_id, state)
    if terminal_diagnostic is not None:
        return terminal_diagnostic
    _assert_controlled_run_publishable(run_id)

    def authorize_published_fork(derived_run_id: str) -> None:
        _require_run_owner_before_lookup(derived_run_id, request)
        _assert_controlled_run_source_usable(derived_run_id, mode="artifact")

    published_fork = _published_candidate_fork_artifact(
        run_id,
        authorize_run=authorize_published_fork,
    )
    if published_fork is not None:
        return GenerateResponse(
            message=Message(
                id=f"msg_{run_id}",
                role="assistant",
                text="Published selected attempt.",
                ts=int(time.time() * 1000),
                run_id=run_id,
                artifact_id=published_fork.artifact_id,
                status="done",
            ),
            artifact=published_fork,
        )
    if state is None:
        # Could be a run started in a previous process; try the on-disk
        # artifact directly. Detect the real artifact type from disk so
        # video runs with final/index.html plus hyperframes/*.mp4 recover
        # as video, not as the intermediate landing page.
        detected_type = _detect_artifact_type_for_run(run_id) or "poster"
        artifact = _build_artifact_response(
            _BOOT_OUT_DIR / "runs" / run_id, run_id, detected_type,
            baseline_artifact_json=None,
        )
        if artifact is None:
            raise HTTPException(404, detail=f"run not found: {run_id}")
        return GenerateResponse(
            message=Message(
                id=f"msg_{run_id}", role="assistant",
                text=f"Recovered artifact for run {run_id[:8]} from disk.",
                ts=int(time.time() * 1000),
                run_id=run_id,
                artifact_id=artifact.artifact_id, status="done",
            ),
            artifact=artifact,
        )

    # If task isn't finished yet, wait a bounded amount. The SSE stream
    # should have already signaled `run.done`, but a slow disk write or
    # an over-eager client could race us here.
    if state.task and not state.task.done():
        try:
            await asyncio.wait_for(asyncio.shield(state.task), timeout=120)
        except asyncio.TimeoutError:
            raise HTTPException(504, detail="run still in progress; retry shortly")

    # state.error means the run errored or was cancelled — the supervised
    # completion monitor synthesizes a result_message carrying a Failure so
    # the frontend can render a structured
    # FailureCard. Only fall back to 500 when there's truly nothing to
    # hand back (shouldn't happen, but guard for safety).
    if state.result_message is None:
        if state.error:
            raise HTTPException(500, detail=state.error)
        raise HTTPException(500, detail="run completed without producing a message")
    return GenerateResponse(
        message=state.result_message, artifact=state.result_artifact,
    )


# ---------- /api/edits/apply ----------

# This is the no-LLM round-trip: the user tweaks a few text-layer fields in
# the right-rail Sidebar, hits Apply, and we re-run the existing
# `apply_edits` pipeline against a copy of the poster/landing HTML with
# `data-*` attributes patched in place. The result is a *new* run_dir
# (lineage chain via `parent_artifact_id`) — fast (5–15 s) and cheap (no
# API spend).
@app.post("/api/edits/apply", response_model=GenerateResponse)
async def edits_apply(
    request: Request,
    run_id: str = Form(...),
    artifact_type: str = Form(...),
    palette_id: str | None = Form(None),
    edits_json: str = Form(...),
    conversation_id: str | None = Form(None),
) -> GenerateResponse:
    _require_run_owner_before_lookup(run_id, request)
    _assert_controlled_run_source_usable(run_id, mode="mutation")
    req_settings = _settings_for_request(request)
    a_type: ArtifactType = _coerce_artifact_type(artifact_type)
    src_dir = req_settings.out_dir / "runs" / run_id / "final"
    source_candidate_lineage = _read_json_file(
        req_settings.out_dir / "runs" / run_id / "candidate_draft_lineage.json"
    )
    lineage_artifact_type = (
        str(source_candidate_lineage.get("artifact_type") or "")
        if isinstance(source_candidate_lineage, dict)
        and source_candidate_lineage.get("status") in {"draft", "published"}
        else ""
    )
    detected_type = (
        lineage_artifact_type
        if lineage_artifact_type in {"poster", "deck", "landing", "video"}
        else _artifact_type_from_final_dir(src_dir)
    )
    if detected_type is not None and detected_type != a_type:
        log(
            "web.edits.artifact_type.corrected",
            run_id=run_id,
            requested=a_type,
            detected=detected_type,
        )
        a_type = detected_type
    effective_palette_id = palette_id
    if (
        a_type == "poster"
        and not str(effective_palette_id or "").strip()
        and isinstance(source_candidate_lineage, dict)
    ):
        effective_palette_id = str(
            source_candidate_lineage.get("poster_palette_id") or ""
        ).strip() or None
        source_run_id = str(
            source_candidate_lineage.get("source_run_id") or ""
        ).strip()
        if effective_palette_id is None and source_run_id:
            effective_palette_id = (
                _poster_palette_id_for_run(
                    req_settings.out_dir / "runs" / source_run_id
                )
                or None
            )
    normalized_palette_id = _validated_web_palette_id(
        a_type,
        effective_palette_id,
    )
    required_color_system: dict[str, Any] | None = None
    if normalized_palette_id:
        try:
            required_color_system = require_academic_color_system(normalized_palette_id)
        except AcademicPaletteCatalogError as exc:
            raise HTTPException(
                503,
                detail={
                    "code": "palette_catalog_unavailable",
                    "message": str(exc),
                },
            ) from exc
    if _DEMO_MODE:
        if a_type != "poster":
            raise HTTPException(
                403,
                detail={
                    "code": "demo_poster_only",
                    "message": "Demo mode only supports paper poster artifacts.",
                },
            )
    try:
        edits = json.loads(edits_json)
    except json.JSONDecodeError as e:
        raise HTTPException(400, detail=f"edits_json must be JSON: {e}") from e
    if not isinstance(edits, dict):
        raise HTTPException(400, detail="edits_json must be an edit payload object")

    fname = _FILE_FOR[a_type][0]
    src_path = src_dir / fname
    if not src_path.exists():
        # Fall back to whatever HTML lives in final/ — same logic as
        # _build_artifact_response below.
        alt = _fallback_artifact(src_dir, "html")
        if alt is None:
            raise HTTPException(
                404,
                detail=f"source artifact missing: {src_path}. "
                "Has the run dir been cleared?",
            )
        src_path = alt

    new_run_id_str = new_run_id()
    event_conversation_id = _event_conversation_id(conversation_id, run_id)
    baseline_artifact_json = json.dumps({"artifact_id": f"art_{run_id}"})
    access_user_id = _demo_register_derived_run_access(
        new_run_id_str,
        request,
        parent_run_id=run_id,
    )
    state = _RunState(
        artifact_type=a_type,
        brief="apply structured artifact edits",
        baseline_artifact_json=baseline_artifact_json,
        conversation_id=event_conversation_id,
        palette_id=normalized_palette_id,
    )
    state.demo_user_id = access_user_id
    async with _RUNS_LOCK:
        _RUNS[new_run_id_str] = state

    input_path = RUNS_DIR / new_run_id_str / "uploads" / "artifact_edit.json"
    staged_html = input_path.parent / src_path.name
    worker_request = ArtifactEditWorkerRequest(
        job_kind="artifact_edit",
        run_id=new_run_id_str,
        parent_run_id=run_id,
        input_path=str(input_path),
        conversation_id=event_conversation_id,
        settings=req_settings,
    )
    descriptor = {
        "job_kind": "artifact_edit",
        "source_artifact_id": f"art_{run_id}",
        "artifact_name": src_path.stem,
        "source_relative_path": _run_relative_path(src_path),
    }
    try:
        start_token = await _start_supervised_derived_job(
            request=worker_request,
            state=state,
            descriptor=descriptor,
            start_immediately=False,
        )
    except Exception as exc:
        raise _web_run_service_error(exc) from exc

    before_layers: list[Any]
    after_layers: list[Any]
    try:
        async with _derived_tree_locks(run_id):
            _require_derived_ancestors_not_cancelled(run_id)
            _require_derived_source_ready(run_id, "artifact_edit")
            _web_run_runtime().control_store.assert_writable(new_run_id_str)
            input_path.parent.mkdir(parents=True, exist_ok=True)
            before_layers = parse_html_layers(src_path)
            try:
                _patch_html_for_apply_edits(
                    src_path,
                    staged_html,
                    edits,
                    source_run_id=run_id,
                )
            except Exception as exc:  # noqa: BLE001
                raise HTTPException(
                    400,
                    detail=f"failed to patch HTML: {exc}",
                ) from exc
            if a_type == "poster" and required_color_system:
                _validate_required_poster_palette_html(
                    staged_html,
                    required_color_system,
                )
            after_layers = parse_html_layers(staged_html)
            candidate_lineage_input = (
                source_candidate_lineage
                if isinstance(source_candidate_lineage, dict)
                and source_candidate_lineage.get("status") in {"draft", "published"}
                else {}
            )
            durable_replace_json(input_path, {
                "version": 1,
                "artifact_type": a_type,
                "source_relative_path": _run_relative_path(src_path),
                "edited_html_relative_path": _run_relative_path(staged_html),
                "edits": edits,
                "required_color_system": required_color_system or {},
                "candidate_lineage": candidate_lineage_input,
            })
            _require_derived_ancestors_not_cancelled(run_id)
            _web_run_runtime().control_store.assert_writable(new_run_id_str)
            await _start_reserved_derived_job(
                run_id=new_run_id_str,
                token=start_token,
                state=state,
                descriptor={
                    "version": _DERIVED_JOB_VERSION,
                    **descriptor,
                    "run_id": new_run_id_str,
                    "parent_run_id": run_id,
                    "artifact_type": a_type,
                    "conversation_id": event_conversation_id,
                    "baseline_artifact_json": baseline_artifact_json,
                },
            )
    except HTTPException:
        await _web_run_runtime().services.cancel(
            new_run_id_str,
            "artifact_edit_input_rejected",
        )
        raise
    except Exception as exc:  # noqa: BLE001
        await _web_run_runtime().services.cancel(
            new_run_id_str,
            "artifact_edit_start_failed",
        )
        if isinstance(exc, RunNotReady):
            raise _web_run_service_error(exc) from exc
        log("web.edits.error", error=type(exc).__name__, msg=str(exc)[:200])
        raise HTTPException(500, detail=f"apply_edits failed: {exc}") from exc

    log(
        "web.edits.apply",
        parent_run=run_id,
        child_run=new_run_id_str,
        edits=_html_edit_count(edits),
        artifact_type=a_type,
        palette_id=normalized_palette_id or "",
    )
    if state.task is None:
        raise HTTPException(500, detail="artifact edit monitor was not started")
    await asyncio.shield(state.task)

    async with _derived_tree_locks(run_id):
        try:
            _require_derived_ancestors_not_cancelled(run_id)
        except RunNotReady as exc:
            raise _web_run_service_error(exc) from exc
        _assert_controlled_run_source_usable(run_id, mode="mutation")
        child_record = _web_run_runtime().control_store.read(new_run_id_str)
        worker_outcome = _recovered_worker_outcome(
            new_run_id_str,
            expected_job_kind="artifact_edit",
        )
        if child_record.state != "completed" or not child_record.publishable:
            error_detail: Any = state.error or "artifact edit worker failed"
            worker_payload = _read_json_file(
                RUNS_DIR / new_run_id_str / "worker_result.json"
            )
            if isinstance(worker_payload, dict):
                worker_error = worker_payload.get("error")
                if isinstance(worker_error, dict):
                    raw_message = worker_error.get("message")
                    try:
                        decoded_detail = json.loads(str(raw_message or ""))
                    except json.JSONDecodeError:
                        decoded_detail = None
                    if isinstance(decoded_detail, dict):
                        error_detail = decoded_detail
            if (
                isinstance(error_detail, dict)
                and error_detail.get("code") == "poster_palette_validation_failed"
            ):
                failure = _read_json_file(
                    RUNS_DIR
                    / new_run_id_str
                    / "apply_edits_palette_validation_failure.json"
                )
                _append_event(
                    req_settings,
                    event_conversation_id,
                    "edits.apply_failed",
                    run_id=new_run_id_str,
                    data=failure if isinstance(failure, dict) else error_detail,
                )
                raise HTTPException(422, detail=error_detail)
            status_code = (
                409
                if child_record.state in {"cancelling", "cancelled"}
                else 500
            )
            raise HTTPException(status_code, detail=error_detail)
        if (
            worker_outcome is None
            or not worker_outcome.ok
            or not isinstance(worker_outcome.result, dict)
        ):
            raise HTTPException(500, detail="artifact edit result is unavailable")
        result_payload = worker_outcome.result
        restored_layer_ids = list(result_payload.get("restored_layer_ids") or [])
        skipped = list(result_payload.get("skipped") or [])
        artifact = state.result_artifact
        edit_diffs = layer_edit_events(before_layers, after_layers)
        for diff in edit_diffs:
            _append_event(
                req_settings, event_conversation_id, "layer.edited",
                run_id=new_run_id_str,
                artifact_id=f"art_{run_id}",
                data={
                    **diff,
                    "parent_run_id": run_id,
                    "new_run_id": new_run_id_str,
                    "artifact_type": a_type,
                    "palette_id": normalized_palette_id,
                    "required_color_system": required_color_system,
                },
            )
        _append_event(
            req_settings, event_conversation_id, "edits.applied",
            run_id=new_run_id_str,
            artifact_id=artifact.artifact_id if artifact else None,
            data={
                "parent_run_id": run_id,
                "artifact_type": a_type,
                "palette_id": normalized_palette_id,
                "required_color_system": required_color_system,
                "n_requested_edits": len(edits),
                "n_layer_diffs": len(edit_diffs),
                "restored_layer_ids": restored_layer_ids,
                "skipped": skipped,
            },
        )
        msg_bits: list[str] = []
        if artifact:
            msg_bits.append(
                f"Applied {len(edits)} edit{'s' if len(edits) != 1 else ''}. "
                f"Re-rendered {len(restored_layer_ids)} layer"
                f"{'s' if len(restored_layer_ids) != 1 else ''}."
            )
        else:
            msg_bits.append(
                f"Edits applied but no final artifact was produced "
                f"(run {new_run_id_str})."
            )
        msg_text = " ".join(msg_bits)
        response = GenerateResponse(
            message=Message(
                id=f"msg_{new_run_id_str}",
                role="assistant",
                text=msg_text,
                ts=int(time.time() * 1000),
                run_id=new_run_id_str,
                artifact_id=artifact.artifact_id if artifact else None,
                status="done" if artifact else "error",
            ),
            artifact=artifact,
        )
    return response


# ---------- helpers ----------


def _event_conversation_id(conversation_id: str | None, run_id: str) -> str:
    clean = (conversation_id or "").strip()
    return clean or f"web_{run_id}"


def _settings_or_boot() -> Path:
    return SETTINGS.out_dir if SETTINGS else _BOOT_OUT_DIR


def _demo_client_ip(request: Request) -> str:
    forwarded = (request.headers.get("x-forwarded-for", "") or "").split(",", 1)[0].strip()
    if forwarded:
        return forwarded[:128]
    if request.client and request.client.host:
        return request.client.host[:128]
    return "unknown"


def _demo_user_id(request: Request) -> str:
    raw = (
        request.headers.get(_DEMO_USER_HEADER, "")
        or request.cookies.get("autodesign_demo_user", "")
        or request.cookies.get("designanything_demo_user", "")
        or ""
    ).strip()
    if raw:
        clean = re.sub(r"[^A-Za-z0-9_.:-]+", "_", raw)[:128]
        if clean:
            return f"user:{clean}"
    return f"ip:{_demo_client_ip(request)}"


def _demo_today() -> str:
    return datetime.now(timezone.utc).date().isoformat()


def _demo_load_json(path: Path, default: dict[str, Any]) -> dict[str, Any]:
    if not path.exists():
        return default
    try:
        data = json.loads(path.read_text(encoding="utf-8"))
    except Exception:  # noqa: BLE001
        return default
    return data if isinstance(data, dict) else default


def _demo_write_json(path: Path, data: dict[str, Any]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    tmp = path.with_suffix(".tmp")
    tmp.write_text(json.dumps(data, indent=2, ensure_ascii=False), encoding="utf-8")
    os.replace(tmp, path)


def _demo_load_access() -> dict[str, Any]:
    data = _demo_load_json(DEMO_ACCESS_PATH, {"v": 1, "runs": {}})
    if not isinstance(data.get("runs"), dict):
        data["runs"] = {}
    return data


def _demo_write_access(data: dict[str, Any]) -> None:
    _demo_write_json(DEMO_ACCESS_PATH, data)


def _demo_require_run_owner(
    run_id: str,
    request: Request,
    *,
    detail: str | None = None,
) -> str:
    if not _RUN_ACCESS_CONTROL:
        return ""
    user_id = _demo_user_id(request)
    if not _demo_user_owns_run(run_id, user_id):
        raise HTTPException(404, detail=detail or f"run not found: {run_id}")
    return user_id


@dataclass(frozen=True)
class _DemoDerivedRunAccessLease:
    run_id: str
    user_id: str
    parent_run_id: str
    token: str


def _demo_acquire_derived_run_access_lease(
    run_id: str,
    request: Request,
    *,
    parent_run_id: str,
    missing_detail: str | None = None,
) -> _DemoDerivedRunAccessLease:
    parent = parent_run_id.strip()
    if not _RUN_ACCESS_CONTROL:
        return _DemoDerivedRunAccessLease(
            run_id=run_id,
            user_id="",
            parent_run_id=parent,
            token="",
        )
    user_id = _demo_user_id(request)
    token = secrets.token_urlsafe(24)
    with _DEMO_ACCESS_LOCK:
        data = _demo_load_access()
        runs = data.setdefault("runs", {})
        parent_entry = runs.get(parent)
        if (
            not isinstance(parent_entry, dict)
            or str(parent_entry.get("owner") or "") != user_id
        ):
            raise HTTPException(
                404,
                detail=missing_detail or f"run not found: {parent}",
            )
        runs[run_id] = {
            "owner": user_id,
            "token": token,
            "created_at": time.time(),
            "parent_run_id": parent,
        }
        _demo_write_access(data)
    return _DemoDerivedRunAccessLease(
        run_id=run_id,
        user_id=user_id,
        parent_run_id=parent,
        token=token,
    )


def _demo_register_derived_run_access(
    run_id: str,
    request: Request,
    *,
    parent_run_id: str | None = None,
    missing_detail: str | None = None,
) -> str:
    parent = (parent_run_id or "").strip()
    if parent:
        return _demo_acquire_derived_run_access_lease(
            run_id,
            request,
            parent_run_id=parent,
            missing_detail=missing_detail,
        ).user_id
    if not _RUN_ACCESS_CONTROL:
        return ""
    user_id = _demo_user_id(request)
    _demo_register_run(run_id, user_id)
    return user_id


def _demo_register_run(run_id: str, user_id: str) -> None:
    if not _RUN_ACCESS_CONTROL:
        return
    with _DEMO_ACCESS_LOCK:
        data = _demo_load_access()
        runs = data.setdefault("runs", {})
        runs[run_id] = {
            "owner": user_id,
            "token": secrets.token_urlsafe(24),
            "created_at": time.time(),
        }
        _demo_write_access(data)


def _demo_release_derived_run_access(
    lease: _DemoDerivedRunAccessLease,
) -> None:
    if not _RUN_ACCESS_CONTROL:
        return
    with _DEMO_ACCESS_LOCK:
        data = _demo_load_access()
        runs = data.setdefault("runs", {})
        entry = runs.get(lease.run_id)
        if (
            not isinstance(entry, dict)
            or str(entry.get("owner") or "") != lease.user_id
            or str(entry.get("parent_run_id") or "") != lease.parent_run_id
            or not lease.token
            or not secrets.compare_digest(
                str(entry.get("token") or ""),
                lease.token,
            )
        ):
            return
        runs.pop(lease.run_id, None)
        _demo_write_access(data)


def _demo_run_access(run_id: str) -> dict[str, Any] | None:
    with _DEMO_ACCESS_LOCK:
        entry = _demo_load_access().get("runs", {}).get(run_id)
    return entry if isinstance(entry, dict) else None


def _demo_user_owns_run(run_id: str, user_id: str) -> bool:
    if not _RUN_ACCESS_CONTROL:
        return True
    entry = _demo_run_access(run_id)
    return bool(entry and str(entry.get("owner") or "") == user_id)


def _demo_run_token(run_id: str) -> str:
    if not _RUN_ACCESS_CONTROL:
        return ""
    entry = _demo_run_access(run_id)
    return str(entry.get("token") or "") if entry else ""


def _demo_run_file_token_cookie_name(run_id: str) -> str:
    safe = re.sub(r"[^A-Za-z0-9_-]+", "_", run_id).strip("_")
    return f"{_RUN_FILE_TOKEN_COOKIE_PREFIX}{safe}"


def _legacy_demo_run_file_token_cookie_name(run_id: str) -> str:
    safe = re.sub(r"[^A-Za-z0-9_-]+", "_", run_id).strip("_")
    return f"{_LEGACY_RUN_FILE_TOKEN_COOKIE_PREFIX}{safe}"


def _demo_run_file_token_cookie_path(run_id: str) -> str:
    return f"/api/files/runs/{run_id}/"


def _run_file_url(run_id: str, rel: str) -> str:
    url = f"/api/files/runs/{run_id}/{rel.lstrip('/')}"
    if not _RUN_ACCESS_CONTROL:
        return url
    token = _demo_run_token(run_id)
    return f"{url}?token={token}" if token else url


def _assert_controlled_run_source_usable(
    run_id: str,
    *,
    mode: Literal["artifact", "mutation", "snapshot"],
) -> None:
    control_path = RUNS_DIR / run_id / "run_control.json"
    if not control_path.is_file():
        return
    try:
        record = RunControlStore(RUNS_DIR).read(run_id)
    except RunControlError as exc:
        raise HTTPException(409, detail="run lifecycle state is unavailable") from exc
    if record.state in {"cancelling", "cancelled"}:
        raise HTTPException(410, detail="cancelled run output is diagnostic only")
    if mode == "snapshot":
        if record.state not in {"running", "completing", "completed", "failed"}:
            raise HTTPException(409, detail=f"run snapshot is unavailable ({record.state})")
        return
    if mode == "mutation":
        if record.state in {"running", "completing"}:
            return
        if record.state == "completed" and record.publishable:
            return
        raise HTTPException(409, detail=f"run cannot be mutated ({record.state})")
    if record.state != "completed" or not record.publishable:
        raise HTTPException(409, detail=f"run artifact is not publishable ({record.state})")


def _controlled_run_is_cancellation_frozen(run_id: str) -> bool:
    control_path = RUNS_DIR / run_id / "run_control.json"
    if not control_path.is_file():
        return False
    try:
        record = RunControlStore(RUNS_DIR).read(run_id)
    except RunControlError:
        return True
    return record.state in {"cancelling", "cancelled"} or record.writes_frozen


def _assert_controlled_run_publishable(run_id: str) -> None:
    _assert_controlled_run_source_usable(run_id, mode="artifact")


class _OpenedPublicRunFileResponse(FileResponse):
    """FileResponse range semantics backed by one already-verified handle."""

    def __init__(self, opened: _OpenedPublicRunFile) -> None:
        self._opened_handle = opened.handle
        super().__init__(opened.path, stat_result=opened.stat_result)

    async def __call__(self, scope: Any, receive: Any, send: Any) -> None:
        try:
            await super().__call__(scope, receive, send)
        finally:
            await asyncio.to_thread(self._opened_handle.close)

    async def _read(self, size: int) -> bytes:
        return await asyncio.to_thread(self._opened_handle.read, size)

    async def _seek(self, offset: int) -> None:
        await asyncio.to_thread(self._opened_handle.seek, offset)

    async def _handle_simple(
        self,
        send: Any,
        send_header_only: bool,
        _send_pathsend: bool,
    ) -> None:
        await send({
            "type": "http.response.start",
            "status": self.status_code,
            "headers": self.raw_headers,
        })
        if send_header_only:
            await send({"type": "http.response.body", "body": b"", "more_body": False})
            return
        await self._seek(0)
        more_body = True
        while more_body:
            chunk = await self._read(self.chunk_size)
            more_body = len(chunk) == self.chunk_size
            await send({"type": "http.response.body", "body": chunk, "more_body": more_body})

    async def _handle_single_range(
        self,
        send: Any,
        start: int,
        end: int,
        file_size: int,
        send_header_only: bool,
    ) -> None:
        headers = MutableHeaders(raw=list(self.raw_headers))
        headers["content-range"] = f"bytes {start}-{end - 1}/{file_size}"
        headers["content-length"] = str(end - start)
        await send({"type": "http.response.start", "status": 206, "headers": headers.raw})
        if send_header_only:
            await send({"type": "http.response.body", "body": b"", "more_body": False})
            return
        await self._seek(start)
        while start < end:
            chunk = await self._read(min(self.chunk_size, end - start))
            if not chunk:
                break
            start += len(chunk)
            await send({
                "type": "http.response.body",
                "body": chunk,
                "more_body": start < end,
            })

    async def _handle_multiple_ranges(
        self,
        send: Any,
        ranges: list[tuple[int, int]],
        file_size: int,
        send_header_only: bool,
    ) -> None:
        boundary = secrets.token_hex(13)
        content_length, header_generator = self.generate_multipart(
            ranges,
            boundary,
            file_size,
            self.headers["content-type"],
        )
        headers = MutableHeaders(raw=list(self.raw_headers))
        headers["content-type"] = f"multipart/byteranges; boundary={boundary}"
        headers["content-length"] = str(content_length)
        await send({"type": "http.response.start", "status": 206, "headers": headers.raw})
        if send_header_only:
            await send({"type": "http.response.body", "body": b"", "more_body": False})
            return
        for start, end in ranges:
            await send({
                "type": "http.response.body",
                "body": header_generator(start, end),
                "more_body": True,
            })
            await self._seek(start)
            while start < end:
                chunk = await self._read(min(self.chunk_size, end - start))
                if not chunk:
                    break
                start += len(chunk)
                await send({"type": "http.response.body", "body": chunk, "more_body": True})
            await send({"type": "http.response.body", "body": b"\r\n", "more_body": True})
        await send({
            "type": "http.response.body",
            "body": f"--{boundary}--".encode("latin-1"),
            "more_body": False,
        })


def _public_run_file_parts(rel_path: str) -> list[str]:
    try:
        return list(canonical_run_file_parts(rel_path))
    except RunFileAccessError:
        raise HTTPException(404, detail="file not found")


def _captured_attempt_snapshot_paths(run_id: str) -> frozenset[str]:
    run_dir = RUNS_DIR / run_id
    try:
        candidates = load_attempt_candidates(run_dir)
    except (OSError, ValueError, ValidationError):
        return frozenset()
    return frozenset(
        path
        for candidate in candidates
        for path in (
            candidate.source_relative_path,
            *candidate.preview_relative_paths,
            *(
                candidate.browser_resource_relative_paths
                if candidate.browser_resource_relative_paths is not None
                else [
                    dependency
                    for dependency in candidate.dependency_relative_paths
                    if is_browser_preview_resource_path(dependency)
                ]
            ),
        )
    )


def _is_captured_attempt_snapshot_file(run_id: str, parts: list[str]) -> bool:
    return "/".join(parts[1:]) in _captured_attempt_snapshot_paths(run_id)


def _open_public_run_file(
    rel_path: str,
    *,
    expected_run_id: str | None = None,
) -> _OpenedPublicRunFile:
    parts = _public_run_file_parts(rel_path)
    run_id = parts[0]
    if expected_run_id is not None and run_id != expected_run_id:
        raise HTTPException(404, detail="file not found")
    if _is_captured_attempt_snapshot_file(run_id, parts):
        _assert_controlled_run_source_usable(run_id, mode="snapshot")
    else:
        _assert_controlled_run_publishable(run_id)
    try:
        return open_run_file(
            RUNS_DIR,
            rel_path,
            expected_run_id=expected_run_id,
        )
    except RunFileAccessError:
        raise HTTPException(404, detail="file not found") from None


def _resolve_public_run_file(rel_path: str) -> Path:
    opened = _open_public_run_file(rel_path)
    try:
        return opened.path
    finally:
        opened.handle.close()


def _demo_run_file_response(rel_path: str, request: Request) -> FileResponse:
    parts = [p for p in rel_path.split("/") if p]
    if not parts:
        raise HTTPException(404, detail="file not found")
    run_id = parts[0]
    entry = _demo_run_access(run_id)
    token = (request.query_params.get("token", "") or "").strip()
    token_from_query = bool(token)
    if not token:
        token = (
            request.cookies.get(_demo_run_file_token_cookie_name(run_id), "")
            or request.cookies.get(_legacy_demo_run_file_token_cookie_name(run_id), "")
            or ""
        ).strip()
    expected = str(entry.get("token") or "") if entry else ""
    if not expected or not token or not secrets.compare_digest(token, expected):
        raise HTTPException(403, detail="invalid file token")
    response = _OpenedPublicRunFileResponse(_open_public_run_file(rel_path))
    if token_from_query:
        response.set_cookie(
            _demo_run_file_token_cookie_name(run_id),
            expected,
            path=_demo_run_file_token_cookie_path(run_id),
            max_age=max(60, int(_DEMO_RUN_TTL_HOURS * 3600)),
            secure=request.url.scheme == "https",
            httponly=True,
            samesite="lax",
        )
    return response


def _demo_queue_full() -> bool:
    return bool(_DEMO_MODE and _DEMO_RUN_QUEUE is not None and _DEMO_RUN_QUEUE.full())


def _demo_require_host_safety() -> None:
    geteuid = getattr(os, "geteuid", None)
    if callable(geteuid) and geteuid() == 0 and not _truthy_env("DEMO_ALLOW_ROOT"):
        raise HTTPException(
            503,
            detail={
                "code": "demo_unsafe_host_user",
                "message": "Demo mode must run as a non-root, low-privilege service user.",
            },
        )


async def _enqueue_demo_run(job: _DemoQueuedRun) -> None:
    global _DEMO_RUN_QUEUE
    async with _web_run_start_guard(job.run_id):
        if _DEMO_RUN_QUEUE is None:
            _DEMO_RUN_QUEUE = asyncio.Queue(maxsize=_DEMO_QUEUE_MAX)
        if _DEMO_RUN_QUEUE.full():
            raise HTTPException(
                429,
                detail={
                    "code": "demo_queue_full",
                    "message": "Demo queue is full. Try again later.",
                },
            )
        _DEMO_RUN_QUEUE.put_nowait(job)
        log("demo.queue.enqueued", run_id=job.run_id, queued=_DEMO_RUN_QUEUE.qsize())


async def _admit_demo_run(request: Request, job: _DemoQueuedRun) -> bool:
    """Atomically charge and admit one demo run, or reuse its admission."""

    global _DEMO_RUN_QUEUE
    async with _RUNS_LOCK:
        if job.state.queued or job.state.task is not None:
            return False
        if _DEMO_RUN_QUEUE is None:
            _DEMO_RUN_QUEUE = asyncio.Queue(maxsize=_DEMO_QUEUE_MAX)
        if _DEMO_RUN_QUEUE.full():
            raise HTTPException(
                429,
                detail={
                    "code": "demo_queue_full",
                    "message": "Demo queue is full. Try again later.",
                },
            )
        _demo_consume_generation_quota(request, job.state.demo_user_id)
        job.state.queued = True
        _DEMO_RUN_QUEUE.put_nowait(job)
        log(
            "demo.queue.enqueued",
            run_id=job.run_id,
            queued=_DEMO_RUN_QUEUE.qsize(),
        )
        return True


def _demo_consume_generation_quota(request: Request, user_id: str) -> None:
    today = _demo_today()
    ip_id = f"ip:{_demo_client_ip(request)}"
    with _DEMO_USAGE_LOCK:
        data = _demo_load_json(DEMO_USAGE_PATH, {"v": 1, "days": {}})
        days = data.setdefault("days", {})
        day = days.setdefault(today, {"users": {}, "ips": {}})
        users = day.setdefault("users", {})
        ips = day.setdefault("ips", {})
        user_count = int(users.get(user_id, 0) or 0)
        ip_count = int(ips.get(ip_id, 0) or 0)
        if user_count >= _DEMO_DAILY_LIMIT:
            raise HTTPException(
                429,
                detail={
                    "code": "demo_daily_limit",
                    "message": f"Demo limit reached: {_DEMO_DAILY_LIMIT} poster runs per user per day.",
                },
            )
        if ip_count >= _DEMO_DAILY_LIMIT:
            raise HTTPException(
                429,
                detail={
                    "code": "demo_ip_daily_limit",
                    "message": f"Demo limit reached: {_DEMO_DAILY_LIMIT} poster runs per IP per day.",
                },
            )
        users[user_id] = user_count + 1
        ips[ip_id] = ip_count + 1
        _demo_write_json(DEMO_USAGE_PATH, data)
        log(
            "demo.quota.consume",
            user=user_id,
            ip=ip_id,
            user_count=users[user_id],
            ip_count=ips[ip_id],
            limit=_DEMO_DAILY_LIMIT,
        )


def _demo_ttl_protected_run_ids() -> set[str]:
    """Return active derived-publication inputs that TTL must not remove."""
    control_store = RunControlStore(RUNS_DIR)
    seeds: set[str] = set()
    terminal_states = {"completed", "failed", "cancelled"}
    for control_path in sorted(RUNS_DIR.glob("*/run_control.json")):
        run_id = control_path.parent.name
        try:
            record = control_store.read(run_id)
        except RunControlError:
            continue
        if (
            record.state in terminal_states
            and record.terminal_reconciliation_status != "pending"
        ):
            continue
        try:
            descriptor = _read_derived_job_descriptor(run_id)
        except (OSError, ValueError):
            continue
        if descriptor is not None:
            seeds.add(run_id)

    protected: set[str] = set()
    pending = list(seeds)
    while pending:
        run_id = pending.pop()
        try:
            validate_run_id(run_id)
        except RunControlError:
            continue
        if run_id in protected:
            continue
        protected.add(run_id)
        try:
            descriptor = _read_derived_job_descriptor(run_id)
        except (OSError, ValueError):
            continue
        if descriptor is None:
            continue
        dependencies = [descriptor["parent_run_id"]]
        if descriptor["job_kind"] == "candidate_publish":
            try:
                direct_publish = _read_direct_candidate_publish_descriptor(run_id)
            except (OSError, ValueError):
                direct_publish = None
            if direct_publish is not None:
                dependencies.append(direct_publish["source_run_id"])
                if direct_publish["version"] == 3:
                    dependencies.append(direct_publish["source_draft_run_id"])
        for dependency_run_id in dependencies:
            try:
                validate_run_id(dependency_run_id)
            except RunControlError:
                continue
            if dependency_run_id not in protected:
                pending.append(dependency_run_id)
    return protected


def _demo_ttl_run_can_expire(
    run_id: str,
    *,
    protected_run_ids: set[str],
    control_store: RunControlStore,
) -> bool:
    if run_id in _RUNS or run_id in protected_run_ids:
        return False
    try:
        validate_run_id(run_id)
    except RunControlError:
        return True
    control_path = RUNS_DIR / run_id / "run_control.json"
    if not control_path.is_file():
        return True
    try:
        record = control_store.read(run_id)
    except RunControlError:
        return False
    return (
        record.state in {"completed", "failed", "cancelled"}
        and record.terminal_reconciliation_status != "pending"
    )


def _demo_ttl_access_lineage_protected_run_ids(
    runs: dict[str, Any],
    deletable_run_ids: set[str],
) -> set[str]:
    """Protect ancestors of every access child not deleted in this pass."""
    protected: set[str] = set()
    pending = [run_id for run_id in runs if run_id not in deletable_run_ids]
    visited: set[str] = set()
    while pending:
        child_run_id = pending.pop()
        if child_run_id in visited:
            continue
        visited.add(child_run_id)
        child = runs.get(child_run_id)
        if not isinstance(child, dict):
            continue
        parent_run_id = str(child.get("parent_run_id") or "")
        try:
            validate_run_id(child_run_id)
            validate_run_id(parent_run_id)
        except RunControlError:
            continue
        parent = runs.get(parent_run_id)
        child_owner = str(child.get("owner") or "")
        if (
            not isinstance(parent, dict)
            or not child_owner
            or str(parent.get("owner") or "") != child_owner
        ):
            continue
        if parent_run_id not in protected:
            protected.add(parent_run_id)
            pending.append(parent_run_id)
    return protected


def _demo_cleanup_expired_runs() -> None:
    if not _DEMO_MODE:
        return
    cutoff = datetime.now(timezone.utc) - timedelta(hours=_DEMO_RUN_TTL_HOURS)
    cutoff_ts = cutoff.timestamp()
    with _DEMO_ACCESS_LOCK:
        access_snapshot = _demo_load_access()
        runs_snapshot = access_snapshot.setdefault("runs", {})
        expired_candidates: dict[str, float] = {}
        for run_id, entry in list(runs_snapshot.items()):
            if not isinstance(entry, dict):
                expired_candidates[run_id] = 0.0
                continue
            try:
                created_at = float(entry.get("created_at") or 0)
            except (TypeError, ValueError):
                created_at = 0
            if created_at < cutoff_ts:
                expired_candidates[run_id] = created_at

    control_store = RunControlStore(RUNS_DIR)
    first_protected_run_ids = _demo_ttl_protected_run_ids()
    eligible = {
        run_id
        for run_id in expired_candidates
        if _demo_ttl_run_can_expire(
            run_id,
            protected_run_ids=first_protected_run_ids,
            control_store=control_store,
        )
    }
    final_protected_run_ids = _demo_ttl_protected_run_ids()
    eligible = {
        run_id
        for run_id in eligible
        if _demo_ttl_run_can_expire(
            run_id,
            protected_run_ids=final_protected_run_ids,
            control_store=control_store,
        )
    }

    removed = 0
    with _DEMO_ACCESS_LOCK:
        data = _demo_load_access()
        runs = data.setdefault("runs", {})
        missing_entry = object()
        deletable_now: set[str] = set()
        for run_id in eligible:
            entry = runs.get(run_id, missing_entry)
            if run_id in _RUNS or entry is missing_entry:
                continue
            try:
                current_created_at = (
                    float(entry.get("created_at") or 0)
                    if isinstance(entry, dict)
                    else 0.0
                )
            except (TypeError, ValueError):
                current_created_at = 0.0
            if (
                current_created_at >= cutoff_ts
                or current_created_at != expired_candidates[run_id]
            ):
                continue
            deletable_now.add(run_id)
        deletable_now.difference_update(
            _demo_ttl_access_lineage_protected_run_ids(runs, deletable_now)
        )
        for run_id in deletable_now:
            try:
                validate_run_id(run_id)
            except RunControlError:
                runs.pop(run_id, None)
                removed += 1
                continue
            run_dir = (RUNS_DIR / run_id).resolve()
            if _path_inside(run_dir, RUNS_DIR.resolve()) and run_dir.exists():
                shutil.rmtree(run_dir, ignore_errors=True)
            runs.pop(run_id, None)
            removed += 1
        if removed:
            _demo_write_access(data)
            log(
                "demo.ttl.cleanup",
                expired=removed,
                ttl_hours=_DEMO_RUN_TTL_HOURS,
            )

    upload_cutoff = cutoff_ts
    if UPLOADS_DIR.exists():
        for path in UPLOADS_DIR.iterdir():
            try:
                if path.stat().st_mtime >= upload_cutoff:
                    continue
                if path.is_dir():
                    shutil.rmtree(path, ignore_errors=True)
                elif path.is_file() and path.name != UPLOADS_INDEX_PATH.name:
                    path.unlink(missing_ok=True)
            except OSError:
                continue


def _append_event(
    settings_or_path: Settings | Path,
    conversation_id: str,
    event: str,
    *,
    run_id: str | None = None,
    artifact_id: str | None = None,
    data: dict[str, Any] | None = None,
) -> None:
    out_dir = (
        settings_or_path.out_dir
        if isinstance(settings_or_path, Settings)
        else settings_or_path
    )
    try:
        append_design_event(
            out_dir,
            conversation_id,
            event,
            run_id=run_id,
            artifact_id=artifact_id,
            data=data,
        )
    except Exception as e:  # noqa: BLE001
        log("design_event.write_failed", design_event=event, error=f"{type(e).__name__}: {e}")


def _load_web_history() -> dict[str, Any]:
    if not WEB_HISTORY_PATH.exists():
        return {"v": 1, "conversations": {}}
    try:
        data = json.loads(WEB_HISTORY_PATH.read_text(encoding="utf-8"))
    except Exception:  # noqa: BLE001
        return {"v": 1, "conversations": {}}
    if not isinstance(data, dict):
        return {"v": 1, "conversations": {}}
    if not isinstance(data.get("conversations"), dict):
        data["conversations"] = {}
    return data


def _write_web_history(data: dict[str, Any]) -> None:
    WEB_HISTORY_PATH.parent.mkdir(parents=True, exist_ok=True)
    tmp = WEB_HISTORY_PATH.with_suffix(".tmp")
    tmp.write_text(json.dumps(data, indent=2, ensure_ascii=False), encoding="utf-8")
    os.replace(tmp, WEB_HISTORY_PATH)
    try:
        conversations = data.get("conversations", {})
        summaries = _history_conversation_summaries(
            conversations if isinstance(conversations, dict) else {},
        )
        _write_web_history_summary_index(summaries)
    except OSError:
        # The index is an optional read cache; the full snapshot is committed.
        pass


def _web_history_summary_index_path() -> Path:
    return WEB_HISTORY_PATH.with_name(f"{WEB_HISTORY_PATH.stem}.index.json")


def _web_history_file_fingerprint() -> dict[str, int] | None:
    try:
        stat = WEB_HISTORY_PATH.stat()
    except OSError:
        return None
    return {"mtime_ns": stat.st_mtime_ns, "size": stat.st_size}


def _load_web_history_summaries() -> dict[str, Any]:
    """Load the lightweight index, rebuilding it once for legacy snapshots."""
    with _WEB_HISTORY_LOCK:
        fingerprint = _web_history_file_fingerprint()
        if fingerprint is None:
            return {}
        index_path = _web_history_summary_index_path()
        try:
            index = json.loads(index_path.read_text(encoding="utf-8"))
        except (OSError, json.JSONDecodeError):
            index = None
        if (
            isinstance(index, dict)
            and index.get("source") == fingerprint
            and isinstance(index.get("conversations"), dict)
        ):
            return index["conversations"]

        data = _load_web_history()
        conversations = data.get("conversations", {})
        summaries = _history_conversation_summaries(
            conversations if isinstance(conversations, dict) else {},
        )
        _write_web_history_summary_index(summaries, fingerprint=fingerprint)
        return summaries


def _write_web_history_summary_index(
    conversations: dict[str, Any],
    *,
    fingerprint: dict[str, int] | None = None,
) -> None:
    source = fingerprint if fingerprint is not None else _web_history_file_fingerprint()
    if source is None:
        return
    path = _web_history_summary_index_path()
    path.parent.mkdir(parents=True, exist_ok=True)
    tmp = path.with_suffix(".tmp")
    tmp.write_text(
        json.dumps({"v": 1, "source": source, "conversations": conversations}, ensure_ascii=False),
        encoding="utf-8",
    )
    os.replace(tmp, path)


def _sanitize_history_conversation(
    raw: dict[str, Any],
    *,
    preserve_run_state: bool = False,
) -> dict[str, Any] | None:
    if not isinstance(raw, dict):
        return None
    conv_id = str(raw.get("id") or "").strip()
    if not conv_id:
        return None
    try:
        created_at = int(raw.get("created_at") or int(time.time() * 1000))
        updated_at = int(raw.get("updated_at") or created_at)
    except (TypeError, ValueError):
        created_at = updated_at = int(time.time() * 1000)
    messages = raw.get("messages") if isinstance(raw.get("messages"), list) else []
    artifacts = raw.get("artifacts") if isinstance(raw.get("artifacts"), dict) else {}
    active = raw.get("active_artifact_id")
    if not isinstance(active, str) or active not in artifacts:
        active = next(iter(artifacts.keys()), None)
    conversation = {
        "id": conv_id,
        "title": str(raw.get("title") or "New chat")[:200],
        "created_at": created_at,
        "updated_at": updated_at,
        "messages": messages,
        "artifacts": artifacts,
        "active_artifact_id": active,
        "pending_edits": {},
    }
    poster_palette_id = str(raw.get("poster_palette_id") or "").strip()
    if poster_palette_id:
        conversation["poster_palette_id"] = poster_palette_id
    poster_canvas_preset_id = str(raw.get("poster_canvas_preset_id") or "").strip()
    if poster_canvas_preset_id:
        conversation["poster_canvas_preset_id"] = poster_canvas_preset_id
    if preserve_run_state and raw.get("pending") is True:
        run_id = str(raw.get("run_id") or "").strip()
        if run_id:
            conversation["pending"] = True
            conversation["run_id"] = run_id
    try:
        message_count = int(raw.get("message_count"))
    except (TypeError, ValueError):
        message_count = 0
    if message_count > 0:
        conversation["message_count"] = message_count
    pending_artifact_type = _history_artifact_type(raw.get("pending_artifact_type"))
    if pending_artifact_type is not None:
        conversation["pending_artifact_type"] = pending_artifact_type
    return conversation


def _merge_history_conversations(
    a: dict[str, Any],
    b: dict[str, Any],
) -> dict[str, Any]:
    merged: dict[str, Any] = {}
    for source_index, source in enumerate((a, b)):
        for conv_id, raw in source.items():
            if not isinstance(raw, dict):
                continue
            conv = _sanitize_history_conversation(
                raw,
                preserve_run_state=source_index == 1,
            )
            if conv is None:
                continue
            existing = merged.get(conv_id)
            if not existing:
                merged[conv_id] = conv
                continue
            merged[conv_id] = _merge_one_history_conversation(existing, conv)
    return merged


def _merge_one_history_conversation(
    left: dict[str, Any],
    right: dict[str, Any],
) -> dict[str, Any]:
    left_updated = int(left.get("updated_at") or 0)
    right_updated = int(right.get("updated_at") or 0)
    primary, secondary = (right, left) if right_updated >= left_updated else (left, right)
    left_artifacts = left.get("artifacts") if isinstance(left.get("artifacts"), dict) else {}
    right_artifacts = right.get("artifacts") if isinstance(right.get("artifacts"), dict) else {}
    # _merge_history_conversations passes stored browser history first and
    # freshly imported run artifacts second. Prefer the imported artifact when
    # ids overlap so old local/native cache entries don't mask newly parsed
    # editable deck/video artifacts.
    artifacts = {
        **left_artifacts,
        **right_artifacts,
    }
    messages: list[dict[str, Any]] = []
    seen: set[str] = set()
    for msg in [
        *(secondary.get("messages") if isinstance(secondary.get("messages"), list) else []),
        *(primary.get("messages") if isinstance(primary.get("messages"), list) else []),
    ]:
        if not isinstance(msg, dict):
            continue
        msg_id = str(msg.get("id") or "")
        run_id = str(msg.get("run_id") or "")
        key = msg_id or f"{msg.get('role')}:{run_id}:{msg.get('text')}"
        if key in seen:
            continue
        seen.add(key)
        messages.append(msg)
    messages.sort(key=lambda m: int(m.get("ts") or 0))
    active = primary.get("active_artifact_id")
    if not isinstance(active, str) or active not in artifacts:
        active = secondary.get("active_artifact_id")
    if not isinstance(active, str) or active not in artifacts:
        active = next(reversed(artifacts), None) if artifacts else None
    poster_palette_id = str(
        primary.get("poster_palette_id")
        or secondary.get("poster_palette_id")
        or ""
    ).strip()
    poster_canvas_preset_id = str(
        primary.get("poster_canvas_preset_id")
        or secondary.get("poster_canvas_preset_id")
        or ""
    ).strip()
    active_run = right if right.get("pending") is True else left if left.get("pending") is True else None
    message_count = max(
        len(messages),
        int(left.get("message_count") or 0),
        int(right.get("message_count") or 0),
    )
    pending_artifact_type = (
        _history_artifact_type(primary.get("pending_artifact_type"))
        or _history_artifact_type(secondary.get("pending_artifact_type"))
    )
    merged = {
        **primary,
        "created_at": min(
            int(left.get("created_at") or left_updated or int(time.time() * 1000)),
            int(right.get("created_at") or right_updated or int(time.time() * 1000)),
        ),
        "updated_at": max(left_updated, right_updated),
        "messages": messages,
        "artifacts": artifacts,
        "active_artifact_id": active,
        "pending_edits": {},
        **({"poster_palette_id": poster_palette_id} if poster_palette_id else {}),
        **(
            {"poster_canvas_preset_id": poster_canvas_preset_id}
            if poster_canvas_preset_id
            else {}
        ),
        **(
            {"pending": True, "run_id": active_run.get("run_id")}
            if active_run and active_run.get("run_id")
            else {}
        ),
    }
    if message_count:
        merged["message_count"] = message_count
    if pending_artifact_type is not None:
        merged["pending_artifact_type"] = pending_artifact_type
    return merged


def _history_artifact_type(value: Any) -> ArtifactType | None:
    artifact_type = str(value or "").strip()
    if artifact_type in {"poster", "landing", "deck", "video"}:
        return artifact_type  # type: ignore[return-value]
    return None


def _history_artifact_preview(raw: Any, *, fallback_id: str) -> dict[str, Any] | None:
    if not isinstance(raw, dict):
        return None
    artifact_type = _history_artifact_type(raw.get("artifact_type"))
    artifact_id = str(raw.get("artifact_id") or fallback_id).strip()
    if artifact_type is None or not artifact_id:
        return None
    _filename, _format, (default_w, default_h) = _FILE_FOR[artifact_type]
    canvas_raw = raw.get("canvas") if isinstance(raw.get("canvas"), dict) else {}
    try:
        width = int(canvas_raw.get("w") or default_w)
    except (TypeError, ValueError):
        width = default_w
    try:
        height = int(canvas_raw.get("h") or default_h)
    except (TypeError, ValueError):
        height = default_h
    preview = {
        "artifact_id": artifact_id,
        "name": str(raw.get("name") or f"{artifact_type.title()} artifact").strip()[:200],
        "artifact_type": artifact_type,
        "canvas": {
            "w": width if width > 0 else default_w,
            "h": height if height > 0 else default_h,
        },
    }
    preview_url = raw.get("preview_url")
    if (
        isinstance(preview_url, str)
        and preview_url
        and not preview_url.lstrip().lower().startswith("data:")
        and len(preview_url) <= 4096
    ):
        preview["preview_url"] = preview_url
    card_preview_url = raw.get("card_preview_url")
    if (
        isinstance(card_preview_url, str)
        and card_preview_url
        and not card_preview_url.lstrip().lower().startswith("data:")
        and len(card_preview_url) <= 4096
    ):
        preview["card_preview_url"] = card_preview_url
    quality_status = raw.get("quality_status")
    if quality_status in {"ready", "ready_with_warnings"}:
        preview["quality_status"] = quality_status
        diagnostics = raw.get("quality_diagnostics")
        if isinstance(diagnostics, list):
            preview["quality_diagnostics"] = [
                value.strip()
                for value in diagnostics
                if isinstance(value, str) and value.strip()
            ]
    return preview


def _history_pending_artifact_type(conversation: dict[str, Any]) -> ArtifactType | None:
    pending_artifact_type = _history_artifact_type(conversation.get("pending_artifact_type"))
    if pending_artifact_type is not None:
        return pending_artifact_type
    messages = conversation.get("messages")
    if isinstance(messages, list):
        for message in reversed(messages):
            if not isinstance(message, dict):
                continue
            payload = message.get("task_payload")
            if not isinstance(payload, dict):
                continue
            pending_artifact_type = _history_artifact_type(payload.get("artifact_type"))
            if pending_artifact_type is not None:
                return pending_artifact_type
    return None


def _history_last_run_metadata(conversation: dict[str, Any]) -> dict[str, Any] | None:
    run_id = str(conversation.get("run_id") or "").strip()
    matched_message: dict[str, Any] | None = None
    messages = conversation.get("messages")
    if isinstance(messages, list):
        for message in reversed(messages):
            if not isinstance(message, dict):
                continue
            message_run_id = str(message.get("run_id") or "").strip()
            if not message_run_id:
                continue
            if not run_id:
                run_id = message_run_id
            if message_run_id == run_id:
                matched_message = message
                break
    if not run_id:
        return None
    metadata: dict[str, Any] = {"run_id": run_id}
    status = str((matched_message or {}).get("status") or "").strip()
    if conversation.get("pending") is True:
        status = "streaming"
    if status in {"streaming", "done", "error"}:
        metadata["status"] = status
    artifact_id = str(
        (matched_message or {}).get("artifact_id")
        or conversation.get("active_artifact_id")
        or ""
    ).strip()
    if artifact_id:
        metadata["artifact_id"] = artifact_id
    return metadata


def _history_conversation_summary(raw: Any) -> dict[str, Any] | None:
    conversation = _sanitize_history_conversation(raw, preserve_run_state=True)
    if conversation is None:
        return None
    previews: dict[str, Any] = {}
    artifacts = conversation.get("artifacts")
    if isinstance(artifacts, dict):
        for artifact_id, artifact in artifacts.items():
            preview = _history_artifact_preview(artifact, fallback_id=str(artifact_id))
            if preview is not None:
                previews[preview["artifact_id"]] = preview
    active_artifact_id = conversation.get("active_artifact_id")
    if active_artifact_id not in previews:
        active_artifact_id = next(reversed(previews), None) if previews else None
    messages = conversation.get("messages")
    message_count = len(messages) if isinstance(messages, list) else 0
    message_count = max(message_count, int(conversation.get("message_count") or 0))
    summary = {
        "id": conversation["id"],
        "title": conversation["title"],
        "created_at": conversation["created_at"],
        "updated_at": conversation["updated_at"],
        "message_count": message_count,
        "artifacts": previews,
        "active_artifact_id": active_artifact_id,
    }
    poster_palette_id = conversation.get("poster_palette_id")
    if isinstance(poster_palette_id, str) and poster_palette_id:
        summary["poster_palette_id"] = poster_palette_id
    poster_canvas_preset_id = conversation.get("poster_canvas_preset_id")
    if isinstance(poster_canvas_preset_id, str) and poster_canvas_preset_id:
        summary["poster_canvas_preset_id"] = poster_canvas_preset_id
    if conversation.get("pending") is True and isinstance(conversation.get("run_id"), str):
        summary["pending"] = True
        summary["run_id"] = conversation["run_id"]
        pending_artifact_type = _history_pending_artifact_type(conversation)
        if pending_artifact_type is not None:
            summary["pending_artifact_type"] = pending_artifact_type
        if isinstance(messages, list):
            for message in reversed(messages):
                if not isinstance(message, dict):
                    continue
                if str(message.get("run_id") or "").strip() != conversation["run_id"]:
                    continue
                task_type = str(message.get("task_type") or "").strip()
                task_payload = message.get("task_payload")
                if task_type in {
                    "generate",
                    "poster_code_edit",
                    "artifact_export_pptx",
                    "candidate_publish",
                }:
                    summary["pending_task_type"] = task_type
                if isinstance(task_payload, dict):
                    recoverable_payload = {
                        key: task_payload[key]
                        for key in (
                            "artifact_type",
                            "template",
                            "canvas_preset_id",
                            "source_artifact_id",
                            "source_run_id",
                            "source_candidate_id",
                            "export_format",
                        )
                        if key in task_payload
                    }
                    if recoverable_payload:
                        summary["pending_task_payload"] = recoverable_payload
                break
    last_run = _history_last_run_metadata(conversation)
    if last_run is not None:
        summary["last_run"] = last_run
    return summary


def _history_conversation_summaries(conversations: dict[str, Any]) -> dict[str, Any]:
    summaries: dict[str, Any] = {}
    for conversation_id, raw in conversations.items():
        summary = _history_conversation_summary(raw)
        if summary is not None:
            summaries[str(conversation_id)] = summary
    return summaries


def _limit_history_conversations(conversations: dict[str, Any], limit: int) -> dict[str, Any]:
    ordered = [
        (conversation_id, conversation)
        for conversation_id, conversation in conversations.items()
        if isinstance(conversation, dict)
    ]
    ordered.sort(
        key=lambda item: (int(item[1].get("updated_at") or 0), item[0]),
        reverse=True,
    )
    return dict(ordered[:limit])


def _history_summary_cutoff(conversations: dict[str, Any], limit: int) -> int | None:
    updated_at: list[int] = []
    for conversation in conversations.values():
        if not isinstance(conversation, dict):
            continue
        try:
            updated_at.append(int(conversation.get("updated_at") or 0))
        except (TypeError, ValueError):
            continue
    if len(updated_at) < limit:
        return None
    updated_at.sort(reverse=True)
    return updated_at[limit - 1]


def _history_detail_conversation(conversation: dict[str, Any]) -> dict[str, Any]:
    detail = dict(conversation)
    detail.pop("message_count", None)
    detail.pop("pending_artifact_type", None)
    return detail


def _count_history_artifacts(conversations: dict[str, Any]) -> int:
    count = 0
    for conv in conversations.values():
        if isinstance(conv, dict) and isinstance(conv.get("artifacts"), dict):
            count += len(conv["artifacts"])
    return count


def _import_history_from_server_events(
    *,
    limit: int | None,
    include_design_sessions: bool,
    demo_user_id: str | None,
    compact: bool = False,
    summary_limit: int | None = None,
    summary_min_updated_at: int | None = None,
    existing_conversation_ids: set[str] | None = None,
) -> dict[str, Any]:
    events = [] if demo_user_id else (_read_design_session_events() if include_design_sessions else [])
    conversations: dict[str, Any] = {}
    run_ids_seen: set[str] = set()
    events_by_conv: dict[str, list[dict[str, Any]]] = {}
    for event in events:
        conv_id = str(event.get("conversation_id") or "").strip()
        if not conv_id:
            continue
        events_by_conv.setdefault(conv_id, []).append(event)

    for conv_id, conv_events in events_by_conv.items():
        conv = _conversation_from_design_events(
            conv_id,
            conv_events,
            run_ids_seen,
            compact=compact,
        )
        if conv is not None:
            conversations[conv_id] = conv

    if compact and summary_limit is not None:
        # Disk-imported conversations use their run directory mtime as
        # `updated_at`. Scan every directory that can enter the global top-N
        # after the stored summaries are merged, including timestamp ties.
        run_paths: list[tuple[Path, int]] = []
        owned_run_ids: set[str] | None = None
        if demo_user_id:
            with _DEMO_ACCESS_LOCK:
                runs = _demo_load_access().get("runs", {})
            owned_run_ids = {
                run_id
                for run_id, entry in runs.items()
                if isinstance(entry, dict) and str(entry.get("owner") or "") == demo_user_id
            }
        for path in RUNS_DIR.iterdir() if RUNS_DIR.exists() else []:
            if not path.is_dir():
                continue
            if owned_run_ids is not None and path.name not in owned_run_ids:
                continue
            try:
                updated_at = int(path.stat().st_mtime * 1000)
            except OSError:
                continue
            run_paths.append((path, updated_at))
        run_paths.sort(key=lambda item: (item[1], item[0].name), reverse=True)
        known_conversation_ids = set(existing_conversation_ids or ()) | set(conversations)
        candidate_cutoff: int | None = None
        for path, updated_at in run_paths:
            if summary_min_updated_at is not None and updated_at < summary_min_updated_at:
                break
            if candidate_cutoff is not None and updated_at < candidate_cutoff:
                break
            run_id = path.name
            if run_id in run_ids_seen:
                continue
            conv = _conversation_from_disk_run(run_id, compact=True)
            if conv is None:
                continue
            conversations[conv["id"]] = conv
            known_conversation_ids.add(conv["id"])
            if (
                summary_min_updated_at is None
                and len(known_conversation_ids) >= summary_limit
                and candidate_cutoff is None
            ):
                candidate_cutoff = updated_at
    else:
        for run_id in _recent_disk_run_ids(limit=limit, demo_user_id=demo_user_id):
            if run_id in run_ids_seen:
                continue
            conv = _conversation_from_disk_run(run_id, compact=compact)
            if conv is not None:
                conversations[conv["id"]] = conv
    return conversations


def _read_design_session_events() -> list[dict[str, Any]]:
    design_dir = _BOOT_OUT_DIR / "design_sessions"
    if not design_dir.exists():
        return []
    events: list[dict[str, Any]] = []
    for path in sorted(design_dir.glob("*.jsonl")):
        try:
            lines = path.read_text(encoding="utf-8").splitlines()
        except OSError:
            continue
        for line in lines:
            if not line.strip():
                continue
            try:
                event = json.loads(line)
            except json.JSONDecodeError:
                continue
            if not isinstance(event, dict):
                continue
            event["_ts_ms"] = _event_ms(event, fallback_path=path)
            events.append(event)
    events.sort(key=lambda e: int(e.get("_ts_ms") or 0))
    return events


def _event_ms(event: dict[str, Any], *, fallback_path: Path | None = None) -> int:
    raw = str(event.get("ts") or "").strip()
    if raw:
        try:
            return int(datetime.fromisoformat(raw).timestamp() * 1000)
        except ValueError:
            pass
    if fallback_path and fallback_path.exists():
        return int(fallback_path.stat().st_mtime * 1000)
    return int(time.time() * 1000)


def _conversation_from_design_events(
    conv_id: str,
    events: list[dict[str, Any]],
    run_ids_seen: set[str],
    *,
    compact: bool = False,
) -> dict[str, Any] | None:
    if not events:
        return None
    events = sorted(events, key=lambda e: int(e.get("_ts_ms") or 0))
    created_at = int(events[0].get("_ts_ms") or int(time.time() * 1000))
    updated_at = int(events[-1].get("_ts_ms") or created_at)
    attachments_by_run = _attachments_by_run(events)
    messages: list[dict[str, Any]] = []
    artifacts: dict[str, Any] = {}
    active_artifact_id: str | None = None
    title = "New chat"
    seen_messages: set[str] = set()
    poster_palette_id: str | None = None
    poster_canvas_preset_id: str | None = None

    for event in events:
        name = str(event.get("event") or "")
        run_id = str(event.get("run_id") or "").strip()
        if run_id:
            run_ids_seen.add(run_id)
        ts = int(event.get("_ts_ms") or updated_at)
        data = event.get("data") if isinstance(event.get("data"), dict) else {}

        if name == "message.user_submitted":
            brief = str(data.get("brief") or "").strip() or "Submitted a request."
            event_artifact_type = str(data.get("artifact_type") or "").strip()
            event_palette_id = str(data.get("palette_id") or "").strip() or None
            event_template = str(data.get("template") or "").strip() or None
            event_canvas_preset_id = (
                str(data.get("canvas_preset_id") or "").strip() or None
            )
            attachments = attachments_by_run.get(run_id, [])
            content_refs = [
                item for item in attachments
                if item.get("role") != "style_reference"
            ]
            reference_ref = next(
                (
                    item for item in attachments
                    if item.get("role") == "style_reference"
                ),
                None,
            )
            if reference_ref is not None and not reference_ref.get("reference_handle"):
                reference_ref = {
                    key: reference_ref[key]
                    for key in ("name", "size", "kind", "role")
                    if key in reference_ref
                }
            if event_artifact_type == "poster" and event_palette_id:
                poster_palette_id = event_palette_id
            if event_artifact_type == "poster" and event_canvas_preset_id:
                poster_canvas_preset_id = event_canvas_preset_id
            if title == "New chat":
                title = brief[:50]
            msg_id = f"msg_import_user_{run_id or ts}"
            if msg_id not in seen_messages:
                messages.append({
                    "id": msg_id,
                    "role": "user",
                    "text": brief,
                    "ts": ts,
                    "status": "done",
                    "attachments": attachments,
                    "task_type": "generate",
                    "task_payload": {
                        "artifact_type": event_artifact_type or None,
                        "palette_id": event_palette_id,
                        "template": event_template,
                        "canvas_preset_id": event_canvas_preset_id,
                        **(
                            {"authoring_max_attempts": data["authoring_max_attempts"]}
                            if isinstance(data.get("authoring_max_attempts"), int)
                            else {}
                        ),
                        "attachment_refs": content_refs or None,
                        "reference_poster_ref": reference_ref,
                    },
                })
                seen_messages.add(msg_id)
            continue

        if name == "artifact.generated":
            artifact = (
                _history_artifact_preview_for_event(run_id, data)
                if compact
                else _history_artifact_for_event(run_id, data)
            )
            if artifact:
                artifacts[artifact["artifact_id"]] = artifact
                active_artifact_id = artifact["artifact_id"]
                msg_id = f"msg_{run_id}"
                if msg_id not in seen_messages:
                    message = {
                        "id": msg_id,
                        "role": "assistant",
                        "text": f"Generated {artifact['name']}.",
                        "ts": ts,
                        "run_id": run_id,
                        "artifact_id": artifact["artifact_id"],
                        "status": "done",
                    }
                    messages.append(message)
                    seen_messages.add(msg_id)
            continue

        if name == "edits.applied":
            event_palette_id = str(data.get("palette_id") or "").strip()
            if event_palette_id:
                poster_palette_id = event_palette_id
            artifact = (
                _history_artifact_preview_for_event(run_id, data)
                if compact
                else _history_artifact_for_event(run_id, data)
            )
            if artifact:
                artifacts[artifact["artifact_id"]] = artifact
                active_artifact_id = artifact["artifact_id"]
                msg_id = f"msg_{run_id}"
                if msg_id not in seen_messages:
                    messages.append({
                        "id": msg_id,
                        "role": "assistant",
                        "text": f"Applied edits and saved {artifact['name']}.",
                        "ts": ts,
                        "run_id": run_id,
                        "artifact_id": artifact["artifact_id"],
                        "status": "done",
                    })
                    seen_messages.add(msg_id)
            continue

        if name == "artifact.opened":
            artifact_id = str(event.get("artifact_id") or "")
            if artifact_id in artifacts:
                active_artifact_id = artifact_id
            continue

        if name in {
            "openresearch.project_requested",
            "openresearch.project_ready",
            "openresearch.project_failed",
        }:
            artifact_id = str(event.get("artifact_id") or data.get("artifact_id") or "")
            if not artifact_id and run_id:
                artifact_id = f"art_{run_id}"
            state = _openresearch_state_from_event(data, name)
            if artifact_id in artifacts and state:
                artifacts[artifact_id] = {
                    **artifacts[artifact_id],
                    "openresearch": state,
                }
            continue

        if name == "artifact.generation_failed":
            msg_id = f"msg_{run_id}"
            if run_id and msg_id not in seen_messages:
                status = str(data.get("status") or "error")
                failure_data = data.get("failure") if isinstance(data.get("failure"), dict) else {}
                artifact = (
                    _history_artifact_preview_for_event(run_id, data)
                    if compact
                    else _history_artifact_for_event(run_id, data)
                )
                if artifact and (compact or _artifact_dict_has_rendered_files(artifact)):
                    artifacts[artifact["artifact_id"]] = artifact
                    active_artifact_id = artifact["artifact_id"]
                    message = {
                        "id": msg_id,
                        "role": "assistant",
                        "text": f"Generated {artifact['name']}.",
                        "ts": ts,
                        "run_id": run_id,
                        "artifact_id": artifact["artifact_id"],
                        "status": "done",
                    }
                    messages.append(message)
                    seen_messages.add(msg_id)
                    continue
                messages.append({
                    "id": msg_id,
                    "role": "assistant",
                    "text": f"Run {run_id[:8]} failed: {status}",
                    "ts": ts,
                    "run_id": run_id,
                    "status": "error",
                    "failure": {
                        "status": status,
                        "phase": (
                            failure_data.get("phase")
                            or data.get("phase")
                            or _phase_from_disk(RUNS_DIR / run_id)
                        ),
                        "error_code": (
                            str(failure_data.get("error_code") or "")[:100]
                            or None
                        ),
                        "error_message": _redacted_error_detail(
                            str(failure_data.get("error_message") or ""),
                            max_chars=500,
                        ),
                        "error_detail": _redacted_error_detail(
                            str(failure_data.get("error_detail") or ""),
                            max_chars=1_200,
                        ),
                        "retry_route": (
                            failure_data.get("retry_route")
                            or data.get("retry_route")
                        ),
                        "parent_run_id": (
                            failure_data.get("parent_run_id")
                            or data.get("parent_run_id")
                        ),
                        "agent_last_note": failure_data.get("agent_last_note"),
                        "pointer_cleanup_warnings": list(
                            _validated_json_pointer_cleanup_warnings(
                                failure_data.get("pointer_cleanup_warnings")
                            )
                        ),
                        "produced_files": failure_data.get("produced_files") or _list_produced_artifacts(run_id),
                        "suggested_designer": (
                            failure_data.get("suggested_designer")
                            or failure_data.get("suggested_planner")
                        ),
                        "suggested_planner": (
                            failure_data.get("suggested_planner")
                            or failure_data.get("suggested_designer")
                        ),
                        "elapsed_ms": failure_data.get("elapsed_ms"),
                    },
                })
                seen_messages.add(msg_id)

    if not messages and not artifacts:
        return None
    if active_artifact_id is None and artifacts:
        active_artifact_id = next(reversed(artifacts))
    messages.sort(key=lambda m: int(m.get("ts") or 0))
    return {
        "id": conv_id,
        "title": title,
        "created_at": created_at,
        "updated_at": max(updated_at, *(int(m.get("ts") or 0) for m in messages), 0),
        "messages": messages,
        "artifacts": artifacts,
        "active_artifact_id": active_artifact_id,
        "pending_edits": {},
        "poster_palette_id": poster_palette_id,
        "poster_canvas_preset_id": poster_canvas_preset_id,
    }


def _attachments_by_run(events: list[dict[str, Any]]) -> dict[str, list[dict[str, Any]]]:
    out: dict[str, list[dict[str, Any]]] = {}
    for event in events:
        if event.get("event") != "attachment.added":
            continue
        run_id = str(event.get("run_id") or "").strip()
        data = event.get("data") if isinstance(event.get("data"), dict) else {}
        name = str(data.get("name") or "").strip()
        if not run_id or not name:
            continue
        suffix = str(data.get("suffix") or Path(name).suffix).lower()
        try:
            size = int(data.get("size") or 0)
        except (TypeError, ValueError):
            size = 0
        items = out.setdefault(run_id, [])
        attachment = {
            "id": f"att_{run_id}_{len(items)}",
            "name": name,
            "size": size,
            "kind": _attachment_kind_from_suffix(suffix),
        }
        if data.get("role") == "style_reference":
            attachment["role"] = "style_reference"
            reference_handle = str(data.get("reference_handle") or "").strip()
            if reference_handle:
                attachment["reference_handle"] = reference_handle
        items.append(attachment)
    return out


def _attachment_kind_from_suffix(suffix: str) -> str:
    ext = suffix.lower()
    if ext == ".pdf":
        return "pdf"
    if ext in {".png", ".jpg", ".jpeg", ".webp", ".gif"}:
        return "image"
    if ext in {".doc", ".docx", ".ppt", ".pptx", ".md", ".txt"}:
        return "doc"
    return "other"


def _history_run_control(run_id: str) -> tuple[bool, Any | None]:
    control_path = RUNS_DIR / run_id / "run_control.json"
    if not control_path.is_file():
        return False, None
    try:
        return True, RunControlStore(RUNS_DIR).read(run_id)
    except RunControlError:
        return True, None


def _history_control_allows_artifact(run_id: str) -> bool:
    controlled, record = _history_run_control(run_id)
    if not controlled:
        return True
    return bool(
        record is not None
        and record.state == "completed"
        and record.publishable
    )


def _history_artifact_for_event(run_id: str, data: dict[str, Any]) -> dict[str, Any] | None:
    if not run_id:
        return None
    if not _history_control_allows_artifact(run_id):
        return None
    a_type = str(data.get("artifact_type") or "")
    if a_type not in ("poster", "landing", "deck", "video"):
        detected = _detect_artifact_type_for_run(run_id)
        if detected is None:
            return None
        a_type = detected
    parent_id = data.get("parent_artifact_id") or data.get("parent_run_id")
    if isinstance(parent_id, str) and parent_id and not parent_id.startswith("art_"):
        parent_id = f"art_{parent_id}"
    baseline_json = (
        json.dumps({"artifact_id": parent_id})
        if isinstance(parent_id, str) and parent_id
        else None
    )
    artifact = _build_artifact_response(
        RUNS_DIR / run_id,
        run_id,
        a_type,  # type: ignore[arg-type]
        baseline_artifact_json=baseline_json,
    )
    if artifact is None:
        return None
    out = _dump_model(artifact)
    name = data.get("name")
    if isinstance(name, str) and name.strip():
        out["name"] = name.strip()
    return out


def _history_artifact_preview_for_event(
    run_id: str,
    data: dict[str, Any],
) -> dict[str, Any] | None:
    if not _history_control_allows_artifact(run_id):
        return None
    artifact_type = _history_artifact_type(data.get("artifact_type"))
    if artifact_type is None:
        artifact_type = _detect_artifact_type_for_run(run_id)
    if artifact_type is None:
        return None
    artifact = _history_artifact_preview_from_run(run_id, artifact_type)
    if artifact is None:
        return None
    name = data.get("name")
    if isinstance(name, str) and name.strip():
        artifact["name"] = name.strip()[:200]
    return artifact


def _history_artifact_preview_from_run(
    run_id: str,
    artifact_type: ArtifactType,
) -> dict[str, Any] | None:
    if not _history_control_allows_artifact(run_id):
        return None
    run_dir = RUNS_DIR / run_id
    final_dir = run_dir / "final"
    filename, _format, (width, height) = _FILE_FOR[artifact_type]
    if artifact_type == "video":
        if not _validated_video_delivery(run_dir).is_passed:
            return None
    else:
        target = final_dir / filename
        if not target.exists():
            if artifact_type == "deck" and (final_dir / "deck.pptx").exists():
                target = final_dir / "deck.pptx"
            else:
                target = _fallback_artifact(final_dir, _format)
        if target is None or not target.exists():
            return None
        fallback_type = _artifact_type_from_final_filename(target.name)
        if fallback_type is not None:
            artifact_type = fallback_type
            _filename, _format, (width, height) = _FILE_FOR[artifact_type]
    preview = {
        "artifact_id": f"art_{run_id}",
        "name": {
            "poster": "Poster",
            "landing": "Landing page",
            "deck": "Slide deck",
            "video": "Video",
        }[artifact_type] + f" - {run_id[:8]}",
        "artifact_type": artifact_type,
        "canvas": {"w": width, "h": height},
    }
    preview_path = final_dir / "preview.png"
    if preview_path.exists():
        preview["preview_url"] = _run_file_url(run_id, "final/preview.png")
    if artifact_type in {"poster", "deck", "landing"}:
        metadata = _validated_authored_delivery_metadata(
            run_dir,
            run_id,
            target,
            artifact_type,
        )
        preview.update(metadata)
    return preview


def _history_run_completion_from_disk(run_dir: Path) -> dict[str, Any]:
    completion: dict[str, Any] = {}
    telemetry = _read_json_file(run_dir / "run_telemetry_summary.json")
    if isinstance(telemetry, dict):
        for key in ("terminal_status", "run_done_wall_s"):
            value = telemetry.get(key)
            if value is not None:
                completion[key] = value

    events_path = run_dir / "run_events.jsonl"
    if not events_path.exists():
        return completion
    try:
        with events_path.open("r", encoding="utf-8") as handle:
            for line in handle:
                if not line.strip():
                    continue
                try:
                    event = json.loads(line)
                except json.JSONDecodeError:
                    continue
                if not isinstance(event, dict) or event.get("event") != "run.done":
                    continue
                for key in (
                    "terminal_status",
                    "critic_verdict",
                    "critic_score",
                    "wall_s",
                    "designer_model",
                ):
                    value = event.get(key)
                    if value is not None:
                        completion[key] = value
    except OSError:
        pass
    return completion


def _openresearch_state_from_event(data: dict[str, Any], event_name: str) -> dict[str, Any] | None:
    job_id = str(data.get("job_id") or "").strip()
    if not job_id:
        return None
    status = str(data.get("status") or "").strip()
    if not status:
        status = "running" if event_name.endswith("_requested") else "error"
    return {
        "status": status,
        "job_id": job_id,
        "result_url": data.get("result_url"),
        "api_log_url": data.get("api_log_url"),
        "project_id": data.get("project_id"),
        "project_url": data.get("project_url"),
        "org_id": data.get("org_id"),
        "paper_id": data.get("paper_id"),
        "repo_full_name": data.get("repo_full_name"),
        "gui_submitter_status": data.get("gui_submitter_status"),
        "gui_submitter_reason": data.get("gui_submitter_reason"),
        "gui_submitter_error": data.get("gui_submitter_error"),
        "gui_submitter_session_url": data.get("gui_submitter_session_url"),
        "agent_prompt_url": data.get("agent_prompt_url"),
        "submitter_log_url": data.get("submitter_log_url"),
        "latest_report_id": data.get("latest_report_id"),
        "latest_report_url": data.get("latest_report_url"),
        "error": data.get("error"),
    }


def _history_control_diagnostic_conversation(
    run_id: str,
    control: Any,
) -> dict[str, Any]:
    status = {
        "cancelled": "cancelled",
        "cancelling": "cancelling",
    }.get(control.state, "error")
    artifact_type = (
        control.artifact_type
        if control.artifact_type in {"poster", "landing", "deck", "video"}
        else "artifact"
    )
    label = {
        "poster": "Poster",
        "landing": "Landing page",
        "deck": "Slide deck",
        "video": "Video",
        "artifact": "Run",
    }[artifact_type]
    ts = int((control.terminal_at or control.updated_at) * 1000)
    text = {
        "cancelled": "Run cancelled. Its frozen output is diagnostic only.",
        "cancelling": "Cancellation is still being verified; no artifact is available.",
        "error": "Run ended without a publishable artifact.",
    }[status]
    persisted_diagnostics = (
        _failure_diagnostics_from_disk(RUNS_DIR / run_id)
        if status == "error"
        else {}
    )
    error_code = {
        "cancelled": "run_cancelled",
        "cancelling": "run_cancellation_pending",
        "error": persisted_diagnostics.get("error_code") or "run_failed",
    }[status]
    pointer_cleanup_warnings = (
        _derived_worker_pointer_cleanup_warnings(run_id)
        if status in {"cancelled", "cancelling"}
        else ()
    )
    return {
        "id": f"server_run_{run_id}",
        "title": f"{label} - {status}",
        "created_at": int(control.created_at * 1000),
        "updated_at": ts,
        "messages": [{
            "id": f"msg_{run_id}",
            "role": "assistant",
            "text": text,
            "ts": ts,
            "run_id": run_id,
            "status": "error",
            "failure": {
                "status": status,
                "phase": persisted_diagnostics.get("phase"),
                "error_code": error_code,
                "error_message": persisted_diagnostics.get("error_message"),
                "error_detail": persisted_diagnostics.get("error_detail"),
                "pointer_cleanup_warnings": list(pointer_cleanup_warnings),
                "produced_files": [],
            },
        }],
        "artifacts": {},
        "active_artifact_id": None,
        "pending_edits": {},
    }


def _conversation_from_disk_run(
    run_id: str,
    *,
    compact: bool = False,
) -> dict[str, Any] | None:
    run_dir = RUNS_DIR / run_id
    try:
        derived_descriptor = _read_derived_job_descriptor(run_id)
    except ValueError:
        derived_descriptor = None
    try:
        direct_publish_descriptor = _read_direct_candidate_publish_descriptor(
            run_id
        )
    except ValueError:
        direct_publish_descriptor = None
    candidate_publish = bool(
        derived_descriptor is not None
        and derived_descriptor.get("job_kind") == "candidate_publish"
    )
    candidate_lineage = _read_json_file(
        run_dir / "candidate_draft_lineage.json"
    )
    if (
        isinstance(candidate_lineage, dict)
        and candidate_lineage.get("source_run_id")
        and candidate_lineage.get("source_candidate_id")
        and not candidate_publish
    ):
        return None
    controlled, control = _history_run_control(run_id)
    if controlled and control is None:
        return None
    if (
        control is not None
        and control.state in {"completed", "failed", "cancelled", "cancelling"}
        and not (control.state == "completed" and control.publishable)
    ):
        return _history_control_diagnostic_conversation(run_id, control)
    recoverable = _disk_run_is_recoverable(run_dir) or bool(
        candidate_publish
        and control is not None
        and control.state in {
            "reserved",
            "uploading",
            "queued",
            "running",
            "completing",
        }
    )
    terminal, _owner_pid, started = _disk_run_lifecycle(run_dir)
    declared_type = _declared_artifact_type_for_run(run_id)
    if recoverable:
        state = _RUNS.get(run_id)
        state_type = str(getattr(state, "artifact_type", "") or "")
        descriptor_type = str(
            derived_descriptor.get("artifact_type") or ""
        ) if derived_descriptor is not None else ""
        a_type = (
            state_type
            if state_type in {"poster", "landing", "deck", "video"}
            else descriptor_type
            if descriptor_type in {"poster", "landing", "deck", "video"}
            else declared_type or _detect_artifact_type_for_run(run_id)
        )
        if a_type is None:
            return None
        ts = int(run_dir.stat().st_mtime * 1000)
        label = {
            "poster": "Poster",
            "landing": "Landing page",
            "deck": "Slide deck",
            "video": "Video",
        }[a_type]
        conversation_id = (
            str(derived_descriptor.get("conversation_id") or "").strip()
            if candidate_publish and derived_descriptor is not None
            else ""
        ) or f"server_run_{run_id}"
        task_type = "generate"
        task_payload: dict[str, Any] = {"artifact_type": a_type}
        run_brief = _read_json_file(run_dir / "run_brief.json")
        if isinstance(run_brief, dict):
            persisted_template = str(
                run_brief.get("effective_template") or ""
            ).strip()
            persisted_canvas_preset_id = str(
                run_brief.get("canvas_preset_id") or ""
            ).strip()
            if persisted_template:
                task_payload["template"] = persisted_template
            if persisted_canvas_preset_id:
                task_payload["canvas_preset_id"] = persisted_canvas_preset_id
        if candidate_publish and derived_descriptor is not None:
            task_type = "candidate_publish"
            if direct_publish_descriptor is not None:
                task_payload.update({
                    "source_artifact_id": str(
                        derived_descriptor.get("source_artifact_id") or ""
                    ),
                    "source_run_id": direct_publish_descriptor["source_run_id"],
                    "source_candidate_id": direct_publish_descriptor[
                        "source_candidate_id"
                    ],
                })
            else:
                source_lineage = candidate_lineage
                if not isinstance(source_lineage, dict):
                    source_lineage = _read_json_file(
                        RUNS_DIR
                        / str(derived_descriptor.get("parent_run_id") or "")
                        / "candidate_draft_lineage.json"
                    )
                source_lineage = (
                    source_lineage if isinstance(source_lineage, dict) else {}
                )
                task_payload.update({
                    "source_artifact_id": str(
                        derived_descriptor.get("source_artifact_id") or ""
                    ),
                    "source_run_id": str(
                        source_lineage.get("source_run_id") or ""
                    ),
                    "source_candidate_id": str(
                        source_lineage.get("source_candidate_id") or ""
                    ),
                })
        return {
            "id": conversation_id,
            "title": f"{label} - running",
            "created_at": ts,
            "updated_at": ts,
            "messages": [{
                "id": f"msg_{run_id}",
                "role": "assistant",
                "text": "",
                "ts": ts,
                "run_id": run_id,
                "status": "streaming",
                "task_type": task_type,
                "task_payload": task_payload,
            }],
            "artifacts": {},
            "active_artifact_id": None,
            "pending": True,
            "run_id": run_id,
            "pending_edits": {},
        }

    if (
        control is not None
        and not (control.state == "completed" and control.publishable)
    ):
        return _history_control_diagnostic_conversation(run_id, control)

    # A started but interrupted run can leave an intermediate artifact in
    # final/. Honor its declared target type so a video scaffold is not
    # imported as a completed landing page after the owner exits.
    a_type = (
        declared_type
        if started and terminal is None and declared_type is not None
        else _detect_artifact_type_for_run(run_id)
    )
    if a_type is None:
        return None
    if compact:
        art = _history_artifact_preview_from_run(run_id, a_type)
    else:
        artifact = _build_artifact_response(
            run_dir,
            run_id,
            a_type,
            baseline_artifact_json=None,
        )
        art = _dump_model(artifact) if artifact is not None else None
    if art is None:
        return None
    run_brief = _read_json_file(run_dir / "run_brief.json")
    palette_id = (
        str(run_brief.get("palette_id") or "").strip()
        if isinstance(run_brief, dict)
        else ""
    )
    canvas_preset_id = (
        str(run_brief.get("canvas_preset_id") or "").strip()
        if isinstance(run_brief, dict)
        else ""
    )
    if not palette_id:
        for manifest_name in (
            "code_editor_revision_manifest.json",
            "authored_poster_edit_manifest.json",
            "apply_edits_palette_manifest.json",
        ):
            manifest = _read_json_file(run_dir / "final" / manifest_name)
            if isinstance(manifest, dict):
                palette_id = str(manifest.get("palette_id") or "").strip()
            if palette_id:
                break
    ts = int(run_dir.stat().st_mtime * 1000)
    message = {
        "id": f"msg_{run_id}",
        "role": "assistant",
        "text": f"Imported {art['name']} from out/runs.",
        "ts": ts,
        "run_id": run_id,
        "artifact_id": art["artifact_id"],
        "status": "done",
    }
    conversation = {
        "id": f"server_run_{run_id}",
        "title": art["name"],
        "created_at": ts,
        "updated_at": ts,
        "messages": [message],
        "artifacts": {art["artifact_id"]: art},
        "active_artifact_id": art["artifact_id"],
        "pending_edits": {},
    }
    if palette_id:
        conversation["poster_palette_id"] = palette_id
    if canvas_preset_id:
        conversation["poster_canvas_preset_id"] = canvas_preset_id
    return conversation


def _recent_disk_run_ids(*, limit: int | None, demo_user_id: str | None = None) -> list[str]:
    if not RUNS_DIR.exists():
        return []
    dirs = [p for p in RUNS_DIR.iterdir() if p.is_dir()]
    dirs.sort(key=lambda p: p.stat().st_mtime, reverse=True)
    eligible = [
        path for path in dirs
        if not demo_user_id or _demo_user_owns_run(path.name, demo_user_id)
    ]
    active_run_ids = [
        path.name for path in eligible
        if _disk_run_is_recoverable(path)
    ]
    run_ids: list[str] = list(active_run_ids)
    for path in dirs:
        if demo_user_id and not _demo_user_owns_run(path.name, demo_user_id):
            continue
        if path.name not in run_ids:
            run_ids.append(path.name)
        if limit is not None and len(run_ids) >= len(active_run_ids) + limit:
            break
    return run_ids


def _history_run_id_from_conversation_id(conversation_id: str) -> str | None:
    prefix = "server_run_"
    if not conversation_id.startswith(prefix):
        return None
    run_id = conversation_id[len(prefix):]
    if (
        not run_id
        or run_id in {".", ".."}
        or Path(run_id).name != run_id
    ):
        return None
    return run_id


def _history_conversation_from_server_sources(
    conversation_id: str,
    *,
    include_design_sessions: bool,
    demo_user_id: str | None,
) -> dict[str, Any] | None:
    conversations: dict[str, Any] = {}
    run_id = _history_run_id_from_conversation_id(conversation_id)
    if run_id is not None:
        if demo_user_id and not _demo_user_owns_run(run_id, demo_user_id):
            return None
        conversation = _conversation_from_disk_run(run_id)
        if conversation is not None:
            conversations[conversation_id] = conversation
    if include_design_sessions and demo_user_id is None:
        events = [
            event
            for event in _read_design_session_events()
            if str(event.get("conversation_id") or "").strip() == conversation_id
        ]
        conversation = _conversation_from_design_events(conversation_id, events, set())
        if conversation is not None:
            conversations[conversation_id] = conversation
    merged = _merge_history_conversations({}, conversations)
    return merged.get(conversation_id)


def _detect_artifact_type_for_run(run_id: str) -> ArtifactType | None:
    run_dir = RUNS_DIR / run_id
    final_dir = run_dir / "final"
    if _validated_video_delivery(run_dir).is_passed:
        return "video"
    final_type = _artifact_type_from_final_dir(final_dir)
    if final_type is not None:
        return final_type
    return _declared_artifact_type_for_run(run_id)


def _declared_artifact_type_for_run(run_id: str) -> ArtifactType | None:
    run_dir = RUNS_DIR / run_id
    for name in ("canvas_plan.json", "resume_state.json"):
        data = _read_json_file(run_dir / name)
        artifact_type = str(data.get("artifact_type") or "") if isinstance(data, dict) else ""
        if artifact_type in {"poster", "landing", "deck", "video"}:
            return artifact_type  # type: ignore[return-value]
    return None


def _artifact_type_from_final_dir(final_dir: Path) -> ArtifactType | None:
    if (final_dir / "poster.html").exists():
        return "poster"
    if (final_dir / "index.html").exists():
        return "landing"
    if (final_dir / "deck.html").exists() or (final_dir / "deck.pptx").exists():
        return "deck"
    html_files = list(final_dir.glob("*.html")) if final_dir.exists() else []
    if html_files:
        return "poster"
    return None


def _artifact_type_from_final_filename(name: str) -> ArtifactType | None:
    lowered = name.lower()
    if lowered == "poster.html":
        return "poster"
    if lowered == "index.html":
        return "landing"
    if lowered in {"deck.html", "deck.pptx"}:
        return "deck"
    return None


def _dump_model(model: BaseModel) -> dict[str, Any]:
    if hasattr(model, "model_dump"):
        return model.model_dump()  # type: ignore[no-any-return]
    return model.dict()  # type: ignore[no-any-return]


def _read_json_file(path: Path) -> Any:
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except Exception:
        return None


def _run_id_from_artifact_id(artifact_id: str | None) -> str | None:
    if not artifact_id or not artifact_id.startswith("art_"):
        return None
    return artifact_id[len("art_"):] or None


def _run_id_from_maybe_artifact_ref(value: Any) -> str | None:
    if not isinstance(value, str) or not value.strip():
        return None
    clean = value.strip()
    return clean[len("art_"):] if clean.startswith("art_") else clean


def _source_poster_html_path(run_id: str) -> Path | None:
    final_dir = RUNS_DIR / run_id / "final"
    target = final_dir / "poster.html"
    if target.exists():
        return target
    return _fallback_artifact(final_dir, "html")


def _poster_revision_context_run_dirs(
    source_run_id: str,
    artifact: dict[str, Any],
    *,
    authorize_run: Callable[[str], None] | None = None,
) -> list[Path]:
    seen: set[str] = set()
    out: list[Path] = []
    current: str | None = source_run_id
    first = True
    while current and current not in seen and len(out) < 12:
        seen.add(current)
        if authorize_run is not None:
            authorize_run(current)
        run_dir = RUNS_DIR / current
        if run_dir.exists():
            out.append(run_dir)
        current = _poster_revision_parent_run_id(
            run_dir,
            artifact if first else {},
        )
        first = False
    return out


def _poster_revision_parent_run_id(run_dir: Path, artifact: dict[str, Any]) -> str | None:
    for raw in (
        artifact.get("parent_run_id"),
        artifact.get("parent_artifact_id"),
    ):
        parent = _run_id_from_maybe_artifact_ref(raw)
        if parent:
            return parent
    for name in ("code_editor_revision_manifest.json", "authored_poster_edit_manifest.json"):
        manifest = _read_json_file(run_dir / "final" / name)
        if not isinstance(manifest, dict):
            continue
        for key in ("parent_run_id", "parent_artifact_id"):
            parent = _run_id_from_maybe_artifact_ref(manifest.get(key))
            if parent:
                return parent
    return None


_ARTIFACT_ASSET_EXTS = {".png", ".jpg", ".jpeg", ".webp"}
_ARTIFACT_ASSET_EXCLUDE_RE = re.compile(
    r"(?:preview|contact[_-]?sheet|identity|logo|brand|overlay|candidate|"
    r"poster[_-]?check|attempt[_-]?preview|validation)",
    re.I,
)
_ARTIFACT_ASSET_INCLUDE_RE = re.compile(
    r"(?:^img_)?ingest_(?:fig|figure|table)|(?:^|_)tbl\d*|(?:^|_)xref\d+",
    re.I,
)


def _collect_artifact_assets(run_dirs: list[Path]) -> list[ArtifactAsset]:
    assets: list[ArtifactAsset] = []
    seen: set[str] = set()
    for run_dir in run_dirs:
        for asset_dir in (run_dir / "final" / "layers", run_dir / "layers"):
            if not asset_dir.exists():
                continue
            for path in sorted(asset_dir.iterdir(), key=lambda p: _asset_natural_key(p.name)):
                if not _is_paper_asset_file(path):
                    continue
                key = _asset_dedupe_key(path)
                if key in seen:
                    continue
                seen.add(key)
                run_id = _run_id_for_path(path)
                if not run_id:
                    continue
                rel = _run_relative_path(path)
                kind = _artifact_asset_kind(path)
                assets.append(
                    ArtifactAsset(
                        asset_id=f"{run_id}:{rel}",
                        name=_artifact_asset_name(path, kind),
                        kind=kind,
                        url=_run_file_url(run_id, rel),
                        filename=path.name,
                        run_id=run_id,
                        source=rel,
                        size=path.stat().st_size,
                    )
                )
    return assets[:96]


def _is_paper_asset_file(path: Path) -> bool:
    if not path.is_file() or path.suffix.lower() not in _ARTIFACT_ASSET_EXTS:
        return False
    name = path.name.lower()
    if _ARTIFACT_ASSET_EXCLUDE_RE.search(name):
        return False
    parts = {part.lower() for part in path.parts}
    if "layers" not in parts:
        return False
    return bool(_ARTIFACT_ASSET_INCLUDE_RE.search(path.stem))


def _artifact_asset_kind(path: Path) -> Literal["figure", "table", "image"]:
    stem = path.stem.lower()
    if "table" in stem or re.search(r"(?:^|_)tbl\d*", stem):
        return "table"
    if "fig" in stem or "figure" in stem or re.search(r"(?:^|_)xref\d+", stem):
        return "figure"
    return "image"


def _artifact_asset_name(path: Path, kind: Literal["figure", "table", "image"]) -> str:
    stem = path.stem
    if match := re.search(r"(?:^img_)?ingest_fig_(\d+)(?:_([a-z]))?$", stem, re.I):
        suffix = match.group(2) or ""
        return f"Figure {int(match.group(1))}{suffix}"
    if match := re.search(r"(?:^img_)?ingest_table_(\d+)(?:_([a-z]))?$", stem, re.I):
        suffix = match.group(2) or ""
        return f"Table {int(match.group(1))}{suffix}"
    if match := re.search(r"p(\d+)_tbl(\d+)", stem, re.I):
        return f"Table crop p{match.group(1)}-{match.group(2)}"
    if match := re.search(r"p(\d+)_xref(\d+)", stem, re.I):
        return f"Figure crop p{match.group(1)}-{match.group(2)}"
    label = {"figure": "Figure", "table": "Table", "image": "Paper crop"}[kind]
    clean = re.sub(r"[_-]+", " ", stem).strip()
    return f"{label} {clean}" if clean else label


def _asset_dedupe_key(path: Path) -> str:
    stem = re.sub(r"^img_", "", path.stem.lower())
    return f"{stem}{path.suffix.lower()}"


def _asset_natural_key(name: str) -> list[int | str]:
    return [int(part) if part.isdigit() else part.lower() for part in re.split(r"(\d+)", name)]


def _final_render_complete_on_disk(run_id: str, a_type: str | None) -> bool:
    """Return True when the run's primary end-of-pipeline file is on
    disk. Used to recover artifacts whose run was marked abort/max_turns
    but actually finished rendering (planner forgot the final finalize
    call). Per-type signal:

    - video: exact MP4 bound to the newest passed delivery manifest
    - poster: final/preview.png (last file _composite_poster writes) or
      final/poster.svg
    - landing: final/index.html
    - deck: final/deck.html (canonical) or legacy final/deck.pptx
    """
    if not run_id:
        return False
    run_dir = RUNS_DIR / run_id
    final_dir = run_dir / "final"
    if a_type == "video":
        return _validated_video_delivery(run_dir).is_passed
    if not final_dir.exists():
        return False
    if a_type == "poster":
        return (final_dir / "preview.png").exists() or (final_dir / "poster.svg").exists()
    if a_type == "landing":
        return (final_dir / "index.html").exists()
    if a_type == "deck":
        return (final_dir / "deck.html").exists() or (final_dir / "deck.pptx").exists()
    return False


def _artifact_dict_has_rendered_files(artifact: dict[str, Any]) -> bool:
    a_type = artifact.get("artifact_type")
    if a_type == "video":
        # URL-only fast path (avoids disk hit when the artifact was
        # built with a fully-formed mp4 url already).
        native_url = artifact.get("native_file_url")
        if artifact.get("native_format") == "mp4" and isinstance(native_url, str):
            return True
        download_url = artifact.get("download_url")
        if isinstance(download_url, str) and ".mp4" in download_url.lower():
            return True
        project = artifact.get("video_project")
        latest = project.get("latest_render") if isinstance(project, dict) else None
        if isinstance(latest, dict):
            mp4_url = latest.get("mp4_url")
            if isinstance(mp4_url, str) and ".mp4" in mp4_url.lower():
                return True
    run_id = _run_id_from_artifact_id(artifact.get("artifact_id"))
    return _final_render_complete_on_disk(run_id or "", a_type)


# Hint to the planner what kind of output the user asked for. The agent
# is brief-driven (no artifact_type parameter on PipelineRunner.run), so
# we encode the click on a quick-action pill as a leading line.
TYPE_PROLOGUE: dict[ArtifactType, str] = {
    "poster": "Type: poster (single-page, print-ready, fixed-size canvas).",
    "landing": "Type: landing page (responsive HTML, web hero with sections).",
    "deck": "Type: slide deck (HTML presentation with native editable text and paper visuals).",
    "video": (
        "Type: video (HTML-first 1920x1080 H.264 conference video, 300-600 s, "
        "with duration selected for paper complexity).\n"
        "Build 10-14 source-grounded scenes directly in the video DesignSpec, "
        "then call export_video. Include English narration and SRT/VTT subtitles; "
        "do not create an intermediate landing page."
    ),
}


def _apply_type_prologue(brief: str, a_type: ArtifactType) -> str:
    return f"{TYPE_PROLOGUE[a_type]}\n\n{brief.strip()}\n"


# ---------- Conversation memory: prior turns + artifacts → preamble ----------

# Hard cap on how many turns we replay. Each turn ~80–200 chars after
# truncation, so 6 turns ≈ 600–1200 chars (~250 tokens). Cheap, but
# enough that "make slides for the same content" works.
_HISTORY_TURNS_MAX = 6
# Per-message char cap for the preamble — the planner doesn't need
# verbatim assistant text, just enough to identify what came before.
_HISTORY_TEXT_MAX = 220


def _apply_conversation_prologue(
    brief: str,
    *,
    history_json: str | None,
    artifacts_json: str | None,
) -> str:
    """Prepend a small "what came before in this thread" block to the
    brief. Schema (set by the frontend, see web/src/lib/store.ts):
        history_json: '[{"role":"user","text":"..."},
                        {"role":"assistant","text":"...","artifact_id":"art_..."}]'
        artifacts_json: '[{"artifact_id":"art_...","name":"...","type":"poster",
                          "canvas":{"w":1200,"h":1600},"native_format":"html"}]'
    Both fields optional. We swallow malformed JSON silently — memory is
    a quality-of-life feature, not a correctness one.
    """
    history = _safe_load_json_list(history_json)
    artifacts = _safe_load_json_list(artifacts_json)
    if not history and not artifacts:
        return brief
    lines: list[str] = []
    lines.append("[Conversation context — your prior turns in this thread:]")
    if artifacts:
        # Surface produced artifacts even when the user's chat text was
        # short ("ok", "make it nicer") — the artifact list is what
        # actually shows shared style across turns.
        lines.append("Artifacts already produced in this thread:")
        for a in artifacts[-_HISTORY_TURNS_MAX:]:
            if not isinstance(a, dict):
                continue
            name = str(a.get("name", "?"))[:60]
            atype = str(a.get("type", "?"))
            canvas = a.get("canvas") or {}
            cw = canvas.get("w") if isinstance(canvas, dict) else None
            ch = canvas.get("h") if isinstance(canvas, dict) else None
            geom = f" {cw}x{ch}" if cw and ch else ""
            fmt = a.get("native_format")
            tag = f".{fmt}" if fmt else ""
            lines.append(f"  - {atype}{tag}: \"{name}\"{geom}")
    if history:
        lines.append("Recent turns:")
        for t in history[-_HISTORY_TURNS_MAX:]:
            if not isinstance(t, dict):
                continue
            role = str(t.get("role", "?")).lower()
            text = str(t.get("text", "")).strip().replace("\n", " ")
            if not text:
                continue
            if len(text) > _HISTORY_TEXT_MAX:
                text = text[:_HISTORY_TEXT_MAX - 1] + "…"
            who = "User" if role == "user" else "You" if role == "assistant" else role
            lines.append(f"  • {who}: {text}")
    lines.append(
        "Use the above as background only. Treat the user's NEXT message as "
        "the actual request; reuse style/content from prior turns where it "
        "makes sense (e.g. matching palette across artifacts).")
    lines.append("")
    lines.append("[User's current request:]")
    preamble = "\n".join(lines)
    return f"{preamble}\n{brief.strip()}\n"


def _safe_load_json_list(s: str | None) -> list[Any]:
    if not s:
        return []
    try:
        v = json.loads(s)
        return v if isinstance(v, list) else []
    except (json.JSONDecodeError, TypeError):
        return []


def _count_history_turns(s: str | None) -> int:
    return len(_safe_load_json_list(s))


def _count_prior_artifacts(s: str | None) -> int:
    return len(_safe_load_json_list(s))


_ARTIFACT_TYPES: tuple[ArtifactType, ...] = ("poster", "landing", "deck", "video")


def _coerce_artifact_type(
    raw: str | None,
    *,
    brief: str | None = None,
) -> ArtifactType:
    if raw not in _ARTIFACT_TYPES:
        inferred = _infer_artifact_type_from_brief(brief or "")
        if inferred is not None:
            return inferred
        # No pill clicked — default to poster (cheapest, fastest).
        return "poster"
    return raw  # type: ignore[return-value]


def _coerce_result_artifact_type(
    result: RunResult,
    *,
    fallback: ArtifactType,
) -> ArtifactType:
    if fallback == "video":
        return "video"
    raw = str(getattr(result, "artifact_type", "") or "")
    return raw if raw in _ARTIFACT_TYPES else fallback  # type: ignore[return-value]


def _infer_artifact_type_from_brief(brief: str) -> ArtifactType | None:
    text = (brief or "").strip().lower()
    if not text:
        return None
    compact = re.sub(r"\s+", " ", text)
    if any(k in compact for k in (
        "video", "mp4", "movie", "动画", "视频", "讲解视频", "演示视频",
    )):
        return "video"
    if any(k in compact for k in (
        "slides", "slide deck", "deck", "ppt", "pptx", "keynote",
        "presentation", "talk deck", "幻灯片", "演示文稿", "汇报", "路演",
    )):
        return "deck"
    landing_markers = (
        "landing page", "project page", "paper page", "web page",
        "homepage", "website", "site", "网页", "页面", "网站", "官网",
        "主页", "项目页", "介绍页",
    )
    if any(k in compact for k in landing_markers):
        return "landing"
    # In this app, users often say "page" as shorthand for a generated
    # project/landing page. Keep explicit poster requests on poster.
    if "page" in compact and not any(k in compact for k in ("poster", "海报")):
        return "landing"
    if any(k in compact for k in ("poster", "海报")):
        return "poster"
    return None


# Per-type artifact filenames inside <run_dir>/final/. The composite tools
# write these via stable symlinks; see autodesign/tools/composite.py.
# Video is a sentinel — the actual MP4 lives under
# `out/runs/<id>/hyperframes-<vid>/renders/*.mp4`, resolved via glob in
# `_build_artifact_response` rather than a fixed path.
_FILE_FOR: dict[ArtifactType, tuple[str, NativeFormat, tuple[int, int]]] = {
    "poster":  ("poster.html",   "html", (1200, 1600)),
    "landing": ("index.html",    "html", (1440, 900)),
    "deck":    ("deck.html",     "html", (1920, 1080)),
    "video":   ("__video__.mp4", "mp4",  (1920, 1080)),
}


def _ensure_authored_html_edit_contract(
    final_dir: Path,
    target: Path,
    artifact_type: Literal["deck", "landing"],
) -> None:
    """Migrate authored HTML once while keeping direct-final hashes valid."""

    if artifact_type == "deck":
        canonical_path = (
            final_dir / "deck.html"
            if (final_dir / "deck.html").is_file()
            else target
        )
        changed = ensure_editable_html_contract(canonical_path, "deck").changed
        for alias_name in ("deck.html", "slides.html"):
            alias_path = final_dir / alias_name
            if alias_path == canonical_path:
                continue
            if not alias_path.is_file() or sha256_file(alias_path) != sha256_file(canonical_path):
                shutil.copy2(canonical_path, alias_path)
                changed = True
    else:
        changed = ensure_editable_html_contract(target, "landing").changed
    if not changed:
        return

    _sync_authored_html_manifest_hashes(final_dir, target, artifact_type)


def _sync_authored_html_manifest_hashes(
    final_dir: Path,
    target: Path,
    artifact_type: Literal["deck", "landing"],
) -> None:
    """Keep copied or migrated direct-final manifests tied to current HTML."""

    canonical_path = (
        final_dir / "deck.html"
        if artifact_type == "deck" and (final_dir / "deck.html").is_file()
        else target
    )
    current_hash = sha256_file(canonical_path)
    manifest_names = (
        ("slides_author_manifest.json", "authored_html_edit_manifest.json")
        if artifact_type == "deck"
        else ("landing_author_manifest.json", "authored_html_edit_manifest.json")
    )
    for manifest_name in manifest_names:
        manifest_path = final_dir / manifest_name
        manifest = _read_json_file(manifest_path)
        if not isinstance(manifest, dict) or "html_sha256" not in manifest:
            continue
        if manifest.get("html_sha256") == current_hash:
            continue
        manifest["html_sha256"] = current_hash
        atomic_write_json(manifest_path, manifest)


def _build_artifact_response(
    run_dir: Path, run_id: str, a_type: ArtifactType,
    *, baseline_artifact_json: str | None,
) -> Artifact | None:
    final_dir = run_dir / "final"
    fname, fmt, (cw, ch) = _FILE_FOR[a_type]
    canvas_plan = _read_canvas_plan(run_dir)
    deck_plan = _read_deck_plan(run_dir)

    # Video is special: the rendered MP4 lives under
    # `<run_dir>/hyperframes-<vid>/renders/*.mp4`, not in final/.
    if a_type == "video":
        return _build_video_artifact(
            run_dir, run_id, baseline_artifact_json=baseline_artifact_json,
        )

    target = final_dir / fname
    if not target.exists():
        # Designer did not produce the expected artifact (e.g. wrote a
        # different type than requested, or aborted before composite).
        # Fall back: if a single .html / .pptx exists in final/, use that.
        target = _fallback_artifact(final_dir, fmt)
        if target is None and a_type == "deck" and (final_dir / "deck.pptx").exists():
            target = final_dir / "deck.pptx"
            fmt = "pptx"
        if target is None:
            log("web.artifact.missing", run_id=run_id, expected=str(final_dir / fname))
            return None
        fname = target.name
        fallback_type = _artifact_type_from_final_filename(fname)
        if fallback_type is not None and fallback_type != a_type:
            log(
                "web.artifact_type.corrected",
                run_id=run_id,
                requested=a_type,
                detected=fallback_type,
                filename=fname,
            )
            a_type = fallback_type
            _, fmt, (cw, ch) = _FILE_FOR[a_type]
            if fname.lower().endswith(".pptx"):
                fmt = "pptx"

    # Supervised artifacts must already carry their editable contract inside
    # the worker-owned promotion. Web response construction is read-only so it
    # cannot race cancellation or mutate a frozen diagnostic snapshot.
    if (
        fmt == "html"
        and a_type in {"deck", "landing"}
        and not (run_dir / "run_control.json").is_file()
    ):
        _ensure_authored_html_edit_contract(final_dir, target, a_type)

    preview_path = final_dir / "preview.png"
    preview_url = (
        _run_file_url(run_id, "final/preview.png")
        if preview_path.exists()
        else None
    )
    delivery_metadata = _validated_authored_delivery_metadata(
        run_dir,
        run_id,
        target,
        a_type,
    )

    parent_id: str | None = None
    if baseline_artifact_json:
        try:
            parent_id = json.loads(baseline_artifact_json).get("artifact_id")
        except (json.JSONDecodeError, AttributeError):
            parent_id = None

    file_url = _run_file_url(run_id, f"final/{fname}")
    if a_type == "poster" and fmt == "html" and _is_authored_paper_poster_html(target):
        cw, ch = _authored_paper_poster_size(target)
    pdf_url = (
        _run_file_url(run_id, "final/deck.pdf")
        if a_type == "deck" and (final_dir / "deck.pdf").exists()
        else None
    )
    downloads = {fmt: file_url}
    if pdf_url:
        downloads["pdf"] = pdf_url
    view_file_url = file_url if fmt != "pptx" else None
    view_format: ViewFormat | None = fmt if fmt != "pptx" else None
    if a_type == "deck" and fmt == "html":
        deck_html = final_dir / "deck.html"
        if deck_html.exists():
            editable = _build_deck_layer_artifact(
                deck_html=deck_html,
                run_id=run_id,
                file_url=file_url,
                preview_url=preview_url,
                card_preview_url=delivery_metadata.get("card_preview_url"),
                quality_status=delivery_metadata.get("quality_status"),
                quality_diagnostics=delivery_metadata.get("quality_diagnostics", []),
                parent_id=parent_id,
                pdf_url=pdf_url,
                downloads=downloads,
                canvas_plan=canvas_plan,
                deck_plan=deck_plan,
            )
            if editable is not None:
                return editable
            view_file_url = _run_file_url(run_id, "final/deck.html")
            view_format = "html"
    pretty = {"poster": "Poster", "landing": "Landing page",
              "deck": "Slide deck", "video": "Video"}[a_type]

    # HTML artifacts carry an authoritative layer manifest in their data-*
    # attributes (the same contract apply_edits.py reads). Parsing it here
    # lets the right-hand Sidebar drive real per-layer edits without us
    # adding a sidecar JSON. Deck (.pptx) has no such substrate — leave [].
    layers: list[dict[str, Any]] = (
        parse_html_layers(target) if fmt == "html" else []
    )

    return Artifact(
        artifact_id=f"art_{run_id}",
        name=f"{pretty} — {run_id[:8]}",
        artifact_type=a_type,
        canvas=Canvas(w=cw, h=ch),
        canvas_plan=canvas_plan,
        deck_plan=deck_plan,
        native_file_url=file_url,
        native_format=fmt,
        view_file_url=view_file_url,
        view_format=view_format,
        download_url=file_url,
        pdf_url=pdf_url,
        downloads=downloads,
        preview_url=preview_url,
        card_preview_url=delivery_metadata.get("card_preview_url"),
        quality_status=delivery_metadata.get("quality_status"),
        quality_diagnostics=delivery_metadata.get("quality_diagnostics", []),
        layers=layers,
        openresearch=_latest_openresearch_artifact_state(run_dir, run_id),
        parent_artifact_id=parent_id,
    )


def _latest_openresearch_artifact_state(run_dir: Path, source_run_id: str) -> dict[str, Any] | None:
    root = run_dir / "openresearch"
    if not root.exists():
        return None
    run_results = sorted(
        root.glob("*/openresearch_project_result.json"),
        key=lambda p: p.stat().st_mtime if p.exists() else 0,
        reverse=True,
    )
    for path in run_results:
        payload = _read_json_file(path)
        if not isinstance(payload, dict):
            continue
        job_id = str(payload.get("job_id") or path.parent.name)
        payload.setdefault("job_id", job_id)
        payload.setdefault("source_run_id", source_run_id)
        payload.setdefault("artifact_id", f"art_{source_run_id}")
        payload.setdefault("result_url", _openresearch_file_url(source_run_id, job_id, "openresearch_project_result.json"))
        payload.setdefault("api_log_url", _openresearch_file_url(source_run_id, job_id, "openresearch_api.jsonl"))
        payload.setdefault("agent_prompt_url", _openresearch_file_url(source_run_id, job_id, OPENRESEARCH_AGENT_PROMPT_FILE))
        payload.setdefault("submitter_log_url", _openresearch_file_url(source_run_id, job_id, OPENRESEARCH_GUI_PROCESS_FILE))
        return {
            "status": payload.get("status") or "submitted",
            "job_id": job_id,
            "result_url": payload.get("result_url"),
            "api_log_url": payload.get("api_log_url"),
            "project_id": payload.get("project_id"),
            "project_url": payload.get("project_url"),
            "org_id": payload.get("org_id"),
            "paper_id": payload.get("paper_id"),
            "repo_full_name": payload.get("repo_full_name"),
            "gui_submitter_status": payload.get("gui_submitter_status"),
            "gui_submitter_reason": payload.get("gui_submitter_reason"),
            "gui_submitter_error": payload.get("gui_submitter_error"),
            "gui_submitter_session_url": payload.get("gui_submitter_session_url"),
            "agent_prompt_url": payload.get("agent_prompt_url"),
            "submitter_log_url": payload.get("submitter_log_url"),
            "latest_report_id": payload.get("latest_report_id"),
            "latest_report_url": payload.get("latest_report_url"),
            "error": payload.get("error"),
        }
    return None


def _build_deck_layer_artifact(
    *,
    deck_html: Path,
    run_id: str,
    file_url: str,
    preview_url: str | None,
    card_preview_url: str | None,
    quality_status: Literal["ready", "ready_with_warnings"] | None,
    quality_diagnostics: list[str],
    parent_id: str | None,
    pdf_url: str | None,
    downloads: dict[str, str],
    canvas_plan: dict[str, Any],
    deck_plan: dict[str, Any],
) -> Artifact | None:
    parsed = parse_deck_html_as_layer_mode(deck_html)
    if not parsed:
        return None
    layers = parsed.get("layers")
    canvas = parsed.get("canvas")
    if not isinstance(layers, list) or not layers or not isinstance(canvas, dict):
        return None
    try:
        canvas_model = Canvas(**canvas)
    except ValidationError:
        return None
    return Artifact(
        artifact_id=f"art_{run_id}",
        name=f"Editable Slide deck — {run_id[:8]}",
        artifact_type="deck",
        canvas=canvas_model,
        canvas_plan=canvas_plan,
        deck_plan=deck_plan,
        native_file_url=file_url,
        native_format="html",
        view_file_url=file_url,
        view_format="html",
        download_url=file_url,
        pdf_url=pdf_url,
        downloads=downloads,
        preview_url=preview_url,
        card_preview_url=card_preview_url,
        quality_status=quality_status,
        quality_diagnostics=quality_diagnostics,
        layers=layers,
        parent_artifact_id=parent_id,
    )


def _validated_authored_delivery_metadata(
    run_dir: Path,
    run_id: str,
    target: Path,
    artifact_type: ArtifactType,
) -> dict[str, Any]:
    """Expose delivery metadata only when it is bound to the current artifact."""

    manifest_names = {
        "poster": (
            "authored_poster_edit_manifest.json",
            "code_editor_revision_manifest.json",
            "apply_edits_palette_manifest.json",
            "designer_author_direct_manifest.json",
        ),
        "deck": (
            "authored_html_edit_manifest.json",
            "slides_author_manifest.json",
        ),
        "landing": (
            "authored_html_edit_manifest.json",
            "landing_author_manifest.json",
        ),
    }.get(artifact_type)
    if manifest_names is None:
        return {}
    author_manifest_name = {
        "poster": "designer_author_direct_manifest.json",
        "deck": "slides_author_manifest.json",
        "landing": "landing_author_manifest.json",
    }[artifact_type]
    target_hash = sha256_file(target)
    matching_manifests: list[tuple[str, dict[str, Any]]] = []
    for manifest_name in manifest_names:
        candidate = _read_json_file(run_dir / "final" / manifest_name)
        if not isinstance(candidate, dict):
            continue
        recorded_hash = str(candidate.get("html_sha256") or "").strip().lower()
        if len(recorded_hash) == 64 and recorded_hash == target_hash:
            matching_manifests.append((manifest_name, candidate))
    if not matching_manifests:
        return {}

    result: dict[str, Any] = {}
    quality_manifest = next(
        (
            candidate
            for manifest_name, candidate in matching_manifests
            if manifest_name == author_manifest_name
        ),
        None,
    )
    if quality_manifest is not None:
        quality_status = str(quality_manifest.get("quality_status") or "").strip()
        if quality_status in {"ready", "ready_with_warnings"}:
            result["quality_status"] = quality_status
            diagnostics = quality_manifest.get("quality_diagnostics")
            if isinstance(diagnostics, list):
                result["quality_diagnostics"] = [
                    value.strip()
                    for value in diagnostics
                    if isinstance(value, str) and value.strip()
                ]

    if artifact_type != "landing":
        return result
    for _manifest_name, manifest in matching_manifests:
        relative = str(manifest.get("card_preview_relative_path") or "").strip()
        recorded_preview_hash = str(
            manifest.get("card_preview_sha256") or ""
        ).strip().lower()
        relative_path = Path(relative)
        if (
            not relative
            or relative_path.is_absolute()
            or ".." in relative_path.parts
            or len(recorded_preview_hash) != 64
        ):
            continue
        preview_path = (run_dir / relative_path).resolve()
        if (
            not preview_path.is_relative_to(run_dir.resolve())
            or not preview_path.is_file()
            or sha256_file(preview_path) != recorded_preview_hash
        ):
            continue
        result["card_preview_url"] = _run_file_url(
            run_id,
            relative_path.as_posix(),
        )
        break
    return result


def _fallback_artifact(final_dir: Path, expected_fmt: NativeFormat) -> Path | None:
    if not final_dir.exists():
        return None
    ext = {
        "html": ".html",
        "svg": ".svg",
        "pptx": ".pptx",
        "mp4": ".mp4",
        "png": ".png",
    }[expected_fmt]
    candidates = sorted(final_dir.glob(f"*{ext}"))
    return candidates[0] if candidates else None


def _read_canvas_plan(run_dir: Path) -> dict[str, Any]:
    path = run_dir / "canvas_plan.json"
    if not path.exists():
        return {}
    try:
        data = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError):
        return {}
    return data if isinstance(data, dict) else {}


def _read_deck_plan(run_dir: Path) -> dict[str, Any]:
    path = run_dir / "deck_plan.json"
    if not path.exists():
        return {}
    try:
        data = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError):
        return {}
    return data if isinstance(data, dict) else {}


def _build_video_artifact(
    run_dir: Path, run_id: str, *, baseline_artifact_json: str | None,
) -> Artifact | None:
    """Publish only the MP4 accepted by the current delivery manifest."""
    if not run_dir.exists():
        log(
            "web.artifact.missing",
            run_id=run_id,
            expected=f"{run_dir}/hyperframes-*/delivery_manifest.json",
        )
        return None
    delivery = _validated_video_delivery(run_dir)
    if not delivery.is_passed or not revalidate_current_video_delivery_snapshots(
        run_dir,
        delivery,
    ):
        log("web.video.no_passed_delivery", run_id=run_id)
        return None
    mp4_snapshot = delivery.snapshots["mp4"]
    subtitle_url = _run_file_url(
        run_id,
        delivery.public_paths["vtt"],
    )
    return _build_video_artifact_from_path(
        run_dir,
        run_id,
        run_dir / mp4_snapshot.relative_path,
        baseline_artifact_json=baseline_artifact_json,
        subtitle_url=subtitle_url,
        validated_snapshot=mp4_snapshot,
    )


def _build_video_artifact_from_path(
    run_dir: Path,
    run_id: str,
    mp4_path: Path,
    *,
    baseline_artifact_json: str | None,
    subtitle_url: str | None = None,
    validated_snapshot: VideoDeliverySnapshot | None = None,
) -> Artifact | None:
    """Build a response for one already-validated MP4 inside the run."""
    run_root = run_dir.resolve()
    if validated_snapshot is None:
        mp4_path = mp4_path.resolve()
        if (
            not mp4_path.is_relative_to(run_root)
            or not mp4_path.is_file()
            or mp4_path.stat().st_size <= 0
        ):
            return None
        rel = mp4_path.relative_to(run_root)
        rendered_at_ms = int(mp4_path.stat().st_mtime * 1000)
    else:
        rel = validated_snapshot.relative_path
        if (
            rel.is_absolute()
            or ".." in rel.parts
            or validated_snapshot.size <= 0
            or mp4_path != run_dir / rel
        ):
            return None
        rendered_at_ms = int((validated_snapshot.mtime_ns or 0) / 1_000_000)
    canvas_plan = _read_canvas_plan(run_dir)
    deck_plan = _read_deck_plan(run_dir)
    file_url = _run_file_url(run_id, rel.as_posix())
    downloads = {"vtt": subtitle_url} if subtitle_url else {}

    # Use the landing's preview.png as the chat thumbnail when available
    # (the agent renders this as part of composite); otherwise fall back
    # to the .mp4's own first-frame loaded by the browser.
    preview_path = run_dir / "final" / "preview.png"
    preview_url = (
        _run_file_url(run_id, "final/preview.png")
        if preview_path.exists()
        else None
    )

    parent_id: str | None = None
    if baseline_artifact_json:
        try:
            parent_id = json.loads(baseline_artifact_json).get("artifact_id")
        except (json.JSONDecodeError, AttributeError):
            parent_id = None

    deck_html = run_dir / "final" / "deck.html"
    editable_video = _build_video_layer_artifact(
        deck_html=deck_html,
        run_id=run_id,
        mp4_url=file_url,
        subtitle_url=subtitle_url,
        preview_url=preview_url,
        parent_id=parent_id,
        rendered_at_ms=rendered_at_ms,
        durations=_video_scene_durations(run_dir),
        canvas_plan=canvas_plan,
        deck_plan=deck_plan,
    )
    if editable_video is not None:
        return editable_video

    return Artifact(
        artifact_id=f"art_{run_id}",
        name=f"Video — {run_id[:8]}",
        artifact_type="video",
        canvas=Canvas(w=1920, h=1080),
        canvas_plan=canvas_plan,
        deck_plan=deck_plan,
        native_file_url=file_url,
        native_format="mp4",
        view_file_url=file_url,
        view_format="mp4",
        download_url=file_url,
        downloads=downloads,
        preview_url=preview_url,
        layers=[],
        parent_artifact_id=parent_id,
    )


def _validated_video_delivery(
    run_dir: Path,
) -> CurrentVideoDeliveryValidation:
    """Return the shared secure current-delivery validation contract."""

    return _validate_current_video_delivery(run_dir)


def _build_video_layer_artifact(
    *,
    deck_html: Path,
    run_id: str,
    mp4_url: str,
    subtitle_url: str | None,
    preview_url: str | None,
    parent_id: str | None,
    rendered_at_ms: int,
    durations: list[float],
    canvas_plan: dict[str, Any],
    deck_plan: dict[str, Any],
) -> Artifact | None:
    parsed = parse_deck_html_as_layer_mode(deck_html)
    if not parsed:
        return None
    layers = parsed.get("layers")
    canvas = parsed.get("canvas")
    frames = parsed.get("frames")
    if (
        not isinstance(layers, list)
        or not layers
        or not isinstance(canvas, dict)
        or not isinstance(frames, list)
        or not frames
    ):
        return None
    try:
        canvas_model = Canvas(**canvas)
    except ValidationError:
        return None

    scenes: list[dict[str, Any]] = []
    total = 0.0
    for idx, frame in enumerate(frames):
        if not isinstance(frame, dict):
            continue
        frame_id = str(frame.get("layer_id") or "")
        if not frame_id:
            continue
        duration = durations[idx] if idx < len(durations) else 6.0
        duration = max(0.5, min(120.0, float(duration or 6.0)))
        total += duration
        scenes.append({
            "scene_id": f"scene_{idx + 1}",
            "name": f"Scene {idx + 1}",
            "frame_layer_id": frame_id,
            "duration_s": round(duration, 3),
            "transition": "cut",
        })
    if not scenes:
        return None

    return Artifact(
        artifact_id=f"art_{run_id}",
        name=f"Editable Video — {run_id[:8]}",
        artifact_type="video",
        canvas=canvas_model,
        canvas_plan=canvas_plan,
        deck_plan=deck_plan,
        native_file_url=None,
        native_format=None,
        view_file_url=None,
        view_format=None,
        download_url=mp4_url,
        downloads={"vtt": subtitle_url} if subtitle_url else {},
        preview_url=preview_url,
        layers=layers,
        video_project={
            "duration_s": round(total, 3),
            "fps": 30,
            "scenes": scenes,
            "latest_render": {
                "run_id": run_id,
                "mp4_url": mp4_url,
                "subtitle_url": subtitle_url,
                "rendered_at": rendered_at_ms,
            },
        },
        parent_artifact_id=parent_id,
    )


def _video_scene_durations(run_dir: Path) -> list[float]:
    """Read authored HyperFrames scene durations when a real video run has them."""
    indexes = sorted(run_dir.glob("hyperframes-*/index.html"), key=lambda p: p.stat().st_mtime, reverse=True)
    for index in indexes:
        try:
            text = index.read_text(encoding="utf-8")
        except OSError:
            continue
        durations: list[float] = []
        for match in re.finditer(r"<section\b[^>]*\bclass=\"[^\"]*\bstage\b[^\"]*\"[^>]*>", text):
            tag = match.group(0)
            if "data-layout-ignore" in tag:
                continue
            duration_match = re.search(r"\bdata-duration=\"([0-9.]+)\"", tag)
            if not duration_match:
                continue
            try:
                durations.append(float(duration_match.group(1)))
            except ValueError:
                continue
        if durations:
            return durations
    return []


_EDITABLE_VIDEO_COMPOSITION_ID = "editable-video-demo"
_EDITABLE_VIDEO_W = 1920
_EDITABLE_VIDEO_H = 1080
_EDITABLE_VIDEO_HYPERFRAMES_JSON = {
    "$schema": "https://hyperframes.heygen.com/schema/hyperframes.json",
    "registry": "https://raw.githubusercontent.com/heygen-com/hyperframes/main/registry",
    "paths": {
        "blocks": "compositions",
        "components": "compositions/components",
        "assets": "assets",
    },
}


async def _render_editable_video_in_background(
    *,
    run_id: str,
    artifact: dict[str, Any],
    state: _RunState,
    conversation_id: str,
    baseline_artifact_json: str | None,
) -> None:
    """No-agent editable video render task used by /api/video/render."""
    started_at = time.time()
    run_dir = RUNS_DIR / run_id
    proj_dir = run_dir / "hyperframes-editable-demo"
    settings_path = _settings_or_boot()
    try:
        _persisted_run_log(
            "run.start",
            run_id,
            pid=os.getpid(),
            artifact_type="video",
            source="editable_video_demo",
        )
        _append_event(
            settings_path,
            conversation_id,
            "artifact.render_started",
            run_id=run_id,
            artifact_id=str(artifact.get("artifact_id") or ""),
            data={"artifact_type": "video", "source": "editable_video_demo"},
        )
        await asyncio.to_thread(
            _write_editable_video_project,
            artifact,
            run_id,
            run_dir,
            proj_dir,
        )
        project = artifact.get("video_project") if isinstance(artifact.get("video_project"), dict) else {}
        fps = int(_positive_float(project.get("fps"), 30, min_value=1, max_value=120))
        log("export_video.render.start", run_id=run_id, fps=fps, project=str(proj_dir))
        render_output, render_ok, mp4_path = await asyncio.to_thread(
            _run_editable_video_render,
            proj_dir,
            fps,
        )
        if not render_ok or mp4_path is None or not mp4_path.exists():
            message = (render_output or "HyperFrames render did not produce an MP4.").strip()
            state.error = message[:1000]
            log("export_video.render.error", run_id=run_id, error=message[:500])
            state.result_message = Message(
                id=f"msg_{run_id}",
                role="assistant",
                text=f"Video render failed: {state.error}",
                ts=int(time.time() * 1000),
                run_id=run_id,
                status="error",
                failure=Failure(
                    status="error",
                    phase="rendering",
                    agent_last_note=state.error,
                    produced_files=_list_produced_artifacts(run_id),
                    elapsed_ms=int((time.time() - started_at) * 1000),
                ),
            )
            _append_event(
                settings_path,
                conversation_id,
                "artifact.generation_failed",
                run_id=run_id,
                data={
                    "status": "error",
                    "artifact_type": "video",
                    "source": "editable_video_demo",
                    "error": state.error,
                    "canvas_plan": _read_canvas_plan(run_dir),
                    "deck_plan": _read_deck_plan(run_dir),
                },
            )
            _persisted_run_log("run.error", run_id, msg=state.error[:200])
            return

        log("export_video.render.done", run_id=run_id, mp4=str(mp4_path))
        artifact_resp = _build_video_artifact_from_path(
            run_dir,
            run_id,
            mp4_path,
            baseline_artifact_json=baseline_artifact_json,
        )
        if artifact_resp is None:
            raise RuntimeError("MP4 exists but artifact response could not be built")
        state.result_artifact = artifact_resp
        state.result_message = Message(
            id=f"msg_{run_id}",
            role="assistant",
            text="Rendered editable video MP4.",
            ts=int(time.time() * 1000),
            run_id=run_id,
            artifact_id=artifact_resp.artifact_id,
            status="done",
        )
        _append_event(
            settings_path,
            conversation_id,
            "artifact.generated",
            run_id=run_id,
            artifact_id=artifact_resp.artifact_id,
            data={
                "artifact_type": "video",
                "name": artifact_resp.name,
                "native_format": "mp4",
                "source": "editable_video_demo",
                "canvas_plan": artifact_resp.canvas_plan,
                "deck_plan": artifact_resp.deck_plan,
            },
        )
        _persisted_run_log("run.done", run_id)
    except Exception as e:  # noqa: BLE001
        state.error = f"{type(e).__name__}: {e}"
        log("web.video.render.error", run_id=run_id, error=type(e).__name__, msg=str(e)[:200])
        state.result_message = Message(
            id=f"msg_{run_id}",
            role="assistant",
            text=f"Video render failed: {state.error}",
            ts=int(time.time() * 1000),
            run_id=run_id,
            status="error",
            failure=Failure(
                status="error",
                phase="rendering",
                agent_last_note=state.error,
                produced_files=_list_produced_artifacts(run_id),
                elapsed_ms=int((time.time() - started_at) * 1000),
            ),
        )
        _append_event(
            settings_path,
            conversation_id,
            "artifact.generation_failed",
            run_id=run_id,
            data={
                "status": "error",
                "artifact_type": "video",
                "source": "editable_video_demo",
                "error": state.error[:500],
                "canvas_plan": _read_canvas_plan(run_dir),
                "deck_plan": _read_deck_plan(run_dir),
            },
        )
        _persisted_run_log("run.error", run_id, msg=state.error[:200])


def _write_editable_video_project(
    artifact: dict[str, Any],
    run_id: str,
    run_dir: Path,
    proj_dir: Path,
) -> None:
    if proj_dir.exists():
        shutil.rmtree(proj_dir)
    (proj_dir / "assets").mkdir(parents=True, exist_ok=True)
    (proj_dir / "renders").mkdir(parents=True, exist_ok=True)
    (run_dir / "final").mkdir(parents=True, exist_ok=True)

    manifest = _editable_video_manifest(artifact, proj_dir)
    (proj_dir / "meta.json").write_text(
        json.dumps(
            {
                "id": _EDITABLE_VIDEO_COMPOSITION_ID,
                "name": "Editable Video Demo",
                "run_id": run_id,
            },
            indent=2,
            ensure_ascii=False,
        ),
        encoding="utf-8",
    )
    (proj_dir / "hyperframes.json").write_text(
        json.dumps(_EDITABLE_VIDEO_HYPERFRAMES_JSON, indent=2),
        encoding="utf-8",
    )
    (proj_dir / "scene_manifest.json").write_text(
        json.dumps(manifest, indent=2, ensure_ascii=False),
        encoding="utf-8",
    )
    (proj_dir / "index.html").write_text(
        _editable_video_index_html(manifest),
        encoding="utf-8",
    )


def _editable_video_manifest(artifact: dict[str, Any], proj_dir: Path) -> dict[str, Any]:
    project = artifact.get("video_project") if isinstance(artifact.get("video_project"), dict) else {}
    raw_scenes = project.get("scenes") if isinstance(project.get("scenes"), list) else []
    raw_layers = artifact.get("layers") if isinstance(artifact.get("layers"), list) else []
    layers = [l for l in raw_layers if isinstance(l, dict)]
    layer_by_id = {
        str(l.get("layer_id")): l
        for l in layers
        if l.get("layer_id") is not None
    }

    scenes: list[dict[str, Any]] = []
    start_s = 0.0
    for idx, raw_scene in enumerate(raw_scenes):
        if not isinstance(raw_scene, dict):
            continue
        frame_layer_id = str(raw_scene.get("frame_layer_id") or "")
        frame_layer = layer_by_id.get(frame_layer_id)
        frame_bbox = _editable_layer_bbox(frame_layer)
        if frame_layer is None or frame_bbox is None:
            continue
        duration_s = _positive_float(
            raw_scene.get("duration_s"),
            4.0,
            min_value=0.5,
            max_value=120.0,
        )
        scene_layers: list[dict[str, Any]] = []
        for layer in layers:
            if layer.get("visible") is False:
                continue
            bbox = _editable_layer_bbox(layer)
            if bbox is None:
                continue
            if str(layer.get("layer_id")) != frame_layer_id:
                if not _bbox_intersects(bbox, frame_bbox):
                    continue
                if layer.get("kind") == "background" and (
                    bbox["w"] > frame_bbox["w"] * 1.05
                    or bbox["h"] > frame_bbox["h"] * 1.05
                ):
                    continue
            scene_layers.append(
                _export_editable_video_layer(layer, bbox, frame_bbox, proj_dir)
            )
        scenes.append({
            "scene_id": str(raw_scene.get("scene_id") or f"scene-{idx + 1}"),
            "name": str(raw_scene.get("name") or f"Scene {idx + 1}"),
            "frame_layer_id": frame_layer_id,
            "transition": str(raw_scene.get("transition") or "cut"),
            "track_index": idx + 1,
            "start_s": round(start_s, 3),
            "duration_s": round(duration_s, 3),
            "layers": sorted(scene_layers, key=lambda l: float(l.get("z_index") or 0)),
        })
        start_s += duration_s

    if not scenes:
        raise ValueError("editable video has no valid scenes to render")

    fps = int(_positive_float(project.get("fps"), 30, min_value=1, max_value=120))
    return {
        "composition_id": _EDITABLE_VIDEO_COMPOSITION_ID,
        "width": _EDITABLE_VIDEO_W,
        "height": _EDITABLE_VIDEO_H,
        "fps": fps,
        "duration_s": round(start_s, 3),
        "scenes": scenes,
    }


def _export_editable_video_layer(
    layer: dict[str, Any],
    bbox: dict[str, float],
    frame_bbox: dict[str, float],
    proj_dir: Path,
) -> dict[str, Any]:
    sx = _EDITABLE_VIDEO_W / frame_bbox["w"]
    sy = _EDITABLE_VIDEO_H / frame_bbox["h"]
    scale = min(sx, sy)
    out = dict(layer)
    out["bbox"] = {
        "x": round((bbox["x"] - frame_bbox["x"]) * sx, 3),
        "y": round((bbox["y"] - frame_bbox["y"]) * sy, 3),
        "w": round(bbox["w"] * sx, 3),
        "h": round(bbox["h"] * sy, 3),
    }
    out["z_index"] = float(layer.get("z_index") or 0)
    if isinstance(layer.get("font_size_px"), (int, float)):
        out["font_size_px"] = round(float(layer["font_size_px"]) * scale, 3)
    if isinstance(layer.get("letter_spacing"), (int, float)):
        out["letter_spacing"] = round(float(layer["letter_spacing"]) * scale, 3)
    if isinstance(layer.get("corner_radius"), (int, float)):
        out["corner_radius"] = round(float(layer["corner_radius"]) * scale, 3)
    if isinstance(layer.get("stroke_width"), (int, float)):
        out["stroke_width"] = round(float(layer["stroke_width"]) * scale, 3)
    shadow = layer.get("shadow")
    if isinstance(shadow, dict):
        out["shadow"] = {
            **shadow,
            "dx": round(float(shadow.get("dx") or 0) * scale, 3),
            "dy": round(float(shadow.get("dy") or 0) * scale, 3),
            "blur": round(float(shadow.get("blur") or 0) * scale, 3),
        }
    if isinstance(layer.get("src"), str):
        out["src"] = _resolve_editable_video_src(layer["src"], proj_dir)
    return out


def _editable_video_index_html(manifest: dict[str, Any]) -> str:
    comp_id = _attr(manifest["composition_id"])
    total = _css_num(float(manifest["duration_s"]))
    scenes_html = "\n".join(
        _render_editable_video_scene(scene)
        for scene in manifest["scenes"]
    )
    timeline_script = _editable_video_timeline_script(manifest)
    return f"""<!doctype html>
<html lang="en">
<head>
  <meta charset="utf-8">
  <meta name="viewport" content="width=device-width, initial-scale=1">
  <title>Editable Video Demo</title>
  <style>
    html, body {{
      margin: 0;
      width: 100%;
      height: 100%;
      overflow: hidden;
      background: #17130f;
      font-family: Inter, ui-sans-serif, system-ui, -apple-system, BlinkMacSystemFont, "Segoe UI", sans-serif;
    }}
    #root {{
      position: relative;
      width: {_EDITABLE_VIDEO_W}px;
      height: {_EDITABLE_VIDEO_H}px;
      overflow: hidden;
      background: #17130f;
    }}
    .scene {{
      position: absolute;
      inset: 0;
      width: {_EDITABLE_VIDEO_W}px;
      height: {_EDITABLE_VIDEO_H}px;
      overflow: hidden;
      opacity: 1;
    }}
    .layer {{
      box-sizing: border-box;
      position: absolute;
    }}
    .text-layer {{
      overflow: hidden;
      white-space: pre-wrap;
    }}
    img.layer {{
      display: block;
    }}
  </style>
</head>
<body>
  <div id="root" data-composition-id="{comp_id}" data-start="0" data-width="{_EDITABLE_VIDEO_W}" data-height="{_EDITABLE_VIDEO_H}" data-duration="{total}">
{scenes_html}
  </div>
  <script src="https://cdn.jsdelivr.net/npm/gsap@3/dist/gsap.min.js"></script>
  <script>
{timeline_script}
  </script>
</body>
</html>
"""


def _render_editable_video_scene(scene: dict[str, Any]) -> str:
    scene_id = _attr(scene["scene_id"])
    start = _css_num(float(scene["start_s"]))
    duration = _css_num(float(scene["duration_s"]))
    track = int(_positive_float(scene.get("track_index"), 1, min_value=1, max_value=999))
    layers_html = "\n".join(
        _render_editable_video_layer(layer)
        for layer in scene.get("layers", [])
        if isinstance(layer, dict)
    )
    return f"""    <section id="{scene_id}" class="clip scene" data-start="{start}" data-duration="{duration}" data-track-index="{track}">
{layers_html}
    </section>"""


def _render_editable_video_layer(layer: dict[str, Any]) -> str:
    kind = str(layer.get("kind") or "shape")
    layer_id = _attr(layer.get("layer_id") or "")
    name = _attr(layer.get("name") or layer.get("layer_id") or "")
    style = _base_layer_style(layer)
    if kind == "text":
        effects = layer.get("effects") if isinstance(layer.get("effects"), dict) else {}
        text_style = [
            style,
            f"font-family:{_safe_css(layer.get('font_family'), 'Inter, Arial, sans-serif')}",
            f"font-size:{_css_num(layer.get('font_size_px') or 36)}px",
            f"font-weight:{int(_positive_float(layer.get('font_weight'), 400, min_value=100, max_value=1000))}",
            f"font-style:{'italic' if layer.get('font_style') == 'italic' else 'normal'}",
            f"line-height:{_css_num(layer.get('line_height') or 1.12)}",
            f"letter-spacing:{_css_num(layer.get('letter_spacing') or 0)}px",
            f"text-align:{_safe_choice(layer.get('align'), {'left', 'center', 'right'}, 'left')}",
            f"text-transform:{'uppercase' if layer.get('text_transform') == 'uppercase' else 'none'}",
            f"color:{_safe_css(effects.get('fill'), '#17130f')}",
        ]
        return (
            f'      <div class="layer text-layer" data-layer-id="{layer_id}" '
            f'data-name="{name}" style="{_attr(";".join(text_style))}">'
            f"{html.escape(str(layer.get('text') or ''))}</div>"
        )
    if kind == "image":
        src = _attr(layer.get("src") or "")
        fit = _safe_choice(layer.get("fit"), {"cover", "contain", "fill"}, "cover")
        pos = layer.get("object_position") if isinstance(layer.get("object_position"), dict) else {}
        ox = _positive_float(pos.get("x"), 0.5, min_value=0, max_value=1) * 100
        oy = _positive_float(pos.get("y"), 0.5, min_value=0, max_value=1) * 100
        img_style = [
            style,
            f"object-fit:{fit}",
            f"object-position:{_css_num(ox)}% {_css_num(oy)}%",
            f"border-radius:{_css_num(layer.get('corner_radius') or 0)}px",
            f"opacity:{_css_num(_positive_float(layer.get('opacity'), 1, min_value=0, max_value=1))}",
            _shadow_style(layer),
        ]
        return (
            f'      <img class="layer image-layer" data-layer-id="{layer_id}" '
            f'data-name="{name}" src="{src}" alt="{name}" '
            f'style="{_attr(";".join(p for p in img_style if p))}">'
        )
    shape_style = [
        style,
        f"background:{_safe_css(layer.get('fill_color'), 'transparent')}",
        _shape_radius_style(layer),
        f"opacity:{_css_num(_positive_float(layer.get('opacity'), 1, min_value=0, max_value=1))}",
        _border_style(layer),
        _shadow_style(layer),
    ]
    return (
        f'      <div class="layer shape-layer" data-layer-id="{layer_id}" '
        f'data-name="{name}" style="{_attr(";".join(p for p in shape_style if p))}"></div>'
    )


def _editable_video_timeline_script(manifest: dict[str, Any]) -> str:
    lines = [
        '    window.__timelines = window.__timelines || {};',
        '    if (window.gsap) {',
        '      const tl = window.gsap.timeline({ paused: true });',
        '      tl.set(".scene", { autoAlpha: 0 }, 0);',
    ]
    for scene in manifest["scenes"]:
        scene_id = json.dumps(str(scene["scene_id"]))
        start = _css_num(float(scene["start_s"]))
        end = _css_num(float(scene["start_s"]) + float(scene["duration_s"]))
        lines.append(f"      tl.set(document.getElementById({scene_id}), {{ autoAlpha: 1 }}, {start});")
        lines.append(f"      tl.set(document.getElementById({scene_id}), {{ autoAlpha: 0 }}, {end});")
    lines.extend([
        f'      window.__timelines[{json.dumps(manifest["composition_id"])}] = tl;',
        "    }",
    ])
    return "\n".join(lines)


def _run_editable_video_render(
    proj_dir: Path,
    fps: int,
) -> tuple[str, bool, Path | None]:
    from autodesign.tools.export_video import resolve_hyperframes_binary

    started_ns = time.time_ns()
    expected_mp4 = proj_dir / "renders" / f"editable-{started_ns}.mp4"
    expected_mp4.parent.mkdir(parents=True, exist_ok=True)
    try:
        proc = subprocess.run(
            [
                str(resolve_hyperframes_binary()),
                "render",
                "--fps", str(fps),
                "--resolution", "landscape",
                "--strict",
                "--no-best-effort",
                "--output", str(expected_mp4.relative_to(proj_dir)),
                ".",
            ],
            cwd=proj_dir,
            env={**os.environ, "HYPERFRAMES_PYTHON": sys.executable},
            capture_output=True,
            text=True,
            timeout=300,
        )
        combined = (proc.stdout or "") + (proc.stderr or "")
        ok = proc.returncode == 0
        log(
            "web.video.hyperframes.render",
            returncode=proc.returncode,
            ok=ok,
            fps=fps,
            output_chars=len(combined),
        )
        mp4_path: Path | None = None
        if ok and expected_mp4.is_file():
            stat = expected_mp4.stat()
            if stat.st_size > 0 and stat.st_mtime_ns + 2_000_000_000 >= started_ns:
                mp4_path = expected_mp4
        ok = ok and mp4_path is not None
        return combined.strip(), ok, mp4_path
    except FileNotFoundError:
        return "pinned HyperFrames CLI not found; run npm install in web/", False, None
    except subprocess.TimeoutExpired:
        return "HyperFrames render timed out after 300 s", False, None
    except Exception as e:  # noqa: BLE001
        return f"{type(e).__name__}: {e}", False, None


def _editable_layer_bbox(layer: dict[str, Any] | None) -> dict[str, float] | None:
    if not isinstance(layer, dict):
        return None
    raw = layer.get("bbox")
    if not isinstance(raw, dict):
        return None
    try:
        x = float(raw.get("x"))
        y = float(raw.get("y"))
        w = float(raw.get("w"))
        h = float(raw.get("h"))
    except (TypeError, ValueError):
        return None
    if w <= 0 or h <= 0:
        return None
    return {"x": x, "y": y, "w": w, "h": h}


def _bbox_intersects(a: dict[str, float], b: dict[str, float]) -> bool:
    return (
        a["x"] < b["x"] + b["w"]
        and a["x"] + a["w"] > b["x"]
        and a["y"] < b["y"] + b["h"]
        and a["y"] + a["h"] > b["y"]
    )


def _base_layer_style(layer: dict[str, Any]) -> str:
    bbox = layer["bbox"]
    return ";".join([
        f"left:{_css_num(bbox['x'])}px",
        f"top:{_css_num(bbox['y'])}px",
        f"width:{_css_num(bbox['w'])}px",
        f"height:{_css_num(bbox['h'])}px",
        f"z-index:{int(float(layer.get('z_index') or 0))}",
    ])


def _shape_radius_style(layer: dict[str, Any]) -> str:
    if layer.get("shape_kind") == "ellipse":
        return "border-radius:9999px"
    return f"border-radius:{_css_num(layer.get('corner_radius') or 0)}px"


def _border_style(layer: dict[str, Any]) -> str:
    width = _positive_float(layer.get("stroke_width"), 0, min_value=0, max_value=1000)
    color = layer.get("stroke_color")
    if width <= 0 or not color:
        return ""
    dash = _safe_choice(layer.get("stroke_dash"), {"solid", "dashed", "dotted"}, "solid")
    return f"border:{_css_num(width)}px {dash} {_safe_css(color, '#17130f')}"


def _shadow_style(layer: dict[str, Any]) -> str:
    shadow = layer.get("shadow")
    if not isinstance(shadow, dict):
        return ""
    dx = _css_num(shadow.get("dx") or 0)
    dy = _css_num(shadow.get("dy") or 0)
    blur = _css_num(shadow.get("blur") or 0)
    opacity = _positive_float(shadow.get("opacity"), 0.18, min_value=0, max_value=1)
    color = _rgba_from_hex(str(shadow.get("color") or "#17130f"), opacity)
    return f"box-shadow:{dx}px {dy}px {blur}px {color}"


def _rgba_from_hex(hex_color: str, opacity: float) -> str:
    clean = hex_color.strip().lstrip("#")
    if len(clean) != 6:
        return f"rgba(23, 19, 15, {_css_num(opacity)})"
    try:
        n = int(clean, 16)
    except ValueError:
        return f"rgba(23, 19, 15, {_css_num(opacity)})"
    return f"rgba({(n >> 16) & 255}, {(n >> 8) & 255}, {n & 255}, {_css_num(opacity)})"


def _resolve_editable_video_src(src: str, proj_dir: Path) -> str:
    clean = src.split("?", 1)[0]
    if clean.startswith("/api/files/runs/"):
        raise ValueError(
            "legacy editable-video assembly cannot authorize run-scoped assets"
        )
    mappings = [
        ("/api/files/editor-assets/", EDITOR_ASSETS_DIR, "editor-assets"),
    ]
    for prefix, root, folder in mappings:
        if not clean.startswith(prefix):
            continue
        rel = clean[len(prefix):].lstrip("/")
        source = (root / rel).resolve()
        root_resolved = root.resolve()
        if not _path_inside(source, root_resolved) or not source.is_file():
            raise ValueError("editable video image source is unavailable")
        dest_name = "-".join(Path(rel).parts)[-96:] or source.name
        dest_dir = proj_dir / "assets" / folder
        dest_dir.mkdir(parents=True, exist_ok=True)
        dest = dest_dir / dest_name
        shutil.copy2(source, dest)
        return f"assets/{folder}/{dest.name}"
    raise ValueError("editable video image source must use an approved local asset URL")


def _path_inside(path: Path, root: Path) -> bool:
    try:
        path.relative_to(root)
        return True
    except ValueError:
        return False


def _positive_float(
    value: Any,
    default: float,
    *,
    min_value: float | None = None,
    max_value: float | None = None,
) -> float:
    try:
        out = float(value)
    except (TypeError, ValueError):
        out = float(default)
    if min_value is not None:
        out = max(min_value, out)
    if max_value is not None:
        out = min(max_value, out)
    return out


def _safe_choice(value: Any, allowed: set[str], default: str) -> str:
    s = str(value or "")
    return s if s in allowed else default


def _safe_css(value: Any, default: str) -> str:
    s = str(value or "").strip()
    if not s or any(c in s for c in "<>{};"):
        return default
    return s


def _css_num(value: Any) -> str:
    try:
        n = float(value)
    except (TypeError, ValueError):
        n = 0.0
    s = f"{n:.3f}".rstrip("0").rstrip(".")
    return s or "0"


def _attr(value: Any) -> str:
    return html.escape(str(value), quote=True)


def _build_assistant_text(result: RunResult, artifact: Artifact | None, attach_paths: list[Path]) -> str:
    """Compose the assistant message that closes a generate run.

    Surfaces what the run produced, terminal status, critic signal, and any
    final designer note without exposing model-decision traces.
    """
    designer_summary = (result.finalize_notes or "").strip()

    if artifact is None:
        # Failed-to-produce-the-asked-for-artifact path. List what *did*
        # land on disk so the user can decide whether to keep, retry, or
        # cancel — instead of a one-line "fail".
        produced = _list_produced_artifacts(result.run_id)
        bits: list[str] = []
        bits.append(
            f"Run finished without producing the requested artifact. "
            f"Status: {result.terminal_status or 'unknown'}."
        )
        if produced:
            bits.append(f"On disk: {', '.join(produced)}.")
        if designer_summary:
            bits.append(f"\n\nAgent's note:\n{designer_summary}")
        bits.append(f"\n\nDiagnostics: out/runs/{result.run_id}/")
        return "".join(bits) if any("\n" in b for b in bits) else " ".join(bits)

    bits = []
    if attach_paths:
        bits.append(f"Read {len(attach_paths)} attached file(s).")
    # Drop the "(video render didn't complete)" tag — already in the
    # artifact name visible in the chat card. Avoid double-stating.
    if "(" in artifact.name and "didn't complete" in artifact.name:
        bits.append(
            f"Generated a {artifact.artifact_type} draft, "
            f"but the final render step didn't complete."
        )
    else:
        bits.append(f"Generated a {artifact.artifact_type} draft.")
    if result.terminal_status and result.terminal_status != "pass":
        bits.append(f"(Terminal status: {result.terminal_status}.)")
    if result.critic_score is not None:
        bits.append(f"Critic score: {result.critic_score:.2f}.")
    if designer_summary:
        # Two-paragraph layout: first the engineering signal, then the
        # designer's natural-language self-report. Newlines preserved.
        head = " ".join(bits)
        return f"{head}\n\n{designer_summary}"
    bits.append("Open in Canvas, refine in chat, or download the source.")
    return " ".join(bits)


def _list_produced_artifacts(run_id: str) -> list[str]:
    """Best-effort enumeration of files in the run dir, so the user can
    see what *was* produced when the run technically failed. Names are
    relative to run_dir for readability."""
    out: list[str] = []
    run_dir = _settings_or_boot() / "runs" / run_id
    if not run_dir.exists():
        return out
    final_dir = run_dir / "final"
    if final_dir.exists():
        for p in sorted(final_dir.iterdir()):
            if p.is_file() and p.suffix in (".html", ".pptx", ".pdf", ".svg", ".mp4"):
                out.append(f"final/{p.name}")
    for hf in sorted(run_dir.glob("hyperframes-*")):
        renders = hf / "renders"
        mp4s = sorted(renders.glob("*.mp4")) if renders.exists() else []
        if mp4s:
            out.append(f"{hf.name}/renders/{mp4s[-1].name}")
        elif (hf / "index.html").exists():
            out.append(f"{hf.name}/index.html (video scaffolded; MP4 not produced)")
    return out


# ---------- Failure metadata (Pass 1 — Failure Card UX) ----------

# Match Kimi K2.6 family on the designer role; paper-poster runs can stall
# on bbox geometry for long PDF inputs.
_KIMI_PATTERN = "kimi"
_OPUS_FALLBACK = "anthropic/claude-opus-4-7"


def _resume_checkpoint_from_disk(run_dir: Path) -> dict[str, int] | None:
    """Return the validated authoring checkpoint exposed by runner resume."""
    try:
        resume_state = _load_resume_state(run_dir)
    except Exception:  # noqa: BLE001
        return None
    if not isinstance(resume_state, dict):
        return None
    previous_attempt_dir = resume_state.get("previous_attempt_dir")
    if not previous_attempt_dir:
        return None
    match = re.fullmatch(r"attempt_(\d+)", Path(previous_attempt_dir).name)
    if match is None:
        return None
    try:
        prior_attempts = int(resume_state.get("prior_attempts") or 0)
    except (TypeError, ValueError):
        return None
    if prior_attempts < 1:
        return None
    return {
        "resume_from_attempt": int(match.group(1)),
        "next_attempt": prior_attempts + 1,
    }


def _event_phase(event_name: str) -> str | None:
    if (
        "_author." in event_name
        or event_name.startswith("designer_author.")
        or event_name.startswith("external_author.")
    ):
        return "authoring"
    if event_name.startswith(
        (
            "render",
            "composite",
            "finalize",
            "export_video",
            "video_render",
            "artifact_export",
        )
    ):
        return "rendering"
    if event_name.startswith(("ingest", "paper_memory", "claim_graph")):
        return "ingest"
    if event_name.startswith(("designer.", "planner.")):
        return "planning"
    return None


def _redacted_error_detail(
    value: str,
    *,
    max_chars: int | None = 1200,
) -> str | None:
    detail = value.strip()
    if not detail:
        return None
    detail = re.sub(
        r"(?i)(authorization:\s*bearer\s+)\S+",
        r"\1[redacted]",
        detail,
    )
    detail = re.sub(
        r"(?i)(api[_-]?key\s*[=:]\s*)\S+",
        r"\1[redacted]",
        detail,
    )
    detail = re.sub(
        r"\bsk-[A-Za-z0-9_-]{12,}\b",
        "[redacted]",
        detail,
    )
    internal_paths = {str(Path.home()), str(RUNS_DIR)}
    if SETTINGS is not None:
        internal_paths.update({
            str(SETTINGS.out_dir),
            str(SETTINGS.repo_root),
        })
    for path in sorted(
        (item for item in internal_paths if item),
        key=len,
        reverse=True,
    ):
        detail = detail.replace(path, "[internal-path]")
    return detail if max_chars is None else detail[-max_chars:]


def _validated_worker_exit_diagnostic(
    event: dict[str, Any],
) -> dict[str, str] | None:
    if event.get("event") != "worker.exit" or event.get("version") != 1:
        return None
    returncode = event.get("returncode")
    error_code = event.get("error_code")
    error_message = event.get("error_message")
    error_detail = event.get("error_detail")
    if (
        type(returncode) is not int
        or error_code not in {
            "worker_result_missing",
            "worker_result_invalid",
            "worker_exit_contradiction",
        }
        or not isinstance(error_message, str)
        or not isinstance(error_detail, str)
        or len(error_code) > 64
        or len(error_message) > 500
        or len(error_detail) > 1_200
    ):
        return None
    phase = event.get("last_phase")
    if phase is not None and (
        not isinstance(phase, str) or not phase or len(phase) > 240
    ):
        phase = None
    safe_message = _redacted_error_detail(error_message, max_chars=500)
    safe_detail = _redacted_error_detail(error_detail, max_chars=1_200)
    if safe_message is None or safe_detail is None:
        return None
    return {
        "phase": phase or "",
        "error_code": error_code,
        "error_message": safe_message,
        "error_detail": safe_detail,
    }


def _worker_exit_failure_fields(outcome: WorkerOutcome) -> dict[str, str]:
    diagnostic = outcome.exit_diagnostic
    if not isinstance(diagnostic, WorkerExitDiagnostic):
        return {}
    validated = _validated_worker_exit_diagnostic(
        diagnostic.event_payload(
            run_id=outcome.run_id,
            job_kind=outcome.job_kind,
        )
    )
    return validated or {}


def _failure_diagnostics_from_disk(run_dir: Path) -> dict[str, Any]:
    """Extract a safe root cause and checkpoint summary from persisted events."""
    _offset, events = _read_disk_run_events(run_dir / "run_events.jsonl", 0)
    phase: str | None = None
    last_author_output: dict[str, Any] | None = None
    last_structured_failure: dict[str, Any] | None = None
    last_run_error: dict[str, Any] | None = None
    last_worker_exit: dict[str, str] | None = None
    for event in events:
        event_name = str(event.get("event") or "")
        event_phase = _event_phase(event_name)
        if event_phase is not None:
            phase = event_phase
        if event_name.endswith(".agent_output"):
            last_author_output = event
        if event_name.endswith(".fail"):
            last_structured_failure = event
        if event_name == "run.error":
            last_run_error = event
        worker_exit = _validated_worker_exit_diagnostic(event)
        if worker_exit is not None:
            last_worker_exit = worker_exit

    error_code: str | None = None
    error_message: str | None = None
    error_detail: str | None = None
    if (
        last_author_output is not None
        and str(last_author_output.get("status") or "").lower() == "error"
    ):
        reason = str(last_author_output.get("reason") or "")
        output = " ".join(
            str(last_author_output.get(key) or "")
            for key in ("stdout_excerpt", "stderr_excerpt")
        )
        error_detail = _redacted_error_detail(output)
        normalized = f"{reason} {output}".lower()
        if (
            "429" in normalized
            or "rate limit" in normalized
            or "rate_limit" in normalized
            or "每分钟请求次数超过限制" in normalized
        ):
            error_code = "provider_rate_limit"
            error_message = (
                "The coding-harness provider rejected the request because its "
                "per-minute rate limit was exceeded. Wait briefly, then resume "
                "from the saved checkpoint."
            )
        elif "timeout" in normalized:
            error_code = "coding_harness_timeout"
            error_message = (
                "The coding harness timed out before completing this attempt. "
                "Resume from the saved checkpoint to continue."
            )
        elif "process_exit" in normalized:
            error_code = "coding_harness_process_exit"
            error_message = (
                "The coding harness exited before producing a usable artifact. "
                "Review the author output in diagnostics, then resume."
            )
        elif "no_output" in normalized:
            error_code = "coding_harness_no_output"
            error_message = (
                "The coding harness finished without producing a usable artifact. "
                "Review the author output in diagnostics, then resume."
            )
    if error_code is None and last_structured_failure is not None:
        error_code = str(
            last_structured_failure.get("reason") or "authoring_failed"
        )
        error_message = str(
            last_structured_failure.get("message") or "Artifact authoring failed."
        )
        diagnostic_payload: dict[str, Any] = {
            "reason": error_code,
            "message": error_message,
        }
        attempt_dir_value = str(
            last_structured_failure.get("attempt_dir") or ""
        ).strip()
        if (
            str(last_structured_failure.get("event") or "")
            == "video_author.fail"
            and attempt_dir_value
        ):
            delivery_errors_path = (
                Path(attempt_dir_value) / "video_author_delivery_errors.json"
            )
            if delivery_errors_path.is_file():
                try:
                    persisted_diagnostics = json.loads(
                        delivery_errors_path.read_text(encoding="utf-8")
                    )
                except (OSError, json.JSONDecodeError):
                    persisted_diagnostics = None
                if isinstance(persisted_diagnostics, dict):
                    diagnostic_payload = persisted_diagnostics
                    error_message = str(
                        persisted_diagnostics.get("error_message")
                        or error_message
                    )
        error_detail = _redacted_error_detail(
            json.dumps(
                diagnostic_payload,
                ensure_ascii=False,
                sort_keys=True,
            ),
            max_chars=None,
        )
    if error_code is None:
        try:
            descriptor = _read_derived_job_descriptor(run_dir.name)
        except (OSError, ValueError):
            descriptor = None
        decoded = (
            _decoded_worker_result_file(
                run_dir.name,
                str(descriptor["job_kind"]),
            )
            if descriptor is not None
            else None
        )
        worker_error = (
            decoded.get("error")
            if isinstance(decoded, dict) and decoded.get("ok") is False
            else None
        )
        if isinstance(worker_error, dict):
            raw_code = str(worker_error.get("type") or "derived_worker_failed")
            error_code = (
                raw_code
                if re.fullmatch(r"[A-Za-z0-9_.-]{1,100}", raw_code)
                else "derived_worker_failed"
            )
            error_message = _redacted_error_detail(
                str(worker_error.get("message") or raw_code),
                max_chars=500,
            )
            error_detail = _redacted_error_detail(
                str(worker_error.get("traceback") or ""),
                max_chars=1_200,
            )
            worker_phase = str(worker_error.get("phase") or "").strip()
            if worker_phase:
                phase = worker_phase[:100]
    if error_code is None and last_run_error is not None:
        raw_runtime_error = str(
            last_run_error.get("msg")
            or last_run_error.get("error")
            or ""
        )
        if raw_runtime_error:
            error_code = "runtime_error"
            error_message = (
                "The run stopped because of an unexpected runtime error. "
                "Open diagnostics for the captured error detail before retrying."
            )
            error_detail = _redacted_error_detail(raw_runtime_error)
    if error_code is None and last_worker_exit is not None:
        error_code = last_worker_exit["error_code"]
        error_message = last_worker_exit["error_message"]
        error_detail = last_worker_exit["error_detail"]
        phase = phase or last_worker_exit["phase"] or None

    checkpoint = _resume_checkpoint_from_disk(run_dir)
    return {
        "phase": phase,
        "error_code": error_code,
        "error_message": error_message,
        "error_detail": error_detail,
        "resume_available": checkpoint is not None,
        "resume_from_attempt": (
            checkpoint["resume_from_attempt"] if checkpoint is not None else None
        ),
        "next_attempt": checkpoint["next_attempt"] if checkpoint is not None else None,
    }


def _suggest_designer_for_retry(
    *, a_type: ArtifactType, designer_model: str, has_pdf: bool,
) -> str | None:
    """If we know a model likely to do better, return it. None means
    we have no concrete suggestion (frontend will show a generic Retry).
    The asymmetric default switches away from Kimi for paper-poster runs
    because that combination can stall on bbox geometry for long PDFs."""
    if a_type == "poster" and has_pdf and _KIMI_PATTERN in designer_model.lower():
        return _OPUS_FALLBACK
    return None


def _phase_from_disk(run_dir: Path) -> str | None:
    """Coarse "how far did we get" inference for failed runs. Walks the
    run dir top-down and reports the latest-stage marker that exists."""
    if not run_dir.exists():
        return None
    if (run_dir / "final").exists():
        return "rendering"
    if any(run_dir.glob("composites/iter_*")):
        return "rendering"
    if any(run_dir.glob("*_author/attempt_*")):
        return "authoring"
    if (run_dir / "layers").exists():
        return "ingest"
    return None


def _failure_for_no_artifact(
    *, result: RunResult, a_type: ArtifactType, designer_model: str, has_pdf: bool,
    elapsed_ms: int | None,
) -> Failure:
    """Build a Failure object when the run finished but produced no
    artifact (terminal_status reached, e.g. max_turns)."""
    run_dir = _settings_or_boot() / "runs" / result.run_id
    diagnostics = _failure_diagnostics_from_disk(run_dir)
    return Failure(
        status=str(result.terminal_status or "error"),
        phase=diagnostics["phase"] or _phase_from_disk(run_dir),
        error_code=diagnostics["error_code"],
        error_message=diagnostics["error_message"],
        error_detail=diagnostics["error_detail"],
        resume_available=diagnostics["resume_available"],
        resume_from_attempt=diagnostics["resume_from_attempt"],
        next_attempt=diagnostics["next_attempt"],
        agent_last_note=result.finalize_notes or None,
        produced_files=_list_produced_artifacts(result.run_id),
        suggested_designer=_suggest_designer_for_retry(
            a_type=a_type, designer_model=designer_model, has_pdf=has_pdf,
        ),
        suggested_planner=_suggest_designer_for_retry(
            a_type=a_type, designer_model=designer_model, has_pdf=has_pdf,
        ),
        elapsed_ms=elapsed_ms,
        critic_verdict=result.critic_verdict,
        critic_score=result.critic_score,
    )


def _should_publish_artifact(result: RunResult, artifact: Artifact | None) -> bool:
    if artifact is None:
        return False
    if result.terminal_status == "pass":
        return True
    # Degraded-but-rendered: when the agent itself graded the run as
    # fail/revise/abort/max_turns BUT the final files actually landed on
    # disk, surface the artifact anyway so the user can see what was
    # produced. The accompanying `failure` metadata stays attached to
    # the message so the frontend renders a quality-warning banner.
    if result.terminal_status in ("fail", "revise", "abort", "max_turns") and _is_rendered_artifact(artifact):
        return True
    return False


def _is_rendered_artifact(artifact: Artifact) -> bool:
    """Same recovery semantics as `_artifact_dict_has_rendered_files`,
    but works on the pydantic Artifact (used at run-completion time
    before the artifact is dumped to dict for SSE/history)."""
    a_type = artifact.artifact_type
    if a_type == "video":
        if artifact.native_format == "mp4" and artifact.native_file_url:
            return True
        if artifact.download_url and ".mp4" in artifact.download_url.lower():
            return True
        project = artifact.video_project if isinstance(artifact.video_project, dict) else {}
        latest = project.get("latest_render") if isinstance(project, dict) else None
        if isinstance(latest, dict):
            mp4_url = latest.get("mp4_url")
            if isinstance(mp4_url, str) and ".mp4" in mp4_url.lower():
                return True
    run_id = _run_id_from_artifact_id(artifact.artifact_id)
    return _final_render_complete_on_disk(run_id or "", a_type)


def _recover_video_artifact_after_exception(
    *,
    run_id: str,
    a_type: ArtifactType,
    baseline_artifact_json: str | None,
    error: str,
) -> Artifact | None:
    """Publish a completed MP4 when a later planner/LLM call crashes."""
    if a_type != "video":
        return None
    if not _final_render_complete_on_disk(run_id, "video"):
        return None
    artifact = _build_artifact_response(
        RUNS_DIR / run_id,
        run_id,
        "video",
        baseline_artifact_json=baseline_artifact_json,
    )
    if artifact is None or not _is_rendered_artifact(artifact):
        log("web.run.error_recovery_failed", run_id=run_id, error=error[:200])
        return None
    return artifact


def _failure_from_disk(
    *, run_id: str, a_type: ArtifactType, status: str,
    designer_model: str, has_pdf: bool, elapsed_ms: int | None,
) -> Failure:
    """Build a Failure when a run was cancelled or threw an exception.
    Walks the run dir for produced files."""
    run_dir = (SETTINGS.out_dir / "runs" / run_id) if SETTINGS else Path()
    produced = _list_produced_artifacts(run_id)
    diagnostics = (
        _failure_diagnostics_from_disk(run_dir)
        if run_dir.exists() and status != "cancelled"
        else {}
    )
    return Failure(
        status=status,
        phase=(
            diagnostics.get("phase")
            or (_phase_from_disk(run_dir) if run_dir.exists() else None)
        ),
        error_code=diagnostics.get("error_code"),
        error_message=diagnostics.get("error_message"),
        error_detail=diagnostics.get("error_detail"),
        resume_available=bool(diagnostics.get("resume_available")),
        resume_from_attempt=diagnostics.get("resume_from_attempt"),
        next_attempt=diagnostics.get("next_attempt"),
        agent_last_note=None,
        produced_files=produced,
        suggested_designer=_suggest_designer_for_retry(
            a_type=a_type, designer_model=designer_model, has_pdf=has_pdf,
        ) if status != "cancelled" else None,
        suggested_planner=_suggest_designer_for_retry(
            a_type=a_type, designer_model=designer_model, has_pdf=has_pdf,
        ) if status != "cancelled" else None,
        elapsed_ms=elapsed_ms,
    )


# ---------- HTML patching for /api/edits/apply ----------

# Frontend Layer fields that map cleanly to apply_edits' read surface
# (`autodesign.apply_edits._restore_text` reads these data-* attrs +
# the div's text content). Anything outside this set is silently dropped
# — the agent's edit_layer tool has the same scope (text-only) so this
# is the contract, not a v1 limitation.
from bs4 import BeautifulSoup as _BS  # noqa: E402  — local import keeps top tidy

_FLOW_TEXT_TAGS = {
    "h1", "h2", "h3", "h4", "h5", "h6",
    "p", "li", "figcaption", "blockquote", "td", "th",
}
_FLOW_TEXT_SELECTOR = ",".join([
    *sorted(_FLOW_TEXT_TAGS),
    ".identity-badge",
    ".comparison-item",
    ".formula",
    ".footer-note",
    ".lead",
    ".mechanism-side-callout",
    ".metric",
    ".muted",
    ".native-row",
    ".readout",
    ".stage",
])
_FLOW_TEXT_SCOPED_SELECTOR = ",".join(
    f".paper-poster {selector.strip()}"
    for selector in _FLOW_TEXT_SELECTOR.split(",")
)


def _is_authored_paper_poster_html(path: Path) -> bool:
    try:
        doc = _BS(path.read_text(encoding="utf-8"), "html.parser")
    except Exception:  # noqa: BLE001
        return False
    return doc.select_one(".paper-poster") is not None and doc.select_one(".canvas") is None


def _authored_paper_poster_size(path: Path) -> tuple[int, int]:
    try:
        doc = _BS(path.read_text(encoding="utf-8"), "html.parser")
    except Exception:  # noqa: BLE001
        return (3072, 1536)
    root = doc.select_one(".paper-poster")
    if root is None:
        return (3072, 1536)
    data_w = _num_attr(root, "data-w")
    data_h = _num_attr(root, "data-h")
    if data_w and data_h:
        return (data_w, data_h)
    style = _style_map(str(root.get("style") or ""))
    w = _px_int(style.get("width"))
    h = _px_int(style.get("height"))
    if w and h:
        return (w, h)
    text = path.read_text(encoding="utf-8", errors="ignore")
    css_w = _css_px_for_selector(text, ".paper-poster", "width")
    css_h = _css_px_for_selector(text, ".paper-poster", "height")
    if css_w and css_h:
        return (css_w, css_h)
    return (3072, 1536)


def _split_html_edit_payload(edits: dict[str, Any]) -> tuple[dict[str, dict[str, Any]], list[dict[str, Any]]]:
    if "layers" in edits or "layout" in edits:
        layers = edits.get("layers") if isinstance(edits.get("layers"), dict) else {}
        layout = edits.get("layout") if isinstance(edits.get("layout"), list) else []
        return (
            {str(k): v for k, v in layers.items() if isinstance(v, dict)},
            [p for p in layout if isinstance(p, dict)],
        )
    return ({str(k): v for k, v in edits.items() if isinstance(v, dict)}, [])


def _html_edit_count(edits: dict[str, Any]) -> int:
    layer_edits, layout_edits = _split_html_edit_payload(edits)
    return len(layer_edits) + len(layout_edits)


def _html_edit_manifest_ids(edits: dict[str, Any]) -> list[str]:
    layer_edits, layout_edits = _split_html_edit_payload(edits)
    ids = [str(k) for k in layer_edits]
    for idx, patch in enumerate(layout_edits, start=1):
        kind = str(patch.get("kind") or "layout")
        target = (
            patch.get("section_id")
            or patch.get("columns_id")
            or f"patch_{idx}"
        )
        ids.append(f"{kind}:{target}")
    return sorted(ids)


def _css_px_for_selector(css_text: str, selector: str, prop: str) -> int | None:
    pattern = re.compile(
        rf"{re.escape(selector)}\s*\{{[^}}]*\b{re.escape(prop)}\s*:\s*([0-9.]+)px",
        re.IGNORECASE | re.DOTALL,
    )
    m = pattern.search(css_text)
    if not m:
        return None
    try:
        return max(1, int(float(m.group(1))))
    except (TypeError, ValueError):
        return None


def _px_int(raw: str | None) -> int | None:
    if not raw:
        return None
    clean = str(raw).strip().lower().removesuffix("px").strip()
    try:
        value = int(float(clean))
    except (TypeError, ValueError):
        return None
    return value if value > 0 else None


def _num_attr(node: Any, attr: str) -> int | None:
    raw = node.get(attr)
    if raw is None:
        return None
    try:
        value = int(float(str(raw)))
    except (TypeError, ValueError):
        return None
    return value if value > 0 else None


def _validate_required_poster_palette_html(
    html_path: Path,
    required_color_system: dict[str, Any],
) -> None:
    from autodesign.tools.propose_paper_poster_html import (
        authored_palette_diagnostics,
        required_palette_diagnostic_is_blocking,
    )

    try:
        html_text = html_path.read_text(encoding="utf-8")
    except OSError as exc:
        raise HTTPException(
            422,
            detail={
                "code": "poster_palette_validation_failed",
                "message": f"Unable to validate edited Poster HTML: {exc}",
            },
        ) from exc
    diagnostics = authored_palette_diagnostics(
        html_text,
        "",
        required_color_system,
        require_selected=True,
    )
    blocking = [
        diagnostic
        for diagnostic in diagnostics
        if required_palette_diagnostic_is_blocking(diagnostic)
    ]
    if blocking:
        raise HTTPException(
            422,
            detail={
                "code": "poster_palette_validation_failed",
                "message": "Edited Poster HTML does not use the required palette.",
                "palette_diagnostics": blocking,
            },
        )


def _persist_apply_edits_palette_manifest(
    run_dir: Path,
    palette_id: str,
    required_color_system: dict[str, Any],
) -> None:
    final_dir = run_dir / "final"
    pending_path = run_dir / "authored_poster_edit_manifest.pending.json"
    payload = _read_json_file(pending_path)
    if isinstance(payload, dict):
        manifest_path = final_dir / "authored_poster_edit_manifest.json"
    else:
        manifest_path = final_dir / "apply_edits_palette_manifest.json"
        payload = {
            "artifact_type": "poster",
            "render_mode": "apply_edits",
        }
    payload["palette_id"] = palette_id
    payload["required_color_system"] = required_color_system
    atomic_write_json(manifest_path, payload)
    pending_path.unlink(missing_ok=True)


def _apply_authored_paper_poster_edits(
    src_path: Path,
    staged_html: Path,
    settings: Settings,
    parent_run_id: str,
    edits: dict[str, Any],
    *,
    required_color_system: dict[str, Any],
) -> ApplyEditsResult:
    """Compatibility helper; endpoint execution is owned by artifact_edit worker."""

    new_id = new_run_id()
    run_dir = settings.out_dir / "runs" / new_id
    final_dir = run_dir / "final"
    final_dir.mkdir(parents=True, exist_ok=True)
    for child in src_path.parent.iterdir():
        if child.is_dir():
            shutil.copytree(child, final_dir / child.name, dirs_exist_ok=True)
    html_path = final_dir / "poster.html"
    shutil.copy2(staged_html, html_path)
    ensure_poster_katex_document(
        html_path,
        Path(getattr(settings, "repo_root", _REPO_ROOT)),
        root_selector=".paper-poster",
    )
    width, height = _authored_paper_poster_size(html_path)
    preview_path = final_dir / "preview.png"
    browser_result = screenshot_html(
        html_path,
        preview_path,
        viewport_width=width,
        viewport_height=height,
        selector=".paper-poster",
        max_edge=settings.poster_preview_max_edge,
        timeout_ms=30_000,
    )
    if browser_result.warnings and (src_path.parent / "preview.png").is_file():
        shutil.copy2(src_path.parent / "preview.png", preview_path)
    restored = _html_edit_manifest_ids(edits)
    atomic_write_json(run_dir / "authored_poster_edit_manifest.pending.json", {
        "artifact_type": "poster",
        "render_mode": "authored_paper_poster_edit",
        "parent_run_id": parent_run_id,
        "palette_id": str(required_color_system.get("palette_id") or ""),
        "required_color_system": required_color_system,
        "edits": restored,
        "canvas": {"w_px": width, "h_px": height},
        "preview": {
            "backend": browser_result.backend,
            "warnings": browser_result.warnings,
            "scale": browser_result.scale,
            "width_px": browser_result.width_px,
            "height_px": browser_result.height_px,
        },
    })
    return ApplyEditsResult(
        run_id=new_id,
        run_dir=str(run_dir),
        parent_run_id=parent_run_id,
        restored_layer_ids=restored,
        skipped=[],
        artifact_type="poster",
    )


def _apply_authored_html_edits(*_args: Any, **_kwargs: Any) -> ApplyEditsResult:
    """Legacy test seam; authored HTML edits now execute only in the worker."""

    raise RuntimeError("authored HTML edits require the supervised artifact_edit worker")


def _patch_html_for_apply_edits(
    src: Path,
    dst: Path,
    edits: dict[str, Any],
    *,
    source_run_id: str | None = None,
) -> None:
    """Read `src`, apply each edit's partial fields to the matching
    `.layer[data-layer-id="…"]`, write to `dst`.

    Text layers patch content/style attrs. Image layers patch geometry and
    replacement sources so real paper figures edited in the web UI survive
    the apply-edits round trip.
    """
    doc = _BS(src.read_text(encoding="utf-8"), "html.parser")
    layer_edits, layout_edits = _split_html_edit_payload(edits)
    for layer_id, patch in layer_edits.items():
        if not isinstance(patch, dict):
            continue
        div = _find_editable_html_node(doc, layer_id, patch)
        if div is None:
            log("web.edits.miss", layer_id=layer_id)
            continue
        kind = div.get("data-kind") or _infer_editable_html_kind(div)
        if kind == "text":
            _patch_text_layer(div, patch)
        elif kind in {"image", "brand_asset", "background"}:
            _patch_image_layer(
                div,
                patch,
                source_run_id=source_run_id,
            )
    for patch in layout_edits:
        _patch_flow_layout(doc, patch)
    dst.write_text(str(doc), encoding="utf-8")


def _find_editable_html_node(doc: Any, layer_id: str, patch: dict[str, Any]) -> Any | None:
    if doc.select_one(".paper-poster") is not None:
        found = _find_authored_poster_edit_node(doc, layer_id, patch)
        if found is not None:
            return found
    found = doc.find(attrs={"data-layer-id": layer_id})
    if found is not None:
        return found
    return _find_authored_poster_edit_node(doc, layer_id, patch)


def _find_authored_poster_edit_node(doc: Any, layer_id: str, patch: dict[str, Any]) -> Any | None:
    wants_image = (
        isinstance(patch.get("src"), str)
        or patch.get("fit") in {"cover", "contain", "fill"}
        or isinstance(patch.get("object_position"), dict)
    )
    if wants_image:
        for node in doc.select(".paper-poster img"):
            if _node_layer_id(node, "image") == layer_id:
                return node
    wants_text = any(
        key in patch
        for key in (
            "text",
            "font_family",
            "font_size_px",
            "font_weight",
            "font_style",
            "line_height",
            "letter_spacing",
            "text_transform",
            "align",
            "effects",
        )
    )
    if wants_text:
        for node in doc.select(_FLOW_TEXT_SCOPED_SELECTOR):
            if _node_layer_id(node, "text") == layer_id:
                return node

    if isinstance(patch.get("bbox"), dict) or isinstance(patch.get("flow_offset"), dict):
        for node in doc.select(_FLOW_TEXT_SCOPED_SELECTOR):
            if _node_layer_id(node, "text") == layer_id:
                return node
        for node in doc.select(".paper-poster img"):
            if _node_layer_id(node, "image") == layer_id:
                return node

    for node in doc.select(
        ".paper-poster .poster-header[data-block-id], "
        ".paper-poster .poster-section[data-block-id]",
    ):
        if _node_layer_id(node, "section") == layer_id:
            return node
    for node in doc.select(_FLOW_TEXT_SCOPED_SELECTOR):
        if _node_layer_id(node, "text") == layer_id:
            return node
    for node in doc.select(".paper-poster img"):
        if _node_layer_id(node, "image") == layer_id:
            return node

    for attr in ("data-layer-id", "data-block-id", "id"):
        found = doc.find(attrs={attr: layer_id})
        if found is not None:
            return found
    return None


def _node_layer_id(node: Any, kind: str) -> str:
    raw = str(node.get("data-layer-id") or "").strip()
    if raw:
        return raw
    if kind == "section":
        block_id = str(node.get("data-block-id") or "section").strip()
        idx = _poster_node_index(node, (
            ".paper-poster .poster-header[data-block-id], "
            ".paper-poster .poster-section[data-block-id]"
        ))
        return f"flow_section_{_slug_id(block_id)}_{idx}"
    if kind == "image":
        base = (
            str(node.get("data-source-id") or "").strip()
            or str(node.get("alt") or "").strip()
            or _nearest_block_id(node)
            or "image"
        )
        return f"flow_image_{_slug_id(base)}_{_poster_node_index(node, '.paper-poster img')}"
    block_id = str(node.get("data-block-id") or "").strip()
    if block_id:
        return block_id
    base = _nearest_block_id(node) or "text"
    return f"flow_text_{_slug_id(base)}_{_poster_node_index(node, _FLOW_TEXT_SCOPED_SELECTOR)}"


def _poster_node_index(node: Any, selector: str) -> int:
    root = node
    while getattr(root, "parent", None) is not None:
        root = root.parent
    for idx, candidate in enumerate(root.select(selector), start=1):
        if candidate is node:
            return idx
    return 1


def _nearest_block_id(node: Any) -> str | None:
    cur = node
    while cur is not None and getattr(cur, "name", None):
        raw = cur.get("data-block-id") or cur.get("data-column-id")
        if raw:
            return str(raw)
        cur = cur.parent
    return None


def _slug_id(raw: str) -> str:
    slug = re.sub(r"[^a-z0-9]+", "_", str(raw).lower()).strip("_")[:72]
    return slug or "item"


def _patch_flow_layout(doc: Any, patch: dict[str, Any]) -> None:
    kind = str(patch.get("kind") or "")
    if kind == "dom_delete":
        target = _find_poster_dom_delete_target(doc, patch)
        target_label = (
            str(patch.get("target_id") or "").strip()
            or str(patch.get("block_id") or "").strip()
            or str(patch.get("selector") or "").strip()
        )
        if target is None:
            log("web.edits.layout_miss", kind=kind, target=target_label or "dom")
            return
        if _is_protected_poster_delete_target(target):
            log("web.edits.layout_skip", kind=kind, target=target_label or getattr(target, "name", "dom"))
            return
        target.extract()
        return

    if kind == "section_height":
        section_id = str(patch.get("section_id") or "").strip()
        height_px = _positive_int(patch.get("height_px"), minimum=80)
        if not section_id or height_px is None:
            return
        section = _find_poster_section(doc, section_id)
        if section is None:
            log("web.edits.layout_miss", kind=kind, target=section_id)
            return
        style = _style_map(section.get("style") or "")
        style["height"] = f"{height_px}px"
        style["min-height"] = f"{height_px}px"
        section["style"] = _style_text(style)
        section["data-layout-height-px"] = str(height_px)
        return

    if kind == "section_size":
        section_id = str(patch.get("section_id") or "").strip()
        width_px = _positive_int(patch.get("width_px"), minimum=120)
        height_px = _positive_int(patch.get("height_px"), minimum=80)
        offset_x_px = _int_or_none(patch.get("offset_x_px"))
        offset_y_px = _int_or_none(patch.get("offset_y_px"))
        if not section_id or (width_px is None and height_px is None):
            return
        section = _find_poster_section(doc, section_id)
        if section is None:
            log("web.edits.layout_miss", kind=kind, target=section_id)
            return
        style = _style_map(section.get("style") or "")
        style["box-sizing"] = "border-box"
        if width_px is not None:
            style["width"] = f"{width_px}px"
            section["data-layout-width-px"] = str(width_px)
        if height_px is not None:
            style["height"] = f"{height_px}px"
            style["min-height"] = f"{height_px}px"
            section["data-layout-height-px"] = str(height_px)
        if offset_x_px is not None or offset_y_px is not None:
            style["position"] = "relative"
            if offset_x_px is not None:
                style["left"] = f"{offset_x_px}px"
                section["data-layout-offset-x-px"] = str(offset_x_px)
            if offset_y_px is not None:
                style["top"] = f"{offset_y_px}px"
                section["data-layout-offset-y-px"] = str(offset_y_px)
        section["style"] = _style_text(style)
        return

    if kind == "poster_style":
        scope = str(patch.get("scope") or "global")
        styles = patch.get("styles")
        if not isinstance(styles, dict):
            return
        clean = {
            key: color
            for key in ("accent", "accent2", "background", "ink")
            if (color := _clean_hex_color(styles.get(key))) is not None
        }
        if not clean:
            return
        if scope == "section":
            section_id = str(patch.get("section_id") or "").strip()
            section = _find_poster_section(doc, section_id)
            if section is None:
                log("web.edits.layout_miss", kind=kind, target=section_id)
                return
            _apply_section_style_patch(section, clean)
            return
        _apply_global_poster_style_patch(doc, clean)
        return

    if kind == "column_widths":
        columns_id = str(patch.get("columns_id") or "").strip()
        raw_widths = patch.get("widths")
        if not isinstance(raw_widths, list) or len(raw_widths) != 3:
            return
        widths = [_positive_float(v) for v in raw_widths]
        if any(v is None for v in widths):
            return
        clamped = _normalize_column_widths([float(v) for v in widths if v is not None])
        columns = _find_poster_columns(doc, columns_id)
        if columns is None:
            log("web.edits.layout_miss", kind=kind, target=columns_id or "poster-columns")
            return
        style = _style_map(columns.get("style") or "")
        style["display"] = "grid"
        style["grid-template-columns"] = " ".join(f"{_css_num(v)}%" for v in clamped)
        columns["style"] = _style_text(style)
        columns["data-layout-column-widths"] = ",".join(_css_num(v) for v in clamped)
        return

    if kind == "section_order":
        raw_columns = patch.get("columns")
        if not isinstance(raw_columns, list):
            return
        for idx, column_patch in enumerate(raw_columns):
            if not isinstance(column_patch, dict):
                continue
            column_id = str(column_patch.get("column_id") or "").strip()
            section_ids = column_patch.get("section_ids")
            if not isinstance(section_ids, list):
                continue
            column = _find_poster_column(doc, column_id, idx)
            if column is None:
                log("web.edits.layout_miss", kind=kind, target=column_id or f"column_{idx + 1}")
                continue
            for section_id in [str(v).strip() for v in section_ids if str(v).strip()]:
                section = _find_poster_section(doc, section_id)
                if section is not None:
                    column.append(section.extract())


def _positive_int(raw: Any, *, minimum: int = 1) -> int | None:
    try:
        value = int(float(raw))
    except (TypeError, ValueError):
        return None
    return value if value >= minimum else minimum


def _int_or_none(raw: Any) -> int | None:
    try:
        return int(float(raw))
    except (TypeError, ValueError):
        return None


def _positive_float(raw: Any) -> float | None:
    try:
        value = float(raw)
    except (TypeError, ValueError):
        return None
    return value if value > 0 else None


def _normalize_column_widths(widths: list[float]) -> list[float]:
    min_pct = 22.0
    clamped = [max(min_pct, v) for v in widths[:3]]
    total = sum(clamped) or 100.0
    if abs(total - 100.0) < 0.01:
        return [round(v, 2) for v in clamped]
    if total > 100.0:
        excess = total - 100.0
        adjustable = sum(max(0.0, v - min_pct) for v in clamped)
        if adjustable > 0:
            clamped = [
                max(min_pct, v - excess * (max(0.0, v - min_pct) / adjustable))
                for v in clamped
            ]
    else:
        deficit = 100.0 - total
        clamped = [v + deficit / len(clamped) for v in clamped]
    rounded = [round(v, 2) for v in clamped]
    rounded[-1] = round(100.0 - sum(rounded[:-1]), 2)
    return rounded


def _find_poster_columns(doc: Any, columns_id: str) -> Any | None:
    if columns_id:
        for attr in ("data-block-id", "data-layout-region", "id"):
            found = doc.find(attrs={attr: columns_id})
            if found is not None and _has_class(found, "poster-columns"):
                return found
    return doc.select_one(".paper-poster .poster-columns")


def _find_poster_column(doc: Any, column_id: str, idx: int) -> Any | None:
    if column_id:
        for attr in ("data-column-id", "data-block-id", "id"):
            found = doc.find(attrs={attr: column_id})
            if found is not None and (_has_class(found, "poster-column") or found.get("data-column-id")):
                return found
    columns = doc.select(".paper-poster .poster-column")
    if 0 <= idx < len(columns):
        return columns[idx]
    return None


def _find_poster_section(doc: Any, section_id: str) -> Any | None:
    for attr in ("data-block-id", "data-layer-id", "id"):
        found = doc.find(attrs={attr: section_id})
        if found is not None and _has_class(found, "poster-section"):
            return found
    return None


def _find_poster_dom_delete_target(doc: Any, patch: dict[str, Any]) -> Any | None:
    selector = str(patch.get("selector") or "").strip()
    if selector:
        try:
            found = doc.select_one(selector)
        except Exception:  # noqa: BLE001 - invalid CSS selector from browser state
            found = None
        if found is not None and found.select_one(".paper-poster") is None:
            return found

    target_id = str(patch.get("target_id") or "").strip()
    block_id = str(patch.get("block_id") or "").strip()
    target_kind = str(patch.get("target_kind") or "").strip()
    for raw_id in [block_id, target_id]:
        if not raw_id:
            continue
        for attr in ("data-layer-id", "data-block-id", "id"):
            found = doc.find(attrs={attr: raw_id})
            if found is not None:
                return found

    if not target_id:
        return None
    if target_kind == "image":
        for node in doc.select(".paper-poster img"):
            if _node_layer_id(node, "image") == target_id:
                return node
    if target_kind == "section":
        for node in doc.select(
            ".paper-poster .poster-header[data-block-id], "
            ".paper-poster .poster-section[data-block-id]"
        ):
            if _node_layer_id(node, "section") == target_id:
                return node
    if target_kind == "text":
        for node in doc.select(_FLOW_TEXT_SCOPED_SELECTOR):
            if _node_layer_id(node, "text") == target_id:
                return node
    for node in doc.select(".paper-poster [data-layer-id], .paper-poster [data-block-id], .paper-poster img"):
        inferred = _infer_editable_html_kind(node)
        if _node_layer_id(node, inferred if inferred in {"text", "image", "section"} else "text") == target_id:
            return node
    return None


def _is_protected_poster_delete_target(node: Any) -> bool:
    if node is None:
        return True
    if (
        _has_class(node, "paper-poster")
        or _has_class(node, "poster-columns")
        or _has_class(node, "poster-column")
        or _has_class(node, "deck-slide")
        or node.get("data-autodesign-artifact-root")
    ):
        return True
    cur = node
    inside_editable_root = False
    while cur is not None and getattr(cur, "name", None):
        if _has_class(cur, "paper-poster") or cur.get("data-autodesign-artifact-root"):
            inside_editable_root = True
            break
        cur = cur.parent
    return not inside_editable_root


def _clean_hex_color(raw: Any) -> str | None:
    text = str(raw or "").strip()
    if re.fullmatch(r"#[0-9a-fA-F]{6}", text):
        return text.lower()
    return None


def _apply_section_style_patch(section: Any, styles: dict[str, str]) -> None:
    section_style = _style_map(section.get("style") or "")
    if background := styles.get("background"):
        section_style["background"] = background
    if ink := styles.get("ink"):
        section_style["color"] = ink
    if accent := styles.get("accent"):
        section_style["border-color"] = accent
    section["style"] = _style_text(section_style)
    section["data-style-tweaked"] = "section"

    title = section.select_one(".section-title, .section-heading, h2, h3")
    if title is not None:
        title_style = _style_map(title.get("style") or "")
        if accent := styles.get("accent"):
            title_style["background"] = accent
            title_style["border-color"] = accent
        if accent2 := styles.get("accent2"):
            title_style["box-shadow"] = f"inset 6px 0 0 {accent2}"
        title_style["color"] = "#ffffff"
        title["style"] = _style_text(title_style)


def _apply_global_poster_style_patch(doc: Any, styles: dict[str, str]) -> None:
    root = doc.select_one(".paper-poster")
    if root is None:
        return
    root_style = _style_map(root.get("style") or "")
    if accent := styles.get("accent"):
        root_style["--da-accent"] = accent
    if accent2 := styles.get("accent2"):
        root_style["--da-accent2"] = accent2
    if background := styles.get("background"):
        root_style["--da-background"] = background
        root_style["background"] = background
    if ink := styles.get("ink"):
        root_style["--da-ink"] = ink
        root_style["color"] = ink
    root["style"] = _style_text(root_style)
    root["data-style-tweaked"] = "global"

    style = doc.find(id="autodesign-style-tweaks") or doc.find(id="designanything-style-tweaks")
    if style is None:
        style = doc.new_tag("style", id="autodesign-style-tweaks")
        head = doc.head or doc.new_tag("head")
        if doc.head is None:
            html = doc.html or doc
            html.insert(0, head)
        head.append(style)
    else:
        style["id"] = "autodesign-style-tweaks"
    style.string = _poster_style_css(styles)


def _poster_style_css(styles: dict[str, str]) -> str:
    accent = styles.get("accent", "var(--da-accent)")
    accent2 = styles.get("accent2", "var(--da-accent2)")
    background = styles.get("background", "var(--da-background)")
    ink = styles.get("ink", "var(--da-ink)")
    return f"""
.paper-poster {{
  --da-accent: {accent};
  --da-accent2: {accent2};
  --da-background: {background};
  --da-ink: {ink};
}}
.paper-poster {{
  background: var(--da-background) !important;
  color: var(--da-ink) !important;
}}
.paper-poster .poster-header,
.paper-poster .section-title,
.paper-poster .section-heading {{
  background: var(--da-accent) !important;
  border-color: var(--da-accent) !important;
}}
.paper-poster .metric,
.paper-poster .stat,
.paper-poster .highlight,
.paper-poster .formula,
.paper-poster .readout,
.paper-poster .native-row strong {{
  color: var(--da-accent2) !important;
}}
.paper-poster table th,
.paper-poster .comparison-item,
.paper-poster .identity-badge {{
  border-color: var(--da-accent2) !important;
}}
""".strip()


def _has_class(node: Any, class_name: str) -> bool:
    classes = node.get("class") if node is not None else None
    if isinstance(classes, str):
        return class_name in classes.split()
    if isinstance(classes, list):
        return class_name in classes
    return False


def _infer_editable_html_kind(node: Any) -> str:
    tag = str(getattr(node, "name", "") or "").lower()
    if tag == "img":
        return "image"
    if tag in _FLOW_TEXT_TAGS or _is_flow_text_unit(node):
        return "text"
    if _has_class(node, "poster-section") or _has_class(node, "poster-header"):
        return "section"
    if node.get("contenteditable") == "true":
        return "text"
    return str(node.get("data-block-kind") or "")


def _is_flow_text_unit(node: Any) -> bool:
    return any(
        _has_class(node, class_name)
        for class_name in (
            "identity-badge",
            "comparison-item",
            "formula",
            "footer-note",
            "lead",
            "mechanism-side-callout",
            "metric",
            "muted",
            "native-row",
            "readout",
            "stage",
        )
    )


def _patch_text_layer(div: Any, patch: dict[str, Any]) -> None:
    # apply_edits._restore_text strips ld-drag-handle spans before
    # reading text via div.get_text(). We replace the div's contents
    # entirely when the user edited `text`, otherwise leave the original
    # text + handle in place.
    if "text" in patch and isinstance(patch["text"], str):
        div.clear()
        div.append(patch["text"])

    # Scalar attrs map 1:1.
    if "font_size_px" in patch and patch["font_size_px"] is not None:
        try:
            font_size = int(float(patch["font_size_px"]))
            div["data-font-size-px"] = str(font_size)
            _patch_style_decl(div, "font-size", f"{font_size}px")
        except (TypeError, ValueError):
            pass
    if isinstance(patch.get("font_family"), str):
        div["data-font-family"] = patch["font_family"]
        _patch_style_decl(div, "font-family", f"'{patch['font_family']}'")
    if "font_weight" in patch and patch["font_weight"] is not None:
        try:
            weight = int(float(patch["font_weight"]))
            weight = max(100, min(900, weight))
            div["data-font-weight"] = str(weight)
            _patch_style_decl(div, "font-weight", str(weight))
        except (TypeError, ValueError):
            pass
    if patch.get("font_style") in {"normal", "italic"}:
        div["data-font-style"] = patch["font_style"]
        _patch_style_decl(div, "font-style", patch["font_style"])
    if "line_height" in patch and patch["line_height"] is not None:
        try:
            line_height = float(patch["line_height"])
            div["data-line-height"] = _css_num(line_height)
            _patch_style_decl(div, "line-height", _css_num(line_height))
        except (TypeError, ValueError):
            pass
    if "letter_spacing" in patch and patch["letter_spacing"] is not None:
        try:
            letter_spacing = float(patch["letter_spacing"])
            div["data-letter-spacing"] = _css_num(letter_spacing)
            _patch_style_decl(div, "letter-spacing", f"{_css_num(letter_spacing)}px")
        except (TypeError, ValueError):
            pass
    if patch.get("text_transform") in {"none", "uppercase"}:
        div["data-text-transform"] = patch["text_transform"]
        _patch_style_decl(div, "text-transform", patch["text_transform"])
    if isinstance(patch.get("align"), str):
        div["data-align"] = patch["align"]
        _patch_style_decl(div, "text-align", patch["align"])
    if "z_index" in patch and patch["z_index"] is not None:
        try:
            div["data-z-index"] = str(int(float(patch["z_index"])))
        except (TypeError, ValueError):
            pass

    effects = patch.get("effects")
    if isinstance(effects, dict) and isinstance(effects.get("fill"), str):
        div["data-fill"] = effects["fill"]
        _patch_style_decl(div, "color", effects["fill"])

    bbox = patch.get("bbox")
    flow_offset = patch.get("flow_offset")
    if isinstance(bbox, dict) and isinstance(flow_offset, dict):
        _patch_flow_layer_bbox(div, bbox, flow_offset)
    elif isinstance(bbox, dict):
        for k, attr in (("x", "data-bbox-x"), ("y", "data-bbox-y"),
                         ("w", "data-bbox-w"), ("h", "data-bbox-h")):
            v = bbox.get(k)
            if v is None:
                continue
            try:
                div[attr] = str(int(float(v)))
            except (TypeError, ValueError):
                pass


def _patch_style_decl(div: Any, key: str, value: str) -> None:
    style: dict[str, str] = {}
    for chunk in str(div.get("style") or "").split(";"):
        if ":" not in chunk:
            continue
        k, v = chunk.split(":", 1)
        k = k.strip()
        if k:
            style[k.lower()] = v.strip()
    style[key] = value
    div["style"] = ";".join(f"{k}:{v}" for k, v in style.items())


def _patch_image_layer(
    div: Any,
    patch: dict[str, Any],
    *,
    source_run_id: str | None,
) -> None:
    if "z_index" in patch and patch["z_index"] is not None:
        try:
            div["data-z-index"] = str(int(float(patch["z_index"])))
        except (TypeError, ValueError):
            pass

    bbox = patch.get("bbox")
    flow_offset = patch.get("flow_offset")
    if isinstance(bbox, dict) and isinstance(flow_offset, dict):
        _patch_flow_layer_bbox(div, bbox, flow_offset)
    elif isinstance(bbox, dict):
        _patch_layer_bbox(div, bbox)

    src = patch.get("src")
    if isinstance(src, str) and src:
        img = div if str(getattr(div, "name", "") or "").lower() == "img" else div.find("img")
        if img is not None:
            img["src"] = _image_src_for_apply_edits(
                src,
                source_run_id=source_run_id,
            )

    image = div if str(getattr(div, "name", "") or "").lower() == "img" else div.find("img")
    if image is None:
        return
    fit = patch.get("fit")
    if fit in {"cover", "contain", "fill"}:
        _patch_style_decl(image, "object-fit", fit)
    object_position = patch.get("object_position")
    if isinstance(object_position, dict):
        try:
            x = max(0.0, min(1.0, float(object_position.get("x"))))
            y = max(0.0, min(1.0, float(object_position.get("y"))))
            _patch_style_decl(
                image,
                "object-position",
                f"{_css_num(x * 100)}% {_css_num(y * 100)}%",
            )
        except (TypeError, ValueError):
            pass


def _patch_flow_layer_bbox(node: Any, bbox: dict[str, Any], flow_offset: dict[str, Any]) -> None:
    values: dict[str, int] = {}
    for key in ("w", "h"):
        raw = bbox.get(key)
        if raw is None:
            continue
        try:
            values[key] = max(1, int(float(raw)))
        except (TypeError, ValueError):
            continue
    offsets: dict[str, int] = {}
    for key in ("dx", "dy"):
        raw = flow_offset.get(key)
        if raw is None:
            continue
        try:
            offsets[key] = int(float(raw))
        except (TypeError, ValueError):
            continue

    style = _style_map(node.get("style") or "")
    style["position"] = "relative"
    if "dx" in offsets:
        style["left"] = f"{offsets['dx']}px"
        node["data-flow-offset-x"] = str(offsets["dx"])
    if "dy" in offsets:
        style["top"] = f"{offsets['dy']}px"
        node["data-flow-offset-y"] = str(offsets["dy"])
    if "w" in values:
        style["width"] = f"{values['w']}px"
    if "h" in values:
        style["height"] = f"{values['h']}px"
    node["style"] = _style_text(style)


def _patch_layer_bbox(div: Any, bbox: dict[str, Any]) -> None:
    values: dict[str, int] = {}
    for key in ("x", "y", "w", "h"):
        raw = bbox.get(key)
        if raw is None:
            continue
        try:
            values[key] = int(float(raw))
        except (TypeError, ValueError):
            continue
    if not values:
        return

    attr_map = {
        "x": "data-bbox-x",
        "y": "data-bbox-y",
        "w": "data-bbox-w",
        "h": "data-bbox-h",
    }
    for key, attr in attr_map.items():
        if key in values:
            div[attr] = str(values[key])
    if "x" in values:
        div["data-bbox-tx"] = str(values["x"])
    if "y" in values:
        div["data-bbox-ty"] = str(values["y"])

    style = _style_map(div.get("style") or "")
    if "x" in values:
        style["left"] = f"{values['x']}px"
    if "y" in values:
        style["top"] = f"{values['y']}px"
    if "w" in values:
        style["width"] = f"{max(1, values['w'])}px"
    if "h" in values:
        style["height"] = f"{max(1, values['h'])}px"
    div["style"] = _style_text(style)


def _image_src_for_apply_edits(
    src: str,
    *,
    source_run_id: str | None,
) -> str:
    if src.startswith("data:image/"):
        return src
    clean = src.split("?", 1)[0]
    run_prefix = "/api/files/runs/"
    if clean.startswith(run_prefix):
        if not source_run_id:
            raise ValueError("run image source requires explicit source-run authority")
        relative = clean[len(run_prefix):].lstrip("/")
        try:
            opened = open_run_file(
                RUNS_DIR,
                relative,
                expected_run_id=source_run_id,
            )
        except RunFileAccessError as exc:
            raise ValueError("run image source is unavailable") from exc
        with opened:
            mime = _mime_for_image_path(opened.path)
            b64 = base64.b64encode(opened.handle.read()).decode("ascii")
        return f"data:{mime};base64,{b64}"
    editor_prefix = "/api/files/editor-assets/"
    if clean.startswith(editor_prefix):
        relative = clean[len(editor_prefix):].lstrip("/")
        source = (EDITOR_ASSETS_DIR / relative).resolve()
        if not _path_inside(source, EDITOR_ASSETS_DIR.resolve()) or not source.is_file():
            return src
        mime = _mime_for_image_path(source)
        b64 = base64.b64encode(source.read_bytes()).decode("ascii")
        return f"data:{mime};base64,{b64}"
    return src


def _mime_for_image_path(path: Path) -> str:
    suffix = path.suffix.lower()
    return {
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".webp": "image/webp",
        ".gif": "image/gif",
        ".png": "image/png",
    }.get(suffix, "image/png")


def _style_map(raw_style: str) -> dict[str, str]:
    out: dict[str, str] = {}
    for part in str(raw_style).split(";"):
        key, sep, value = part.partition(":")
        if not sep:
            continue
        out[key.strip().lower()] = value.strip()
    return out


def _style_text(style: dict[str, str]) -> str:
    return "; ".join(f"{key}:{value}" for key, value in style.items() if value)


# ---------- Local built frontend ----------
#
# Hosted deployments can serve web/dist with Caddy or another static server. For
# local general users, serving the built SPA from FastAPI keeps startup to a
# single backend process and removes Node/npm from the runtime path.
WEB_DIST_DIR = Path.cwd() / "web" / "dist"
WEB_DIST_INDEX = WEB_DIST_DIR / "index.html"


@app.get("/{spa_path:path}", include_in_schema=False)
async def local_web_app(spa_path: str) -> FileResponse:
    if not WEB_DIST_INDEX.exists():
        raise HTTPException(
            404,
            detail=(
                "Built web UI not found. Developers can run `cd web && npm run dev`, "
                "or build it with `cd web && npm ci && npm run build`."
            ),
        )

    clean = unquote(spa_path or "").lstrip("/")
    if clean:
        candidate = (WEB_DIST_DIR / clean).resolve()
        root = WEB_DIST_DIR.resolve()
        if _path_inside(candidate, root) and candidate.is_file():
            return FileResponse(candidate)
    return FileResponse(WEB_DIST_INDEX)
